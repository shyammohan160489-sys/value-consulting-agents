"""
Nordic FinTech Forum — Helsinki, May 2026
AI Roundtable Opening — PPTX BUILDER (Google Slides compatible)

Builds nordic_fintech_forum_opening_pov.pptx with:
- 13 slides matching the v2 visual design (light typography, color palette,
  L-motif corner mark, Backbase wordmark + page number footer)
- 16:9 canvas at 20" × 11.25" (Google Slides native — matches v2 HTML)
- Speaker notes verbatim from SPEAKER_NOTES (the WHY / SAY / PAUSE / DROP-INS
  / NEXT format Barry uses in presenter mode)

Run:  python3 build_pptx.py
Output: nordic_fintech_forum_opening_pov.pptx (in same folder)
"""

import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent

# Pull speaker notes from build_deck.py (single source of truth)
sys.path.insert(0, str(HERE))
from build_deck import SPEAKER_NOTES  # noqa: E402

# ──────────────────────────────────────────────────────────────────────────
# Design tokens (verbatim from the v2 HTML engine)
# ──────────────────────────────────────────────────────────────────────────

NAVY = RGBColor(0x04, 0x13, 0x26)        # #041326
BLUE = RGBColor(0x33, 0x66, 0xFF)        # #3366FF
RED = RGBColor(0xFF, 0x50, 0x3C)         # #FF503C
LIGHT_BLUE = RGBColor(0xE5, 0xEB, 0xFF)  # #E5EBFF
OFF_WHITE = RGBColor(0xF3, 0xF6, 0xF9)   # #F3F6F9
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MAIN = NAVY
TEXT_MUTED = RGBColor(0x5C, 0x6E, 0x84)  # #5C6E84
LINE_GRAY = RGBColor(0xCE, 0xD2, 0xD7)   # #CED2D7
RED_SOFT = RGBColor(0xFF, 0xE6, 0xE2)    # #FFE6E2
NAVY_SOFT = RGBColor(0xE5, 0xEA, 0xF0)   # #E5EAF0

FONT = "Libre Franklin"   # Fallback to Helvetica/Arial via PPTX system substitution
MONO = "Consolas"

# Canvas: 20" × 11.25" (Google Slides native widescreen, matches v2 HTML)
CANVAS_W = Inches(20)
CANVAS_H = Inches(11.25)

# ──────────────────────────────────────────────────────────────────────────
# Helpers
# ──────────────────────────────────────────────────────────────────────────

def set_slide_background(slide, rgb):
    """Set slide background fill colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_text(slide, x, y, w, h, text, *, font=FONT, size=18, bold=False,
             color=TEXT_MAIN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             italic=False, letter_spacing=None, line_spacing=1.2):
    """Add a textbox with given styling. x/y/w/h in Inches."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None  # Disable autofit — required for Google Slides import
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    # Split on \n for multi-line
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def add_rich_text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=1.2):
    """Add a textbox where each 'run' is a dict with {text, ...font_props}.
    Used for inline highlighting (mixed colours/weights in one line).
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for run in runs:
        r = p.add_run()
        r.text = run["text"]
        f = r.font
        f.name = run.get("font", FONT)
        f.size = Pt(run.get("size", 18))
        f.bold = run.get("bold", False)
        f.italic = run.get("italic", False)
        f.color.rgb = run.get("color", TEXT_MAIN)
    return tb


def add_box(slide, x, y, w, h, *, fill=None, line_color=None, line_width=0.75,
            corner_radius=None):
    """Add a rectangle (or rounded rectangle) shape. Returns the shape."""
    use_rounded = bool(corner_radius)
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if use_rounded else MSO_SHAPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, x, y, w, h)
    if use_rounded:
        # Adjustment factor on rounded rect — fraction of half-shortest-side
        box.adjustments[0] = corner_radius
    box.shadow.inherit = False
    if fill is None:
        box.fill.background()
    else:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    if line_color is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line_color
        box.line.width = Pt(line_width)
    # Remove default text frame margins
    tf = box.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    return box


def add_pill(slide, x, y, text, *, bg, fg, size=9):
    """Small uppercase pill label (used at top of cards, etc.).
    We approximate width generously so bold uppercase doesn't wrap.
    """
    # ~0.012 inches per character at 9pt — empirically tuned for bold uppercase.
    char_w_in = size * 0.012
    pad_in = 0.35
    w = Inches(char_w_in * len(text) + pad_in)
    h = Inches(0.32)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    pill.adjustments[0] = 0.5
    pill.shadow.inherit = False
    pill.fill.solid()
    pill.fill.fore_color.rgb = bg
    pill.line.fill.background()
    tf = pill.text_frame
    tf.word_wrap = False  # never wrap a pill — render as single line even if it overflows
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = fg
    return pill


def add_line(slide, x1, y1, x2, y2, *, color=LINE_GRAY, width=0.5):
    """Subtle line (used for grid lines, dividers)."""
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_chrome(slide, page_num, *, motif_color=BLUE, on_dark=False):
    """Add the v2 chrome on every slide:
    - L-shape motif top-left
    - Subtle grid lines
    - Backbase wordmark bottom-right + page number
    """
    text_color = WHITE if on_dark else NAVY
    grid_color = RGBColor(0xFF, 0xFF, 0xFF) if on_dark else LINE_GRAY
    grid_alpha_color = RGBColor(0x6B, 0x77, 0x86) if on_dark else LINE_GRAY

    # L-shape motif (top-left, navy + blue)
    motif_x = Inches(0.55)
    motif_y = Inches(0.50)
    s = Inches(0.18)
    half = Inches(0.09)
    # Top-right block
    b1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, motif_x + half, motif_y, half, half)
    b1.fill.solid(); b1.fill.fore_color.rgb = motif_color; b1.line.fill.background(); b1.shadow.inherit = False
    # Bottom row (full width)
    b2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, motif_x, motif_y + half, s, half)
    b2.fill.solid(); b2.fill.fore_color.rgb = motif_color; b2.line.fill.background(); b2.shadow.inherit = False

    # Grid lines — vertical at left margin + right margin; horizontal at top
    # under motif and bottom above footer
    add_line(slide, Inches(0.58), Inches(0.0), Inches(0.58), CANVAS_H, color=grid_alpha_color, width=0.5)
    add_line(slide, Inches(19.42), Inches(0.0), Inches(19.42), CANVAS_H, color=grid_alpha_color, width=0.5)
    add_line(slide, Inches(0.0), Inches(0.95), CANVAS_W, Inches(0.95), color=grid_alpha_color, width=0.5)
    add_line(slide, Inches(0.0), Inches(10.5), CANVAS_W, Inches(10.5), color=grid_alpha_color, width=0.5)

    # Backbase wordmark + page number (bottom-right)
    add_text(slide, Inches(17.5), Inches(10.65), Inches(1.6), Inches(0.4),
             "Backbase", size=14, bold=True, color=text_color, align=PP_ALIGN.RIGHT)
    add_text(slide, Inches(19.2), Inches(10.65), Inches(0.5), Inches(0.4),
             str(page_num), size=10, color=text_color, align=PP_ALIGN.RIGHT)


def add_speaker_notes(slide, key):
    """Inject speaker notes from SPEAKER_NOTES dict, preserving exact format."""
    notes = SPEAKER_NOTES.get(str(key), "")
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    p = notes_tf.paragraphs[0]
    r = p.add_run()
    r.text = notes
    r.font.size = Pt(11)
    r.font.name = FONT


# ──────────────────────────────────────────────────────────────────────────
# Per-slide builders
# ──────────────────────────────────────────────────────────────────────────

def slide_cover(prs):
    """Slide 1 — Cover."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide, NAVY)

    # L-motif (white on navy)
    motif_x = Inches(0.55); motif_y = Inches(0.50); s = Inches(0.18); half = Inches(0.09)
    b1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, motif_x + half, motif_y, half, half)
    b1.fill.solid(); b1.fill.fore_color.rgb = BLUE; b1.line.fill.background(); b1.shadow.inherit = False
    b2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, motif_x, motif_y + half, s, half)
    b2.fill.solid(); b2.fill.fore_color.rgb = BLUE; b2.line.fill.background(); b2.shadow.inherit = False

    # Backbase wordmark top-left
    add_text(slide, Inches(1.4), Inches(0.45), Inches(4.0), Inches(0.6),
             "Backbase", size=24, bold=True, color=WHITE)

    # Label
    add_text(slide, Inches(1.4), Inches(3.0), Inches(15.0), Inches(0.5),
             "AI ROUNDTABLE · NORDIC FINTECH FORUM", size=14, color=RGBColor(0xCC, 0xCC, 0xCC),
             letter_spacing=2)

    # Title (two lines, large light type)
    add_text(slide, Inches(1.4), Inches(3.7), Inches(15.0), Inches(3.5),
             "When every bank ships AI,\nwhere does difference live?",
             size=64, bold=False, color=WHITE, line_spacing=1.15)

    # Date
    add_text(slide, Inches(1.4), Inches(7.5), Inches(15.0), Inches(0.5),
             "Helsinki · May 2026 · Chatham House", size=18, color=RGBColor(0xCC, 0xCC, 0xCC))

    # Page number bottom-right
    add_text(slide, Inches(19.2), Inches(10.65), Inches(0.5), Inches(0.4),
             "1 / 13", size=10, color=RGBColor(0x99, 0xA8, 0xB9), align=PP_ALIGN.RIGHT)

    add_speaker_notes(slide, 1)


def _slide_header(slide, page_num, label, title, subtitle=None):
    """Shared header for content slides: label + title + optional subtitle."""
    add_chrome(slide, f"{page_num} / 13")
    # Label (small caps)
    add_text(slide, Inches(1.4), Inches(1.4), Inches(17.0), Inches(0.4),
             label, size=14, color=TEXT_MUTED, letter_spacing=2)
    # Title (big, light weight)
    add_text(slide, Inches(1.4), Inches(1.9), Inches(17.0), Inches(1.1),
             title, size=44, bold=False, color=NAVY, line_spacing=1.15)
    # Subtitle (optional, blue)
    if subtitle:
        add_text(slide, Inches(1.4), Inches(3.1), Inches(17.0), Inches(0.8),
                 subtitle, size=18, color=BLUE, line_spacing=1.45)


def _column_box(slide, x, y, w, h, *, subtitle=None, stat=None, body=None, source=None,
                divider_color=None):
    """Render a single column of a content-columns slide.
    Used for slides 2, 3, 7, 9, 11.
    All x/y/w/h are EMU integers (from Inches()).
    """
    if divider_color:
        add_line(slide, x + w, y, x + w, y + h, color=divider_color, width=0.5)
    cursor_y = y
    if subtitle:
        add_text(slide, x, cursor_y, w, Inches(0.5),
                 subtitle, size=12, bold=True, color=BLUE, letter_spacing=2)
        cursor_y += Inches(0.55)
    if stat:
        add_text(slide, x, cursor_y, w, Inches(1.4),
                 stat, size=58, bold=True, color=BLUE, line_spacing=1.0)
        cursor_y += Inches(1.6)
    if body:
        body_h = h - (cursor_y - y) - Inches(0.8)
        add_text(slide, x, cursor_y, w, body_h,
                 body, size=15, color=NAVY, line_spacing=1.5)
    if source:
        add_text(slide, x, y + h - Inches(0.7), w, Inches(0.6),
                 source, size=10, italic=True, color=TEXT_MUTED, line_spacing=1.45)


def slide_trap(prs):
    """Slide 2 — The Trap (3 columns of stats)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    _slide_header(slide, 2, "THE TRAP", "Intelligence is commoditising")

    col_y = Inches(4.4)
    col_h = Inches(5.3)
    col_gap = Inches(0.3)
    col_w = Inches((17.0 - 0.3 * 2) / 3)

    cols = [
        {
            "stat": "<5%",
            "body": "Performance gap between frontier models on banking-relevant tasks.\nClosing fast. The model is not the moat.",
            "source": "Stanford AI Index 2025; HELM benchmarks",
        },
        {
            "stat": "~100%",
            "body": "Of global top-50 banks have shipped a GenAI assistant since 2024.\nThe feature is saturated. Yours looks like theirs.",
            "source": "Evident AI Index, 2024",
        },
        {
            "stat": "~18 mo",
            "body": "Until feature parity — your copilot will look like your competitor's.\nDifferentiation has a clock.",
            "source": "Gartner AI Hype Cycle 2024; BCG GenAI Banking Pulse",
        },
    ]

    x = Inches(1.4)
    for i, c in enumerate(cols):
        _column_box(slide, x, col_y, col_w, col_h,
                    stat=c["stat"], body=c["body"], source=c["source"],
                    divider_color=LINE_GRAY if i < 2 else None)
        x += col_w + col_gap

    add_speaker_notes(slide, 2)


def slide_pattern(prs):
    """Slide 3 — The Pattern (2 columns of proof cases)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    _slide_header(slide, 3, "THE PATTERN · IT HAS HAPPENED BEFORE",
                  "The winners own what sits above the commodity.")

    col_y = Inches(4.4)
    col_h = Inches(5.3)
    col_gap = Inches(0.4)
    col_w = Inches((17.0 - 0.4) / 2)

    cols = [
        {
            "subtitle": "AWS · 2008 → 2020",
            "stat": "$100B+",
            "headline": "Compute commoditised. Orchestration became the moat.",
            "body": "EC2 reached parity with Azure and GCP within ~24 months. AWS pulled ahead on primitives, developer experience, and a trust layer regulators could verify. The silicon became a line item; the orchestration layer became the largest infrastructure business in history.",
            "source": "Synergy Research Cloud Market Share 2024; AWS segment results",
        },
        {
            "subtitle": "Spotify · 2008 → 2024",
            "stat": "675M",
            "headline": "Every catalogue became identical. Discovery became the product.",
            "body": "Rivals streamed the same songs. Spotify won on intent orchestration — Discover Weekly, Release Radar — and turned a commoditised catalogue into a 675M-user moat no competitor has closed.",
            "source": "Spotify MAU, Q4 2024; IFPI Global Music Report 2024",
        },
    ]

    x = Inches(1.4)
    for i, c in enumerate(cols):
        cursor_y = col_y
        add_text(slide, x, cursor_y, col_w, Inches(0.4),
                 c["subtitle"], size=12, bold=True, color=BLUE, letter_spacing=2)
        cursor_y += Inches(0.5)
        add_text(slide, x, cursor_y, col_w, Inches(1.2),
                 c["stat"], size=52, bold=True, color=BLUE, line_spacing=1.0)
        cursor_y += Inches(1.4)
        add_text(slide, x, cursor_y, col_w, Inches(0.7),
                 c["headline"], size=17, bold=True, color=NAVY, line_spacing=1.35)
        cursor_y += Inches(0.85)
        add_text(slide, x, cursor_y, col_w, Inches(2.4),
                 c["body"], size=14, color=NAVY, line_spacing=1.55)
        add_text(slide, x, col_y + col_h - Inches(0.5), col_w, Inches(0.5),
                 c["source"], size=10, italic=True, color=TEXT_MUTED, line_spacing=1.45)
        if i < len(cols) - 1:
            add_line(slide, x + col_w, col_y, x + col_w, col_y + col_h, color=LINE_GRAY)
        x += col_w + col_gap

    add_speaker_notes(slide, 3)


def _arch_card(slide, x, y, w, h, *, pill_text, pill_bg, pill_fg, title,
               title_color, body, customer_moment, border_color, body_top_color):
    """One of the three Architecture / Chain cards. Top: bank-architecture
    spine. Bottom (italic, separated): customer moment."""
    box = add_box(slide, x, y, w, h, fill=WHITE, line_color=border_color,
                  line_width=1.0, corner_radius=0.05)
    # Pill (top-left inside)
    pill_y = y + Inches(0.25)
    add_pill(slide, x + Inches(0.3), pill_y, pill_text, bg=pill_bg, fg=pill_fg, size=10)
    # Title (big, colored)
    add_text(slide, x + Inches(0.3), pill_y + Inches(0.5), w - Inches(0.6), Inches(0.9),
             title, size=36, bold=True, color=title_color, line_spacing=1.0)
    # Body (bank spine)
    add_text(slide, x + Inches(0.3), pill_y + Inches(1.55), w - Inches(0.6), Inches(2.0),
             body, size=13, color=NAVY, line_spacing=1.55)
    # Divider line + customer moment (italic, muted)
    div_y = y + h - Inches(1.4)
    add_line(slide, x + Inches(0.3), div_y, x + w - Inches(0.3), div_y, color=pill_bg)
    add_text(slide, x + Inches(0.3), div_y + Inches(0.15), w - Inches(0.6), Inches(1.1),
             customer_moment, size=12, italic=True, color=TEXT_MUTED, line_spacing=1.5)


def _arrow_between(slide, x, y_center, color):
    """Small right-arrow between cards."""
    arrow_w = Inches(0.4)
    arrow_h = Inches(0.4)
    add_text(slide, x - arrow_w / 2, y_center - arrow_h / 2, arrow_w, arrow_h,
             "→", size=28, bold=True, color=color, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


def slide_architecture(prs):
    """Slide 4 — The Architecture is shifting (3 cards, hybrid voice)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    _slide_header(slide, 4, "OUR OBSERVATION",
                  "The architecture of banking is shifting.",
                  "Most banks have built strong products — vertically. The customer doesn't see the silos. They see one bank, about to start asking it for advice.")

    # Three cards with arrows
    card_y = Inches(4.6)
    card_h = Inches(5.4)
    arrow_w = Inches(0.5)
    n = 3
    card_w = Inches((17.0 - arrow_w.inches * (n - 1)) / n)

    cards = [
        {
            "pill_text": "STRONG BUT ISOLATED",
            "pill_bg": RED_SOFT, "pill_fg": RED,
            "title": "Products", "title_color": RED,
            "body": "Vertically built, in silos. Onboarding lives in one system. Lending in another. Wealth in a third. The products work — the connections between them are manual.",
            "customer_moment": "The customer feels it: one bank, three apps, advice that's only ever a slice.",
            "border": RED,
        },
        {
            "pill_text": "CONNECTIVE TISSUE",
            "pill_bg": NAVY_SOFT, "pill_fg": NAVY,
            "title": "Trust", "title_color": NAVY,
            "body": "Identity, consent, risk, compliance — not as checkboxes, but as shared infrastructure that connects everything. Reusable across products, markets and channels.",
            "customer_moment": "When the AI advises, the bank that can explain HOW — provenance, peers, history — earns trust. The one that can't, doesn't.",
            "border": NAVY,
        },
        {
            "pill_text": "EMERGING FRONTIER",
            "pill_bg": LIGHT_BLUE, "pill_fg": BLUE,
            "title": "Intent", "title_color": BLUE,
            "body": "AI that doesn't just report or respond — it acts. Within governed boundaries.",
            "customer_moment": "Approves a top-up. Books a wealth review. Settles a dispute. Trust converts to the next product sale.",
            "border": BLUE,
        },
    ]

    x = Inches(1.4)
    arrow_color = [RED, BLUE]
    for i, c in enumerate(cards):
        _arch_card(slide, x, card_y, card_w, card_h,
                   pill_text=c["pill_text"], pill_bg=c["pill_bg"], pill_fg=c["pill_fg"],
                   title=c["title"], title_color=c["title_color"],
                   body=c["body"], customer_moment=c["customer_moment"],
                   border_color=c["border"], body_top_color=NAVY)
        if i < n - 1:
            _arrow_between(slide, x + card_w + arrow_w / 2,
                           card_y + card_h / 2, arrow_color[i])
        x += card_w + arrow_w

    # Footer band
    band_y = Inches(10.05)
    band = add_box(slide, Inches(1.4), band_y, Inches(17.0), Inches(0.4),
                   fill=LIGHT_BLUE, corner_radius=0.2)
    tf = band.text_frame
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Value is migrating from the product to the trust layer, and from trust to intent. The banks that connect trust to real-time decisioning at scale earn the customer's next move."
    r.font.name = FONT; r.font.size = Pt(11); r.font.color.rgb = NAVY

    add_speaker_notes(slide, 4)


def _chain_card(slide, x, y, w, h, *, pill_text, title, title_color, body,
                border_color, fill=WHITE):
    """One of the four Chain cards (slide 5). Simpler than arch_card —
    no customer-moment bottom block."""
    box = add_box(slide, x, y, w, h, fill=fill, line_color=border_color,
                  line_width=1.0, corner_radius=0.05)
    pill_y = y + Inches(0.25)
    add_pill(slide, x + Inches(0.25), pill_y, pill_text, bg=LIGHT_BLUE, fg=BLUE, size=9)
    add_text(slide, x + Inches(0.25), pill_y + Inches(0.5), w - Inches(0.5), Inches(0.9),
             title, size=32, bold=True, color=title_color, line_spacing=1.0)
    add_text(slide, x + Inches(0.25), pill_y + Inches(1.55), w - Inches(0.5), h - Inches(2.0),
             body, size=12, color=NAVY, line_spacing=1.55)


def slide_chain(prs):
    """Slide 5 — The Chain (4 cards: Trust → Intent → Orchestration → Wallet)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    _slide_header(slide, 5, "THE CHAIN · WHERE THE NEXT MOAT IS",
                  "Differentiation moves up the stack.",
                  "Banks competed bottom-up — Wallet, then Mind. AI flips it: Trust first, and the chain compounds back to a bigger wallet.")

    card_y = Inches(4.6)
    card_h = Inches(5.0)
    arrow_w = Inches(0.4)
    n = 4
    card_w = Inches((17.0 - arrow_w.inches * (n - 1)) / n)

    cards = [
        {"pill": "FOUNDATION", "title": "Trust", "color": BLUE, "border": BLUE,
         "body": "Digitised, provable, audited. Not a brand claim — plumbing. Explainability, provenance, consent, oversight become product features."},
        {"pill": "SIGNAL", "title": "Intent", "color": BLUE, "border": BLUE,
         "body": "Captured, delegated, contracted. Customers (and their agents) hand banks a mandate to act. Banks that cannot accept one cannot compete."},
        {"pill": "EXECUTION", "title": "Orchestration", "color": BLUE, "border": BLUE,
         "body": "Hand-offs across customer × employee × agent — without dropped context. The moat lives in the seams, not in any single channel."},
        {"pill": "OUTCOME", "title": "Share of wallet", "color": NAVY, "border": NAVY,
         "body": "The conversion that compounds. When the first three line up, the agent picks you — and the customer's next product follows."},
    ]

    x = Inches(1.4)
    for i, c in enumerate(cards):
        _chain_card(slide, x, card_y, card_w, card_h,
                    pill_text=c["pill"], title=c["title"], title_color=c["color"],
                    body=c["body"], border_color=c["border"],
                    fill=LIGHT_BLUE if i == n - 1 else WHITE)
        if i < n - 1:
            _arrow_between(slide, x + card_w + arrow_w / 2,
                           card_y + card_h / 2, BLUE)
        x += card_w + arrow_w

    # Footer band
    band_y = Inches(10.05)
    band = add_box(slide, Inches(1.4), band_y, Inches(17.0), Inches(0.4),
                   fill=LIGHT_BLUE, corner_radius=0.2)
    tf = band.text_frame
    tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Trust → Intent → Orchestration → Share of wallet. The order matters. Skip a link and the chain breaks."
    r.font.name = FONT; r.font.size = Pt(11); r.font.bold = True; r.font.color.rgb = NAVY

    add_speaker_notes(slide, 5)


def slide_statement(prs, page_num, label, text, highlight):
    """Statement layout — used for slides 6, 8, 10, 13 (Q slides + close).
    Big text with one phrase highlighted in blue."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    add_chrome(slide, f"{page_num} / 13")

    # Light blue band behind the text
    band_y = Inches(3.5)
    band_h = Inches(4.5)
    add_box(slide, Inches(1.0), band_y, Inches(18.0), band_h,
            fill=LIGHT_BLUE, corner_radius=0.0)

    # Label (top of band)
    add_text(slide, Inches(1.4), band_y + Inches(0.4), Inches(17.0), Inches(0.5),
             label, size=14, color=TEXT_MUTED, letter_spacing=2)

    # Text with highlighted phrase
    # Split text on the highlight phrase
    before, _, after = text.partition(highlight)
    after_text = after  # may include trailing punctuation
    add_rich_text(slide, Inches(1.4), band_y + Inches(1.2),
                  Inches(17.0), Inches(3.0),
                  [
                      {"text": before, "size": 36, "color": NAVY, "bold": False},
                      {"text": highlight, "size": 36, "color": BLUE, "bold": False},
                      {"text": after_text, "size": 36, "color": NAVY, "bold": False},
                  ],
                  line_spacing=1.4)

    add_speaker_notes(slide, page_num)


def _insight_cols(slide, columns):
    """3-column insight layout — used for I1 (7), I2 (9), I3 (11)."""
    col_y = Inches(4.6)
    col_h = Inches(5.4)
    col_gap = Inches(0.3)
    col_w = Inches((17.0 - 0.3 * 2) / 3)
    x = Inches(1.4)
    half_gap = Inches(0.15)
    for i, c in enumerate(columns):
        cursor_y = col_y
        if c.get("subtitle"):
            add_text(slide, x, cursor_y, col_w, Inches(0.5),
                     c["subtitle"], size=12, bold=True, color=BLUE, letter_spacing=2)
            cursor_y += Inches(0.55)
        if c.get("stat"):
            add_text(slide, x, cursor_y, col_w, Inches(1.4),
                     c["stat"], size=52, bold=True, color=BLUE, line_spacing=1.0)
            cursor_y += Inches(1.4)
        if c.get("headline"):
            add_text(slide, x, cursor_y, col_w, Inches(0.8),
                     c["headline"], size=16, bold=True, color=NAVY, line_spacing=1.4)
            cursor_y += Inches(0.9)
        if c.get("body"):
            body_h = col_h - (cursor_y - col_y) - Inches(0.8)
            add_text(slide, x, cursor_y, col_w, body_h,
                     c["body"], size=13, color=NAVY, line_spacing=1.55)
        if c.get("source"):
            add_text(slide, x, col_y + col_h - Inches(0.6), col_w, Inches(0.5),
                     c["source"], size=10, italic=True, color=TEXT_MUTED, line_spacing=1.45)
        if i < len(columns) - 1:
            add_line(slide, x + col_w + half_gap, col_y,
                     x + col_w + half_gap, col_y + col_h, color=LINE_GRAY)
        x += col_w + col_gap


def slide_i1(prs):
    """Slide 7 — I1: Trust is moving from claim to KPI."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    _slide_header(slide, 7, "INSIGHT 1 · TRUST IS THE NEW INFRASTRUCTURE",
                  "Trust is moving from claim to KPI.")
    _insight_cols(slide, [
        {
            "subtitle": "THE GAP · 1 IN 3",
            "stat": "1 in 3",
            "headline": "Firms at AI-governance maturity L3+ today.",
            "body": "Two-thirds of the industry cannot yet prove how their AI decides. That gap is the window — and it is closing.",
            "source": "McKinsey, The State of AI 2024; BCG Responsible AI Survey 2024",
        },
        {
            "subtitle": "THE SHIFT · DIGITISING THE RM",
            "headline": "Trust becomes a measurable KPI — a way to digitise what the RM did one-on-one.",
            "body": "Explainability, provenance, consent and human oversight are what gave the relationship manager their authority. Now they become product features that scale: Apple priced privacy (ATT), Mastercard sells AI-graded trust as infrastructure (~165B txns/yr scored).",
            "source": "EU AI Act Articles 13–15; NIST AI RMF 1.0; Mastercard Decision Intelligence",
        },
        {
            "subtitle": "THE EDGE · NORDICS",
            "headline": "Best-positioned region globally.",
            "body": "BankID and Swish — cooperated trust infrastructure with >85% adult reach. Nordic banks built shared identity and shared payment rails, then competed on the experience above. Trust-as-infrastructure here is an export, not a cost.",
            "source": "Finansiell ID-Teknik BID AB; Getswish AB, 2024",
        },
    ])
    add_speaker_notes(slide, 7)


def slide_i2(prs):
    """Slide 9 — I2: The moat has moved into the seams."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    _slide_header(slide, 9, "INSIGHT 2 · THE MOAT IS IN THE SEAMS",
                  "The moat has moved into the seams.")
    _insight_cols(slide, [
        {
            "subtitle": "WHERE IT BREAKS",
            "headline": "Context dies between systems.",
            "body": "The customer re-explains themselves to the bot, the call-centre agent, the relationship manager and the underwriter. Each had what was needed — none of them had it together.\n\n\"Faster cars on the same broken roads.\"",
        },
        {
            "subtitle": "WHAT GOOD LOOKS LIKE",
            "headline": "Context travels with the customer.",
            "body": "Every party — bot, banker, underwriter, the customer's own agent — opens the same authoritative view. The call becomes confirmation, not discovery. Hand-off measured in seconds, not days.",
        },
        {
            "subtitle": "WHERE TO ATTACK",
            "headline": "Start where revenue is decided.",
            "body": "Pick the seam where the next product or the next renewal is won or lost. Make context travel across that seam first. Compound from there.\n\nThe moat is in the seams, not in any single channel.",
        },
    ])
    add_speaker_notes(slide, 9)


def slide_i3(prs):
    """Slide 11 — I3: When someone else owns the intent."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    _slide_header(slide, 11, "INSIGHT 3 · WHEN SOMEONE ELSE OWNS THE INTENT",
                  "When someone else owns the intent, they capture the relationship.")
    _insight_cols(slide, [
        {
            "subtitle": "KLARNA · SWEDEN, 2005 →",
            "stat": "150M",
            "headline": "The shopping agent owned the credit.",
            "body": "Klarna owned the checkout moment — the intent layer of consumer commerce. 150M+ consumers globally. The bank extending credit became invisible; Klarna became the brand on the screen.",
            "source": "Klarna investor materials, 2024; BCG Global Payments Report 2024",
        },
        {
            "subtitle": "PLAID · US, 2013 →",
            "stat": "1 in 4",
            "headline": "The API between customers and their banks.",
            "body": "One in four US adults have connected an account through Plaid. Every consumer-fintech app routes through it; banks became the backend. The intent layer moved to a startup.",
            "source": "Plaid public statements, 2024; CFPB Section 1033 filings",
        },
        {
            "subtitle": "M-PESA · KENYA, 2007 →",
            "stat": "~60%",
            "headline": "A telco became the bank.",
            "body": "Safaricom's M-Pesa captured trust + intent before Kenyan banks modernised. It now moves ~60% of Kenya's GDP. Banks remain — they are no longer the primary layer the customer interacts with.",
            "source": "Safaricom FY24 annual report; CBK National Payments Strategy 2022–25",
        },
    ])
    add_speaker_notes(slide, 11)


def slide_choice(prs):
    """Slide 12 — The Choice: bolt-on vs re-platform."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, WHITE)
    _slide_header(slide, 12, "THE CHOICE · TWO RESPONSES TO AI",
                  "Bolt-on programmes lose. Re-platforms win.",
                  "An identity shift, not a tech upgrade. The decision is not whether to use AI — it is what AI plugs into.")

    card_y = Inches(4.6)
    card_h = Inches(4.6)
    card_gap = Inches(0.4)
    card_w = Inches((17.0 - 0.4) / 2)

    # Left card — Pattern A · Bolt-on (red)
    left = add_box(slide, Inches(1.4), card_y, card_w, card_h,
                   fill=RGBColor(0xFA, 0xE0, 0xDE), line_color=None, corner_radius=0.03)
    # Left red bar
    add_box(slide, Inches(1.4), card_y, Inches(0.12), card_h, fill=RED)

    inner_x = Inches(1.4) + Inches(0.45)
    inner_w = card_w - Inches(0.7)
    add_text(slide, inner_x, card_y + Inches(0.4), inner_w, Inches(0.4),
             "PATTERN A · BOLT-ON", size=12, bold=True, color=RED, letter_spacing=2)
    add_text(slide, inner_x, card_y + Inches(0.95), inner_w, Inches(0.7),
             "AI added to product silos.", size=22, bold=True, color=NAVY, line_spacing=1.3)
    bullets_a = [
        "Channel-by-channel agents, each owned by a different P&L",
        "Five copies of the customer; none authoritative",
        "AI fluent at the front, broken at the back",
        "NPS plateaus; cost-to-serve doesn't move",
    ]
    add_text(slide, inner_x, card_y + Inches(1.85), inner_w, Inches(1.8),
             "▸ " + "\n▸ ".join(bullets_a),
             size=13, color=TEXT_MUTED, line_spacing=1.7)
    add_text(slide, inner_x, card_y + card_h - Inches(0.65), inner_w, Inches(0.5),
             "\"Faster cars on the same broken roads.\"",
             size=11, italic=True, color=TEXT_MUTED)

    # Right card — Pattern B · Re-platform (blue)
    right_x = Inches(1.4) + card_w + card_gap
    add_box(slide, right_x, card_y, card_w, card_h,
            fill=LIGHT_BLUE, line_color=None, corner_radius=0.03)
    add_box(slide, right_x, card_y, Inches(0.12), card_h, fill=BLUE)

    inner_x = right_x + Inches(0.45)
    add_text(slide, inner_x, card_y + Inches(0.4), inner_w, Inches(0.4),
             "PATTERN B · RE-PLATFORM", size=12, bold=True, color=BLUE, letter_spacing=2)
    add_text(slide, inner_x, card_y + Inches(0.95), inner_w, Inches(0.7),
             "A single fabric for the customer life.", size=22, bold=True, color=NAVY, line_spacing=1.3)
    bullets_b = [
        "One identity, one consent, one customer record",
        "Orchestration layer any AI agent can plug into",
        "Consistent CX from app to branch to advisor",
        "Lift compounds — NPS, NTB, cost and ROE move together",
    ]
    add_text(slide, inner_x, card_y + Inches(1.85), inner_w, Inches(1.8),
             "▸ " + "\n▸ ".join(bullets_b),
             size=13, color=NAVY, line_spacing=1.7)
    add_text(slide, inner_x, card_y + card_h - Inches(0.65), inner_w, Inches(0.5),
             "\"One operating fabric for the customer life.\"",
             size=11, italic=True, color=NAVY)

    # Footer callout
    foot_y = Inches(9.55)
    add_box(slide, Inches(1.4), foot_y, Inches(17.0), Inches(0.6),
            fill=WHITE, line_color=LINE_GRAY, line_width=0.75, corner_radius=0.1)
    add_box(slide, Inches(1.4), foot_y, Inches(0.12), Inches(0.6), fill=BLUE)
    add_text(slide, Inches(1.85), foot_y + Inches(0.15), Inches(16.5), Inches(0.4),
             "The decision is not whether to use AI. It is what AI plugs into. — A bolt-on programme adds intelligence to silos. A re-platform turns the bank into the platform.",
             size=12, color=NAVY, italic=True)

    add_speaker_notes(slide, 12)


# ──────────────────────────────────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = CANVAS_W
    prs.slide_height = CANVAS_H

    # Build slides in order
    slide_cover(prs)                                                       # 1
    slide_trap(prs)                                                        # 2
    slide_pattern(prs)                                                     # 3
    slide_architecture(prs)                                                # 4
    slide_chain(prs)                                                       # 5
    slide_statement(prs, 6, "VALIDATE · QUESTION 1 · TRUST",
                    "When the model is the same as your competitor's, what makes you the trusted counterparty?",
                    "trusted counterparty")                                # 6
    slide_i1(prs)                                                          # 7
    slide_statement(prs, 8, "VALIDATE · QUESTION 2 · HAND-OFF",
                    "Where in your bank would the customer × employee × agent hand-off break today?",
                    "break today")                                         # 8
    slide_i2(prs)                                                          # 9
    slide_statement(prs, 10, "VALIDATE · QUESTION 3 · AGENT'S CHOICE",
                    "If an agent shopped this mortgage for your customer, would it pick you?",
                    "pick you")                                            # 10
    slide_i3(prs)                                                          # 11
    slide_choice(prs)                                                      # 12
    slide_statement(prs, 13, "THE TAKE-HOME",
                    "The bank that orchestrates trust and intent wins the agent's choice — and the customer's wallet.",
                    "trust and intent")                                    # 13

    out_path = HERE / "nordic_fintech_forum_opening_pov.pptx"
    prs.save(out_path)
    print(f"Wrote: {out_path}  ({out_path.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
