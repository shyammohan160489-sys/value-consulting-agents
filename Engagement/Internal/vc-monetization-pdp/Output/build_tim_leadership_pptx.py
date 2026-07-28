#!/usr/bin/env python3
"""Shyam × Tim — the leadership case, v2 (visual / infographic edition).

Nine slides on the locked exhibit-slides-pptx engine, one exhibit per question:
cover · strategy house · why-now convergence · the arsenal · the ammunition belt ·
the Backbase leverage stack · the skills ladder · the flywheel → revenue · the ask.
All exhibits composed from engine primitives plus flat autoshapes (chevron,
pentagon) under the same flat/strip discipline. Every claim traces to a repo
artifact; nothing invented.

Run:  python3 build_tim_leadership_pptx.py [out.pptx]
"""
import os
import sys

_HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.abspath(os.path.join(
    _HERE, '..', '..', '..', '..', '.claude', 'skills', 'exhibit-slides-pptx', 'scripts')))
from exhibit_pptx import (ExhibitDeck, NAVY, BLUE, BLUE2, BLUE3, BLUE4, TINT, TINT2,
                          CYAN, CORAL, WHITE, MUT, FN, HAIR, SUB_D, W, H)
from pptx.util import Inches as I, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = sys.argv[1] if len(sys.argv) > 1 else os.path.join(_HERE, 'Shyam_Talent_Programme_Tim.pptx')
d = ExhibitDeck()


def shape(s, kind, x, y, w, h, fill, line=None, line_w=0.75):
    sp = s.shapes.add_shape(kind, I(x), I(y), I(w), I(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    d.flat(sp)
    return sp


# ═════════ S1 · COVER (dark) ═════════
s = d.slide(dark=True)
d.step_glyph(s, 0.406, 0.406, 0.167, 0.166, CYAN)
d.txt(s, 1.0, 2.0, 11.0, 0.35, 'TALENT PROGRAMME · SHYAM × TIM RUTTNER · ONE HOUR', size=12.5, color=CYAN, bold=True)
d.txt(s, 1.0, 2.45, 11.4, 1.9, [[('The operator behind', 44, WHITE, False)],
                                [('the playbook.', 44, WHITE, True)]], line_sp=1.05)
d.txt(s, 1.0, 4.35, 9.4, 0.9, 'Seven questions, one exhibit each: the strategy, the why, the arsenal, the '
                              'ammunition, what I leverage, what I build, and how it all feeds revenue.',
      size=15, color=SUB_D, line_sp=1.25)
d.txt(s, 1.0, 6.35, 9.0, 0.3, 'Shyam · July 2026 · companion to the Product Factory execution playbook',
      size=11, color=MUT)
d.txt(s, 11.7, 7.05, 1.06, 0.3, 'Backbase', size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
d.notes(s, 'FRAME THE HOUR: 15 min this deck, 10 min the playbook short version, 35 min working the asks.')

# ═════════ S2 · THE STRATEGY (house) ═════════
s = d.slide()
d.chrome(s, 'Q1 · what is the strategy', 'The strategy: turn value consulting into a product-led P&L')
hx, hw = 1.0, 7.45
d.rect(s, hx, 1.95, hw, 0.66, fill=NAVY, round_=True)
d.txt(s, hx + 0.25, 2.05, hw - 0.5, 0.2, 'AMBITION', size=8.5, color=CYAN, bold=True)
d.txt(s, hx + 0.25, 2.25, hw - 0.5, 0.3, 'The VC function becomes a revenue line that pre-installs Banking OS',
      size=11.5, color=WHITE, bold=True)
d.rect(s, hx, 2.73, 3.66, 1.92, fill=TINT, round_=True)
d.txt(s, hx + 0.2, 2.86, 3.3, 0.2, 'ENGINE A · PRODUCTS', size=9, color=BLUE2, bold=True)
d.txt(s, hx + 0.2, 3.10, 3.3, 1.4, [[('· Four paid diagnostic SKUs, €15-100K', 9.5, NAVY, False)],
                                    [('· Each installs a layer: Missions, Nexus, Sentinel', 9.5, NAVY, False)],
                                    [('· Fees credit forward into the POC', 9.5, NAVY, False)]], line_sp=1.18, sp_after=4)
d.rect(s, hx + 3.79, 2.73, 3.66, 1.92, fill=TINT, round_=True)
d.txt(s, hx + 3.99, 2.86, 3.3, 0.2, 'ENGINE B · SERVICES', size=9, color=BLUE2, bold=True)
d.txt(s, hx + 3.99, 3.10, 3.3, 1.4, [[('· AI-native operating model + org design', 9.5, NAVY, False)],
                                     [('· Sold downstream, on Engine A evidence', 9.5, NAVY, False)],
                                     [('· Opens the CEO and CHRO door', 9.5, NAVY, False)]], line_sp=1.18, sp_after=4)
d.rect(s, hx, 4.77, hw, 0.44, fill=BLUE, round_=True)
d.txt(s, hx + 0.25, 4.86, hw - 0.5, 0.26,
      'THE WEDGE LADDER · free Inspire  →  €60K wedge  →  €60K Mission POC  →  recurring assurance',
      size=9.5, color=WHITE, bold=True)
d.rect(s, hx, 5.33, hw, 0.44, fill=TINT2, round_=True)
d.txt(s, hx + 0.25, 5.42, hw - 0.5, 0.26,
      'FOUNDATIONS · the Cortex machine · evidence discipline · the locked design language',
      size=9.5, color=MUT, bold=True)
d.txt(s, 8.85, 1.98, 3.85, 0.22, 'WHAT SUCCESS LOOKS LIKE', size=9, color=MUT, bold=True)
targets = [('Year one', '3-4 installs', 'the proof year, €0.4-0.5M'),
           ('Year two', '10 a year', 'the cost-neutral rate, ~€1.2M covered'),
           ('Then', 'recurring', '€30-50K per account per year + gain share')]
for i, (tag, big, lab) in enumerate(targets):
    y = 2.28 + i * 1.20
    d.rect(s, 8.85, y, 3.85, 1.06, fill=TINT2, round_=True)
    d.txt(s, 9.05, y + 0.10, 3.5, 0.2, tag.upper(), size=8.5, color=BLUE, bold=True)
    d.txt(s, 9.05, y + 0.30, 3.5, 0.34, big, size=17, color=BLUE2, bold=True)
    d.txt(s, 9.05, y + 0.68, 3.5, 0.3, lab, size=9, color=MUT, line_sp=1.1)
d.footnote(s, 'The full execution behind each element is the 57-slide playbook; this house is its one-page compression.')
d.notes(s, 'STRATEGY HOUSE. Roof = ambition, pillars = the two engines, beam = the commercial ladder, '
           'foundation = the machine. Right rail = the three numbers that define success.')

# ═════════ S3 · WHY (convergence) ═════════
s = d.slide()
d.chrome(s, 'Q2 · why are we doing this', 'Why: three currents converge on this move, right now')
whys = [
    ('THE MARKET', 'The decade shifts from channels to operating models; whoever measures outcomes owns the renewal.'),
    ('THE FUNCTION', 'Value consulting must earn its keep: cost-neutral makes it protected, revenue makes it grow.'),
    ('ME', 'I am ready for commercial ownership, and the systems I build already scale beyond my own hands.'),
]
for i, (tag, body) in enumerate(whys):
    y = 2.05 + i * 1.32
    shape(s, MSO_SHAPE.PENTAGON, 1.0, y, 6.55, 1.12, TINT if i == 1 else TINT2)
    d.txt(s, 1.30, y + 0.16, 1.6, 0.24, tag, size=10, color=BLUE2, bold=True)
    d.txt(s, 2.95, y + 0.14, 4.15, 0.9, body, size=10.5, color=NAVY, line_sp=1.2)
d.rect(s, 8.10, 2.05, 4.62, 3.96, fill=NAVY, round_=True)
d.txt(s, 8.40, 2.30, 4.0, 0.24, 'THE JOIN', size=10.5, color=CYAN, bold=True)
d.txt(s, 8.40, 2.62, 4.05, 0.68, 'A product-led value line serves all three at once.', size=14.5, color=WHITE,
      bold=True, line_sp=1.15)
for i, (who, what) in enumerate([('Backbase gets', 'pipeline that arrives pre-installed and pre-evidenced'),
                                 ('The team gets', 'growth funded by revenue it generates itself'),
                                 ('I get', 'the leadership step: a line to own and be measured on')]):
    y = 3.55 + i * 0.74
    d.oval(s, 8.42, y + 0.02, 0.12, 0.12, CYAN)
    d.txt(s, 8.64, y - 0.05, 3.9, 0.24, who.upper(), size=8.5, color=CYAN, bold=True)
    d.txt(s, 8.64, y + 0.16, 3.95, 0.44, what, size=10, color=SUB_D, line_sp=1.12)
d.footnote(s, 'The market thesis is the Banking OS canon; the function thesis is the PDP north star; the third is the talent-programme conversation itself.')
d.notes(s, 'CONVERGENCE. Three arrows, one join. If Tim pushes on any single arrow, the other two still carry the case.')

# ═════════ S4 · WHAT WE HAVE (the arsenal) ═════════
s = d.slide()
d.chrome(s, 'Q3 · what do we have at our disposal', 'At my disposal: evidence, a machine, and a shipped record')
arsenal = [
    ('65', 'processes in the APA V3 catalog', 'EVIDENCE', BLUE,
     ['283 steps, six layers deep, canon-labelled', 'The conversational cost model (with Mayur + Deepak)',
      'The autonomy framework: the value-side twin', 'Benchmarks + ROI models from the validation cohort']),
    ('24h', 'per turn of the Pursuit Loop', 'THE MACHINE', BLUE2,
     ['Cortex: the VC AgenticOS, launched and governed', 'The Flywheel: telemetry that ships its own fixes',
      'The locked exhibit engine, ratified team default', 'Publish-without-git, so consultants never stall']),
    ('10', 'accounts shaped or validated', 'THE RECORD', NAVY,
     ['SNB Capital · BACB · ABSA ×5 workstreams', 'HSBC · Schroders · SEB · NFIS',
      'BECU, WSFS, MyState pipeline validations', 'Nordic FinTech Forum 2026: the public POV']),
]
for i, (num, numlab, tag, c, items) in enumerate(arsenal):
    x = 1.0 + i * 3.97
    d.rect(s, x, 1.95, 3.77, 3.85, fill=TINT2, round_=True)
    d.rect(s, x, 1.95, 3.77, 0.34, fill=c, round_=True)
    d.txt(s, x + 0.16, 2.01, 3.4, 0.24, tag, size=9.5, color=WHITE, bold=True)
    d.txt(s, x + 0.20, 2.44, 3.4, 0.62, num, size=38, color=BLUE2, bold=True)
    d.txt(s, x + 0.20, 3.14, 3.4, 0.3, numlab, size=10, color=MUT, bold=True, line_sp=1.05)
    for j, it in enumerate(items):
        d.oval(s, x + 0.22, 3.62 + j * 0.52 + 0.05, 0.09, 0.09, c)
        d.txt(s, x + 0.42, 3.56 + j * 0.52, 3.2, 0.48, it, size=9, color=NAVY, line_sp=1.1)
d.footnote(s, 'Counts are auditable in the Cortex repository: the APA catalog, the pursuit-loop method file, engagement and validation outputs.')
d.notes(s, 'ARSENAL. Three big numbers, each defensible: 65 (apa_v3.json), 24h (pursuit-loop-method.md), 10 (engagement + tests folders).')

# ═════════ S5 · THE AMMUNITION (chevron belt) ═════════
s = d.slide()
d.chrome(s, 'Q4 · what ammunition do we fire', 'The ammunition: five loaded offers and eight EAP slots')
rounds = [
    ('Process X-Ray', '€50-75K', BLUE, 'Prices the leakage; lands the Mission backlog'),
    ('Value Telemetry', '€15-50K+', BLUE, 'Cost per outcome; opens the recurring line'),
    ('Guardrail Studio', '€50-100K', BLUE3, 'Policy into rules; loads Sentinel'),
    ('Cartographer', '€75-100K', BLUE4, 'Data into truth; scopes Nexus'),
    ('Org Design', '€80-150K', BLUE4, 'The inverted-T org; the CHRO door'),
]
d.txt(s, 1.0, 1.95, 6.0, 0.22, 'THE FIVE ROUNDS · SHADE = FEASIBILITY TODAY', size=9, color=MUT, bold=True)
for i, (name, price, c, what) in enumerate(rounds):
    x = 1.0 + i * 2.24
    shape(s, MSO_SHAPE.CHEVRON, x, 2.25, 2.40, 0.92, c)
    tc = WHITE if c in (BLUE, BLUE2, BLUE3) else NAVY
    tx0 = x + (0.18 if i == 0 else 0.52)
    d.txt(s, tx0, 2.38, 1.80, 0.24, name, size=9.5, color=tc, bold=True)
    d.txt(s, tx0, 2.62, 1.75, 0.22, price, size=9.5, color=CYAN if c in (BLUE, BLUE3) else BLUE2, bold=True)
    d.txt(s, x + 0.06, 3.30, 2.14, 0.65, what, size=8.5, color=NAVY, line_sp=1.15)
shape(s, MSO_SHAPE.OVAL, 11.76, 2.24, 0.94, 0.94, NAVY)
d.txt(s, 11.76, 2.50, 0.94, 0.5, 'the\nPOC', size=10, color=CYAN, bold=True, align=PP_ALIGN.CENTER, line_sp=1.0)
d.txt(s, 1.0, 4.25, 6.0, 0.22, 'THE PROPELLANT · WHAT MAKES EACH ROUND FIRE', size=9, color=MUT, bold=True)
prop = [('8 EAP slots', 'nominated, never tendered; co-creation as the route past procurement'),
        ('Specimen packs', 'the exact artifacts shown before any signature; the bank buys a known object'),
        ('Credit-forward fees', '100% of the wedge credits into the POC; budget is never spent on a maybe')]
for i, (h_, b_) in enumerate(prop):
    x = 1.0 + i * 3.97
    d.rect(s, x, 4.55, 3.77, 0.92, fill=TINT2, round_=True)
    d.rect(s, x, 4.61, 0.045, 0.80, fill=BLUE)
    d.txt(s, x + 0.2, 4.65, 3.4, 0.22, h_, size=10.5, color=NAVY, bold=True)
    d.txt(s, x + 0.2, 4.90, 3.45, 0.5, b_, size=8.5, color=MUT, line_sp=1.12)
d.takeaway_band(s, 'Every round is aimed at the same target: ', 'a live Mission POC that installs Banking OS.', y=5.75)
d.footnote(s, 'Shades: solid blue = high feasibility, mid = gated, light = a named dependency. Detail: playbook short version, slide 4.', y=6.48)
d.notes(s, 'AMMUNITION BELT. Five chevrons pointing at one target. The propellant row answers “why would they buy it” in three objects.')

# ═════════ S6 · LEVERAGE FROM BACKBASE (two-layer stack) ═════════
s = d.slide()
d.chrome(s, 'Q5 · what do I leverage from Backbase', 'From Backbase I leverage the platform; I add the packaging')
d.rect(s, 1.0, 2.00, 11.708, 1.28, fill=NAVY, round_=True)
d.txt(s, 1.25, 2.12, 6.0, 0.22, 'WHAT I ADD ON TOP', size=9.5, color=CYAN, bold=True)
adds = ['Productized packaging + pricing', 'The value-evidence method', 'The org-design practice', 'Category collateral, with your team']
for i, a_ in enumerate(adds):
    x = 1.25 + i * 2.87
    d.rect(s, x, 2.42, 2.70, 0.62, fill=NAVY, line=CYAN, line_w=0.9, round_=True)
    d.txt(s, x + 0.14, 2.48, 2.45, 0.5, a_, size=9.5, color=WHITE, bold=True, line_sp=1.1, anchor=MSO_ANCHOR.MIDDLE)
for i in range(3):
    d.txt(s, 3.3 + i * 3.1, 3.36, 0.4, 0.3, '↑', size=16, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
d.rect(s, 1.0, 3.74, 11.708, 2.10, fill=TINT, round_=True)
d.txt(s, 1.25, 3.88, 8.0, 0.22, 'WHAT BACKBASE ALREADY GIVES ME · NOTHING HERE NEEDS BUILDING', size=9.5, color=BLUE2, bold=True)
gives = [('Factory tooling', 'Designer, Modeler, Connector Studio, Simulation'),
         ('Nexus + Sentinel', 'the product the wedges pre-install'),
         ('The EAP', '8 slots of co-creation currency'),
         ('AE channel + brand', 'accounts, trust, the 120+ logo base'),
         ('The cost model', 'Mayur + Deepak’s supply-side truth'),
         ('FDE capacity', 'Mission Sprints staffed at install')]
for i, (h_, b_) in enumerate(gives):
    x = 1.25 + (i % 3) * 3.80
    y = 4.20 + (i // 3) * 0.80
    d.rect(s, x, y, 3.60, 0.68, fill=WHITE, round_=True)
    d.txt(s, x + 0.16, y + 0.07, 3.3, 0.22, h_, size=10, color=NAVY, bold=True)
    d.txt(s, x + 0.16, y + 0.32, 3.35, 0.3, b_, size=8.5, color=MUT, line_sp=1.08)
d.takeaway_band(s, 'The ask is leverage, never construction: ', 'I package what the company already built.', y=6.00)
d.footnote(s, 'This is why the build effort stays at person-weeks: every SKU productizes an existing Factory tool.', y=6.52)
d.notes(s, 'LEVERAGE STACK. Bottom layer = given (six tiles). Top layer = my four additions. The arrows carry the whole argument: thin layer on a deep platform.')

# ═════════ S7 · SKILLS I AM BUILDING (from→to ladder) ═════════
s = d.slide()
d.chrome(s, 'Q6 · what am I building in myself', 'The skills I am building: from consultant to line owner')
skills = [
    ('Commercial ownership', 'shaping deals', 'owning a quota and a P&L', 0.45),
    ('Product management', 'delivery kits', 'SKU lifecycle: price, launch, iterate', 0.40),
    ('Category building', 'deck-level narrative', 'launch collateral with your team', 0.35),
    ('Executive altitude', 'COO and ops rooms', 'CEO and CHRO conversations', 0.55),
    ('Org-design practice', 'method v0.1 on paper', 'a flagship practice, evidence-fed', 0.30),
]
d.txt(s, 7.05, 1.95, 2.0, 0.2, 'TODAY', size=8.5, color=BLUE3, bold=True)
d.txt(s, 10.9, 1.95, 1.8, 0.2, 'TARGET', size=8.5, color=BLUE2, bold=True)
for i, (skill, frm, to, pos) in enumerate(skills):
    y = 2.25 + i * 0.72
    d.txt(s, 1.0, y + 0.08, 2.4, 0.4, skill, size=10.5, color=NAVY, bold=True, line_sp=1.05)
    lx0, lx1 = 3.55, 12.35
    d.hline(s, lx0, y + 0.20, lx1, y + 0.20, color=HAIR, wpt=1.4)
    tx = lx0 + (lx1 - lx0) * pos
    d.oval(s, tx - 0.08, y + 0.12, 0.16, 0.16, BLUE3)
    d.txt(s, tx - 1.1, y + 0.32, 2.2, 0.22, frm, size=8, color=MUT, align=PP_ALIGN.CENTER)
    d.oval(s, lx1 - 0.10, y + 0.10, 0.20, 0.20, BLUE2)
    d.txt(s, lx1 - 2.6, y - 0.10, 2.5, 0.2, to, size=8, color=BLUE2, bold=True, align=PP_ALIGN.RIGHT)
d.takeaway_band(s, 'This is the development ask: ', 'the gap between the dots is what the programme and your sponsorship close.', y=5.95)
d.footnote(s, 'Category building is the row where the CMO track matters most: I want to learn launching from the people who launch.', y=6.48)
d.notes(s, 'SKILLS LADDER. Five rows, today-dot vs target-dot. Row three (category building) is the explicit Tim ask: mentorship where his team is world-class.')

# ═════════ S8 · THE FLYWHEEL → REVENUE ═════════
s = d.slide()
d.chrome(s, 'Q7 · how does it feed revenue', 'The flywheel: every install makes the next sale cheaper')
cx0, cy0, r = 4.05, 4.05, 1.72
import math
nodes = [
    ('1', 'Free wedge qualifies', TINT2),
    ('2', 'Paid install lands €', TINT),
    ('3', 'Evidence harvested', TINT2),
    ('4', 'Kits sharpen; next sale cheaper', TINT),
    ('5', 'Platform pull: POC + licence', TINT2),
    ('6', 'Telemetry recurs €, funds next', TINT),
]
pos = []
for i in range(6):
    ang = math.radians(90 - i * 60)
    pos.append((cx0 + r * math.cos(ang), cy0 - r * math.sin(ang)))
for i in range(6):
    x1, y1 = pos[i]
    x2, y2 = pos[(i + 1) % 6]
    d.hline(s, x1, y1, x2, y2, color=HAIR, wpt=1.2)
shape(s, MSO_SHAPE.OVAL, cx0 - 0.72, cy0 - 0.72, 1.44, 1.44, NAVY)
d.txt(s, cx0 - 0.72, cy0 - 0.42, 1.44, 0.5, '↻', size=22, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
d.txt(s, cx0 - 0.70, cy0 + 0.02, 1.40, 0.4, 'compounds', size=9.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
for i, (n_, lab, fill) in enumerate(nodes):
    x, y = pos[i]
    d.rect(s, x - 1.02, y - 0.31, 2.04, 0.62, fill=fill, round_=True)
    d.oval(s, x - 0.96, y - 0.13, 0.26, 0.26, BLUE, label=n_, fs=9)
    d.txt(s, x - 0.62, y - 0.24, 1.58, 0.5, lab, size=8.5, color=NAVY, bold=True, line_sp=1.02,
          anchor=MSO_ANCHOR.MIDDLE)
d.txt(s, 8.30, 1.98, 4.3, 0.22, 'HOW THE WHEEL PAYS OUT', size=9, color=MUT, bold=True)
pays = [('Year one', '€0.4-0.5M', '3-4 installs, one per family', TINT2),
        ('Year two', '~€1.2M', 'ten installs a year: cost-neutral', TINT),
        ('Layered on top', '€30-50K /acct /yr', 'telemetry recurrence + gain share', TINT)]
for i, (tag, big, lab, fill) in enumerate(pays):
    y = 2.30 + i * 1.06
    d.rect(s, 8.30, y, 4.42, 0.94, fill=fill, round_=True)
    d.txt(s, 8.50, y + 0.08, 4.0, 0.2, tag.upper(), size=8.5, color=BLUE, bold=True)
    d.txt(s, 8.50, y + 0.27, 2.6, 0.34, big, size=16, color=BLUE2, bold=True)
    d.txt(s, 8.50, y + 0.62, 4.05, 0.28, lab, size=8.5, color=MUT)
d.open_badge(s, 8.30, 5.52, 4.42, 'Unsized · ', 'licence pull-through from installs; sized with finance once the pilot lands.', h=0.66)
d.footnote(s, 'The wheel’s discipline: node 3 is the moat (evidence nobody else holds) and node 6 pays for node 1’s next turn.')
d.notes(s, 'THE FLYWHEEL. Walk it once clockwise, then point right: the wheel is the mechanism, the rail is the money. '
           'DEFENSE: year figures are the playbook’s capacity-checked numbers, not aspirations.')

# ═════════ S9 · THE ASK ═════════
s = d.slide()
d.chrome(s, 'The ask · where this goes', 'I am asking for ownership, sponsorship and a seed, in order')
d.txt(s, 1.0, 1.95, 5.0, 0.24, 'THE MOVE I AM MAKING', size=10, color=MUT, bold=True, track='100')
cards = [('senior consultant', 'Line P&L', 'owner of the products + services line, with a real quota'),
         ('one pair of hands', 'Cortex', 'the team runs what I build: the AgenticOS, kits, methods'),
         ('deliverables', 'Installs', 'every engagement leaves Banking OS components behind'),
         ('€5,000 budget', 'Seed', 'product one packaged, demoed and launched with marketing')]
cx = 1.0
for f_, v_, l_ in cards:
    d.stat_card(s, cx, 2.28, f_, v_, l_)
    cx += 2.975
asks = [
    ('1 · A named path to owning the line', 'The leadership development I want is commercial ownership: the product + services P&L, reviewed on revenue.'),
    ('2 · ExCo-altitude sponsorship', 'CMO air cover on the category story and launch collateral; a door to the EAP nominations.'),
    ('3 · The €5,000 as seed, not training', 'It funds the Process X-Ray launch kit; the course it replaces would have taught me less than January will.'),
]
for i, (h_, b_) in enumerate(asks):
    y = 3.80 + i * 0.72
    d.rect(s, 1.0, y, 11.708, 0.62, fill=TINT2, round_=True)
    d.rect(s, 1.0, y + 0.06, 0.045, 0.50, fill=BLUE)
    d.txt(s, 1.24, y + 0.07, 3.85, 0.5, h_, size=10.5, color=NAVY, bold=True, line_sp=1.05)
    d.txt(s, 5.25, y + 0.10, 7.3, 0.46, b_, size=10, color=NAVY, line_sp=1.12)
d.takeaway_band(s, 'Judge me on January: ', 'one product in market, sold exactly the way the playbook says.', y=6.02)
d.footnote(s, 'The measurable commitments behind this sit in the playbook: proof-year targets, gates and the honest capacity math (chapters 05-06).', y=6.55)
d.notes(s, 'CLOSE OF THE HOUR. Ownership is the development goal, sponsorship is Tim’s currency, the seed is the smallest yes.')

d.save(OUT)
print(f'Wrote {OUT} ({len(d.prs.slides._sldIdLst)} slides)')
