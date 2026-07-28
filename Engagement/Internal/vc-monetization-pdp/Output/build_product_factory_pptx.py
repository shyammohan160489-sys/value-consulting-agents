#!/usr/bin/env python3
"""Product Factory execution plan — built on the exhibit-slides-pptx engine (v3.1, locked).

Per the ratified default (Shyam, 28 Jul 2026): exhibit lane, chrome v3.1 with right
rail, palette #071224/#4066F5, one exhibit per slide, live slidenum page fields.
Content unchanged from deck v3 (41 slides). The engine is never modified here;
custom exhibits are composed from its primitives in the same grammar.

Run:  python3 build_product_factory_pptx.py [out.pptx]
"""
import os
import sys

_HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.abspath(os.path.join(
    _HERE, '..', '..', '..', '..', '.claude', 'skills', 'exhibit-slides-pptx', 'scripts')))
from exhibit_pptx import (ExhibitDeck, NAVY, BLUE, BLUE2, BLUE3, BLUE4, TINT, TINT2,
                          CYAN, CORAL, WHITE, MUT, FN, HAIR, HAIR_ROW, SUB_D, W, H)
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

GREEN = RGBColor(0x2E, 0x8B, 0x57)   # interface-mode dot only (API-ready)
AMBER = RGBColor(0xB4, 0x53, 0x09)   # interface-mode dot only (file)

OUT = sys.argv[1] if len(sys.argv) > 1 else os.path.join(_HERE, 'Product_Factory_Execution_Plan_Exhibit.pptx')
d = ExhibitDeck()

ILLUS = 'Illustrative specimen: every figure is a placeholder, re-based on the bank’s own data during the engagement.'


# ── composed helpers (primitives only; chrome untouched) ─────
def chip_card(s, x, y, w, h, head, body, fill=TINT2, accent=BLUE, head_c=NAVY, body_c=NAVY):
    d.rect(s, x, y, w, h, fill=fill, round_=True)
    d.rect(s, x, y + 0.08, 0.045, h - 0.16, fill=accent)
    d.txt(s, x + 0.22, y + 0.10, w - 0.40, 0.26, head, size=12, color=head_c, bold=True)
    d.txt(s, x + 0.22, y + 0.38, w - 0.40, h - 0.46, body, size=10.5, color=body_c, line_sp=1.12)


def hero(s, x, y, w, h, eyebrow, big, body):
    d.rect(s, x, y, w, h, fill=NAVY, round_=True)
    d.txt(s, x + 0.28, y + 0.22, w - 0.56, 0.24, eyebrow.upper(), size=10.5, color=CYAN, bold=True)
    d.txt(s, x + 0.28, y + 0.52, w - 0.56, 1.0, big, size=15.5, color=WHITE, bold=True, line_sp=1.12)
    d.txt(s, x + 0.28, y + h - 1.15, w - 0.56, 1.05, body, size=10.5, color=SUB_D, line_sp=1.2)


def col_card(s, x, y, w, h, head, bullets):
    d.rect(s, x, y, w, h, fill=TINT2, round_=True)
    d.rect(s, x, y + 0.08, 0.045, h - 0.16, fill=BLUE)
    d.txt(s, x + 0.20, y + 0.12, w - 0.36, 0.24, head.upper(), size=10, color=BLUE, bold=True)
    d.txt(s, x + 0.20, y + 0.40, w - 0.36, h - 0.5,
          [[('· ' + b, 10, NAVY, False)] for b in bullets], line_sp=1.16, sp_after=4)


def panel(s, x, y, w, h, eyebrow):
    d.rect(s, x, y, w, h, fill=WHITE, line=HAIR, line_w=0.9, round_=True)
    d.txt(s, x + 0.18, y + 0.11, w - 0.36, 0.22, eyebrow.upper(), size=9.5, color=BLUE, bold=True)
    return x + 0.18, y + 0.40, w - 0.36


def deep_dive(kicker, title, activities, entails, outputs, addup_lead, addup_rest, footnote, notes=None):
    s = d.slide()
    d.chrome(s, kicker, title)
    col_card(s, 1.0, 1.95, 4.60, 3.45, 'The activities, week by week', activities)
    col_card(s, 5.80, 1.95, 3.40, 3.45, 'What it entails', entails)
    col_card(s, 9.40, 1.95, 3.30, 3.45, 'What comes out', outputs)
    d.takeaway_band(s, addup_lead, addup_rest, y=5.62)
    d.footnote(s, footnote)
    if notes:
        d.notes(s, notes)
    return s


def product_plan(kicker, title, milestones, stats, footnote, notes=None):
    s = d.slide()
    d.chrome(s, kicker, title)
    n = len(milestones)
    x0, x1, line_y = 1.95, 11.75, 2.95
    d.hline(s, x0, line_y, x1, line_y, color=HAIR, wpt=1.6)
    span = (x1 - x0) / (n - 1)
    for i, (date, label, body, gate) in enumerate(milestones):
        cx = x0 + i * span
        if gate == 'decision':
            d.diamond(s, cx - 0.11, line_y - 0.11, 0.22, 0.22, fill=CORAL, line=WHITE, line_w=1.0)
        elif gate == 'ga':
            d.diamond(s, cx - 0.11, line_y - 0.11, 0.22, 0.22, fill=BLUE, line=WHITE, line_w=1.0)
        else:
            d.oval(s, cx - 0.085, line_y - 0.085, 0.17, 0.17, BLUE)
    for i, (date, label, body, gate) in enumerate(milestones):
        cx = x0 + i * span
        d.txt(s, cx - 1.0, line_y - 0.45, 2.0, 0.24, date, size=10.5, color=BLUE, bold=True,
              align=PP_ALIGN.CENTER)
        d.txt(s, cx - 1.0, line_y + 0.24, 2.0, 0.26, label, size=11.5, color=NAVY, bold=True,
              align=PP_ALIGN.CENTER)
        d.txt(s, cx - 1.0, line_y + 0.52, 2.0, 1.0, body, size=9, color=MUT,
              align=PP_ALIGN.CENTER, line_sp=1.12)
    band_y = 4.75
    d.rect(s, 1.0, band_y, 11.708, 1.0, fill=NAVY, round_=True)
    cw = 11.708 / len(stats)
    for i, (lab, val) in enumerate(stats):
        bx = 1.0 + i * cw
        d.txt(s, bx + 0.30, band_y + 0.16, cw - 0.45, 0.22, lab.upper(), size=9.5, color=CYAN, bold=True)
        d.txt(s, bx + 0.30, band_y + 0.44, cw - 0.45, 0.48, val, size=12.5, color=WHITE, bold=True,
              line_sp=1.05)
    d.footnote(s, footnote)
    if notes:
        d.notes(s, notes)
    return s


# ═════════ S1 · COVER (dark) ═════════
s = d.slide(dark=True)
d.step_glyph(s, 0.406, 0.406, 0.167, 0.166, CYAN)
d.txt(s, 1.0, 2.0, 11.0, 0.35, 'VALUE CONSULTING · PDP · INTERNAL', size=12.5, color=CYAN, bold=True)
d.txt(s, 1.0, 2.45, 11.0, 1.9, [[('Product Factory:', 44, WHITE, False)],
                                [('from concept to revenue.', 44, WHITE, True)]], line_sp=1.05)
d.txt(s, 1.0, 4.35, 9.4, 0.9, 'Four paid product wedges that pre-install Banking OS: prove the model in '
                              'year one, reach cost neutrality in year two, open a recurring line.',
      size=15, color=SUB_D, line_sp=1.25)
d.txt(s, 1.0, 6.35, 10.0, 0.3, 'Shyam · July 2026 · for the talent programme (Tim Ruttner) and the PDP track (Mayur)',
      size=11, color=MUT)
d.txt(s, 11.7, 7.05, 1.06, 0.3, 'Backbase', size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
d.notes(s, 'COVER. One line and move: the function stops being a cost centre; here is the twelve-month execution.')

# ═════════ EXEC CHAPTER · THE SHORT VERSION (for Tim) ═════════
d.divider('00', 'The short version',
          'Six slides for a ten-minute read. The full playbook behind them starts at chapter 01.')

# E1 · the play on a page
s = d.slide()
d.chrome(s, 'The short version · the play', 'We sell paid discovery that leaves Banking OS installed')
ladder_mini = [
    ('Rung 0 · free', 'Ignite Inspire: ~4 meetings surface the evidence and the wedge choice.'),
    ('Rung 1 · €60K wedge', 'A product diagnostic on the bank’s own data: X-Ray, Cartographer, Guardrail or Telemetry.'),
    ('Rung 2 · €60K Mission POC', 'The top candidate goes live in 6-12 weeks via the Factory. Wedge fee credits in.'),
    ('Rung 3 · recurring', 'Value Assurance: quarterly cost-per-outcome proof that funds the next domain.'),
]
for i, (h_, b_) in enumerate(ladder_mini):
    chip_card(s, 1.0, 1.95 + i * 0.80, 5.9, 0.70, h_, b_)
hero(s, 7.15, 1.95, 5.55, 3.10, 'Two engines, one wedge',
     'Products install the technical layers; services install the human layer.',
     'Engine A (four product SKUs) pre-binds Nexus, Sentinel and Missions. Engine B (AI-native operating '
     'model) redesigns the org that runs them. B never sells cold.')
d.takeaway_band(s, 'The one-liner: ', 'paid discovery that leaves the platform installed.')
d.footnote(s, 'Full detail: chapters 01-02. Canon: knowledge/product/banking-os.md.')
d.notes(s, 'TIM OPENER. One breath: we stop selling reports; every paid engagement installs a component of Banking OS and pre-writes the next purchase.')

# E2 · the line-up
s = d.slide()
d.chrome(s, 'The short version · the line-up', 'Five offers, one motion: each installs part of Banking OS')
lineup = [
    ('Process X-Ray · €50-75K', 'Prices the leakage in real work; installs the Mission backlog + integration map. First revenue Jan 2027.'),
    ('Value Telemetry · €15-50K + share', 'Cost per resolved outcome on live AI; the recurring line. First revenue Jan 2027.'),
    ('Guardrail Studio · €50-100K', 'Turns policy into executable rules; installs Sentinel readiness. First revenue Jun 2027.'),
    ('Ontology Cartographer · €75-100K', 'Maps scattered data to shared truth; installs the Nexus blueprint. First revenue Aug 2027.'),
    ('AI-Native Org Design · €80-150K', 'Engine B flagship: the inverted-T target org, sold downstream of a wedge. First revenue May 2027.'),
]
for i, (h_, b_) in enumerate(lineup):
    chip_card(s, 1.0, 1.95 + i * 0.80, 11.708, 0.70, h_, b_)
d.footnote(s, 'Tickets are price points to pressure-test (chapter 05 assumptions). Every wedge ships all four lenses; the SKU sets which one goes deep.')
d.notes(s, 'Each offer passes the three-part filter: earns a fee, harvests proprietary evidence, pre-binds a Banking OS layer.')

# E3 · the numbers
s = d.slide()
d.chrome(s, 'The short version · the numbers', 'A proof year, then a cost-neutral line with recurring revenue')
d.txt(s, 1.0, 2.05, 5.0, 0.24, 'WHAT THE PLAN MOVES', size=10, color=MUT, bold=True, track='100')
cards = [('Year one', '3-4', 'paid installs, one per product family (€0.4-0.5M)'),
         ('Year two', '10+', 'installs a year: the cost-neutral rate (~€1.2M covered)'),
         ('First revenue', 'Jan 27', 'X-Ray and Telemetry first signatures'),
         ('Recurring', '€30-50K', 'a year per account, plus 15-35% gain share')]
cx = 1.0
for f_, v_, l_ in cards:
    d.stat_card(s, cx, 2.40, f_, v_, l_)
    cx += 2.975
d.rect(s, 1.0, 4.10, 5.6, 0.82, fill=NAVY, round_=True)
d.txt(s, 1.25, 4.10, 5.1, 0.82, 'Cost neutral in year two; beyond dot ten it is a revenue line',
      size=14, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, line_sp=1.12)
d.txt(s, 6.95, 4.14, 5.75, 0.85,
      [[('Capacity-checked: 1.5 FTE pod caps year one at 4-5 installs', 11, NAVY, False)],
       [('Year-two scale is funded by year-one revenue, delivering from kits', 11, NAVY, False)]],
      sp_after=6, line_sp=1.15)
d.footnote(s, 'Team cost ~€1.2M is the PDP working figure, owned by finance to validate. Full math: chapter 06, honest critique.')
d.notes(s, 'DEFENSE: the year-one number is deliberately conservative (capacity math, chapter 06). Overpromising year one burns the line’s credibility.')

# E4 · how banks buy it
s = d.slide()
d.chrome(s, 'The short version · the route to yes', 'Banks buy it by nomination, with the fee crediting forward')
buy3 = [
    ('Early Access nomination', '8 slots remain in the mid-year cohort. Nominated against program criteria; nothing for procurement to tender.'),
    ('The fee credits forward, 100%', 'On conversion the wedge fee credits into the Mission POC. The budget is never spent on a maybe.'),
    ('The bank owns the evidence', 'Vendor-neutral artifacts, a 10×-value-or-free floor, and the specimen pack shown before any signature.'),
]
for i, (h_, b_) in enumerate(buy3):
    chip_card(s, 1.0, 1.95 + i * 1.02, 5.9, 0.90, h_, b_)
hero(s, 7.15, 1.95, 5.55, 3.00, 'Nominated, never tendered',
     'Co-creation is the answer to “why you, you are new”.',
     'Roadmap influence, engineering at the table, joint proof. Scarcity is honest: when the slots are '
     'gone, the cohort closes.')
d.open_badge(s, 7.15, 5.15, 5.55, 'Open · ', 'EAP eligibility (new logos vs existing customers) to confirm with the program owner.', h=0.62)
d.footnote(s, 'Full buying-criteria and pre-RFP mechanics: chapter 06.')
d.notes(s, 'For Tim: this is also the category story — co-innovation with named design partners, not another vendor diagnostic.')

# E5 · who needs it now
s = d.slide()
d.chrome(s, 'The short version · who needs it now', 'Must-have accounts show the pressure and the stated ambition')
qx, qy, qw, qh = 1.0, 2.05, 6.3, 3.3
d.rect(s, qx, qy, qw / 2, qh / 2, fill=TINT2)
d.rect(s, qx + qw / 2, qy, qw / 2, qh / 2, fill=TINT)
d.rect(s, qx, qy + qh / 2, qw / 2, qh / 2, fill=TINT2)
d.rect(s, qx + qw / 2, qy + qh / 2, qw / 2, qh / 2, fill=TINT2)
d.txt(s, qx + qw / 2 + 0.15, qy + 0.10, qw / 2 - 0.3, 0.22, 'MUST-HAVE NOW', size=10, color=BLUE2, bold=True)
d.txt(s, qx + qw / 2 + 0.15, qy + 0.34, qw / 2 - 0.3, 1.0, 'High cost pressure + stated AI ambition. Lead with X-Ray + EAP nomination.', size=9.5, color=NAVY, line_sp=1.15)
d.txt(s, qx + 0.15, qy + 0.10, qw / 2 - 0.3, 0.22, 'MUST-HAVE SOON', size=10, color=NAVY, bold=True)
d.txt(s, qx + 0.15, qy + 0.34, qw / 2 - 0.3, 1.0, 'High pressure, ambition not yet articulated. Sell the ambition with the free wedge.', size=9.5, color=NAVY, line_sp=1.15)
d.txt(s, qx + qw / 2 + 0.15, qy + qh / 2 + 0.10, qw / 2 - 0.3, 0.22, 'GOOD-TO-HAVE', size=10, color=NAVY, bold=True)
d.txt(s, qx + qw / 2 + 0.15, qy + qh / 2 + 0.34, qw / 2 - 0.3, 1.0, 'Ambition without pressure: innovation budgets, ideal EAP candidates.', size=9.5, color=NAVY, line_sp=1.15)
d.txt(s, qx + 0.15, qy + qh / 2 + 0.10, qw / 2 - 0.3, 0.22, 'LATER', size=10, color=MUT, bold=True)
d.txt(s, qx + 0.15, qy + qh / 2 + 0.34, qw / 2 - 0.3, 1.0, 'Neither yet visible. Keep in the free-wedge funnel.', size=9.5, color=MUT, line_sp=1.15)
d.txt(s, qx, qy + qh + 0.10, qw, 0.22, 'STATED AI / DIGITAL AMBITION →', size=9, color=MUT, bold=True, align=PP_ALIGN.CENTER)
d.txt(s, 1.0, 1.75, 4.0, 0.22, '↑ STRATEGIC COST PRESSURE', size=9, color=MUT, bold=True)
d.rect(s, 7.75, 2.05, 4.95, 3.3, fill=TINT2, round_=True)
d.txt(s, 7.95, 2.22, 4.5, 0.22, 'SIGNALS TO READ (5-YEAR STRATEGY)', size=9.5, color=BLUE, bold=True)
for i, sig in enumerate(['Cost-income guidance and named cost programs', 'AI or agentic ambition stated in the strategy',
                         'Branch and headcount reduction targets', 'Digital attacker pressure in home market',
                         'M&A integration or regulatory remediation load']):
    d.txt(s, 7.95, 2.52 + i * 0.34, 4.55, 0.3, '· ' + sig, size=10, color=NAVY, line_sp=1.1)
d.takeaway_band(s, 'Week one with the AEs: ', 'classify the account list against these signals.')
d.footnote(s, 'Archetypes, never named accounts on slides; classify real accounts from annual-report signals with the AE. Full version: chapter 04.')
d.notes(s, 'Must-have = the strategy already needs this to be true (cost program + AI ambition). Good-to-have = innovation-led, perfect EAP design partners.')

# E6 · the asks
s = d.slide()
d.chrome(s, 'The short version · the asks', 'Three decisions this week put product one in market')
asks3 = [
    ('1 · The Mayur 1:1', 'Quota treatment, pilot green-light, the named SE at 0.5 FTE, and the five open decisions of chapter 06.'),
    ('2 · EAP eligibility', 'Confirm with the program owner whether the 8 slots take new logos, existing customers, or both; shortlist nominations.'),
    ('3 · The Tim track', 'Category story and launch collateral with marketing; the €5,000 talent budget becomes seed capital for Process X-Ray.'),
]
for i, (h_, b_) in enumerate(asks3):
    chip_card(s, 1.0, 1.95 + i * 1.02, 11.708, 0.90, h_, b_)
d.takeaway_band(s, 'Everything else in this deck ', 'exists to earn these three decisions.', y=5.30)
d.footnote(s, 'The full ask set and the honest critique behind the revised targets: chapters 05-06.')
d.notes(s, 'End the ten minutes here. If Tim asks for depth on any slide, the playbook chapter is the backup.')

# ═════════ S2 · DIVIDER 01 ═════════
d.divider('01', 'The frameworks',
          'The thesis, the filter, the economics, the ladder and lifecycle, and the pod that runs it.')

# ═════════ S3 · THESIS ═════════
s = d.slide()
d.chrome(s, 'Frameworks · the thesis', 'Every product earns a fee, harvests reality, installs the OS')
tests = [
    ('Test 1 · It earns a fee', 'Each install is paid work at €15-100K, priced as a product with a fixed scope and a delivery kit.'),
    ('Test 2 · It harvests proprietary data', 'The engagement captures the bank’s real processes, data structures and policies. Evidence only Backbase holds.'),
    ('Test 3 · It pre-binds a Banking OS layer', 'The output is a working component: Nexus, Sentinel or a Mission candidate, already present in the account.'),
]
for i, (h_, b_) in enumerate(tests):
    chip_card(s, 1.0, 1.95 + i * 1.24, 5.9, 1.10, h_, b_)
hero(s, 7.15, 1.95, 5.55, 3.55, 'The one-liner',
     'Paid discovery that leaves the platform installed.',
     'Each engagement ends with a live Banking OS component in the account, a quantified expansion case '
     'for the AE, and account intelligence competitors never see.')
d.footnote(s, 'Filter rule: a candidate that fails any one test is consulting and stays out of the line. Aligned to the Banking OS canon (banking-os.md).')
d.notes(s, 'The three-part design filter from the 29 Jun PDP session. Product-proximate is the moat; customer-proximate is space we lose.')

# ═════════ S4 · TWO ENGINES ═════════
s = d.slide()
d.chrome(s, 'Frameworks · the two engines', 'The products harvest evidence; the services redesign the bank')
d.rect(s, 1.0, 1.95, 5.75, 3.15, fill=TINT, round_=True)
d.txt(s, 1.26, 2.13, 5.2, 0.24, 'ENGINE A · PRODUCT FACTORY', size=10.5, color=BLUE2, bold=True)
d.txt(s, 1.26, 2.40, 5.2, 0.30, 'Installs the technical layers', size=14.5, color=NAVY, bold=True)
for i, line in enumerate(['Four SKUs: X-Ray, Cartographer, Guardrail, Telemetry',
                          'Pre-binds Nexus, Sentinel and Mission candidates',
                          'Earns €15-100K per install; fees credit forward',
                          'Harvests the account evidence nobody else holds']):
    d.txt(s, 1.26, 2.80 + i * 0.32, 5.3, 0.28, '· ' + line, size=10.5, color=NAVY, line_sp=1.1)
d.txt(s, 1.26, 4.72, 5.2, 0.24, '→ the evidence engine', size=10, color=BLUE, bold=True)
d.rect(s, 6.95, 1.95, 5.75, 3.15, fill=NAVY, round_=True)
d.txt(s, 7.21, 2.13, 5.2, 0.24, 'ENGINE B · AI-NATIVE SERVICES', size=10.5, color=CYAN, bold=True)
d.txt(s, 7.21, 2.40, 5.2, 0.30, 'Installs the human layer', size=14.5, color=WHITE, bold=True)
for i, line in enumerate(['Maturity Assessment, Workforce Optimization, and the flagship AI-Native Org Design (v0.1 exists)',
                          'Redesigns the org chart on the A1-A5 autonomy curve: pyramid to inverted-T',
                          'Pre-binds the roles that run Banking OS: Mission Owners, AgentOps, Exception Desks',
                          'Earns €80-150K per engagement (to validate)']):
    d.txt(s, 7.21, 2.80 + i * 0.47, 5.3, 0.44, '· ' + line, size=10.5, color=SUB_D, line_sp=1.12)
d.txt(s, 7.21, 4.72, 5.2, 0.24, '→ the operating-model engine', size=10, color=CYAN, bold=True)
d.takeaway_band(s, 'Sequencing rule: ', 'Engine B never sells cold; it runs on Engine A’s harvested evidence.')
d.footnote(s, 'Engine B passes the same three tests when fed by Engine A. Ticket range is an assumption to validate against market org-design pricing.')
d.notes(s, 'PDP workstream 4. The moat: generic consultancies can run org design; only we can run it on the bank’s own autonomy evidence (X-Ray scores, Guardrail authority maps, telemetry).')

# ═════════ S5 · ECONOMICS DOT GRID ═════════
s = d.slide()
d.chrome(s, 'Frameworks · the economics', 'Ten installs cover the team cost; every one beyond is revenue')
d.txt(s, 1.0, 2.05, 6.4, 0.28, 'One dot = one paid install a year', size=11, color=MUT)
dd, gap = 0.44, 0.31
for i in range(16):
    row, col = divmod(i, 8)
    d.oval(s, 1.0 + col * (dd + gap), 2.45 + row * (dd + gap), dd, dd, BLUE if i < 10 else TINT)
d.oval(s, 1.0, 4.30, 0.16, 0.16, BLUE)
d.txt(s, 1.24, 4.26, 3.2, 0.24, 'Installs to cost neutrality (10)', size=10.5, color=MUT)
d.oval(s, 4.55, 4.30, 0.16, 0.16, TINT)
d.txt(s, 4.79, 4.26, 3.4, 0.24, 'Year-two expansion range (14-16)', size=10.5, color=MUT)
d.rect(s, 8.05, 2.15, 4.67, 2.35, fill=TINT2, round_=True)
d.txt(s, 8.40, 2.40, 2.0, 0.75, '10', size=48, color=BLUE, bold=True)
d.txt(s, 8.40, 3.30, 4.0, 1.0, 'installs a year at an average package of €120K cover the ~€1.2M annual team cost.',
      size=12.5, color=NAVY, line_sp=1.22)
d.takeaway_band(s, 'Cost neutrality is the floor: ', 'the objective is the revenue line beyond dot ten.')
d.footnote(s, 'Package = Assess & Solution €60K + Mission POC €60K. Team cost ~€1.2M is the PDP working figure, to validate with finance. Chapter 04 phases year one at 3-4 installs.')
d.notes(s, 'DOT GRID. DEFENSE: the €1.2M is a working figure, owned by finance to confirm; the 10-install bar moves with it. Year-one phasing is in the critique chapter.')

# ═════════ S6 · LADDER (segmented bar) ═════════
s = d.slide()
d.chrome(s, 'Frameworks · the ladder', 'The wedge stays free; depth and recurrence carry the price')
chip_card(s, 1.0, 2.35, 2.50, 1.40, 'Rung 0 · free wedge',
          'Ignite Inspire: ~4 meetings, benchmark and use-case shortlist. Grows the funnel.')
d.txt(s, 3.58, 2.85, 0.3, 0.35, '→', size=16, color=MUT, align=PP_ALIGN.CENTER)
d.txt(s, 3.95, 2.05, 5.8, 0.24, 'RUNGS 1 + 2 · THE €120K PACKAGE (TO SCALE)', size=10, color=MUT, bold=True)
d.rect(s, 3.95, 2.35, 2.90, 0.80, fill=BLUE)
d.txt(s, 4.13, 2.47, 2.6, 0.56, [[('Assess & Solution', 11.5, WHITE, True)], [('€60K', 11.5, WHITE, False)]], line_sp=1.1)
d.rect(s, 6.85, 2.35, 2.90, 0.80, fill=BLUE2)
d.txt(s, 7.03, 2.47, 2.6, 0.56, [[('Mission POC', 11.5, WHITE, True)], [('€60K', 11.5, WHITE, False)]], line_sp=1.1)
d.txt(s, 9.83, 2.85, 0.3, 0.35, '→', size=16, color=MUT, align=PP_ALIGN.CENTER)
d.rect(s, 10.20, 2.35, 2.50, 1.40, fill=NAVY, round_=True)
d.txt(s, 10.42, 2.47, 2.1, 0.24, 'RUNG 3 · RECURRING', size=9.5, color=CYAN, bold=True)
d.txt(s, 10.42, 2.74, 2.1, 0.9, 'Value Assurance\n€30-50K a year\n+ gain share', size=11, color=WHITE, bold=True, line_sp=1.15)
chip_card(s, 1.0, 4.35, 5.75, 0.95, 'The free wedge feeds the funnel',
          'Four meetings surface the leakage evidence that qualifies the paid wedge. Zero friction to start.')
chip_card(s, 6.95, 4.35, 5.75, 0.95, 'Recurrence is the growth flywheel',
          'Telemetry converts every live deployment into quarterly revenue that compounds with AI consumption.')
d.footnote(s, 'Recurring range anchored on SAP and ServiceNow success plans (10-30% of licence) and FinOps gain-share practice (15-35% of realised savings).')
d.notes(s, 'T03 segmented bar, to scale on the €120K package. Rung 0 free = funnel; rung 3 = the flywheel.')

# ═════════ S7 · LADDER ↔ LIFECYCLE BRIDGE ═════════
s = d.slide()
d.chrome(s, 'Frameworks · ladder meets lifecycle', 'The bank buys the ladder; the pod delivers it as one lifecycle')
rungs = [
    ('RUNG 0 · FREE', 'Ignite Inspire', 'Hypotheses, wedge choice, specimen pack', TINT2, NAVY, MUT),
    ('RUNG 1 · €60K', 'The wedge product', 'X-Ray, Cartographer, Guardrail or Telemetry on the bank’s data', TINT, NAVY, BLUE2),
    ('RUNG 2 · €60K', 'Mission POC', 'Factory Mission Sprint: the top candidate live in 6-12 weeks', BLUE, WHITE, WHITE),
    ('RUNG 3 · RECURRING', 'Value Assurance', 'Quarterly cost-per-outcome proof funds the next domain', NAVY, WHITE, CYAN),
]
rw, rgap, ry = 2.62, 0.40, 2.18
d.txt(s, 1.0, 1.92, 4.0, 0.22, 'WHAT THE BANK BUYS', size=9.5, color=MUT, bold=True)
for i, (tag, name, desc, fill, tc, tagc) in enumerate(rungs):
    x = 1.0 + i * (rw + rgap)
    d.rect(s, x, ry, rw, 1.30, fill=fill, round_=True)
    d.txt(s, x + 0.18, ry + 0.10, rw - 0.32, 0.20, tag, size=9, color=tagc, bold=True)
    d.txt(s, x + 0.18, ry + 0.32, rw - 0.32, 0.26, name, size=12.5, color=tc, bold=True)
    d.txt(s, x + 0.18, ry + 0.62, rw - 0.32, 0.62, desc, size=9,
          color=tc if fill in (BLUE, NAVY) else NAVY, line_sp=1.12)
    if i < 3:
        ax = x + rw + 0.02
        d.txt(s, ax, ry + 0.38, rgap - 0.04, 0.3, '→', size=13, color=MUT, align=PP_ALIGN.CENTER)
        d.txt(s, ax - 0.06, ry + 0.72, rgap + 0.08, 0.24, '05 · 06', size=8, color=BLUE,
              bold=True, align=PP_ALIGN.CENTER)
d.txt(s, 1.0, 3.86, 4.0, 0.22, 'HOW THE POD DELIVERS', size=9.5, color=MUT, bold=True)
stage_groups = ['01 Concept · 02 Design', '03 Prototype · 04 Proof', '07 Install', 'Run · measure · expand']
for i, sg in enumerate(stage_groups):
    x = 1.0 + i * (rw + rgap)
    d.hline(s, x + rw / 2, ry + 1.30, x + rw / 2, 4.12, color=HAIR)
    d.rect(s, x, 4.12, rw, 0.50, fill=WHITE, line=BLUE, line_w=1.1, round_=True)
    d.txt(s, x + 0.1, 4.12, rw - 0.2, 0.50, sg, size=10.5, color=BLUE2, bold=True,
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
d.takeaway_band(s, 'Read it both ways: ', 'each rung is a paid lifecycle stage, and each exit artifact is the next purchase’s entry contract.')
d.footnote(s, 'Stages 05 Proposal and 06 Signature are the arrows between rungs: each rung ends by proposing and signing the next one.')
d.notes(s, 'THE BRIDGE SLIDE. Vertically: rung = lifecycle stage the bank pays for. Horizontally: exit artifact = entry contract. The account never buys a maybe.')

# ═════════ S8 · LIFECYCLE STEP FLOW ═════════
s = d.slide()
d.chrome(s, 'Frameworks · the lifecycle', 'One motion takes each product from concept to install')
steps = [
    ('01', 'Concept', 'VC · 1 wk'), ('02', 'Design', 'VC + SE · 1-2 wks'),
    ('03', 'Prototype', 'SE · 2-3 wks'), ('04', 'Proof', 'VC · 2-4 wks'),
    ('05', 'Proposal', 'VC + AE · 1 wk'), ('06', 'Signature', 'AE + VC · 2-4 wks'),
    ('07', 'Install', 'FDE + SE · 6-12 wks'),
]
sw, sgap = 1.585, 0.10
for i, (n, name, who) in enumerate(steps):
    x = 1.0 + i * (sw + sgap)
    d.rect(s, x, 2.10, sw, 1.55, fill=TINT2, round_=True)
    d.oval(s, x + 0.16, 2.26, 0.36, 0.36, BLUE, label=n, fs=10.5)
    d.txt(s, x + 0.16, 2.74, sw - 0.3, 0.26, name, size=12, color=NAVY, bold=True)
    d.txt(s, x + 0.16, 3.04, sw - 0.3, 0.5, who, size=9, color=MUT, line_sp=1.12)
    if i < 6:
        d.txt(s, x + sw - 0.02, 2.68, 0.18, 0.3, '→', size=12, color=MUT)
d.takeaway_band(s, 'It runs twice over: ', 'the first pass builds the product; every later pass is a 4-8 week cycle per install.', y=4.15)
d.footnote(s, 'Durations are per-engagement once the product exists, ±30% tolerance. The first pass through stages 1-4 is the one-time build.')
d.notes(s, 'Exit gates: Concept = qualified sponsor. Design = signed SoW + data access. Prototype = demo on client data. '
           'Proof = CFO-grade evidence. Proposal = offer + named Mission. Signature = order form. Install = layer live + expansion roadmap.')

# ═════════ S9 · ARTIFACT CHAIN + GATES ═════════
s = d.slide()
d.chrome(s, 'Frameworks · connecting the dots', 'Each deliverable is the entry contract of the next purchase')
arts = [
    ('AFTER INSPIRE (FREE)', 'The Inspire pack', ['Leakage hypotheses', 'Wedge choice + scope', 'Specimen output pack']),
    ('AFTER THE WEDGE', 'The evidence pack', ['€ leakage, evidenced', 'Integration landscape map', 'Ranked Mission backlog']),
    ('AFTER THE POC', 'The live loop', ['Resolution loop in production', 'Telemetry baseline', 'Sentinel-governed, auditable']),
    ('EVERY QUARTER', 'The value proof', ['Cost per outcome, measured', 'Savings register', 'Next Mission ranked + funded']),
]
aw_, agap, ay = 2.62, 0.40, 1.95
for i, (eyebrow, name, bullets) in enumerate(arts):
    x = 1.0 + i * (aw_ + agap)
    d.rect(s, x, ay, aw_, 2.05, fill=TINT2, round_=True)
    d.txt(s, x + 0.18, ay + 0.12, aw_ - 0.32, 0.2, eyebrow, size=8.5, color=BLUE, bold=True)
    d.txt(s, x + 0.18, ay + 0.36, aw_ - 0.32, 0.26, name, size=12.5, color=NAVY, bold=True)
    d.txt(s, x + 0.18, ay + 0.68, aw_ - 0.32, 1.3,
          [[('· ' + b, 9.5, NAVY, False)] for b in bullets], line_sp=1.2, sp_after=3)
    if i < 3:
        gx_ = x + aw_ + agap / 2
        d.diamond(s, gx_ - 0.11, ay + 0.85, 0.22, 0.22, fill=BLUE, line=WHITE, line_w=1.0)
        d.txt(s, gx_ - 0.2, ay + 1.10, 0.4, 0.2, f'G{i + 1}', size=9, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
gates = [
    ('G1', 'Sponsor signs the wedge SoW, inside their own delegated authority.'),
    ('G2', 'ExCo funds Mission #1; the wedge fee credits into the POC.'),
    ('G3', 'Measured value meets the case; the licence conversation opens.'),
    ('G4', 'The next domain is funded; the loop repeats on the same Banking OS.'),
]
for i, (g, t_) in enumerate(gates):
    x = 1.0 + (i % 2) * 5.95
    y = 4.35 + (i // 2) * 0.52
    d.diamond(s, x, y + 0.05, 0.18, 0.18, fill=BLUE)
    d.txt(s, x + 0.28, y, 0.4, 0.24, g, size=10.5, color=BLUE, bold=True)
    d.txt(s, x + 0.62, y, 5.1, 0.44, t_, size=10.5, color=NAVY, line_sp=1.1)
d.takeaway_band(s, 'Nothing is a shelf report: ', 'every artifact is authored as the input of the next purchase order.')
d.footnote(s, 'G1-G4 are the account-level gates. Each product also has its own build gates (SE named, pilot green-light, R&D alignment), shown on the roadmap.')
d.notes(s, 'THE CHAIN. This is what makes the wedge convert instead of inform.')

# ═════════ S10 · THE POD ═════════
s = d.slide()
d.chrome(s, 'Frameworks · resourcing', 'A 2.5-person pod runs the line; hiring follows revenue')
pods = [
    ('You · product owner · 1.0 FTE', 'SKU design, value narrative, proof economics, proposal. Owns the P&L of the line.'),
    ('Solution engineer · 0.5 FTE, named', 'Connectors, data extraction, Factory tooling at prototype and install. The technical spine.'),
    ('FDE · on demand at install', 'Runs the Factory Mission Sprint once an install is signed. 6-12 weeks per Mission.'),
    ('SME · on call', 'Risk and compliance review for Guardrail Studio; R&D liaison for the Cartographer.'),
]
for i, (h_, b_) in enumerate(pods):
    chip_card(s, 1.0, 1.95 + i * 0.92, 5.9, 0.80, h_, b_)
hero(s, 7.15, 1.95, 5.55, 3.60, 'Why the solution engineer matters',
     'One named SE across all four products, so skill compounds.',
     'Every product runs on the same Factory tools: Connector Studio, Semantic Modeler, Simulation & Testing. '
     'Each install makes the next one cheaper and faster.')
d.footnote(s, 'FTEs are steady-state allocations. FDE capacity comes from standard Mission Sprint staffing at install and is charged to the delivery.')
d.notes(s, 'An allocation ask. The SE is the critical unlock; naming one person is ask #1.')

# ═════════ S11 · OBJECTIONS ═════════
s = d.slide()
d.chrome(s, 'Frameworks · the objections', 'The three obvious objections have concrete answers')
qa = [
    ('“Banks will not pay for discovery.”',
     'The market already does: AIB paid ~€80K for a weaker assessment. AWS, SAP and ServiceNow run the diagnostic-wedge motion; ours installs product.'),
    ('“We have no capacity to build products.”',
     'Build effort is 6-12 person-weeks per product because the Factory tools, the APA V3 catalog and the cost model already exist. One product at a time.'),
    ('“This distracts from platform sales.”',
     'Every install pre-binds Nexus, Sentinel or a Mission and hands the AE a quantified expansion case. The wedge is the land motion.'),
]
for i, (c, a) in enumerate(qa):
    y = 2.0 + i * 1.28
    d.rect(s, 1.0, y, 4.35, 1.12, fill=NAVY, round_=True)
    d.txt(s, 1.24, y + 0.16, 3.9, 0.8, c, size=12.5, color=WHITE, bold=True, line_sp=1.15)
    d.rect(s, 5.55, y, 7.15, 1.12, fill=TINT2, round_=True)
    d.rect(s, 5.55, y + 0.08, 0.045, 0.96, fill=BLUE)
    d.txt(s, 5.80, y + 0.14, 6.7, 0.86, a, size=11, color=NAVY, line_sp=1.18)
d.footnote(s, 'AIB reference: ~€80K external assessment (ADS), PDP session record, 29 Jun 2026. Wedge patterns: AWS OLA, SAP Outside-In, ServiceNow Inspire.')
d.notes(s, 'T10 concern-answer. Use verbatim in the Tim and Mayur conversations; every answer is evidence-backed.')

# ═════════ S12 · DIVIDER 02 ═════════
d.divider('02', 'The execution plans',
          'Four products and the services flagship: what each does, its outputs, and the plan to market.')

# ═════════ S13 · ONE WEDGE, FOUR LENSES ═════════
s = d.slide()
d.chrome(s, 'Execution plans · how the four fit', 'A bank buys a Mission, and every wedge surveys all four lenses')
doors = [
    ('Ops pain is the trigger', 'Enter via Process X-Ray. Deep: work + leakage. Thin: ontology, guardrails, telemetry.'),
    ('Data ambition (RI) is the trigger', 'Enter via Ontology Cartographer. Deep: shared truth. Thin: work, guardrails, telemetry.'),
    ('Risk is blocking AI', 'Enter via Guardrail Studio. Deep: authority + policy. Thin: work, ontology, telemetry.'),
    ('AI is live, ROI is doubted', 'Enter via Value Telemetry. Deep: measured value. Thin: work, ontology, guardrails.'),
]
for i, (h_, b_) in enumerate(doors):
    chip_card(s, 1.0, 1.95 + i * 0.92, 5.9, 0.80, h_, b_)
hero(s, 7.15, 1.95, 5.55, 3.60, 'One deep lens, three thin lenses',
     'Every engagement ships all four lenses; the wedge sets the depth.',
     'Each wedge adds an ontology appendix, an authority snapshot and telemetry hooks beside its deep dive. '
     'The account always ends holding a full-stack blueprint.')
d.footnote(s, 'The entry door follows the account’s live buying trigger. All-lens coverage is standard because a Mission needs truth, guardrails and telemetry at once.')
d.notes(s, 'Answer to “why would a bank buy one layer”: it buys a Mission. The wedge only decides which lens goes deep.')

# ═════════ S14-16 · PROCESS X-RAY ═════════
deep_dive(
    'Execution plan 1 of 4 · Process X-Ray', 'Process X-Ray reconstructs real work and prices the leakage',
    ['Wk 1 · scope 2-3 journeys from the APA matrix; agree access and baselines',
     'Wk 1-2 · inventory every system, interface and handoff per journey (the integration landscape)',
     'Wk 2-4 · extract logs and queues; shadow the ops floor (6-8 sessions); rebuild actual vs designed flows',
     'Wk 4-5 · price the leakage per step: time × volume × cost, abandonment × revenue',
     'Wk 5-6 · rank Mission candidates; exec readout with CFO-grade evidence'],
    ['Sponsor: COO or head of operations',
     'Bank effort: SME time ~6-10 hours a week',
     'Lite mode: ops MI packs + sampled cases; no InfoSec cycle, prospect-safe',
     'Instrumented mode: read-only logs, for existing customers',
     'Pod: VC 1.0 + SE 0.5'],
    ['€ p.a. leakage per journey, evidenced',
     'Integration landscape: systems, interfaces, API readiness',
     'Ranked Mission backlog, top 3, with ROI case each',
     'Thin-lens appendices: ontology + authority snapshots',
     'The Rung-2 POC proposal, pre-written'],
    'Adds up: ', 'the systems map scopes Connector Studio before anyone commits.',
    'Lite mode exists because prospects rarely grant system access pre-contract; it trades precision for zero InfoSec friction.',
    notes='The productized Process & Workspace Designer. ~60% of bank work lives between systems; the X-Ray prices that whitespace.')

s = d.slide()
d.chrome(s, 'Process X-Ray · sample outputs', 'The X-Ray hands over the heatmap, systems map and Mission card')
p1x, p1y, p1w = panel(s, 1.0, 1.95, 3.77, 4.15, 'Value-leakage heatmap · €M p.a.')
cols_ = ['Onb', 'Serv', 'Disp', 'Lend']
vals = [('Retail', [0.8, 1.1, 2.4, 1.6]), ('SME', [0.5, 0.9, 1.2, 1.4]), ('Comm', [0.3, 0.6, 0.8, 1.1])]
for j, c in enumerate(cols_):
    d.txt(s, p1x + 0.86 + j * 0.64, p1y + 0.02, 0.62, 0.2, c, size=8.5, color=MUT, bold=True, align=PP_ALIGN.CENTER)
for i, (rname, rv) in enumerate(vals):
    y = p1y + 0.26 + i * 0.44
    d.txt(s, p1x, y + 0.10, 0.8, 0.2, rname, size=9, color=MUT, bold=True)
    for j, v in enumerate(rv):
        fill = CORAL if v >= 2.0 else (BLUE2 if v >= 1.3 else (BLUE if v >= 0.9 else (BLUE4 if v >= 0.6 else TINT)))
        d.rect(s, p1x + 0.86 + j * 0.64, y, 0.62, 0.40, fill=fill)
        d.txt(s, p1x + 0.86 + j * 0.64, y + 0.09, 0.62, 0.2, f'{v}', size=9.5,
              color=WHITE if v >= 0.6 else NAVY, bold=True, align=PP_ALIGN.CENTER)
d.txt(s, p1x, p1y + 1.75, p1w, 0.65, 'Disputes · Retail leads: €2.4M p.a. across 41% manual touches and a 9.2-day cycle.',
      size=9.5, color=NAVY, line_sp=1.2)
d.txt(s, p1x, p1y + 3.35, p1w, 0.24, '→ picks the domain the Mission starts in', size=9, color=BLUE, bold=True)
p2x, p2y, p2w = panel(s, 4.97, 1.95, 3.77, 4.15, 'Integration landscape · dispute journey')
systems = [('Core (T24)', 'API', GREEN), ('CRM (SFDC)', 'API', GREEN), ('Fraud (Falcon)', 'file', AMBER),
           ('Cases (Pega)', 'manual', CORAL), ('Mail (O365)', 'manual', CORAL), ('Pay hub', 'file', AMBER)]
for i, (name, mode, mc) in enumerate(systems):
    x = p2x + (i % 2) * 1.76
    y = p2y + 0.04 + (i // 2) * 0.60
    d.rect(s, x, y, 1.64, 0.50, fill=TINT2, line=HAIR, line_w=0.75, round_=True)
    d.txt(s, x + 0.1, y + 0.05, 1.46, 0.2, name, size=9, color=NAVY, bold=True)
    d.oval(s, x + 0.12, y + 0.31, 0.09, 0.09, mc)
    d.txt(s, x + 0.26, y + 0.26, 1.3, 0.2, mode + ' interface', size=8, color=MUT)
d.txt(s, p2x, p2y + 2.00, p2w, 0.55, '8 systems (+2 not shown) · 3 manual swivels · 2 API-ready · 4 to build via Connector Studio',
      size=9.5, color=NAVY, line_sp=1.2)
d.txt(s, p2x, p2y + 3.35, p2w, 0.24, '→ scopes the integration work before commitment', size=9, color=BLUE, bold=True)
d.rect(s, 8.94, 1.95, 3.77, 4.15, fill=NAVY, round_=True)
d.txt(s, 9.12, 2.10, 3.4, 0.2, 'MISSION CANDIDATE #1', size=9.5, color=CYAN, bold=True)
d.txt(s, 9.12, 2.36, 3.4, 0.55, 'Card dispute resolution loop', size=14, color=WHITE, bold=True, line_sp=1.1)
for i, line in enumerate(['Value: €1.9M p.a., evidenced', 'Autonomy target: A4, execute by exception',
                          'Feasibility: high · 12 integration points', '8 live via APIs, 4 built in the Sprint',
                          'Mission length: 8 weeks']):
    d.txt(s, 9.12, 3.02 + i * 0.32, 3.4, 0.28, '· ' + line, size=10, color=SUB_D, line_sp=1.1)
d.txt(s, 9.12, 5.30, 3.4, 0.5, '→ becomes the Rung-2 POC scope, fee credited', size=9.5, color=CYAN, bold=True, line_sp=1.15)
d.footnote(s, ILLUS)
d.notes(s, 'The specimen pack: shown in the free wedge so the bank buys a known object. DEFENSE: all figures placeholders until re-based on client data.')

product_plan(
    'Execution plan 1 of 4 · Process X-Ray', 'X-Ray ships first: proof by December, revenue in January',
    [
        ('Aug 2026', 'Design', 'Delivery kit, pricing one-pager, demo script. 3 pw, VC.', None),
        ('Sep-Oct 2026', 'Prototype', 'Demo on the APA V3 catalog (65 processes) plus one friendly dataset. 5 pw.', None),
        ('Oct 2026', 'Pilot gate', 'Prospect pilot green-light with Mayur. AIB-profile account.', 'decision'),
        ('Oct-Dec 2026', 'Proof', 'Paid pilot: leakage heatmap, top-3 Mission candidates with ROI. 6 pw.', None),
        ('Jan 2027', 'First signature', 'Wedge sold at €50-75K; install starts as a Mission Sprint.', 'ga'),
    ],
    [('Build effort', '8-10 person-wks'), ('Delivery per install', '4-6 wks · 12-14 pw'),
     ('Pod', 'VC 1.0 + SE 0.5'), ('Pre-binds', 'Mission + ROI')],
    'Effort ±30%. Build is low-risk: the APA V3 catalog (65 processes, 283 steps), journey builder and leakage method already exist.',
    notes='Coral diamond = decision Mayur owns (pilot green-light). Blue diamond = first revenue.')

# ═════════ S17-19 · ONTOLOGY CARTOGRAPHER ═════════
deep_dive(
    'Execution plan 2 of 4 · Ontology Cartographer', 'The Cartographer maps scattered data to one shared truth',
    ['Wk 1-2 · inventory the data landscape: core, CRM, LMS, credit, workflow',
     'Wk 1-2 · harvest schemas and data dictionaries (metadata only)',
     'Wk 3-5 · map entities to the banking ontology; surface overlaps and conflicts',
     'Wk 5-7 · score truth-gaps per entity; price them in ops time and error cost',
     'Wk 7-8 · phased Nexus binding blueprint with a precise entity scope'],
    ['Sponsor: CDO or CIO',
     'Bank effort: data stewards 4-6 hours a week',
     'Access: schema metadata only; customer data never leaves the bank',
     'R&D liaison: Semantic Modeler and Nexus teams',
     'Pod: SE 1.0 + VC 0.5 (SE-led)'],
    ['Ontology coverage map, entity by entity',
     'Truth-gap report, priced in ops time',
     'Data-source integration inventory',
     'Phased Nexus blueprint: what binds first, and why',
     'Thin-lens appendices: leakage + authority snapshots'],
    'Adds up: ', 'shared truth is the precondition for RI and any A4-plus Mission.',
    'Metadata-only access keeps the InfoSec review at schema level. Nexus is the system of truth; systems of record stay where they are.',
    notes='Productizes the Semantic Modeler. Ontology is key to every use case, so the lens ships thin in every wedge; this SKU is the deep version.')

s = d.slide()
d.chrome(s, 'Ontology Cartographer · sample outputs', 'The Cartographer shows where the bank’s truth fragments')
p1x, p1y, p1w = panel(s, 1.0, 1.95, 3.77, 4.15, 'Ontology coverage map · by entity')
ents = [('Customer', '6 sources · 4 conflicts', CORAL), ('Account', '3 sources · aligned', GREEN),
        ('Product', '4 sources · 2 conflicts', AMBER), ('Transaction', '2 sources · aligned', GREEN),
        ('Collateral', 'fragmented', AMBER), ('Case', 'absent as an entity', CORAL)]
for i, (e, s_, c) in enumerate(ents):
    x = p1x + (i % 2) * 1.76
    y = p1y + 0.04 + (i // 2) * 0.64
    d.rect(s, x, y, 1.64, 0.54, fill=TINT2, round_=True)
    d.rect(s, x, y + 0.06, 0.04, 0.42, fill=c)
    d.txt(s, x + 0.14, y + 0.06, 1.44, 0.2, e, size=10, color=NAVY, bold=True)
    d.txt(s, x + 0.14, y + 0.28, 1.44, 0.2, s_, size=8, color=MUT)
d.txt(s, p1x, p1y + 2.25, p1w, 0.55, 'Four definitions of “active customer” is the exhibit that lands.',
      size=9.5, color=NAVY, line_sp=1.2)
d.txt(s, p1x, p1y + 3.35, p1w, 0.24, '→ the conversation starter for the CDO', size=9, color=BLUE, bold=True)
p2x, p2y, p2w = panel(s, 4.97, 1.95, 3.77, 4.15, 'Truth-gap report · fragmentation cost')
d.txt(s, p2x, p2y + 0.05, p2w, 0.6, '11 min', size=34, color=BLUE, bold=True)
d.txt(s, p2x, p2y + 0.70, p2w, 0.45, 'for an RM to assemble one customer view across 5 screens',
      size=10.5, color=NAVY, line_sp=1.2)
for i, line in enumerate(['23% of fields disagree between systems', '€1.1M p.a. of RM time reassembling truth',
                          '31% of complaints touch a data mismatch']):
    d.oval(s, p2x + 0.02, p2y + 1.42 + i * 0.34, 0.09, 0.09, CORAL)
    d.txt(s, p2x + 0.2, p2y + 1.34 + i * 0.34, p2w - 0.2, 0.3, line, size=10, color=NAVY)
d.txt(s, p2x, p2y + 3.35, p2w, 0.24, '→ the € case for shared truth, before licence talk', size=9, color=BLUE, bold=True)
p3x, p3y, p3w = panel(s, 8.94, 1.95, 3.77, 4.15, 'Nexus binding blueprint · phased')
phases = [('Phase 1 · 12 wks', 'Core + CRM · 40 entities · customer, account, product', BLUE, WHITE, WHITE),
          ('Phase 2 · 10 wks', 'Credit + LMS · +25 entities · exposure, collateral', BLUE3, NAVY, NAVY),
          ('Phase 3 · next', 'Cases + documents · the service layer of truth', TINT, NAVY, NAVY)]
for i, (ph, de, fill, tc, bc) in enumerate(phases):
    y = p3y + 0.05 + i * 0.80
    d.rect(s, p3x, y, p3w, 0.68, fill=fill, round_=True)
    d.txt(s, p3x + 0.16, y + 0.08, p3w - 0.3, 0.22, ph, size=10.5, color=tc, bold=True)
    d.txt(s, p3x + 0.16, y + 0.33, p3w - 0.3, 0.3, de, size=8.5, color=bc, line_sp=1.1)
d.txt(s, p3x, p3y + 3.20, p3w, 0.45, '→ makes the Nexus licence scope precise, so pricing stops being a guess',
      size=9, color=BLUE, bold=True, line_sp=1.15)
d.footnote(s, ILLUS)
d.notes(s, 'Coverage map opens the door; the truth-gap card carries the money argument. DEFENSE: placeholders until re-based.')

product_plan(
    'Execution plan 2 of 4 · Ontology Cartographer', 'The Cartographer is the deepest build and lands on X-Ray demand',
    [
        ('Feb 2027', 'R&D gate', 'Alignment with the Semantic Modeler and Nexus teams on tooling reuse.', 'decision'),
        ('Mar-Apr 2027', 'Prototype', 'Ontology mapping demo on two sample cores. 6 pw, SE-led.', None),
        ('May-Jul 2027', 'Proof', 'Paid pilot: coverage map, truth-gap report, Nexus blueprint. 8 pw.', None),
        ('Aug 2027', 'First signature', 'Sold at €75-100K, attached to Relationship Intelligence deals.', 'ga'),
    ],
    [('Build effort', '10-12 person-wks'), ('Delivery per install', '6-8 wks · 12-16 pw'),
     ('Pod', 'SE 1.0 + VC 0.5'), ('Pre-binds', 'Nexus')],
    'Sequenced last as a full SKU: highest technical depth, needs R&D alignment. The ontology lens ships thin in every other wedge from day one.',
    notes='RI needs shared truth, so every RI opportunity is a Cartographer prospect.')

# ═════════ S20-22 · GUARDRAIL STUDIO ═════════
deep_dive(
    'Execution plan 3 of 4 · Guardrail Studio', 'Guardrail Studio turns written policy into executable rules',
    ['Wk 1-2 · harvest delegation of authority, credit policy, SOPs, entitlement matrices',
     'Wk 2-3 · map authority per process step: who may recommend, approve, execute',
     'Wk 3-4 · score agent-readiness per step on the A1-A5 autonomy scale',
     'Wk 4-5 · codification workshops (4-5) with risk and compliance: policy to executable rules',
     'Wk 5-6 · audit-trail design; Sentinel readiness score and blueprint'],
    ['Sponsor: CRO or COO',
     'Bank effort: risk, compliance and ops SMEs in 4-5 workshops',
     'Access: documents and workshops; no system access needed',
     'InfoSec stays off the critical path',
     'Pod: VC 1.0 + SE 0.5 + compliance SME 0.2'],
    ['Authority map: current vs agent-ready state',
     'Guardrail backlog: 20-30 executable rules, each traced to policy',
     'Audit-trail design for immutable evidence',
     'Sentinel readiness score + install blueprint',
     'Thin-lens appendices: leakage + ontology snapshots'],
    'Adds up: ', 'this is the approval pack a risk committee needs before any agent acts.',
    'Timed to the Conversational Banking pipeline: governed execution is the prerequisite for Transact and Resolve.',
    notes='Productizes Decision & Policy. Sentinel = authority, policies, entitlements, with immutable auditability.')

s = d.slide()
d.chrome(s, 'Guardrail Studio · sample outputs', 'Guardrail Studio shows where an agent can safely act')
p1x, p1y, p1w = panel(s, 1.0, 1.95, 3.77, 4.15, 'Authority map · card dispute process')
amap = [('Freeze card', 'A4 · agent executes', BLUE), ('Provisional credit ≤ €500', 'A4 · executes with audit', BLUE),
        ('Provisional credit > €500', 'A3 · human approves', BLUE3), ('Account closure', 'A1 · human only', HAIR)]
for i, (step, aut, c) in enumerate(amap):
    y = p1y + 0.05 + i * 0.62
    d.rect(s, p1x, y, p1w, 0.52, fill=TINT2, round_=True)
    d.txt(s, p1x + 0.14, y + 0.06, 2.4, 0.2, step, size=9.5, color=NAVY, bold=True)
    d.rect(s, p1x + 0.14, y + 0.30, 1.0, 0.15, fill=c, round_=True)
    d.txt(s, p1x + 1.22, y + 0.27, p1w - 1.35, 0.2, aut, size=8.5, color=MUT)
d.txt(s, p1x, p1y + 3.35, p1w, 0.24, '→ the boundary, drawn from the bank’s own policy', size=9, color=BLUE, bold=True)
p2x, p2y, p2w = panel(s, 4.97, 1.95, 3.77, 4.15, 'Guardrail rule card · of the 20-30 backlog')
d.rect(s, p2x - 0.04, p2y + 0.04, p2w + 0.08, 1.95, fill=NAVY, round_=True)
for i, line in enumerate(['GR-014 · Provisional credit', 'IF dispute.amount ≤ €500', 'AND fraud_score < 0.2',
                          'THEN agent EXECUTES, immutable audit log', 'ELSE route to human approver']):
    d.txt(s, p2x + 0.14, p2y + 0.18 + i * 0.33, p2w - 0.24, 0.28, line, size=10,
          color=CYAN if i == 0 else WHITE, bold=(i == 0), line_sp=1.1)
d.txt(s, p2x, p2y + 2.20, p2w, 0.24, 'Source: Card Policy §4.2 · Status: codified · Owner: dispute ops',
      size=8.5, color=MUT)
d.txt(s, p2x, p2y + 3.35, p2w, 0.24, '→ each rule loads straight into Sentinel at install', size=9, color=BLUE, bold=True)
p3x, p3y, p3w = panel(s, 8.94, 1.95, 3.77, 4.15, 'Sentinel readiness · scored')
d.txt(s, p3x, p3y + 0.02, p3w, 0.55, '34 / 100', size=32, color=CORAL, bold=True)
bars = [('Authority mapped', 45), ('Policies codified', 20), ('Audit trails in place', 38)]
for i, (lab, v) in enumerate(bars):
    y = p3y + 0.72 + i * 0.56
    d.txt(s, p3x, y, p3w, 0.2, f'{lab} · {v}%', size=9.5, color=NAVY, bold=True)
    d.rect(s, p3x, y + 0.25, p3w, 0.11, fill=TINT2)
    d.rect(s, p3x, y + 0.25, p3w * v / 100, 0.11, fill=BLUE)
d.txt(s, p3x, p3y + 2.60, p3w, 0.5, 'Target: 80+ once the backlog is codified.', size=9.5, color=NAVY, line_sp=1.2)
d.txt(s, p3x, p3y + 3.35, p3w, 0.24, '→ the risk committee’s go/no-go dashboard', size=9, color=BLUE, bold=True)
d.footnote(s, ILLUS)
d.notes(s, 'The rule card changes the conversation: policy as executable, auditable logic. DEFENSE: placeholders until re-based.')

product_plan(
    'Execution plan 3 of 4 · Guardrail Studio', 'Guardrail Studio attaches to Conversational deals in flight',
    [
        ('Jan 2027', 'Design', 'Decision & Policy alignment; risk and compliance SME briefed. 3 pw.', None),
        ('Feb-Mar 2027', 'Prototype', 'Authority-map demo: recommend, approve, execute per process. 5 pw.', None),
        ('Apr-May 2027', 'Proof', 'Paid pilot: guardrail backlog, entitlements map, Sentinel blueprint. 6 pw.', None),
        ('Jun 2027', 'First signature', 'Sold at €50-100K, riding Transact and Resolve deals.', 'ga'),
    ],
    [('Build effort', '8-10 person-wks'), ('Delivery per install', '5-6 wks · 8-10 pw'),
     ('Pod', 'VC 1.0 + SE 0.5 + SME'), ('Pre-binds', 'Sentinel')],
    'Governed execution is the prerequisite for Transact and Resolve, which makes the Sentinel-readiness wedge a natural attach on active deals.',
    notes='Sells best where a Conversational Banking deal is already in flight and risk is the blocker.')

# ═════════ S23-25 · VALUE TELEMETRY ═════════
deep_dive(
    'Execution plan 4 of 4 · Value Telemetry', 'Value Telemetry measures the true cost per resolved outcome',
    ['Wk 1 · instrument the live deployment: interaction logs, model calls, consumption',
     'Wk 1-2 · tag outcomes: resolved, escalated, abandoned; define the taxonomy',
     'Wk 2 · compute cost per outcome per use case (never per interaction)',
     'Wk 2-3 · scan for waste: PTU over-provision, model routing, dead agents, retry loops',
     'Wk 3 · re-prove ROI vs the original case; install the quarterly cadence'],
    ['Prerequisite: a live AI deployment (Backbase or third party)',
     'Access: read-only consumption + interaction logs',
     'Bank effort: finance partner to validate cost lines',
     'Pod: VC 0.5 + SE 0.25',
     'Dependency: the keystone cost-per-outcome model (with Deepak)'],
    ['Cost-per-outcome baseline per use case',
     'Waste register with € values and named owners',
     'ROI re-proof vs the original business case',
     'A quarterly value cadence, installed',
     'Thin-lens appendices: leakage + guardrail observations'],
    'Adds up: ', 'the baseline becomes the yardstick every later Mission is measured against.',
    'Only the vendor in the execution path measures cost per resolved outcome; external FinOps tools see the cloud bill. Gain-share per market practice.',
    notes='Backlog items 2.2 (keystone) + 2.3 (SKU). Partner with Aayushi on value assurance.')

s = d.slide()
d.chrome(s, 'Value Telemetry · sample outputs', 'Telemetry turns AI worth into a number the CFO can sign')
p1x, p1y, p1w = panel(s, 1.0, 1.95, 3.77, 4.15, 'Cost per outcome · by use case')
tiles = [('Card freeze', '€0.11 / outcome', '94% contained, zero human touches', GREEN, BLUE2),
         ('Dispute intake', '€1.87 / outcome', '3.2 human touches; 68% of simple intents on the premium model', CORAL, CORAL)]
for i, (uc, cost, note, ac, vc_) in enumerate(tiles):
    y = p1y + 0.05 + i * 1.42
    d.rect(s, p1x, y, p1w, 1.28, fill=TINT2, round_=True)
    d.rect(s, p1x, y + 0.08, 0.045, 1.12, fill=ac)
    d.txt(s, p1x + 0.18, y + 0.10, p1w - 0.3, 0.22, uc, size=10.5, color=NAVY, bold=True)
    d.txt(s, p1x + 0.18, y + 0.36, p1w - 0.3, 0.34, cost, size=16, color=vc_, bold=True)
    d.txt(s, p1x + 0.18, y + 0.76, p1w - 0.34, 0.45, note, size=9, color=MUT, line_sp=1.15)
d.txt(s, p1x, p1y + 3.35, p1w, 0.24, '→ the unit economics of every use case, measured', size=9, color=BLUE, bold=True)
p2x, p2y, p2w = panel(s, 4.97, 1.95, 3.77, 4.15, 'Waste register · €312K annualized')
wastes = [('PTU over-provision', 140), ('Model routing', 95), ('Dead agents', 41), ('Retry loops', 36)]
for i, (lab, v) in enumerate(wastes):
    y = p2y + 0.08 + i * 0.56
    d.txt(s, p2x, y, p2w, 0.2, f'{lab} · €{v}K', size=9.5, color=NAVY, bold=True)
    d.rect(s, p2x, y + 0.25, p2w, 0.12, fill=TINT2)
    d.rect(s, p2x, y + 0.25, p2w * v / 140, 0.12, fill=CORAL if i == 0 else BLUE)
d.txt(s, p2x, p2y + 2.45, p2w, 0.6, 'Each line has an owner and a fix; savings are machine-measurable, which is what makes gain-share defensible.',
      size=9.5, color=NAVY, line_sp=1.2)
d.txt(s, p2x, p2y + 3.35, p2w, 0.24, '→ usually pays for the diagnostic by itself', size=9, color=BLUE, bold=True)
p3x, p3y, p3w = panel(s, 8.94, 1.95, 3.77, 4.15, 'ROI re-proof · case vs measured')
rois = [('Original business case', 1.10, BLUE3), ('Measured today', 0.82, CORAL), ('After routing + PTU fixes', 1.05, BLUE)]
for i, (lab, v, c) in enumerate(rois):
    y = p3y + 0.08 + i * 0.68
    d.txt(s, p3x, y, p3w, 0.2, f'{lab} · €{v:.2f}M a year', size=9.5, color=NAVY, bold=True)
    d.rect(s, p3x, y + 0.25, p3w, 0.14, fill=TINT2)
    d.rect(s, p3x, y + 0.25, p3w * v / 1.10, 0.14, fill=c)
d.rect(s, p3x, p3y + 2.30, 2.1, 0.32, fill=TINT, round_=True)
d.txt(s, p3x + 0.14, p3y + 2.36, 1.9, 0.2, 'Case re-proved at 95%', size=9.5, color=BLUE2, bold=True)
d.txt(s, p3x, p3y + 3.35, p3w, 0.24, '→ renewal argued with the bank’s own numbers', size=9, color=BLUE, bold=True)
d.footnote(s, ILLUS)
d.notes(s, 'The waste register is the fastest trust-builder. DEFENSE: placeholders until re-based; gain-share is only offered where telemetry is installed.')

product_plan(
    'Execution plan 4 of 4 · Value Telemetry', 'Telemetry rides the keystone model and opens the recurring line',
    [
        ('Sep 2026', 'Keystone gate', 'Extend the cost model to cost-per-outcome, with Deepak. ~5 use cases.', 'decision'),
        ('Sep 2026', 'Design', 'SKU spec: diagnostic scope, telemetry set, gain-share terms. 4 pw.', None),
        ('Oct 2026', 'Prototype', 'Cost-per-outcome dashboard on one live deployment’s data. 3 pw.', None),
        ('Nov-Dec 2026', 'Proof', '2-3 week diagnostic at a live customer: waste quantified, ROI re-proved.', None),
        ('Jan 2027', 'First signature', '€15-50K diagnostic + 15-35% gain share; recurrence begins.', 'ga'),
    ],
    [('Build effort', '6-8 person-wks'), ('Delivery per install', '2-3 wks · 3-4 pw'),
     ('Pod', 'VC 0.5 + SE 0.25'), ('Model', 'Recurring + gain share')],
    'Gain share of 15-35% of realised savings is market-validated (ProsperOps, nOps, Vantage) with a “10% or free” floor. Cost-model sample today: n=2 use cases.',
    notes='The keystone study is the single hardest dependency; it gates the whole recurring line.')

# ═════════ ENGINE B CHAPTER · THE AI-NATIVE OPERATING MODEL ═════════
d.divider('03', 'The AI-native operating model',
          'Engine B in full: loop postures, preconditions, the maturity path, the workforce shift, and how to run it.')

deep_dive(
    'Engine B · AI-native operating model · what it is', 'Engine B redesigns the org chart for a world where agents work',
    ['Maturity assessment (2-3 wks) · score every LOB × ops area on the A1-A5 curve; readiness on data, guardrails, skills',
     'Workforce optimization (3-4 wks) · which roles rise, which collapse into agents; loop boundaries per function',
     'Org design flagship (4-6 wks) · the inverted-T target: a thin judgment layer over an agent workforce',
     'Transition path · function by function up the curve; the org chart follows the highest-autonomy function'],
    ['Sponsor: CEO, COO or CHRO; ExCo workshops',
     'Inputs: Engine A evidence: X-Ray scores, authority maps, telemetry',
     'HR data: role inventory, spans and layers',
     'Pod: VC 1.0 + org/HR SME on call',
     'Method: AI-Native Org Design draft v0.1 exists'],
    ['A1-A5 maturity heatmap by function, current vs target',
     'Workforce transition map: roles rising, collapsing, new',
     'Target org chart: the inverted-T with the governance spine',
     'New-role definitions: Mission Owners, AgentOps, Exception Desks',
     'Function-by-function transition roadmap'],
    'Adds up: ', 'Engine B opens the CEO and CHRO door Engine A cannot.',
    'Sold only downstream of a wedge install, so the design runs on the bank’s own evidence. Ticket €80-150K is an assumption to validate.',
    notes='PDP workstream 4. New roles per the backlog: Chief Agentic Operations Officer, Mission Owners, AgentOps, '
          'Nexus Stewards, Sentinel Governance, Exception Desks, Workforce Transition office.')

# B2 · loop postures
s = d.slide()
d.chrome(s, 'Engine B · the loop postures', 'Three loop postures span the curve; one zone stays human')
posts = [
    ('IN THE LOOP · A1-A3', BLUE4, 'Human decides', 'AI drafts, recommends, or executes only on approval.',
     ['Credit memos pre-assembled', 'Next-best-action for RMs', 'Payment repair suggestions'],
     'The human becomes: the decider, with better inputs.'),
    ('ON THE LOOP · A4', BLUE, 'Agent acts, human supervises', 'Execute by exception: the agent completes the loop and escalates edge cases.',
     ['Dispute intake and triage', 'KYC periodic refresh', 'Routine servicing requests'],
     'The human becomes: exception handler and supervisor.'),
    ('ABOVE THE LOOP · A5', BLUE2, 'Agent runs the process', 'Self-directed within guardrails; humans govern, design and audit.',
     ['Document collection chase', 'Status communication loops', 'Recon break clearing'],
     'The human becomes: governor and Mission designer.'),
    ('NOT DELEGATED', CORAL, 'Human only, by design', 'Risk keeps these out of agent hands regardless of capability.',
     ['Final credit above thresholds', 'Vulnerable-customer journeys', 'Regulatory attestations'],
     'The human stays: accountable judgment.'),
]
cw_, cg_ = 2.85, 0.11
for i, (tag, c, head, body, exs, becomes) in enumerate(posts):
    x = 1.0 + i * (cw_ + cg_)
    d.rect(s, x, 1.95, cw_, 3.55, fill=TINT2, round_=True)
    d.rect(s, x, 1.95, cw_, 0.34, fill=c, round_=True)
    d.txt(s, x + 0.14, 2.01, cw_ - 0.28, 0.24, tag, size=9,
          color=WHITE if c in (BLUE, BLUE2, CORAL) else NAVY, bold=True)
    d.txt(s, x + 0.16, 2.42, cw_ - 0.32, 0.26, head, size=11.5, color=NAVY, bold=True)
    d.txt(s, x + 0.16, 2.72, cw_ - 0.32, 0.62, body, size=9.5, color=NAVY, line_sp=1.15)
    for j, e_ in enumerate(exs):
        d.txt(s, x + 0.16, 3.42 + j * 0.28, cw_ - 0.32, 0.26, '· ' + e_, size=9, color=MUT, line_sp=1.05)
    d.txt(s, x + 0.16, 4.42 + 0.55, cw_ - 0.32, 0.5, becomes, size=9, color=BLUE2 if c != CORAL else CORAL,
          bold=True, line_sp=1.1)
d.takeaway_band(s, 'The autonomy scale is the operating-model axis: ', 'the posture, the role and the org all follow it.')
d.footnote(s, 'Terminology per the PDP autonomy framework (A1-A5). Market reference: a large US bank runs the same three cases as in, on and out of the loop.')
d.notes(s, 'Citibank framing: in the loop, on the loop, out of the loop. Ours: in (A1-A3), on (A4), above (A5), plus the not-delegated zone which risk owns. Detail nodes per posture.')

# B3 · preconditions: possible vs not
s = d.slide()
d.chrome(s, 'Engine B · the permission map', 'Each step up the curve has preconditions; some moves stay off')
lads = [
    ('A2 → A3 · execute on approval', ['Policy codified for the step', 'Immutable audit trail live', 'Approval UX in the workspace']),
    ('A3 → A4 · execute by exception', ['Guardrails live in Sentinel', 'Nexus truth for the journey', 'Exception desk stood up', 'Model performance ≥ agreed floor']),
    ('A4 → A5 · self-directed', ['Regulator comfort documented', 'Two quarters of clean A4 evidence', 'Rollback and kill-switch tested', 'Telemetry mature (cost per outcome)']),
]
for i, (h_, conds) in enumerate(lads):
    y = 1.95 + i * 1.22
    d.rect(s, 1.0, y, 7.6, 1.08, fill=TINT2, round_=True)
    d.rect(s, 1.0, y + 0.07, 0.045, 0.94, fill=BLUE)
    d.txt(s, 1.22, y + 0.09, 7.2, 0.24, h_, size=11.5, color=NAVY, bold=True)
    d.txt(s, 1.22, y + 0.38, 7.25, 0.65,
          [[('· ' + c_, 9.5, NAVY, False)] for c_ in conds[:2]], line_sp=1.12, sp_after=2)
    d.txt(s, 4.85, y + 0.38, 3.6, 0.65,
          [[('· ' + c_, 9.5, NAVY, False)] for c_ in conds[2:]], line_sp=1.12, sp_after=2)
d.rect(s, 8.85, 1.95, 3.87, 3.62, line=CORAL, line_w=1.2, round_=True, dash='dash')
d.txt(s, 9.05, 2.10, 3.5, 0.24, 'NOT PERMITTED FOR NOW', size=10, color=CORAL, bold=True)
for i, n_ in enumerate(['Autonomous final credit decisions above policy thresholds',
                        'Autonomous action on vulnerable-customer journeys',
                        'Attestations and regulatory filings',
                        'Any action without an immutable audit trail',
                        'Cross-border data pooling outside residency rules']):
    d.txt(s, 9.05, 2.44 + i * 0.60, 3.5, 0.56, '· ' + n_, size=9.5, color=NAVY, line_sp=1.12)
d.takeaway_band(s, 'The permission map is drawn per bank: ', 'Guardrail Studio’s output is exactly this slide, on their policy.')
d.footnote(s, 'Preconditions compound: nothing reaches A4 without Sentinel and Nexus for that journey, which is why the wedges precede the org design.')
d.notes(s, 'What is possible vs not possible. The not-permitted list is a starting position drawn from typical bank policy; each bank redraws it from its own delegation of authority.')

# B4 · the maturity path
s = d.slide()
d.chrome(s, 'Engine B · the maturity path', 'The org follows the highest-autonomy function up the curve')
phases = [
    ('EARLY · A1-A3', BLUE4, 'Assist and approve',
     ['Existing teams keep the work', 'AI literacy + approval discipline', 'First guardrails codified'],
     'Metric: % of steps AI-assisted'),
    ('MID · A4', BLUE, 'Execute by exception',
     ['Exception desks form from the processing floor', 'AgentOps stands up (home of FinOps)', 'Mission Owners named per LOB'],
     'Metric: containment % · cost per outcome'),
    ('LATE · A5', NAVY, 'Self-directed in guardrails',
     ['The governance spine is the centre of gravity', 'Thin judgment layer over the agent workforce', 'Workforce Transition office completes moves'],
     'Metric: outcomes per FTE · value per Mission'),
]
cw2, cg2 = 3.77, 0.20
for i, (tag, c, head, items, metric) in enumerate(phases):
    x = 1.0 + i * (cw2 + cg2)
    d.rect(s, x, 1.95, cw2, 3.15, fill=TINT2, round_=True)
    d.rect(s, x, 1.95, cw2, 0.34, fill=c, round_=True)
    d.txt(s, x + 0.16, 2.01, cw2 - 0.3, 0.24, tag, size=9.5, color=WHITE if c in (BLUE, NAVY) else NAVY, bold=True)
    d.txt(s, x + 0.18, 2.42, cw2 - 0.36, 0.26, head, size=12.5, color=NAVY, bold=True)
    for j, it in enumerate(items):
        d.txt(s, x + 0.18, 2.78 + j * 0.52, cw2 - 0.36, 0.5, '· ' + it, size=9.5, color=NAVY, line_sp=1.12)
    d.txt(s, x + 0.18, 4.62, cw2 - 0.36, 0.4, metric, size=9.5, color=BLUE2, bold=True, line_sp=1.1)
    if i < 2:
        d.txt(s, x + cw2 - 0.02, 3.35, 0.24, 0.3, '→', size=14, color=MUT)
d.takeaway_band(s, 'Function by function: ', 'a bank moves domains up the curve one at a time, and the chart follows the leader.')
d.footnote(s, 'Phase gates are the B3 preconditions. A function can be LATE while the bank average is EARLY; the org chart follows the most autonomous function.')
d.notes(s, 'The transition path from the PDP backlog (4.3). This is the slide that turns the org chart from a poster into a plan.')

# B5 · workforce and skills
s = d.slide()
d.chrome(s, 'Engine B · the workforce shift', 'Skills shift from processing to judgment, supervision, design')
shifts = [
    ('AUGMENTED ROLES', BLUE4, 'AI sharpens human judgment',
     ['RMs: next-best-action, pre-briefed meetings', 'Credit officers: pre-assembled cases', 'Risk analysts: surfaced anomalies'],
     'Skills to build: interrogating AI output, override discipline, data literacy.'),
    ('AUTONOMOUS ZONES', BLUE, 'Humans move above the work',
     ['Servicing, disputes intake, KYC refresh', 'Processing floor re-skills into exception desks', 'Supervisors read agent telemetry'],
     'Skills to build: exception triage, root-cause analysis, agent performance review.'),
    ('NEW ROLES', BLUE2, 'Humans design the system',
     ['Mission Owners and Mission designers', 'AgentOps: deploy, monitor, tune, FinOps', 'Nexus Stewards and Sentinel Governance'],
     'Skills to build: mission design, guardrail authoring, data-as-product stewardship.'),
]
for i, (tag, c, head, items, skills) in enumerate(shifts):
    x = 1.0 + i * (cw2 + cg2)
    d.rect(s, x, 1.95, cw2, 3.15, fill=TINT2, round_=True)
    d.rect(s, x, 1.95, cw2, 0.34, fill=c, round_=True)
    d.txt(s, x + 0.16, 2.01, cw2 - 0.3, 0.24, tag, size=9.5, color=WHITE if c in (BLUE, BLUE2) else NAVY, bold=True)
    d.txt(s, x + 0.18, 2.42, cw2 - 0.36, 0.26, head, size=12.5, color=NAVY, bold=True)
    for j, it in enumerate(items):
        d.txt(s, x + 0.18, 2.78 + j * 0.46, cw2 - 0.36, 0.44, '· ' + it, size=9.5, color=NAVY, line_sp=1.12)
    d.txt(s, x + 0.18, 4.28, cw2 - 0.36, 0.75, skills, size=9.5, color=BLUE2, bold=True, line_sp=1.15)
d.takeaway_band(s, 'The transition motion: ', 'evidence, co-design with the floor, skills academy, exception desks, then measure judgment time.')
d.footnote(s, 'Sequenced to avoid a layoffs-first read: the processing floor re-skills into exception desks before any structural change is announced.')
d.notes(s, 'Answers “how do we transition people to the right-hand side”: show the evidence (X-Ray), co-design roles with the floor, build the three skill sets, form exception desks, track % time on judgment work as the success metric.')

s = d.slide()
d.chrome(s, 'Engine B · sample outputs', 'Engine B draws the target org from the bank’s own evidence')
p1x, p1y, p1w = panel(s, 1.0, 1.95, 3.77, 4.15, 'Autonomy maturity heatmap · current → target')
funcs = [('Disputes', 'A2', 'A4'), ('Onboarding', 'A1', 'A3'), ('KYC review', 'A2', 'A4'),
         ('Contact centre', 'A2', 'A4'), ('Credit servicing', 'A1', 'A3')]
for i, (f_, cur, tgt) in enumerate(funcs):
    y = p1y + 0.05 + i * 0.52
    d.rect(s, p1x, y, p1w, 0.42, fill=TINT2, round_=True)
    d.txt(s, p1x + 0.14, y + 0.10, 1.6, 0.2, f_, size=9.5, color=NAVY, bold=True)
    d.rect(s, p1x + 1.80, y + 0.08, 0.46, 0.26, fill=BLUE3, round_=True)
    d.txt(s, p1x + 1.80, y + 0.11, 0.46, 0.2, cur, size=9, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    d.txt(s, p1x + 2.30, y + 0.09, 0.26, 0.2, '→', size=10, color=MUT, align=PP_ALIGN.CENTER)
    d.rect(s, p1x + 2.60, y + 0.08, 0.46, 0.26, fill=BLUE, round_=True)
    d.txt(s, p1x + 2.60, y + 0.11, 0.46, 0.2, tgt, size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
d.txt(s, p1x, p1y + 3.30, p1w, 0.45, '→ scored from X-Ray and Guardrail evidence, never from interviews alone',
      size=9, color=BLUE, bold=True, line_sp=1.12)
p2x, p2y, p2w = panel(s, 4.97, 1.95, 3.77, 4.15, 'The org shift · pyramid → inverted-T')
pyr_cx = p2x + 0.85
for i, w_ in enumerate([0.62, 1.14, 1.66]):
    d.rect(s, pyr_cx - w_ / 2, p2y + 0.20 + i * 0.36, w_, 0.30, fill=BLUE3)
d.txt(s, pyr_cx - 0.85, p2y + 1.40, 1.7, 0.45, 'Today: layers of human processing', size=8.5, color=MUT,
      align=PP_ALIGN.CENTER, line_sp=1.12)
d.txt(s, p2x + 1.62, p2y + 0.60, 0.3, 0.3, '→', size=14, color=MUT, align=PP_ALIGN.CENTER)
it_cx = p2x + 2.72
d.rect(s, it_cx - 0.85, p2y + 0.20, 1.70, 0.30, fill=BLUE)
d.txt(s, it_cx - 0.83, p2y + 0.25, 1.66, 0.2, 'Human judgment layer', size=7.5, color=WHITE, bold=True,
      align=PP_ALIGN.CENTER)
d.rect(s, it_cx - 0.21, p2y + 0.54, 0.42, 0.72, fill=NAVY)
d.txt(s, it_cx - 0.85, p2y + 1.40, 1.7, 0.45, 'Governance + AgentOps spine', size=8.5, color=MUT,
      align=PP_ALIGN.CENTER, line_sp=1.12)
d.rect(s, p2x, p2y + 2.00, p2w, 0.36, fill=TINT, round_=True)
d.txt(s, p2x + 0.12, p2y + 2.08, p2w - 0.24, 0.2, 'Agent workforce running Missions underneath',
      size=9.5, color=BLUE2, bold=True)
d.txt(s, p2x, p2y + 3.35, p2w, 0.24, '→ the flagship exhibit for the CEO and CHRO', size=9, color=BLUE, bold=True)
p3x, p3y, p3w = panel(s, 8.94, 1.95, 3.77, 4.15, 'New roles on the target chart')
roles = [('Mission Owners', 'own agent-delivered outcomes per LOB'),
         ('AgentOps', 'deploy, monitor, optimize the fleet; home of FinOps'),
         ('Nexus Stewards', 'data-as-product; keepers of shared truth'),
         ('Sentinel Governance', 'AI assurance, authority, audit'),
         ('Exception Desks', 'the re-skilled processing floor')]
for i, (r_, de) in enumerate(roles):
    y = p3y + 0.05 + i * 0.54
    d.rect(s, p3x, y, p3w, 0.44, fill=TINT2, round_=True)
    d.rect(s, p3x, y + 0.05, 0.04, 0.34, fill=BLUE)
    d.txt(s, p3x + 0.13, y + 0.04, 1.55, 0.2, r_, size=9.5, color=NAVY, bold=True)
    d.txt(s, p3x + 1.62, y + 0.05, p3w - 1.72, 0.36, de, size=8, color=MUT, line_sp=1.05)
d.txt(s, p3x, p3y + 3.35, p3w, 0.24, '→ every role binds the org to a Banking OS construct', size=9, color=BLUE, bold=True)
d.footnote(s, ILLUS)
d.notes(s, 'The pyramid-to-inverted-T exhibit is the flagship visual: the org chart follows the autonomy curve.')

# B7 · how to execute
s = d.slide()
d.chrome(s, 'Engine B · how to run it', 'Six workshops take an ExCo from evidence to a signed target org')
wshops = [
    ('W1', 'Evidence readout', 'X-Ray leakage + Guardrail authority data on the bank’s own journeys', 'ExCo + COO team'),
    ('W2', 'Loop postures', 'Map every function to in, on, above the loop; agree the not-delegated zone', 'COO + CRO + LOB heads'),
    ('W3', 'Permission map', 'Preconditions per step; what is possible now vs preconditioned vs off', 'CRO + compliance + IT'),
    ('W4', 'Target org draft', 'The inverted-T: judgment layer, governance spine, agent workforce', 'CEO + CHRO + COO'),
    ('W5', 'Workforce plan', 'Role moves, skills academy, exception-desk formation', 'CHRO + LOB heads'),
    ('W6', 'Transition roadmap', 'Function-by-function sequence, metrics, sign-off', 'ExCo'),
]
for i, (n_, name, out_, who) in enumerate(wshops):
    x = 1.0 + (i % 3) * 3.97
    y = 2.00 + (i // 3) * 1.65
    d.rect(s, x, y, 3.77, 1.50, fill=TINT2, round_=True)
    d.oval(s, x + 0.16, y + 0.14, 0.40, 0.40, BLUE, label=n_, fs=10)
    d.txt(s, x + 0.68, y + 0.20, 3.0, 0.26, name, size=12, color=NAVY, bold=True)
    d.txt(s, x + 0.18, y + 0.62, 3.4, 0.55, out_, size=9.5, color=NAVY, line_sp=1.15)
    d.txt(s, x + 0.18, y + 1.20, 3.4, 0.24, who, size=8.5, color=BLUE2, bold=True)
d.takeaway_band(s, 'Cadence: ', 'six workshops over 6-8 weeks, each producing a signed artifact, ending in the transition roadmap.')
d.footnote(s, 'W1 requires a completed wedge (the sequencing rule). Attendee lists are the minimum quorum; the CHRO joins from W4 at the latest.')
d.notes(s, 'How to execute Engine B. The signed W6 roadmap is the deliverable that converts into the Workforce Transition office and the next Missions.')

product_plan(
    'Engine B · the plan', 'Engine B hardens v0.1 now and sells downstream of X-Ray',
    [
        ('Sep 2026', 'Method hardening', 'Pressure-test org-design v0.1 with Mayur; codify the kit. 4 pw.', None),
        ('Oct-Dec 2026', 'Evidence binding', 'Wire the method to X-Ray and Guardrail outputs.', None),
        ('Q1 2027', 'Sequencing gate', 'First sale only downstream of a wedge install; never sold cold.', 'decision'),
        ('Mar-May 2027', 'First engagement', 'Maturity + workforce + org design at the X-Ray pilot account.', None),
        ('Jun 2027', 'Flagship scale', 'CEO/CHRO-altitude offer; Early Access cohort accounts first.', 'ga'),
    ],
    [('Build effort', '6-8 pw (v0.1 exists)'), ('Delivery', '6-8 wks · 10-14 pw'),
     ('Pod', 'VC 1.0 + org SME'), ('Pre-binds', 'The human layer of the OS')],
    'Ticket €80-150K is an assumption to validate against market org-design pricing. The moat is the evidence feed, never the method alone.',
    notes='Coral diamond = the sequencing rule decision (never sold cold). This is the fifth row of the critique chapter.')

# ═════════ S29 · REGISTER (the ONE table) ═════════
s = d.slide()
d.chrome(s, 'Execution plans · the register', 'Twelve months of product build fits inside the existing pod')
rows = [('Product', 'One-time build', 'Delivery per install', 'Pod', 'First revenue'),
        ('Process X-Ray', '8-10 pw', '4-6 wks · 12-14 pw', 'VC 1.0 + SE 0.5', 'Jan 2027'),
        ('Ontology Cartographer', '10-12 pw', '6-8 wks · 12-16 pw', 'SE 1.0 + VC 0.5', 'Aug 2027'),
        ('Guardrail Studio', '8-10 pw', '5-6 wks · 8-10 pw', 'VC 1.0 + SE 0.5 + SME', 'Jun 2027'),
        ('Value Telemetry', '6-8 pw', '2-3 wks · 3-4 pw', 'VC 0.5 + SE 0.25', 'Jan 2027'),
        ('AI-Native Org Design (Engine B)', '6-8 pw (v0.1)', '6-8 wks · 10-14 pw', 'VC 1.0 + org SME', 'May 2027')]
x = 1.0
widths = [3.6, 1.9, 2.4, 2.6, 1.2]
y = 2.15
for ri, row in enumerate(rows):
    cx = x
    for ci, cell in enumerate(row):
        if ri == 0:
            d.txt(s, cx, y, widths[ci], 0.3, cell.upper(), size=10.5, color=BLUE, bold=True)
        else:
            d.txt(s, cx, y, widths[ci], 0.3, cell, size=11.5, color=NAVY, bold=(ci == 0))
        cx += widths[ci]
    d.hline(s, x, y + 0.44, x + sum(widths), y + 0.44, color=HAIR_ROW if ri else HAIR, wpt=0.75)
    y += 0.60
d.takeaway_band(s, 'Never two builds at once: ', 'the sequence keeps total build load at 38-48 person-weeks over twelve months.')
d.footnote(s, 'pw = person-weeks. Build = one-time productization. Delivery = per client install once productized. All efforts carry ±30% tolerance.')
d.notes(s, 'THE one table (T13). Presentation order by layer logic; build order optimizes risk (see roadmap).')

# ═════════ BUSINESS CONNECTION CHAPTER ═════════
d.divider('04', 'The business connection',
          'Every POC lands under a named business objective, and the account list splits into must-have and good-to-have.')

# C1 · ambition → objective → Mission bridge
s = d.slide()
d.chrome(s, 'Business connection · the bridge', 'Every POC lands under a named business objective, or it waits')
nodes = [
    ('BUSINESS AMBITION', 'The 5-year strategy', 'Cost programs, growth targets, AI ambition, license to operate', TINT2, NAVY),
    ('INVESTMENT OBJECTIVE', 'What the CFO funds', 'Cut cost-to-serve · grow revenue · de-risk and comply · scale without headcount', TINT2, NAVY),
    ('VALUE POOL', 'The Banking OS math', 'Cost-to-serve down 20-40% · conversion up 10-25% · AI to production 3-5× faster', TINT, BLUE2),
    ('MISSION POC', 'Go-live, in production', 'The top wedge candidate live in 6-12 weeks, telemetry from day one', BLUE, WHITE),
    ('MEASURED PROOF', 'The renewal engine', 'Cost per outcome vs the case, quarterly; funds the next objective', NAVY, WHITE),
]
nw, ng = 2.20, 0.17
for i, (tag, head, body, fill, tc) in enumerate(nodes):
    x = 1.0 + i * (nw + ng)
    d.rect(s, x, 2.10, nw, 2.60, fill=fill, round_=True)
    d.txt(s, x + 0.14, 2.24, nw - 0.28, 0.4, tag, size=8.5, color=BLUE if fill in (TINT2, TINT) else CYAN, bold=True)
    d.txt(s, x + 0.14, 2.62, nw - 0.28, 0.5, head, size=11, color=tc if fill in (BLUE, NAVY) else NAVY, bold=True, line_sp=1.05)
    d.txt(s, x + 0.14, 3.14, nw - 0.28, 1.4, body, size=9,
          color=tc if fill in (BLUE, NAVY) else NAVY, line_sp=1.18)
    if i < 4:
        d.txt(s, x + nw - 0.01, 3.15, 0.19, 0.3, '→', size=12, color=MUT)
d.takeaway_band(s, 'A use case without an objective is a demo: ', 'under an objective it is an investment with a payback date.')
d.footnote(s, 'Value pools per the Banking OS canon (three pools, $150M-300M+ annual unlock, bank-size dependent). The POC inherits its target from the objective, never the reverse.')
d.notes(s, 'The missing connection the consultant flagged: selling POCs with go-live is powerful only when the POC lands under a funded objective. This chain is the qualification test in the free wedge.')

# C2 · must-have vs good-to-have (full version)
s = d.slide()
d.chrome(s, 'Business connection · the account lens', 'Must-have shows in the strategy before it shows in the pipeline')
segs = [
    ('MUST-HAVE NOW', BLUE2, 'High cost pressure and stated AI ambition',
     ['Named cost or CIR program in flight', 'AI or agentic ambition in the strategy', 'Digital attacker in the home market'],
     'Motion: X-Ray + EAP nomination, ExCo altitude, this quarter.'),
    ('MUST-HAVE SOON', NAVY, 'The pressure is visible, the ambition is not yet',
     ['Margin compression, no stated response', 'Headcount or branch targets without a how', 'Regulatory remediation load'],
     'Motion: free wedge sells the ambition; the evidence writes their case.'),
    ('GOOD-TO-HAVE', BLUE, 'Ambition without acute pressure',
     ['Innovation budget and appetite', 'Digital leadership positioning', 'Early agentic experiments running'],
     'Motion: ideal EAP design partners; Telemetry on whatever already runs.'),
]
for i, (tag, c, head, sigs, motion) in enumerate(segs):
    x = 1.0 + i * (cw2 + cg2)
    d.rect(s, x, 1.95, cw2, 3.30, fill=TINT2, round_=True)
    d.rect(s, x, 1.95, cw2, 0.34, fill=c, round_=True)
    d.txt(s, x + 0.16, 2.01, cw2 - 0.3, 0.24, tag, size=9.5, color=WHITE, bold=True)
    d.txt(s, x + 0.18, 2.42, cw2 - 0.36, 0.5, head, size=11.5, color=NAVY, bold=True, line_sp=1.1)
    for j, sg in enumerate(sigs):
        d.txt(s, x + 0.18, 3.02 + j * 0.42, cw2 - 0.36, 0.4, '· ' + sg, size=9.5, color=NAVY, line_sp=1.1)
    d.txt(s, x + 0.18, 4.42, cw2 - 0.36, 0.72, motion, size=9.5, color=BLUE2, bold=True, line_sp=1.15)
d.takeaway_band(s, 'Classify with the AEs in week one: ', 'annual-report signals sort the account list before any outreach.')
d.footnote(s, 'Archetypes on the slide, named accounts in the working session. Signals read from public strategy documents; classification validated with the AE.')
d.notes(s, 'The five-year-strategy lens the consultant asked for. Must-have = the strategy already needs this to be true. DEFENSE: classification is evidence-based (public signals), refreshed quarterly.')

# ═════════ DIVIDER 05 · MOBILIZATION ═════════
d.divider('05', 'Mobilization',
          'The twelve-month sequence, how the Delivery Factory carries it, and the asks that unlock it.')

# ═════════ S31 · ROADMAP ═════════
s = d.slide()
d.chrome(s, 'Mobilization · the sequence', 'One product at a time reaches four SKUs inside twelve months')
months = ['Aug', 'Sep', 'Oct', 'Nov', 'Dec', 'Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul']
gx0, gy0, colw, rowh = 2.72, 2.55, 0.828, 0.60
d.txt(s, gx0, gy0 - 0.62, 2.0, 0.2, '2026', size=9, color=MUT)
d.txt(s, gx0 + 5 * colw, gy0 - 0.62, 2.0, 0.2, '2027', size=9, color=MUT)
for i, m in enumerate(months):
    d.txt(s, gx0 + i * colw, gy0 - 0.36, colw, 0.22, m, size=9, color=MUT, bold=True, align=PP_ALIGN.CENTER)
    d.hline(s, gx0 + i * colw, gy0 - 0.08, gx0 + i * colw, gy0 + 6 * rowh - 0.15, color=HAIR_ROW)
grows = [
    ('Keystone study', [(1, 2, BLUE2, 'With Deepak')], [(1.0, CORAL)]),
    ('Process X-Ray', [(0, 1, BLUE4, ''), (1, 2, BLUE, 'Prototype'), (3, 2, BLUE2, 'Proof'), (5, 7, NAVY, 'Sell + install')], [(2.0, CORAL), (5.0, BLUE)]),
    ('Value Telemetry', [(1, 1, BLUE4, ''), (2, 1, BLUE, ''), (3, 2, BLUE2, 'Proof'), (5, 7, NAVY, 'Recurring')], [(5.0, BLUE)]),
    ('Guardrail Studio', [(5, 1, BLUE4, ''), (6, 2, BLUE, 'Prototype'), (8, 2, BLUE2, 'Proof'), (10, 2, NAVY, 'Sell')], [(10.0, BLUE)]),
    ('Ontology Cartographer', [(6, 1, BLUE4, ''), (7, 2, BLUE, 'Prototype'), (9, 3, BLUE2, 'Proof')], [(6.0, CORAL)]),
    ('Org Design (Engine B)', [(1, 1, BLUE4, ''), (2, 3, BLUE, 'Evidence binding'), (7, 3, BLUE2, 'First engagement')], [(7.0, CORAL), (10.0, BLUE)]),
]
for r, (label, bars, gates) in enumerate(grows):
    y = gy0 + r * rowh
    d.txt(s, 1.0, y + 0.06, 1.68, 0.45, label, size=10, color=NAVY, bold=True, line_sp=1.05)
    for (start, dur, color, blabel) in bars:
        bx = gx0 + start * colw + 0.02
        bw = dur * colw - 0.04
        d.rect(s, bx, y + 0.06, bw, 0.28, fill=color, round_=True)
        if blabel and bw > 1.2:
            d.txt(s, bx + 0.1, y + 0.10, bw - 0.16, 0.2, blabel, size=8.5, color=WHITE, bold=True)
    for (gm, gcol) in gates:
        d.diamond(s, gx0 + gm * colw - 0.08, y + 0.12, 0.16, 0.16, fill=gcol, line=WHITE, line_w=0.8)
d.footnote(s, 'Coral diamonds = decisions an owner must make (SE named · pilot green-light · R&D alignment · sequencing). Blue diamonds = first revenue per product.')
d.notes(s, 'T04 waves. Sequencing: X-Ray first (lowest risk), Telemetry parallel (shares the keystone), Guardrail rides the Conversational pipeline, Cartographer lands on X-Ray demand.')

# ═════════ S32 · DELIVERY FACTORY ═════════
s = d.slide()
d.chrome(s, 'Mobilization · the Delivery Factory', 'The Delivery Factory carries every stage through to install')
stages = [
    ('Prototype', 'Semantic Modeler + Connector Studio', 'A working demo on the bank’s own data, in weeks.'),
    ('Proof', 'Simulation & Testing + Ops telemetry', 'CFO-grade evidence runs with zero production risk.'),
    ('Install', 'Mission Contract, Nexus + Sentinel binding', 'The Banking OS layer live in 6-12 weeks via the agentic SDLC.'),
    ('Run', 'Deployment & Ops Control', 'Live telemetry that feeds Value Telemetry and the expansion case.'),
]
cw, cgap = 2.78, 0.20
for i, (stage, tools, out_) in enumerate(stages):
    x = 1.0 + i * (cw + cgap)
    d.rect(s, x, 2.05, cw, 2.90, fill=TINT2, round_=True)
    d.txt(s, x + 0.20, 2.25, cw - 0.4, 0.22, stage.upper(), size=10, color=BLUE, bold=True)
    d.txt(s, x + 0.20, 2.55, cw - 0.4, 0.75, tools, size=12.5, color=NAVY, bold=True, line_sp=1.12)
    d.hline(s, x + 0.20, 3.45, x + cw - 0.20, 3.45, color=HAIR)
    d.txt(s, x + 0.20, 3.58, cw - 0.4, 1.2, out_, size=10.5, color=NAVY, line_sp=1.2)
    if i < 3:
        d.txt(s, x + cw + 0.005, 3.25, 0.19, 0.3, '→', size=13, color=MUT)
d.takeaway_band(s, 'Nothing built from scratch: ', 'the products are the Factory’s discovery tools, productized.')
d.footnote(s, 'Factory toolchain per the Banking OS canon: Process & Workspace Designer · Semantic Modeler · Agent Builder · Decision & Policy · Connector Studio · Simulation & Testing · Ops Control.')
d.notes(s, 'The PDP insight: the Factory already contains the discovery-end tools. Delivery cost stays at person-weeks.')

# ═════════ S33 · THE ASKS ═════════
s = d.slide()
d.chrome(s, 'Mobilization · the ask', 'Four decisions unlock execution; the €5,000 seeds product one')
asks = [
    ('1 · Name the solution engineer', '0.5 FTE from the regional pool, committed by September. The single critical-path resource.'),
    ('2 · Green-light one prospect pilot', 'AIB-profile account for the X-Ray proof, by October. A prospect; existing customers stay out of scope.'),
    ('3 · Back the keystone study', 'Cost-per-outcome model with Deepak in September. Value Telemetry and the FinOps line depend on it.'),
    ('4 · Agree the commercial treatment', 'Quota recognition of ~€200K so the line carries a real target from day one.'),
]
for i, (h_, b_) in enumerate(asks):
    chip_card(s, 1.0, 1.95 + i * 0.92, 5.9, 0.80, h_, b_)
hero(s, 7.15, 1.95, 5.55, 3.60, 'The €5,000 seed',
     'The talent budget launches product number one.',
     'Packaging and delivery kit for Process X-Ray, a demo on the APA V3 catalog, launch collateral with '
     'marketing, and pricing pressure-tests with three friendly accounts.')
d.footnote(s, 'Quota (~€200K) and pilot-account choice are open PDP items (1.3, 1.4) for the next Mayur 1:1. Keystone scope sits with Deepak (2.2).')
d.notes(s, 'For Tim: the talent budget is seed capital for a product launch. For Mayur: the same asks already on the PDP backlog.')

# ═════════ S34 · ASSUMPTIONS ═════════
s = d.slide()
d.chrome(s, 'Mobilization · planning basis', 'Six explicit assumptions carry the plan, each with an owner')
assumps = [
    ('Team cost ~€1.2M a year', 'PDP working figure', 'Low', 'Mayur + finance'),
    ('Effort estimates within ±30%', 'Analogous VC engagements; first build recalibrates', 'Medium', 'First build'),
    ('SE at 0.5 FTE is available', 'Regional pool capacity, unconfirmed', 'Low', 'SE leadership'),
    ('Tickets of €15-100K clear the market', 'AIB paid ~€80K for a weaker ADS; SAP + AWS wedges', 'Medium', 'First 3 deals'),
    ('Wedge-to-Mission attach 50%+ in 6 months', 'SAP and ServiceNow expansion patterns', 'Low', 'Pilot cohort'),
    ('EAP slots open to our nominations', '8 slots remain; eligibility unconfirmed', 'Low', 'Program owner'),
]
conf_color = {'Low': CORAL, 'Medium': AMBER, 'High': GREEN}
for i, (a, basis, conf, owner) in enumerate(assumps):
    y = 1.95 + i * 0.70
    d.rect(s, 1.0, y, 11.708, 0.60, fill=TINT2, round_=True)
    d.txt(s, 1.24, y + 0.08, 4.1, 0.45, a, size=11, color=NAVY, bold=True, line_sp=1.05)
    d.txt(s, 5.45, y + 0.08, 4.3, 0.45, basis, size=10, color=NAVY, line_sp=1.1)
    d.oval(s, 9.95, y + 0.24, 0.12, 0.12, conf_color[conf])
    d.txt(s, 10.13, y + 0.17, 0.75, 0.24, conf, size=10, color=MUT, bold=True)
    d.txt(s, 10.95, y + 0.08, 1.7, 0.45, 'Validate: ' + owner, size=9, color=BLUE, bold=True, line_sp=1.05)
d.footnote(s, 'Repo standard: every assumption documented, conservative, sensitivity-tested before any figure becomes client-facing or budget-committed.')
d.notes(s, 'No hidden assumptions. Validating the first three in September de-risks the sequence; EAP eligibility gates nominations.')

# ═════════ DIVIDER 06 · SELLING INTO APPREHENSION ═════════
d.divider('06', 'Selling into apprehension',
          'The bank’s buying criteria, the Early Access route around the RFP, and the honest critique.')

# ═════════ S36 · BUYING CRITERIA ═════════
s = d.slide()
d.chrome(s, 'Apprehension · the buying criteria', 'Five questions decide a pre-RFP purchase; the wedge answers all')
crits = [
    ('“Can one sponsor sign it?”', 'Fixed price, sized inside the sponsor’s delegated authority; anything larger phases into rungs.'),
    ('“What is our InfoSec exposure?”', 'Lite modes run on MI packs and schema metadata; customer data never leaves the bank.'),
    ('“What if we stop after this?”', 'The bank owns every artifact outright; the evidence is vendor-neutral and usable in any later tender.'),
    ('“Is the spend defensible to audit?”', 'Fixed scope, named deliverables, a value floor (10× fee or free), full fee credit on conversion.'),
    ('“Why you? You are new to this.”', 'Co-creation instead of procurement: an Early Access nomination with roadmap influence and joint proof.'),
]
for i, (c, a) in enumerate(crits):
    y = 1.95 + i * 0.87
    d.rect(s, 1.0, y, 3.95, 0.77, fill=NAVY, round_=True)
    d.txt(s, 1.22, y + 0.12, 3.55, 0.55, c, size=11, color=WHITE, bold=True, line_sp=1.12)
    d.rect(s, 5.15, y, 7.55, 0.77, fill=TINT2, round_=True)
    d.rect(s, 5.15, y + 0.06, 0.045, 0.65, fill=BLUE)
    d.txt(s, 5.40, y + 0.12, 7.1, 0.55, a, size=10.5, color=NAVY, line_sp=1.15)
d.footnote(s, 'Delegated-authority thresholds vary by bank; price each rung inside one named sponsor’s delegation and validate per account in the free wedge.')
d.notes(s, 'The free wedge exists partly to discover the sponsor’s threshold and the bank’s procurement rules before anything is priced.')

# ═════════ S37 · EARLY ACCESS PROGRAM ═════════
s = d.slide()
d.chrome(s, 'Apprehension · the Early Access route', 'Early Access is the pre-RFP route: nominated, never tendered')
eap = [
    ('Named roadmap influence', 'The bank co-shapes Banking OS: direct input on Nexus, Sentinel and Mission patterns.'),
    ('Fee becomes participation', 'The wedge fee is framed as program participation and credits forward into the Mission POC.'),
    ('Engineering at the table', 'A direct line to Factory engineering during the Mission; issues resolve in days.'),
    ('Joint proof, joint story', 'Co-authored case study and reference rights; the bank is a design partner.'),
]
for i, (h_, b_) in enumerate(eap):
    chip_card(s, 1.0, 1.95 + i * 0.80, 5.9, 0.70, h_, b_)
hero(s, 7.15, 1.95, 5.55, 3.00, 'Nominated, never tendered',
     '8 slots remain in the mid-year cohort.',
     'A bank joins by nomination against program criteria, so there is nothing to RFP. The scarcity is honest: '
     'when the slots are gone, the cohort closes.')
d.open_badge(s, 7.15, 5.15, 5.55, 'Open · ', 'eligibility (new logos vs existing customers) unconfirmed; verify with the program owner first.', h=0.62)
d.footnote(s, 'Slot count (8 remaining) per the mid-year program status. Positioning: new to the category and hungry to co-innovate; the EAP makes that a strength.')
d.notes(s, 'A nomination against program criteria is a selection, so procurement has nothing to tender. Scarcity creates urgency without discounting.')

# ═════════ S38 · PRE-RFP MECHANICS ═════════
s = d.slide()
d.chrome(s, 'Apprehension · the pre-RFP mechanics', 'Five mechanics keep the wedge inside the sponsor’s pen')
mechs = [
    ('1 · A product SKU, never a services engagement', 'Fixed price, fixed scope, catalogued deliverables. Procurement treats it like a software purchase.'),
    ('2 · Sized to the sponsor’s pen', 'Each rung sits inside one executive’s delegated authority; larger ambitions phase across rungs.'),
    ('3 · The fee credits forward, 100%', 'On conversion the wedge fee credits into the Mission POC or licence. The budget is never spent on a maybe.'),
    ('4 · The specimen pack sells before the SoW', 'The free wedge ends by showing the exact artifacts. The bank buys a known object.'),
    ('5 · The bank owns the evidence, floor-guaranteed', 'All artifacts are the bank’s property, vendor-neutral, with a value floor: 10× fee documented, or free.'),
]
for i, (h_, b_) in enumerate(mechs):
    y = 1.95 + i * 0.87
    d.rect(s, 1.0, y, 11.708, 0.77, fill=TINT2, round_=True)
    d.rect(s, 1.0, y + 0.06, 0.045, 0.65, fill=BLUE)
    d.txt(s, 1.24, y + 0.10, 4.35, 0.6, h_, size=11, color=NAVY, bold=True, line_sp=1.1)
    d.txt(s, 5.75, y + 0.12, 6.8, 0.55, b_, size=10.5, color=NAVY, line_sp=1.15)
d.footnote(s, 'The one thing never to do: dodge procurement by making the work free. Free diagnostics re-break the economics; the credit-forward fee is the middle path.')
d.notes(s, 'Mechanics 1-2 prevent the RFP trigger; 3-4 defuse the budget-on-a-maybe fear; 5 defuses the captive-diagnostic objection.')

# ═════════ S39 · CRITIQUE A ═════════
s = d.slide()
d.chrome(s, 'Apprehension · honest critique 1 of 2', 'The math caps year one at four installs; call it a proof year')
tiles = [
    ('Demand of ten installs', '120-140 person-weeks', 'Delivery effort alone, per the register; excludes the build.'),
    ('Pod supply', '≈ 70 person-weeks', '1.5 FTE across build, sell, deliver and the PDP itself.'),
    ('Realistic ceiling', '4-5 installs', 'The honest year-one number, with zero slack for surprises.'),
]
for i, (h_, big, note) in enumerate(tiles):
    x = 1.0 + i * 3.97
    d.rect(s, x, 2.05, 3.77, 1.95, fill=TINT2, round_=True)
    d.txt(s, x + 0.22, 2.25, 3.35, 0.22, h_.upper(), size=9.5, color=BLUE, bold=True)
    d.txt(s, x + 0.22, 2.55, 3.35, 0.6, big, size=19, color=NAVY if i < 2 else CORAL, bold=True, line_sp=1.05)
    d.txt(s, x + 0.22, 3.30, 3.35, 0.55, note, size=10, color=MUT, line_sp=1.18)
d.rect(s, 1.0, 4.35, 11.708, 1.05, fill=TINT, round_=True)
d.txt(s, 1.25, 4.50, 2.7, 0.6, 'The revised plan.', size=12, color=BLUE2, bold=True)
d.txt(s, 4.0, 4.48, 8.5, 0.85, 'Year one proves the model: 3-4 paid installs (€0.4-0.5M), one per product family, plus the '
                               'Early Access cohort. Year two reaches the ten-install cost-neutral rate with two delivery '
                               'consultants funded by year-one revenue.', size=11, color=NAVY, line_sp=1.2)
d.footnote(s, 'Revision to the chapter-1 economics: ten a year remains the cost-neutral bar, reached in year two with delivery consultants and the kits.')
d.notes(s, 'Self-critique per the repo’s conservative-bias standard. Overpromising year one burns the credibility the whole line depends on.')

# ═════════ S40 · CRITIQUE B ═════════
s = d.slide()
d.chrome(s, 'Apprehension · honest critique 2 of 2', 'Five contradictions need a named decision before we sell')
decs = [
    ('Prospect-only rule vs Telemetry', 'The pilot must be a prospect (Mayur’s ask), yet Telemetry needs a live deployment, usually an existing customer.', 'Mayur: scope the rule'),
    ('Credited fees vs the revenue KPI', 'A credited fee is pipeline; an uncredited install is revenue. Pick the line’s primary KPI before quota is set.', 'Mayur + finance: fix the KPI'),
    ('The wedge reads as a tollbooth', 'A €60K gate in front of a platform deal pushes AEs to give discovery away free unless wedge revenue credits the account team.', 'Sales: align incentives'),
    ('EAP eligibility is unconfirmed', 'Whether the 8 slots take new logos, existing customers, or both changes which accounts get nominated first.', 'Program owner: confirm'),
    ('Engine B fights MBB on their turf', 'Generic org design is customer-proximate space we lose; it is defensible only on Engine A’s harvested evidence.', 'You: enforce sequencing'),
]
for i, (h_, b_, owner) in enumerate(decs):
    y = 1.95 + i * 0.87
    d.rect(s, 1.0, y, 11.708, 0.77, fill=TINT2, round_=True)
    d.rect(s, 1.0, y + 0.06, 0.045, 0.65, fill=CORAL)
    d.txt(s, 1.24, y + 0.09, 2.95, 0.6, h_, size=10.5, color=NAVY, bold=True, line_sp=1.08)
    d.txt(s, 4.35, y + 0.09, 5.85, 0.62, b_, size=9.5, color=NAVY, line_sp=1.12)
    d.rect(s, 10.35, y + 0.12, 2.2, 0.52, fill=TINT, round_=True)
    d.txt(s, 10.48, y + 0.17, 1.95, 0.45, owner, size=9, color=BLUE2, bold=True, line_sp=1.08)
d.footnote(s, 'Each row is a real tension surfaced deliberately: resolving them on paper now is cheaper than discovering them inside a live account.')
d.notes(s, 'These five go on the Mayur 1:1 agenda with the four asks. The EAP row also gates the Tim conversation.')

# ═════════ S41 · CLOSE (dark) ═════════
s = d.slide(dark=True)
d.step_glyph(s, 0.406, 0.406, 0.167, 0.166, CYAN)
d.txt(s, 1.0, 2.0, 11.0, 0.35, 'THE MOTION', size=12.5, color=CYAN, bold=True)
d.txt(s, 1.0, 2.45, 11.0, 1.4, [[('Land fast. Prove value. Expand.', 40, WHITE, True)]], line_sp=1.05)
d.txt(s, 1.0, 3.55, 9.0, 0.8, 'One product in market by January. A proof year of 3-4 installs. Cost neutral in '
                              'year two. The platform pre-installed in every account we touch.',
      size=14, color=SUB_D, line_sp=1.3)
steps_now = [
    'This week · book the Mayur 1:1: quota, pilot green-light, SE allocation, the five decisions.',
    'This week · confirm EAP eligibility with the program owner; shortlist nominations.',
    'This week · follow up with Tim: category story, launch collateral, talent-budget plan.',
    'September · scope the keystone cost-per-outcome study with Deepak.',
]
for i, s_ in enumerate(steps_now):
    d.oval(s, 1.02, 4.72 + i * 0.44, 0.11, 0.11, CYAN)
    d.txt(s, 1.28, 4.64 + i * 0.44, 10.5, 0.4, s_, size=11.5, color=WHITE, line_sp=1.15)
d.txt(s, 11.7, 7.05, 1.06, 0.3, 'Backbase', size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
d.notes(s, 'CLOSE. The three this-week actions are the whole ask; everything else in the deck exists to earn them.')

d.save(OUT)
print(f'Wrote {OUT} ({len(d.prs.slides._sldIdLst)} slides)')
