#!/usr/bin/env python3
"""Product Factory execution plan — exhibit-style PPTX (internal, PDP). v3.

Style: McKinsey-exhibit content rules on the Backbase Master Template 2026
chrome (hairline grid, stepped-square mark, master margins), rendered in
Frontline 2026 tokens. Page numbers are `slidenum` fields (auto-renumber).

v3 adds: ladder↔lifecycle bridge · artifact chain with decision gates ·
"one wedge, four lenses" · per-product deep-dives + illustrative sample
outputs (incl. integration landscape) · buying criteria · Early Access
Program as the pre-RFP vehicle · honest critique chapter (capacity math,
open decisions). All specimen figures are placeholders, marked as such.
"""
import copy
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Frontline tokens ─────────────────────────────────────────
NAVY = RGBColor(0x04, 0x13, 0x26)
BLUE = RGBColor(0x33, 0x67, 0xFF)
BLUE_DARK = RGBColor(0x26, 0x4E, 0xC7)
BLUE_LIGHT = RGBColor(0xE5, 0xEB, 0xFF)
OFF = RGBColor(0xF3, 0xF6, 0xF9)
BORDER = RGBColor(0xCE, 0xD2, 0xD7)
MUTED = RGBColor(0x6B, 0x77, 0x86)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CYAN = RGBColor(0x69, 0xFE, 0xFF)
RED = RGBColor(0xFF, 0x50, 0x3C)
AMBER = RGBColor(0xB4, 0x53, 0x09)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
LIGHT_ON_NAVY = RGBColor(0xC5, 0xCF, 0xDE)
RAIL_DARK = RGBColor(0x2E, 0x3D, 0x52)
BODY_DARK = RGBColor(0x2A, 0x38, 0x4A)
T75 = RGBColor(0x66, 0x8D, 0xFF)
T50 = RGBColor(0x99, 0xB3, 0xFF)
T25 = RGBColor(0xCC, 0xD9, 0xFF)

FONT = 'Libre Franklin'
HERE = os.path.dirname(os.path.abspath(__file__))
LOGO_DARK = os.path.join(HERE, '../../../knowledge/design-system/claude-design-exhibit-kit/assets/logo/backbase-wordmark-dark.png')
LOGO_WHITE = os.path.join(HERE, '../../../knowledge/design-system/claude-design-exhibit-kit/assets/logo/backbase-wordmark-white.png')

ML = 54
MR = 1226
AW = MR - ML
RAIL_L, RAIL_R = 37, 1243
TOP_Y, FOOT_Y = 37, 684
CV_RAIL_L, CV_RAIL_R = 53, 1227

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def px(v):
    return Emu(int(v / 96.0 * 914400))


def pt_of(v_px):
    return Pt(v_px * 0.75)


def new_slide(bg=WHITE):
    sl = prs.slides.add_slide(BLANK)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = bg
    return sl


def txt(sl, s, x, y, w, h, size=18, color=NAVY, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, sp=1.15, font=FONT):
    tb = sl.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, line in enumerate(str(s).split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = sp
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = pt_of(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def rect(sl, x, y, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = sl.shapes.add_shape(shape, px(x), px(y), px(w), px(h))
    if round_:
        try:
            sp.adjustments[0] = 0.08
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def hairline(sl, x, y, w, h=1.4, color=BORDER):
    return rect(sl, x, y, w, h, fill=color)


def step_mark(sl, cross_x, cross_y, s=16, color=BLUE, orient='above'):
    if orient == 'above':
        rect(sl, cross_x - s / 2, cross_y - s, s / 2, s, fill=color)
        rect(sl, cross_x - s, cross_y - s / 2, s / 2, s / 2, fill=color)
    else:
        rect(sl, cross_x - s / 2, cross_y, s / 2, s, fill=color)
        rect(sl, cross_x - s, cross_y, s / 2, s / 2, fill=color)


def diamond(sl, cx, cy, size, fill):
    sp = sl.shapes.add_shape(MSO_SHAPE.DIAMOND, px(cx - size / 2), px(cy - size / 2),
                             px(size), px(size))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = WHITE
    sp.line.width = Pt(1.2)
    sp.shadow.inherit = False
    return sp


def dot(sl, cx, cy, d, fill, line=None):
    sp = sl.shapes.add_shape(MSO_SHAPE.OVAL, px(cx - d / 2), px(cy - d / 2), px(d), px(d))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.5)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def logo(sl, dark_bg=False, x=1128, y=690, w=72):
    path = LOGO_WHITE if dark_bg else LOGO_DARK
    if os.path.exists(path):
        sl.shapes.add_picture(path, px(x), px(y), width=px(w))


def page_field(sl, color=MUTED, x=1216, y=692, w=48, h=22, size=12):
    """Auto-updating slide number (`slidenum` field)."""
    tb = sl.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = '1'
    r.font.name = FONT
    r.font.size = pt_of(size)
    r.font.color.rgb = color
    rPr = r._r.find(qn('a:rPr'))
    fld = r._r.makeelement(qn('a:fld'), {
        'id': '{1F4E2DE4-8ADA-4D4E-9951-90A1D26586E7}', 'type': 'slidenum'})
    if rPr is not None:
        fld.append(copy.deepcopy(rPr))
    t = fld.makeelement(qn('a:t'), {})
    t.text = '1'
    fld.append(t)
    r._r.getparent().replace(r._r, fld)
    return tb


def chrome(sl, kicker, title, footnote=None, notes=None, title_w=1080):
    hairline(sl, 0, TOP_Y, 1280)
    hairline(sl, RAIL_L, 0, 1.4, h=FOOT_Y)
    hairline(sl, RAIL_R, 0, 1.4, h=FOOT_Y)
    hairline(sl, 0, FOOT_Y, 1280)
    step_mark(sl, RAIL_L, TOP_Y, 16, BLUE, 'above')
    txt(sl, kicker.upper(), ML, 54, 900, 24, size=12.5, color=BLUE, bold=True)
    txt(sl, title, ML, 82, title_w, 110, size=29, bold=True, sp=1.08)
    logo(sl, dark_bg=False)
    page_field(sl)
    if footnote:
        hairline(sl, ML, FOOT_Y - 36, MR - ML)
        txt(sl, footnote, ML, FOOT_Y - 29, MR - ML, 26, size=11.5, color=MUTED, sp=1.05)
    if notes:
        sl.notes_slide.notes_text_frame.text = notes


def cover_chrome(sl):
    hairline(sl, CV_RAIL_L, 0, 1.4, h=720, color=RAIL_DARK)
    hairline(sl, CV_RAIL_R, 0, 1.4, h=720, color=RAIL_DARK)
    hairline(sl, CV_RAIL_L, 155, CV_RAIL_R - CV_RAIL_L, color=RAIL_DARK)
    hairline(sl, 829, 0, 1.4, h=720, color=RAIL_DARK)
    hairline(sl, CV_RAIL_L, 499, 829 - CV_RAIL_L, color=RAIL_DARK)
    step_mark(sl, 829, 155, 16, WHITE, 'below')


def chapter_chrome(sl):
    hairline(sl, 0, 155, 1280, color=RAIL_DARK)
    hairline(sl, CV_RAIL_L, 0, 1.4, h=720, color=RAIL_DARK)
    hairline(sl, CV_RAIL_R, 0, 1.4, h=720, color=RAIL_DARK)
    hairline(sl, CV_RAIL_L, 467, CV_RAIL_R - CV_RAIL_L, color=RAIL_DARK)
    step_mark(sl, CV_RAIL_L, 155, 16, WHITE, 'above')
    logo(sl, dark_bg=True)
    page_field(sl, color=RGBColor(0x5A, 0x6B, 0x80))


def chapter(number, title, subtitle):
    sl = new_slide(NAVY)
    chapter_chrome(sl)
    txt(sl, f'{number} · CHAPTER', 107, 212, 700, 24, size=12.5, color=CYAN, bold=True)
    txt(sl, title, 105, 250, 1000, 80, size=40, color=WHITE, bold=True)
    txt(sl, subtitle, 107, 350, 980, 60, size=15, color=LIGHT_ON_NAVY, sp=1.25)
    return sl


def chip(sl, x, y, w, h, head, body, fill=OFF, head_color=NAVY, body_color=None,
         accent=None):
    rect(sl, x, y, w, h, fill=fill, line=BORDER, line_w=0.75, round_=True)
    if accent:
        rect(sl, x, y + 8, 4, h - 16, fill=accent)
    txt(sl, head, x + 20, y + 12, w - 36, 26, size=14, bold=True, color=head_color)
    txt(sl, body, x + 20, y + 40, w - 36, h - 50, size=12, color=body_color or BODY_DARK, sp=1.12)


def hero_card(sl, x, y, w, h, eyebrow, big, body):
    rect(sl, x, y, w, h, fill=NAVY, round_=True)
    txt(sl, eyebrow.upper(), x + 30, y + 26, w - 60, 24, size=12, color=CYAN, bold=True)
    txt(sl, big, x + 30, y + 58, w - 60, 120, size=21, color=WHITE, bold=True, sp=1.12)
    txt(sl, body, x + 30, y + h - 118, w - 60, 104, size=12.5, color=LIGHT_ON_NAVY, sp=1.2)


def col_card(sl, x, y, w, h, head, bullets, accent=BLUE):
    rect(sl, x, y, w, h, fill=OFF, line=BORDER, line_w=0.75, round_=True)
    rect(sl, x, y + 8, 4, h - 16, fill=accent)
    txt(sl, head.upper(), x + 22, y + 14, w - 40, 24, size=11.5, color=BLUE, bold=True)
    txt(sl, '\n'.join('· ' + b for b in bullets), x + 22, y + 42, w - 40, h - 56,
        size=11.5, color=BODY_DARK, sp=1.22)


def band(sl, y, lead, body, h=76):
    rect(sl, ML, y, AW, h, fill=BLUE_LIGHT, round_=True)
    txt(sl, lead, ML + 24, y + 14, 300, h - 26, size=13, bold=True, color=BLUE_DARK, sp=1.1)
    txt(sl, body, ML + 336, y + 12, AW - 360, h - 22, size=12, sp=1.18)


def panel(sl, x, y, w, h, eyebrow):
    rect(sl, x, y, w, h, fill=WHITE, line=BORDER, line_w=1.0, round_=True)
    txt(sl, eyebrow.upper(), x + 18, y + 12, w - 36, 22, size=10.5, color=BLUE, bold=True)
    return x + 18, y + 40, w - 36


ILLUSTRATIVE = ('Illustrative specimen: every figure on this slide is a placeholder, '
                're-based on the bank’s own data during the engagement.')

# ═════════════════════════════════════════════════════════════
# S1 — Cover
# ═════════════════════════════════════════════════════════════
sl = new_slide(NAVY)
cover_chrome(sl)
logo(sl, dark_bg=True, x=96, y=64, w=120)
txt(sl, 'VALUE CONSULTING · PDP · INTERNAL', 105, 220, 700, 24, size=12.5, color=CYAN, bold=True)
txt(sl, 'Product Factory:\nfrom concept to revenue', 102, 262, 720, 170, size=44,
    color=WHITE, bold=True, sp=1.04)
txt(sl, 'The execution plan for four paid product wedges that pre-install Banking OS:\n'
        'prove the model in year one, reach cost neutrality in year two, open a recurring line.',
    102, 528, 710, 80, size=14, color=LIGHT_ON_NAVY, sp=1.3)
txt(sl, 'Shyam · July 2026 · prepared for the talent programme (Tim Ruttner) and the PDP track (Mayur)',
    102, 645, 700, 26, size=11.5, color=MUTED)
txt(sl, 'One product in market\nby January 2027.', 880, 220, 320, 90, size=18,
    color=WHITE, bold=True, sp=1.2)
txt(sl, 'Four sellable SKUs, the Engine B services\nflagship, an Early Access cohort, and a\nrecurring line by mid 2027.',
    880, 320, 330, 80, size=13, color=LIGHT_ON_NAVY, sp=1.25)

# ═════════════════════════════════════════════════════════════
# S2 — Chapter 01
# ═════════════════════════════════════════════════════════════
chapter('01', 'The frameworks',
        'The thesis, the product filter, the economics, the ladder and lifecycle, and the pod that runs it.')

# ═════════════════════════════════════════════════════════════
# S3 — Thesis
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Frameworks · the thesis',
       'Every product earns a fee, harvests the bank’s reality, and installs the '
       'first component of Banking OS.',
       footnote='Filter rule: a candidate that fails any one of the three tests is consulting, '
                'and stays out of the product line. Aligned to the Banking OS canon (knowledge/product/banking-os.md).',
       notes='The three-part design filter from the 29 Jun PDP session. Product-proximate = moat; '
             'customer-proximate = competitive space we lose.')
for i, (h_, b_) in enumerate([
    ('Test 1 · It earns a fee', 'Each install is paid work at €15-100K, priced as a product with a fixed scope and a delivery kit.'),
    ('Test 2 · It harvests proprietary data', 'The engagement captures the bank’s real processes, data structures and policies. Evidence only Backbase holds.'),
    ('Test 3 · It pre-binds a Banking OS layer', 'The output is a working component: Nexus, Sentinel or a Mission candidate, already present in the account.'),
]):
    chip(sl, ML, 222 + i * 122, 560, 108, h_, b_, accent=BLUE)
hero_card(sl, 650, 222, 576, 348, 'The one-liner',
          'Paid discovery that leaves the platform installed.',
          'Each engagement ends with a live Banking OS component in the account, a quantified '
          'expansion case for the AE, and account intelligence competitors and SIs never see.')

# ═════════════════════════════════════════════════════════════
# S3b — Two engines, one wedge (Engine A + Engine B)
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Frameworks · the two engines',
       'The products harvest the evidence; the services redesign the bank around it.',
       footnote='Engine B passes the same three tests when fed by Engine A: it earns a fee, harvests the '
                'organisation’s reality, and pre-binds the human layer of Banking OS. Ticket range is an '
                'assumption to validate against market org-design pricing.',
       notes='Engine B (from the PDP backlog, workstream 4): AI Maturity Assessment, Workforce Optimization, '
             'and the flagship AI-Native Org Design (draft v0.1 exists). The sequencing rule is the moat: '
             'MBB can run generic org design; only we can run it on the bank’s own autonomy evidence.')
rect(sl, ML, 208, 560, 300, fill=BLUE_LIGHT, round_=True)
txt(sl, 'ENGINE A · PRODUCT FACTORY', ML + 26, 230, 500, 22, size=11.5, color=BLUE_DARK, bold=True)
txt(sl, 'Installs the technical layers', ML + 26, 256, 500, 30, size=16, bold=True)
for i, line in enumerate(['Four SKUs: X-Ray, Cartographer, Guardrail, Telemetry',
                          'Pre-binds Nexus, Sentinel and Mission candidates',
                          'Earns €15-100K per install; fees credit forward',
                          'Harvests the account evidence nobody else holds']):
    txt(sl, '· ' + line, ML + 26, 296 + i * 30, 510, 26, size=11.5, color=BODY_DARK, sp=1.1)
txt(sl, '→ the evidence engine', ML + 26, 470, 500, 22, size=11, color=BLUE, bold=True)
rect(sl, 650, 208, 576, 300, fill=NAVY, round_=True)
txt(sl, 'ENGINE B · AI-NATIVE SERVICES', 676, 230, 520, 22, size=11.5, color=CYAN, bold=True)
txt(sl, 'Installs the human layer', 676, 256, 520, 30, size=16, color=WHITE, bold=True)
for i, line in enumerate(['Three offerings: Maturity Assessment, Workforce Optimization, and the flagship AI-Native Org Design (draft v0.1 exists)',
                          'Redesigns the org chart around the A1-A5 autonomy curve: pyramid → inverted-T',
                          'Pre-binds the roles that run Banking OS: Mission Owners, AgentOps, Exception Desks',
                          'Earns €80-150K per engagement (to validate)']):
    txt(sl, '· ' + line, 676, 296 + i * 42, 526, 40, size=11.5, color=LIGHT_ON_NAVY, sp=1.12)
txt(sl, '→ the operating-model engine', 676, 472, 500, 22, size=11, color=CYAN, bold=True)
band(sl, 532, 'The sequencing rule.',
     'Engine B never sells cold: it runs on Engine A’s harvested evidence (autonomy scores, authority maps, '
     'cost per outcome), which is exactly what a generic consultancy cannot copy.', h=72)

# ═════════════════════════════════════════════════════════════
# S4 — Economics dot grid
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Frameworks · the economics',
       'Ten paid installs a year cover the full team cost; every install beyond is new revenue.',
       footnote='Package = Assess & Solution €60K + Mission POC €60K. Team cost of ~€1.2M a year is the PDP working '
                'figure, to be validated with finance. Chapter 04 phases this: year one proves with 3-4 installs; '
                'ten a year is the cost-neutral run rate from year two.',
       notes='Cost-neutrality is the floor, revenue generation is the objective. See critique chapter for the '
             'capacity-checked phasing.')
gx, gy, d, gap = 100, 290, 42, 30
for i in range(16):
    row, col = divmod(i, 8)
    dot(sl, gx + col * (d + gap) + d / 2, gy + row * (d + gap) + d / 2, d,
        BLUE if i < 10 else T25)
dot(sl, gx + 10, 480, 16, BLUE)
txt(sl, 'Installs to cost neutrality (10)', gx + 28, 470, 300, 24, size=12, color=MUTED)
dot(sl, gx + 370, 480, 16, T25)
txt(sl, 'Year-two expansion range (14-16)', gx + 388, 470, 320, 24, size=12, color=MUTED)
rect(sl, 766, 240, 460, 260, fill=OFF, round_=True)
txt(sl, '10', 806, 268, 180, 90, size=64, color=BLUE, bold=True)
txt(sl, 'installs a year at an average package of €120K cover the ~€1.2M annual team cost.',
    806, 370, 380, 100, size=14.5, sp=1.25)

# ═════════════════════════════════════════════════════════════
# S5 — Package ladder
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Frameworks · the ladder',
       'The wedge stays free; depth and recurrence carry the price.',
       footnote='Recurring range is a target shape anchored on SAP and ServiceNow success-plan pricing of 10-30% of '
                'licence value, and on FinOps gain-share market practice (15-35% of realised savings).',
       notes='Rung 0 free = funnel. Rungs 1+2 = the €120K SKVC package. Rung 3 = the growth flywheel.')
y0 = 290
chip(sl, ML, y0 - 34, 240, 130, 'Rung 0 · free wedge',
     'Ignite Inspire\n~4 meetings, benchmark and use-case shortlist. Grows the funnel.', fill=OFF)
txt(sl, '→', 306, y0 + 10, 30, 40, size=20, color=MUTED, align=PP_ALIGN.CENTER)
bar_x, bar_w = 350, 560
rect(sl, bar_x, y0, int(bar_w / 2), 76, fill=BLUE)
txt(sl, 'Assess & Solution\n€60K', bar_x + 18, y0 + 12, 240, 56, size=13, color=WHITE, bold=True, sp=1.1)
rect(sl, bar_x + int(bar_w / 2), y0, int(bar_w / 2), 76, fill=BLUE_DARK)
txt(sl, 'Mission POC\n€60K', bar_x + int(bar_w / 2) + 18, y0 + 12, 240, 56, size=13, color=WHITE, bold=True, sp=1.1)
txt(sl, 'Rungs 1 + 2 · the €120K package (to scale)', bar_x, y0 - 32, 560, 24, size=12.5, color=MUTED, bold=True)
txt(sl, '→', 922, y0 + 10, 30, 40, size=20, color=MUTED, align=PP_ALIGN.CENTER)
rect(sl, 964, y0 - 34, 262, 130, fill=NAVY, round_=True)
txt(sl, 'Rung 3 · recurring', 986, y0 - 22, 220, 24, size=12, color=CYAN, bold=True)
txt(sl, 'Value Assurance\n€30-50K a year + gain share', 986, y0 + 4, 220, 84, size=13, color=WHITE, bold=True, sp=1.15)
chip(sl, ML, 460, 572, 92, 'The free wedge feeds the funnel',
     'Four meetings surface the leakage evidence that qualifies the paid wedge. Zero friction to start.', accent=BLUE)
chip(sl, 654, 460, 572, 92, 'Recurrence is the growth flywheel',
     'Telemetry converts every live deployment into quarterly revenue that compounds with AI consumption.', accent=BLUE)

# ═════════════════════════════════════════════════════════════
# S6 — NEW · Ladder ↔ lifecycle bridge
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Frameworks · ladder meets lifecycle',
       'The ladder is what the bank buys; the lifecycle is how the pod delivers it.',
       footnote='Stages 05 Proposal and 06 Signature are the arrows between rungs: each rung ends by proposing and '
                'signing the next one. The wedge product fulfils Rung 1; the Factory Mission Sprint fulfils Rung 2.',
       notes='The bridge slide: read vertically, each rung IS a lifecycle stage the bank pays for. '
             'Read horizontally, each rung’s exit artifact is the entry contract of the next.')
rungs = [
    ('RUNG 0 · FREE', 'Ignite Inspire', '~4 meetings: hypotheses, wedge choice, specimen pack', OFF, NAVY, MUTED),
    ('RUNG 1 · €60K', 'The wedge product', 'X-Ray, Cartographer, Guardrail or Telemetry, on the bank’s own data', BLUE_LIGHT, NAVY, BLUE_DARK),
    ('RUNG 2 · €60K', 'Mission POC', 'Factory Mission Sprint: the top candidate goes live in 6-12 weeks', BLUE, WHITE, WHITE),
    ('RUNG 3 · RECURRING', 'Value Assurance', 'Quarterly cost-per-outcome proof; funds the next domain', NAVY, WHITE, CYAN),
]
rw, rgap, ry = 250, 57, 218
txt(sl, 'WHAT THE BANK BUYS', ML, 192, 400, 20, size=10.5, color=MUTED, bold=True)
for i, (tag, name, desc, fill, tcol, tagcol) in enumerate(rungs):
    x = ML + i * (rw + rgap)
    rect(sl, x, ry, rw, 126, fill=fill, round_=True,
         line=BORDER if fill in (OFF, BLUE_LIGHT) else None, line_w=0.75)
    txt(sl, tag, x + 18, ry + 10, rw - 32, 20, size=10, color=tagcol, bold=True)
    txt(sl, name, x + 18, ry + 32, rw - 32, 26, size=14.5, color=tcol, bold=True)
    txt(sl, desc, x + 18, ry + 60, rw - 32, 60, size=10.5,
        color=tcol if fill in (BLUE, NAVY) else BODY_DARK, sp=1.15)
    if i < 3:
        ax = x + rw + 6
        txt(sl, '→', ax, ry + 40, rgap - 12, 30, size=16, color=MUTED, align=PP_ALIGN.CENTER)
        txt(sl, '05 propose\n06 sign', ax, ry + 74, rgap - 12, 40, size=8.5, color=BLUE,
            bold=True, align=PP_ALIGN.CENTER, sp=1.1)
txt(sl, 'HOW THE POD DELIVERS', ML, 388, 400, 20, size=10.5, color=MUTED, bold=True)
stage_groups = ['01 Concept · 02 Design', '03 Prototype · 04 Proof', '07 Install', 'Run · measure · expand']
for i, sg in enumerate(stage_groups):
    x = ML + i * (rw + rgap)
    hairline(sl, x + rw / 2, ry + 126, 1.4, h=412 - ry - 126, color=BORDER)
    rect(sl, x, 412, rw, 52, fill=WHITE, line=BLUE, line_w=1.2, round_=True)
    txt(sl, sg, x + 14, 412, rw - 28, 52, size=11.5, color=BLUE_DARK, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
band(sl, 496,
     'Read it both ways.',
     'Vertically: each rung of the ladder is the lifecycle stage the bank pays for. Horizontally: each rung’s '
     'exit artifact is written as the entry contract of the next purchase, so the account never buys a maybe.')

# ═════════════════════════════════════════════════════════════
# S7 — Lifecycle step flow
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Frameworks · the lifecycle',
       'One motion carries every product from concept to installed platform.',
       footnote='Durations are per-engagement figures once the product is built, with a ±30% planning tolerance. '
                'The first pass through stages 1-4 is the one-time product build.',
       notes='Stage exit gates: Concept = qualified sponsor + pain. Design = signed SoW, agreed metrics, data access. '
             'Prototype = working demo on the client’s own data. Proof = CFO-grade evidence pack. '
             'Proposal = commercial offer + named Mission candidate. Signature = signed order form. '
             'Install = layer live + expansion roadmap handed to the AE.')
steps = [
    ('01', 'Concept', 'VC · 1 wk'), ('02', 'Design', 'VC + SE · 1-2 wks'),
    ('03', 'Prototype', 'SE · 2-3 wks'), ('04', 'Proof', 'VC · 2-4 wks'),
    ('05', 'Proposal', 'VC + AE · 1 wk'), ('06', 'Signature', 'AE + VC · 2-4 wks'),
    ('07', 'Install', 'FDE + SE · 6-12 wks'),
]
sw, sgap, sx, sy = 156, 13, ML, 258
for i, (n, name, who) in enumerate(steps):
    x = sx + i * (sw + sgap)
    rect(sl, x, sy, sw, 150, fill=OFF, line=BORDER, line_w=0.75, round_=True)
    dot(sl, x + 32, sy + 34, 34, BLUE)
    txt(sl, n, x + 15, sy + 23, 34, 24, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, name, x + 16, sy + 62, sw - 30, 26, size=14, bold=True)
    txt(sl, who, x + 16, sy + 92, sw - 30, 46, size=10.5, color=MUTED, sp=1.15)
    if i < 6:
        txt(sl, '→', x + sw - 2, sy + 58, 20, 30, size=14, color=MUTED)
band(sl, 460, 'The lifecycle runs twice over.',
     'The first pass through stages 1-4 builds the product itself. From the second account onward, the same '
     'motion runs as a 4-8 week sales-and-delivery cycle per install.', h=84)

# ═════════════════════════════════════════════════════════════
# S8 — NEW · Artifact chain + decision gates
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Frameworks · connecting the dots',
       'Every deliverable is the entry contract of the next purchase; four gates carry the account to install.',
       footnote='Nothing in the chain is a shelf report: each artifact is authored as the input of the next purchase '
                'order, which is what makes the wedge convert instead of inform.',
       notes='G1-G4 are the account-level decision gates. Internally each product also has its own build gates '
             '(SE named, pilot green-light, R&D alignment) shown on the roadmap slide.')
arts = [
    ('AFTER INSPIRE (FREE)', 'The Inspire pack',
     ['Leakage hypotheses', 'Wedge choice + scope', 'Specimen output pack']),
    ('AFTER THE WEDGE', 'The evidence pack',
     ['€ leakage, evidenced', 'Integration landscape map', 'Ranked Mission backlog']),
    ('AFTER THE POC', 'The live loop',
     ['Resolution loop in production', 'Telemetry baseline', 'Sentinel-governed, auditable']),
    ('EVERY QUARTER', 'The value proof',
     ['Cost per outcome, measured', 'Savings register', 'Next Mission ranked + funded']),
]
aw_, agap, ay = 250, 57, 210
for i, (eyebrow, name, bullets) in enumerate(arts):
    x = ML + i * (aw_ + agap)
    rect(sl, x, ay, aw_, 200, fill=OFF, line=BORDER, line_w=0.75, round_=True)
    txt(sl, eyebrow, x + 18, ay + 12, aw_ - 32, 20, size=9.5, color=BLUE, bold=True)
    txt(sl, name, x + 18, ay + 36, aw_ - 32, 26, size=14.5, bold=True)
    txt(sl, '\n'.join('· ' + b for b in bullets), x + 18, ay + 68, aw_ - 32, 120,
        size=11, color=BODY_DARK, sp=1.25)
    if i < 3:
        gx_ = x + aw_ + agap / 2
        diamond(sl, gx_, ay + 100, 22, BLUE)
        txt(sl, f'G{i + 1}', gx_ - 20, ay + 118, 40, 20, size=10, color=BLUE, bold=True,
            align=PP_ALIGN.CENTER)
gates = [
    ('G1', 'Sponsor signs the wedge SoW, inside their own delegated authority.'),
    ('G2', 'ExCo funds Mission #1; the wedge fee credits into the POC.'),
    ('G3', 'Measured value meets the case; the licence conversation opens.'),
    ('G4', 'The next domain is funded; the loop repeats on the same Banking OS.'),
]
for i, (g, t_) in enumerate(gates):
    x = ML + (i % 2) * (AW / 2 + 12)
    y = 452 + (i // 2) * 52
    diamond(sl, x + 12, y + 20, 16, BLUE)
    txt(sl, g, x + 28, y + 10, 36, 22, size=11.5, color=BLUE, bold=True)
    txt(sl, t_, x + 62, y + 10, AW / 2 - 80, 40, size=11.5, color=BODY_DARK, sp=1.1)

# ═════════════════════════════════════════════════════════════
# S9 — The pod
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Frameworks · resourcing',
       'A pod of two and a half people runs the whole line; hiring follows revenue.',
       footnote='FTE figures are steady-state allocations during build and delivery phases; FDE capacity is drawn '
                'from standard Mission Sprint staffing at install, so it is charged to the delivery, and not to this pod.',
       notes='Allocation ask, specifically NOT a hiring ask. SE is the critical unlock.')
pods = [
    ('You · product owner · 1.0 FTE', 'SKU design, value narrative, proof economics, proposal. Owns the P&L of the line.'),
    ('Solution engineer · 0.5 FTE, named', 'Connectors, data extraction, Factory tooling at prototype and install. The technical spine.'),
    ('FDE · on demand at install', 'Runs the Factory Mission Sprint once an install is signed. 6-12 weeks per Mission.'),
    ('SME · on call', 'Risk and compliance review for Guardrail Studio; R&D liaison for the Cartographer.'),
]
for i, (h_, b_) in enumerate(pods):
    chip(sl, ML, 216 + i * 92, 560, 80, h_, b_, accent=BLUE)
hero_card(sl, 650, 216, 576, 356, 'Why the solution engineer matters',
          'One named SE across all four products, so skill compounds.',
          'Every product runs on the same Factory tools: Connector Studio, Semantic Modeler, '
          'Simulation & Testing. Each install makes the next one cheaper and faster to deliver.')

# ═════════════════════════════════════════════════════════════
# S10 — Objections
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Frameworks · the objections',
       'The three obvious objections have concrete answers.',
       footnote='AIB reference: paid an estimated ~€80K for an external current/future-state assessment (ADS), PDP session '
                'record, 29 Jun 2026. Wedge pattern references: AWS OLA, SAP Outside-In, ServiceNow Inspire.',
       notes='Use these in the Tim and Mayur conversations.')
qa = [
    ('“Banks will not pay for discovery.”',
     'The market already does: AIB paid ~€80K for a weaker assessment. AWS, SAP and ServiceNow all run '
     'the diagnostic-wedge motion; ours goes further by installing product.'),
    ('“We have no capacity to build products.”',
     'Build effort is 6-12 person-weeks per product because the Factory tools, the APA V3 catalog and the '
     'cost model already exist. The sequence builds one product at a time.'),
    ('“This distracts from platform sales.”',
     'Every install pre-binds Nexus, Sentinel or a Mission and hands the AE a quantified expansion case. '
     'The wedge is the land motion of land-and-expand.'),
]
for i, (c, a) in enumerate(qa):
    y = 222 + i * 122
    rect(sl, ML, y, 420, 108, fill=NAVY, round_=True)
    txt(sl, c, ML + 24, y + 18, 372, 76, size=14, color=WHITE, bold=True, sp=1.15)
    rect(sl, 494, y, 732, 108, fill=OFF, round_=True)
    rect(sl, 494, y + 8, 4, 92, fill=BLUE)
    txt(sl, a, 518, y + 14, 684, 84, size=12.5, sp=1.18)

# ═════════════════════════════════════════════════════════════
# S11 — Chapter 02
# ═════════════════════════════════════════════════════════════
chapter('02', 'The execution plans',
        'Four products and the services flagship: what each does, what its outputs look like, and the plan to market.')

# ═════════════════════════════════════════════════════════════
# S12 — NEW · One wedge, four lenses
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Execution plans · how the four fit together',
       'A bank never buys one layer; it buys a Mission, and every wedge surveys all four lenses.',
       footnote='The entry door follows the account’s live buying trigger. Coverage of all four lenses is standard '
                'in every engagement, because a Mission needs truth, guardrails and telemetry at once.',
       notes='Direct answer to “why would a bank choose only one of these layers, that too from us”: they choose a '
             'Mission. The Blueprint always covers the stack; the wedge only decides which lens goes deep. Ontology '
             'and guardrail lenses ship as thin appendices in every engagement from day one.')
doors = [
    ('Ops pain is the trigger', 'Enter via Process X-Ray. Deep: work + leakage. Thin: ontology, guardrails, telemetry.'),
    ('Data ambition (RI) is the trigger', 'Enter via Ontology Cartographer. Deep: shared truth. Thin: work, guardrails, telemetry.'),
    ('Risk is blocking AI', 'Enter via Guardrail Studio. Deep: authority + policy. Thin: work, ontology, telemetry.'),
    ('AI is live, ROI is doubted', 'Enter via Value Telemetry. Deep: measured value. Thin: work, ontology, guardrails.'),
]
for i, (h_, b_) in enumerate(doors):
    chip(sl, ML, 216 + i * 92, 560, 80, h_, b_, accent=BLUE)
hero_card(sl, 650, 216, 576, 356, 'One deep lens, three thin lenses',
          'Every engagement ships all four lenses; the wedge sets the depth.',
          'Each wedge includes an ontology appendix, an authority snapshot and telemetry hooks alongside its '
          'deep dive. The account ends every engagement holding a full-stack Banking OS blueprint, at the '
          'depth its first Mission needs.')

# ─────────────────────────────────────────────────────────────
# Product blocks
# ─────────────────────────────────────────────────────────────


def deep_dive(kicker, title, activities, entails, outputs, addup, footnote, notes=None):
    sl = new_slide()
    chrome(sl, kicker, title, footnote=footnote, notes=notes)
    col_card(sl, ML, 200, 450, 320, 'The activities, week by week', activities)
    col_card(sl, ML + 474, 200, 340, 320, 'What it entails', entails)
    col_card(sl, ML + 838, 200, 334, 320, 'What comes out', outputs)
    band(sl, 540, 'How it adds up.', addup, h=72)
    return sl


def product_plan(kicker, title, milestones, stats, footnote, notes=None):
    sl = new_slide()
    chrome(sl, kicker, title, footnote=footnote, notes=notes)
    n = len(milestones)
    line_y = 290
    x0, x1 = 190, 1090
    hairline(sl, x0, line_y, x1 - x0, 2.4, color=BORDER)
    span = (x1 - x0) / (n - 1)
    for i, (date, label, body, gate) in enumerate(milestones):
        cx = x0 + i * span
        if gate == 'decision':
            diamond(sl, cx, line_y + 1, 22, RED)
        elif gate == 'ga':
            diamond(sl, cx, line_y + 1, 22, BLUE)
        else:
            dot(sl, cx, line_y + 1, 16, BLUE, line=WHITE)
        txt(sl, date, cx - 100, line_y - 40, 200, 24, size=12, color=BLUE, bold=True,
            align=PP_ALIGN.CENTER)
        txt(sl, label, cx - 100, line_y + 26, 200, 26, size=13.5, bold=True,
            align=PP_ALIGN.CENTER)
        txt(sl, body, cx - 100, line_y + 54, 200, 100, size=10.5, color=MUTED,
            align=PP_ALIGN.CENTER, sp=1.12)
    band_y = 462
    rect(sl, ML, band_y, AW, 96, fill=NAVY, round_=True)
    cw = AW / len(stats)
    for i, (lab, val) in enumerate(stats):
        bx = ML + i * cw
        txt(sl, lab.upper(), bx + 28, band_y + 16, cw - 40, 22, size=10.5, color=CYAN, bold=True)
        txt(sl, val, bx + 28, band_y + 42, cw - 40, 46, size=14.5, color=WHITE, bold=True, sp=1.05)
    return sl


# ══════════ Process X-Ray ══════════
deep_dive(
    'Execution plan 1 of 4 · Process X-Ray · what it is',
    'Process X-Ray reconstructs how work really flows, and prices the leakage in it.',
    ['Wk 1 · scope 2-3 journeys from the APA matrix; agree access and baselines',
     'Wk 1-2 · inventory every system, interface and handoff per journey (the integration landscape)',
     'Wk 2-4 · extract logs and queue data; shadow the ops floor (6-8 sessions); rebuild actual flows vs designed flows',
     'Wk 4-5 · price the leakage at every step: time × volume × cost, abandonment × revenue',
     'Wk 5-6 · rank Mission candidates; exec readout with the CFO-grade evidence'],
    ['Sponsor: COO or head of operations',
     'Bank effort: SME time of ~6-10 hours a week',
     'Access, lite mode: ops MI packs + sampled case files; no InfoSec cycle, prospect-safe',
     'Access, instrumented mode: read-only logs; for existing customers',
     'Pod: VC 1.0 + SE 0.5'],
    ['€ p.a. leakage per journey, evidenced',
     'Integration landscape map: systems, interfaces, API readiness per journey',
     'Ranked Mission backlog, top 3, with ROI case each',
     'Thin-lens appendices: ontology snapshot + authority snapshot',
     'The Rung-2 POC proposal, pre-written'],
    'The systems map scopes the Connector Studio work before anyone commits; the heatmap becomes the account’s '
    'demand map for every later product and Mission.',
    'Lite mode exists because prospects rarely grant system access pre-contract: the workshop-and-MI-pack variant '
    'trades precision for zero InfoSec friction. Aligned to the value-leakage method in banking-os.md §9.',
    notes='The productized Process & Workspace Designer. ~60% of bank work lives between systems; the X-Ray '
          'makes that whitespace visible and priceable.')

sl = new_slide()
chrome(sl, 'Execution plan 1 of 4 · Process X-Ray · sample outputs',
       'The X-Ray hands the sponsor three artifacts: the heatmap, the systems map, and the Mission card.',
       footnote=ILLUSTRATIVE,
       notes='Show these in the free wedge as the specimen pack: the bank buys a known object, never a maybe.')
p1x, p1y, p1w = panel(sl, ML, 200, 374, 330, 'Value-leakage heatmap · € M p.a.')
cols_ = ['Onb', 'Serv', 'Disp', 'Lend']
rows_ = ['Retail', 'SME', 'Comm']
vals = [[0.8, 1.1, 2.4, 1.6], [0.5, 0.9, 1.2, 1.4], [0.3, 0.6, 0.8, 1.1]]
cw_, ch_ = 60, 38
for j, c in enumerate(cols_):
    txt(sl, c, p1x + 86 + j * (cw_ + 2), p1y + 6, cw_, 18, size=9.5, color=MUTED, bold=True,
        align=PP_ALIGN.CENTER)
for i, rname in enumerate(rows_):
    y = p1y + 28 + i * (ch_ + 3)
    txt(sl, rname, p1x, y + 9, 80, 20, size=10, color=MUTED, bold=True)
    for j in range(4):
        v = vals[i][j]
        fill = RED if v >= 2.0 else (BLUE_DARK if v >= 1.3 else (BLUE if v >= 0.9 else (T50 if v >= 0.6 else T25)))
        rect(sl, p1x + 86 + j * (cw_ + 2), y, cw_, ch_, fill=fill)
        txt(sl, f'{v}', p1x + 86 + j * (cw_ + 2), y + 8, cw_, 20, size=10.5,
            color=WHITE if v >= 0.6 else NAVY, bold=True, align=PP_ALIGN.CENTER)
txt(sl, 'Disputes · Retail leads: €2.4M p.a. across 41% manual touches and a 9.2-day cycle.',
    p1x, p1y + 170, p1w, 60, size=10.5, color=BODY_DARK, sp=1.2)
txt(sl, '→ picks the domain the Mission starts in', p1x, p1y + 244, p1w, 20, size=10,
    color=BLUE, bold=True)
p2x, p2y, p2w = panel(sl, ML + 399, 200, 374, 330, 'Integration landscape · card dispute journey')
systems = [('Core (T24)', 'API'), ('CRM (SFDC)', 'API'), ('Fraud (Falcon)', 'file'),
           ('Cases (Pega)', 'manual'), ('Mail (O365)', 'manual'), ('Pay hub', 'file')]
for i, (name, mode) in enumerate(systems):
    x = p2x + (i % 2) * 172
    y = p2y + 6 + (i // 2) * 56
    rect(sl, x, y, 160, 46, fill=OFF, line=BORDER, line_w=0.75, round_=True)
    txt(sl, name, x + 10, y + 5, 144, 20, size=10.5, bold=True)
    mcol = GREEN if mode == 'API' else (AMBER if mode == 'file' else RED)
    dot(sl, x + 16, y + 33, 8, mcol)
    txt(sl, mode + ' interface', x + 26, y + 25, 120, 18, size=9, color=MUTED)
txt(sl, '8 systems (+2 not shown) · 3 manual swivels · 2 API-ready · 4 to build via Connector Studio',
    p2x, p2y + 184, p2w, 40, size=10.5, color=BODY_DARK, sp=1.2)
txt(sl, '→ scopes the integration work before commitment', p2x, p2y + 244, p2w, 20,
    size=10, color=BLUE, bold=True)
rect(sl, ML + 798, 200, 374, 330, fill=NAVY, round_=True)
txt(sl, 'MISSION CANDIDATE #1', ML + 816, 214, 340, 20, size=10.5, color=CYAN, bold=True)
txt(sl, 'Card dispute resolution loop', ML + 816, 240, 340, 52, size=16, color=WHITE, bold=True, sp=1.1)
for i, line in enumerate(['Value: €1.9M p.a., evidenced', 'Autonomy target: A4, execute by exception',
                          'Feasibility: high · 12 integration points', '8 live via APIs, 4 built in the Sprint',
                          'Mission length: 8 weeks']):
    txt(sl, '· ' + line, ML + 816, 300 + i * 30, 340, 26, size=11.5, color=LIGHT_ON_NAVY, sp=1.1)
txt(sl, '→ becomes the Rung-2 POC scope, fee credited', ML + 816, 490, 340, 24, size=10.5,
    color=CYAN, bold=True)

product_plan(
    'Execution plan 1 of 4 · Process X-Ray · the plan',
    'Process X-Ray ships first: demo by October, paid proof by December, first revenue in January.',
    [
        ('Aug 2026', 'Design', 'Delivery kit, pricing one-pager, demo script. 3 pw, VC.', None),
        ('Sep-Oct 2026', 'Prototype', 'Demo on the APA V3 catalog (65 processes) plus one friendly dataset. 5 pw, SE-led.', None),
        ('Oct 2026', 'Pilot gate', 'Prospect pilot green-light with Mayur. AIB-profile account.', 'decision'),
        ('Oct-Dec 2026', 'Proof', 'Paid pilot: leakage heatmap, top-3 Mission candidates with ROI evidence. 6 pw.', None),
        ('Jan 2027', 'First signature', 'Wedge sold at €50-75K; install starts as a Factory Mission Sprint.', 'ga'),
    ],
    [('Build effort', '8-10 person-wks'), ('Delivery per install', '4-6 wks · 12-14 pw'),
     ('Pod', 'VC 1.0 + SE 0.5'), ('Pre-binds', 'Mission candidate + ROI')],
    'Effort carries a ±30% tolerance. Build is low-risk because the APA V3 catalog (65 processes, 283 steps), the '
    'journey builder and the value-leakage method already exist in the repo.')

# ══════════ Ontology Cartographer ══════════
deep_dive(
    'Execution plan 2 of 4 · Ontology Cartographer · what it is',
    'The Cartographer maps scattered data to one shared truth, and prices the gap.',
    ['Wk 1-2 · inventory the data landscape: core, CRM, LMS, credit, workflow',
     'Wk 1-2 · harvest schemas and data dictionaries (metadata only)',
     'Wk 3-5 · map entities to the banking ontology; surface overlaps and conflicts',
     'Wk 5-7 · score truth-gaps per entity; price them in ops time and error cost',
     'Wk 7-8 · phased Nexus binding blueprint with a precise entity scope'],
    ['Sponsor: CDO or CIO',
     'Bank effort: data stewards 4-6 hours a week',
     'Access: schema-level metadata; customer data never leaves the bank, which keeps InfoSec light',
     'R&D liaison: Semantic Modeler and Nexus teams',
     'Pod: SE 1.0 + VC 0.5 (SE-led)'],
    ['Ontology coverage map, entity by entity',
     'Truth-gap report, priced in ops time',
     'Data-source integration inventory',
     'Phased Nexus blueprint: what binds first, and why',
     'Thin-lens appendices: leakage + authority snapshots'],
    'Shared truth is the precondition for Relationship Intelligence and any A4-plus Mission, which is why every '
    'other wedge carries an ontology appendix; this SKU takes the same lens to build depth.',
    'Metadata-only access is the buying-criteria answer: no customer data leaves the bank, so the InfoSec review '
    'stays at the schema level. Nexus is the system of truth, never the system of record.',
    notes='Productizes the Semantic Modeler. The ontology is key for any use case, so the lens ships thin in '
          'every wedge from day one; the full SKU is the deep version.')

sl = new_slide()
chrome(sl, 'Execution plan 2 of 4 · Ontology Cartographer · sample outputs',
       'The Cartographer shows the bank where its truth fragments, and what binding it first is worth.',
       footnote=ILLUSTRATIVE,
       notes='The coverage map is the conversation starter; the truth-gap card carries the money argument.')
p1x, p1y, p1w = panel(sl, ML, 200, 374, 330, 'Ontology coverage map · by entity')
ents = [('Customer', '6 sources · 4 conflicts', RED), ('Account', '3 sources · aligned', GREEN),
        ('Product', '4 sources · 2 conflicts', AMBER), ('Transaction', '2 sources · aligned', GREEN),
        ('Collateral', 'fragmented', AMBER), ('Case', 'absent as an entity', RED)]
for i, (e, s_, c) in enumerate(ents):
    x = p1x + (i % 2) * 172
    y = p1y + 6 + (i // 2) * 62
    rect(sl, x, y, 160, 52, fill=OFF, round_=True)
    rect(sl, x, y + 6, 4, 40, fill=c)
    txt(sl, e, x + 14, y + 6, 140, 20, size=11, bold=True)
    txt(sl, s_, x + 14, y + 27, 140, 20, size=9, color=MUTED)
txt(sl, '→ four definitions of “active customer” is the exhibit that lands', p1x, p1y + 214,
    p1w, 36, size=10, color=BLUE, bold=True, sp=1.15)
p2x, p2y, p2w = panel(sl, ML + 399, 200, 374, 330, 'Truth-gap report · the cost of fragmentation')
txt(sl, '11 min', p2x, p2y + 8, p2w, 54, size=40, color=BLUE, bold=True)
txt(sl, 'for a relationship manager to assemble one customer view across 5 screens',
    p2x, p2y + 66, p2w, 40, size=11.5, color=BODY_DARK, sp=1.2)
for i, line in enumerate(['23% of fields disagree between systems', '€1.1M p.a. of RM time spent reassembling truth',
                          '31% of complaints touch a data mismatch']):
    dot(sl, p2x + 6, p2y + 132 + i * 32, 8, RED)
    txt(sl, line, p2x + 20, p2y + 122 + i * 32, p2w - 20, 28, size=11, color=BODY_DARK)
txt(sl, '→ the € case for shared truth, before any licence talk', p2x, p2y + 240, p2w, 20,
    size=10, color=BLUE, bold=True)
p3x, p3y, p3w = panel(sl, ML + 798, 200, 374, 330, 'Nexus binding blueprint · phased')
phases = [('Phase 1 · 12 wks', 'Core + CRM · 40 entities · customer, account, product', BLUE, WHITE),
          ('Phase 2 · 10 wks', 'Credit + LMS · +25 entities · exposure, collateral', T50, NAVY),
          ('Phase 3 · next', 'Cases + documents · the service layer of truth', T25, NAVY)]
for i, (ph, de, fill, tcol) in enumerate(phases):
    y = p3y + 6 + i * 74
    rect(sl, p3x, y, p3w, 64, fill=fill, round_=True)
    txt(sl, ph, p3x + 16, y + 8, p3w - 30, 22, size=11.5, color=tcol, bold=True)
    txt(sl, de, p3x + 16, y + 32, p3w - 30, 26, size=10,
        color=WHITE if fill == BLUE else BODY_DARK, sp=1.1)
txt(sl, '→ makes the Nexus licence scope precise, so pricing stops being a guess',
    p3x, p3y + 240, p3w, 30, size=10, color=BLUE, bold=True, sp=1.1)

product_plan(
    'Execution plan 2 of 4 · Ontology Cartographer · the plan',
    'The Cartographer is the deepest build and lands where X-Ray has created demand.',
    [
        ('Feb 2027', 'R&D gate', 'Alignment with the Semantic Modeler and Nexus teams on tooling reuse.', 'decision'),
        ('Mar-Apr 2027', 'Prototype', 'Ontology mapping demo on two sample cores (core, CRM, LMS, credit). 6 pw, SE-led.', None),
        ('May-Jul 2027', 'Proof', 'Paid pilot: ontology coverage map, data-to-truth gap report, Nexus blueprint. 8 pw.', None),
        ('Aug 2027', 'First signature', 'Sold at €75-100K, attached to Relationship Intelligence opportunities.', 'ga'),
    ],
    [('Build effort', '10-12 person-wks'), ('Delivery per install', '6-8 wks · 12-16 pw'),
     ('Pod', 'SE 1.0 + VC 0.5 (SE-led)'), ('Pre-binds', 'Nexus')],
    'Sequenced last as a full SKU because it has the highest technical depth and needs R&D alignment; the ontology '
    'lens itself ships thin inside every other wedge from day one (see “one wedge, four lenses”).')

# ══════════ Guardrail Studio ══════════
deep_dive(
    'Execution plan 3 of 4 · Guardrail Studio · what it is',
    'Guardrail Studio turns written policy into executable rules an agent can be trusted with.',
    ['Wk 1-2 · harvest delegation of authority, credit policy, SOPs, entitlement matrices',
     'Wk 2-3 · map authority per process step: who may recommend, approve, execute',
     'Wk 3-4 · score agent-readiness per step on the A1-A5 autonomy scale',
     'Wk 4-5 · codification workshops (4-5) with risk and compliance: policy → executable rules',
     'Wk 5-6 · audit-trail design; Sentinel readiness score and deployment blueprint'],
    ['Sponsor: CRO or COO',
     'Bank effort: risk, compliance and ops SMEs in 4-5 workshops',
     'Access: documents and workshops; no system access needed, so InfoSec stays out of the critical path',
     'Pod: VC 1.0 + SE 0.5 + compliance SME 0.2'],
    ['Authority map: current state vs agent-ready state',
     'Guardrail backlog: 20-30 executable rules, each traced to its policy source',
     'Audit-trail design for immutable evidence',
     'Sentinel readiness score + install blueprint',
     'Thin-lens appendices: leakage + ontology snapshots'],
    'This is the approval pack a risk committee needs before ANY agent acts; it unblocks Transact and Resolve '
    'deals in flight, and every codified rule pre-loads Sentinel.',
    'Timed to the Conversational Banking pipeline: governed execution is the prerequisite for Transact and '
    'Resolve, which makes this wedge a natural attach on active deals.',
    notes='Productizes Decision & Policy. Sentinel: authority (who may decide) · policies (under which rules) · '
          'entitlements (multi-entity access), with observability and immutable auditability.')

sl = new_slide()
chrome(sl, 'Execution plan 3 of 4 · Guardrail Studio · sample outputs',
       'Guardrail Studio answers the question every risk committee asks: where can an agent safely act?',
       footnote=ILLUSTRATIVE,
       notes='The rule card is the artifact that changes the conversation: policy as executable, auditable logic.')
p1x, p1y, p1w = panel(sl, ML, 200, 374, 330, 'Authority map · card dispute process')
amap = [('Freeze card', 'A4 · agent executes', BLUE), ('Provisional credit ≤ €500', 'A4 · executes with audit', BLUE),
        ('Provisional credit > €500', 'A3 · human approves', T50), ('Account closure', 'A1 · human only', BORDER)]
for i, (step, aut, c) in enumerate(amap):
    y = p1y + 6 + i * 58
    rect(sl, p1x, y, p1w, 48, fill=OFF, round_=True)
    txt(sl, step, p1x + 14, y + 6, 220, 20, size=10.5, bold=True)
    rect(sl, p1x + 14, y + 28, 96, 14, fill=c, round_=True)
    txt(sl, aut, p1x + 118, y + 25, p1w - 130, 20, size=9.5, color=MUTED)
txt(sl, '→ the boundary line, drawn from the bank’s own policy', p1x, p1y + 246, p1w, 20,
    size=10, color=BLUE, bold=True)
p2x, p2y, p2w = panel(sl, ML + 399, 200, 374, 330, 'Guardrail rule card · from the backlog of 20-30')
rect(sl, p2x - 4, p2y + 4, p2w + 8, 180, fill=NAVY, round_=True)
for i, line in enumerate(['GR-014 · Provisional credit', 'IF dispute.amount ≤ €500', 'AND fraud_score < 0.2',
                          'THEN agent EXECUTES, immutable audit log', 'ELSE route to human approver']):
    txt(sl, line, p2x + 14, p2y + 18 + i * 30, p2w - 24, 26, size=11,
        color=CYAN if i == 0 else WHITE, bold=(i == 0), sp=1.1)
txt(sl, 'Source: Card Policy §4.2 · Status: codified · Owner: dispute ops',
    p2x, p2y + 196, p2w, 24, size=10, color=MUTED)
txt(sl, '→ each rule loads straight into Sentinel at install', p2x, p2y + 240, p2w, 20,
    size=10, color=BLUE, bold=True)
p3x, p3y, p3w = panel(sl, ML + 798, 200, 374, 330, 'Sentinel readiness · scored')
txt(sl, '34 / 100', p3x, p3y + 4, p3w, 50, size=36, color=RED, bold=True)
bars = [('Authority mapped', 45), ('Policies codified', 20), ('Audit trails in place', 38)]
for i, (lab, v) in enumerate(bars):
    y = p3y + 66 + i * 52
    txt(sl, f'{lab} · {v}%', p3x, y, p3w, 20, size=10.5, color=BODY_DARK, bold=True)
    rect(sl, p3x, y + 24, p3w, 10, fill=OFF)
    rect(sl, p3x, y + 24, p3w * v / 100, 10, fill=BLUE)
txt(sl, '→ target: 80+ once the backlog is codified; the risk committee’s go/no-go dashboard',
    p3x, p3y + 228, p3w, 36, size=10, color=BLUE, bold=True, sp=1.15)

product_plan(
    'Execution plan 3 of 4 · Guardrail Studio · the plan',
    'Guardrail Studio attaches to every Conversational Banking deal in flight.',
    [
        ('Jan 2027', 'Design', 'Decision & Policy tooling alignment; risk and compliance SME briefed. 3 pw.', None),
        ('Feb-Mar 2027', 'Prototype', 'Authority-map demo: recommend, approve, execute per process. 5 pw.', None),
        ('Apr-May 2027', 'Proof', 'Paid pilot: guardrail backlog, entitlements map, Sentinel install blueprint. 6 pw.', None),
        ('Jun 2027', 'First signature', 'Sold at €50-100K, riding Transact and Resolve deals in flight.', 'ga'),
    ],
    [('Build effort', '8-10 person-wks'), ('Delivery per install', '5-6 wks · 8-10 pw'),
     ('Pod', 'VC 1.0 + SE 0.5 + SME 0.2'), ('Pre-binds', 'Sentinel')],
    'Governed execution is the prerequisite for Transact and Resolve, which makes the Sentinel-readiness '
    'wedge a natural attach on active Conversational Banking deals.')

# ══════════ Value Telemetry ══════════
deep_dive(
    'Execution plan 4 of 4 · Value Telemetry · what it is',
    'Value Telemetry measures what live AI really costs per resolved outcome, and re-proves the ROI.',
    ['Wk 1 · instrument the live deployment: interaction logs, model calls, consumption',
     'Wk 1-2 · tag outcomes: resolved, escalated, abandoned; define the outcome taxonomy',
     'Wk 2 · compute cost per outcome per use case (never per interaction)',
     'Wk 2-3 · scan for waste: PTU over-provision, model routing, dead agents, retry loops',
     'Wk 3 · re-prove ROI vs the original case; install the quarterly cadence'],
    ['Prerequisite: a live AI deployment (Backbase or third party)',
     'Access: read-only consumption + interaction logs',
     'Bank effort: finance partner to validate cost lines',
     'Pod: VC 0.5 + SE 0.25',
     'Dependency: the keystone cost-per-outcome model (with Deepak)'],
    ['Cost-per-outcome baseline per use case',
     'Waste register with € values and named owners',
     'ROI re-proof vs the original business case',
     'A quarterly value cadence, installed',
     'Thin-lens appendices: leakage + guardrail observations'],
    'The baseline becomes the yardstick every later Mission is measured against, and the identified waste '
    'often funds the next wedge by itself.',
    'Only the vendor in the execution path can measure cost per resolved outcome; external FinOps tools see '
    'the cloud bill. Gain-share terms (15-35% of realised savings) per market practice: ProsperOps, nOps, Vantage.',
    notes='Backlog items 2.2 (keystone) + 2.3 (SKU). Partner with Aayushi on the value-assurance line.')

sl = new_slide()
chrome(sl, 'Execution plan 4 of 4 · Value Telemetry · sample outputs',
       'Telemetry turns “is the AI worth it” into a measured number the CFO can sign.',
       footnote=ILLUSTRATIVE,
       notes='The waste register is the fastest trust-builder: it usually pays for the diagnostic several times over.')
p1x, p1y, p1w = panel(sl, ML, 200, 374, 330, 'Cost per outcome · by use case')
tiles = [('Card freeze', '€0.11 / outcome', '94% contained, zero human touches', GREEN),
         ('Dispute intake', '€1.87 / outcome', '3.2 human touches; 68% of simple intents routed to the premium model', RED)]
for i, (uc, cost, note, c) in enumerate(tiles):
    y = p1y + 6 + i * 122
    rect(sl, p1x, y, p1w, 112, fill=OFF, round_=True)
    rect(sl, p1x, y + 8, 4, 96, fill=c)
    txt(sl, uc, p1x + 16, y + 10, p1w - 30, 22, size=11.5, bold=True)
    txt(sl, cost, p1x + 16, y + 34, p1w - 30, 30, size=18,
        color=BLUE_DARK if c == GREEN else RED, bold=True)
    txt(sl, note, p1x + 16, y + 68, p1w - 30, 40, size=10, color=MUTED, sp=1.15)
txt(sl, '→ the unit economics of every use case, measured', p1x, p1y + 250, p1w, 20,
    size=10, color=BLUE, bold=True)
p2x, p2y, p2w = panel(sl, ML + 399, 200, 374, 330, 'Waste register · €312K annualized')
wastes = [('PTU over-provision', 140), ('Model routing', 95), ('Dead agents', 41), ('Retry loops', 36)]
for i, (lab, v) in enumerate(wastes):
    y = p2y + 10 + i * 50
    txt(sl, f'{lab} · €{v}K', p2x, y, p2w, 20, size=10.5, color=BODY_DARK, bold=True)
    rect(sl, p2x, y + 24, p2w, 12, fill=OFF)
    rect(sl, p2x, y + 24, p2w * v / 140, 12, fill=RED if i == 0 else BLUE)
txt(sl, '→ each line has an owner and a fix; savings are machine-measurable, which is what '
        'makes gain-share defensible', p2x, p2y + 218, p2w, 48, size=10, color=BLUE, bold=True, sp=1.15)
p3x, p3y, p3w = panel(sl, ML + 798, 200, 374, 330, 'ROI re-proof · case vs measured')
rois = [('Original business case', 1.10, T50), ('Measured today', 0.82, RED),
        ('After routing + PTU fixes', 1.05, BLUE)]
for i, (lab, v, c) in enumerate(rois):
    y = p3y + 10 + i * 62
    txt(sl, f'{lab} · €{v:.2f}M a year', p3x, y, p3w, 20, size=10.5, color=BODY_DARK, bold=True)
    rect(sl, p3x, y + 24, p3w, 14, fill=OFF)
    rect(sl, p3x, y + 24, p3w * v / 1.10, 14, fill=c)
rect(sl, p3x, p3y + 200, 210, 30, fill=BLUE_LIGHT, round_=True)
txt(sl, 'Case re-proved at 95%', p3x + 14, p3y + 206, 190, 20, size=10.5, color=BLUE_DARK, bold=True)
txt(sl, '→ renewal and expansion, argued with the bank’s own numbers', p3x, p3y + 244, p3w, 24,
    size=10, color=BLUE, bold=True)

product_plan(
    'Execution plan 4 of 4 · Value Telemetry · the plan',
    'Value Telemetry rides the keystone cost-per-outcome model and opens the recurring line.',
    [
        ('Sep 2026', 'Keystone gate', 'Extend the cost model from cost-per-interaction to cost-per-outcome, with Deepak. ~5 use cases.', 'decision'),
        ('Sep 2026', 'Design', 'SKU spec: diagnostic scope, telemetry set, gain-share terms. 4 pw, VC.', None),
        ('Oct 2026', 'Prototype', 'Cost-per-outcome dashboard on one live deployment’s data. 3 pw.', None),
        ('Nov-Dec 2026', 'Proof', '2-3 week diagnostic at a live customer: waste quantified, ROI re-proved. 3 pw.', None),
        ('Jan 2027', 'First signature', '€15-50K diagnostic + 15-35% gain share; quarterly recurrence begins.', 'ga'),
    ],
    [('Build effort', '6-8 person-wks'), ('Delivery per install', '2-3 wks · 3-4 pw'),
     ('Pod', 'VC 0.5 + SE 0.25'), ('Model', 'Recurring + gain share')],
    'Gain-share pricing of 15-35% of realised savings is market-validated (ProsperOps, nOps, Vantage), with a '
    '“10% or free” floor. Current cost-model sample is n=2 use cases; the keystone study raises it to ~5.')

# ══════════ Engine B · AI-Native Services ══════════
deep_dive(
    'Execution plan 5 · Engine B · AI-Native Services · what it is',
    'Engine B redesigns the bank’s org chart for a world where agents do the work.',
    ['Maturity assessment (2-3 wks) · score every LOB × ops area on the A1-A5 autonomy curve; readiness on data, guardrails, skills',
     'Workforce optimization (3-4 wks) · which roles rise up the value chain, which collapse into agents; in-loop, on-loop, above-loop boundaries per function',
     'Org design flagship (4-6 wks) · target inverted-T org: a thin high-judgment human layer over an agent workforce, bound by a governance spine',
     'Transition path · function by function up the curve; the org chart follows the highest-autonomy function'],
    ['Sponsor: CEO, COO or CHRO; ExCo workshops',
     'Inputs: Engine A evidence: X-Ray autonomy scores, Guardrail authority map, telemetry baselines',
     'HR data: role inventory, spans and layers',
     'Pod: VC 1.0 + org/HR SME on call',
     'Method: AI-Native Org Design draft v0.1 exists'],
    ['A1-A5 maturity heatmap by function, current vs target',
     'Workforce transition map: roles rising, collapsing, new',
     'Target org chart: the inverted-T, with the governance spine',
     'New-role definitions: Mission Owners, AgentOps, Exception Desks',
     'Function-by-function transition roadmap'],
    'Engine B opens the door Engine A cannot: the CEO and CHRO at operating-model altitude. An organisation '
    'designed around Missions renews the platform that runs them.',
    'Sequencing rule: sold only downstream of a wedge install, so the design runs on the bank’s own evidence. '
    'Ticket of €80-150K is an assumption to validate against market org-design pricing.',
    notes='PDP backlog workstream 4. The autonomy scale is the operating-model axis: A1 assist (human does) '
          'through A5 self-directed (human governs). New roles per the backlog: Chief Agentic Operations '
          'Officer, Mission Owners, AgentOps, Nexus Stewards, Sentinel Governance, Exception Desks, '
          'Workforce Transition office.')

sl = new_slide()
chrome(sl, 'Execution plan 5 · Engine B · sample outputs',
       'Engine B shows the bank the org it will need, drawn from its own autonomy evidence.',
       footnote=ILLUSTRATIVE,
       notes='The pyramid-to-inverted-T exhibit is the flagship visual: the org chart follows the autonomy curve.')
p1x, p1y, p1w = panel(sl, ML, 200, 374, 330, 'Autonomy maturity heatmap · current → target')
funcs = [('Disputes', 'A2', 'A4'), ('Onboarding', 'A1', 'A3'), ('KYC review', 'A2', 'A4'),
         ('Contact centre', 'A2', 'A4'), ('Credit servicing', 'A1', 'A3')]
for i, (f, cur, tgt) in enumerate(funcs):
    y = p1y + 6 + i * 46
    rect(sl, p1x, y, p1w, 38, fill=OFF, round_=True)
    txt(sl, f, p1x + 14, y + 9, 150, 20, size=10.5, bold=True)
    rect(sl, p1x + 168, y + 8, 42, 22, fill=T50, round_=True)
    txt(sl, cur, p1x + 168, y + 10, 42, 18, size=10, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, '→', p1x + 216, y + 9, 24, 20, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    rect(sl, p1x + 244, y + 8, 42, 22, fill=BLUE, round_=True)
    txt(sl, tgt, p1x + 244, y + 10, 42, 18, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(sl, '→ scored from X-Ray + Guardrail evidence, never from interviews alone', p1x, p1y + 244,
    p1w, 36, size=10, color=BLUE, bold=True, sp=1.12)
p2x, p2y, p2w = panel(sl, ML + 399, 200, 374, 330, 'The org shift · pyramid → inverted-T')
pyr_cx = p2x + 80
for i, w_ in enumerate([60, 110, 160]):
    rect(sl, pyr_cx - w_ / 2, p2y + 20 + i * 34, w_, 28, fill=T50)
txt(sl, 'Today: layers of\nhuman processing', pyr_cx - 80, p2y + 130, 160, 40, size=9.5,
    color=MUTED, align=PP_ALIGN.CENTER, sp=1.1)
txt(sl, '→', p2x + 158, p2y + 60, 30, 30, size=16, color=MUTED, align=PP_ALIGN.CENTER)
it_cx = p2x + 262
rect(sl, it_cx - 80, p2y + 20, 160, 28, fill=BLUE)
txt(sl, 'Human judgment layer', it_cx - 78, p2y + 25, 156, 18, size=8.5, color=WHITE,
    bold=True, align=PP_ALIGN.CENTER)
rect(sl, it_cx - 20, p2y + 52, 40, 68, fill=NAVY)
txt(sl, 'Governance + AgentOps spine', it_cx - 78, p2y + 130, 160, 40, size=9.5,
    color=MUTED, align=PP_ALIGN.CENTER, sp=1.1)
rect(sl, p2x, p2y + 180, p2w, 34, fill=BLUE_LIGHT, round_=True)
txt(sl, 'Agent workforce running Missions underneath', p2x + 12, p2y + 188, p2w - 24, 20,
    size=10, color=BLUE_DARK, bold=True)
txt(sl, '→ the flagship exhibit for the CEO and CHRO conversation', p2x, p2y + 244, p2w, 24,
    size=10, color=BLUE, bold=True)
p3x, p3y, p3w = panel(sl, ML + 798, 200, 374, 330, 'New roles on the target chart')
roles = [('Mission Owners', 'own agent-delivered outcomes per LOB'),
         ('AgentOps', 'deploy, monitor, optimize the fleet; home of FinOps'),
         ('Nexus Stewards', 'data-as-product; keepers of shared truth'),
         ('Sentinel Governance', 'AI assurance, authority, audit'),
         ('Exception Desks', 'the re-skilled processing floor')]
for i, (r_, d_) in enumerate(roles):
    y = p3y + 6 + i * 48
    rect(sl, p3x, y, p3w, 40, fill=OFF, round_=True)
    rect(sl, p3x, y + 5, 4, 30, fill=BLUE)
    txt(sl, r_, p3x + 14, y + 4, 150, 20, size=10.5, bold=True)
    txt(sl, d_, p3x + 150, y + 5, p3w - 160, 34, size=9, color=MUTED, sp=1.05)
txt(sl, '→ every role binds the org to a Banking OS construct', p3x, p3y + 254, p3w, 24,
    size=10, color=BLUE, bold=True)

product_plan(
    'Execution plan 5 · Engine B · the plan',
    'Engine B hardens the v0.1 method now, and sells its first engagement downstream of X-Ray.',
    [
        ('Sep 2026', 'Method hardening', 'Pressure-test org-design v0.1 with Mayur; codify the delivery kit. 4 pw.', None),
        ('Oct-Dec 2026', 'Evidence binding', 'Wire the method to X-Ray and Guardrail outputs: autonomy scores feed the org model.', None),
        ('Q1 2027', 'Sequencing gate', 'First sale only downstream of a wedge install; never sold cold.', 'decision'),
        ('Mar-May 2027', 'First engagement', 'Maturity + workforce + org design at the X-Ray pilot account. 10-14 pw.', None),
        ('Jun 2027', 'Flagship scale', 'CEO/CHRO-altitude offer; Early Access cohort accounts first.', 'ga'),
    ],
    [('Build effort', '6-8 pw (v0.1 exists)'), ('Delivery per engagement', '6-8 wks · 10-14 pw'),
     ('Pod', 'VC 1.0 + org SME'), ('Pre-binds', 'The human layer of Banking OS')],
    'Ticket of €80-150K is an assumption to validate: anchor against market org-design pricing and the value '
    'of the workforce decisions it informs. The moat is the evidence feed, never the method alone.')

# ═════════════════════════════════════════════════════════════
# S25 — Register table
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Execution plans · the register',
       'Twelve months of product build fits inside the existing pod.',
       footnote='pw = person-weeks. Build = one-time productization effort. Delivery = elapsed and effort figures per '
                'client install once productized. All efforts carry a ±30% planning tolerance.',
       notes='Presentation order is by layer logic (work, truth, authority, proof); build order optimizes risk and '
             'is shown on the roadmap.')
cols = [('Product', 350), ('One-time build', 190), ('Delivery per install', 230),
        ('Pod', 250), ('First revenue', 150)]
rows = [
    ('Process X-Ray', '8-10 pw', '4-6 wks · 12-14 pw', 'VC 1.0 + SE 0.5', 'Jan 2027'),
    ('Ontology Cartographer', '10-12 pw', '6-8 wks · 12-16 pw', 'SE 1.0 + VC 0.5', 'Aug 2027'),
    ('Guardrail Studio', '8-10 pw', '5-6 wks · 8-10 pw', 'VC 1.0 + SE 0.5 + SME', 'Jun 2027'),
    ('Value Telemetry', '6-8 pw', '2-3 wks · 3-4 pw', 'VC 0.5 + SE 0.25', 'Jan 2027'),
    ('AI-Native Org Design (Engine B)', '6-8 pw (v0.1 exists)', '6-8 wks · 10-14 pw', 'VC 1.0 + org SME', 'May 2027'),
]
tx, ty, rh = ML, 230, 62
x = tx
for name, w in cols:
    txt(sl, name.upper(), x, ty, w - 16, 24, size=11.5, color=MUTED, bold=True)
    x += w
hairline(sl, tx, ty + 30, sum(w for _, w in cols), 2.2, color=NAVY)
for r, row in enumerate(rows):
    y = ty + 44 + r * rh
    x = tx
    for c, val in enumerate(row):
        w = cols[c][1]
        txt(sl, val, x, y + 12, w - 16, 40, size=13.5, bold=(c == 0),
            color=NAVY if c == 0 else BODY_DARK)
        x += w
    hairline(sl, tx, y + rh - 6, sum(w for _, w in cols))

# ═════════════════════════════════════════════════════════════
# S26 — Chapter 03
# ═════════════════════════════════════════════════════════════
chapter('03', 'Mobilization',
        'The twelve-month sequence, how the Delivery Factory carries it, and the asks that unlock it.')

# ═════════════════════════════════════════════════════════════
# S27 — Portfolio wave roadmap
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Mobilization · the sequence',
       'One product at a time reaches four sellable SKUs and recurring revenue inside twelve months.',
       footnote='Red diamonds are decision gates (SE named · pilot green-light · R&D alignment). Blue diamonds mark first '
                'revenue per product. The ontology and guardrail lenses ship thin inside X-Ray from the first pilot.',
       notes='Sequencing: X-Ray first (lowest build risk, front-door product), Telemetry parallel (shares the '
             'keystone), Guardrail rides the Conversational pipeline, Cartographer lands on X-Ray demand.')
months = ['Aug', 'Sep', 'Oct', 'Nov', 'Dec', 'Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul']
gx0, gy0, colw, rowh = 260, 242, 80, 58
for i, m in enumerate(months):
    txt(sl, m, gx0 + i * colw, gy0 - 30, colw, 22, size=10.5, color=MUTED, bold=True,
        align=PP_ALIGN.CENTER)
    hairline(sl, gx0 + i * colw, gy0 - 6, 1, h=6 * rowh, color=RGBColor(0xE4, 0xE8, 0xEE))
txt(sl, '2026', gx0, gy0 - 52, 200, 20, size=10, color=MUTED)
txt(sl, '2027', gx0 + 5 * colw, gy0 - 52, 200, 20, size=10, color=MUTED)
grows = [
    ('Keystone study', [(1, 2, BLUE_DARK, 'With Deepak')], [(1.0, RED)]),
    ('Process X-Ray', [(0, 1, T50, 'Design'), (1, 2, BLUE, 'Prototype'), (3, 2, BLUE_DARK, 'Proof'), (5, 7, NAVY, 'Sell + install')], [(2.0, RED), (5.0, BLUE)]),
    ('Value Telemetry', [(1, 1, T50, ''), (2, 1, BLUE, ''), (3, 2, BLUE_DARK, 'Proof'), (5, 7, NAVY, 'Recurring')], [(5.0, BLUE)]),
    ('Guardrail Studio', [(5, 1, T50, ''), (6, 2, BLUE, 'Prototype'), (8, 2, BLUE_DARK, 'Proof'), (10, 2, NAVY, 'Sell')], [(10.0, BLUE)]),
    ('Ontology Cartographer', [(6, 1, T50, ''), (7, 2, BLUE, 'Prototype'), (9, 3, BLUE_DARK, 'Proof')], [(6.0, RED)]),
    ('Org Design (Engine B)', [(1, 1, T50, ''), (2, 3, BLUE, 'Evidence binding'), (7, 3, BLUE_DARK, 'First engagement')], [(7.0, RED), (10.0, BLUE)]),
]
for r, (label, bars, gates) in enumerate(grows):
    y = gy0 + r * rowh
    txt(sl, label, ML, y + 10, 190, 44, size=12, bold=True, sp=1.05)
    for (start, dur, color, blabel) in bars:
        bx = gx0 + start * colw + 2
        bw = dur * colw - 4
        rect(sl, bx, y + 8, bw, 26, fill=color, round_=True)
        if blabel and bw > 120:
            txt(sl, blabel, bx + 10, y + 11, bw - 16, 20, size=9.5, color=WHITE, bold=True)
    for (gm, gcol) in gates:
        diamond(sl, gx0 + gm * colw, y + 21, 16, gcol)

# ═════════════════════════════════════════════════════════════
# S28 — Delivery Factory usage
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Mobilization · the Delivery Factory',
       'The Delivery Factory carries every stage from prototype to installed layer.',
       footnote='Factory toolchain per the Banking OS canon: Process & Workspace Designer · Semantic Modeler · Agent '
                'Builder · Decision & Policy · Connector Studio · Simulation & Testing · Deployment & Ops Control.',
       notes='The Factory already contains the discovery-end tools; the products ARE those tools productized. '
             'Delivery cost stays at person-weeks because nothing is built from scratch.')
stages = [
    ('Prototype', 'Semantic Modeler +\nConnector Studio', 'A working demo on the bank’s own data, in weeks.'),
    ('Proof', 'Simulation & Testing +\nOps Control telemetry', 'CFO-grade evidence runs with zero production risk.'),
    ('Install', 'Mission Contract with\nNexus + Sentinel binding', 'The Banking OS layer live in 6-12 weeks, via the agentic SDLC.'),
    ('Run', 'Deployment &\nOps Control', 'Live telemetry that feeds Value Telemetry and the expansion case.'),
]
cw, cgap = 275, 24
for i, (stage, tools, out_) in enumerate(stages):
    x = ML + i * (cw + cgap)
    rect(sl, x, 230, cw, 280, fill=OFF, line=BORDER, line_w=0.75, round_=True)
    txt(sl, stage.upper(), x + 22, 252, cw - 40, 24, size=11.5, color=BLUE, bold=True)
    txt(sl, tools, x + 22, 282, cw - 40, 70, size=14.5, bold=True, sp=1.12)
    hairline(sl, x + 22, 368, cw - 44)
    txt(sl, out_, x + 22, 382, cw - 40, 110, size=12, color=BODY_DARK, sp=1.2)
    if i < 3:
        txt(sl, '→', x + cw + 2, 350, 20, 30, size=15, color=MUTED)

# ═════════════════════════════════════════════════════════════
# S29 — The asks
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Mobilization · the ask',
       'Four decisions unlock execution; the €5,000 talent budget seeds product number one.',
       footnote='Quota figure of ~€200K recognition and pilot-account choice are open PDP items (1.3, 1.4) for the '
                'next 1:1 with Mayur. Keystone study scope sits with Deepak (backlog item 2.2).',
       notes='Framing for Tim: the talent budget is seed capital for a product launch. Framing for Mayur: these '
             'are the same asks already on the PDP backlog.')
asks = [
    ('1 · Name the solution engineer', '0.5 FTE from the regional pool, committed by September. The single critical-path resource.'),
    ('2 · Green-light one prospect pilot', 'AIB-profile account for the X-Ray proof, by October. A prospect, keeping existing customers out of scope.'),
    ('3 · Back the keystone study', 'Cost-per-outcome model with Deepak in September. Value Telemetry and the FinOps line depend on it.'),
    ('4 · Agree the commercial treatment', 'Quota recognition of ~€200K so the line carries a real target from day one.'),
]
for i, (h_, b_) in enumerate(asks):
    chip(sl, ML, 216 + i * 92, 560, 80, h_, b_, accent=BLUE)
hero_card(sl, 650, 216, 576, 356, 'The €5,000 seed',
          'The talent budget launches product number one.',
          'Packaging and delivery kit for Process X-Ray, a demo built on the APA V3 catalog, launch '
          'collateral produced with marketing, and pricing pressure-tests with three friendly accounts.')

# ═════════════════════════════════════════════════════════════
# S30 — Planning basis / assumptions
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Mobilization · planning basis',
       'The plan stands on six explicit assumptions, each with an owner to validate.',
       footnote='Confidence ratings follow the repo standard: every assumption is documented, conservative, and '
                'sensitivity-tested before any figure becomes client-facing or budget-committed.',
       notes='Per CLAUDE.md governance: no hidden assumptions. Validating the first three in September de-risks '
             'the sequence. EAP eligibility must be confirmed before any account conversation.')
assumps = [
    ('Team cost ~€1.2M a year', 'PDP working figure', 'Low', 'Mayur + finance'),
    ('Effort estimates within ±30%', 'Analogous VC engagements; first build actuals recalibrate', 'Medium', 'First build'),
    ('SE at 0.5 FTE is available', 'Regional pool capacity, unconfirmed', 'Low', 'SE leadership'),
    ('Tickets of €15-100K clear the market', 'AIB paid ~€80K for a weaker ADS; SAP and AWS wedge benchmarks', 'Medium', 'First 3 deals'),
    ('Wedge-to-Mission attach of 50%+ in 6 months', 'SAP and ServiceNow expansion patterns', 'Low', 'Pilot cohort'),
    ('EAP slots are open to our nominations', 'Mid-year program; 8 slots remain; eligibility (new logos vs existing customers) unconfirmed', 'Low', 'Program owner'),
]
conf_color = {'Low': RED, 'Medium': AMBER, 'High': GREEN}
for i, (a, basis, conf, owner) in enumerate(assumps):
    y = 210 + i * 68
    rect(sl, ML, y, AW, 58, fill=OFF, round_=True)
    txt(sl, a, ML + 24, y + 8, 400, 44, size=12.5, bold=True, sp=1.05)
    txt(sl, basis, 494, y + 8, 430, 44, size=11, color=BODY_DARK, sp=1.1)
    dot(sl, 1000, y + 29, 12, conf_color[conf])
    txt(sl, conf, 1014, y + 18, 76, 24, size=11.5, color=MUTED, bold=True)
    txt(sl, 'Validate: ' + owner, 1092, y + 10, 128, 46, size=10, color=BLUE, bold=True, sp=1.05)

# ═════════════════════════════════════════════════════════════
# S31 — Chapter 04
# ═════════════════════════════════════════════════════════════
chapter('04', 'Selling into apprehension',
        'The bank’s buying criteria, the Early Access route around the RFP, and the honest critique of this plan.')

# ═════════════════════════════════════════════════════════════
# S32 — Buying criteria
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Apprehension · the buying criteria',
       'Five questions decide a pre-RFP purchase; the wedge is shaped to answer every one.',
       footnote='Delegated-authority thresholds vary by bank; the design rule is to price each rung inside one named '
                'sponsor’s delegation, and phase anything larger. Validate per account in the free wedge.',
       notes='These are the real buying criteria of a bank buying outside an RFP. The free wedge exists partly to '
             'discover the sponsor’s threshold and the bank’s procurement rules before anything is priced.')
crits = [
    ('“Can one sponsor sign it?”', 'Fixed price, sized inside the sponsor’s delegated authority; anything larger phases into rungs.'),
    ('“What is our InfoSec exposure?”', 'Lite modes run on MI packs and schema metadata; customer data never leaves the bank.'),
    ('“What if we stop after this?”', 'The bank owns every artifact outright; the evidence is vendor-neutral and usable in any later tender.'),
    ('“Is the spend defensible to audit?”', 'Fixed scope, named deliverables, a value floor (documented value ≥ 10× fee or it is free), full fee credit on conversion.'),
    ('“Why you? You are new to this.”', 'Co-creation, never procurement: an Early Access nomination with roadmap influence, engineering access and joint proof.'),
]
for i, (c, a) in enumerate(crits):
    y = 208 + i * 86
    rect(sl, ML, y, 380, 76, fill=NAVY, round_=True)
    txt(sl, c, ML + 22, y + 12, 336, 56, size=12.5, color=WHITE, bold=True, sp=1.12)
    rect(sl, 454, y, 772, 76, fill=OFF, round_=True)
    rect(sl, 454, y + 6, 4, 64, fill=BLUE)
    txt(sl, a, 478, y + 12, 724, 56, size=11.5, sp=1.15)

# ═════════════════════════════════════════════════════════════
# S33 — Early Access Program
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Apprehension · the Early Access route',
       'The Early Access Program is the pre-RFP vehicle: slots are nominated, never tendered.',
       footnote='OPEN ITEM: eligibility (new logos vs existing customers) is unconfirmed; verify with the program '
                'owner before the first account conversation. Slot count (8 remaining) per the mid-year program status.',
       notes='Positioning: we are new to this category and hungry to co-innovate; the EAP makes that a strength. '
             'A nomination against program criteria is a selection, so there is nothing for procurement to tender. '
             'Scarcity (8 slots) creates urgency without discounting.')
eap = [
    ('Named roadmap influence', 'The bank co-shapes Banking OS: direct input on Nexus, Sentinel and Mission patterns during the program.'),
    ('Fee becomes participation', 'The wedge fee is framed as program participation and credits forward into the Mission POC.'),
    ('Engineering at the table', 'A direct line to Factory engineering during the Mission; issues resolve in days, never through support queues.'),
    ('Joint proof, joint story', 'Co-authored case study and reference rights; the bank is a design partner, never a logo on a slide.'),
]
for i, (h_, b_) in enumerate(eap):
    chip(sl, ML, 216 + i * 92, 560, 80, h_, b_, accent=BLUE)
hero_card(sl, 650, 216, 576, 356, 'Nominated, never tendered',
          '8 slots remain in the mid-year cohort.',
          'A bank joins by nomination against program criteria, so there is nothing to RFP. The scarcity is '
          'honest: when the slots are gone, the cohort closes. Confirm eligibility for new logos vs existing '
          'customers with the program owner before nominating.')

# ═════════════════════════════════════════════════════════════
# S34 — Pre-RFP mechanics
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Apprehension · the pre-RFP mechanics',
       'Five mechanics keep the wedge inside the sponsor’s pen and outside the tender process.',
       footnote='The one thing never to do: dodge procurement by making the work free. Free diagnostics re-break the '
                'economics and devalue the evidence. The credit-forward fee is the middle path.',
       notes='Mechanics 1-2 prevent the RFP trigger; 3-4 defuse the “budget on a maybe” fear; 5 defuses the '
             '“captive diagnostic” objection.')
mechs = [
    ('1 · A product SKU, never a services engagement', 'Fixed price, fixed scope, catalogued deliverables. Procurement treats it like a software purchase; there is nothing to tender against.'),
    ('2 · Sized to the sponsor’s pen', 'Each rung sits inside one executive’s delegated authority. Larger ambitions phase across rungs; the threshold is never jumped.'),
    ('3 · The fee credits forward, 100%', 'On conversion the wedge fee credits into the Mission POC or licence. The budget is never spent on a maybe.'),
    ('4 · The specimen pack sells before the SoW', 'The free wedge ends by showing the exact artifacts (this deck’s sample outputs). The bank buys a known object.'),
    ('5 · The bank owns the evidence, floor-guaranteed', 'All artifacts are the bank’s property, vendor-neutral, with a value floor: documented value ≥ 10× fee, or it is free.'),
]
for i, (h_, b_) in enumerate(mechs):
    y = 208 + i * 86
    rect(sl, ML, y, AW, 76, fill=OFF, round_=True)
    rect(sl, ML, y + 6, 4, 64, fill=BLUE)
    txt(sl, h_, ML + 24, y + 10, 420, 60, size=12.5, bold=True, sp=1.1)
    txt(sl, b_, 510, y + 12, 700, 56, size=11.5, sp=1.15)

# ═════════════════════════════════════════════════════════════
# S35 — Critique A · capacity
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Apprehension · the honest critique 1 of 2',
       'The capacity math does not support ten installs in year one; call year one a proof year.',
       footnote='Revision to the chapter-1 economics: ten a year remains the cost-neutral bar, reached in year two '
                'with two delivery consultants funded by year-one revenue and the delivery kits.',
       notes='Self-critique per the repo’s conservative-bias standard. Overpromising year one would burn the '
             'credibility the whole line depends on.')
tiles = [
    ('Demand of ten installs', '120-140 person-weeks a year', 'Delivery effort alone, per the register; excludes the build.'),
    ('Pod supply', '≈ 70 person-weeks a year', '1.5 FTE across build, sell, deliver and the PDP itself.'),
    ('Realistic ceiling', '4-5 installs alongside the build', 'The honest year-one number, with zero slack for surprises.'),
]
for i, (h_, big, note) in enumerate(tiles):
    x = ML + i * 399
    rect(sl, x, 210, 374, 190, fill=OFF, round_=True)
    txt(sl, h_.upper(), x + 22, 228, 330, 22, size=10.5, color=BLUE, bold=True)
    txt(sl, big, x + 22, 256, 330, 60, size=19, color=NAVY if i < 2 else RED, bold=True, sp=1.1)
    txt(sl, note, x + 22, 330, 330, 56, size=11, color=MUTED, sp=1.2)
band(sl, 432, 'The revised plan.',
     'Year one proves the model: 3-4 paid installs (€0.4-0.5M), one per product family, plus the Early Access '
     'cohort. Year two reaches the ten-install cost-neutral rate with two delivery consultants funded by '
     'year-one revenue, delivering from the kits.', h=96)

# ═════════════════════════════════════════════════════════════
# S36 — Critique B · open decisions
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Apprehension · the honest critique 2 of 2',
       'Five contradictions need a named decision before the first account conversation.',
       footnote='Each row is a real tension inside the current plan, surfaced deliberately: resolving them on paper '
                'now is cheaper than discovering them inside a live account.',
       notes='These four go on the Mayur 1:1 agenda alongside the four asks. The EAP row also gates the Tim '
             'conversation: nominations need confirmed eligibility.')
decs = [
    ('Prospect-only rule vs Telemetry', 'The pilot must be a prospect (Mayur’s ask), yet Telemetry needs a live deployment, which usually means an existing customer.', 'Mayur: scope the rule per product'),
    ('Credited fees vs the revenue KPI', 'A credited fee is pipeline; an uncredited install is revenue. The line carries one primary KPI; pick it before quota is set.', 'Mayur + finance: fix the KPI'),
    ('The wedge reads as a tollbooth', 'A €60K gate in front of a platform deal will push AEs to give discovery away free unless wedge revenue credits the account team.', 'Sales leadership: align incentives'),
    ('EAP eligibility is unconfirmed', 'Whether the 8 remaining slots take new logos, existing customers, or both changes which accounts can be nominated first.', 'Program owner: confirm eligibility'),
    ('Engine B fights MBB on their turf', 'Generic org design is customer-proximate space we lose; it is defensible only when it runs on Engine A’s harvested evidence.', 'You: enforce the sequencing rule'),
]
for i, (h_, b_, owner) in enumerate(decs):
    y = 204 + i * 88
    rect(sl, ML, y, AW, 80, fill=OFF, round_=True)
    rect(sl, ML, y + 8, 4, 64, fill=RED)
    txt(sl, h_, ML + 24, y + 10, 330, 64, size=12.5, bold=True, sp=1.1)
    txt(sl, b_, 420, y + 10, 540, 64, size=11, sp=1.15)
    rect(sl, 984, y + 18, 226, 46, fill=BLUE_LIGHT, round_=True)
    txt(sl, owner, 998, y + 24, 200, 40, size=10.5, color=BLUE_DARK, bold=True, sp=1.1)

# ═════════════════════════════════════════════════════════════
# S37 — Close
# ═════════════════════════════════════════════════════════════
sl = new_slide(NAVY)
cover_chrome(sl)
txt(sl, 'THE MOTION', 105, 220, 700, 24, size=12.5, color=CYAN, bold=True)
txt(sl, 'Land fast. Prove value.\nExpand.', 102, 262, 720, 160, size=42, color=WHITE, bold=True, sp=1.05)
txt(sl, 'One product in market by January. A proof year of 3-4 installs.\n'
        'Cost neutral in year two. The platform pre-installed in every account we touch.',
    102, 528, 710, 80, size=14, color=LIGHT_ON_NAVY, sp=1.3)
steps_now = [
    'This week · book the Mayur 1:1: quota, pilot green-light, SE allocation, the four decisions.',
    'This week · confirm EAP eligibility with the program owner; shortlist nominations.',
    'This week · follow up with Tim: category story, launch collateral, talent-budget plan.',
    'September · scope the keystone cost-per-outcome study with Deepak.',
]
for i, s_ in enumerate(steps_now):
    dot(sl, 890, 226 + i * 58, 10, CYAN)
    txt(sl, s_, 912, 212 + i * 58, 310, 52, size=11, color=WHITE, sp=1.18)
logo(sl, dark_bg=True)
page_field(sl, color=RGBColor(0x5A, 0x6B, 0x80))

out = os.path.join(HERE, 'Product_Factory_Execution_Plan_Exhibit.pptx')
prs.save(out)
import pathlib
print(f'Saved {out} ({pathlib.Path(out).stat().st_size // 1024} KB, {len(prs.slides._sldIdLst)} slides)')
