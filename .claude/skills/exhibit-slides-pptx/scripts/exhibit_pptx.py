#!/usr/bin/env python3
"""exhibit_pptx — the Backbase exhibit-style PPTX engine.

Extracted VERBATIM from the validated production builders:
  - Engagement/SNB Capital/Output/build_snbc_vc_pptx.py   (21 Jul 2026, chrome v3 FINAL)
  - Engagement/BACB/Output/build_scripts/bacb_close_exhibit_pptx.py (16 Jul 2026)

v3.1 (28 Jul 2026, ratified by Shyam as the DEFAULT for all PPTX decks):
  - right hairline rail standard (mirrors the left rail at x=12.760) — per the SNB 22 Jul deck
  - title default 25pt, ONE line, <=63 chars (28.5 only for short punch titles)
  - chip() primitive added (swim-lane cells, channel tags, coral STEER prompts)
  - page numbers are LIVE fields (a:fld type="slidenum") — auto-renumber on insert/
    reorder in PowerPoint & Google Slides; mechanism from the Product Factory deck
    (Mayur PDP session), styling unchanged (12.75pt bold black on the divider)
  Reference builds: build_snbc_vc_pptx.py + build_snbc_journey_maps_pptx.py (SNB Capital).

Every geometry number, color, and XML fix in here was verified against decks that
survived a Google Slides round-trip and were presented to clients. DO NOT retune
values here per deck — the whole point of this module is that the style never moves.

Usage (per-deck build script):

    from exhibit_pptx import ExhibitDeck, NAVY, BLUE, TINT, ...

    d = ExhibitDeck()                      # 13.333 x 7.5 in, blank layouts
    s = d.slide()
    d.chrome(s, "Kicker · section", "Action title as a full sentence")
    ... compose the exhibit with d.rect / d.txt / d.oval / d.hline ...
    d.takeaway_band(s, "Bold lead: ", "one-line remainder.")
    d.footnote(s, "Source: ...")
    d.notes(s, "Speaker notes incl. DEFENSE lines.")
    d.save("out.pptx")                     # runs the mandatory flat/strip pass
"""
import os
import copy

from pptx import Presentation
from pptx.util import Inches as I, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette (LOCKED)
NAVY  = RGBColor(0x07, 0x12, 0x24)   # ink + dark surfaces
BLUE  = RGBColor(0x40, 0x66, 0xF5)   # primary accent (kicker, lead bars, glyph)
BLUE2 = RGBColor(0x1F, 0x37, 0x99)   # deep secondary
BLUE3 = RGBColor(0x5F, 0x7D, 0xF7)   # mid secondary
BLUE4 = RGBColor(0xB5, 0xC1, 0xF1)   # light secondary
TINT  = RGBColor(0xE6, 0xEB, 0xFE)   # blue tint fill
TINT2 = RGBColor(0xF5, 0xF6, 0xFA)   # neutral tint fill
CYAN  = RGBColor(0x93, 0xFB, 0xFE)   # highlight on dark only — never a gradient step
CORAL = RGBColor(0xEC, 0x5E, 0x48)   # gates, risks, placeholders, ILLUSTRATIVE only
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUT   = RGBColor(0x6A, 0x71, 0x7C)   # muted body
GRAY55= RGBColor(0x77, 0x7D, 0x87)   # micro-labels ("PROVEN HERE")
FN    = RGBColor(0x8F, 0x94, 0x9C)   # footnotes
HAIR  = RGBColor(0xD2, 0xD4, 0xD8)   # hairlines
HAIR_ROW = RGBColor(0xE6, 0xE7, 0xE9)  # table row rules
DASHC = RGBColor(0xA8, 0xAC, 0xB2)   # dashed "yours/not built" outlines
DIVGR = RGBColor(0x9A, 0x9E, 0xA6)   # footer page-number divider
# dark-slide companions
HAIR_D = RGBColor(0x3E, 0x46, 0x54)
SUB_D  = RGBColor(0xC1, 0xC4, 0xC8)

FONT = "Libre Franklin"
W, H = 13.333, 7.5

_ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "assets")
LOGO_BLACK = os.path.abspath(os.path.join(_ASSETS, "backbase_logo_black.png"))
LOGO_WHITE = os.path.abspath(os.path.join(_ASSETS, "backbase_wordmark_white.png"))


class ExhibitDeck:
    """One deck, exhibit chrome baked in. All coordinates in inches on 13.333x7.5."""

    def __init__(self, logo=LOGO_BLACK):
        self.prs = Presentation()
        self.prs.slide_width = I(W)
        self.prs.slide_height = I(H)
        self._blank = self.prs.slide_layouts[6]
        self.logo = logo
        self.page = 0

    # ------------------------------------------------------------ slide factory
    def slide(self, dark=False):
        self.page += 1
        s = self.prs.slides.add_slide(self._blank)
        for ph in list(s.placeholders):
            ph._element.getparent().remove(ph._element)
        if dark:
            self.dark_bg(s)
        return s

    # ------------------------------------------------------------ primitives
    @staticmethod
    def flat(shp):
        """Explicit empty a:effectLst — LibreOffice/Google re-add theme shadows otherwise."""
        spPr = shp._element.spPr
        for el in spPr.findall(qn('a:effectLst')):
            spPr.remove(el)
        spPr.append(spPr.makeelement(qn('a:effectLst'), {}))

    def hline(self, s, x1, y1, x2, y2, color=HAIR, wpt=0.75):
        ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
        ln.line.color.rgb = color
        ln.line.width = Pt(wpt)
        ln.shadow.inherit = False
        spPr = ln._element.spPr
        spPr.append(spPr.makeelement(qn('a:effectLst'), {}))
        # strip the theme style reference (its effectRef re-adds a shadow in some renderers)
        el = ln._element
        for st in el.findall(qn('p:style')):
            el.remove(st)
        return ln

    def dashed_conn(self, s, x1, y1, x2, y2, color=BLUE, wpt=1.6, dash='dash'):
        conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
        conn.line.color.rgb = color
        conn.line.width = Pt(wpt)
        conn.shadow.inherit = False
        lnel = conn.line._get_or_add_ln()
        lnel.append(lnel.makeelement(qn('a:prstDash'), {'val': dash}))
        el = conn._element
        for st in el.findall(qn('p:style')):
            el.remove(st)
        return conn

    def step_glyph(self, s, x, y, w, h, fill):
        """The authentic Backbase step glyph (custGeom from the June master deck):
        a square with its top-left quadrant removed. Path units 9168 x 9096.
        NOT a plain square, NOT two stacked rects."""
        sx = w * 914400 / 9168.0
        sy = h * 914400 / 9096.0
        fb = s.shapes.build_freeform(19, 4762, scale=(sx, sy))
        fb.add_line_segments([(4567, 4762), (4566, 0), (9168, 0), (9168, 9096), (0, 9096)], close=True)
        shp = fb.convert_to_shape(I(x), I(y))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.fill.background()
        shp.shadow.inherit = False
        self.flat(shp)
        return shp

    def rect(self, s, x, y, w, h, fill=None, line=None, line_w=0.75, round_=False, dash=None):
        shp = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
            I(x), I(y), I(w), I(h))
        if round_:
            try:
                shp.adjustments[0] = 0.08
            except Exception:
                pass
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = Pt(line_w)
            if dash:
                ln = shp.line._get_or_add_ln()
                d = ln.makeelement(qn('a:prstDash'), {'val': dash})
                ln.append(d)
        shp.shadow.inherit = False
        self.flat(shp)
        return shp

    def diamond(self, s, x, y, w, h, fill=None, line=None, line_w=1.0):
        """Gate marker for wave timelines / milestone strips (T04, T08).
        Blue = platform release, coral = a decision the client owns."""
        shp = s.shapes.add_shape(MSO_SHAPE.DIAMOND, I(x), I(y), I(w), I(h))
        if fill is None:
            shp.fill.background()
        else:
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill
        if line is None:
            shp.line.fill.background()
        else:
            shp.line.color.rgb = line
            shp.line.width = Pt(line_w)
        shp.shadow.inherit = False
        self.flat(shp)
        return shp

    def oval(self, s, x, y, w, h, fill, label=None, tc=WHITE, fs=11, bold=True):
        ov = s.shapes.add_shape(MSO_SHAPE.OVAL, I(x), I(y), I(w), I(h))
        ov.shadow.inherit = False
        ov.line.fill.background()
        ov.fill.solid()
        ov.fill.fore_color.rgb = fill
        self.flat(ov)
        if label is not None:
            tf = ov.text_frame
            tf.word_wrap = True
            tf.margin_left = 0
            tf.margin_right = 0
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = label
            f = r.font
            f.name = FONT
            f.size = Pt(fs)
            f.bold = bold
            f.color.rgb = tc
        return ov

    def txt(self, s, x, y, w, h, runs, size=14, color=NAVY, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp_after=2, line_sp=1.0,
            wrap=True, track=None):
        """runs: plain string, OR list of paragraphs, each a list of (text, size, color, bold)."""
        tb = s.shapes.add_textbox(I(x), I(y), I(w), I(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = anchor
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        if isinstance(runs, str):
            runs = [[(runs, size, color, bold)]]
        first = True
        for para in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            p.space_after = Pt(sp_after)
            p.line_spacing = line_sp
            for (t, sz, c, b) in para:
                r = p.add_run()
                r.text = t
                f = r.font
                f.name = FONT
                f.size = Pt(sz)
                f.color.rgb = c
                f.bold = b
                if track is not None:
                    r._r.get_or_add_rPr().set('spc', str(track))
        return tb

    def page_field(self, s, page, x=12.872, y=7.118, w=0.42, h=0.249):
        """LIVE slide-number field (<a:fld type="slidenum">): auto-renumbers when
        slides are added, removed or reordered in PowerPoint/Google Slides — no more
        hardcoded-number renumbering passes. Mechanism ported from the Product
        Factory deck (Mayur PDP session, 28 Jul 2026); styled to the exhibit footer
        spec. The literal text is the build-time fallback for renderers that don't
        evaluate fields."""
        tb = s.shapes.add_textbox(I(x), I(y), I(w), I(h))
        tf = tb.text_frame
        tf.word_wrap = False
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        p = tf.paragraphs[0]._p
        fld = p.makeelement(qn('a:fld'), {
            'id': '{B7A1C0DE-2026-4728-8000-%012X}' % int(page), 'type': 'slidenum'})
        rPr = p.makeelement(qn('a:rPr'), {'lang': 'en-US', 'sz': '1275', 'b': '1', 'dirty': '0'})
        fill = p.makeelement(qn('a:solidFill'), {})
        clr = p.makeelement(qn('a:srgbClr'), {'val': '000000'})
        fill.append(clr)
        rPr.append(fill)
        rPr.append(p.makeelement(qn('a:latin'), {'typeface': FONT}))
        t = p.makeelement(qn('a:t'), {})
        t.text = str(int(page))
        fld.append(rPr)
        fld.append(t)
        p.append(fld)
        return tb

    # ------------------------------------------------------------ chrome (v3.1)
    def chrome(self, s, kicker, title, page=None, title_size=25, marker=None, right_rail=True):
        """Standard white-slide chrome. page defaults to the running counter.
        Title law (v3.1): ONE line, <=63 chars at 25pt; 28.5 only for short punch
        titles (<=48 chars). Never let a title wrap.
        marker: optional (labels, active) tuple for a top-right session strip,
        e.g. (["OPEN","DEMO","CLOSE"], "OPEN"). Off by default."""
        page = self.page if page is None else page
        # full-bleed hairlines (#D2D4D8, 0.75pt); right rail standard since v3.1
        self.hline(s, 0, 0.573, W, 0.573)
        self.hline(s, 0.573, 0, 0.573, 7.042)
        if right_rail:
            self.hline(s, 12.760, 0, 12.760, 7.042)
        self.hline(s, 0, 7.042, W, 7.042)
        # authentic step glyph hugging the hairline crossing; corner ends at (0.573, 0.573)
        self.step_glyph(s, 0.406, 0.406, 0.167, 0.166, BLUE)
        if marker:
            labels, active = marker
            for i, l in enumerate(labels):
                on = (l == active)
                b = self.rect(s, 10.9 + i * 0.58, 0.20, 0.54, 0.26,
                              fill=BLUE if on else None, line=None if on else HAIR, line_w=0.75)
                tf = b.text_frame
                tf.word_wrap = False
                tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = l
                f = r.font; f.name = FONT; f.size = Pt(8.5); f.bold = on
                f.color.rgb = WHITE if on else MUT
        # kicker 13.5 regular blue at 0.812; title 28.5 bold navy at 1.104, trailing period
        self.txt(s, 1.0, 0.812, 10.42, 0.25, kicker.upper(), size=13.5, color=BLUE, bold=False, track="120")
        t = title if title.rstrip().endswith((".", "?", ":")) else title + "."
        self.txt(s, 1.0, 1.104, 10.94, 1.05, t, size=title_size, color=NAVY, bold=True, line_sp=1.04)
        # white footer: large black wordmark │ page number
        if self.logo and os.path.exists(self.logo):
            s.shapes.add_picture(self.logo, I(11.633), I(7.185), I(1.067), I(0.173))
        else:
            self.txt(s, 10.6, 7.16, 2.1, 0.25, "Backbase", size=13, color=RGBColor(0, 0, 0),
                     bold=True, align=PP_ALIGN.RIGHT)
        self.hline(s, 12.802, 7.118, 12.802, 7.367, color=DIVGR, wpt=0.9)
        self.page_field(s, page)

    def footnote(self, s, text, y=6.50):
        """Every numeric slide ends with one. Hairline + 9.5pt source line."""
        self.hline(s, 1.0, y + 0.06, 12.708, y + 0.06)
        self.txt(s, 1.0, y + 0.14, 11.708, 0.40, text, size=9.5, color=FN, line_sp=1.15)

    def notes(self, s, text):
        s.notes_slide.notes_text_frame.text = text

    def takeaway_band(self, s, lead, rest, y=5.62, x=1.0, w=11.708):
        """Navy punchline strip: bold cyan lead-in + white remainder.
        EXACTLY one sentence, fitting ONE line — cut words until it does."""
        self.rect(s, x, y, w, 0.479, fill=NAVY, round_=True)
        self.txt(s, x + 0.25, y + 0.125, w - 0.5, 0.28,
                 [[(lead, 13.5, CYAN, True), (rest, 13.5, WHITE, False)]])

    def open_badge(self, s, x, y, w, lead, rest, h=0.72):
        """Coral dashed box = still open / placeholder / illustrative-outside-in."""
        self.rect(s, x, y, w, h, line=CORAL, line_w=1.2, round_=True, dash='dash')
        self.txt(s, x + 0.25, y + 0.11, w - 0.5, h - 0.2,
                 [[(lead, 11.5, CORAL, True), (rest, 11.5, NAVY, False)]], line_sp=1.12)

    def stat_card(self, s, x, y, from_text, value, label, w=2.775, h=1.21):
        """Impact stat card (ratified 28 Jul 2026 — Shyam prefers this over the
        3-hero-tile close: 'a lot cleaner'). Tint card: muted "from →" line,
        30pt deep-blue landing value, one/two-line muted label. Standard row =
        FOUR across the content width at x=1.0, pitch 2.975 (gap 0.2)."""
        self.rect(s, x, y, w, h, fill=TINT, round_=True)
        self.txt(s, x + 0.18, y + 0.13, w - 0.36, 0.24, from_text + "  →", size=11.5, color=MUT)
        self.txt(s, x + 0.18, y + 0.40, w - 0.36, 0.48, value, size=30, color=BLUE2, bold=True)
        self.txt(s, x + 0.18, y + 0.85, w - 0.36, 0.34, label, size=10, color=MUT, line_sp=1.05)

    def proven_band(self, s, items, y=5.62, x=1.0, w=11.708, h=0.34, lead="Proven here:  "):
        """Tint credibility strip above the footnote: bold blue lead + where it
        runs live (Backbase client names allowed per anonymization rules)."""
        self.rect(s, x, y, w, h, fill=TINT2, round_=True)
        self.txt(s, x + 0.30, y + 0.055, w - 0.60, h - 0.10,
                 [[(lead, 11.5, BLUE, True), (items, 11.5, NAVY, False)]])

    def chip(self, s, x, y, w, h, text, fill=TINT2, tc=NAVY, fs=8.5, bold=False,
             line=None, dash=None, align=PP_ALIGN.LEFT):
        """Small rounded label chip (v3.1): swim-lane cells, stage headers, channel
        tags, workshop STEER prompts. Text vertically centered, 0.09in side inset.
        fill=None + line=CORAL + dash='dash' = the coral STEER/open prompt.
        Keep the chip's bottom edge >=0.06in clear of the footnote hairline."""
        self.rect(s, x, y, w, h, fill=fill, line=line, dash=dash, round_=True)
        self.txt(s, x + 0.09, y + 0.04, w - 0.18, h - 0.08, text, size=fs, color=tc,
                 bold=bold, line_sp=0.98, align=align, anchor=MSO_ANCHOR.MIDDLE)

    # ------------------------------------------------------------ dark slides
    def dark_bg(self, s):
        self.rect(s, 0, 0, W, H, fill=NAVY)
        self.rect(s, 6.0, -2.2, 9.5, 5.2, fill=RGBColor(0x14, 0x2A, 0x6E))  # glow approximation
        self.rect(s, 8.2, -1.4, 6.2, 3.2, fill=RGBColor(0x1B, 0x38, 0x94))
        return s

    def divider(self, number, title, subtitle=None):
        """T14 chapter divider: dark, big light-weight number + title."""
        s = self.slide(dark=True)
        self.step_glyph(s, 0.406, 0.406, 0.167, 0.166, CYAN)
        self.rect(s, 0.57, 0.57, 12.19, 0.013, fill=RGBColor(0x2E, 0x3A, 0x52))
        self.txt(s, 1.0, 2.55, 3.0, 1.6, str(number), size=96, color=RGBColor(0x3A, 0x4A, 0x6B), bold=False)
        self.txt(s, 1.0, 4.15, 10.8, 0.9, title, size=34, color=WHITE, bold=True, line_sp=1.05)
        if subtitle:
            self.txt(s, 1.0, 4.95, 9.5, 0.7, subtitle, size=14, color=SUB_D, line_sp=1.25)
        self.rect(s, 0.57, 6.95, 12.19, 0.013, fill=RGBColor(0x2E, 0x3A, 0x52))
        self.txt(s, 11.7, 7.05, 1.06, 0.3, "Backbase", size=11, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
        return s

    # ------------------------------------------------------------ save
    def _strip_theme_styles(self):
        """Mandatory save pass: remove <p:style> from every shape and guarantee an
        explicit empty <a:effectLst> — otherwise LibreOffice/Google Slides re-apply
        theme drop-shadows (found the hard way on the BACB close deck)."""
        for slide_ in self.prs.slides:
            for shp in slide_.shapes:
                el = shp._element
                for st in el.findall(qn('p:style')):
                    el.remove(st)
                spPr = getattr(el, 'spPr', None)
                if spPr is not None and not spPr.findall(qn('a:effectLst')):
                    spPr.append(spPr.makeelement(qn('a:effectLst'), {}))

    def save(self, path):
        self._strip_theme_styles()
        self.prs.save(path)
        return path
