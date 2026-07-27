#!/usr/bin/env python3
"""Product Factory execution plan — exhibit-style PPTX (internal, PDP).

Style: McKinsey-exhibit patterns from the claude-design-exhibit-kit
(action-title sentences, one exhibit per slide, hairline chrome, source
footnotes), rendered in Frontline 2026 tokens per the July 2026 adoption
decision (patterns adopted, kit palette / L·E·C chrome deliberately not).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE

# ── Frontline tokens ─────────────────────────────────────────
NAVY = RGBColor(0x04, 0x13, 0x26)
BLUE = RGBColor(0x33, 0x67, 0xFF)
BLUE_DARK = RGBColor(0x26, 0x4E, 0xC7)
BLUE_LIGHT = RGBColor(0xE5, 0xEB, 0xFF)
OFF = RGBColor(0xF3, 0xF6, 0xF9)
BORDER = RGBColor(0xCE, 0xD2, 0xD7)
MUTED = RGBColor(0x6B, 0x77, 0x86)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CYAN = RGBColor(0x69, 0xFE, 0xFF)
RED = RGBColor(0xFF, 0x50, 0x3C)
AMBER = RGBColor(0xB4, 0x53, 0x09)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
LIGHT_ON_NAVY = RGBColor(0xC5, 0xCF, 0xDE)
# Blue intensity ramp (alpha steps of token blue on white)
T75 = RGBColor(0x66, 0x8D, 0xFF)
T50 = RGBColor(0x99, 0xB3, 0xFF)
T25 = RGBColor(0xCC, 0xD9, 0xFF)

FONT = 'Libre Franklin'
KIT = os.path.dirname(os.path.abspath(__file__))
LOGO_DARK = os.path.join(KIT, '../../knowledge/design-system/claude-design-exhibit-kit/assets/logo/backbase-wordmark-dark.png')
LOGO_WHITE = os.path.join(KIT, '../../knowledge/design-system/claude-design-exhibit-kit/assets/logo/backbase-wordmark-white.png')

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
PAGE = [0]


def px(v):
    """Design px (1280×720 canvas) → EMU."""
    return Emu(int(v / 96.0 * 914400))


def pt_of(v_px):
    return Pt(v_px * 0.75)


def new_slide(bg=WHITE):
    sl = prs.slides.add_slide(BLANK)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = bg
    PAGE[0] += 1
    return sl


def txt(sl, s, x, y, w, h, size=18, color=NAVY, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, sp=1.15, font=FONT):
    tb = sl.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = str(s).split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = sp
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = pt_of(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def rect(sl, x, y, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = sl.shapes.add_shape(shape, px(x), px(y), px(w), px(h))
    if round_:
        try:
            sp.adjustments[0] = 0.08
        except Exception:
            pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def hairline(sl, x, y, w, h=1.6, color=BORDER):
    return rect(sl, x, y, w, h, fill=color)


def diamond(sl, cx, cy, size, fill):
    sp = sl.shapes.add_shape(MSO_SHAPE.DIAMOND, px(cx - size / 2), px(cy - size / 2),
                             px(size), px(size))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = WHITE
    sp.line.width = Pt(1.2)
    sp.shadow.inherit = False
    return sp


def dot(sl, cx, cy, d, fill, line=None):
    sp = sl.shapes.add_shape(MSO_SHAPE.OVAL, px(cx - d / 2), px(cy - d / 2), px(d), px(d))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line:
        sp.line.color.rgb = line
        sp.line.width = Pt(1.5)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def logo(sl, dark_bg=False, x=1128, y=684, w=72):
    path = LOGO_WHITE if dark_bg else LOGO_DARK
    if os.path.exists(path):
        sl.shapes.add_picture(path, px(x), px(y), width=px(w))


def chrome(sl, kicker, title, footnote=None, notes=None, title_w=980):
    hairline(sl, 55, 55, 1280 - 55 - 60)
    hairline(sl, 55, 55, 1.6, h=720 - 55 - 64)
    txt(sl, kicker.upper(), 96, 64, 900, 26, size=13, color=BLUE, bold=True)
    txt(sl, title, 96, 94, title_w, 110, size=29, bold=True, sp=1.08)
    # footer
    hairline(sl, 55, 720 - 44, 1280 - 55 - 60)
    txt(sl, str(PAGE[0]), 96, 720 - 36, 60, 22, size=12, color=MUTED)
    logo(sl, dark_bg=False)
    if footnote:
        hairline(sl, 96, 720 - 44 - 34, 1280 - 96 - 60)
        txt(sl, footnote, 96, 720 - 44 - 27, 1280 - 96 - 60, 26, size=11.5,
            color=MUTED, sp=1.05)
    if notes:
        sl.notes_slide.notes_text_frame.text = notes


def chip(sl, x, y, w, h, head, body, fill=OFF, head_color=NAVY, body_color=None,
         accent=None):
    rect(sl, x, y, w, h, fill=fill, line=BORDER, line_w=0.75, round_=True)
    if accent:
        rect(sl, x, y + 8, 4, h - 16, fill=accent)
    txt(sl, head, x + 20, y + 12, w - 36, 26, size=14, bold=True, color=head_color)
    txt(sl, body, x + 20, y + 40, w - 36, h - 50, size=12,
        color=body_color or RGBColor(0x2A, 0x38, 0x4A), sp=1.12)


def hero_card(sl, x, y, w, h, eyebrow, big, body):
    rect(sl, x, y, w, h, fill=NAVY, round_=True)
    txt(sl, eyebrow.upper(), x + 30, y + 26, w - 60, 24, size=12, color=CYAN, bold=True)
    txt(sl, big, x + 30, y + 58, w - 60, 120, size=21, color=WHITE, bold=True, sp=1.12)
    txt(sl, body, x + 30, y + h - 118, w - 60, 104, size=12.5, color=LIGHT_ON_NAVY, sp=1.2)


# ═════════════════════════════════════════════════════════════
# S1 — Cover
# ═════════════════════════════════════════════════════════════
sl = new_slide(NAVY)
logo(sl, dark_bg=True, x=96, y=64, w=120)
rect(sl, 96, 250, 56, 5, fill=BLUE)
txt(sl, 'VALUE CONSULTING · PDP · INTERNAL', 96, 276, 800, 26, size=13, color=CYAN, bold=True)
txt(sl, 'Product Factory:\nfrom concept to revenue', 92, 312, 1050, 170, size=52,
    color=WHITE, bold=True, sp=1.02)
txt(sl, 'The execution plan for four paid product wedges that pre-install Banking OS,\n'
        'make the function cost neutral inside a year, and open a recurring revenue line.',
    96, 500, 980, 70, size=16, color=LIGHT_ON_NAVY, sp=1.3)
txt(sl, 'Shyam · July 2026 · prepared for the talent programme (Tim Ruttner) and the PDP track (Mayur)',
    96, 640, 1000, 26, size=12, color=MUTED)

# ═════════════════════════════════════════════════════════════
# S2 — Chapter 01
# ═════════════════════════════════════════════════════════════
sl = new_slide(NAVY)
txt(sl, '01', 96, 170, 400, 170, size=110, color=T50, bold=False)
txt(sl, 'The frameworks', 96, 348, 900, 70, size=40, color=WHITE, bold=True)
txt(sl, 'The thesis, the product filter, the economics, the lifecycle, and the pod that runs it.',
    96, 430, 860, 56, size=15, color=LIGHT_ON_NAVY, sp=1.25)
logo(sl, dark_bg=True)
txt(sl, str(PAGE[0]), 96, 720 - 36, 60, 22, size=12, color=RGBColor(0x5A, 0x6B, 0x80))

# ═════════════════════════════════════════════════════════════
# S3 — Thesis (chips + hero split)
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Frameworks · the thesis',
       'Every product earns a fee, harvests the bank’s reality, and installs the '
       'first component of Banking OS.',
       footnote='Filter rule: a candidate that fails any one of the three tests is consulting, '
                'and stays out of the product line. Aligned to the Banking OS canon (knowledge/product/banking-os.md).',
       notes='The three-part design filter from the 29 Jun PDP session. Product-proximate = moat; '
             'customer-proximate = competitive space we lose. Every SKU must pass all three tests.')
chips_y = 232
for i, (h_, b_) in enumerate([
    ('Test 1 · It earns a fee', 'Each install is paid work at €15-100K, priced as a product with a fixed scope and a delivery kit.'),
    ('Test 2 · It harvests proprietary data', 'The engagement captures the bank’s real processes, data structures and policies. Evidence only Backbase holds.'),
    ('Test 3 · It pre-binds a Banking OS layer', 'The output is a working component: Nexus, Sentinel or a Mission candidate, already present in the account.'),
]):
    chip(sl, 96, chips_y + i * 122, 560, 108, h_, b_, accent=BLUE)
hero_card(sl, 692, 232, 528, 348, 'The one-liner',
          'Paid discovery that leaves the platform installed.',
          'Each engagement ends with a live Banking OS component in the account, a quantified '
          'expansion case for the AE, and account intelligence competitors and SIs never see.')

# ═════════════════════════════════════════════════════════════
# S4 — Cost-neutrality dot grid
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Frameworks · the economics',
       'Ten paid installs a year cover the full team cost; every install beyond is new revenue.',
       footnote='Package = Assess & Solution €60K + Mission POC €60K. Team cost of ~€1.2M a year is the PDP '
                'working figure and is to be validated with finance. Year-two target of 14-16 installs is a planning range.',
       notes='Cost-neutrality is the floor, revenue generation is the objective. 10 × €120K ≈ €1.2M. '
             'Year two adds recurring Value Telemetry revenue on top of installs.')
gx, gy, d, gap = 140, 300, 42, 30
for i in range(16):
    row, col = divmod(i, 8)
    fill = BLUE if i < 10 else T25
    dot(sl, gx + col * (d + gap) + d / 2, gy + row * (d + gap) + d / 2, d, fill)
dot(sl, 140 + 10, 490, 16, BLUE)
txt(sl, 'Installs to cost neutrality (10)', 168, 480, 300, 24, size=12, color=MUTED)
dot(sl, 500 + 10, 490, 16, T25)
txt(sl, 'Year-two expansion range (14-16)', 528, 480, 320, 24, size=12, color=MUTED)
rect(sl, 760, 250, 460, 260, fill=OFF, round_=True)
txt(sl, '10', 800, 278, 180, 90, size=64, color=BLUE, bold=True)
txt(sl, 'installs a year at an average package of €120K cover the ~€1.2M annual team cost.',
    800, 380, 380, 100, size=14.5, sp=1.25)

# ═════════════════════════════════════════════════════════════
# S5 — Package ladder (segmented to-scale bar)
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Frameworks · the ladder',
       'The wedge stays free; depth and recurrence carry the price.',
       footnote='Recurring range is a target shape anchored on SAP and ServiceNow success-plan pricing of 10-30% of '
                'licence value, and on FinOps gain-share market practice (15-35% of realised savings).',
       notes='Rung 0 free = funnel. Rungs 1+2 = the €120K SKVC package = cost neutrality. '
             'Rung 3 = the growth flywheel (Value Telemetry + gain-share).')
y0 = 300
chip(sl, 96, y0 - 34, 240, 130, 'Rung 0 · free wedge',
     'Ignite Inspire\n~4 meetings, benchmark and use-case shortlist. Grows the funnel.', fill=OFF)
txt(sl, '→', 348, y0 + 10, 30, 40, size=20, color=MUTED, align=PP_ALIGN.CENTER)
# to-scale €120K bar
bar_x, bar_w = 392, 560
rect(sl, bar_x, y0, int(bar_w / 2), 76, fill=BLUE)
txt(sl, 'Assess & Solution\n€60K', bar_x + 18, y0 + 12, 240, 56, size=13, color=WHITE, bold=True, sp=1.1)
rect(sl, bar_x + int(bar_w / 2), y0, int(bar_w / 2), 76, fill=BLUE_DARK)
txt(sl, 'Mission POC\n€60K', bar_x + int(bar_w / 2) + 18, y0 + 12, 240, 56, size=13, color=WHITE, bold=True, sp=1.1)
txt(sl, 'Rungs 1 + 2 · the €120K package (to scale)', bar_x, y0 - 32, 560, 24, size=12.5, color=MUTED, bold=True)
txt(sl, '→', 964, y0 + 10, 30, 40, size=20, color=MUTED, align=PP_ALIGN.CENTER)
rect(sl, 1006, y0 - 34, 214, 130, fill=NAVY, round_=True)
txt(sl, 'Rung 3 · recurring', 1026, y0 - 22, 180, 24, size=12, color=CYAN, bold=True)
txt(sl, 'Value Assurance\n€30-50K a year\n+ gain share', 1026, y0 + 4, 180, 84, size=13, color=WHITE, bold=True, sp=1.15)
chip(sl, 96, 470, 548, 92, 'The free wedge feeds the funnel',
     'Four meetings surface the leakage evidence that qualifies the paid wedge. Zero friction to start.', accent=BLUE)
chip(sl, 672, 470, 548, 92, 'Recurrence is the growth flywheel',
     'Telemetry converts every live deployment into quarterly revenue that compounds with AI consumption.', accent=BLUE)

# ═════════════════════════════════════════════════════════════
# S6 — Lifecycle step flow
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Frameworks · the lifecycle',
       'One motion carries every product from concept to installed platform.',
       footnote='Durations are per-engagement figures once the product is built, with a ±30% planning tolerance. '
                'The first pass through stages 1-4 is the one-time product build.',
       notes='Stage exit gates: Concept = qualified sponsor + pain. Design = signed SoW, agreed metrics, data access. '
             'Prototype = working demo on the client’s own data. Proof = CFO-grade evidence pack. '
             'Proposal = commercial offer + named Mission candidate. Signature = signed order form. '
             'Install = layer live + expansion roadmap handed to the AE.')
steps = [
    ('01', 'Concept', 'VC · 1 wk'),
    ('02', 'Design', 'VC + SE · 1-2 wks'),
    ('03', 'Prototype', 'SE · 2-3 wks'),
    ('04', 'Proof', 'VC · 2-4 wks'),
    ('05', 'Proposal', 'VC + AE · 1 wk'),
    ('06', 'Signature', 'AE + VC · 2-4 wks'),
    ('07', 'Install', 'FDE + SE · 6-12 wks'),
]
sw, sgap, sx, sy = 148, 12, 96, 268
for i, (n, name, who) in enumerate(steps):
    x = sx + i * (sw + sgap)
    rect(sl, x, sy, sw, 150, fill=OFF, line=BORDER, line_w=0.75, round_=True)
    dot(sl, x + 32, sy + 34, 34, BLUE)
    txt(sl, n, x + 15, sy + 23, 34, 24, size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, name, x + 16, sy + 62, sw - 30, 26, size=14, bold=True)
    txt(sl, who, x + 16, sy + 92, sw - 30, 46, size=10.5, color=MUTED, sp=1.15)
    if i < 6:
        txt(sl, '→', x + sw - 3, sy + 58, 20, 30, size=14, color=MUTED)
rect(sl, 96, 470, 1124, 84, fill=BLUE_LIGHT, round_=True)
txt(sl, 'The lifecycle runs twice over.', 120, 486, 300, 52, size=13.5, bold=True, color=BLUE_DARK, sp=1.1)
txt(sl, 'The first pass through stages 1-4 builds the product itself. From the second account onward, '
        'the same motion runs as a 4-8 week sales-and-delivery cycle per install.',
    430, 482, 766, 64, size=12.5, sp=1.2)

# ═════════════════════════════════════════════════════════════
# S7 — The pod
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Frameworks · resourcing',
       'A pod of two and a half people runs the whole line; hiring follows revenue.',
       footnote='FTE figures are steady-state allocations during build and delivery phases; FDE capacity is drawn '
                'from standard Mission Sprint staffing at install, so it is charged to the delivery, and not to this pod.',
       notes='This is an allocation ask, and specifically NOT a hiring ask. The SE is the critical unlock: '
             'same person across all four products so tooling skill compounds.')
pods = [
    ('You · product owner · 1.0 FTE', 'SKU design, value narrative, proof economics, proposal. Owns the P&L of the line.'),
    ('Solution engineer · 0.5 FTE, named', 'Connectors, data extraction, Factory tooling at prototype and install. The technical spine.'),
    ('FDE · on demand at install', 'Runs the Factory Mission Sprint once an install is signed. 6-12 weeks per Mission.'),
    ('SME · on call', 'Risk and compliance review for Guardrail Studio; R&D liaison for the Cartographer.'),
]
for i, (h_, b_) in enumerate(pods):
    chip(sl, 96, 226 + i * 92, 560, 80, h_, b_, accent=BLUE)
hero_card(sl, 692, 226, 528, 356, 'Why the solution engineer matters',
          'One named SE across all four products, so skill compounds.',
          'Every product runs on the same Factory tools: Connector Studio, Semantic Modeler, '
          'Simulation & Testing. Each install makes the next one cheaper and faster to deliver.')

# ═════════════════════════════════════════════════════════════
# S8 — Concern → answer stack
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Frameworks · the objections',
       'The three obvious objections have concrete answers.',
       footnote='AIB reference: paid an estimated ~€80K for an external current/future-state assessment (ADS), PDP session '
                'record, 29 Jun 2026. Wedge pattern references: AWS OLA, SAP Outside-In, ServiceNow Inspire.',
       notes='Use these in the Tim and Mayur conversations. Each answer is evidence-backed, from the PDP research.')
qa = [
    ('“Banks will not pay for discovery.”',
     'The market already does: AIB paid ~€80K for a weaker assessment. AWS, SAP and ServiceNow all run '
     'the diagnostic-wedge motion; ours goes further by installing product.'),
    ('“We have no capacity to build products.”',
     'Build effort is 6-12 person-weeks per product because the Factory tools, the APA V3 catalog and the '
     'cost model already exist. The sequence builds one product at a time.'),
    ('“This distracts from platform sales.”',
     'Every install pre-binds Nexus, Sentinel or a Mission and hands the AE a quantified expansion case. '
     'The wedge is the land motion of land-and-expand.'),
]
for i, (c, a) in enumerate(qa):
    y = 232 + i * 122
    rect(sl, 96, y, 420, 108, fill=NAVY, round_=True)
    txt(sl, c, 120, y + 18, 372, 76, size=14, color=WHITE, bold=True, sp=1.15)
    rect(sl, 536, y, 684, 108, fill=OFF, round_=True)
    rect(sl, 536, y + 8, 4, 92, fill=BLUE)
    txt(sl, a, 560, y + 14, 636, 84, size=12.5, sp=1.18)

# ═════════════════════════════════════════════════════════════
# S9 — Chapter 02
# ═════════════════════════════════════════════════════════════
sl = new_slide(NAVY)
txt(sl, '02', 96, 170, 400, 170, size=110, color=T50)
txt(sl, 'Four execution plans', 96, 348, 900, 70, size=40, color=WHITE, bold=True)
txt(sl, 'Product by product: milestones, effort, pod, ticket, and the gate that unlocks each one.',
    96, 430, 900, 56, size=15, color=LIGHT_ON_NAVY, sp=1.25)
logo(sl, dark_bg=True)
txt(sl, str(PAGE[0]), 96, 720 - 36, 60, 22, size=12, color=RGBColor(0x5A, 0x6B, 0x80))


def product_plan(kicker, title, milestones, stats, footnote, notes=None):
    """Milestone strip + dark stat band. milestone = (date, label, body, gate)
    gate: None | 'decision' (red diamond) | 'ga' (blue diamond)."""
    sl = new_slide()
    chrome(sl, kicker, title, footnote=footnote, notes=notes)
    n = len(milestones)
    line_y = 300
    x0, x1 = 230, 1090
    hairline(sl, x0, line_y, x1 - x0, 2.4, color=BORDER)
    span = (x1 - x0) / (n - 1)
    for i, (date, label, body, gate) in enumerate(milestones):
        cx = x0 + i * span
        if gate == 'decision':
            diamond(sl, cx, line_y + 1, 22, RED)
        elif gate == 'ga':
            diamond(sl, cx, line_y + 1, 22, BLUE)
        else:
            dot(sl, cx, line_y + 1, 16, BLUE, line=WHITE)
        txt(sl, date, cx - 100, line_y - 40, 200, 24, size=12, color=BLUE, bold=True,
            align=PP_ALIGN.CENTER)
        txt(sl, label, cx - 100, line_y + 26, 200, 26, size=13.5, bold=True,
            align=PP_ALIGN.CENTER)
        txt(sl, body, cx - 100, line_y + 54, 200, 100, size=10.5, color=MUTED,
            align=PP_ALIGN.CENTER, sp=1.12)
    band_y = 470
    rect(sl, 96, band_y, 1124, 96, fill=NAVY, round_=True)
    cw = 1124 / len(stats)
    for i, (lab, val) in enumerate(stats):
        bx = 96 + i * cw
        txt(sl, lab.upper(), bx + 28, band_y + 16, cw - 40, 22, size=10.5, color=CYAN, bold=True)
        txt(sl, val, bx + 28, band_y + 42, cw - 40, 46, size=14.5, color=WHITE, bold=True, sp=1.05)
    return sl


product_plan(
    'Execution plan 1 of 4 · Process X-Ray',
    'Process X-Ray ships first: demo by October, paid proof by December, first revenue in January.',
    [
        ('Aug 2026', 'Design', 'Delivery kit, pricing one-pager, demo script. 3 pw, VC.', None),
        ('Sep-Oct 2026', 'Prototype', 'Demo on the APA V3 catalog (65 processes) plus one friendly dataset. 5 pw, SE-led.', None),
        ('Oct 2026', 'Pilot gate', 'Prospect pilot green-light with Mayur. AIB-profile account.', 'decision'),
        ('Oct-Dec 2026', 'Proof', 'Paid pilot: leakage heatmap, top-3 Mission candidates with ROI evidence. 6 pw.', None),
        ('Jan 2027', 'First signature', 'Wedge sold at €50-75K; install starts as a Factory Mission Sprint.', 'ga'),
    ],
    [('Build effort', '8-10 person-wks'), ('Delivery per install', '4-6 wks · 12-14 pw'),
     ('Pod', 'VC 1.0 + SE 0.5'), ('Pre-binds', 'Mission candidate + ROI')],
    'Effort carries a ±30% tolerance. Build is low-risk because the APA V3 catalog (65 processes, 283 steps), the '
    'journey builder and the value-leakage method already exist in the repo.',
    notes='X-Ray productizes the Process & Workspace Designer. It observes real flows across the 6-12 systems per '
          'journey and quantifies leakage in the messy middle, where ~60% of bank work lives.')

product_plan(
    'Execution plan 2 of 4 · Value Telemetry',
    'Value Telemetry rides the keystone cost-per-outcome model and opens the recurring line.',
    [
        ('Sep 2026', 'Keystone gate', 'Extend the cost model from cost-per-interaction to cost-per-outcome, with Deepak. ~5 use cases.', 'decision'),
        ('Sep 2026', 'Design', 'SKU spec: diagnostic scope, telemetry set, gain-share terms. 4 pw, VC.', None),
        ('Oct 2026', 'Prototype', 'Cost-per-outcome dashboard on one live deployment’s data. 3 pw.', None),
        ('Nov-Dec 2026', 'Proof', '2-3 week diagnostic at a live customer: waste quantified, ROI re-proved. 3 pw.', None),
        ('Jan 2027', 'First signature', '€15-50K diagnostic + 15-35% gain share; quarterly recurrence begins.', 'ga'),
    ],
    [('Build effort', '6-8 person-wks'), ('Delivery per install', '2-3 wks · 3-4 pw'),
     ('Pod', 'VC 0.5 + SE 0.25'), ('Model', 'Recurring + gain share')],
    'Gain-share pricing of 15-35% of realised savings is market-validated (ProsperOps, nOps, Vantage), with a '
    '“10% or free” floor. Current cost-model sample is n=2 use cases; the keystone study raises it to ~5.',
    notes='This is backlog item 2.2 (keystone, do first) + 2.3 (the SKU). Partner with Aayushi on the value-assurance '
          'line. Only Backbase sees cost-per-resolved-outcome; external FinOps tools see the cloud bill.')

product_plan(
    'Execution plan 3 of 4 · Guardrail Studio',
    'Guardrail Studio attaches to every Conversational Banking deal in flight.',
    [
        ('Jan 2027', 'Design', 'Decision & Policy tooling alignment; risk and compliance SME briefed. 3 pw.', None),
        ('Feb-Mar 2027', 'Prototype', 'Authority-map demo: recommend, approve, execute per process. 5 pw.', None),
        ('Apr-May 2027', 'Proof', 'Paid pilot: guardrail backlog, entitlements map, Sentinel install blueprint. 6 pw.', None),
        ('Jun 2027', 'First signature', 'Sold at €50-100K, riding Transact and Resolve deals in flight.', 'ga'),
    ],
    [('Build effort', '8-10 person-wks'), ('Delivery per install', '5-6 wks · 8-10 pw'),
     ('Pod', 'VC 1.0 + SE 0.5 + SME 0.2'), ('Pre-binds', 'Sentinel')],
    'Timed to the Conversational Banking pipeline: Transact and Resolve capabilities require governed execution, '
    'which makes the Sentinel-readiness wedge a natural attach on active deals.',
    notes='Guardrail Studio productizes Decision & Policy. It codifies authority, policies, entitlements and audit '
          'needs into executable guardrails, and the output is the Sentinel deployment blueprint.')

product_plan(
    'Execution plan 4 of 4 · Ontology Cartographer',
    'Ontology Cartographer is the deepest build and lands where X-Ray has created demand.',
    [
        ('Feb 2027', 'R&D gate', 'Alignment with the Semantic Modeler and Nexus teams on tooling reuse.', 'decision'),
        ('Mar-Apr 2027', 'Prototype', 'Ontology mapping demo on two sample cores (core, CRM, LMS, credit). 6 pw, SE-led.', None),
        ('May-Jul 2027', 'Proof', 'Paid pilot: ontology coverage map, data-to-truth gap report, Nexus blueprint. 8 pw.', None),
        ('Aug 2027', 'First signature', 'Sold at €75-100K, attached to Relationship Intelligence opportunities.', 'ga'),
    ],
    [('Build effort', '10-12 person-wks'), ('Delivery per install', '6-8 wks · 12-16 pw'),
     ('Pod', 'SE 1.0 + VC 0.5 (SE-led)'), ('Pre-binds', 'Nexus')],
    'Sequenced last deliberately: it has the highest technical depth, needs R&D alignment, and sells best into '
    'accounts where an X-Ray install has already evidenced the cost of fragmented data.',
    notes='Cartographer productizes the Semantic Modeler. Relationship Intelligence (Grow plays, share of wallet) '
          'requires shared truth, so every RI opportunity is a Cartographer prospect.')

# ═════════════════════════════════════════════════════════════
# S14 — The one table: effort & pod register
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Execution plans · the register',
       'Twelve months of product build fits inside the existing pod.',
       footnote='pw = person-weeks. Build = one-time productization effort. Delivery = elapsed and effort figures per '
                'client install once productized. All efforts carry a ±30% planning tolerance.',
       notes='Total build load Aug 2026 - Jul 2027: 32-40 person-weeks across four products, sequenced so the pod '
             'never carries two builds at once.')
cols = [('Product', 300), ('One-time build', 190), ('Delivery per install', 230),
        ('Pod', 250), ('First revenue', 150)]
rows = [
    ('Process X-Ray', '8-10 pw', '4-6 wks · 12-14 pw', 'VC 1.0 + SE 0.5', 'Jan 2027'),
    ('Value Telemetry', '6-8 pw', '2-3 wks · 3-4 pw', 'VC 0.5 + SE 0.25', 'Jan 2027'),
    ('Guardrail Studio', '8-10 pw', '5-6 wks · 8-10 pw', 'VC 1.0 + SE 0.5 + SME', 'Jun 2027'),
    ('Ontology Cartographer', '10-12 pw', '6-8 wks · 12-16 pw', 'SE 1.0 + VC 0.5', 'Aug 2027'),
]
tx, ty, rh = 96, 240, 66
x = tx
for name, w in cols:
    txt(sl, name.upper(), x, ty, w - 16, 24, size=11.5, color=MUTED, bold=True)
    x += w
hairline(sl, tx, ty + 30, sum(w for _, w in cols), 2.2, color=NAVY)
for r, row in enumerate(rows):
    y = ty + 44 + r * rh
    x = tx
    for c, (val) in enumerate(row):
        w = cols[c][1]
        txt(sl, val, x, y + 12, w - 16, 40, size=13.5, bold=(c == 0),
            color=NAVY if c == 0 else RGBColor(0x2A, 0x38, 0x4A))
        x += w
    hairline(sl, tx, y + rh - 6, sum(w for _, w in cols))

# ═════════════════════════════════════════════════════════════
# S15 — Chapter 03
# ═════════════════════════════════════════════════════════════
sl = new_slide(NAVY)
txt(sl, '03', 96, 170, 400, 170, size=110, color=T50)
txt(sl, 'Mobilization and the ask', 96, 348, 900, 70, size=40, color=WHITE, bold=True)
txt(sl, 'The twelve-month sequence, how the Delivery Factory carries it, and the four decisions that unlock it.',
    96, 430, 940, 56, size=15, color=LIGHT_ON_NAVY, sp=1.25)
logo(sl, dark_bg=True)
txt(sl, str(PAGE[0]), 96, 720 - 36, 60, 22, size=12, color=RGBColor(0x5A, 0x6B, 0x80))

# ═════════════════════════════════════════════════════════════
# S16 — Portfolio wave roadmap
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Mobilization · the sequence',
       'One product at a time reaches four sellable SKUs and recurring revenue inside twelve months.',
       footnote='Red diamonds are decision gates (SE named · pilot green-light · R&D alignment). Blue diamonds mark first '
                'revenue per product. Bars show elapsed calendar time, and effort within them is per the register.',
       notes='Sequencing logic: X-Ray first (lowest build risk, front-door product), Telemetry parallel (shares the '
             'keystone cost model work), Guardrail rides the Conversational pipeline, Cartographer lands on X-Ray demand.')
months = ['Aug', 'Sep', 'Oct', 'Nov', 'Dec', 'Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul']
gx0, gy0, colw, rowh = 300, 252, 76, 62
for i, m in enumerate(months):
    txt(sl, m, gx0 + i * colw, gy0 - 30, colw, 22, size=10.5, color=MUTED, bold=True,
        align=PP_ALIGN.CENTER)
    hairline(sl, gx0 + i * colw, gy0 - 6, 1, h=5 * rowh, color=RGBColor(0xE4, 0xE8, 0xEE))
txt(sl, "2026", gx0, gy0 - 52, 200, 20, size=10, color=MUTED)
txt(sl, "2027", gx0 + 5 * colw, gy0 - 52, 200, 20, size=10, color=MUTED)
grows = [
    ('Keystone study', [(1, 2, BLUE_DARK, 'With Deepak')], [(1.0, RED)]),
    ('Process X-Ray', [(0, 1, T50, 'Design'), (1, 2, BLUE, 'Prototype'), (3, 2, BLUE_DARK, 'Proof'), (5, 7, NAVY, 'Sell + install')], [(2.0, RED), (5.0, BLUE)]),
    ('Value Telemetry', [(1, 1, T50, ''), (2, 1, BLUE, ''), (3, 2, BLUE_DARK, 'Proof'), (5, 7, NAVY, 'Recurring')], [(5.0, BLUE)]),
    ('Guardrail Studio', [(5, 1, T50, ''), (6, 2, BLUE, 'Prototype'), (8, 2, BLUE_DARK, 'Proof'), (10, 2, NAVY, 'Sell')], [(10.0, BLUE)]),
    ('Ontology Cartographer', [(6, 1, T50, ''), (7, 2, BLUE, 'Prototype'), (9, 3, BLUE_DARK, 'Proof')], [(6.0, RED)]),
]
for r, (label, bars, gates) in enumerate(grows):
    y = gy0 + r * rowh
    txt(sl, label, 96, y + 10, 196, 44, size=12, bold=True, sp=1.05)
    for (start, dur, color, blabel) in bars:
        bx = gx0 + start * colw + 2
        bw = dur * colw - 4
        rect(sl, bx, y + 8, bw, 26, fill=color, round_=True)
        if blabel and bw > 120:
            txt(sl, blabel, bx + 10, y + 11, bw - 16, 20, size=9.5, color=WHITE, bold=True)
    for (gm, gcol) in gates:
        diamond(sl, gx0 + gm * colw, y + 21, 16, gcol)

# ═════════════════════════════════════════════════════════════
# S17 — Delivery Factory usage
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Mobilization · the Delivery Factory',
       'The Delivery Factory carries every stage from prototype to installed layer.',
       footnote='Factory toolchain per the Banking OS canon: Process & Workspace Designer · Semantic Modeler · Agent '
                'Builder · Decision & Policy · Connector Studio · Simulation & Testing · Deployment & Ops Control.',
       notes='The insight from the PDP session: the Factory already contains the discovery-end tools. The products ARE '
             'those tools, productized as standalone paid wedges. Delivery cost stays at person-weeks because nothing '
             'is built from scratch.')
stages = [
    ('Prototype', 'Semantic Modeler +\nConnector Studio', 'A working demo on the bank’s own data, in weeks.'),
    ('Proof', 'Simulation & Testing +\nOps Control telemetry', 'CFO-grade evidence runs with zero production risk.'),
    ('Install', 'Mission Contract with\nNexus + Sentinel binding', 'The Banking OS layer live in 6-12 weeks, via the agentic SDLC.'),
    ('Run', 'Deployment &\nOps Control', 'Live telemetry that feeds Value Telemetry and the expansion case.'),
]
cw, cgap = 262, 24
for i, (stage, tools, out) in enumerate(stages):
    x = 96 + i * (cw + cgap)
    rect(sl, x, 240, cw, 280, fill=OFF, line=BORDER, line_w=0.75, round_=True)
    txt(sl, stage.upper(), x + 22, 262, cw - 40, 24, size=11.5, color=BLUE, bold=True)
    txt(sl, tools, x + 22, 292, cw - 40, 70, size=14.5, bold=True, sp=1.12)
    hairline(sl, x + 22, 378, cw - 44)
    txt(sl, out, x + 22, 392, cw - 40, 110, size=12, color=RGBColor(0x2A, 0x38, 0x4A), sp=1.2)
    if i < 3:
        txt(sl, '→', x + cw + 1, 360, 22, 30, size=15, color=MUTED)

# ═════════════════════════════════════════════════════════════
# S18 — The asks + €5K seed
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Mobilization · the ask',
       'Four decisions unlock execution; the €5,000 talent budget seeds product number one.',
       footnote='Quota figure of ~€200K recognition and pilot-account choice are open PDP items (1.3, 1.4) for the '
                'next 1:1 with Mayur. Keystone study scope sits with Deepak (backlog item 2.2).',
       notes='Framing for Tim: the talent budget becomes seed capital for a product launch, and specifically not '
             'a training spend. Framing for Mayur: these four asks are the same asks already on the PDP backlog.')
asks = [
    ('1 · Name the solution engineer', '0.5 FTE from the regional pool, committed by September. The single critical-path resource.'),
    ('2 · Green-light one prospect pilot', 'AIB-profile account for the X-Ray proof, by October. A prospect, keeping existing customers out of scope.'),
    ('3 · Back the keystone study', 'Cost-per-outcome model with Deepak in September. Value Telemetry and the FinOps line depend on it.'),
    ('4 · Agree the commercial treatment', 'Quota recognition of ~€200K so the line carries a real target from day one.'),
]
for i, (h_, b_) in enumerate(asks):
    chip(sl, 96, 226 + i * 92, 560, 80, h_, b_, accent=BLUE)
hero_card(sl, 692, 226, 528, 356, 'The €5,000 seed',
          'The talent budget launches product number one.',
          'Packaging and delivery kit for Process X-Ray, a demo built on the APA V3 catalog, launch '
          'collateral produced with marketing, and pricing pressure-tests with three friendly accounts.')

# ═════════════════════════════════════════════════════════════
# S19 — Planning basis / assumptions
# ═════════════════════════════════════════════════════════════
sl = new_slide()
chrome(sl, 'Mobilization · planning basis',
       'The plan stands on five explicit assumptions, each with an owner to validate.',
       footnote='Confidence ratings follow the repo standard: every assumption is documented, conservative, and '
                'sensitivity-tested before any figure becomes client-facing or budget-committed.',
       notes='Per CLAUDE.md governance: no hidden assumptions. These five carry the plan; validating the first three '
             'in September de-risks the whole sequence.')
assumps = [
    ('Team cost ~€1.2M a year', 'PDP working figure', 'Low', 'Mayur + finance'),
    ('Effort estimates within ±30%', 'Analogous VC engagements; first product build actuals will recalibrate', 'Medium', 'First build'),
    ('SE at 0.5 FTE is available', 'Regional pool capacity, unconfirmed', 'Low', 'SE leadership'),
    ('Tickets of €15-100K clear the market', 'AIB paid ~€80K for a weaker ADS; SAP and AWS wedge benchmarks', 'Medium', 'First 3 deals'),
    ('Wedge-to-Mission attach of 50%+ in 6 months', 'SAP and ServiceNow expansion patterns', 'Low', 'Pilot cohort'),
]
conf_color = {'Low': RED, 'Medium': AMBER, 'High': GREEN}
for i, (a, basis, conf, owner) in enumerate(assumps):
    y = 232 + i * 72
    rect(sl, 96, y, 1124, 62, fill=OFF, round_=True)
    txt(sl, a, 120, y + 10, 400, 44, size=13, bold=True, sp=1.05)
    txt(sl, basis, 536, y + 10, 420, 44, size=11.5, color=RGBColor(0x2A, 0x38, 0x4A), sp=1.1)
    dot(sl, 986, y + 31, 12, conf_color[conf])
    txt(sl, conf, 1000, y + 20, 80, 24, size=11.5, color=MUTED, bold=True)
    txt(sl, 'Validate: ' + owner, 1078, y + 20, 140, 40, size=10.5, color=BLUE, bold=True, sp=1.0)

# ═════════════════════════════════════════════════════════════
# S20 — Close
# ═════════════════════════════════════════════════════════════
sl = new_slide(NAVY)
rect(sl, 96, 170, 56, 5, fill=BLUE)
txt(sl, 'Land fast. Prove value. Expand.', 96, 200, 1000, 80, size=42, color=WHITE, bold=True)
txt(sl, 'One product in market by January. Cost neutral inside year one.\n'
        'Four installed-platform wedges and a recurring line by mid 2027.',
    96, 300, 980, 80, size=17, color=LIGHT_ON_NAVY, sp=1.35)
steps_now = [
    'This week · book the Mayur 1:1: quota treatment, pilot green-light, SE allocation.',
    'This week · follow up with Tim: category story, launch collateral, talent-budget plan.',
    'September · scope the keystone cost-per-outcome study with Deepak.',
]
for i, s_ in enumerate(steps_now):
    dot(sl, 106, 436 + i * 52, 10, CYAN)
    txt(sl, s_, 130, 424 + i * 52, 1000, 36, size=14, color=WHITE, sp=1.1)
logo(sl, dark_bg=True)
txt(sl, str(PAGE[0]), 96, 720 - 36, 60, 22, size=12, color=RGBColor(0x5A, 0x6B, 0x80))

out = os.path.join(KIT, 'Product_Factory_Execution_Plan_Exhibit.pptx')
prs.save(out)
import pathlib
print(f'Saved {out} ({pathlib.Path(out).stat().st_size // 1024} KB, {len(prs.slides.slides._sldIdLst) if hasattr(prs.slides, "slides") else len(prs.slides._sldIdLst)} slides)')
