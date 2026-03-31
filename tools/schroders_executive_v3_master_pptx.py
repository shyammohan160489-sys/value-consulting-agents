#!/usr/bin/env python3
"""
Schroders Executive Read-Out v3 — PPTX Generator (Backbase Master Template)

10-slide simplified deck for COO meeting (March 9, 2026).
Uses the Backbase Master Template (20" x 11.25").
All coordinates scaled 1.5x from the standard 13.333" x 7.5" version.

Usage:
    python3 tools/schroders_executive_v3_master_pptx.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Path to the Backbase Master Template
TEMPLATE_PATH = str(Path(__file__).parent.parent
                    / 'knowledge/learnings/ignite/raw/engagement-plans'
                    / 'Backbase Master Template - 2022 (1).pptx')


class SchrordersV3MasterPptxGenerator:
    """Generates a 10-slide executive read-out PPTX using Backbase Master Template."""

    # ── Dimensions (20" x 11.25" from template) ──────────────
    SLIDE_W = Inches(20.0)
    SLIDE_H = Inches(11.25)

    # ── Colors (Backbase Brand) ────────────────────────────
    DARK_BG    = RGBColor(0x09, 0x1C, 0x35)
    CARD_BG    = RGBColor(0x0E, 0x22, 0x40)
    BLUE       = RGBColor(0x33, 0x66, 0xFF)
    PURPLE     = RGBColor(0x7B, 0x2F, 0xFF)
    RED        = RGBColor(0xDC, 0x26, 0x26)
    GREEN      = RGBColor(0x05, 0x96, 0x69)
    AMBER      = RGBColor(0xD9, 0x77, 0x06)
    CYAN       = RGBColor(0x08, 0x91, 0xB2)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BG   = RGBColor(0xF5, 0xFA, 0xFF)
    LIGHT_CARD = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_TEXT   = RGBColor(0x09, 0x1C, 0x35)
    SUB_TEXT    = RGBColor(0x64, 0x74, 0x8B)
    MUTED      = RGBColor(0x94, 0xA3, 0xB8)
    BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
    BORDER_DK  = RGBColor(0x1A, 0x3A, 0x5E)
    BLUE_LIGHT = RGBColor(0xEF, 0xF6, 0xFF)
    GREEN_LIGHT = RGBColor(0xEC, 0xFD, 0xF5)
    AMBER_LIGHT = RGBColor(0xFF, 0xFB, 0xEB)
    PURPLE_LIGHT = RGBColor(0xF5, 0xF3, 0xFF)

    FONT = 'Libre Franklin'
    # Margins and content area scaled 1.5x
    ML = Inches(0.9)      # 0.6 * 1.5
    MR = Inches(0.9)      # 0.6 * 1.5
    CW = Inches(18.15)    # 12.1 * 1.5
    CT = Inches(2.4)      # 1.6 * 1.5

    # ── HELPERS ─────────────────────────────────────────────

    def _new_slide(self, dark=False):
        slide = self.prs.slides.add_slide(self.blank_layout)
        # Remove any inherited placeholders
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)
        # Add background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.DARK_BG if dark else self.LIGHT_BG
        return slide

    def _txt(self, slide, text, left, top, width, height, size=Pt(18),
             color=None, bold=False, align=PP_ALIGN.LEFT, font=None):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color or self.DARK_TEXT
        p.font.name = font or self.FONT
        p.alignment = align
        return tf

    def _section_label(self, slide, text, top=Inches(0.6), color=None, dark=False):
        # top: 0.4*1.5=0.6, width CW, height Pt(24)=16*1.5, size Pt(14)=9*1.5
        self._txt(slide, text.upper(), self.ML, top, self.CW, Pt(24),
                  size=Pt(14), color=color or (self.BLUE if not dark else RGBColor(0x69, 0xFE, 0xFF)),
                  bold=True, align=PP_ALIGN.LEFT)

    def _heading(self, slide, main, accent='', top=Inches(1.125), size=Pt(48),
                 dark=False):
        # top: 0.75*1.5=1.125, height: 0.7*1.5=1.05, size: 32*1.5=48
        tb = slide.shapes.add_textbox(self.ML, top, self.CW, Inches(1.05))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = main
        r1.font.size = size
        r1.font.bold = True
        r1.font.color.rgb = self.WHITE if dark else self.DARK_TEXT
        r1.font.name = self.FONT
        if accent:
            r2 = p.add_run()
            r2.text = accent
            r2.font.size = size
            r2.font.bold = True
            r2.font.color.rgb = self.BLUE if not dark else RGBColor(0x69, 0xFE, 0xFF)
            r2.font.name = self.FONT

    def _subtitle(self, slide, text, top=Inches(1.875), dark=False):
        # top: 1.25*1.5=1.875, height: 0.35*1.5=0.525, size: 12*1.5=18
        self._txt(slide, text, self.ML, top, self.CW, Inches(0.525),
                  size=Pt(18), color=self.MUTED if dark else self.SUB_TEXT)

    def _card(self, slide, left, top, w, h, fill=None, border=None, dark=False):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = fill or (self.CARD_BG if dark else self.LIGHT_CARD)
        s.line.color.rgb = border or (self.BORDER_DK if dark else self.BORDER)
        s.line.width = Pt(1.5)  # 1*1.5
        return s

    def _stat_box(self, slide, value, label, left, top, w, h,
                  val_color=None, bg=None):
        """Compact stat box — value on top, label below."""
        self._card(slide, left, top, w, h, fill=bg or self.BLUE_LIGHT)
        # value: top+0.08*1.5=+0.12, height 0.35*1.5=0.525, size 20*1.5=30
        self._txt(slide, str(value), left, top + Inches(0.12), w, Inches(0.525),
                  size=Pt(30), color=val_color or self.BLUE, bold=True,
                  align=PP_ALIGN.CENTER)
        # label: top+0.42*1.5=+0.63, height Pt(18)=12*1.5, size Pt(11)=7*1.5 rounded
        self._txt(slide, label.upper(), left, top + Inches(0.63), w, Pt(18),
                  size=Pt(11), color=self.SUB_TEXT, bold=True,
                  align=PP_ALIGN.CENTER)

    def _multi_text(self, slide, runs, left, top, width, height, align=PP_ALIGN.LEFT):
        """Add textbox with multiple styled runs."""
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        for text, size, color, bold in runs:
            r = p.add_run()
            r.text = str(text)
            r.font.size = size
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = self.FONT
        return tf

    def _divider(self, slide, top):
        """Thin horizontal line."""
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   self.ML, top, self.CW, Pt(1.5))  # 1*1.5
        s.fill.solid()
        s.fill.fore_color.rgb = self.BORDER
        s.line.fill.background()

    # ── SLIDE BUILDERS ──────────────────────────────────────

    def _slide_01_cover(self):
        """Dark hero cover slide."""
        s = self._new_slide(dark=True)
        # Section label — top: 0.5*1.5=0.75
        self._section_label(s, 'Executive Read-Out \u00b7 March 9th, 2026',
                           top=Inches(0.75), dark=True)
        # Mega title — top: 1.3*1.5=1.95, height: 2.0*1.5=3.0
        tb = s.shapes.add_textbox(self.ML, Inches(1.95), self.CW, Inches(3.0))
        tf = tb.text_frame
        tf.word_wrap = True
        for line, color in [('Accelerating the', self.WHITE),
                            ('Schroders Advantage.', self.BLUE)]:
            p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = line
            r.font.size = Pt(72)   # 48*1.5
            r.font.bold = True
            r.font.color.rgb = color
            r.font.name = self.FONT
        # Subtitle — top: 3.2*1.5=4.8, width: 8*1.5=12, height: 0.6*1.5=0.9, size: 14*1.5=21
        self._txt(s, 'A digitally enabled, insight-led client and advisor platform.\nFrom relationship depth to digital scale.',
                  self.ML, Inches(4.8), Inches(12.0), Inches(0.9),
                  size=Pt(21), color=self.MUTED)

        # 4 pill tags
        tags = ['Advisor Productivity', 'Client Retention & AUM Growth',
                'Innovation Agility', 'No AI Lock-In']
        tx = self.ML
        for tag in tags:
            pw = Inches(len(tag) * 0.15 + 0.45)  # 0.1*1.5=0.15, 0.3*1.5=0.45
            self._card(s, tx, Inches(6.15), pw, Inches(0.45),  # 4.1*1.5=6.15, 0.3*1.5=0.45
                      fill=RGBColor(0x0E, 0x22, 0x40), border=self.BORDER_DK,
                      dark=True)
            self._txt(s, tag, tx + Inches(0.15), Inches(6.18), pw - Inches(0.3),  # 0.1*1.5=0.15, 4.12*1.5=6.18, 0.2*1.5=0.3
                     Inches(0.375),  # 0.25*1.5
                     size=Pt(12), color=self.MUTED, bold=True)  # 8*1.5=12
            tx += pw + Inches(0.225)  # 0.15*1.5=0.225

        # 3 hero stats — stat_w: 3.4*1.5=5.1, gap: 0.4*1.5=0.6
        stat_w = Inches(5.1)
        gap = Inches(0.6)
        sx = self.ML
        for val, label, vcol in [
            ('\u00a3148M', 'Value at Stake', self.GREEN),
            ('~7.0\u00d7', 'Return on Investment', self.BLUE),
            ('Month 8', 'Breakeven', self.WHITE),
        ]:
            # top: 5.0*1.5=7.5, height: 1.2*1.5=1.8
            self._card(s, sx, Inches(7.5), stat_w, Inches(1.8),
                      fill=self.CARD_BG, dark=True)
            # value: 5.1*1.5=7.65, height: 0.5*1.5=0.75, size: 28*1.5=42
            self._txt(s, val, sx, Inches(7.65), stat_w, Inches(0.75),
                     size=Pt(42), color=vcol, bold=True, align=PP_ALIGN.CENTER)
            # label: 5.65*1.5=8.475, height: Pt(21)=14*1.5, size: 8*1.5=12
            self._txt(s, label.upper(), sx, Inches(8.475), stat_w, Pt(21),
                     size=Pt(12), color=self.MUTED, bold=True, align=PP_ALIGN.CENTER)
            sx += stat_w + gap

        # Footer — top: 6.9*1.5=10.35, height: Pt(18)=12*1.5, size: 8*1.5=12
        self._txt(s, 'v3.0 \u00b7 Schroders \u00d7 Backbase',
                  self.ML, Inches(10.35), self.CW, Pt(18),
                  size=Pt(12), color=RGBColor(0x3A, 0x49, 0x5D))

    def _slide_02_context(self):
        """The Case Is Clear — context and maturity."""
        s = self._new_slide()
        self._section_label(s, 'Context \u00b7 Where We Stand')
        self._heading(s, 'The Case Is Clear. ', 'The Moment Is Now.')
        self._subtitle(s, '10 workshops \u00b7 12 capabilities assessed \u00b7 Master analysis validated')

        # Today stats row — top: 1.7*1.5=2.55, size: 8*1.5=12
        self._txt(s, 'TODAY \u00b7 SCHRODERS WEALTH', self.ML, Inches(2.55), self.CW, Pt(18),
                 size=Pt(12), color=self.BLUE, bold=True)

        # stats — top: 2.0*1.5=3.0, w: 3.4*1.5=5.1, h: 0.7*1.5=1.05, gap: 3.8*1.5=5.7
        stats = [('225', 'Client Advisors'), ('35K+', 'Clients'), ('\u00a372B', 'AUM')]
        sx = self.ML
        for val, label in stats:
            self._stat_box(s, val, label, sx, Inches(3.0), Inches(5.1), Inches(1.05))
            sx += Inches(5.7)

        # Maturity assessment — divider: 3.0*1.5=4.5
        self._divider(s, Inches(4.5))
        # label: 3.2*1.5=4.8, size: 8*1.5=12
        self._txt(s, 'CAPABILITY MATURITY ASSESSMENT', self.ML, Inches(4.8), self.CW, Pt(18),
                 size=Pt(12), color=self.BLUE, bold=True)

        # Current vs Target card — top: 3.5*1.5=5.25, w: 5.5*1.5=8.25, h: 1.6*1.5=2.4
        self._card(s, self.ML, Inches(5.25), Inches(8.25), Inches(2.4))
        # "12" — left+0.3*1.5=+0.45, top: 3.55*1.5=5.325, w: 1*1.5=1.5, h: 0.5*1.5=0.75, size: 32*1.5=48
        self._txt(s, '12', self.ML + Inches(0.45), Inches(5.325), Inches(1.5), Inches(0.75),
                 size=Pt(48), color=self.BLUE, bold=True)
        # "Core Capabilities" — left+1.3*1.5=+1.95, top: 3.65*1.5=5.475, w: 3*1.5=4.5, h: Pt(21)=14*1.5, size: 12*1.5=18
        self._txt(s, 'Core Capabilities', self.ML + Inches(1.95), Inches(5.475), Inches(4.5), Pt(21),
                 size=Pt(18), color=self.DARK_TEXT, bold=True)
        # Current score — left+0.3*1.5=+0.45, top: 4.2*1.5=6.3, w: 4*1.5=6, h: Pt(18)=12*1.5, size: 11*1.5=17 (round)
        self._txt(s, 'Current: 0.6 / 4.0', self.ML + Inches(0.45), Inches(6.3), Inches(6.0), Pt(18),
                 size=Pt(17), color=self.RED, bold=True)
        # Target — top: 4.55*1.5=6.825
        self._txt(s, 'Target: 3.5 / 4.0', self.ML + Inches(0.45), Inches(6.825), Inches(6.0), Pt(18),
                 size=Pt(17), color=self.GREEN, bold=True)

        # Value at stake card — left: 6.8*1.5=10.2, top: 3.5*1.5=5.25, w: 5.5*1.5=8.25, h: 1.6*1.5=2.4
        self._card(s, Inches(10.2), Inches(5.25), Inches(8.25), Inches(2.4),
                  fill=self.BLUE_LIGHT, border=self.BLUE)
        # "148M" — left: 7.0*1.5=10.5, top: 3.6*1.5=5.4, w: 4*1.5=6, h: 0.5*1.5=0.75, size: 36*1.5=54
        self._txt(s, '\u00a3148M', Inches(10.5), Inches(5.4), Inches(6.0), Inches(0.75),
                 size=Pt(54), color=self.GREEN, bold=True)
        # subtitle — left: 7.0*1.5=10.5, top: 4.15*1.5=6.225, w: 5*1.5=7.5, h: Pt(21)=14*1.5, size: 11*1.5=17
        self._txt(s, 'Five-Year Value (AI from Year 2)', Inches(10.5), Inches(6.225), Inches(7.5), Pt(21),
                 size=Pt(17), color=self.DARK_TEXT)
        # ROI — left: 7.0*1.5=10.5, top: 4.55*1.5=6.825, w: 2*1.5=3, size: 14*1.5=21
        self._txt(s, '~7.0\u00d7', Inches(10.5), Inches(6.825), Inches(3.0), Pt(21),
                 size=Pt(21), color=self.BLUE, bold=True)
        # "Return on Investment" — left: 8.2*1.5=12.3, top: 4.58*1.5=6.87, w: 3*1.5=4.5, size: 10*1.5=15
        self._txt(s, 'Return on Investment', Inches(12.3), Inches(6.87), Inches(4.5), Pt(21),
                 size=Pt(15), color=self.SUB_TEXT)

        # Bottom insight — left: ML, top: 5.6*1.5=8.4, w: CW, h: 1.0*1.5=1.5
        self._card(s, self.ML, Inches(8.4), self.CW, Inches(1.5),
                  fill=self.LIGHT_CARD, border=self.BORDER)
        # text — left+0.3*1.5=+0.45, top: 5.7*1.5=8.55, w: CW-0.6*1.5=CW-0.9, h: 0.8*1.5=1.2, size: 10*1.5=15
        self._txt(s, 'Every peer is investing. SJP: \u00a3150M+ tech transformation. Quilter: \u00a340M+ platform modernisation. JP Morgan: $2B+ AI investment. The question is not whether, but how quickly.',
                  self.ML + Inches(0.45), Inches(8.55), self.CW - Inches(0.9), Inches(1.2),
                  size=Pt(15), color=self.SUB_TEXT)

    def _slide_03_platform(self):
        """What You Get — 3-card platform overview."""
        s = self._new_slide()
        self._section_label(s, 'The Platform')
        self._heading(s, 'One Platform. ', 'Every Channel. Every Data Source.')
        self._subtitle(s, 'Not a channel replacement. A platform transformation \u2014 connecting advisors, clients, systems and intelligence.')

        # 3 cards — cw: 3.7*1.5=5.55, gap: 0.3*1.5=0.45, ch: 4.4*1.5=6.6
        cw = Inches(5.55)
        gap = Inches(0.45)
        ch = Inches(6.6)
        cx = self.ML

        cards = [
            ('Digital Channels', self.BLUE, 'Client & Employee Experiences', [
                'Client portal (web & mobile)',
                'Prospect portal & onboarding',
                'RM Workspace & Digital Assist',
                'Unified employee dashboard',
                'Replaces Temenos e-Services',
            ]),
            ('Grand Central', self.GREEN, 'Data Connectivity Backbone', [
                'T24, Salesforce, Folio connectivity',
                'Configuration-based, not code',
                'API-first, event-driven',
                'Pre-built wealth connectors',
                'Backbase Marketplace',
                'Pre-existing connectors with partners',
                'Group-wide potential',
            ]),
            ('Intelligence Fabric', self.PURPLE, 'AI Orchestration Layer', [
                'Multi-model AI (Claude, GPT, Gemini)',
                'Custom Agent Factory',
                'Productised agent library',
                'RAG, observability, auditability',
                'Multi-agent governance',
            ]),
        ]

        for title, color, sub, bullets in cards:
            # card top: 1.8*1.5=2.7
            self._card(s, cx, Inches(2.7), cw, ch, border=color)
            # Colored top bar — top: 1.8*1.5=2.7, height: Pt(6)=4*1.5
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     cx, Inches(2.7), cw, Pt(6))
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()
            # Title — +0.2*1.5=+0.3, top: 2.0*1.5=3.0, w: cw-0.4*1.5=cw-0.6, h: Pt(21)=14*1.5, size: 10*1.5=15
            self._txt(s, title.upper(), cx + Inches(0.3), Inches(3.0), cw - Inches(0.6), Pt(21),
                     size=Pt(15), color=color, bold=True)
            # Subtitle — top: 2.3*1.5=3.45, size: 11*1.5=17
            self._txt(s, sub, cx + Inches(0.3), Inches(3.45), cw - Inches(0.6), Pt(18),
                     size=Pt(17), color=self.DARK_TEXT, bold=True)
            # Bullets — start: 2.7*1.5=4.05, step: 0.3*1.5=0.45
            by = Inches(4.05)
            for b in bullets:
                # size: 9*1.5=14
                self._txt(s, '\u2192 ' + b, cx + Inches(0.3), by, cw - Inches(0.6), Pt(21),
                         size=Pt(14), color=self.SUB_TEXT)
                by += Inches(0.45)
            cx += cw + gap

        # Bottom insight — top: 6.4*1.5=9.6, h: 0.7*1.5=1.05
        self._card(s, self.ML, Inches(9.6), self.CW, Inches(1.05),
                  fill=self.BLUE_LIGHT, border=self.BLUE)
        # text — +0.3*1.5=+0.45, top: 6.5*1.5=9.75, w: CW-0.6*1.5=CW-0.9, h: 0.5*1.5=0.75, size: 10*1.5=15
        self._txt(s, 'Salesforce sees CRM data. Temenos sees core data. Backbase sees everything \u2014 and makes it available to every channel, every agent, every decision.',
                  self.ML + Inches(0.45), Inches(9.75), self.CW - Inches(0.9), Inches(0.75),
                  size=Pt(15), color=self.DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)

    def _slide_04_ai(self):
        """What You Get: AI — Agent Factory + Productised Library."""
        s = self._new_slide()
        self._section_label(s, 'Intelligence Fabric')
        self._heading(s, 'Intelligent by Design.')
        self._subtitle(s, 'Two paths to AI. Build your own with Agent Factory. Activate ours from the productised library. Same infrastructure, same governance.')

        # 2-column — hw: 5.8*1.5=8.7, gap: 0.3*1.5=0.45
        hw = Inches(8.7)
        gap = Inches(0.45)

        # Left — Agent Factory — top: 1.8*1.5=2.7, h: 2.3*1.5=3.45
        self._card(s, self.ML, Inches(2.7), hw, Inches(3.45), border=self.PURPLE)
        # label — +0.2*1.5=+0.3, top: 1.9*1.5=2.85, size: 8*1.5=12
        self._txt(s, 'CUSTOM AGENTS \u00b7 AGENT FACTORY', self.ML + Inches(0.3), Inches(2.85), hw, Pt(18),
                 size=Pt(12), color=self.PURPLE, bold=True)
        # subtitle — top: 2.15*1.5=3.225, size: 11*1.5=17
        self._txt(s, 'Design. Test. Deploy. Iterate.', self.ML + Inches(0.3), Inches(3.225), hw, Pt(18),
                 size=Pt(17), color=self.DARK_TEXT, bold=True)
        bullets_l = ['Build agents from Schroders\' own policies and data',
                     'Full development lifecycle within the platform',
                     'Your IP, your logic, your competitive advantage',
                     'Available from Day 1 on Enterprise']
        # bullets start: 2.55*1.5=3.825, step: 0.25*1.5=0.375
        by = Inches(3.825)
        for b in bullets_l:
            # +0.2*1.5=+0.3, w: hw-0.4*1.5=hw-0.6, size: 9*1.5=14
            self._txt(s, '\u2192 ' + b, self.ML + Inches(0.3), by, hw - Inches(0.6), Pt(18),
                     size=Pt(14), color=self.SUB_TEXT)
            by += Inches(0.375)

        # Right — Productised
        rx = self.ML + hw + gap
        self._card(s, rx, Inches(2.7), hw, Inches(3.45), border=self.BLUE)
        self._txt(s, 'PRODUCTISED AGENTS \u00b7 BACKBASE LIBRARY', rx + Inches(0.3), Inches(2.85), hw, Pt(18),
                 size=Pt(12), color=self.BLUE, bold=True)
        self._txt(s, 'Ready-made. Continuously Improved.', rx + Inches(0.3), Inches(3.225), hw, Pt(18),
                 size=Pt(17), color=self.DARK_TEXT, bold=True)
        bullets_r = ['Pre-built agents for wealth management tasks',
                     'Maintained and improved by Backbase R&D',
                     'Low-effort agents bundled with the platform',
                     'Complex use cases via BIC blocks as library expands']
        by = Inches(3.825)
        for b in bullets_r:
            self._txt(s, '\u2192 ' + b, rx + Inches(0.3), by, hw - Inches(0.6), Pt(18),
                     size=Pt(14), color=self.SUB_TEXT)
            by += Inches(0.375)

        # 7 validated use cases — top: 4.3*1.5=6.45, size: 8*1.5=12
        self._txt(s, 'VALIDATED USE CASES \u00b7 SCOPED WITH YOUR TEAM', self.ML, Inches(6.45), self.CW, Pt(18),
                 size=Pt(12), color=self.BLUE, bold=True)

        use_cases = [
            ('Review Prep', '\u00a32.5M/yr manual cost', self.PURPLE),
            ('Meeting Summary', 'Automated capture', self.PURPLE),
            ('Sentiment Analysis', 'Message intelligence', self.PURPLE),
            ('Auto Reply', 'Contextual drafting', self.BLUE),
            ('Doc Processing', 'Extraction & routing', self.BLUE),
            ('Smart Signals', 'Proactive alerts', self.BLUE),
            ('Personalised Outreach', 'Tailored comms', self.BLUE),
        ]
        # ucw: 1.6*1.5=2.4
        ucw = Inches(2.4)
        ucx = self.ML
        for name, desc, color in use_cases:
            # top: 4.55*1.5=6.825, h: 0.8*1.5=1.2
            self._card(s, ucx, Inches(6.825), ucw, Inches(1.2), border=color)
            # name — +0.08*1.5=+0.12, top: 4.6*1.5=6.9, w: ucw-0.16*1.5=ucw-0.24, h: Pt(15)=10*1.5, size: 7*1.5=11
            self._txt(s, name.upper(), ucx + Inches(0.12), Inches(6.9), ucw - Inches(0.24), Pt(15),
                     size=Pt(11), color=color, bold=True)
            # desc — top: 4.85*1.5=7.275, size: 7*1.5=11
            self._txt(s, desc, ucx + Inches(0.12), Inches(7.275), ucw - Inches(0.24), Pt(15),
                     size=Pt(11), color=self.SUB_TEXT)
            ucx += ucw + Inches(0.15)  # 0.1*1.5=0.15

        # Bottom cards: No AI Lock-In + Multi-Agent Governance
        # left: ML, top: 5.6*1.5=8.4, w: 5.8*1.5=8.7, h: 1.1*1.5=1.65
        self._card(s, self.ML, Inches(8.4), Inches(8.7), Inches(1.65), border=self.CYAN)
        # label — +0.2*1.5=+0.3, top: 5.65*1.5=8.475, size: 8*1.5=12
        self._txt(s, 'NO AI LOCK-IN', self.ML + Inches(0.3), Inches(8.475), Inches(7.5), Pt(15),
                 size=Pt(12), color=self.CYAN, bold=True)
        # text — top: 5.9*1.5=8.85, w: 5.4*1.5=8.1, h: 0.6*1.5=0.9, size: 9*1.5=14
        self._txt(s, 'Multi-model architecture \u2014 Claude, GPT, Gemini, open-source. Switch providers without re-engineering agents.',
                  self.ML + Inches(0.3), Inches(8.85), Inches(8.1), Inches(0.9),
                  size=Pt(14), color=self.SUB_TEXT)

        # right card — left: 6.7*1.5=10.05, top: 5.6*1.5=8.4
        self._card(s, Inches(10.05), Inches(8.4), Inches(8.7), Inches(1.65), border=self.PURPLE)
        # label — left: 6.9*1.5=10.35, top: 5.65*1.5=8.475
        self._txt(s, 'MULTI-AGENT GOVERNANCE', Inches(10.35), Inches(8.475), Inches(7.5), Pt(15),
                 size=Pt(12), color=self.PURPLE, bold=True)
        # text — left: 6.9*1.5=10.35, top: 5.9*1.5=8.85, w: 5.4*1.5=8.1
        self._txt(s, 'Backbase as the control layer for all AI agents \u2014 including Salesforce and other systems. One governance framework. Full auditability.',
                  Inches(10.35), Inches(8.85), Inches(8.1), Inches(0.9),
                  size=Pt(14), color=self.SUB_TEXT)

    def _slide_05_commercial(self):
        """Enterprise or Pay As You Grow — side by side."""
        s = self._new_slide()
        self._section_label(s, 'Commercial Model')
        self._heading(s, 'Enterprise or ', 'Pay As You Grow.')
        self._subtitle(s, 'Same platform. Same delivery. Same value. Two commercial structures.')

        # hw: 5.8*1.5=8.7, gap: 0.3*1.5=0.45
        hw = Inches(8.7)
        gap = Inches(0.45)

        # -- Enterprise (left) -- top: 1.7*1.5=2.55, h: 4.0*1.5=6.0
        self._card(s, self.ML, Inches(2.55), hw, Inches(6.0),
                  fill=self.BLUE_LIGHT, border=self.BLUE)
        # Recommended badge — left: ML+hw-1.5*1.5=ML+hw-2.25, top: 1.75*1.5=2.625, w: 1.4*1.5=2.1, h: 0.25*1.5=0.375
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   self.ML + hw - Inches(2.25), Inches(2.625),
                                   Inches(2.1), Inches(0.375))
        badge.fill.solid()
        badge.fill.fore_color.rgb = self.BLUE
        badge.line.fill.background()
        # badge text — size: 7*1.5=11
        self._txt(s, 'RECOMMENDED', self.ML + hw - Inches(2.25), Inches(2.64),
                 Inches(2.1), Inches(0.33),
                 size=Pt(11), color=self.WHITE, bold=True, align=PP_ALIGN.CENTER)

        # label — +0.2*1.5=+0.3, top: 1.85*1.5=2.775, size: 8*1.5=12
        self._txt(s, 'OPTION A \u00b7 ENTERPRISE', self.ML + Inches(0.3), Inches(2.775),
                 hw, Pt(18), size=Pt(12), color=self.BLUE, bold=True)

        # Year boxes — bw: 1.0*1.5=1.5, bx start: ML+0.2*1.5=ML+0.3
        ent_years = [('Y1', '\u00a32.3M', '60% discount'), ('Y2', '\u00a32.6M', '55% discount'),
                     ('Y3', '\u00a33.6M', '55% discount'), ('Y4', '\u00a34.0M', '50% discount'),
                     ('Y5', '\u00a34.4M', '45% discount')]
        bw = Inches(1.5)
        bx = self.ML + Inches(0.3)
        for yr, price, disc in ent_years:
            # top: 2.2*1.5=3.3, h: 0.8*1.5=1.2
            self._card(s, bx, Inches(3.3), bw, Inches(1.2), fill=self.LIGHT_CARD)
            # yr — top: 2.22*1.5=3.33, h: Pt(12)=8*1.5, size: 7*1.5=11
            self._txt(s, yr, bx, Inches(3.33), bw, Pt(12),
                     size=Pt(11), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)
            # price — top: 2.4*1.5=3.6, h: Pt(30)=20*1.5, size: 16*1.5=24
            self._txt(s, price, bx, Inches(3.6), bw, Pt(30),
                     size=Pt(24), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)
            # disc — top: 2.72*1.5=4.08, h: Pt(12)=8*1.5, size: 6*1.5=9
            self._txt(s, disc, bx, Inches(4.08), bw, Pt(12),
                     size=Pt(9), color=self.GREEN, bold=True, align=PP_ALIGN.CENTER)
            bx += bw + Inches(0.12)  # 0.08*1.5=0.12

        # Description — +0.2*1.5=+0.3, top: 3.15*1.5=4.725, w: hw-0.4*1.5=hw-0.6, h: 0.5*1.5=0.75, size: 9*1.5=14
        self._txt(s, 'Full platform from Day 1. Scaling entry \u2014 pay less upfront, more as value compounds. All AI capabilities immediately.',
                  self.ML + Inches(0.3), Inches(4.725), hw - Inches(0.6), Inches(0.75),
                  size=Pt(14), color=self.SUB_TEXT)

        # Enterprise summary stats — top: 3.75*1.5=5.625, step: 1.8*1.5=2.7
        for i, (val, label) in enumerate([('\u00a316.9M', '5-YR LICENCE'), ('~7.0\u00d7', 'ROI'), ('Month 8', 'BREAKEVEN')]):
            ex = self.ML + Inches(0.3) + i * Inches(2.7)
            # val — w: 1.6*1.5=2.4, h: Pt(27)=18*1.5, size: 14*1.5=21
            self._txt(s, val, ex, Inches(5.625), Inches(2.4), Pt(27),
                     size=Pt(21), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)
            # label — top: 4.05*1.5=6.075, h: Pt(15)=10*1.5, size: 7*1.5=11
            self._txt(s, label, ex, Inches(6.075), Inches(2.4), Pt(15),
                     size=Pt(11), color=self.SUB_TEXT, bold=True, align=PP_ALIGN.CENTER)

        # -- PAG (right) --
        rx = self.ML + hw + gap
        self._card(s, rx, Inches(2.55), hw, Inches(6.0),
                  fill=self.AMBER_LIGHT, border=self.AMBER)
        self._txt(s, 'OPTION B \u00b7 PAY AS YOU GROW', rx + Inches(0.3), Inches(2.775),
                 hw, Pt(18), size=Pt(12), color=self.AMBER, bold=True)

        # PAG year boxes — bw_pag: 1.6*1.5=2.4, step: +0.15*1.5=+0.225
        pag_years = [('Y1', '\u00a32.7M', '30% discount'), ('Y2', '\u00a33.5M', '30% discount'),
                     ('Y3+', '\u00a34.8M', '30% discount')]
        bx = rx + Inches(0.3)
        bw_pag = Inches(2.4)
        for yr, price, disc in pag_years:
            self._card(s, bx, Inches(3.3), bw_pag, Inches(1.2), fill=self.LIGHT_CARD)
            self._txt(s, yr, bx, Inches(3.33), bw_pag, Pt(12),
                     size=Pt(11), color=self.AMBER, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, price, bx, Inches(3.6), bw_pag, Pt(30),
                     size=Pt(24), color=self.AMBER, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, disc, bx, Inches(4.08), bw_pag, Pt(12),
                     size=Pt(9), color=self.GREEN, bold=True, align=PP_ALIGN.CENTER)
            bx += bw_pag + Inches(0.225)

        # PAG description
        self._txt(s, 'Modular activation. License grows as modules go live. AI unlocks progressively. Agentic Runtime from Year 1.',
                  rx + Inches(0.3), Inches(4.725), hw - Inches(0.6), Inches(0.75),
                  size=Pt(14), color=self.SUB_TEXT)

        # PAG summary stats
        for i, (val, label) in enumerate([('\u00a320.6M', '5-YR LICENCE'), ('~5.9\u00d7', 'ROI'), ('Month 9', 'BREAKEVEN')]):
            ex = rx + Inches(0.3) + i * Inches(2.7)
            self._txt(s, val, ex, Inches(5.625), Inches(2.4), Pt(27),
                     size=Pt(21), color=self.AMBER, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, label, ex, Inches(6.075), Inches(2.4), Pt(15),
                     size=Pt(11), color=self.SUB_TEXT, bold=True, align=PP_ALIGN.CENTER)

        # Key insight — top: 5.85*1.5=8.775, h: 0.5*1.5=0.75
        self._card(s, self.ML, Inches(8.775), self.CW, Inches(0.75),
                  fill=self.BLUE_LIGHT, border=self.BLUE)
        # text — +0.3*1.5=+0.45, top: 5.92*1.5=8.88, w: CW-0.6*1.5=CW-0.9, h: 0.35*1.5=0.525, size: 10*1.5=15
        self._txt(s, 'Enterprise is \u00a30.4M less in Year 1 and \u00a33.7M less over five years. Every price increase maps to a specific module activation.',
                  self.ML + Inches(0.45), Inches(8.88), self.CW - Inches(0.9), Inches(0.525),
                  size=Pt(15), color=self.DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)

        # Tier swimlane — top: 6.45*1.5=9.675, size: 7*1.5=11
        self._txt(s, 'MODULE ACTIVATION BY YEAR', self.ML, Inches(9.675), self.CW, Pt(15),
                 size=Pt(11), color=self.BLUE, bold=True)

        tier_data = [
            ('Digital Wealth', 'Premium', '\u2192 Signature', 'Signature'),
            ('Digital Onboarding', '\u2014', '\u2605 Essential', 'Essential'),
            ('Digital Assist', 'Premium', 'Premium', 'Premium'),
            ('CLO', 'Essential', '\u2192 Premium', 'Premium'),
            ('RM Workspace', '\u2014', '\u2605 Premium', '\u2192 Signature'),
            ('Intelligence Fabric', '\u2605 Agentic', '+ Foundation AI', 'Full AI + Custom'),
        ]
        # ty: 6.6*1.5=9.9, tw: 2.0*1.5=3.0, col gap: 0.1*1.5=0.15
        ty = Inches(9.9)
        tw = Inches(3.0)
        cx_cols = [self.ML, self.ML + tw + Inches(0.15),
                   self.ML + 2*(tw + Inches(0.15)),
                   self.ML + 3*(tw + Inches(0.15))]
        # Headers — size: 6*1.5=9
        for ci, hdr in enumerate(['Module', 'Year 1', 'Year 2', 'Year 3+']):
            col = [self.DARK_TEXT, self.BLUE, self.GREEN, self.PURPLE][ci]
            self._txt(s, hdr, cx_cols[ci], ty, tw, Pt(12),
                     size=Pt(9), color=col, bold=True)
        ty += Pt(15)  # 10*1.5
        for mod, y1, y2, y3 in tier_data:
            for ci, val in enumerate([mod, y1, y2, y3]):
                col = self.DARK_TEXT if ci == 0 else self.SUB_TEXT
                # size: 5*1.5=8
                self._txt(s, val, cx_cols[ci], ty, tw, Pt(12),
                         size=Pt(8), color=col)
            ty += Pt(12)  # 8*1.5

    def _slide_06_phases(self):
        """Phased Delivery — 3-column merged phases."""
        s = self._new_slide()
        self._section_label(s, 'Delivery Roadmap')
        self._heading(s, 'Three Phases. ', 'One Destination.')
        self._subtitle(s, 'Phase 1 builds the foundation. Phase 2 activates AI and external clients. Phase 3 scales across jurisdictions.')

        # cw: 3.7*1.5=5.55, gap: 0.3*1.5=0.45, ch: 5.0*1.5=7.5
        cw = Inches(5.55)
        gap = Inches(0.45)
        ch = Inches(7.5)

        phases = [
            ('Phase 1', 'Year 1 \u2014 Foundation', self.BLUE,
             '~\u00a32.3\u20133.9M', '\u00a31.5M change', '~\u00a311M value', '7 AI use cases',
             ['Internal release first',
              'Prospect + Client Portal (Web + Mobile)',
              'Employee Portal + basic advisor tools',
              'Digital Wealth, Onboarding, Assist, CLO',
              'Agentic Runtime \u2014 data foundation',
              'Grand Central: T24, Salesforce, Folio',
              'Temenos e-Services continues (external)']),
            ('Phase 2', 'Year 2 \u2014 Expand', self.GREEN,
             '~\u00a32.6\u20135.1M', '\u00a31.0M change', '~\u00a331M cumul.', '+3 AI Agents',
             ['External client migration (go-live)',
              'Temenos e-Services deprecated',
              'RM Workspace Premium for advisors',
              'Full Intelligence Fabric activated',
              'Sentiment Analysis, Auto Reply, Doc Processing',
              'ComplyAdvantage + Onfido replace legacy',
              'Snowflake + ContentStack integrated']),
            ('Phase 3', 'Year 3+ \u2014 Scale', self.PURPLE,
             '~\u00a33.6\u20136.9M', '~\u00a30.6M change', '\u00a3148M cumul.', '7+ AI use cases',
             ['Full platform \u2014 all modules active',
              'Digital Advice, Sales (Lending, Advisory, Discretionary)',
              'Custom Agents, Smart Signals, Outreach',
              'Review Prep Copilot, Meeting Summarisation',
              'Multi-Agent Governance',
              'Aladdin Wealth Connector',
              'Geographic: UK \u2192 Guernsey \u2192 Switzerland \u2192 Singapore \u2192 US']),
        ]

        cx = self.ML
        for phase, subtitle, color, cost, change, value, ai, bullets in phases:
            # card top: 1.75*1.5=2.625
            self._card(s, cx, Inches(2.625), cw, ch, border=color)
            # Phase badge — +0.15*1.5=+0.225, top: 1.85*1.5=2.775, w: 0.8*1.5=1.2, h: 0.22*1.5=0.33
            badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       cx + Inches(0.225), Inches(2.775),
                                       Inches(1.2), Inches(0.33))
            badge.fill.solid()
            badge.fill.fore_color.rgb = color
            badge.line.fill.background()
            # badge text — top: 1.86*1.5=2.79, h: 0.2*1.5=0.3, size: 7*1.5=11
            self._txt(s, phase.upper(), cx + Inches(0.225), Inches(2.79),
                     Inches(1.2), Inches(0.3),
                     size=Pt(11), color=self.WHITE, bold=True, align=PP_ALIGN.CENTER)
            # Subtitle — +1.05*1.5=+1.575, top: 1.87*1.5=2.805, w: cw-1.2*1.5=cw-1.8, size: 9*1.5=14
            self._txt(s, subtitle, cx + Inches(1.575), Inches(2.805), cw - Inches(1.8), Pt(18),
                     size=Pt(14), color=color, bold=True)
            # Summary stats (2x2) — sy: 2.2*1.5=3.3, mini_w: 1.6*1.5=2.4
            sy = Inches(3.3)
            mini_w = Inches(2.4)
            for row in [(cost, change), (value, ai)]:
                # mx: cx+0.15*1.5=cx+0.225
                mx = cx + Inches(0.225)
                for val in row:
                    # size: 8*1.5=12
                    self._txt(s, val, mx, sy, mini_w, Pt(18),
                             size=Pt(12), color=self.DARK_TEXT, bold=True)
                    mx += mini_w
                sy += Pt(21)  # 14*1.5
            # Divider — +0.15*1.5=+0.225, top: 2.7*1.5=4.05, w: cw-0.3*1.5=cw-0.45
            div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     cx + Inches(0.225), Inches(4.05),
                                     cw - Inches(0.45), Pt(1.5))
            div.fill.solid()
            div.fill.fore_color.rgb = self.BORDER
            div.line.fill.background()
            # Bullets — start: 2.85*1.5=4.275, step: 0.28*1.5=0.42
            by = Inches(4.275)
            for b in bullets:
                # +0.15*1.5=+0.225, w: cw-0.3*1.5=cw-0.45, size: 8*1.5=12
                self._txt(s, '\u2192 ' + b, cx + Inches(0.225), by, cw - Inches(0.45), Pt(18),
                         size=Pt(12), color=self.SUB_TEXT)
                by += Inches(0.42)
            cx += cw + gap

    def _slide_07_investment(self):
        """Cost vs Value — static bar chart as a table representation."""
        s = self._new_slide()
        self._section_label(s, 'Where Value Meets Investment')
        self._heading(s, 'Year-on-Year: Cost, Value & ', 'AI Amplification.')
        self._subtitle(s, 'Enterprise model shown. Licence from sales pricing. Value from master analysis (validated).')

        # YoY data (Enterprise)
        yoy = [
            {'year': 'Year 1', 'license': 2.3, 'change': 1.5, 'base': 11.0, 'ai': 0},
            {'year': 'Year 2', 'license': 2.6, 'change': 1.0, 'base': 17.6, 'ai': 2.0},
            {'year': 'Year 3', 'license': 3.6, 'change': 0.6, 'base': 27.8, 'ai': 4.0},
            {'year': 'Year 4', 'license': 4.0, 'change': 0.6, 'base': 36.8, 'ai': 5.7},
            {'year': 'Year 5', 'license': 4.4, 'change': 0.6, 'base': 36.1, 'ai': 6.7},
        ]

        # Stacked bar chart using shapes
        # chart_left: 1.2*1.5=1.8, chart_top: 2.0*1.5=3.0, chart_w: 11.0*1.5=16.5, chart_h: 3.8*1.5=5.7
        chart_left = Inches(1.8)
        chart_top = Inches(3.0)
        chart_w = Inches(16.5)
        chart_h = Inches(5.7)
        max_val = 45
        bar_w = Inches(0.675)  # 0.45*1.5
        group_w = chart_w / 5
        colors = [self.RED, self.AMBER, self.GREEN, self.PURPLE]

        # Y-axis labels
        for val in [0, 10, 20, 30, 40]:
            y_pos = chart_top + chart_h - (val / max_val) * chart_h
            # left: 0.3*1.5=0.45, w: 0.8*1.5=1.2, size: 7*1.5=11
            self._txt(s, f'\u00a3{val}M', Inches(0.45), y_pos - Pt(8), Inches(1.2), Pt(18),
                     size=Pt(11), color=self.MUTED, align=PP_ALIGN.RIGHT)
            # Grid line
            line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      chart_left, y_pos, chart_w, Pt(0.75))  # 0.5*1.5
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            line.line.fill.background()

        for i, d in enumerate(yoy):
            # gap between bars: 0.05*1.5=0.075
            gx = chart_left + i * group_w + (group_w - 4 * bar_w - 3 * Inches(0.075)) / 2
            values = [d['license'], d['change'], d['base'], d['ai']]

            for j, (val, color) in enumerate(zip(values, colors)):
                if val == 0:
                    continue
                bar_x = gx + j * (bar_w + Inches(0.075))
                bar_h = (val / max_val) * chart_h
                bar_y = chart_top + chart_h - bar_h
                bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         bar_x, bar_y, bar_w, bar_h)
                bar.fill.solid()
                bar.fill.fore_color.rgb = color
                bar.line.fill.background()
                # Value label above bar — -0.05*1.5=-0.075, h: Pt(15)=10*1.5, w: bar_w+0.1*1.5=bar_w+0.15, size: 6*1.5=9
                self._txt(s, f'\u00a3{val:.1f}M', bar_x - Inches(0.075), bar_y - Pt(15),
                         bar_w + Inches(0.15), Pt(15),
                         size=Pt(9), color=color, bold=True, align=PP_ALIGN.CENTER)

            # X-axis label — w: 2*1.5=3, size: 9*1.5=14
            self._txt(s, d['year'], gx, chart_top + chart_h + Pt(6),
                     Inches(3.0), Pt(18),
                     size=Pt(14), color=self.DARK_TEXT, bold=True)

        # Legend — lx: 1.2*1.5=1.8, ly: 6.2*1.5=9.3
        lx = Inches(1.8)
        ly = Inches(9.3)
        for label, color in [('Licence', self.RED), ('Change', self.AMBER),
                              ('Base Value', self.GREEN), ('AI Value', self.PURPLE)]:
            # dot: Pt(12)=8*1.5
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, lx, ly + Pt(3), Pt(12), Pt(12))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            # label: +Pt(18)=12*1.5, w: 1*1.5=1.5, size: 8*1.5=12
            self._txt(s, label, lx + Pt(18), ly, Inches(1.5), Pt(18),
                     size=Pt(12), color=self.SUB_TEXT)
            lx += Inches(2.25)  # 1.5*1.5

        # Summary stats — left: 7.5*1.5=11.25, top: 6.1*1.5=9.15, w: 5.2*1.5=7.8, h: 0.5*1.5=0.75
        self._card(s, Inches(11.25), Inches(9.15), Inches(7.8), Inches(0.75),
                  fill=self.BLUE_LIGHT, border=self.BLUE)
        for i, (val, label) in enumerate([('\u00a3148M', 'CUMULATIVE VALUE'),
                                           ('\u00a321.2M', 'TOTAL INVESTMENT'),
                                           ('~7.0\u00d7', 'ROI')]):
            # left: 7.7*1.5=11.55, step: 1.7*1.5=2.55, w: 1.5*1.5=2.25
            # val top: 6.12*1.5=9.18, h: Pt(21)=14*1.5, size: 12*1.5=18
            self._txt(s, val, Inches(11.55) + i * Inches(2.55), Inches(9.18), Inches(2.25), Pt(21),
                     size=Pt(18), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)
            # label top: 6.35*1.5=9.525, h: Pt(12)=8*1.5, size: 6*1.5=9
            self._txt(s, label, Inches(11.55) + i * Inches(2.55), Inches(9.525), Inches(2.25), Pt(12),
                     size=Pt(9), color=self.SUB_TEXT, bold=True, align=PP_ALIGN.CENTER)

    def _slide_08_financials(self):
        """Five-Year Economics — summary table."""
        s = self._new_slide()
        self._section_label(s, 'Five-Year Economics')
        self._heading(s, 'The Numbers. ', 'Side by Side.')
        self._subtitle(s, 'Enterprise (recommended) vs Pay As You Grow. All figures in \u00a3M.')

        # Table — top: 1.9*1.5=2.85, w: CW, h: 3.0*1.5=4.5
        rows = 7  # header + 5 years + total
        cols = 8
        tbl_shape = s.shapes.add_table(rows, cols,
                                        self.ML, Inches(2.85),
                                        self.CW, Inches(4.5))
        tbl = tbl_shape.table

        # Column widths (all *1.5)
        col_widths = [Inches(1.2), Inches(1.8), Inches(2.55), Inches(2.55),
                      Inches(1.8), Inches(2.25), Inches(3.0), Inches(3.0)]
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w

        # Headers
        headers = ['Year', 'Phase', 'Ent. Licence', 'PAG Licence', 'Change',
                   'Ent. Total', 'Value (Base+AI)', 'Cumulative']
        for i, h in enumerate(headers):
            cell = tbl.cell(0, i)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)  # 8*1.5
                p.font.bold = True
                p.font.color.rgb = self.WHITE
                p.font.name = self.FONT
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.DARK_BG

        # Data
        data = [
            ['Y1', 'Foundation', '\u00a32.3M', '\u00a32.7M', '\u00a31.5M', '\u00a33.8M', '\u00a311.0M', '\u00a311M'],
            ['Y2', 'Expand', '\u00a32.6M', '\u00a33.5M', '\u00a31.0M', '\u00a33.6M', '\u00a319.6M', '\u00a331M'],
            ['Y3', 'Scale', '\u00a33.6M', '\u00a34.8M', '\u00a30.6M', '\u00a34.2M', '\u00a331.8M', '\u00a362M'],
            ['Y4', 'Optimise', '\u00a34.0M', '\u00a34.8M', '\u00a30.6M', '\u00a34.6M', '\u00a342.5M', '\u00a3105M'],
            ['Y5', 'Compound', '\u00a34.4M', '\u00a34.8M', '\u00a30.6M', '\u00a35.0M', '\u00a342.8M', '\u00a3148M'],
            ['Total', '', '\u00a316.9M', '\u00a320.6M', '\u00a34.3M', '\u00a321.2M', '\u00a3147.7M', '\u00a3148M'],
        ]
        for r, row_data in enumerate(data):
            for c, val in enumerate(row_data):
                cell = tbl.cell(r + 1, c)
                cell.text = val
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(12)  # 8*1.5
                    p.font.name = self.FONT
                    p.font.color.rgb = self.DARK_TEXT
                    if r == 5:  # Total row
                        p.font.bold = True
                    if c == 2:  # Ent licence
                        p.font.color.rgb = self.BLUE
                    elif c == 3:  # PAG licence
                        p.font.color.rgb = self.AMBER
                    elif c == 6 or c == 7:  # Value
                        p.font.color.rgb = self.GREEN
                if r == 5:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.BLUE_LIGHT

        # ROI comparison cards — top: 5.2*1.5=7.8, w: 5.8*1.5=8.7, h: 1.6*1.5=2.4
        cards_data = [
            ('Enterprise', self.BLUE, self.BLUE_LIGHT,
             [('~7.0\u00d7', 'ROI'), ('Month 8', 'Breakeven'), ('\u00a316.9M', '5yr Licence')]),
            ('Pay As You Grow', self.AMBER, self.AMBER_LIGHT,
             [('~5.9\u00d7', 'ROI'), ('Month 9', 'Breakeven'), ('\u00a320.6M', '5yr Licence')]),
        ]
        cx = self.ML
        for name, color, bg, stats in cards_data:
            self._card(s, cx, Inches(7.8), Inches(8.7), Inches(2.4), fill=bg, border=color)
            # label — +0.2*1.5=+0.3, top: 5.25*1.5=7.875, w: 3*1.5=4.5, size: 8*1.5=12
            self._txt(s, name.upper(), cx + Inches(0.3), Inches(7.875), Inches(4.5), Pt(15),
                     size=Pt(12), color=color, bold=True)
            sx = cx + Inches(0.3)
            for val, label in stats:
                # val top: 5.55*1.5=8.325, w: 1.6*1.5=2.4, h: Pt(27)=18*1.5, size: 16*1.5=24
                self._txt(s, val, sx, Inches(8.325), Inches(2.4), Pt(27),
                         size=Pt(24), color=color, bold=True, align=PP_ALIGN.CENTER)
                # label top: 5.9*1.5=8.85, h: Pt(15)=10*1.5, size: 7*1.5=11
                self._txt(s, label, sx, Inches(8.85), Inches(2.4), Pt(15),
                         size=Pt(11), color=self.SUB_TEXT, bold=True, align=PP_ALIGN.CENTER)
                sx += Inches(2.7)  # 1.8*1.5
            cx += Inches(9.15)  # 6.1*1.5

        # Insight — top: 6.95*1.5=10.425, h: 0.35*1.5=0.525
        self._card(s, self.ML, Inches(10.425), self.CW, Inches(0.525),
                  fill=self.BLUE_LIGHT, border=self.BLUE)
        # text — +0.3*1.5=+0.45, top: 6.97*1.5=10.455, w: CW-0.6*1.5=CW-0.9, h: 0.3*1.5=0.45, size: 8*1.5=12
        self._txt(s, 'Enterprise saves \u00a33.7M over 5 years vs PAG. Scaling discount (60%\u219245%). Every price increase maps to a specific module activation.',
                  self.ML + Inches(0.45), Inches(10.455), self.CW - Inches(0.9), Inches(0.45),
                  size=Pt(12), color=self.DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)

    def _slide_09_recommendation(self):
        """Recommendation — 4 expectations mapped."""
        s = self._new_slide()
        self._section_label(s, 'Recommendation')
        self._heading(s, 'Your Expectations. ', 'Our Answers.')
        self._subtitle(s, 'Every recommendation maps to a specific platform capability and proven reference.')

        recs = [
            ('Exit analysis, mobilise immediately.',
             'Enterprise model. Phase 1 Month 1. POC validates T24/Salesforce connectivity in 1 week.',
             'Ila Bank \u2014 6-month mobilisation'),
            ('No AI vendor lock-in.',
             'Multi-model architecture (Claude, GPT, Gemini). Agent Factory for custom agents. Switch LLMs without re-engineering.',
             'EQ Bank \u2014 multi-model AI'),
            ('Full integration, not point solutions.',
             'Grand Central: T24, Salesforce, Folio. 15+ connectors. Config not code. \u00a32M integration cost avoided.',
             'NatWest \u2014 T24 Grand Central'),
            ('Single pane of glass for advisors.',
             'RM Workspace + Intelligence Fabric. Review Prep Copilot saves 225 advisors 2.5hrs/review \u2014 \u00a32.5M/yr.',
             'RM Workspace \u2014 unified advisor experience'),
        ]

        # ty: 1.8*1.5=2.7
        ty = Inches(2.7)
        for i, (expectation, recommendation, reference) in enumerate(recs):
            # row_h: 1.15*1.5=1.725
            row_h = Inches(1.725)
            # Expectation card (left) — w: 4.0*1.5=6.0
            self._card(s, self.ML, ty, Inches(6.0), row_h, border=self.BORDER)
            # label — +0.15*1.5=+0.225, +0.08*1.5=+0.12, w: 3.5*1.5=5.25, size: 6*1.5=9
            self._txt(s, 'EXPECTATION', self.ML + Inches(0.225), ty + Inches(0.12),
                     Inches(5.25), Pt(12),
                     size=Pt(9), color=self.BLUE, bold=True)
            # text — +0.28*1.5=+0.42, w: 3.7*1.5=5.55, h: 0.7*1.5=1.05, size: 9*1.5=14
            self._txt(s, f'"{expectation}"', self.ML + Inches(0.225), ty + Inches(0.42),
                     Inches(5.55), Inches(1.05),
                     size=Pt(14), color=self.DARK_TEXT, bold=True)

            # Arrow — left: 4.8*1.5=7.2, +0.35*1.5=+0.525, w: 0.3*1.5=0.45, h: 0.3*1.5=0.45, size: 18*1.5=27
            self._txt(s, '\u2192', Inches(7.2), ty + Inches(0.525), Inches(0.45), Inches(0.45),
                     size=Pt(27), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)

            # Recommendation card (right) — left: 5.2*1.5=7.8, w: 7.2*1.5=10.8
            self._card(s, Inches(7.8), ty, Inches(10.8), row_h,
                      fill=self.BLUE_LIGHT, border=self.BLUE)
            # label — left: 5.35*1.5=8.025, +0.08*1.5=+0.12, w: 6.5*1.5=9.75, size: 6*1.5=9
            self._txt(s, 'RECOMMENDATION', Inches(8.025), ty + Inches(0.12),
                     Inches(9.75), Pt(12),
                     size=Pt(9), color=self.BLUE, bold=True)
            # text — +0.28*1.5=+0.42, w: 6.8*1.5=10.2, h: 0.45*1.5=0.675, size: 8*1.5=12
            self._txt(s, recommendation, Inches(8.025), ty + Inches(0.42),
                     Inches(10.2), Inches(0.675),
                     size=Pt(12), color=self.DARK_TEXT)
            # ref — +0.78*1.5=+1.17, w: 6.5*1.5=9.75, h: Pt(15)=10*1.5, size: 7*1.5=11
            self._txt(s, f'Ref: {reference}', Inches(8.025), ty + Inches(1.17),
                     Inches(9.75), Pt(15),
                     size=Pt(11), color=self.MUTED)

            ty += row_h + Inches(0.225)  # 0.15*1.5

        # Bottom card — top: 6.2*1.5=9.3, h: 0.9*1.5=1.35
        self._card(s, self.ML, Inches(9.3), self.CW, Inches(1.35),
                  fill=self.LIGHT_CARD, border=self.BLUE)
        # label — +0.15*1.5=+0.225, top: 6.28*1.5=9.42, w: 3*1.5=4.5, size: 6*1.5=9
        self._txt(s, 'COST PREDICTABILITY', self.ML + Inches(0.225), Inches(9.42),
                 Inches(4.5), Pt(12), size=Pt(9), color=self.BLUE, bold=True)
        # text — top: 6.48*1.5=9.72, w: CW-0.3*1.5=CW-0.45, h: 0.5*1.5=0.75, size: 9*1.5=14
        self._txt(s, 'Scaling discount (60%\u219245% of list). Capped BIC for AI. Every price increase maps to a specific module. \u00a316.9M over 5 years \u2014 predictable, auditable.',
                  self.ML + Inches(0.225), Inches(9.72), self.CW - Inches(0.45), Inches(0.75),
                  size=Pt(14), color=self.SUB_TEXT)

    def _slide_10_close(self):
        """Close Plan + CTA — dark hero."""
        s = self._new_slide(dark=True)
        # top: 0.4*1.5=0.6
        self._section_label(s, 'Mutual Close Plan', top=Inches(0.6), dark=True)
        # top: 0.75*1.5=1.125
        self._heading(s, 'Target: Signed by ', 'End of April 2026.',
                     top=Inches(1.125), dark=True)

        # Timeline milestones
        milestones = [
            ('March 9', 'COO Meeting', 'Platform Value & AI Use Cases', self.BLUE, True),
            ('Mid-March', 'Scope & Confidence', 'Product Catalogue refinement', self.WHITE, False),
            ('Late March', 'Exclusivity', 'Formal procurement starts', self.PURPLE, False),
            ('Early April', '1-Week POC', 'T24/Salesforce connectivity proof', self.GREEN, False),
            ('Mid-April', 'Procurement', 'Due Diligence \u00b7 Infosec \u00b7 MSA Review', self.WHITE, False),
            ('End of April', 'Signature', 'Contract signed \u00b7 Implementation begins', self.GREEN, False),
        ]

        mx = self.ML
        # mw: 1.8*1.5=2.7, my: 1.5*1.5=2.25
        mw = Inches(2.7)
        my = Inches(2.25)

        for date, title, sub, color, is_current in milestones:
            # h: 1.3*1.5=1.95
            self._card(s, mx, my, mw, Inches(1.95),
                      fill=self.CARD_BG, border=color, dark=True)
            # Dot indicator
            if is_current:
                # Pt(6)=4*1.5, Pt(12)=8*1.5
                dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                          mx + mw/2 - Pt(6), my - Pt(12),
                                          Pt(12), Pt(12))
                dot.fill.solid()
                dot.fill.fore_color.rgb = self.BLUE
                dot.line.fill.background()

            # date — +0.1*1.5=+0.15, +0.08*1.5=+0.12, w: mw-0.2*1.5=mw-0.3, size: 8*1.5=12
            self._txt(s, date, mx + Inches(0.15), my + Inches(0.12), mw - Inches(0.3), Pt(15),
                     size=Pt(12), color=color, bold=True, align=PP_ALIGN.CENTER)
            # title — +0.35*1.5=+0.525, size: 10*1.5=15
            self._txt(s, title, mx + Inches(0.15), my + Inches(0.525), mw - Inches(0.3), Pt(21),
                     size=Pt(15), color=self.WHITE, bold=True, align=PP_ALIGN.CENTER)
            # sub — +0.65*1.5=+0.975, h: 0.5*1.5=0.75, size: 7*1.5=11
            self._txt(s, sub, mx + Inches(0.15), my + Inches(0.975), mw - Inches(0.3), Inches(0.75),
                     size=Pt(11), color=self.MUTED, align=PP_ALIGN.CENTER)
            mx += mw + Inches(0.225)  # 0.15*1.5

        # Connecting line — top: my-Pt(5)*1.5=my-Pt(7.5), w: 11.0*1.5=16.5, h: Pt(3)=2*1.5
        line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   self.ML + mw/2, my - Pt(7.5),
                                   Inches(16.5), Pt(3))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x5E)
        line.line.fill.background()

        # Hero CTA — top: 3.3*1.5=4.95, h: 1.5*1.5=2.25
        tb = s.shapes.add_textbox(self.ML, Inches(4.95), self.CW, Inches(2.25))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for text, color in [('Build the Foundation.\n', self.WHITE),
                            ('Accelerate the Advantage.', self.BLUE)]:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(54)   # 36*1.5
            r.font.bold = True
            r.font.color.rgb = color
            r.font.name = self.FONT

        # 4 hero stats — stat_w: 2.6*1.5=3.9, gap: 0.3*1.5=0.45, sx start: ML+0.3*1.5=ML+0.45
        stat_w = Inches(3.9)
        gap = Inches(0.45)
        sx = self.ML + Inches(0.45)
        for val, label, vcol in [
            ('~\u00a32.3M', 'Y1 Entry (Enterprise)', self.BLUE),
            ('\u00a316.9M', '5-Year Licence', self.WHITE),
            ('\u00a3148M', 'Value at Stake', self.GREEN),
            ('~7.0\u00d7', 'ROI', self.BLUE),
        ]:
            # top: 5.0*1.5=7.5, h: 1.0*1.5=1.5
            self._card(s, sx, Inches(7.5), stat_w, Inches(1.5),
                      fill=self.CARD_BG, dark=True)
            # val — top: 5.05*1.5=7.575, h: 0.45*1.5=0.675, size: 22*1.5=33
            self._txt(s, val, sx, Inches(7.575), stat_w, Inches(0.675),
                     size=Pt(33), color=vcol, bold=True, align=PP_ALIGN.CENTER)
            # label — top: 5.5*1.5=8.25, h: Pt(18)=12*1.5, size: 7*1.5=11
            self._txt(s, label.upper(), sx, Inches(8.25), stat_w, Pt(18),
                     size=Pt(11), color=self.MUTED, bold=True, align=PP_ALIGN.CENTER)
            sx += stat_w + gap

        # Footer — top: 6.3*1.5=9.45, h: 0.3*1.5=0.45, size: 11*1.5=17
        self._txt(s, 'One platform. Every channel. Every data source. From relationship depth to digital scale.',
                  self.ML, Inches(9.45), self.CW, Inches(0.45),
                  size=Pt(17), color=self.MUTED, align=PP_ALIGN.CENTER)
        # bottom line — top: 6.7*1.5=10.05, h: Pt(18)=12*1.5, size: 8*1.5=12
        self._txt(s, 'Schroders \u00d7 Backbase \u00b7 March 2026',
                  self.ML, Inches(10.05), self.CW, Pt(18),
                  size=Pt(12), color=RGBColor(0x3A, 0x49, 0x5D), align=PP_ALIGN.CENTER)

    # ── GENERATE ──────────────────────────────────────────

    def generate(self, output_path):
        # Load the Backbase Master Template
        self.prs = Presentation(TEMPLATE_PATH)

        # Delete all existing slides from the template
        xml_slides = self.prs.slides._sldIdLst
        for sldId in list(xml_slides):
            rId = sldId.get('r:id')
            if rId is None:
                # Try the full namespace
                rId = sldId.attrib.get(
                    '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
                )
            if rId:
                self.prs.part.drop_rel(rId)
            xml_slides.remove(sldId)

        # Find a blank layout (one with fewest placeholders)
        self.blank_layout = None
        for layout in self.prs.slide_layouts:
            if len(layout.placeholders) == 0:
                self.blank_layout = layout
                break
        if not self.blank_layout:
            self.blank_layout = min(self.prs.slide_layouts,
                                     key=lambda l: len(l.placeholders))

        # Build all 10 slides
        self._slide_01_cover()          # 1  — Dark hero
        self._slide_02_context()        # 2  — Context
        self._slide_03_platform()       # 3  — What You Get: Platform
        self._slide_04_ai()             # 4  — What You Get: AI
        self._slide_05_commercial()     # 5  — Enterprise vs PAG
        self._slide_06_phases()         # 6  — Phased Delivery
        self._slide_07_investment()     # 7  — Cost vs Value (chart)
        self._slide_08_financials()     # 8  — Five-Year Economics
        self._slide_09_recommendation() # 9  — Recommendation
        self._slide_10_close()          # 10 — Close Plan + CTA

        self.prs.save(output_path)
        size_kb = Path(output_path).stat().st_size // 1024
        print(f'Saved: {output_path} ({size_kb} KB, 10 slides)')


if __name__ == '__main__':
    OUTPUT_DIR = Path(__file__).parent.parent / 'Engagement/Schroders Group/Output'
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    output_file = OUTPUT_DIR / 'Schroders_Executive_Readout_v3_Master.pptx'
    SchrordersV3MasterPptxGenerator().generate(str(output_file))
