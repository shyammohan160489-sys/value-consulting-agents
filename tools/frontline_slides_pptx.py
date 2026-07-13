#!/usr/bin/env python3
"""
BackbaseSlidesPresenter — PPTX builder for the Frontline 2026 Slide Engine.

Renders the same 17 layout types as engine.js to Google Slides-compatible PPTX.
Uses the same data schema (layout name + properties) so HTML and PPTX skills
produce visually consistent output from identical input.

Tokens align to Master Template theme1.xml — see knowledge/design-system/frontline-tokens.json.

Usage:
    from tools.frontline_slides_pptx import BackbaseSlidesPresenter

    deck = BackbaseSlidesPresenter('Deck Title')
    deck.add_cover_color_block('BACKBASE', 'Title\\nLine 2', 'April 2026')
    deck.add_content_standard('light', 'TOPIC', 'Title', 'Subtitle', '<ul><li>Point</li></ul>')
    deck.add_thank_you()
    deck.save('output.pptx')
"""

import re
import io
from pathlib import Path
from html.parser import HTMLParser
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE


class _HTMLToText(HTMLParser):
    """Minimal HTML-to-text converter for body fields."""
    def __init__(self):
        super().__init__()
        self._lines = []
        self._current = ''
        self._in_li = False

    def handle_starttag(self, tag, attrs):
        if tag == 'li':
            self._in_li = True
            self._current = '\u2192 '
        elif tag == 'br':
            self._lines.append(self._current)
            self._current = ''
        elif tag in ('p', 'div'):
            if self._current:
                self._lines.append(self._current)
                self._current = ''

    def handle_endtag(self, tag):
        if tag == 'li':
            self._lines.append(self._current)
            self._current = ''
            self._in_li = False
        elif tag in ('p', 'div', 'ul', 'ol'):
            if self._current:
                self._lines.append(self._current)
                self._current = ''

    def handle_data(self, data):
        self._current += data.strip()

    def get_text(self):
        if self._current:
            self._lines.append(self._current)
        return '\n'.join(l for l in self._lines if l)


def html_to_text(html_str):
    """Convert simple HTML body content to plain text with arrow bullets."""
    if not html_str:
        return ''
    parser = _HTMLToText()
    parser.feed(html_str)
    return parser.get_text()


class _HTMLBlocks(HTMLParser):
    """Split body HTML into ordered (kind, inner_html) blocks where kind is
    'header' (<div>/<p>) or 'bullet' (<li>). Inner HTML is preserved so inline
    <b>/<span> emphasis can be re-parsed downstream. Tables are skipped."""
    def __init__(self):
        super().__init__(convert_charrefs=False)
        self.blocks = []
        self._buf = ''
        self._kind = None
        self._depth_skip = 0

    def handle_starttag(self, tag, attrs):
        if tag in ('table',):
            self._depth_skip += 1
            return
        if self._depth_skip:
            return
        if tag == 'li':
            self._flush(); self._kind = 'bullet'; self._buf = ''
        elif tag in ('div', 'p'):
            self._flush(); self._kind = 'header'; self._buf = ''
        elif tag in ('b', 'strong', 'span', 'br'):
            self._buf += self.get_starttag_text() or ''

    def handle_endtag(self, tag):
        if tag in ('table',):
            if self._depth_skip:
                self._depth_skip -= 1
            return
        if self._depth_skip:
            return
        if tag in ('li', 'div', 'p'):
            self._flush()
        elif tag in ('b', 'strong', 'span'):
            self._buf += f'</{tag}>'

    def handle_data(self, data):
        if self._depth_skip:
            return
        self._buf += data

    def handle_entityref(self, name):
        if not self._depth_skip:
            self._buf += f'&{name};'

    def handle_charref(self, name):
        if not self._depth_skip:
            self._buf += f'&#{name};'

    def _flush(self):
        if self._kind and self._buf.strip():
            self.blocks.append((self._kind, self._buf.strip()))
        self._buf = ''
        self._kind = None

    def get_blocks(self):
        self._flush()
        return self.blocks


def html_to_blocks(html_str):
    """Parse body HTML into [(kind, inner_html), ...] for hierarchical rendering."""
    if not html_str:
        return []
    p = _HTMLBlocks()
    p.feed(str(html_str))
    return p.get_blocks()


def strip_html_tags(text):
    """Remove HTML tags, keeping text. Converts <span class="hl">X</span> to X."""
    if not text:
        return ''
    return re.sub(r'<[^>]+>', '', text)


def _hex_to_rgb(hex_str):
    """'#3367FF' / '#fff' -> RGBColor."""
    h = hex_str.lstrip('#')
    if len(h) == 3:
        h = ''.join(c * 2 for c in h)
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class _RichRuns(HTMLParser):
    """Parse a short HTML snippet into styled runs: (text, bold, color|None).

    Honours <b>/<strong>, <span class="hl"> (mapped to hl_color), and
    <span style="color:#xxx">. Entities are converted automatically
    (convert_charrefs defaults to True)."""
    def __init__(self, hl_color=None):
        super().__init__()
        self.runs = []
        self.bold = 0
        self.color_stack = []
        self.hl_color = hl_color

    def handle_starttag(self, tag, attrs):
        ad = dict(attrs)
        if tag in ('b', 'strong'):
            self.bold += 1
        elif tag == 'span':
            col = None
            if 'hl' in ad.get('class', '').split():
                col = self.hl_color
            m = re.search(r'color:\s*(#[0-9A-Fa-f]{3,6})', ad.get('style', ''))
            if m:
                col = _hex_to_rgb(m.group(1))
            self.color_stack.append(col)
        elif tag == 'br':
            self.runs.append(('\n', False, None))

    def handle_endtag(self, tag):
        if tag in ('b', 'strong') and self.bold > 0:
            self.bold -= 1
        elif tag == 'span' and self.color_stack:
            self.color_stack.pop()

    def handle_data(self, data):
        text = re.sub(r'\s+', ' ', data)
        if not text:
            return
        col = next((c for c in reversed(self.color_stack) if c is not None), None)
        self.runs.append((text, self.bold > 0, col))


def parse_rich_runs(html_str, hl_color=None):
    """Convert an inline HTML string to a list of (text, bold, color) runs."""
    if html_str is None:
        return []
    p = _RichRuns(hl_color=hl_color)
    p.feed(str(html_str))
    runs = p.runs
    if runs:  # trim outer whitespace
        t0, b0, c0 = runs[0]
        runs[0] = (t0.lstrip(), b0, c0)
        t1, b1, c1 = runs[-1]
        runs[-1] = (t1.rstrip(), b1, c1)
        runs = [(t, b, c) for (t, b, c) in runs if t]
    return runs


class BackbaseSlidesPresenter:
    """PPTX builder matching the Backbase Slide Template Engine's 17 layouts."""

    # ── Dimensions (Google Slides widescreen) ─────────────
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    # ── Colors (matching engine.js design tokens) ─────────
    NAVY       = RGBColor(0x04, 0x13, 0x26)
    NAVY_DARK  = RGBColor(0x04, 0x13, 0x26)
    BLUE       = RGBColor(0x33, 0x67, 0xFF)
    CYAN       = RGBColor(0x69, 0xFE, 0xFF)
    RED        = RGBColor(0xFF, 0x50, 0x3C)
    LIGHT_BLUE = RGBColor(0xE5, 0xEB, 0xFF)
    OFF_WHITE  = RGBColor(0xF3, 0xF6, 0xF9)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK      = RGBColor(0x00, 0x00, 0x00)
    BORDER     = RGBColor(0xCE, 0xD2, 0xD7)
    BLUE_DARK  = RGBColor(0x26, 0x4E, 0xC7)
    GREEN      = RGBColor(0x2E, 0xCC, 0x71)

    # Text colors
    TEXT_DARK    = RGBColor(0x04, 0x13, 0x26)
    TEXT_MUTED   = RGBColor(0x64, 0x74, 0x8B)
    TEXT_MUTED_W = RGBColor(0x99, 0x99, 0xAA)  # white bg muted
    TEXT_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
    LABEL_DARK   = RGBColor(0x04, 0x13, 0x26)
    LABEL_WHITE  = RGBColor(0xCC, 0xCC, 0xDD)

    # Grid line colors (very subtle — matches HTML rgba opacity)
    GRID_LIGHT = RGBColor(0xE0, 0xE3, 0xE8)  # light slides: barely visible
    GRID_DARK  = RGBColor(0x1A, 0x2D, 0x45)  # dark slides: very subtle against navy

    # Bar colors for roadmap
    BAR_COLORS = {
        'blue': RGBColor(0x33, 0x67, 0xFF),
        'navy': RGBColor(0x04, 0x13, 0x26),
        'cyan': RGBColor(0x2B, 0xBC, 0xC4),
        'red':  RGBColor(0xFF, 0x50, 0x3C),
    }

    # Statement band backgrounds
    BAND_BLUE = RGBColor(0xE5, 0xEB, 0xFF)
    BAND_RED  = RGBColor(0xFF, 0xEF, 0xEC)

    # ── Typography ────────────────────────────────────────
    FONT = 'Libre Franklin'

    # ── Layout constants (percentages → inches at 13.333 × 7.5) ──
    # Content area margins
    ML = Inches(0.77)   # ~5.8% of 13.333
    MR = Inches(0.39)   # ~2.9% of 13.333
    CW = Inches(11.73)  # ~88% content width

    def __init__(self, title='Backbase Presentation'):
        self.title = title
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_W
        self.prs.slide_height = self.SLIDE_H
        self.blank_layout = self.prs.slide_layouts[6]
        self._slide_num = 0

    # ══════════════════════════════════════════════════════
    #  INTERNAL HELPERS
    # ══════════════════════════════════════════════════════

    def _new_slide(self, bg_color=None):
        """Add a blank slide with specified background color."""
        self._slide_num += 1
        slide = self.prs.slides.add_slide(self.blank_layout)
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = bg_color or self.WHITE
        return slide

    def _txt(self, slide, text, left, top, width, height, size=Pt(12),
             color=None, bold=False, align=PP_ALIGN.LEFT, name=None):
        """Add a textbox with auto_size disabled (Google Slides rule)."""
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color or self.TEXT_DARK
        p.font.name = name or self.FONT
        p.alignment = align
        tb.name = f"txt_{self._slide_num}_{left}"
        return tf

    def _multi_line_txt(self, slide, lines, left, top, width, height,
                        size=Pt(12), color=None, bold=False, align=PP_ALIGN.LEFT,
                        spacing=Pt(4)):
        """Multi-line textbox from a list of strings."""
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(line)
            p.font.size = size
            p.font.bold = bold
            p.font.color.rgb = color or self.TEXT_DARK
            p.font.name = self.FONT
            p.alignment = align
            p.space_after = spacing
        return tf

    def _label(self, slide, text, left, top, dark=False):
        """Uppercase section label."""
        color = self.LABEL_WHITE if dark else self.LABEL_DARK
        self._txt(slide, text.upper(), left, top,
                  Inches(4), Pt(20), size=Pt(9),
                  color=color, bold=False)

    def _title(self, slide, text, left, top, width=None, size=Pt(28),
               dark=False):
        """Main title text, supports \\n line breaks."""
        color = self.TEXT_WHITE if dark else self.TEXT_DARK
        lines = text.split('\n') if text else ['']
        w = width or self.CW
        self._multi_line_txt(slide, lines, left, top, w, Inches(1.2),
                             size=size, color=color, bold=False)

    def _subtitle_text(self, slide, text, left, top, width=None, dark=False):
        """Subtitle text."""
        color = self.TEXT_MUTED_W if dark else self.BLUE
        w = width or Inches(9.3)
        self._txt(slide, text, left, top, w, Inches(0.6),
                  size=Pt(14), color=color)

    def _body_text(self, slide, html_body, left, top, width=None, dark=False):
        """Body text from HTML content, with visual hierarchy: <div>/<p> render
        as bold group headers, <li> as indented arrow bullets. Inline <b> and
        colored <span> emphasis is preserved."""
        blocks = html_to_blocks(html_body)
        if not blocks:
            return
        color = self.TEXT_WHITE if dark else self.TEXT_DARK
        w = width or self.CW
        tb = slide.shapes.add_textbox(left, top, w, Inches(3.8))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        first = True
        for kind, raw in blocks:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.line_spacing = 1.2
            if kind == 'header':
                p.space_before = Pt(8)
                p.space_after = Pt(3)
                runs = parse_rich_runs(raw)
                self._rich_paragraph(p, runs, Pt(12), color, bold=True)
            else:  # bullet
                p.space_after = Pt(4)
                r0 = p.add_run(); r0.text = '→  '
                r0.font.size = Pt(11); r0.font.name = self.FONT
                r0.font.color.rgb = self.BLUE
                self._rich_paragraph(p, parse_rich_runs(raw), Pt(11), color)

    def _footer(self, slide, dark=False):
        """Backbase wordmark + slide number at bottom-right."""
        text_color = self.TEXT_MUTED_W if dark else self.TEXT_MUTED
        # "Backbase  |  N"
        tb = slide.shapes.add_textbox(
            self.SLIDE_W - Inches(2.0), self.SLIDE_H - Inches(0.4),
            Inches(1.7), Inches(0.3)
        )
        tf = tb.text_frame
        tf.word_wrap = False
        tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r1 = p.add_run()
        r1.text = 'Backbase'
        r1.font.size = Pt(9)
        r1.font.color.rgb = text_color
        r1.font.name = self.FONT
        r2 = p.add_run()
        r2.text = '   |   '
        r2.font.size = Pt(9)
        r2.font.color.rgb = self.GRID_LIGHT if not dark else self.GRID_DARK
        r2.font.name = self.FONT
        r3 = p.add_run()
        r3.text = str(self._slide_num)
        r3.font.size = Pt(9)
        r3.font.color.rgb = text_color
        r3.font.name = self.FONT
        tb.name = f"footer_{self._slide_num}"

    def _grid_line_v(self, slide, x_pct, dark=False, top_pct=0, bottom_pct=100):
        """Vertical grid line at x% of slide width. Hairline (1px equivalent)."""
        color = self.GRID_DARK if dark else self.GRID_LIGHT
        left = Emu(int(self.SLIDE_W.emu * x_pct / 100))
        top = Emu(int(self.SLIDE_H.emu * top_pct / 100))
        h = Emu(int(self.SLIDE_H.emu * (bottom_pct - top_pct) / 100))
        # Use a line shape for true hairline width
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Emu(6350), h)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        s.name = f"grid_v_{x_pct}"

    def _grid_line_h(self, slide, y_pct, dark=False, left_pct=0, right_pct=100):
        """Horizontal grid line at y% of slide height. Hairline (1px equivalent)."""
        color = self.GRID_DARK if dark else self.GRID_LIGHT
        left = Emu(int(self.SLIDE_W.emu * left_pct / 100))
        top = Emu(int(self.SLIDE_H.emu * y_pct / 100))
        w = Emu(int(self.SLIDE_W.emu * (right_pct - left_pct) / 100))
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Emu(6350))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        s.name = f"grid_h_{y_pct}"

    def _motif(self, slide, left, top, color=None):
        """Small Backbase motif shape — tiny accent square at grid intersections."""
        c = color or self.BLUE
        size = Inches(0.10)  # very small — just a subtle accent mark
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    left - size, top - size, size, size)
        s.fill.solid()
        s.fill.fore_color.rgb = c
        s.line.fill.background()
        s.name = f"motif_{self._slide_num}"

    # Exact clip-path polygon from deck-template.html (statement band).
    _BAND_POLY = [(32.7, 25.13), (97.1, 25.13), (97.1, 65.9), (58, 65.9),
                  (58, 73), (2.9, 73), (2.9, 32.23), (32.7, 32.23)]

    def _band_shape(self, slide, color):
        """Stepped statement band — a single freeform polygon that matches the
        HTML clip-path exactly (replaces the old two-rectangle approximation,
        which produced a visible protruding step artifact)."""
        pts = [(int(self.SLIDE_W.emu * px / 100), int(self.SLIDE_H.emu * py / 100))
               for px, py in self._BAND_POLY]
        fb = slide.shapes.build_freeform(pts[0][0], pts[0][1], scale=1.0)
        fb.add_line_segments(pts[1:], close=True)
        shape = fb.convert_to_shape()
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        shape.shadow.inherit = False
        shape.name = f"band_{self._slide_num}"

    # ══════════════════════════════════════════════════════
    #  17 LAYOUT METHODS
    # ══════════════════════════════════════════════════════

    def add_cover_color_block(self, label='BACKBASE', title='Title',
                               date='', partner=False):
        """Layout 1: Navy background with blue accent. Opening slide."""
        slide = self._new_slide(self.NAVY)
        # Subtle blue accent oval in bottom-right (simulates CSS radial-gradient)
        accent = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(int(self.SLIDE_W.emu * 0.55)),
            Emu(int(self.SLIDE_H.emu * 0.60)),
            Emu(int(self.SLIDE_W.emu * 0.50)),
            Emu(int(self.SLIDE_H.emu * 0.50))
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(0x0D, 0x21, 0x3F)  # very subtle blue tint
        accent.line.fill.background()
        accent.name = "cover_glow"
        # Grid lines
        self._grid_line_v(slide, 4.13, dark=True)
        self._grid_line_v(slide, 54.1, dark=True)
        self._grid_line_v(slide, 95.86, dark=True)
        self._grid_line_h(slide, 21.54, dark=True)
        self._grid_line_h(slide, 64.86, dark=True)
        self._grid_line_h(slide, 92.65, dark=True)
        # Motif
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.541)),
                    Emu(int(self.SLIDE_H.emu * 0.2154)), self.WHITE)
        # Logo text (bold, white, smaller — wordmark style)
        self._txt(slide, 'Backbase', Inches(1.07), Inches(0.53),
                  Inches(2), Inches(0.4), size=Pt(13),
                  color=self.TEXT_WHITE, bold=True)
        # Label
        if label:
            self._label(slide, label, Inches(1.07), Inches(1.91), dark=True)
        # Title (large, light weight feel)
        self._title(slide, title, Inches(1.07), Inches(2.4),
                    width=Inches(6.4), size=Pt(48), dark=True)
        # Date
        if date:
            self._txt(slide, date, Inches(1.07), Inches(5.25),
                      Inches(4), Inches(0.4), size=Pt(14),
                      color=self.TEXT_MUTED_W)

    def add_cover_photo(self, label='', title='', date='',
                        image_path=None, partner=None):
        """Layout 2: Cover with optional photo or navy variant."""
        if image_path:
            # Photo variant — white panel overlay
            slide = self._new_slide(self.NAVY)
            try:
                slide.shapes.add_picture(image_path, 0, 0,
                                          self.SLIDE_W, self.SLIDE_H)
            except Exception:
                pass
            # White panel (left half)
            left = Emu(int(self.SLIDE_W.emu * 0.0413))
            top = Emu(int(self.SLIDE_H.emu * 0.2154))
            w = Emu(int(self.SLIDE_W.emu * 0.50))
            h = Emu(int(self.SLIDE_H.emu * 0.72))
            panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
            panel.fill.solid()
            panel.fill.fore_color.rgb = self.WHITE
            panel.fill.fore_color.brightness = -0.0
            panel.line.fill.background()
            panel.name = "cover_panel"
            # Content on panel
            if label:
                self._label(slide, label, Inches(1.07), Inches(1.91))
            self._title(slide, title, Inches(1.07), Inches(2.4),
                        width=Inches(6.4), size=Pt(42))
            if date:
                self._txt(slide, date, Inches(1.07), Inches(5.25),
                          Inches(4), Inches(0.4), size=Pt(14),
                          color=self.TEXT_MUTED)
        else:
            # Navy variant (no photo)
            slide = self._new_slide(self.NAVY_DARK)
            self._grid_line_v(slide, 4.13, dark=True)
            self._grid_line_v(slide, 54.1, dark=True)
            self._grid_line_v(slide, 95.86, dark=True)
            self._grid_line_h(slide, 21.54, dark=True)
            self._grid_line_h(slide, 64.86, dark=True)
            # Logo
            self._txt(slide, 'Backbase', Inches(1.07), Inches(0.53),
                      Inches(2), Inches(0.4), size=Pt(14),
                      color=self.TEXT_WHITE, bold=True)
            if label:
                self._label(slide, label, Inches(1.07), Inches(1.91), dark=True)
            self._title(slide, title, Inches(1.07), Inches(2.4),
                        width=Inches(6.4), size=Pt(42), dark=True)
            if date:
                self._txt(slide, date, Inches(1.07), Inches(5.25),
                          Inches(4), Inches(0.4), size=Pt(14),
                          color=self.TEXT_MUTED_W)

    def add_chapter_numbered(self, theme='navy', number='01', label='',
                              title='', subtitle=''):
        """Layout 3: Large number left, title right. Section divider."""
        dark = True
        bg = self.NAVY_DARK if theme == 'navy' else self.BLUE
        slide = self._new_slide(bg)
        # Subtle blue accent for navy theme
        if theme == 'navy':
            accent = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Emu(int(self.SLIDE_W.emu * 0.55)),
                Emu(int(self.SLIDE_H.emu * 0.60)),
                Emu(int(self.SLIDE_W.emu * 0.50)),
                Emu(int(self.SLIDE_H.emu * 0.50))
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor(0x08, 0x1A, 0x30)
            accent.line.fill.background()
            accent.name = "chapter_glow"
        # Grid
        self._grid_line_v(slide, 4.13, dark=True)
        self._grid_line_v(slide, 27, dark=True, top_pct=21.54, bottom_pct=64.86)
        self._grid_line_v(slide, 95.86, dark=True)
        self._grid_line_h(slide, 21.54, dark=True)
        self._grid_line_h(slide, 64.86, dark=True)
        # Motif
        mc = self.WHITE if theme == 'blue' else self.BLUE
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.0413)),
                    Emu(int(self.SLIDE_H.emu * 0.2154)), mc)
        # Number
        num_color = self.WHITE if theme == 'blue' else self.BLUE
        self._txt(slide, number, Emu(int(self.SLIDE_W.emu * 0.0413)),
                  Emu(int(self.SLIDE_H.emu * 0.115)),
                  Emu(int(self.SLIDE_W.emu * 0.2287)),
                  Emu(int(self.SLIDE_H.emu * 0.5336)),
                  size=Pt(84), color=num_color, align=PP_ALIGN.CENTER)
        # Label
        if label:
            self._label(slide, label, Emu(int(self.SLIDE_W.emu * 0.29)),
                        Emu(int(self.SLIDE_H.emu * 0.25)), dark=True)
        # Title
        self._title(slide, title, Emu(int(self.SLIDE_W.emu * 0.29)),
                    Emu(int(self.SLIDE_H.emu * 0.32)),
                    width=Emu(int(self.SLIDE_W.emu * 0.62)),
                    size=Pt(34), dark=True)
        # Subtitle
        if subtitle:
            self._txt(slide, subtitle, Emu(int(self.SLIDE_W.emu * 0.29)),
                      Emu(int(self.SLIDE_H.emu * 0.55)),
                      Emu(int(self.SLIDE_W.emu * 0.60)),
                      Inches(1), size=Pt(14), color=self.TEXT_MUTED_W)

    def add_chapter_standard(self, theme='navy', label='', title='',
                              subtitle=''):
        """Layout 4: Full-width title section divider."""
        bg = self.NAVY_DARK if theme == 'navy' else self.BLUE
        slide = self._new_slide(bg)
        self._grid_line_v(slide, 4.13, dark=True)
        self._grid_line_v(slide, 95.86, dark=True)
        self._grid_line_h(slide, 21.54, dark=True)
        self._grid_line_h(slide, 64.86, dark=True)
        mc = self.WHITE if theme == 'blue' else self.BLUE
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.0413)),
                    Emu(int(self.SLIDE_H.emu * 0.2154)), mc)
        if label:
            self._label(slide, label, self.ML, Emu(int(self.SLIDE_H.emu * 0.25)),
                        dark=True)
        self._title(slide, title, self.ML, Emu(int(self.SLIDE_H.emu * 0.32)),
                    width=Emu(int(self.SLIDE_W.emu * 0.85)), size=Pt(34), dark=True)
        if subtitle:
            self._txt(slide, subtitle, self.ML, Emu(int(self.SLIDE_H.emu * 0.55)),
                      Emu(int(self.SLIDE_W.emu * 0.60)), Inches(1),
                      size=Pt(14), color=self.TEXT_MUTED_W)

    def add_toc(self, label='AGENDA', title='Contents', items=None,
                numbered=True):
        """Layout 5: Table of contents with numbered/plain rows."""
        items = items or []
        slide = self._new_slide(self.WHITE)
        self._grid_line_v(slide, 2.9)
        self._grid_line_v(slide, 40)
        self._grid_line_v(slide, 97.1)
        self._grid_line_h(slide, 5.15)
        self._grid_line_h(slide, 94.85)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        if label:
            self._label(slide, label, self.ML, Emu(int(self.SLIDE_H.emu * 0.10)))
        if title:
            self._title(slide, title, self.ML, Emu(int(self.SLIDE_H.emu * 0.15)),
                        width=Inches(4))
        # TOC rows on the right side
        row_left = Emu(int(self.SLIDE_W.emu * 0.40))
        row_width = Emu(int(self.SLIDE_W.emu * 0.57))
        row_top_start = Emu(int(self.SLIDE_H.emu * 0.0515))
        row_area_h = Emu(int(self.SLIDE_H.emu * 0.897))
        n = len(items)
        if n == 0:
            self._footer(slide)
            return
        row_h = row_area_h // n
        for idx, item in enumerate(items):
            y = row_top_start + row_h * idx
            # Divider line between rows
            if idx > 0:
                s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            row_left, y, row_width, Pt(0.5))
                s.fill.solid()
                s.fill.fore_color.rgb = self.GRID_LIGHT
                s.line.fill.background()
            prefix = f"{str(idx + 1).zfill(2)}    " if numbered else ''
            self._txt(slide, f"{prefix}{item}", row_left + Inches(0.2),
                      y + Emu(row_h // 3),
                      row_width - Inches(0.4), Emu(row_h // 2),
                      size=Pt(13), color=self.TEXT_DARK)
        self._footer(slide)

    def add_content_standard(self, theme='light', label='', title='',
                              subtitle='', body=''):
        """Layout 6: General-purpose content slide. Most versatile."""
        dark = theme == 'dark'
        bg = self.NAVY if dark else self.WHITE
        slide = self._new_slide(bg)
        self._grid_line_v(slide, 2.9, dark=dark)
        self._grid_line_v(slide, 97.1, dark=dark)
        self._grid_line_h(slide, 5.15, dark=dark)
        self._grid_line_h(slide, 94.85, dark=dark)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        if label:
            self._label(slide, label, self.ML,
                        Emu(int(self.SLIDE_H.emu * 0.10)), dark=dark)
        if title:
            self._title(slide, title, self.ML,
                        Emu(int(self.SLIDE_H.emu * 0.15)), dark=dark)
        if subtitle:
            self._subtitle_text(slide, subtitle, self.ML,
                                Emu(int(self.SLIDE_H.emu * 0.24)),
                                dark=dark)
        if body:
            self._body_text(slide, body, self.ML,
                            Emu(int(self.SLIDE_H.emu * 0.32)),
                            dark=dark)
        self._footer(slide, dark=dark)

    def add_content_columns(self, label='', title='', columns=None):
        """Layout 7: 2-5 equal columns below a title."""
        columns = columns or []
        slide = self._new_slide(self.WHITE)
        self._grid_line_v(slide, 2.9)
        self._grid_line_v(slide, 97.1)
        self._grid_line_h(slide, 5.15)
        self._grid_line_h(slide, 27.53)
        self._grid_line_h(slide, 32.42)
        self._grid_line_h(slide, 94.85)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        if label:
            self._label(slide, label, self.ML, Emu(int(self.SLIDE_H.emu * 0.10)))
        if title:
            self._title(slide, title, self.ML, Emu(int(self.SLIDE_H.emu * 0.15)))
        # Columns area
        n = len(columns)
        if n == 0:
            self._footer(slide)
            return
        col_left_start = Emu(int(self.SLIDE_W.emu * 0.029))
        col_area_w = Emu(int(self.SLIDE_W.emu * 0.942))
        col_top = Emu(int(self.SLIDE_H.emu * 0.3242))
        col_h = Emu(int(self.SLIDE_H.emu * 0.6243))
        col_w = col_area_w // n
        for i, col in enumerate(columns):
            x = col_left_start + col_w * i
            # Column divider
            if i > 0:
                s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            x, col_top, Pt(0.5), col_h)
                s.fill.solid()
                s.fill.fore_color.rgb = self.GRID_LIGHT
                s.line.fill.background()
            # Column subtitle
            if col.get('subtitle'):
                self._txt(slide, col['subtitle'], x + Inches(0.15),
                          col_top + Inches(0.15),
                          Emu(col_w - Inches(0.3).emu), Inches(0.4),
                          size=Pt(11), color=self.BLUE, bold=True)
            # Column body
            if col.get('body'):
                self._txt(slide, col['body'], x + Inches(0.15),
                          col_top + Inches(0.55),
                          Emu(col_w - Inches(0.3).emu), Emu(col_h - Inches(0.7).emu),
                          size=Pt(9), color=self.TEXT_MUTED)
        self._footer(slide)

    def add_overview_about(self, label='', title='', subtitle='',
                            image_path=None, stats=None):
        """Layout 8: Content left, image right, 5 stats at bottom."""
        stats = stats or []
        slide = self._new_slide(self.WHITE)
        self._grid_line_v(slide, 2.9)
        self._grid_line_v(slide, 59.45)
        self._grid_line_v(slide, 97.1)
        self._grid_line_h(slide, 5.15)
        self._grid_line_h(slide, 74.85)
        self._grid_line_h(slide, 94.85)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        if label:
            self._label(slide, label, self.ML, Emu(int(self.SLIDE_H.emu * 0.10)))
        if title:
            self._title(slide, title, self.ML, Emu(int(self.SLIDE_H.emu * 0.15)),
                        width=Inches(5.3))
        if subtitle:
            self._txt(slide, subtitle, self.ML, Emu(int(self.SLIDE_H.emu * 0.28)),
                      Inches(5.3), Inches(2.5), size=Pt(11), color=self.TEXT_MUTED)
        # Image area (right)
        img_left = Emu(int(self.SLIDE_W.emu * 0.5945))
        img_top = Emu(int(self.SLIDE_H.emu * 0.0515))
        img_w = Emu(int(self.SLIDE_W.emu * 0.3765))
        img_h = Emu(int(self.SLIDE_H.emu * 0.697))
        if image_path:
            try:
                slide.shapes.add_picture(image_path, img_left, img_top, img_w, img_h)
            except Exception:
                self._placeholder_rect(slide, img_left, img_top, img_w, img_h)
        else:
            self._placeholder_rect(slide, img_left, img_top, img_w, img_h)
        # Stats row at bottom
        self._stats_row(slide, stats, Emu(int(self.SLIDE_W.emu * 0.029)),
                        Emu(int(self.SLIDE_H.emu * 0.7485)),
                        Emu(int(self.SLIDE_W.emu * 0.942)),
                        Emu(int(self.SLIDE_H.emu * 0.20)))
        self._footer(slide)

    def add_overview_stats(self, label='', title='', subtitle='',
                            image_path=None, stats=None):
        """Layout 9: Similar to overview-about with 4 stats."""
        stats = stats or []
        slide = self._new_slide(self.WHITE)
        self._grid_line_v(slide, 2.9)
        self._grid_line_v(slide, 50)
        self._grid_line_v(slide, 97.1)
        self._grid_line_h(slide, 5.15)
        self._grid_line_h(slide, 74.85)
        self._grid_line_h(slide, 94.85)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        if label:
            self._label(slide, label, self.ML, Emu(int(self.SLIDE_H.emu * 0.10)))
        if title:
            self._title(slide, title, self.ML, Emu(int(self.SLIDE_H.emu * 0.15)),
                        width=Inches(5.0))
        if subtitle:
            self._txt(slide, subtitle, self.ML, Emu(int(self.SLIDE_H.emu * 0.28)),
                      Inches(5.0), Inches(2.5), size=Pt(11), color=self.TEXT_DARK)
        # Image (right half)
        img_left = Emu(int(self.SLIDE_W.emu * 0.50))
        img_top = Emu(int(self.SLIDE_H.emu * 0.0515))
        img_w = Emu(int(self.SLIDE_W.emu * 0.471))
        img_h = Emu(int(self.SLIDE_H.emu * 0.697))
        if image_path:
            try:
                slide.shapes.add_picture(image_path, img_left, img_top, img_w, img_h)
            except Exception:
                self._placeholder_rect(slide, img_left, img_top, img_w, img_h)
        else:
            self._placeholder_rect(slide, img_left, img_top, img_w, img_h)
        # Stats
        self._stats_row(slide, stats, Emu(int(self.SLIDE_W.emu * 0.029)),
                        Emu(int(self.SLIDE_H.emu * 0.7485)),
                        Emu(int(self.SLIDE_W.emu * 0.942)),
                        Emu(int(self.SLIDE_H.emu * 0.20)))
        self._footer(slide)

    def add_product(self, label='', title='', subtitle='', body='',
                     image_path=None, image_bg=True):
        """Layout 10: Title + bullets left, screenshot right."""
        slide = self._new_slide(self.WHITE)
        self._grid_line_v(slide, 2.9)
        self._grid_line_v(slide, 50)
        self._grid_line_v(slide, 97.1)
        self._grid_line_h(slide, 5.15)
        self._grid_line_h(slide, 94.85)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        if label:
            self._label(slide, label, self.ML, Emu(int(self.SLIDE_H.emu * 0.10)))
        if title:
            self._title(slide, title, self.ML, Emu(int(self.SLIDE_H.emu * 0.15)),
                        width=Inches(5.3))
        if subtitle:
            self._subtitle_text(slide, subtitle, self.ML,
                                Emu(int(self.SLIDE_H.emu * 0.28)),
                                width=Inches(5.3))
        if body:
            self._body_text(slide, body, self.ML,
                            Emu(int(self.SLIDE_H.emu * 0.53)),
                            width=Inches(5.3))
        # Image area (right half)
        img_left = Emu(int(self.SLIDE_W.emu * 0.50))
        img_top = Emu(int(self.SLIDE_H.emu * 0.0515))
        img_w = Emu(int(self.SLIDE_W.emu * 0.471))
        img_h = Emu(int(self.SLIDE_H.emu * 0.897))
        if image_bg and not image_path:
            bg_rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              img_left, img_top, img_w, img_h)
            try:
                bg_rect.adjustments[0] = 0.03
            except Exception:
                pass
            bg_rect.fill.solid()
            bg_rect.fill.fore_color.rgb = self.LIGHT_BLUE
            bg_rect.line.fill.background()
            bg_rect.shadow.inherit = False
            tb = slide.shapes.add_textbox(
                img_left, Emu(img_top + img_h // 2 - Inches(0.2).emu),
                img_w, Inches(0.4))
            tf = tb.text_frame; tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = 'PRODUCT SCREENSHOT'
            r.font.size = Pt(9); r.font.name = self.FONT
            r.font.bold = True; r.font.color.rgb = self.BLUE
        if image_path:
            try:
                slide.shapes.add_picture(image_path, img_left + Inches(0.2),
                                          img_top + Inches(0.2),
                                          Emu(img_w - Inches(0.4).emu),
                                          Emu(img_h - Inches(0.4).emu))
            except Exception:
                pass
        self._footer(slide)

    def add_statement(self, accent='blue', label='', text='', variant=None):
        """Layout 11: Large text on colored band. Impactful quotes."""
        dark = variant == 'dark'
        bg = self.NAVY if dark else self.WHITE
        slide = self._new_slide(bg)
        self._grid_line_v(slide, 2.9, dark=dark)
        self._grid_line_v(slide, 97.1, dark=dark)
        self._grid_line_h(slide, 5.15, dark=dark)
        self._grid_line_h(slide, 94.85, dark=dark)
        mc = self.WHITE if dark else self.BLUE
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)), mc)
        # Band
        if dark:
            band_color = RGBColor(0x14, 0x2A, 0x45)
        elif accent == 'red':
            band_color = self.BAND_RED
        else:
            band_color = self.BAND_BLUE
        self._band_shape(slide, band_color)
        # Label
        if label:
            self._label(slide, label, self.ML,
                        Emu(int(self.SLIDE_H.emu * 0.2727)), dark=dark)
        # Text — preserve <span class="hl"> highlights (blue, or red on accent-red)
        text_color = self.TEXT_WHITE if dark else self.TEXT_DARK
        hl = self.RED if accent == 'red' else self.BLUE
        tb = slide.shapes.add_textbox(self.ML, Emu(int(self.SLIDE_H.emu * 0.37)),
                                      Emu(int(self.SLIDE_W.emu * 0.88)), Inches(2.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.paragraphs[0]
        p.line_spacing = 1.4
        self._rich_paragraph(p, parse_rich_runs(text, hl_color=hl),
                             Pt(18), text_color)
        self._footer(slide, dark=dark)

    def add_statement_stat(self, accent='blue', label='', stat='70%',
                            text='', source=''):
        """Layout 12: Big number + description on colored band."""
        slide = self._new_slide(self.WHITE)
        self._grid_line_v(slide, 2.9)
        self._grid_line_v(slide, 97.1)
        self._grid_line_h(slide, 5.15)
        self._grid_line_h(slide, 94.85)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        # Band
        band_color = self.BAND_RED if accent == 'red' else self.BAND_BLUE
        self._band_shape(slide, band_color)
        # Label
        if label:
            self._label(slide, label, self.ML,
                        Emu(int(self.SLIDE_H.emu * 0.2727)))
        # Stat number (large, left)
        stat_color = self.RED if accent == 'red' else self.BLUE
        self._txt(slide, stat, self.ML,
                  Emu(int(self.SLIDE_H.emu * 0.35)),
                  Inches(4), Inches(1.5), size=Pt(60),
                  color=stat_color, bold=True)
        # Description text (right of number) — preserve hl highlights
        hl = self.RED if accent == 'red' else self.BLUE
        tb = slide.shapes.add_textbox(Inches(5.5), Emu(int(self.SLIDE_H.emu * 0.38)),
                                      Inches(6.5), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE
        p = tf.paragraphs[0]
        p.line_spacing = 1.3
        self._rich_paragraph(p, parse_rich_runs(text, hl_color=hl),
                             Pt(13), self.TEXT_DARK)
        # Source
        if source:
            self._txt(slide, f"Source: {source}",
                      Inches(5.5), Emu(int(self.SLIDE_H.emu * 0.55)),
                      Inches(6), Inches(0.3), size=Pt(9),
                      color=self.TEXT_MUTED)
        self._footer(slide)

    def add_testimonial(self, logo='', quote='', name='', role='',
                         image_path=None):
        """Layout 13: Customer photo left, quote right."""
        slide = self._new_slide(self.WHITE)
        self._grid_line_v(slide, 2.9)
        self._grid_line_v(slide, 33)
        self._grid_line_v(slide, 97.1)
        self._grid_line_h(slide, 5.15)
        self._grid_line_h(slide, 94.85)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        # Image (left)
        img_left = self.ML
        img_top = Emu(int(self.SLIDE_H.emu * 0.20))
        img_w = Emu(int(self.SLIDE_W.emu * 0.21))
        img_h = Emu(int(self.SLIDE_H.emu * 0.48))
        if image_path:
            try:
                slide.shapes.add_picture(image_path, img_left, img_top, img_w, img_h)
            except Exception:
                self._placeholder_rect(slide, img_left, img_top, img_w, img_h)
        else:
            self._placeholder_rect(slide, img_left, img_top, img_w, img_h)
        # Label
        self._label(slide, 'TESTIMONIAL',
                    Emu(int(self.SLIDE_W.emu * 0.33)),
                    Emu(int(self.SLIDE_H.emu * 0.085)))
        # Vertical divider
        div_x = Emu(int(self.SLIDE_W.emu * 0.33))
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, div_x,
                                    Emu(int(self.SLIDE_H.emu * 0.12)),
                                    Pt(0.5),
                                    Emu(int(self.SLIDE_H.emu * 0.76)))
        s.fill.solid()
        s.fill.fore_color.rgb = self.GRID_LIGHT
        s.line.fill.background()
        # Content (right of divider)
        content_left = Emu(int(self.SLIDE_W.emu * 0.35))
        content_w = Emu(int(self.SLIDE_W.emu * 0.60))
        # Logo/company name
        if logo:
            self._txt(slide, logo, content_left,
                      Emu(int(self.SLIDE_H.emu * 0.20)),
                      content_w, Inches(0.5), size=Pt(22),
                      color=self.TEXT_DARK, bold=True)
        # Quote
        self._txt(slide, f'"{quote}"', content_left,
                  Emu(int(self.SLIDE_H.emu * 0.32)),
                  content_w, Inches(2.5), size=Pt(13),
                  color=self.TEXT_DARK)
        # Author divider
        author_y = Emu(int(self.SLIDE_H.emu * 0.72))
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, content_left, author_y,
                                    content_w, Pt(0.5))
        s.fill.solid()
        s.fill.fore_color.rgb = self.GRID_LIGHT
        s.line.fill.background()
        # Author name
        self._txt(slide, name, content_left,
                  Emu(int(self.SLIDE_H.emu * 0.75)),
                  content_w, Inches(0.3), size=Pt(12), color=self.BLUE)
        # Author role
        self._txt(slide, role, content_left,
                  Emu(int(self.SLIDE_H.emu * 0.80)),
                  content_w, Inches(0.3), size=Pt(10), color=self.TEXT_MUTED)
        self._footer(slide)

    def add_team(self, label='', title='', members=None):
        """Layout 14: Auto-layout grid for team introductions."""
        members = members or []
        slide = self._new_slide(self.WHITE)
        self._grid_line_v(slide, 2.9)
        self._grid_line_v(slide, 97.1)
        self._grid_line_h(slide, 5.15)
        self._grid_line_h(slide, 94.85)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        if label:
            self._label(slide, label, self.ML, Emu(int(self.SLIDE_H.emu * 0.10)))
        if title:
            self._title(slide, title, self.ML, Emu(int(self.SLIDE_H.emu * 0.15)))
        # Grid area
        n = len(members)
        if n == 0:
            self._footer(slide)
            return
        grid_left = Emu(int(self.SLIDE_W.emu * 0.029))
        grid_top = Emu(int(self.SLIDE_H.emu * 0.3242))
        grid_w = Emu(int(self.SLIDE_W.emu * 0.942))
        grid_h = Emu(int(self.SLIDE_H.emu * 0.6243))
        # Calculate grid dimensions
        cols = min(n, 5) if n <= 5 else min(n, 4)
        rows = (n + cols - 1) // cols
        cell_w = grid_w // cols
        cell_h = grid_h // rows
        for idx, member in enumerate(members):
            r = idx // cols
            c = idx % cols
            x = grid_left + cell_w * c
            y = grid_top + cell_h * r
            # Cell background
            cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                           x, y, cell_w, cell_h)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.WHITE
            cell.line.color.rgb = self.GRID_LIGHT
            cell.line.width = Pt(0.5)
            # Photo placeholder (left third)
            photo_w = Emu(int(cell_w * 0.35))
            photo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            x, y, photo_w, cell_h)
            photo.fill.solid()
            photo.fill.fore_color.rgb = self.OFF_WHITE
            photo.line.fill.background()
            if member.get('image'):
                try:
                    slide.shapes.add_picture(member['image'], x, y,
                                              photo_w, cell_h)
                except Exception:
                    pass
            # Name and role (right of photo)
            text_x = x + photo_w + Inches(0.15)
            text_w = Emu(cell_w - photo_w - Inches(0.3).emu)
            self._txt(slide, member.get('name', ''), text_x,
                      y + Emu(cell_h // 3), text_w, Inches(0.3),
                      size=Pt(10), color=self.BLUE)
            self._txt(slide, member.get('role', ''), text_x,
                      y + Emu(cell_h // 3) + Inches(0.25),
                      text_w, Inches(0.3),
                      size=Pt(8), color=self.TEXT_MUTED)
        self._footer(slide)

    def add_agenda_table(self, label='', title='', headers=None, rows=None):
        """Layout 15: Navy table with highlighted rows."""
        headers = headers or ['Session', 'Time', 'Location', 'Speaker']
        rows = rows or []
        slide = self._new_slide(self.NAVY_DARK)
        self._grid_line_v(slide, 2.9, dark=True)
        self._grid_line_v(slide, 97.1, dark=True)
        self._grid_line_h(slide, 5.15, dark=True)
        self._grid_line_h(slide, 94.85, dark=True)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)), self.BLUE)
        if label:
            self._label(slide, label, self.ML,
                        Emu(int(self.SLIDE_H.emu * 0.14)), dark=True)
        if title:
            self._title(slide, title, self.ML,
                        Emu(int(self.SLIDE_H.emu * 0.19)), dark=True)
        # Table
        table_left = Emu(int(self.SLIDE_W.emu * 0.029))
        table_top = Emu(int(self.SLIDE_H.emu * 0.30))
        table_w = Emu(int(self.SLIDE_W.emu * 0.942))
        n_rows = len(rows) + 1  # +1 for header
        n_cols = len(headers)
        row_h = Inches(0.45)
        tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                            table_left, table_top,
                                            table_w, row_h * n_rows)
        tbl = tbl_shape.table
        col_w = table_w // n_cols
        for ci in range(n_cols):
            tbl.columns[ci].width = col_w
        # Header row
        for ci, h in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.bold = True
                p.font.name = self.FONT
                p.font.color.rgb = self.TEXT_MUTED_W
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.NAVY_DARK
        # Data rows
        for ri, row in enumerate(rows):
            highlight = row.get('highlight', False)
            cells = row.get('cells', [])
            for ci, val in enumerate(cells):
                if ci >= n_cols:
                    break
                cell = tbl.cell(ri + 1, ci)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(9)
                    p.font.name = self.FONT
                    p.font.color.rgb = RGBColor(0xDD, 0xDD, 0xEE)
                cell.fill.solid()
                if highlight:
                    cell.fill.fore_color.rgb = RGBColor(0x0F, 0x23, 0x3D)
                else:
                    cell.fill.fore_color.rgb = self.NAVY_DARK
        self._footer(slide, dark=True)

    def add_roadmap(self, label='', title='', months=None, rows=None):
        """Layout 16: Gantt-style timeline with colored bars."""
        months = months or ['Q1', 'Q2', 'Q3', 'Q4']
        rows = rows or []
        slide = self._new_slide(self.WHITE)
        self._grid_line_v(slide, 2.9)
        self._grid_line_v(slide, 97.1)
        self._grid_line_h(slide, 5.15)
        self._grid_line_h(slide, 94.85)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        if label:
            self._label(slide, label, self.ML, Emu(int(self.SLIDE_H.emu * 0.10)))
        if title:
            self._title(slide, title, self.ML, Emu(int(self.SLIDE_H.emu * 0.15)))
        # Chart area
        chart_left = Emu(int(self.SLIDE_W.emu * 0.029))
        chart_top = Emu(int(self.SLIDE_H.emu * 0.25))
        chart_w = Emu(int(self.SLIDE_W.emu * 0.942))
        chart_h = Emu(int(self.SLIDE_H.emu * 0.70))
        label_w = Emu(int(chart_w * 0.15))
        track_w = chart_w - label_w
        track_left = chart_left + label_w
        # Month headers
        n_months = len(months)
        month_w = track_w // n_months
        for mi, m in enumerate(months):
            x = track_left + month_w * mi
            self._txt(slide, m, x, chart_top, month_w, Inches(0.35),
                      size=Pt(9), color=self.TEXT_DARK, bold=True,
                      align=PP_ALIGN.CENTER)
        # Month dividers
        for mi in range(n_months + 1):
            x = track_left + month_w * mi
            s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        x, chart_top + Inches(0.35),
                                        Pt(0.5), chart_h - Inches(0.35).emu)
            s.fill.solid()
            s.fill.fore_color.rgb = self.GRID_LIGHT
            s.line.fill.background()
        # Rows
        n_rows = len(rows)
        if n_rows == 0:
            self._footer(slide)
            return
        row_top = chart_top + Inches(0.4)
        row_area_h = chart_h - Inches(0.4).emu
        row_h = Emu(row_area_h) // n_rows
        for ri, row in enumerate(rows):
            y = Emu(row_top) + row_h * ri
            # Row label
            self._txt(slide, row.get('label', ''),
                      chart_left + Inches(0.1), y + Emu(row_h // 4),
                      Emu(label_w - Inches(0.2).emu), Emu(row_h // 2),
                      size=Pt(9), color=self.TEXT_DARK)
            # Row divider
            s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        chart_left, y + row_h,
                                        chart_w, Pt(0.5))
            s.fill.solid()
            s.fill.fore_color.rgb = self.GRID_LIGHT
            s.line.fill.background()
            # Bars
            for bar in row.get('bars', []):
                bar_start = bar.get('start', 0)
                bar_dur = bar.get('duration', 1)
                bar_color = self.BAR_COLORS.get(bar.get('color', 'blue'), self.BLUE)
                bar_label = bar.get('label', '')
                bx = track_left + Emu(int(month_w * bar_start))
                bw = Emu(int(month_w * bar_dur))
                bar_h = Emu(int(row_h * 0.6))
                by = y + Emu(int(row_h * 0.2))
                rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                               bx, by, bw, bar_h)
                rect.fill.solid()
                rect.fill.fore_color.rgb = bar_color
                rect.line.fill.background()
                if bar_label:
                    self._txt(slide, bar_label, bx + Inches(0.1),
                              by + Emu(bar_h // 4),
                              Emu(bw - Inches(0.2).emu), Emu(bar_h // 2),
                              size=Pt(8), color=self.WHITE)
        self._footer(slide)

    def add_thank_you(self):
        """Layout 17: Navy background with 'Thank you' text."""
        slide = self._new_slide(self.NAVY_DARK)
        self._grid_line_v(slide, 2.9, dark=True)
        self._grid_line_v(slide, 97.1, dark=True)
        self._grid_line_h(slide, 26.11, dark=True)
        self._grid_line_h(slide, 73.8, dark=True)
        # Logo
        self._txt(slide, 'Backbase', Inches(1.07), Inches(0.53),
                  Inches(2), Inches(0.4), size=Pt(14),
                  color=self.TEXT_WHITE, bold=True)
        # Motif
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.2611)), self.WHITE)
        # Thank you text
        self._txt(slide, 'Thank you', self.ML,
                  Emu(int(self.SLIDE_H.emu * 0.35)),
                  Emu(int(self.SLIDE_W.emu * 0.85)),
                  Inches(1.5), size=Pt(72),
                  color=self.TEXT_WHITE)

    # ══════════════════════════════════════════════════════
    #  INTERNAL HELPERS (continued)
    # ══════════════════════════════════════════════════════

    def _placeholder_rect(self, slide, left, top, w, h, caption='Screenshot'):
        """Framed image placeholder — a soft rounded panel with a centered
        muted caption, so a screenshot-less deck reads as intentional
        (a labelled frame) rather than an empty grey box."""
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        try:
            s.adjustments[0] = 0.03
        except Exception:
            pass
        s.fill.solid()
        s.fill.fore_color.rgb = self.OFF_WHITE
        s.line.color.rgb = self.BORDER
        s.line.width = Pt(1)
        s.shadow.inherit = False
        s.name = f"placeholder_{self._slide_num}"
        if caption:
            tb = slide.shapes.add_textbox(left, Emu(top + h // 2 - Inches(0.2).emu),
                                          w, Inches(0.4))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = caption.upper()
            r.font.size = Pt(9)
            r.font.name = self.FONT
            r.font.color.rgb = self.TEXT_MUTED
            r.font.bold = True

    def _stats_row(self, slide, stats, left, top, width, height):
        """Row of stat items with value + label."""
        n = len(stats)
        if n == 0:
            return
        # Top border
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(0.5))
        s.fill.solid()
        s.fill.fore_color.rgb = self.GRID_LIGHT
        s.line.fill.background()
        stat_w = width // n
        for i, stat in enumerate(stats):
            x = left + stat_w * i
            # Divider between stats
            if i > 0:
                d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                            x, top, Pt(0.5), height)
                d.fill.solid()
                d.fill.fore_color.rgb = self.GRID_LIGHT
                d.line.fill.background()
            # Value
            self._txt(slide, stat.get('value', ''),
                      x + Inches(0.15), top + Emu(height // 4),
                      Emu(stat_w - Inches(0.3).emu), Inches(0.4),
                      size=Pt(20), color=self.BLUE, bold=True)
            # Label
            self._txt(slide, stat.get('label', ''),
                      x + Inches(0.15), top + Emu(int(height * 0.6)),
                      Emu(stat_w - Inches(0.3).emu), Inches(0.3),
                      size=Pt(9), color=self.TEXT_DARK)

    # ══════════════════════════════════════════════════════
    #  RICH PRIMITIVES (cards, chips, panels, comparison,
    #  options, callout, freeform band) — added 2026-06 to
    #  close the HTML↔PPTX component-vocabulary gap.
    # ══════════════════════════════════════════════════════

    # Tone presets: (fill, default eyebrow color)
    TONE = {
        'off':   (OFF_WHITE,  RGBColor(0x6B, 0x77, 0x86)),  # neutral / muted eyebrow
        'blue':  (LIGHT_BLUE, RGBColor(0x33, 0x67, 0xFF)),  # platform / positive
        'white': (WHITE,      RGBColor(0x33, 0x67, 0xFF)),
        'navy':  (NAVY,       RGBColor(0x69, 0xFE, 0xFF)),  # dark callout, cyan eyebrow
    }
    ACCENT = {
        'red':  RGBColor(0xFF, 0x50, 0x3C),
        'blue': RGBColor(0x33, 0x67, 0xFF),
        'navy': RGBColor(0x04, 0x13, 0x26),
        'muted': RGBColor(0x6B, 0x77, 0x86),
    }

    def _round_rect(self, slide, left, top, w, h, fill, border=None,
                    radius=0.06, line_w=Pt(0.75)):
        """Rounded rectangle with a sane (small) corner radius."""
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        if border is not None:
            s.line.color.rgb = border
            s.line.width = line_w
        else:
            s.line.fill.background()
        s.shadow.inherit = False
        s.name = f"panel_{self._slide_num}"
        return s

    def _accent_bar(self, slide, left, top, h, color, w=Emu(38100)):
        """Thin vertical accent bar on the left edge of a card (~3px)."""
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()
        s.shadow.inherit = False
        return s

    def _eyebrow(self, slide, text, left, top, w, color, size=Pt(8.5)):
        """Uppercase bold eyebrow inside a card."""
        tf = self._txt(slide, text.upper(), left, top, w, Inches(0.3),
                       size=size, color=color, bold=True)
        tf.paragraphs[0].font.bold = True
        return tf

    def _rich_paragraph(self, p, runs, size, default_color, bold=False):
        """Write a list of (text, bold, color) runs into a paragraph."""
        first = True
        for text, run_bold, run_color in runs:
            r = p.add_run()
            r.text = text
            r.font.size = size
            r.font.name = self.FONT
            r.font.bold = bool(run_bold) or bold
            r.font.color.rgb = run_color or default_color
            first = False
        if first:  # no runs — keep an empty run so paragraph isn't blank
            r = p.add_run()
            r.text = ''
            r.font.size = size

    def _card(self, slide, x, y, w, h, tone='off', eyebrow=None,
              eyebrow_color=None, body=None, bullets=None, chips=None,
              accent=None, border=None, body_size=Pt(10), valign='middle'):
        """A content card: rounded panel + uppercase eyebrow + body / bullets / chips.
        valign='middle' vertically centres the content group so a card that is
        taller than its content reads balanced, never marooned at the top."""
        fill, default_eb = self.TONE.get(tone, self.TONE['off'])
        self._round_rect(slide, x, y, w, h, fill, border=border)
        if accent:
            self._accent_bar(slide, x, y, h, self.ACCENT.get(accent, self.BLUE))
        pad_l = x + Inches(0.18)
        pad_w = Emu(w - Inches(0.36).emu)
        # Vertically centre the content group within the card.
        content_emu = int(self._content_h(eyebrow, body, bullets, chips, pad_w,
                                          body_size.pt) * self._EMU_IN)
        if valign == 'middle' and h > content_emu:
            cy = Emu(y + max(Inches(0.16).emu, (h - content_emu) // 2))
        else:
            cy = y + Inches(0.16)
        if eyebrow:
            self._eyebrow(slide, eyebrow, pad_l, cy, pad_w,
                          eyebrow_color or default_eb)
            cy = cy + Inches(0.32)
        if body:
            tb = slide.shapes.add_textbox(pad_l, cy, pad_w,
                                          Emu(h - (cy - y) - Inches(0.16).emu))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            p = tf.paragraphs[0]
            p.line_spacing = 1.2
            self._rich_paragraph(p, parse_rich_runs(body), body_size,
                                 self.TEXT_DARK)
            cy = cy + Inches(0.5)
        if bullets:
            tb = slide.shapes.add_textbox(pad_l, cy, pad_w,
                                          Emu(h - (cy - y) - Inches(0.16).emu))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            for bi, bl in enumerate(bullets):
                p = tf.paragraphs[0] if bi == 0 else tf.add_paragraph()
                p.line_spacing = 1.15
                p.space_after = Pt(3)
                r0 = p.add_run(); r0.text = '→  '
                r0.font.size = body_size; r0.font.name = self.FONT
                r0.font.color.rgb = self.BLUE
                self._rich_paragraph(p, parse_rich_runs(bl), body_size,
                                     self.TEXT_DARK)
            cy = cy + Inches(0.3) * len(bullets)
        if chips:
            self._chip_row(slide, chips, pad_l, cy, pad_w)

    def _chip_row(self, slide, chips, left, top, max_w, size=Pt(8),
                  fill=None, border=None):
        """A wrapping row of pill chips."""
        fill = fill or self.WHITE
        border = border if border is not None else self.BORDER
        x = left
        y = top
        chip_h = Inches(0.26)
        gap = Inches(0.08)
        char_w = Inches(0.062)   # rough per-character advance at 8pt
        pad = Inches(0.26)
        for chip in chips:
            cw = Emu(int(pad.emu + len(chip) * char_w.emu))
            if x + cw > left + max_w:        # wrap
                x = left
                y = y + chip_h + Inches(0.06)
            s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cw, chip_h)
            try:
                s.adjustments[0] = 0.5
            except Exception:
                pass
            s.fill.solid(); s.fill.fore_color.rgb = fill
            s.line.color.rgb = border; s.line.width = Pt(0.5)
            s.shadow.inherit = False
            tf = s.text_frame
            tf.word_wrap = False
            tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
            tf.margin_top = 0; tf.margin_bottom = 0
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = chip
            r.font.size = size; r.font.name = self.FONT; r.font.bold = True
            r.font.color.rgb = self.NAVY
            x = x + cw + gap

    def _callout_strip(self, slide, left, top, w, lead, body, lead_color=None):
        """Full-width navy callout strip: cyan lead phrase + white body."""
        h = Inches(0.75)
        self._round_rect(slide, left, top, w, h, self.NAVY)
        tb = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.12),
                                      Emu(w - Inches(0.5).emu),
                                      Emu(h - Inches(0.24).emu))
        tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.line_spacing = 1.2
        r1 = p.add_run(); r1.text = lead + ' '
        r1.font.size = Pt(11); r1.font.bold = True; r1.font.name = self.FONT
        r1.font.color.rgb = lead_color or self.CYAN
        r2 = p.add_run(); r2.text = body
        r2.font.size = Pt(11); r2.font.name = self.FONT
        r2.font.color.rgb = self.WHITE

    # ── New rich layout methods ───────────────────────────

    def _content_header(self, slide, label, title, subtitle, dark=False):
        """Shared header chrome for content slides (label + title + subtitle)."""
        self._grid_line_v(slide, 2.9, dark=dark)
        self._grid_line_v(slide, 97.1, dark=dark)
        self._grid_line_h(slide, 5.15, dark=dark)
        self._grid_line_h(slide, 94.85, dark=dark)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        if label:
            self._label(slide, label, self.ML,
                        Emu(int(self.SLIDE_H.emu * 0.10)), dark=dark)
        if title:
            self._title(slide, title, self.ML,
                        Emu(int(self.SLIDE_H.emu * 0.15)), dark=dark)
        if subtitle:
            self._subtitle_text(slide, subtitle, self.ML,
                                Emu(int(self.SLIDE_H.emu * 0.255)), dark=dark)

    # ── Content sizing (engine discipline: size cards to content, never stretch) ──
    _EMU_IN = 914400

    def _text_block_h(self, text, width_in, pt):
        """Estimate rendered height (inches) of wrapped text at the given pt size."""
        plain = re.sub(r'<[^>]+>', '', str(text or ''))
        if not plain.strip():
            return 0.0
        char_w = pt * 0.0072          # ~Libre Franklin advance per char, inches (conservative)
        cpl = max(8, int(width_in / char_w))
        n = len(plain)
        lines = (n + cpl - 1) // cpl   # ceil
        return max(1, lines) * (pt * 1.3 / 72.0)

    def _chips_block_h(self, chips, width_in):
        """Estimate height (inches) of a wrapping chip row."""
        if not chips:
            return 0.0
        x, rows = 0.0, 1
        for c in chips:
            cw = 0.26 + len(str(c)) * 0.062
            if x + cw > width_in and x > 0:
                rows += 1
                x = cw + 0.08
            else:
                x += cw + 0.08
        return rows * 0.32

    def _content_h(self, eyebrow, body, bullets, chips, inner_w_emu, body_pt=10):
        """Height (inches) of a card's CONTENT group, excluding outer padding."""
        inner_in = inner_w_emu / self._EMU_IN
        h = 0.0
        if eyebrow:
            h += 0.32
        if body:
            h += self._text_block_h(body, inner_in, body_pt) + 0.04
        for b in (bullets or []):
            h += self._text_block_h(b, inner_in - 0.25, max(body_pt, 11)) + 0.06
        if chips:
            h += self._chips_block_h(chips, inner_in) + 0.06
        return h

    def _estimate_card_h(self, card, inner_w_emu):
        """Content-sized card height (inches) = content + outer padding."""
        ch = self._content_h(card.get('eyebrow'), card.get('body'),
                             card.get('bullets'), card.get('chips'), inner_w_emu)
        return max(ch + 0.34, 1.0)

    # ── Composition rule: anchor + balanced fill (engine discipline) ──────────
    # Content top sits on the engine's body line (32% of slide height) — a fixed,
    # moderate gap under the headline. The content band runs to the footer line.
    BODY_TOP_FRAC = 0.33      # headline→content gap, constant (was drifting 0.34–0.40)
    BAND_BOTTOM_FRAC = 0.90   # bottom of the content band (above footer)
    CARD_COMFORT_IN = 0.85    # breathing room added to content (proportionate, not fill)
    CARD_MIN_IN = 1.5         # a card is never a thin strip

    def _band(self):
        """(top_emu, available_height_emu) for the content band."""
        top = int(self.SLIDE_H.emu * self.BODY_TOP_FRAC)
        avail = int(self.SLIDE_H.emu * self.BAND_BOTTOM_FRAC) - top
        return top, avail

    def _row_card_h(self, cards, widths_emu, avail_emu):
        """Row height = tallest card's CONTENT + comfortable padding (proportionate
        to content, never a fixed fill — so cards breathe without internal float),
        floored to a sensible minimum and capped at the band. Text is centred."""
        content = max(
            self._content_h(c.get('eyebrow'), c.get('body'), c.get('bullets'),
                            c.get('chips'), w)
            for c, w in zip(cards, widths_emu))
        h_emu = int((content + self.CARD_COMFORT_IN) * self._EMU_IN)
        return Emu(min(max(h_emu, int(self.CARD_MIN_IN * self._EMU_IN)), avail_emu))

    def add_content_cards(self, label='', title='', subtitle='', cards=None,
                          callout=None):
        """Content slide whose body is a row of equal cards.
        card = {eyebrow, body|bullets, chips, tone, accent, eyebrow_color}."""
        cards = cards or []
        slide = self._new_slide(self.WHITE)
        self._content_header(slide, label, title, subtitle)
        n = len(cards)
        if n == 0:
            self._footer(slide); return
        area_l = self.ML
        area_w = Emu(int(self.SLIDE_W.emu * 0.942) - (self.ML.emu - int(self.SLIDE_W.emu * 0.029)))
        area_w = Emu(int(self.SLIDE_W.emu * 0.88))
        band_top, band_avail = self._band()
        top = Emu(band_top)
        avail = band_avail - (Inches(1.05).emu if callout else 0)
        gap = Inches(0.2)
        weights = [c.get('weight', 1.0) for c in cards]
        wsum = sum(weights)
        usable = Emu(area_w - gap.emu * (n - 1))
        widths = [Emu(int(usable * wt / wsum)) for wt in weights]
        # Content-sized row height — never stretch a short card to fill.
        h = self._row_card_h(cards, widths, avail)
        x = area_l
        for c, cw in zip(cards, widths):
            self._card(slide, x, top, cw, h,
                       tone=c.get('tone', 'off'),
                       eyebrow=c.get('eyebrow'),
                       eyebrow_color=c.get('eyebrow_color'),
                       body=c.get('body'),
                       bullets=c.get('bullets'),
                       chips=c.get('chips'),
                       accent=c.get('accent'),
                       border=self.BORDER if c.get('border') else None)
            x = Emu(x + cw + gap.emu)
        if callout:
            # Sit the callout just below the cards, not at a fixed low anchor.
            self._callout_strip(slide, area_l,
                                Emu(top + h.emu + Inches(0.25).emu), area_w,
                                callout.get('lead', ''), callout.get('body', ''))
        self._footer(slide)

    def add_comparison(self, label='', title='', subtitle='', left=None,
                       right=None, chips=None, arrow=True, callout=None):
        """Before→after two-tone comparison: red-accented 'from' card, arrow,
        blue-accented 'to' card. Optional chip row + navy callout below.
        left/right = {eyebrow, body|bullets, tone, accent, eyebrow_color, weight}."""
        slide = self._new_slide(self.WHITE)
        self._content_header(slide, label, title, subtitle)
        area_l = self.ML
        area_w = Emu(int(self.SLIDE_W.emu * 0.88))
        band_top, band_avail = self._band()
        top = Emu(band_top)
        if chips and callout:
            avail = band_avail - Inches(1.45).emu
        elif chips or callout:
            avail = band_avail - Inches(0.95).emu
        else:
            avail = band_avail
        arrow_w = Inches(0.45) if arrow else Inches(0.0)
        gap = Inches(0.12)
        lw = left.get('weight', 1.0) if left else 1.0
        rw = right.get('weight', 1.0) if right else 1.0
        usable = Emu(area_w - arrow_w.emu - gap.emu * 2)
        lcw = Emu(int(usable * lw / (lw + rw)))
        rcw = Emu(usable - lcw)
        # Content-sized height — tallest of the two cards, capped.
        present = [(c, w) for c, w in ((left, lcw), (right, rcw)) if c]
        h = self._row_card_h([c for c, _ in present], [w for _, w in present], avail)
        x = area_l
        if left:
            self._card(slide, x, top, lcw, h, tone=left.get('tone', 'off'),
                       eyebrow=left.get('eyebrow'),
                       eyebrow_color=left.get('eyebrow_color', self.RED),
                       body=left.get('body'), bullets=left.get('bullets'),
                       accent=left.get('accent', 'red'))
        x = Emu(x + lcw + gap.emu)
        if arrow:
            ar = self._txt(slide, '→', x, Emu(top + h.emu // 2 - Inches(0.3).emu),
                           arrow_w, Inches(0.6), size=Pt(22),
                           color=self.TEXT_MUTED, align=PP_ALIGN.CENTER)
        x = Emu(x + arrow_w.emu + gap.emu)
        if right:
            self._card(slide, x, top, rcw, h, tone=right.get('tone', 'blue'),
                       eyebrow=right.get('eyebrow'),
                       eyebrow_color=right.get('eyebrow_color', self.BLUE),
                       body=right.get('body'), bullets=right.get('bullets'),
                       accent=right.get('accent', 'blue'))
        cy = Emu(top + h.emu + Inches(0.25).emu)   # follow the cards, don't float low
        if chips:
            self._chip_row(slide, chips, area_l, cy, area_w)
            cy = Emu(cy + Inches(0.4).emu)
        if callout:
            self._callout_strip(slide, area_l, cy, area_w,
                                callout.get('lead', ''), callout.get('body', ''))
        self._footer(slide)

    def add_options(self, label='', title='', subtitle='', options=None):
        """Option cards (2–4). option = {tag, title, body, led_by, led_by_color,
        recommended}. The recommended card gets a blue border + RECOMMENDED badge."""
        options = options or []
        slide = self._new_slide(self.WHITE)
        self._content_header(slide, label, title, subtitle)
        n = len(options)
        if n == 0:
            self._footer(slide); return
        area_l = self.ML
        area_w = Emu(int(self.SLIDE_W.emu * 0.88))
        band_top, band_avail = self._band()
        top = Emu(band_top)
        avail = band_avail
        gap = Inches(0.2)
        cw = Emu((area_w - gap.emu * (n - 1)) // n)
        inner_in = (cw.emu - Inches(0.36).emu) / self._EMU_IN

        def _opt_h(opt):
            hh = 0.18 + 0.30 + 0.62 + 0.18          # pad + tag + title + bottom pad
            hh += self._text_block_h(opt.get('body', ''), inner_in, 9.5) + 0.05
            if opt.get('led_by'):
                hh += 0.25 + 0.70                    # gap + led-by footer block
            return max(hh, 1.6)

        # Option cards carry top (tag/title/body) + bottom (led-by) content, so a
        # little comfort makes them read fleshed-out without internal float.
        est_emu = int((max(_opt_h(o) for o in options) + 0.40) * self._EMU_IN)
        h = Emu(min(max(est_emu, int(self.CARD_MIN_IN * self._EMU_IN)), avail))
        x = area_l
        for opt in options:
            rec = opt.get('recommended')
            tone_fill = self.LIGHT_BLUE if rec else self.OFF_WHITE
            border = self.BLUE if rec else self.BORDER
            self._round_rect(slide, x, top, cw, h, tone_fill, border=border,
                             line_w=Pt(1.5) if rec else Pt(0.75))
            pad_l = x + Inches(0.18)
            pad_w = Emu(cw - Inches(0.36).emu)
            cy = top + Inches(0.18)
            # tag (Option N)
            self._eyebrow(slide, opt.get('tag', ''), pad_l, cy, pad_w,
                          self.BLUE if rec else self.TEXT_MUTED)
            cy = cy + Inches(0.3)
            # title
            self._txt(slide, opt.get('title', ''), pad_l, cy, pad_w, Inches(0.6),
                      size=Pt(12), color=self.NAVY, bold=True)
            cy = cy + Inches(0.62)
            # body
            tb = slide.shapes.add_textbox(pad_l, cy, pad_w, Inches(2.0))
            tf = tb.text_frame; tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            p = tf.paragraphs[0]; p.line_spacing = 1.25
            self._rich_paragraph(p, parse_rich_runs(opt.get('body', '')),
                                 Pt(9.5), self.TEXT_DARK)
            # led-by footer
            led = opt.get('led_by')
            if led:
                fy = Emu(top + h.emu - Inches(0.70).emu)   # anchored to card bottom
                divc = self.BLUE if rec else self.BORDER
                d = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, pad_l, fy,
                                           pad_w, Pt(0.75))
                d.fill.solid(); d.fill.fore_color.rgb = divc
                d.line.fill.background(); d.shadow.inherit = False
                self._eyebrow(slide, 'Upgrade led by', pad_l,
                              fy + Inches(0.08), pad_w,
                              self.BLUE if rec else self.TEXT_MUTED, size=Pt(7))
                self._txt(slide, led, pad_l, fy + Inches(0.34), pad_w,
                          Inches(0.3), size=Pt(11),
                          color=opt.get('led_by_color', self.NAVY), bold=True)
            # RECOMMENDED badge
            if rec:
                bw = Inches(1.5)
                badge = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Emu(x + cw - bw.emu - Inches(0.15).emu),
                    Emu(top - Inches(0.14).emu), bw, Inches(0.32))
                try:
                    badge.adjustments[0] = 0.5
                except Exception:
                    pass
                badge.fill.solid(); badge.fill.fore_color.rgb = self.BLUE
                badge.line.fill.background(); badge.shadow.inherit = False
                btf = badge.text_frame; btf.word_wrap = False
                btf.margin_top = 0; btf.margin_bottom = 0
                bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
                br = bp.add_run(); br.text = 'RECOMMENDED'
                br.font.size = Pt(8); br.font.bold = True
                br.font.name = self.FONT; br.font.color.rgb = self.WHITE
            x = Emu(x + cw + gap.emu)
        self._footer(slide)

    def add_timeline(self, label='', title='', subtitle='', milestones=None,
                     callout=None):
        """Horizontal milestone timeline — a connector line with coloured nodes
        and a card under each (title · body · coloured footer). Cleaner than a
        Gantt for a 'journey'/sequence story.
        milestone = {node, title, body, footer, accent ('blue'|'red'|'green'|
                     'amber'|'navy'), tone ('off'|'blue'|'red')}."""
        milestones = milestones or []
        slide = self._new_slide(self.WHITE)
        self._content_header(slide, label, title, subtitle)
        n = len(milestones)
        if n == 0:
            self._footer(slide); return
        red_tint = RGBColor(0xFA, 0xE0, 0xDE)
        area_l = self.ML.emu
        area_w = int(self.SLIDE_W.emu * 0.88)
        gap = Inches(0.2).emu
        cw = (area_w - gap * (n - 1)) // n
        centers = [area_l + cw // 2 + i * (cw + gap) for i in range(n)]

        # Connector line between first and last node
        line_y = int(self.SLIDE_H.emu * 0.47)
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(centers[0]),
                                    Emu(line_y), Emu(centers[-1] - centers[0]), Pt(2.5))
        ln.fill.solid(); ln.fill.fore_color.rgb = self.BORDER
        ln.line.fill.background(); ln.shadow.inherit = False

        # Node labels + dots
        for i, m in enumerate(milestones):
            acc = self.ACCENT.get(m.get('accent', 'blue'), self.BLUE)
            cx = centers[i]
            self._txt(slide, m.get('node', ''), Emu(cx - cw // 2),
                      Emu(line_y - Inches(0.42).emu), Emu(cw), Inches(0.3),
                      size=Pt(10), color=acc, bold=True, align=PP_ALIGN.CENTER)
            dd = Inches(0.2).emu
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Emu(cx - dd // 2),
                                         Emu(line_y - dd // 2 + Pt(1).emu), Emu(dd), Emu(dd))
            dot.fill.solid(); dot.fill.fore_color.rgb = acc
            dot.line.color.rgb = self.WHITE; dot.line.width = Pt(2.5)
            dot.shadow.inherit = False

        # Cards below — content-sized + comfortable (same discipline as cards)
        card_top = int(self.SLIDE_H.emu * 0.55)
        band_bottom = int(self.SLIDE_H.emu * (0.79 if callout else 0.90))
        inner_in = (cw - Inches(0.36).emu) / self._EMU_IN

        def _mile_h(m):
            hh = 0.18 + 0.42 + 0.18
            hh += self._text_block_h(m.get('body', ''), inner_in, 9.5) + 0.05
            if m.get('footer'):
                hh += 0.46
            return hh
        est_emu = int((max(_mile_h(m) for m in milestones) + 0.25) * self._EMU_IN)
        card_h = min(est_emu, band_bottom - card_top)

        for i, m in enumerate(milestones):
            x = area_l + i * (cw + gap)
            acc = self.ACCENT.get(m.get('accent', 'blue'), self.BLUE)
            tone = m.get('tone')
            fill = (red_tint if tone == 'red' else
                    self.LIGHT_BLUE if tone == 'blue' else self.OFF_WHITE)
            self._round_rect(slide, Emu(x), Emu(card_top), Emu(cw), Emu(card_h), fill)
            pad_l = Emu(x + Inches(0.18).emu)
            pad_w = Emu(cw - Inches(0.36).emu)
            cy = card_top + Inches(0.18).emu
            self._txt(slide, m.get('title', ''), pad_l, Emu(cy), pad_w, Inches(0.4),
                      size=Pt(12), color=self.NAVY, bold=True)
            cy += Inches(0.46).emu
            tb = slide.shapes.add_textbox(pad_l, Emu(cy), pad_w, Inches(1.6))
            tf = tb.text_frame; tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.NONE
            p = tf.paragraphs[0]; p.line_spacing = 1.2
            self._rich_paragraph(p, parse_rich_runs(m.get('body', '')), Pt(9.5),
                                 self.TEXT_DARK)
            if m.get('footer'):
                self._txt(slide, m['footer'], pad_l,
                          Emu(card_top + card_h - Inches(0.40).emu), pad_w,
                          Inches(0.3), size=Pt(9), color=acc, bold=True)

        if callout:
            self._callout_strip(slide, self.ML, Emu(int(self.SLIDE_H.emu * 0.83)),
                                Emu(area_w), callout.get('lead', ''),
                                callout.get('body', ''))
        self._footer(slide)

    # ══════════════════════════════════════════════════════
    #  OUTPUT
    # ══════════════════════════════════════════════════════

    def save(self, output_path):
        """Save the presentation and print summary."""
        self.prs.save(output_path)
        size_kb = Path(output_path).stat().st_size // 1024
        slide_count = len(self.prs.slides)
        print(f'\u2713 Saved {output_path} ({size_kb} KB, {slide_count} slides)')
