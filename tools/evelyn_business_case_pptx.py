#!/usr/bin/env python3
"""
Evelyn Partners — Business Case for Change (15-slide PPTX)

Mirrors the HTML deck at:
  Engagement/Evelyn Partners/Output/Evelyn_Business_Case_for_Change.html

Usage:
    python3 tools/evelyn_business_case_pptx.py
"""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from tools.frontline_slides_pptx import BackbaseSlidesPresenter, strip_html_tags
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE


class EvelynBusinessCase(BackbaseSlidesPresenter):
    """Generates the 15-slide Business Case for Change deck."""

    # Extra accent colors
    GREEN  = RGBColor(0x05, 0x96, 0x69)
    PURPLE = RGBColor(0x7B, 0x2F, 0xFF)
    AMBER  = RGBColor(0xD9, 0x77, 0x06)

    GREEN_LIGHT  = RGBColor(0xEE, 0xF9, 0xF4)
    PURPLE_LIGHT = RGBColor(0xF3, 0xEE, 0xFF)
    AMBER_LIGHT  = RGBColor(0xFD, 0xF5, 0xE6)
    RED_LIGHT    = RGBColor(0xFF, 0xEF, 0xEC)

    def _card(self, slide, left, top, w, h, bg_color=None, border_color=None):
        """Draw a card rectangle with optional border."""
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = bg_color or self.OFF_WHITE
        if border_color:
            s.line.color.rgb = border_color
            s.line.width = Pt(0.75)
        else:
            s.line.fill.background()
        s.name = f"card_{self._slide_num}_{left}"

    def _accent_line(self, slide, left, top, w, color):
        """Colored accent line at top of a card."""
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Pt(3))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()

    def _left_accent(self, slide, left, top, h, color):
        """Colored left border accent on a card."""
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Pt(3), h)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()

    def generate(self, output_path: str):
        self._init()
        self._slide_01_cover()
        self._slide_02_agenda()
        self._slide_03_why_now()
        self._slide_04_quote_imperative()
        self._slide_05_chapter_what_we_heard()
        self._slide_06_challenges()
        self._slide_07_value_at_stake()
        self._slide_08_bridge_statement()
        self._slide_09_business_outcomes()
        self._slide_10_cost_of_standing_still()
        self._slide_11_quote_choice()
        self._slide_12_three_paths()
        self._slide_13_commitment()
        self._slide_14_activity_plan()
        self._slide_15_thank_you()
        self.save(output_path)

    def _init(self):
        """Reset state."""
        pass

    # ══════════════════════════════════════════════════════
    #  SLIDE 1: Cover
    # ══════════════════════════════════════════════════════
    def _slide_01_cover(self):
        self.add_cover_color_block(
            label='EVELYN PARTNERS \u00d7 BACKBASE',
            title='Business Case\nfor Change',
            date='April 2026',
            partner=False,
        )

    # ══════════════════════════════════════════════════════
    #  SLIDE 2: Agenda
    # ══════════════════════════════════════════════════════
    def _slide_02_agenda(self):
        self.add_toc(
            label='AGENDA',
            title='Contents',
            numbered=True,
            items=[
                'Why Now \u2014 Three forces that make this the right moment',
                'What We Heard \u2014 Your challenges, your ambition, the value at stake',
                'Business Outcomes \u2014 What the upgrade delivers for Evelyn',
                'The Cost of Standing Still \u2014 What happens if nothing changes',
                'Three Paths to Consider \u2014 Options with escalating commitment and investment',
                'Our Commitment \u2014 What Backbase will do to help you build the case',
                'Mutual Activity Plan \u2014 Timeline from today to contract signature',
            ],
        )

    # ══════════════════════════════════════════════════════
    #  SLIDE 3: Why Now — 3 forcing functions
    # ══════════════════════════════════════════════════════
    def _slide_03_why_now(self):
        self.add_content_columns(
            label='BUSINESS CASE FOR CHANGE',
            title='A Convergence of Timing and Ambition',
            columns=[
                {
                    'subtitle': 'Renewal Window',
                    'body': (
                        'Contract renewal in May 2027 creates a natural decision '
                        'point. Early action secures preferential terms, avoids '
                        'disruption, and locks in the upgrade path before the '
                        'current stack reaches end-of-support.'
                    ),
                },
                {
                    'subtitle': 'Post-Acquisition Scrutiny',
                    'body': (
                        'NatWest ownership brings new governance, new reporting '
                        'standards, and heightened expectations on operational '
                        'efficiency. The digital platform must be board-ready '
                        '\u2014 not legacy-dependent.'
                    ),
                },
                {
                    'subtitle': 'AI Momentum',
                    'body': (
                        'The successful AI POC proved the art of the possible. '
                        'Waiting risks losing internal champions and falling '
                        'behind peers who are moving to production AI now. '
                        'The window to capitalise is open.'
                    ),
                },
            ],
        )

    # ══════════════════════════════════════════════════════
    #  SLIDE 4: Quote — The Imperative
    # ══════════════════════════════════════════════════════
    def _slide_04_quote_imperative(self):
        self.add_statement(
            accent='blue',
            label='THE IMPERATIVE',
            text=(
                'The question is not whether to modernise \u2014 '
                'it\u2019s whether to do it on your terms or under pressure.'
            ),
        )

    # ══════════════════════════════════════════════════════
    #  SLIDE 5: Chapter — What We Heard
    # ══════════════════════════════════════════════════════
    def _slide_05_chapter_what_we_heard(self):
        self.add_chapter_numbered(
            theme='navy',
            number='02',
            label='YOUR AMBITION. YOUR REALITY.',
            title='What We Heard',
            subtitle=(
                'Since October, we\u2019ve been listening. Before we talk '
                'about what Backbase can do, we want to reflect back what '
                'matters most to Evelyn.'
            ),
        )

    # ══════════════════════════════════════════════════════
    #  SLIDE 6: Challenges You're Navigating — 3 columns
    # ══════════════════════════════════════════════════════
    def _slide_06_challenges(self):
        self.add_content_columns(
            label='WHAT WE HEARD',
            title='The Challenges You\u2019re Navigating',
            columns=[
                {
                    'subtitle': 'The Missing Colleague Layer',
                    'body': (
                        'Relationship managers navigate between Xplan, Avaloq, '
                        'Aladdin, and the portal daily \u2014 with no unified '
                        'workspace that brings it all together. There is no '
                        '\u201csee what the client sees\u201d view for staff. '
                        'The result: system-hopping, fragmented context, and '
                        'time lost before every client interaction.'
                    ),
                },
                {
                    'subtitle': 'Workflow Beyond Onboarding',
                    'body': (
                        'Onboarding has progressed, but broader servicing '
                        'remains email-heavy. Annual reviews, suitability '
                        'evidence, payments, and corporate actions still '
                        'break across legacy systems. Operations orchestrate '
                        'through inboxes, not workflows.'
                    ),
                },
                {
                    'subtitle': 'Growth Channels Untapped',
                    'body': (
                        'Evelyn generates interest through events, networks, '
                        'and introductions \u2014 but there is no digital '
                        'nurture path from prospect to client. The IFA and '
                        'intermediary proposition is strategically important '
                        'but digitally underserved. Leads arrive, but '
                        'conversion is manual.'
                    ),
                },
            ],
        )

    # ══════════════════════════════════════════════════════
    #  SLIDE 7: The Value at Stake — 4 stat cards (custom)
    # ══════════════════════════════════════════════════════
    def _slide_07_value_at_stake(self):
        slide = self._new_slide(self.WHITE)
        # Grid
        self._grid_line_v(slide, 2.9)
        self._grid_line_v(slide, 97.1)
        self._grid_line_h(slide, 5.15)
        self._grid_line_h(slide, 94.85)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        # Label, title, subtitle
        self._label(slide, 'WHAT WE HEARD', self.ML,
                    Emu(int(self.SLIDE_H.emu * 0.10)))
        self._title(slide, 'The Value at Stake', self.ML,
                    Emu(int(self.SLIDE_H.emu * 0.15)))
        self._subtitle_text(
            slide,
            'These are not abstract challenges. They have a measurable cost '
            '\u2014 in time, in missed opportunities, and in competitive position.',
            self.ML, Emu(int(self.SLIDE_H.emu * 0.24)),
        )

        # 2x2 stat cards
        card_w = Inches(5.8)
        card_h = Inches(2.15)
        gap_x = Inches(0.5)
        gap_y = Inches(0.35)
        grid_left = self.ML
        grid_top = Inches(2.55)

        cards = [
            {
                'stat': '5+ Systems',
                'label': 'Daily Advisor System-Hopping',
                'body': (
                    'Every client interaction requires navigating across '
                    'Xplan, Avaloq, Aladdin, reporting tools, and the portal. '
                    'Context is rebuilt manually each time. No single pane of '
                    'glass for the adviser.'
                ),
                'accent': self.BLUE,
                'bg': self.LIGHT_BLUE,
            },
            {
                'stat': 'Email-Led',
                'label': 'Servicing & Operations',
                'body': (
                    'Annual reviews, service requests, and client follow-ups '
                    'are orchestrated through inboxes. No structured workflow, '
                    'no audit trail, no real-time status visibility for clients '
                    'or staff.'
                ),
                'accent': self.RED,
                'bg': self.RED_LIGHT,
            },
            {
                'stat': 'Manual',
                'label': 'Prospect-to-Client Conversion',
                'body': (
                    'Events and networks generate interest, but there is no '
                    'digital journey from prospect to onboarded client. Leads '
                    'are followed up individually. No nurture, no pipeline '
                    'visibility, no conversion analytics.'
                ),
                'accent': self.PURPLE,
                'bg': self.PURPLE_LIGHT,
            },
            {
                'stat': 'Core-Locked',
                'label': 'Payments & Corporate Actions',
                'body': (
                    'Critical client actions \u2014 money movements, corporate '
                    'action elections \u2014 still require direct Avaloq access. '
                    'These workflows cannot be surfaced in the portal or the '
                    'adviser layer today.'
                ),
                'accent': self.AMBER,
                'bg': self.AMBER_LIGHT,
            },
        ]

        for idx, c in enumerate(cards):
            col = idx % 2
            row = idx // 2
            cx = grid_left + col * (card_w + gap_x)
            cy = grid_top + row * (card_h + gap_y)

            # Card background
            self._card(slide, cx, cy, card_w, card_h, bg_color=c['bg'])
            # Left accent
            self._left_accent(slide, cx, cy, card_h, c['accent'])
            # Stat number
            self._txt(slide, c['stat'],
                      cx + Inches(0.3), cy + Inches(0.2),
                      card_w - Inches(0.6), Inches(0.55),
                      size=Pt(36), color=c['accent'], bold=True)
            # Label
            self._txt(slide, c['label'].upper(),
                      cx + Inches(0.3), cy + Inches(0.85),
                      card_w - Inches(0.6), Inches(0.3),
                      size=Pt(9), color=c['accent'], bold=True)
            # Body
            self._txt(slide, c['body'],
                      cx + Inches(0.3), cy + Inches(1.2),
                      card_w - Inches(0.6), Inches(0.85),
                      size=Pt(10), color=self.TEXT_MUTED)

        # Attribution
        self._txt(slide,
                  'Based on discussions with Martin Horton \u2014 '
                  'Chief Transformation & Digital Officer, Evelyn Partners',
                  self.ML, Inches(7.0), self.CW, Inches(0.3),
                  size=Pt(7), color=self.TEXT_MUTED)
        self._footer(slide)

    # ══════════════════════════════════════════════════════
    #  SLIDE 8: Bridge Statement — The Missing Layer
    # ══════════════════════════════════════════════════════
    def _slide_08_bridge_statement(self):
        self.add_statement(
            accent='blue',
            label='THE MISSING LAYER',
            text=(
                'Evelyn has built the target wealth stack. What\u2019s missing '
                'is the experience, workflow, and intelligence layer that makes '
                'it work for clients, colleagues, and partners.'
            ),
        )

    # ══════════════════════════════════════════════════════
    #  SLIDE 9: Business Outcomes — 4 stat cards (custom)
    # ══════════════════════════════════════════════════════
    def _slide_09_business_outcomes(self):
        slide = self._new_slide(self.WHITE)
        self._grid_line_v(slide, 2.9)
        self._grid_line_v(slide, 97.1)
        self._grid_line_h(slide, 5.15)
        self._grid_line_h(slide, 94.85)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        self._label(slide, 'BUSINESS CASE FOR CHANGE', self.ML,
                    Emu(int(self.SLIDE_H.emu * 0.10)))
        self._title(slide, 'What This Means for Evelyn', self.ML,
                    Emu(int(self.SLIDE_H.emu * 0.15)))
        self._subtitle_text(
            slide,
            'Quantified business outcomes from the platform upgrade and AI enablement.',
            self.ML, Emu(int(self.SLIDE_H.emu * 0.24)),
        )

        card_w = Inches(5.8)
        card_h = Inches(2.15)
        gap_x = Inches(0.5)
        gap_y = Inches(0.35)
        grid_left = self.ML
        grid_top = Inches(2.55)

        outcomes = [
            {
                'stat': '40%',
                'label': 'Reduction in Advisor Admin Time',
                'body': (
                    'Advisor Workspace and AI assistants shift advisors from '
                    'data gathering to client engagement. Industry benchmark: '
                    'advisors spend 60%+ of time on admin; best-in-class is '
                    'under 40%.'
                ),
                'accent': self.BLUE,
                'bg': self.LIGHT_BLUE,
            },
            {
                'stat': '4x Faster',
                'label': 'Meeting Preparation',
                'body': (
                    'AI-powered meeting prep \u2014 demonstrated in the Evelyn '
                    'POC \u2014 reduces prep from 2+ hours to under 30 minutes. '
                    'Immediate, measurable savings per advisor, per day.'
                ),
                'accent': self.GREEN,
                'bg': self.GREEN_LIGHT,
            },
            {
                'stat': '6x',
                'label': 'Faster Client Onboarding',
                'body': (
                    'Wealth 2.0 digital onboarding with document vault and '
                    'e-signature compresses onboarding from weeks to days. '
                    'Peer benchmark: 31 days to 5 days.'
                ),
                'accent': self.PURPLE,
                'bg': self.PURPLE_LIGHT,
            },
            {
                'stat': 'Lower TCO',
                'label': 'Platform Consolidation',
                'body': (
                    'Moving from custom legacy to product-standard reduces '
                    'maintenance burden, eliminates custom code risk, and '
                    'positions for continuous innovation.'
                ),
                'accent': self.AMBER,
                'bg': self.AMBER_LIGHT,
            },
        ]

        for idx, o in enumerate(outcomes):
            col = idx % 2
            row = idx // 2
            cx = grid_left + col * (card_w + gap_x)
            cy = grid_top + row * (card_h + gap_y)
            self._card(slide, cx, cy, card_w, card_h, bg_color=o['bg'])
            self._left_accent(slide, cx, cy, card_h, o['accent'])
            self._txt(slide, o['stat'],
                      cx + Inches(0.3), cy + Inches(0.2),
                      card_w - Inches(0.6), Inches(0.55),
                      size=Pt(36), color=o['accent'], bold=True)
            self._txt(slide, o['label'].upper(),
                      cx + Inches(0.3), cy + Inches(0.85),
                      card_w - Inches(0.6), Inches(0.3),
                      size=Pt(9), color=o['accent'], bold=True)
            self._txt(slide, o['body'],
                      cx + Inches(0.3), cy + Inches(1.2),
                      card_w - Inches(0.6), Inches(0.85),
                      size=Pt(10), color=self.TEXT_MUTED)

        self._txt(slide,
                  'Benchmarks from Backbase wealth management engagements '
                  '(Goodbody, HNB, industry analysis)',
                  self.ML, Inches(7.0), self.CW, Inches(0.3),
                  size=Pt(7), color=self.TEXT_MUTED)
        self._footer(slide)

    # ══════════════════════════════════════════════════════
    #  SLIDE 10: Cost of Standing Still — 2 comparison columns (custom)
    # ══════════════════════════════════════════════════════
    def _slide_10_cost_of_standing_still(self):
        slide = self._new_slide(self.WHITE)
        self._grid_line_v(slide, 2.9)
        self._grid_line_v(slide, 97.1)
        self._grid_line_h(slide, 5.15)
        self._grid_line_h(slide, 94.85)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        self._label(slide, 'BUSINESS CASE FOR CHANGE', self.ML,
                    Emu(int(self.SLIDE_H.emu * 0.10)))
        self._title(slide, 'What Happens If Nothing Changes', self.ML,
                    Emu(int(self.SLIDE_H.emu * 0.15)))
        self._subtitle_text(
            slide,
            'A side-by-side view: auto-renewal on the current stack vs. early renewal with upgrade.',
            self.ML, Emu(int(self.SLIDE_H.emu * 0.24)),
        )

        col_w = Inches(5.7)
        col_gap = Inches(0.6)
        col_left = self.ML
        col_top = Inches(2.55)
        col_h = Inches(4.2)

        # ── Left: Stay on Current Stack ──
        lx = col_left
        self._card(slide, lx, col_top, col_w, col_h, bg_color=self.WHITE,
                   border_color=RGBColor(0xE0, 0xC0, 0xC0))
        self._accent_line(slide, lx, col_top, col_w, self.RED)

        self._txt(slide, '\u2718  STAY ON CURRENT STACK',
                  lx + Inches(0.3), col_top + Inches(0.2),
                  col_w - Inches(0.6), Inches(0.35),
                  size=Pt(11), color=self.RED, bold=True)

        left_items = [
            'Custom code = growing maintenance cost',
            'Manual advisor workflows persist',
            'Legacy onboarding = slow client acquisition',
            'No AI foundation = starting from scratch later',
            'Harder to demonstrate value to the board',
        ]
        y = col_top + Inches(0.75)
        for item in left_items:
            self._txt(slide, f'\u2014  {item}', lx + Inches(0.3), y,
                      col_w - Inches(0.6), Inches(0.4),
                      size=Pt(11), color=self.TEXT_MUTED)
            y += Inches(0.5)

        self._txt(slide,
                  'Net result: Increasing technical debt, no competitive '
                  'differentiation, and a harder conversation with '
                  'ownership every quarter.',
                  lx + Inches(0.3), col_top + Inches(3.5),
                  col_w - Inches(0.6), Inches(0.6),
                  size=Pt(9), color=self.TEXT_MUTED)

        # ── Right: Move to Wealth 2.0 + AI ──
        rx = col_left + col_w + col_gap
        self._card(slide, rx, col_top, col_w, col_h, bg_color=self.WHITE,
                   border_color=RGBColor(0xC0, 0xE0, 0xC0))
        self._accent_line(slide, rx, col_top, col_w, self.GREEN)

        self._txt(slide, '\u2714  MOVE TO WEALTH 2.0 + AI',
                  rx + Inches(0.3), col_top + Inches(0.2),
                  col_w - Inches(0.6), Inches(0.35),
                  size=Pt(11), color=self.GREEN, bold=True)

        right_items = [
            'Product-standard = continuous innovation',
            'AI-augmented advisors from day one',
            'Digital-first onboarding = competitive advantage',
            'POC momentum \u2192 production in 2027',
            'Clear modernisation story for new ownership',
        ]
        y = col_top + Inches(0.75)
        for item in right_items:
            self._txt(slide, f'\u2714  {item}', rx + Inches(0.3), y,
                      col_w - Inches(0.6), Inches(0.4),
                      size=Pt(11), color=self.TEXT_MUTED)
            y += Inches(0.5)

        self._txt(slide,
                  'Net result: A modern, AI-ready platform that '
                  'demonstrates digital leadership from day one '
                  'of the new contract term.',
                  rx + Inches(0.3), col_top + Inches(3.5),
                  col_w - Inches(0.6), Inches(0.6),
                  size=Pt(9), color=self.GREEN)

        self._footer(slide)

    # ══════════════════════════════════════════════════════
    #  SLIDE 11: Quote — The Choice
    # ══════════════════════════════════════════════════════
    def _slide_11_quote_choice(self):
        self.add_statement(
            accent='blue',
            label='THE CHOICE',
            text=(
                'Auto-renewal preserves the status quo. Early renewal with '
                'upgrade positions Evelyn as a digital leader under NatWest '
                '\u2014 at preferential economics.'
            ),
        )

    # ══════════════════════════════════════════════════════
    #  SLIDE 12: Three Paths — 3 option cards (custom)
    # ══════════════════════════════════════════════════════
    def _slide_12_three_paths(self):
        slide = self._new_slide(self.WHITE)
        self._grid_line_v(slide, 2.9)
        self._grid_line_v(slide, 97.1)
        self._grid_line_h(slide, 5.15)
        self._grid_line_h(slide, 94.85)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        self._label(slide, 'PATH FORWARD', self.ML,
                    Emu(int(self.SLIDE_H.emu * 0.10)))
        self._title(slide, 'Three Paths to Consider', self.ML,
                    Emu(int(self.SLIDE_H.emu * 0.15)))
        self._subtitle_text(
            slide,
            'Each path reflects a different level of commitment \u2014 '
            'and a different level of Backbase investment.',
            self.ML, Emu(int(self.SLIDE_H.emu * 0.24)),
        )

        card_w = Inches(3.75)
        card_gap = Inches(0.4)
        card_top = Inches(2.5)
        card_h = Inches(4.5)
        start_x = self.ML

        options = [
            {
                'title': 'Auto-Renew',
                'tier': 'Do Nothing',
                'tier_color': self.TEXT_MUTED,
                'accent': self.GRID_LIGHT,
                'bg': self.OFF_WHITE,
                'border': self.GRID_LIGHT,
                'bullets': [
                    '12-month rolling renewal',
                    'Stay on current LTS 2409',
                    'Extended support becomes chargeable (Oct \u201926)',
                    'No Backbase investment',
                ],
                'summary': 'Unsupported Alpha stack core. No innovation story. Reactive.',
                'sum_color': self.TEXT_MUTED,
                'metrics': [('BB Investment', 'None'), ('NatWest Ready', 'Weak')],
                'met_color': self.TEXT_MUTED,
                'recommended': False,
            },
            {
                'title': 'Upgrade Only',
                'tier': 'Stay Current',
                'tier_color': self.AMBER,
                'accent': self.AMBER,
                'bg': self.WHITE,
                'border': RGBColor(0xE5, 0xD0, 0xA0),
                'bullets': [
                    'Upgrade to LTS 26.09',
                    'Short-term renewal (1\u20132 years)',
                    'Digital Wealth (feature parity)',
                    'Evelyn funds the upgrade',
                ],
                'summary': 'Current but not competitive. No advisor workspace, no AI.',
                'sum_color': self.AMBER,
                'metrics': [('BB Investment', 'Limited'), ('NatWest Ready', 'Moderate')],
                'met_color': self.AMBER,
                'recommended': False,
            },
            {
                'title': 'Strategic Partnership',
                'tier': 'Commit & Innovate',
                'tier_color': self.BLUE,
                'accent': self.BLUE,
                'bg': self.WHITE,
                'border': RGBColor(0xA0, 0xB8, 0xFF),
                'bullets': [
                    'Multi-year renewal (3\u20135 years)',
                    'Upgrade to 26.09 + Digital Wealth',
                    'Advisor Workspace + AI capabilities',
                    'Backbase invests in the upgrade',
                    'Incremental discounting on commitment',
                ],
                'summary': 'Positions Evelyn as digital leader. Proactive, not reactive.',
                'sum_color': self.BLUE,
                'metrics': [('BB Investment', 'Material'), ('NatWest Ready', 'Strong')],
                'met_color': self.GREEN,
                'recommended': True,
            },
        ]

        for i, o in enumerate(options):
            cx = start_x + i * (card_w + card_gap)
            self._card(slide, cx, card_top, card_w, card_h,
                       bg_color=o['bg'], border_color=o['border'])
            self._accent_line(slide, cx, card_top, card_w, o['accent'])

            # RECOMMENDED badge
            if o['recommended']:
                badge_w = Inches(1.3)
                badge_h = Inches(0.25)
                bx = cx + card_w - badge_w - Inches(0.15)
                by = card_top + Inches(0.15)
                badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                               bx, by, badge_w, badge_h)
                badge.fill.solid()
                badge.fill.fore_color.rgb = self.BLUE
                badge.line.fill.background()
                self._txt(slide, 'RECOMMENDED', bx + Inches(0.05), by,
                          badge_w - Inches(0.1), badge_h,
                          size=Pt(7), color=self.WHITE, bold=True,
                          align=PP_ALIGN.CENTER)

            # Title
            self._txt(slide, o['title'],
                      cx + Inches(0.25), card_top + Inches(0.3),
                      card_w - Inches(0.5), Inches(0.35),
                      size=Pt(14), color=self.TEXT_DARK, bold=True)
            # Tier label
            self._txt(slide, o['tier'].upper(),
                      cx + Inches(0.25), card_top + Inches(0.7),
                      card_w - Inches(0.5), Inches(0.25),
                      size=Pt(8), color=o['tier_color'], bold=True)

            # Bullets
            by = card_top + Inches(1.2)
            for bullet in o['bullets']:
                self._txt(slide, f'\u2192  {bullet}',
                          cx + Inches(0.25), by,
                          card_w - Inches(0.5), Inches(0.35),
                          size=Pt(9), color=self.TEXT_MUTED)
                by += Inches(0.4)

            # Summary
            self._txt(slide, o['summary'],
                      cx + Inches(0.25), card_top + Inches(3.3),
                      card_w - Inches(0.5), Inches(0.5),
                      size=Pt(9), color=o['sum_color'])

            # Bottom metrics
            met_y = card_top + Inches(3.9)
            for mi, (ml, mv) in enumerate(o['metrics']):
                mx = cx + Inches(0.25) + mi * Inches(1.6)
                self._txt(slide, ml.upper(), mx, met_y,
                          Inches(1.4), Inches(0.2),
                          size=Pt(7), color=self.TEXT_MUTED, bold=True)
                self._txt(slide, mv, mx, met_y + Inches(0.2),
                          Inches(1.4), Inches(0.3),
                          size=Pt(12), color=o['met_color'], bold=True)

        self._footer(slide)

    # ══════════════════════════════════════════════════════
    #  SLIDE 13: Our Commitment — 4 cards (custom)
    # ══════════════════════════════════════════════════════
    def _slide_13_commitment(self):
        slide = self._new_slide(self.WHITE)
        self._grid_line_v(slide, 2.9)
        self._grid_line_v(slide, 97.1)
        self._grid_line_h(slide, 5.15)
        self._grid_line_h(slide, 94.85)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        self._label(slide, 'OUR COMMITMENT TO YOU', self.ML,
                    Emu(int(self.SLIDE_H.emu * 0.10)))
        self._title(slide, 'If You Choose the Strategic Path', self.ML,
                    Emu(int(self.SLIDE_H.emu * 0.15)))
        self._subtitle_text(
            slide,
            'Here is what Backbase will do to help you build the case and make it happen.',
            self.ML, Emu(int(self.SLIDE_H.emu * 0.24)),
        )

        card_w = Inches(5.8)
        card_h = Inches(1.8)
        gap_x = Inches(0.5)
        gap_y = Inches(0.3)
        grid_left = self.ML
        grid_top = Inches(2.55)

        commitments = [
            ('\u2460 Sharpen the Value Story Together',
             'We will run 2\u20133 deep-dive workshops with your team to '
             'validate the solution scope, map integrations, and quantify '
             'the business outcomes. No guesswork \u2014 evidence-based.'),
            ('\u2461 Build Your Executive Decision Paper',
             'We will produce a board-ready business case with ROI modelling, '
             'showing 5\u20137x return on the incremental investment. Designed '
             'to arm you for internal conversations.'),
            ('\u2462 Backbase Services Investment',
             'Michiel and the European services team will make a material '
             'investment into the upgrade delivery \u2014 directly reducing '
             'your cost to execute. Conditional on renewal commitment.'),
            ('\u2463 Present Back by Mid-May',
             'Within one month, we come back with: validated solution scope, '
             'commercial proposal, and the decision paper \u2014 ready for '
             'you to take to the right people when the time is right.'),
        ]

        for idx, (title, body) in enumerate(commitments):
            col = idx % 2
            row = idx // 2
            cx = grid_left + col * (card_w + gap_x)
            cy = grid_top + row * (card_h + gap_y)
            self._card(slide, cx, cy, card_w, card_h, bg_color=self.LIGHT_BLUE)
            self._txt(slide, title, cx + Inches(0.25), cy + Inches(0.2),
                      card_w - Inches(0.5), Inches(0.35),
                      size=Pt(11), color=self.BLUE, bold=True)
            self._txt(slide, body, cx + Inches(0.25), cy + Inches(0.65),
                      card_w - Inches(0.5), Inches(1.0),
                      size=Pt(10), color=self.TEXT_MUTED)

        # CTA banner
        cta_y = grid_top + 2 * (card_h + gap_y) + Inches(0.15)
        cta = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self.ML, cta_y, self.CW, Inches(0.45),
        )
        cta.fill.solid()
        cta.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)
        cta.line.color.rgb = RGBColor(0xCC, 0xD9, 0xFF)
        cta.line.width = Pt(0.5)
        self._txt(slide,
                  'What we need from you: a point of contact for deep-dives  \u2022  '
                  'confirmation of appetite for Option 3  \u2022  '
                  'visibility on who holds budget authority post-June',
                  self.ML + Inches(0.2), cta_y + Inches(0.05),
                  Emu(self.CW.emu - Inches(0.4).emu), Inches(0.35),
                  size=Pt(9), color=self.BLUE, bold=True,
                  align=PP_ALIGN.CENTER)

        self._footer(slide)

    # ══════════════════════════════════════════════════════
    #  SLIDE 14: Mutual Activity Plan — table (custom)
    # ══════════════════════════════════════════════════════
    def _slide_14_activity_plan(self):
        slide = self._new_slide(self.WHITE)
        self._grid_line_v(slide, 2.9)
        self._grid_line_v(slide, 97.1)
        self._grid_line_h(slide, 5.15)
        self._grid_line_h(slide, 94.85)
        self._motif(slide, Emu(int(self.SLIDE_W.emu * 0.029)),
                    Emu(int(self.SLIDE_H.emu * 0.0515)))
        self._label(slide, 'PATH FORWARD', self.ML,
                    Emu(int(self.SLIDE_H.emu * 0.10)))
        self._title(slide, 'Mutual Activity Plan', self.ML,
                    Emu(int(self.SLIDE_H.emu * 0.15)))
        self._subtitle_text(
            slide,
            'From today to upgrade kickoff \u2014 six milestones over six months.',
            self.ML, Emu(int(self.SLIDE_H.emu * 0.24)),
        )

        # Table
        headers = ['#', 'Milestone', 'What happens', 'When']
        col_widths = [Inches(0.6), Inches(3.0), Inches(6.5), Inches(1.9)]
        rows = [
            ('1', 'Strategic Alignment',
             'Confirm appetite for Option 3, identify decision-makers, agree next steps',
             'Today', self.BLUE, True),
            ('2', 'Deep-Dive Workshops',
             '2\u20133 sessions: solution scope, integration mapping, bill of materials',
             'Apr \u2013 May', None, False),
            ('3', 'Decision Paper Delivered',
             'Executive business case with ROI model \u2014 ready to take internally',
             'Mid-May', None, True),
            ('4', 'Internal Green Light',
             'Regulatory approval \u2192 internal communication opens',
             'June', None, False),
            ('5', 'Commercial & Contract',
             'Best & final proposal, legal review, procurement, signature',
             'Jul \u2013 Sep', None, True),
            ('6', 'Upgrade Kickoff',
             'Backbase-led LTS 26.09 upgrade with services investment',
             'October', self.GREEN, False),
        ]

        tbl_left = Emu(int(self.SLIDE_W.emu * 0.029))
        tbl_top = Inches(2.5)
        tbl_w = Emu(int(self.SLIDE_W.emu * 0.942))
        n_rows = len(rows) + 1
        n_cols = len(headers)
        row_h = Inches(0.55)

        tbl_shape = slide.shapes.add_table(n_rows, n_cols,
                                            tbl_left, tbl_top,
                                            tbl_w, row_h * n_rows)
        tbl = tbl_shape.table
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = w

        # Header row
        for ci, h in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.bold = True
                p.font.name = self.FONT
                p.font.color.rgb = self.TEXT_MUTED
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.WHITE

        # Data rows
        for ri, (num, milestone, desc, when, accent, highlight) in enumerate(rows):
            bg = RGBColor(0xF0, 0xF4, 0xFF) if highlight else self.OFF_WHITE
            accent_c = accent or self.BLUE

            cells_data = [num, milestone, desc, when]
            for ci, val in enumerate(cells_data):
                cell = tbl.cell(ri + 1, ci)
                cell.text = val
                for p in cell.text_frame.paragraphs:
                    p.font.name = self.FONT
                    if ci == 0:
                        p.font.size = Pt(11)
                        p.font.bold = True
                        p.font.color.rgb = accent_c
                    elif ci == 1:
                        p.font.size = Pt(11)
                        p.font.bold = True
                        if accent:
                            p.font.color.rgb = accent
                        else:
                            p.font.color.rgb = self.TEXT_DARK
                    elif ci == 2:
                        p.font.size = Pt(10)
                        p.font.color.rgb = self.TEXT_MUTED
                    elif ci == 3:
                        p.font.size = Pt(10)
                        p.font.bold = True if accent else False
                        p.font.color.rgb = accent_c if accent else self.TEXT_MUTED
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg

        # CTA banner
        cta_y = tbl_top + row_h * n_rows + Inches(0.3)
        cta = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            self.ML, cta_y, self.CW, Inches(0.45),
        )
        cta.fill.solid()
        cta.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)
        cta.line.color.rgb = RGBColor(0xCC, 0xD9, 0xFF)
        cta.line.width = Pt(0.5)
        self._txt(slide,
                  'Immediate next step: schedule deep-dive workshop #1 in late April',
                  self.ML + Inches(0.2), cta_y + Inches(0.05),
                  Emu(self.CW.emu - Inches(0.4).emu), Inches(0.35),
                  size=Pt(9), color=self.BLUE, bold=True,
                  align=PP_ALIGN.CENTER)

        self._footer(slide)

    # ══════════════════════════════════════════════════════
    #  SLIDE 15: Thank You
    # ══════════════════════════════════════════════════════
    def _slide_15_thank_you(self):
        self.add_thank_you()

    # ══════════════════════════════════════════════════════
    #  SAVE
    # ══════════════════════════════════════════════════════
    def save(self, output_path: str):
        self.prs.save(output_path)
        print(f'\u2713 Saved {output_path}')


# ══════════════════════════════════════════════════════
#  CLI
# ══════════════════════════════════════════════════════
if __name__ == '__main__':
    out = Path(__file__).resolve().parent.parent / \
        'Engagement' / 'Evelyn Partners' / 'Output' / \
        'Evelyn_Business_Case_for_Change.pptx'
    EvelynBusinessCase('Evelyn Partners — Business Case for Change').generate(str(out))
