#!/usr/bin/env python3
"""
PptxPresenter — Reusable base class for Backbase-branded PPTX generation.

Extracted from the 4 existing Schroders generators. Provides:
  - Backbase brand colors, fonts, margins as class constants
  - 15 helper methods for common slide components
  - Template loading + slide stripping for Backbase Master Template
  - Fresh presentation creation for lightweight output

Subclass this and implement slide builder methods for each engagement.

Usage:
    from tools.pptx_presenter import PptxPresenter

    class MyDeck(PptxPresenter):
        def generate(self, output_path):
            self._init_presentation()  # or self._init_from_template(path)
            self._slide_cover('Title', 'Subtitle')
            self.save(output_path)
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


class PptxPresenter:
    """Base class for Backbase-branded PPTX presentation generators."""

    # ── Dimensions (Google Slides widescreen) ─────────────
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    # ── Colors (Backbase Brand — matches design-system.md) ─
    DARK_BG      = RGBColor(0x09, 0x1C, 0x35)  # --bb-dark / --dark
    CARD_BG      = RGBColor(0x0E, 0x22, 0x40)  # dark card fill
    BLUE         = RGBColor(0x33, 0x66, 0xFF)  # --bb-blue / --blue
    PURPLE       = RGBColor(0x7B, 0x2F, 0xFF)  # --bb-purple / --purple
    RED          = RGBColor(0xDC, 0x26, 0x26)  # --bb-red / --red
    GREEN        = RGBColor(0x05, 0x96, 0x69)  # --bb-green / --green
    AMBER        = RGBColor(0xD9, 0x77, 0x06)  # --bb-amber / --amber
    CYAN         = RGBColor(0x08, 0x91, 0xB2)  # --bb-cyan / --cyan
    CYAN_BRIGHT  = RGBColor(0x69, 0xFE, 0xFF)  # bright cyan for dark slides
    GOLD         = RGBColor(0xB4, 0x53, 0x09)  # --gold
    WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BG     = RGBColor(0xF5, 0xFA, 0xFF)  # --bg
    LIGHT_CARD   = RGBColor(0xFF, 0xFF, 0xFF)  # --card
    DARK_TEXT     = RGBColor(0x09, 0x1C, 0x35)  # --text
    SUB_TEXT      = RGBColor(0x64, 0x74, 0x8B)  # --sub
    MUTED        = RGBColor(0x94, 0xA3, 0xB8)  # --muted
    BORDER       = RGBColor(0xE2, 0xE8, 0xF0)  # --border
    BORDER_DK    = RGBColor(0x1A, 0x3A, 0x5E)  # dark border
    BLUE_LIGHT   = RGBColor(0xEF, 0xF6, 0xFF)  # light blue fill
    GREEN_LIGHT  = RGBColor(0xEC, 0xFD, 0xF5)  # light green fill
    AMBER_LIGHT  = RGBColor(0xFF, 0xFB, 0xEB)  # light amber fill
    PURPLE_LIGHT = RGBColor(0xF5, 0xF3, 0xFF)  # light purple fill
    RED_LIGHT    = RGBColor(0xFE, 0xF2, 0xF2)  # light red fill
    ALT_ROW      = RGBColor(0xF8, 0xFA, 0xFC)  # table alternate row

    # ── Typography ────────────────────────────────────────
    FONT = 'Libre Franklin'

    # ── Layout (standard widescreen) ──────────────────────
    ML = Inches(0.6)      # margin left
    MR = Inches(0.6)      # margin right
    CW = Inches(12.1)     # content width
    CT = Inches(1.6)      # content top (below heading area)

    # ══════════════════════════════════════════════════════
    #  INITIALIZATION
    # ══════════════════════════════════════════════════════

    def _init_presentation(self):
        """Create a fresh blank presentation with Backbase dimensions."""
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_W
        self.prs.slide_height = self.SLIDE_H
        self.blank_layout = self.prs.slide_layouts[6]  # blank layout

    def _init_from_template(self, template_path):
        """Load a PPTX template, strip all existing slides, find blank layout."""
        self.prs = Presentation(template_path)
        # Strip all existing slides
        xml_slides = self.prs.slides._sldIdLst
        for sldId in list(xml_slides):
            rId = sldId.get('r:id')
            if rId is None:
                rId = sldId.attrib.get(
                    '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'
                )
            if rId:
                self.prs.part.drop_rel(rId)
            xml_slides.remove(sldId)
        # Find blank layout (prefer one with 0 placeholders)
        self.blank_layout = None
        for layout in self.prs.slide_layouts:
            if len(layout.placeholders) == 0:
                self.blank_layout = layout
                break
        if not self.blank_layout:
            self.blank_layout = min(self.prs.slide_layouts,
                                     key=lambda l: len(l.placeholders))

    # ══════════════════════════════════════════════════════
    #  HELPER METHODS
    # ══════════════════════════════════════════════════════

    def _new_slide(self, dark=False):
        """Add a blank slide with dark or light background."""
        slide = self.prs.slides.add_slide(self.blank_layout)
        # Remove inherited placeholders (from template layouts)
        for ph in list(slide.placeholders):
            sp = ph._element
            sp.getparent().remove(sp)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.DARK_BG if dark else self.LIGHT_BG
        return slide

    def _txt(self, slide, text, left, top, width, height, size=Pt(12),
             color=None, bold=False, align=PP_ALIGN.LEFT, font=None,
             anchor=None):
        """Add a single-paragraph textbox."""
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        if anchor:
            tf.paragraphs[0].alignment = align
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color or self.DARK_TEXT
        p.font.name = font or self.FONT
        p.alignment = align
        return tf

    def _section_label(self, slide, text, top=Inches(0.4), color=None,
                       dark=False):
        """Uppercase section label (overline text)."""
        self._txt(slide, text.upper(), self.ML, top, self.CW, Pt(16),
                  size=Pt(9), color=color or (self.CYAN_BRIGHT if dark else self.BLUE),
                  bold=True)

    def _heading(self, slide, main, accent='', top=Inches(0.75), size=Pt(28),
                 dark=False):
        """Two-part heading: main text + colored accent text."""
        tb = slide.shapes.add_textbox(self.ML, top, self.CW, Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = main
        r1.font.size = size
        r1.font.bold = True
        r1.font.color.rgb = self.WHITE if dark else self.DARK_TEXT
        r1.font.name = self.FONT
        if accent:
            r2 = p.add_run()
            r2.text = accent
            r2.font.size = size
            r2.font.bold = True
            r2.font.color.rgb = self.CYAN_BRIGHT if dark else self.BLUE
            r2.font.name = self.FONT

    def _subtitle(self, slide, text, top=Inches(1.2), dark=False):
        """Muted subtitle text below heading."""
        self._txt(slide, text, self.ML, top, self.CW, Inches(0.4),
                  size=Pt(11), color=self.MUTED if dark else self.SUB_TEXT)

    def _card(self, slide, left, top, w, h, fill=None, border=None,
              dark=False, border_top_color=None):
        """Rounded rectangle card shape."""
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = fill or (self.CARD_BG if dark else self.LIGHT_CARD)
        s.line.color.rgb = border or (self.BORDER_DK if dark else self.BORDER)
        s.line.width = Pt(1)
        return s

    def _stat_box(self, slide, value, label, left, top, w, h,
                  val_color=None, bg=None):
        """Stat card: large bold value + uppercase label below."""
        self._card(slide, left, top, w, h, fill=bg or self.BLUE_LIGHT)
        self._txt(slide, str(value), left, top + Inches(0.06), w, Inches(0.3),
                  size=Pt(18), color=val_color or self.BLUE, bold=True,
                  align=PP_ALIGN.CENTER)
        self._txt(slide, label.upper(), left, top + Inches(0.36), w, Pt(10),
                  size=Pt(6.5), color=self.SUB_TEXT, bold=True,
                  align=PP_ALIGN.CENTER)

    def _multi_text(self, slide, runs, left, top, width, height,
                    align=PP_ALIGN.LEFT):
        """Textbox with multiple styled runs in a single paragraph.
        runs = [(text, size, color, bold), ...]"""
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        for text, size, color, bold in runs:
            r = p.add_run()
            r.text = str(text)
            r.font.size = size
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = self.FONT
        return tf

    def _divider(self, slide, top, color=None):
        """Thin horizontal line spanning content width."""
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   self.ML, top, self.CW, Pt(1))
        s.fill.solid()
        s.fill.fore_color.rgb = color or self.BORDER
        s.line.fill.background()

    def _bullet_list(self, slide, items, left, top, width, height,
                     size=Pt(9), color=None, spacing=Pt(2)):
        """Multi-line bullet list with arrow prefix."""
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"\u2192 {item}"
            p.font.size = size
            p.font.color.rgb = color or self.SUB_TEXT
            p.font.name = self.FONT
            p.space_after = spacing
        return tf

    def _bar_rect(self, slide, left, top, w, h, fill):
        """Solid rectangle (no border) — for bar charts and progress bars."""
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        s.line.fill.background()
        return s

    def _colored_top_line(self, slide, left, top, w, color):
        """Thin colored accent line at top of a card (3pt tall)."""
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Pt(3))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()

    def _insight_card(self, slide, label, text, left, top, w, h,
                      label_color=None, bg=None, border=None):
        """Card with a bold colored label header and body text below."""
        self._card(slide, left, top, w, h, fill=bg or self.BLUE_LIGHT,
                   border=border or self.BORDER)
        self._txt(slide, label.upper(), left + Inches(0.15), top + Inches(0.08),
                  w - Inches(0.3), Pt(10),
                  size=Pt(7), color=label_color or self.BLUE, bold=True)
        self._txt(slide, text, left + Inches(0.15), top + Inches(0.25),
                  w - Inches(0.3), h - Inches(0.3),
                  size=Pt(8), color=self.SUB_TEXT)

    def _add_table(self, slide, rows, col_widths, left, top,
                   row_height=Inches(0.35), header_bg=None,
                   header_color=None, body_size=Pt(8)):
        """Styled table with header row and alternating body rows.
        rows = list of lists. First row is treated as header."""
        nrows = len(rows)
        ncols = len(rows[0])
        total_w = sum(col_widths)
        tbl_shape = slide.shapes.add_table(nrows, ncols, left, top,
                                           Inches(total_w), row_height * nrows)
        tbl = tbl_shape.table
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(cw)

        for ri, row in enumerate(rows):
            for ci, cell_text in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = str(cell_text)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(7) if ri == 0 else body_size
                    p.font.bold = ri == 0
                    p.font.name = self.FONT
                    p.font.color.rgb = (header_color or self.WHITE) if ri == 0 else self.DARK_TEXT
                cell.fill.solid()
                if ri == 0:
                    cell.fill.fore_color.rgb = header_bg or self.DARK_BG
                else:
                    cell.fill.fore_color.rgb = self.LIGHT_CARD if ri % 2 == 1 else self.ALT_ROW
        return tbl

    def _footer(self, slide, slide_number, dark=False):
        """Add 'Backbase | n' footer at bottom-right."""
        text_color = self.MUTED if dark else self.SUB_TEXT
        tb = slide.shapes.add_textbox(
            self.SLIDE_W - Inches(1.8), self.SLIDE_H - Inches(0.4),
            Inches(1.5), Inches(0.25)
        )
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        # "Backbase" text
        r1 = p.add_run()
        r1.text = 'Backbase'
        r1.font.size = Pt(8)
        r1.font.color.rgb = text_color
        r1.font.name = self.FONT
        # separator
        r2 = p.add_run()
        r2.text = '  |  '
        r2.font.size = Pt(8)
        r2.font.color.rgb = text_color
        r2.font.name = self.FONT
        # slide number
        r3 = p.add_run()
        r3.text = str(slide_number)
        r3.font.size = Pt(8)
        r3.font.color.rgb = text_color
        r3.font.name = self.FONT

    # ══════════════════════════════════════════════════════
    #  COMPOSITE SLIDE BUILDERS (common patterns)
    # ══════════════════════════════════════════════════════

    def _slide_cover(self, title_lines, subtitle='', label='',
                     pills=None):
        """Dark hero cover slide.
        title_lines = [(text, color), ...] for multi-line titles.
        pills = [(text, color), ...] for tag pills."""
        s = self._new_slide(dark=True)
        if label:
            self._section_label(s, label, top=Inches(0.5), dark=True)

        # Multi-line title
        tb = s.shapes.add_textbox(self.ML, Inches(1.4), self.CW, Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        for line_text, line_color in title_lines:
            p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = line_text
            r.font.size = Pt(44)
            r.font.bold = True
            r.font.color.rgb = line_color
            r.font.name = self.FONT

        if subtitle:
            self._txt(s, subtitle,
                      Inches(2), Inches(3.3), Inches(9), Inches(0.6),
                      size=Pt(13), color=self.MUTED, align=PP_ALIGN.CENTER)

        # Pill tags row
        if pills:
            pill_y = Inches(4.1)
            total_pills = len(pills)
            pill_w = Inches(2.2)
            gap = Inches(0.2)
            total_w = pill_w * total_pills + gap * (total_pills - 1)
            start_x = (self.SLIDE_W - total_w) / 2
            for i, (pill_text, pill_color) in enumerate(pills):
                px = start_x + i * (pill_w + gap)
                pill_shape = self._bar_rect(s, px, pill_y, pill_w, Inches(0.3),
                                            fill=pill_color)
                self._txt(s, pill_text, px, pill_y + Inches(0.04),
                          pill_w, Inches(0.22),
                          size=Pt(8), color=self.WHITE, bold=True,
                          align=PP_ALIGN.CENTER)

        self._footer(s, 1, dark=True)
        return s

    def _slide_content(self, label, heading, heading_accent='',
                       subtitle_text='', slide_number=1):
        """Standard content slide scaffold (light background).
        Returns the slide for further composition."""
        s = self._new_slide(dark=False)
        self._section_label(s, label)
        self._heading(s, heading, accent=heading_accent)
        if subtitle_text:
            self._subtitle(s, subtitle_text)
        self._footer(s, slide_number)
        return s

    def _slide_dark_feature(self, label, heading, heading_accent='',
                            subtitle_text='', slide_number=1):
        """Dark feature slide scaffold.
        Returns the slide for further composition."""
        s = self._new_slide(dark=True)
        self._section_label(s, label, dark=True)
        self._heading(s, heading, accent=heading_accent, dark=True)
        if subtitle_text:
            self._subtitle(s, subtitle_text, dark=True)
        self._footer(s, slide_number, dark=True)
        return s

    def _slide_closing(self, title_lines, subtitle='', label='',
                       slide_number=1):
        """Dark closing slide (same structure as cover)."""
        s = self._new_slide(dark=True)
        if label:
            self._section_label(s, label, top=Inches(0.5), dark=True)

        tb = s.shapes.add_textbox(self.ML, Inches(2.0), self.CW, Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        for line_text, line_color in title_lines:
            p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = line_text
            r.font.size = Pt(40)
            r.font.bold = True
            r.font.color.rgb = line_color
            r.font.name = self.FONT

        if subtitle:
            self._txt(s, subtitle,
                      Inches(2), Inches(4.0), Inches(9), Inches(0.6),
                      size=Pt(13), color=self.MUTED, align=PP_ALIGN.CENTER)

        self._footer(s, slide_number, dark=True)
        return s

    # ══════════════════════════════════════════════════════
    #  OUTPUT
    # ══════════════════════════════════════════════════════

    def save(self, output_path):
        """Save the presentation and print summary."""
        self.prs.save(output_path)
        size_kb = Path(output_path).stat().st_size // 1024
        slide_count = len(self.prs.slides)
        print(f'\u2713 Saved {output_path} ({size_kb} KB, {slide_count} slides)')
