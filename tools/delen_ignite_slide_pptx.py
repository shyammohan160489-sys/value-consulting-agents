"""
Delen IGNITE next-step slide — single-slide PPTX at 10" x 5.625" (matches Delen deck canvas).

Replicates the SEB "IGNITE | Let's inspire a way forward" reference layout:
- 2 columns (Day 1 green / Day 2 blue)
- 4 checkmark bullets left, 4 bullet items in bordered box right
- 3 pill chips at the bottom of each column
- Eyebrow label, H1 title, divider, footer

Content tailored to Delen Private Bank (wealth, internal-build-first, unified interface).
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree


# Canvas — matches Delen deck exactly
CANVAS_W = Inches(10.0)
CANVAS_H = Inches(5.625)

# Colors
NAVY = RGBColor(0x09, 0x1C, 0x35)
BLUE = RGBColor(0x33, 0x66, 0xFF)
GREEN = RGBColor(0x21, 0xB5, 0x73)
MUTED = RGBColor(0x5C, 0x6E, 0x84)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_FILL = RGBColor(0x21, 0xB5, 0x73)
BLUE_FILL = RGBColor(0x33, 0x66, 0xFF)

FONT = "Libre Franklin"


def _disable_autofit(tf):
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    for child in list(bodyPr):
        if child.tag in (qn("a:spAutoFit"), qn("a:normAutofit")):
            bodyPr.remove(child)
    bodyPr.append(etree.SubElement(bodyPr, qn("a:noAutofit")))


def _set_run(run, text, size_pt, color, bold=False, font=FONT):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_text(slide, x, y, w, h, text, size_pt, color, bold=False,
              align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    _disable_autofit(tf)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    _set_run(p.add_run(), text, size_pt, color, bold=bold)
    return box


def _add_bullet_list(slide, x, y, w, h, items, size_pt, color, bullet_color,
                     bullet_char="✓", bullet_bold=True, line_gap_pt=6):
    """Render a bullet list. First run per paragraph = colored bullet, rest = body."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    _disable_autofit(tf)

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if i > 0:
            p.space_before = Pt(line_gap_pt)

        r_bullet = p.add_run()
        _set_run(r_bullet, f"{bullet_char}  ", size_pt, bullet_color, bold=bullet_bold)

        # Support bold prefix via "**prefix** rest" convention
        if "**" in item:
            before, rest = item.split("**", 1)
            bold_part, after = rest.split("**", 1)
            if before:
                _set_run(p.add_run(), before, size_pt, color)
            _set_run(p.add_run(), bold_part, size_pt, color, bold=True)
            _set_run(p.add_run(), after, size_pt, color)
        else:
            _set_run(p.add_run(), item, size_pt, color)
    return box


def _add_chip(slide, x, y, w, h, label, fill_color):
    """Pill-shaped chip with colored fill and white label."""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    # Try to maximize rounding
    try:
        shp.adjustments[0] = 0.5
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    shp.shadow.inherit = False

    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    _disable_autofit(tf)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), label, 9, WHITE, bold=True)
    return shp


def _add_box_outline(slide, x, y, w, h, border_color=BORDER):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.background()
    shp.line.color.rgb = border_color
    shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def _add_line(slide, x1, y1, x2, y2, color=BORDER, weight_pt=0.75):
    shp = slide.shapes.add_connector(1, x1, y1, x2, y2)
    shp.line.color.rgb = color
    shp.line.width = Pt(weight_pt)
    return shp


def _add_accent_square(slide, x, y, size):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, size, size)
    shp.fill.solid()
    shp.fill.fore_color.rgb = BLUE
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def build_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # ── Eyebrow (blue accent square + uppercase label)
    _add_accent_square(slide, Inches(0.45), Inches(0.33), Inches(0.1))
    _add_text(slide, Inches(0.62), Inches(0.28), Inches(6.0), Inches(0.25),
              "THE AI-NATIVE TARGET STATE", size_pt=10, color=NAVY, bold=False)

    # ── Title: "IGNITE | Let's inspire a way forward" (IGNITE in blue)
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.6),
                                         Inches(9.1), Inches(0.6))
    ttf = title_box.text_frame
    ttf.word_wrap = True
    ttf.margin_left = 0
    ttf.margin_right = 0
    ttf.margin_top = 0
    ttf.margin_bottom = 0
    _disable_autofit(ttf)
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.LEFT
    _set_run(tp.add_run(), "IGNITE", 26, BLUE, bold=True)
    _set_run(tp.add_run(), " | ", 26, NAVY, bold=False)
    _set_run(tp.add_run(), "Let's inspire a way forward", 26, NAVY, bold=False)

    # ── Horizontal divider under title
    _add_line(slide, Inches(0.45), Inches(1.25), Inches(9.55), Inches(1.25),
              color=BORDER, weight_pt=0.75)

    # ── Vertical divider between columns
    _add_line(slide, Inches(5.0), Inches(1.4), Inches(5.0), Inches(4.95),
              color=BORDER, weight_pt=0.75)

    # ─────────────────────────────────────────
    # LEFT COLUMN — Day 1
    # ─────────────────────────────────────────
    left_x = Inches(0.5)
    left_w = Inches(4.4)

    # Header
    hbox = slide.shapes.add_textbox(left_x, Inches(1.45), left_w, Inches(0.35))
    htf = hbox.text_frame
    htf.word_wrap = True
    htf.margin_left = 0
    htf.margin_right = 0
    htf.margin_top = 0
    htf.margin_bottom = 0
    _disable_autofit(htf)
    hp = htf.paragraphs[0]
    _set_run(hp.add_run(), "Day 1", 16, GREEN, bold=True)
    _set_run(hp.add_run(), " | discovery & validation", 16, NAVY, bold=False)

    day1_bullets = [
        "**Map the current advisor-to-client flow** — where does it break, with whom, how often?",
        "**Assess maturity (L1 → L2 → L3)** — where is Delen today across onboarding, advisory, servicing?",
        "**Identify the unified interface gap** — what internal modules stay, what needs embedding into one workspace?",
        "**Validate priority use cases** — RM workspace, prospect portal, portfolio advisory, client self-service",
    ]
    _add_bullet_list(slide, left_x, Inches(1.95), left_w, Inches(2.55),
                     day1_bullets, size_pt=9.5, color=NAVY, bullet_color=GREEN,
                     bullet_char="✓", line_gap_pt=7)

    # Day 1 chips
    day1_chips = ["Capability Assessment", "Prioritized Use Cases", "Best-in-class Benchmark"]
    chip_y = Inches(4.65)
    chip_h = Inches(0.32)
    chip_w = Inches(1.38)
    chip_gap = Inches(0.08)
    for i, label in enumerate(day1_chips):
        cx = left_x + (chip_w + chip_gap) * i
        _add_chip(slide, cx, chip_y, chip_w, chip_h, label, GREEN_FILL)

    # ─────────────────────────────────────────
    # RIGHT COLUMN — Day 2
    # ─────────────────────────────────────────
    right_x = Inches(5.1)
    right_w = Inches(4.4)

    # Header
    hbox2 = slide.shapes.add_textbox(right_x, Inches(1.45), right_w, Inches(0.35))
    htf2 = hbox2.text_frame
    htf2.word_wrap = True
    htf2.margin_left = 0
    htf2.margin_right = 0
    htf2.margin_top = 0
    htf2.margin_bottom = 0
    _disable_autofit(htf2)
    hp2 = htf2.paragraphs[0]
    _set_run(hp2.add_run(), "Day 2", 16, BLUE, bold=True)
    _set_run(hp2.add_run(), " | architecture, roadmap & value case", 16, NAVY, bold=False)

    # Bordered box around Day 2 bullets
    box_x = right_x
    box_y = Inches(1.9)
    box_w = right_w
    box_h = Inches(2.6)
    _add_box_outline(slide, box_x, box_y, box_w, box_h, border_color=BORDER)

    day2_bullets = [
        "**Design the target architecture** — Backbase platform with seamless Delen internal-module embedding",
        "**Quantify the business case** — AUM per RM uplift, advisor productivity, onboarding acceleration, cost-to-serve",
        "**Build the phased roadmap** — quick wins → foundation → target state",
        "**Live demo deep-dive** — hands-on with the Advisor Workspace & AI-native capabilities",
    ]
    _add_bullet_list(slide, box_x + Inches(0.2), box_y + Inches(0.15),
                     box_w - Inches(0.4), box_h - Inches(0.3),
                     day2_bullets, size_pt=9.5, color=NAVY, bullet_color=BLUE,
                     bullet_char="•", bullet_bold=True, line_gap_pt=7)

    # Day 2 chips
    day2_chips = ["Target Architecture", "Tailored Solutioning", "Business Case"]
    for i, label in enumerate(day2_chips):
        cx = right_x + (chip_w + chip_gap) * i
        _add_chip(slide, cx, chip_y, chip_w, chip_h, label, BLUE_FILL)

    # ── Footer: backbase wordmark (text) + slide number
    _add_text(slide, Inches(7.8), Inches(5.25), Inches(1.4), Inches(0.22),
              "Backbase", size_pt=10, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
    _add_text(slide, Inches(9.25), Inches(5.25), Inches(0.3), Inches(0.22),
              "21", size_pt=9, color=MUTED, bold=False, align=PP_ALIGN.RIGHT)

    return slide


def main():
    repo_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    out_dir = os.path.join(repo_root, "outputs")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "delen_ignite_slide.pptx")

    prs = Presentation()
    prs.slide_width = CANVAS_W
    prs.slide_height = CANVAS_H

    build_slide(prs)

    prs.save(out_path)
    print(f"Wrote: {out_path}")
    print(f"Canvas: {prs.slide_width/914400:.3f}\" x {prs.slide_height/914400:.3f}\"")


if __name__ == "__main__":
    main()
