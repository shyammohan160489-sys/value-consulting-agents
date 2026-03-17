#!/usr/bin/env python3
"""
Schroders Commercial Refresh v3 — PPTX Generator

15-slide commercial proposal deck matching schroders_commercial_refresh_v3.html.
Google Slides compatible (13.333" × 7.5" widescreen). Backbase brand.

v3 changes: Marketplace licensing baked into both models, CLO row restored
without Essential/Premium labels, discount %s removed.

Usage:
    python3 tools/schroders_commercial_v3_pptx.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


class SchrodersCommercialV3Pptx:
    """Generates a 15-slide commercial proposal PPTX."""

    # ── Dimensions ─────────────────────────────────────────
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    # ── Colors (Backbase Brand — matches HTML CSS vars) ────
    DARK_BG    = RGBColor(0x09, 0x1C, 0x35)
    CARD_BG    = RGBColor(0x0E, 0x22, 0x40)
    BLUE       = RGBColor(0x33, 0x66, 0xFF)
    PURPLE     = RGBColor(0x7B, 0x2F, 0xFF)
    RED        = RGBColor(0xDC, 0x26, 0x26)
    GREEN      = RGBColor(0x05, 0x96, 0x69)
    AMBER      = RGBColor(0xD9, 0x77, 0x06)
    CYAN       = RGBColor(0x08, 0x91, 0xB2)
    GOLD       = RGBColor(0xB4, 0x53, 0x09)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BG   = RGBColor(0xF5, 0xFA, 0xFF)
    LIGHT_CARD = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_TEXT   = RGBColor(0x09, 0x1C, 0x35)
    SUB_TEXT    = RGBColor(0x64, 0x74, 0x8B)
    MUTED      = RGBColor(0x94, 0xA3, 0xB8)
    BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
    BORDER_DK  = RGBColor(0x1A, 0x3A, 0x5E)
    BLUE_LIGHT = RGBColor(0xEF, 0xF6, 0xFF)
    GREEN_LIGHT = RGBColor(0xEC, 0xFD, 0xF5)
    AMBER_LIGHT = RGBColor(0xFF, 0xFB, 0xEB)
    PURPLE_LIGHT = RGBColor(0xF5, 0xF3, 0xFF)
    RED_LIGHT  = RGBColor(0xFE, 0xF2, 0xF2)

    FONT = 'Libre Franklin'
    ML = Inches(0.6)     # margin left
    MR = Inches(0.6)     # margin right
    CW = Inches(12.1)    # content width
    CT = Inches(1.6)     # content top

    # ── Data ───────────────────────────────────────────────
    YOY_ENT = [
        {'year': 'Year 1', 'phase': 'Foundation', 'lic': 2.45, 'impl': 1.5, 'base': 11.0, 'ai': 0},
        {'year': 'Year 2', 'phase': 'Expand',     'lic': 2.69, 'impl': 1.0, 'base': 17.6, 'ai': 2.0},
        {'year': 'Year 3', 'phase': 'Scale',      'lic': 3.73, 'impl': 0.6, 'base': 27.8, 'ai': 4.0},
        {'year': 'Year 4', 'phase': 'Optimise',   'lic': 4.14, 'impl': 0.6, 'base': 36.8, 'ai': 5.7},
        {'year': 'Year 5', 'phase': 'Compound',   'lic': 4.55, 'impl': 0.6, 'base': 36.1, 'ai': 6.7},
    ]
    YOY_PAG = [
        {'year': 'Year 1', 'phase': 'Foundation', 'lic': 2.88, 'impl': 1.5, 'base': 11.0, 'ai': 0},
        {'year': 'Year 2', 'phase': 'Expand',     'lic': 3.66, 'impl': 1.0, 'base': 17.6, 'ai': 2.0},
        {'year': 'Year 3', 'phase': 'Scale',      'lic': 4.93, 'impl': 0.6, 'base': 27.8, 'ai': 4.0},
        {'year': 'Year 4', 'phase': 'Optimise',   'lic': 4.94, 'impl': 0.6, 'base': 36.8, 'ai': 5.7},
        {'year': 'Year 5', 'phase': 'Compound',   'lic': 4.95, 'impl': 0.6, 'base': 36.1, 'ai': 6.7},
    ]

    # ── HELPERS ────────────────────────────────────────────

    def _new_slide(self, dark=False):
        layout = self.prs.slide_layouts[6]  # blank
        slide = self.prs.slides.add_slide(layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.DARK_BG if dark else self.LIGHT_BG
        return slide

    def _txt(self, slide, text, left, top, width, height, size=Pt(12),
             color=None, bold=False, align=PP_ALIGN.LEFT, font=None,
             anchor=None):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        if anchor:
            tf.paragraphs[0].alignment = align
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color or self.DARK_TEXT
        p.font.name = font or self.FONT
        p.alignment = align
        return tf

    def _section_label(self, slide, text, top=Inches(0.4), color=None, dark=False):
        self._txt(slide, text.upper(), self.ML, top, self.CW, Pt(16),
                  size=Pt(9), color=color or (RGBColor(0x69, 0xFE, 0xFF) if dark else self.BLUE),
                  bold=True)

    def _heading(self, slide, main, accent='', top=Inches(0.75), size=Pt(28),
                 dark=False):
        tb = slide.shapes.add_textbox(self.ML, top, self.CW, Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = main
        r1.font.size = size
        r1.font.bold = True
        r1.font.color.rgb = self.WHITE if dark else self.DARK_TEXT
        r1.font.name = self.FONT
        if accent:
            r2 = p.add_run()
            r2.text = accent
            r2.font.size = size
            r2.font.bold = True
            r2.font.color.rgb = self.BLUE if not dark else RGBColor(0x69, 0xFE, 0xFF)
            r2.font.name = self.FONT

    def _subtitle(self, slide, text, top=Inches(1.2), dark=False):
        self._txt(slide, text, self.ML, top, self.CW, Inches(0.4),
                  size=Pt(11), color=self.MUTED if dark else self.SUB_TEXT)

    def _card(self, slide, left, top, w, h, fill=None, border=None, dark=False,
              border_top_color=None):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = fill or (self.CARD_BG if dark else self.LIGHT_CARD)
        s.line.color.rgb = border or (self.BORDER_DK if dark else self.BORDER)
        s.line.width = Pt(1)
        return s

    def _stat_box(self, slide, value, label, left, top, w, h,
                  val_color=None, bg=None):
        self._card(slide, left, top, w, h, fill=bg or self.BLUE_LIGHT)
        self._txt(slide, str(value), left, top + Inches(0.06), w, Inches(0.3),
                  size=Pt(18), color=val_color or self.BLUE, bold=True,
                  align=PP_ALIGN.CENTER)
        self._txt(slide, label.upper(), left, top + Inches(0.36), w, Pt(10),
                  size=Pt(6.5), color=self.SUB_TEXT, bold=True,
                  align=PP_ALIGN.CENTER)

    def _multi_text(self, slide, runs, left, top, width, height, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        for text, size, color, bold in runs:
            r = p.add_run()
            r.text = str(text)
            r.font.size = size
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = self.FONT
        return tf

    def _divider(self, slide, top, color=None):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   self.ML, top, self.CW, Pt(1))
        s.fill.solid()
        s.fill.fore_color.rgb = color or self.BORDER
        s.line.fill.background()

    def _bullet_list(self, slide, items, left, top, width, height,
                     size=Pt(9), color=None, spacing=Pt(2)):
        """Multi-line bullet list in a single textbox."""
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"→ {item}"
            p.font.size = size
            p.font.color.rgb = color or self.SUB_TEXT
            p.font.name = self.FONT
            p.space_after = spacing
        return tf

    def _bar_rect(self, slide, left, top, w, h, fill):
        """Solid rectangle for bar chart."""
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        s.line.fill.background()
        return s

    def _colored_top_line(self, slide, left, top, w, color):
        """Thin colored line at top of a card."""
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Pt(3))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()

    def _insight_card(self, slide, label, text, left, top, w, h,
                      label_color=None, bg=None, border=None):
        """Card with a bold label header and body text."""
        self._card(slide, left, top, w, h, fill=bg or self.BLUE_LIGHT,
                   border=border or self.BORDER)
        self._txt(slide, label.upper(), left + Inches(0.15), top + Inches(0.08),
                  w - Inches(0.3), Pt(10),
                  size=Pt(7), color=label_color or self.BLUE, bold=True)
        self._txt(slide, text, left + Inches(0.15), top + Inches(0.25),
                  w - Inches(0.3), h - Inches(0.3),
                  size=Pt(8), color=self.SUB_TEXT)

    # ── TABLE HELPER ───────────────────────────────────────

    def _add_table(self, slide, rows, col_widths, left, top, row_height=Inches(0.35),
                   header_bg=None, header_color=None, body_size=Pt(8)):
        """Add a styled table. rows = list of lists. First row is header."""
        nrows = len(rows)
        ncols = len(rows[0])
        total_w = sum(col_widths)
        tbl_shape = slide.shapes.add_table(nrows, ncols, left, top,
                                           Inches(total_w), row_height * nrows)
        tbl = tbl_shape.table
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(cw)

        for ri, row in enumerate(rows):
            for ci, cell_text in enumerate(row):
                cell = tbl.cell(ri, ci)
                cell.text = str(cell_text)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(7) if ri == 0 else body_size
                    p.font.bold = ri == 0
                    p.font.name = self.FONT
                    p.font.color.rgb = (header_color or self.WHITE) if ri == 0 else self.DARK_TEXT
                cell.fill.solid()
                if ri == 0:
                    cell.fill.fore_color.rgb = header_bg or self.DARK_BG
                else:
                    cell.fill.fore_color.rgb = self.LIGHT_CARD if ri % 2 == 1 else RGBColor(0xF8, 0xFA, 0xFC)
        return tbl

    # ══════════════════════════════════════════════════════
    #  SLIDE BUILDERS
    # ══════════════════════════════════════════════════════

    def _slide_01_cover(self):
        s = self._new_slide(dark=True)
        self._section_label(s, 'Executive Proposal · March 2026', top=Inches(0.5), dark=True)

        # Mega title
        tb = s.shapes.add_textbox(self.ML, Inches(1.4), self.CW, Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        for line, color in [('Accelerating the', self.WHITE),
                            ('Schroders Advantage.', self.BLUE)]:
            p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = line
            r.font.size = Pt(44)
            r.font.bold = True
            r.font.color.rgb = color
            r.font.name = self.FONT

        self._txt(s, 'A digitally enabled, insight-led client and advisor platform.\nFrom relationship depth to digital scale.',
                  Inches(2), Inches(3.3), Inches(9), Inches(0.6),
                  size=Pt(13), color=self.MUTED, align=PP_ALIGN.CENTER)

        # 4 pill tags
        tags = [('Advisor Productivity', self.GREEN),
                ('Client Retention & AUM Growth', self.BLUE),
                ('Innovation Agility', self.PURPLE),
                ('No AI Lock-In', self.MUTED)]
        tag_w = Inches(2.6)
        gap = Inches(0.2)
        total = len(tags) * tag_w + (len(tags) - 1) * gap
        tx = (self.SLIDE_W - total) / 2
        for label, col in tags:
            self._card(s, tx, Inches(4.2), tag_w, Inches(0.28),
                      fill=self.CARD_BG, border=self.BORDER_DK, dark=True)
            self._txt(s, label.upper(), tx, Inches(4.22), tag_w, Inches(0.25),
                     size=Pt(7), color=col, bold=True, align=PP_ALIGN.CENTER)
            tx += tag_w + gap

        # 3 hero stats
        stat_w = Inches(3.0)
        gap = Inches(0.3)
        sx = (self.SLIDE_W - 3 * stat_w - 2 * gap) / 2
        for val, label, vcol in [
            ('£148M', 'Value at Stake', self.GREEN),
            ('~6.8×', 'Return on Investment', self.BLUE),
            ('Month 8', 'Breakeven', self.WHITE),
        ]:
            self._card(s, sx, Inches(5.1), stat_w, Inches(0.65),
                      fill=self.CARD_BG, border=self.BORDER_DK, dark=True)
            self._txt(s, val, sx, Inches(5.12), stat_w, Inches(0.3),
                     size=Pt(20), color=vcol, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, label.upper(), sx, Inches(5.44), stat_w, Pt(10),
                     size=Pt(7), color=self.MUTED, bold=True, align=PP_ALIGN.CENTER)
            sx += stat_w + gap

        self._txt(s, 'CONFIDENTIAL · Press → or click to navigate',
                  Inches(3), Inches(6.6), Inches(7), Inches(0.25),
                  size=Pt(9), color=self.MUTED, align=PP_ALIGN.CENTER)

    def _slide_02_where_we_stand(self):
        s = self._new_slide()
        self._section_label(s, 'Context · Where We Stand')
        self._heading(s, 'The Case Is Clear. ', 'The Moment Is Now.')
        self._subtitle(s, '10 workshops · 12 capabilities assessed · Master analysis validated')

        # 3 stat cards
        cw = Inches(3.8)
        gap = Inches(0.25)
        cx = self.ML
        cy = Inches(1.65)
        ch = Inches(1.3)

        # Card 1: Today
        self._card(s, cx, cy, cw, ch)
        self._txt(s, 'TODAY · SCHRODERS WEALTH', cx, cy + Inches(0.08), cw, Pt(10),
                  size=Pt(7), color=self.AMBER, bold=True, align=PP_ALIGN.CENTER)
        self._multi_text(s, [
            ('225', Pt(18), self.DARK_TEXT, True), ('  Client Advisors     ', Pt(9), self.SUB_TEXT, False),
            ('35K+', Pt(18), self.DARK_TEXT, True), ('  Clients', Pt(9), self.SUB_TEXT, False),
        ], cx + Inches(0.2), cy + Inches(0.35), cw - Inches(0.4), Inches(0.4), align=PP_ALIGN.CENTER)
        self._txt(s, '£72B', cx, cy + Inches(0.8), cw, Inches(0.3),
                  size=Pt(22), color=self.PURPLE, bold=True, align=PP_ALIGN.CENTER)
        self._txt(s, 'Assets Under Management', cx, cy + Inches(1.08), cw, Pt(10),
                  size=Pt(8), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)

        cx += cw + gap
        # Card 2: Maturity
        self._card(s, cx, cy, cw, ch)
        self._txt(s, 'ASSESSED · MATURITY', cx, cy + Inches(0.08), cw, Pt(10),
                  size=Pt(7), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)
        self._txt(s, '12', cx, cy + Inches(0.3), cw, Inches(0.35),
                  size=Pt(28), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)
        self._txt(s, 'Core Capabilities', cx, cy + Inches(0.6), cw, Pt(10),
                  size=Pt(8), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
        self._multi_text(s, [
            ('Current  ', Pt(8), self.SUB_TEXT, False),
            ('0.6', Pt(14), self.RED, True), (' / 4.0     ', Pt(8), self.SUB_TEXT, False),
            ('Target  ', Pt(8), self.SUB_TEXT, False),
            ('3.5', Pt(14), self.GREEN, True), (' / 4.0', Pt(8), self.SUB_TEXT, False),
        ], cx + Inches(0.2), cy + Inches(0.85), cw - Inches(0.4), Inches(0.4), align=PP_ALIGN.CENTER)

        cx += cw + gap
        # Card 3: Value at Stake
        self._card(s, cx, cy, cw, ch)
        self._txt(s, 'VALUE AT STAKE', cx, cy + Inches(0.08), cw, Pt(10),
                  size=Pt(7), color=self.GREEN, bold=True, align=PP_ALIGN.CENTER)
        self._txt(s, '£148M', cx, cy + Inches(0.3), cw, Inches(0.35),
                  size=Pt(28), color=self.GREEN, bold=True, align=PP_ALIGN.CENTER)
        self._txt(s, 'Five-Year Value (AI from Year 2)', cx, cy + Inches(0.65), cw, Pt(10),
                  size=Pt(8), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
        self._txt(s, '~6.8×', cx, cy + Inches(0.88), cw, Inches(0.25),
                  size=Pt(18), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)
        self._txt(s, 'Return on Investment', cx, cy + Inches(1.1), cw, Pt(10),
                  size=Pt(8), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)

        # UK Wealth Peers card
        py = Inches(3.15)
        self._card(s, self.ML, py, self.CW, Inches(2.6))
        self._txt(s, 'UK WEALTH PEERS — AUM GROWTH & DIGITAL INVESTMENT', self.ML + Inches(0.2), py + Inches(0.08),
                  self.CW - Inches(0.4), Pt(10), size=Pt(7), color=self.AMBER, bold=True)

        peers = [
            ['Peer', 'AUM Growth', 'Digital Investment'],
            ['JP Morgan WM', '+18%', '$2B+ (global)'],
            ['Quilter', '+13%', '£40M+'],
            ['SJP', '+12%', '£150M+'],
            ['Evelyn Partners', '+9%', 'Post-merger programme'],
            ['Schroders WM', '+7%', 'This proposal'],
        ]
        self._add_table(s, peers, [3.0, 2.0, 3.0], self.ML + Inches(0.2), py + Inches(0.35),
                        row_height=Inches(0.28), header_bg=self.DARK_BG)

        self._insight_card(s, 'Insight',
            'Every peer is investing. The question is not whether, but how quickly Schroders responds.',
            self.ML, Inches(5.95), self.CW, Inches(0.55), label_color=self.AMBER, bg=self.AMBER_LIGHT)

    def _slide_03_platform_vision(self):
        s = self._new_slide()
        self._section_label(s, 'The Platform')
        self._heading(s, 'One Platform. Every Channel. ', 'Every Data Source.')
        self._subtitle(s, 'Not a channel replacement. A platform transformation — connecting advisors, clients, systems and intelligence into a single orchestrated experience.')

        cw = Inches(3.8)
        gap = Inches(0.25)
        cx = self.ML
        cy = Inches(1.75)
        ch = Inches(3.5)

        cols = [
            ('DIGITAL CHANNELS', self.BLUE, [
                'Client portal (web & mobile)',
                'Prospect portal & onboarding',
                'RM Workspace & Digital Assist',
                'Unified employee dashboard',
                'Replaces Temenos e-Services',
            ]),
            ('GRAND CENTRAL', self.GREEN, [
                'T24, Salesforce, Folio connectivity',
                'Configuration-based, not code',
                'API-first, event-driven',
                'Pre-built wealth connectors',
                'Backbase Marketplace',
                'Snowflake integration (NBAs)',
                'Group-wide potential',
            ]),
            ('INTELLIGENCE FABRIC', self.PURPLE, [
                'Multi-model AI (Claude, GPT, Gemini)',
                'Custom Agent Factory',
                'Productised agent library',
                'RAG, observability, auditability',
                'Multi-agent governance',
            ]),
        ]
        for title, color, bullets in cols:
            self._card(s, cx, cy, cw, ch)
            self._colored_top_line(s, cx, cy, cw, color)
            self._txt(s, title, cx + Inches(0.15), cy + Inches(0.15), cw - Inches(0.3), Pt(10),
                      size=Pt(7), color=color, bold=True)
            self._bullet_list(s, bullets, cx + Inches(0.15), cy + Inches(0.4),
                            cw - Inches(0.3), ch - Inches(0.5), size=Pt(8.5))
            cx += cw + gap

        self._insight_card(s, 'Integration Advantage',
            'Replace fragmentation with orchestration. Salesforce sees CRM data. Temenos sees core data. Backbase, through Grand Central, sees everything — and orchestrates it for every channel.',
            self.ML, Inches(5.45), self.CW, Inches(0.65), label_color=self.BLUE)

    def _slide_04_commercial_models(self):
        s = self._new_slide()
        self._section_label(s, 'Commercial Model')
        self._heading(s, 'Enterprise or ', 'Pay As You Grow.')
        self._subtitle(s, 'Same platform. Same delivery. Same value. Two commercial structures — choose the model that fits your commitment profile.')

        hw = Inches(5.85)
        gap = Inches(0.4)
        cy = Inches(1.7)
        ch = Inches(4.0)

        # Option A: Enterprise
        lx = self.ML
        self._card(s, lx, cy, hw, ch, border=self.BLUE)
        self._txt(s, 'OPTION A · ENTERPRISE', lx + Inches(0.2), cy + Inches(0.1),
                  hw - Inches(0.4), Pt(10), size=Pt(8), color=self.BLUE, bold=True)
        self._txt(s, '★ Recommended', lx + hw - Inches(1.5), cy + Inches(0.1),
                  Inches(1.3), Pt(10), size=Pt(7), color=self.GREEN, bold=True, align=PP_ALIGN.RIGHT)
        ent_years = [
            ('Year 1', '£2.5M', 'incl. Marketplace'),
            ('Year 2', '£2.7M', 'incl. Marketplace'),
            ('Year 3', '£3.7M', 'incl. Marketplace'),
            ('Year 4', '£4.1M', 'incl. Marketplace'),
            ('Year 5', '£4.6M', 'incl. Marketplace'),
        ]
        ey = cy + Inches(0.45)
        for yr, price, disc in ent_years:
            self._txt(s, yr, lx + Inches(0.25), ey, Inches(1), Pt(10), size=Pt(8), color=self.SUB_TEXT)
            self._txt(s, price, lx + Inches(1.3), ey, Inches(1), Pt(10), size=Pt(9), color=self.DARK_TEXT, bold=True)
            self._txt(s, disc, lx + Inches(2.5), ey, Inches(1.5), Pt(10), size=Pt(7), color=self.GREEN)
            ey += Inches(0.22)

        self._txt(s, 'Full platform commitment from Day 1. Scaling entry — you pay less upfront, more as value compounds. All AI capabilities available immediately.',
                  lx + Inches(0.25), cy + Inches(1.7), hw - Inches(0.5), Inches(0.5),
                  size=Pt(8), color=self.SUB_TEXT)

        # Enterprise stats
        sby = cy + Inches(2.3)
        sw = Inches(1.6)
        for i, (val, lbl, col) in enumerate([
            ('£17.6M', '5-yr Licence', self.RED),
            ('~6.8×', 'ROI', self.BLUE),
            ('Month 8', 'Breakeven', self.GREEN),
        ]):
            self._stat_box(s, val, lbl, lx + Inches(0.25) + i * (sw + Inches(0.15)),
                          sby, sw, Inches(0.55), val_color=col)

        # Option B: PAG
        rx = lx + hw + gap
        self._card(s, rx, cy, hw, ch, border=self.AMBER)
        self._txt(s, 'OPTION B · PAY AS YOU GROW', rx + Inches(0.2), cy + Inches(0.1),
                  hw - Inches(0.4), Pt(10), size=Pt(8), color=self.AMBER, bold=True)
        pag_years = [
            ('Year 1', '£2.9M', 'incl. Marketplace'),
            ('Year 2', '£3.7M', 'incl. Marketplace'),
            ('Year 3+', '£4.9M', 'incl. Marketplace'),
        ]
        py = cy + Inches(0.45)
        for yr, price, disc in pag_years:
            self._txt(s, yr, rx + Inches(0.25), py, Inches(1), Pt(10), size=Pt(8), color=self.SUB_TEXT)
            self._txt(s, price, rx + Inches(1.3), py, Inches(1), Pt(10), size=Pt(9), color=self.DARK_TEXT, bold=True)
            self._txt(s, disc, rx + Inches(2.5), py, Inches(1.5), Pt(10), size=Pt(7), color=self.AMBER)
            py += Inches(0.22)

        self._txt(s, 'Modular activation. Licence grows as modules go live. AI capabilities unlock progressively with tier upgrades. Agentic Runtime from Year 1.',
                  rx + Inches(0.25), cy + Inches(1.3), hw - Inches(0.5), Inches(0.5),
                  size=Pt(8), color=self.SUB_TEXT)

        sby = cy + Inches(2.3)
        for i, (val, lbl, col) in enumerate([
            ('£21.4M', '5-yr Licence', self.RED),
            ('~5.8×', 'ROI', self.BLUE),
            ('Month 9', 'Breakeven', self.AMBER),
        ]):
            self._stat_box(s, val, lbl, rx + Inches(0.25) + i * (sw + Inches(0.15)),
                          sby, sw, Inches(0.55), val_color=col)

        # Key Insight
        self._insight_card(s, 'Key Insight',
            'Enterprise is £0.4M less in Year 1 and £3.8M less over five years. Every price increase maps to a specific module or connector activation. Nothing is arbitrary.',
            self.ML, Inches(5.9), self.CW, Inches(0.55))

    def _slide_05_investment_value(self):
        s = self._new_slide()
        self._section_label(s, 'Where Value Meets Investment')
        self._heading(s, 'Year-on-Year: Cost, Value & AI Amplification · ', '£148M Return', size=Pt(22))
        self._subtitle(s, 'Licence: from sales pricing. Change costs: Backbase quoted scope. Value: from master analysis (20-tab model). AI pricing: subject to finalised BIC agreement.',
                       top=Inches(1.15))

        # ── Enterprise YoY bar chart (shape-based) ──
        chart_left = self.ML
        chart_top = Inches(1.55)
        chart_w = Inches(7.2)
        chart_h = Inches(2.5)
        self._card(s, chart_left, chart_top, chart_w, chart_h)
        self._txt(s, 'COST VS VALUE — YEAR ON YEAR · ENTERPRISE',
                  chart_left + Inches(0.15), chart_top + Inches(0.08),
                  chart_w - Inches(0.3), Pt(10), size=Pt(7), color=self.DARK_TEXT, bold=True)

        # Bar chart area
        bar_left = chart_left + Inches(0.6)
        bar_top = chart_top + Inches(0.45)
        bar_area_w = Inches(6.2)
        bar_area_h = Inches(1.7)
        max_val = 42.0  # max Y-axis
        group_w = bar_area_w / 5
        bar_w = group_w * 0.15

        # Y-axis labels
        for yv in [0, 10, 20, 30, 40]:
            y_pos = bar_top + bar_area_h * (1 - yv / max_val)
            self._txt(s, f'£{yv}M', chart_left + Inches(0.05), y_pos - Pt(5),
                      Inches(0.5), Pt(10), size=Pt(6), color=self.MUTED, align=PP_ALIGN.RIGHT)

        colors = [self.RED, self.AMBER, self.GREEN, self.PURPLE]
        for gi, d in enumerate(self.YOY_ENT):
            gx = bar_left + gi * group_w + group_w * 0.1
            vals = [d['lic'], d['impl'], d['base'], d['ai']]
            for bi, (v, c) in enumerate(zip(vals, colors)):
                if v <= 0:
                    continue
                bh = bar_area_h * (v / max_val)
                by = bar_top + bar_area_h - bh
                bx = gx + bi * (bar_w + Pt(2))
                self._bar_rect(s, bx, by, bar_w, bh, c)
                # Value label on top of green bar
                if bi == 2:
                    self._txt(s, f'£{v}M', bx - Inches(0.15), by - Pt(10),
                              Inches(0.6), Pt(10), size=Pt(5.5), color=self.GREEN, bold=True,
                              align=PP_ALIGN.CENTER)

            # X-axis label
            self._txt(s, f"{d['year']}\n{d['phase']}", gx, bar_top + bar_area_h + Pt(3),
                      group_w * 0.8, Inches(0.3), size=Pt(6), color=self.SUB_TEXT,
                      align=PP_ALIGN.CENTER)

        # Legend
        leg_y = chart_top + chart_h - Inches(0.22)
        leg_x = bar_left
        for label, color in [('Licence', self.RED), ('Implementation', self.AMBER),
                             ('Base Value', self.GREEN), ('AI Value', self.PURPLE)]:
            self._bar_rect(s, leg_x, leg_y, Pt(8), Pt(8), color)
            self._txt(s, label, leg_x + Pt(11), leg_y - Pt(1), Inches(1), Pt(10),
                      size=Pt(6), color=self.SUB_TEXT)
            leg_x += Inches(1.3)

        # ── Right side cards ──
        rcx = self.ML + Inches(7.5)
        rcw = Inches(4.5)

        # 5-Year Totals
        self._card(s, rcx, chart_top, rcw, Inches(0.85), border=self.BLUE)
        self._txt(s, '5-YEAR TOTALS', rcx, chart_top + Inches(0.05), rcw, Pt(10),
                  size=Pt(7), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)
        self._multi_text(s, [
            ('£21.9M', Pt(22), self.RED, True), ('  ', Pt(10), self.DARK_TEXT, False),
            ('£148M', Pt(22), self.GREEN, True),
        ], rcx, chart_top + Inches(0.25), rcw, Inches(0.3), align=PP_ALIGN.CENTER)
        self._txt(s, 'Total Investment          Total Value', rcx, chart_top + Inches(0.52), rcw, Pt(10),
                  size=Pt(7), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
        self._txt(s, 'Licence £17.6M (Enterprise) + Change £4.3M', rcx, chart_top + Inches(0.67), rcw, Pt(10),
                  size=Pt(7), color=self.MUTED, align=PP_ALIGN.CENTER)

        # AI Contribution
        ai_top = chart_top + Inches(0.95)
        self._card(s, rcx, ai_top, rcw, Inches(0.8), border=self.PURPLE)
        self._txt(s, 'AI CONTRIBUTION', rcx, ai_top + Inches(0.05), rcw, Pt(10),
                  size=Pt(7), color=self.PURPLE, bold=True, align=PP_ALIGN.CENTER)
        self._txt(s, '£18.4M', rcx, ai_top + Inches(0.2), rcw, Inches(0.3),
                  size=Pt(22), color=self.PURPLE, bold=True, align=PP_ALIGN.CENTER)
        self._txt(s, 'incremental AI uplift (12% of total)', rcx, ai_top + Inches(0.48), rcw, Pt(10),
                  size=Pt(7), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
        self._txt(s, '7 agents, Y2-Y5 ramp. Base £129.3M + £18.4M AI amplification.',
                  rcx, ai_top + Inches(0.6), rcw, Pt(10),
                  size=Pt(6.5), color=self.MUTED, align=PP_ALIGN.CENTER)

        # Enterprise advantage
        adv_top = ai_top + Inches(0.9)
        self._card(s, rcx, adv_top, rcw, Inches(0.55), fill=self.GREEN_LIGHT, border=self.GREEN)
        self._txt(s, 'Enterprise advantage: Enterprise saves £3.8M over 5 years vs PAG. AI from Day 1. Scaling entry rewards commitment.',
                  rcx + Inches(0.1), adv_top + Inches(0.08), rcw - Inches(0.2), Inches(0.45),
                  size=Pt(7.5), color=self.DARK_TEXT)

        # Stat bar
        sb_y = Inches(4.2)
        sb_w = Inches(2.9)
        sb_gap = Inches(0.15)
        for i, (val, lbl, col) in enumerate([
            ('~6.8×', 'ROI', self.BLUE),
            ('Month 8', 'Breakeven', self.GREEN),
            ('~£98M', 'Net NPV (8%)', self.DARK_TEXT),
            ('£18.4M', 'AI Uplift', self.PURPLE),
        ]):
            self._stat_box(s, val, lbl, self.ML + i * (sb_w + sb_gap), sb_y,
                          sb_w, Inches(0.5), val_color=col)

        # ── PAG Section ──
        self._divider(s, Inches(4.9), color=self.AMBER)
        self._txt(s, 'OPTION B · PAY AS YOU GROW', self.ML, Inches(5.0), self.CW, Pt(12),
                  size=Pt(8), color=self.AMBER, bold=True)

        # PAG summary table
        pag_rows = [
            ['', 'Year 1', 'Year 2', 'Year 3', 'Year 4', 'Year 5', '5yr Total'],
            ['Licence', '£2.88M', '£3.66M', '£4.93M', '£4.94M', '£4.95M', '£21.4M'],
            ['Value', '£11.0M', '£19.6M', '£31.8M', '£42.5M', '£42.8M', '£148M'],
        ]
        self._add_table(s, pag_rows, [1.2, 1.4, 1.4, 1.4, 1.4, 1.4, 1.4],
                        self.ML, Inches(5.25), row_height=Inches(0.22), header_bg=self.AMBER)

        # PAG stats
        for i, (val, lbl, col) in enumerate([
            ('£25.7M', 'Total Investment', self.RED),
            ('~5.8×', 'ROI', self.BLUE),
            ('Month 9', 'Breakeven', self.AMBER),
        ]):
            self._stat_box(s, val, lbl, self.ML + Inches(8.5) + i * Inches(1.25), Inches(5.2),
                          Inches(1.15), Inches(0.5), val_color=col)

    def _slide_06_pricing_bridge(self):
        s = self._new_slide()
        self._section_label(s, 'Pricing Bridge')
        self._heading(s, 'September 2025 → ', 'March 2026.')
        self._subtitle(s, 'Five months of workshops, scope refinement, and integration validation. Here is what changed and why.')

        # Table 1: Component comparison
        rows1 = [
            ['Component', 'September 2025', 'March 2026', 'Delta', 'What Changed'],
            ['Annual Licence', '£1.68M – £2.95M\nIndicative range', '£2.29M → £4.41M\nEnterprise, ramped Y1→Y5', '+AI\n+SCOPE', 'AI / Agentic use cases added. Scope refined from 31 workshops. 5yr avg £3.2M vs Sept £2.9M.'],
            ['Implementation', '£1.5M – £2.25M\nMVP one-off', '[TBD]\nRe-estimating with Mikheil', '+MP', 'Marketplace vendors (Jumio, ComplyAdvantage, Smarty) confirmed. Same boundaries.'],
            ['Managed Services', '£675K/yr\n4-pax team', '£675K/yr\nUnchanged', '—', 'Stable. Post-MVP ongoing operations & innovation.'],
            ['Marketplace', 'Mentioned, never priced', '£162K Y1 → £140K Y5\n~£0.7M over 5yr', 'NEW', 'Jumio, Smarty, ComplyAdvantage, Twilio, Docusign. Per-vendor: £10K setup + £7.5K/yr maint + per-txn.'],
        ]
        self._add_table(s, rows1, [1.8, 2.5, 2.5, 0.8, 4.5],
                        self.ML, Inches(1.55), row_height=Inches(0.55))

        # Insight cards
        iy = Inches(3.9)
        self._insight_card(s, 'Key Insight',
            'The 5-year blended average of £3.2M/yr (Enterprise) compares to the September top-of-range of £2.95M. The delta is driven by confirmed AI use cases (~£400-500K) and refined scope — not a price increase.',
            self.ML, iy, Inches(5.9), Inches(0.65))
        self._insight_card(s, 'Important Context',
            'Year 1 entry at £2.29M is below the September indicative range (£2.95M upper). You don\'t pay full platform price until you\'re consuming full platform value.',
            self.ML + Inches(6.2), iy, Inches(5.9), Inches(0.65),
            label_color=self.AMBER, bg=self.AMBER_LIGHT)

        # Table 2: What September Did Not Include
        self._txt(s, 'WHAT SEPTEMBER DID NOT INCLUDE', self.ML, Inches(4.75), self.CW, Pt(12),
                  size=Pt(8), color=self.AMBER, bold=True)
        rows2 = [
            ['Component', 'September 2025', 'March 2026', 'Why It\'s Needed'],
            ['CLO / Digital Engage', 'Not included', 'MVP → Phase 2', 'Real-time engagement: banners, push, secure messaging.'],
            ['Custom Connector Slots', 'Not scoped', 'Folio, T24, ContentStack; Snowflake', 'Portfolio data, core banking, content & analytics.'],
            ['Agentic Runtime', 'Not available', 'Licence from Year 1', 'Enables 6 AI use cases worth £59M over 5 years.'],
            ['Digital Assist Licences', 'Not scoped separately', 'Additional to RM Workspace', 'Separate user pool for client-facing digital assist.'],
            ['Marketplace Vendors', 'Mentioned, never priced', '£162K Y1 → £140K Y5\n~£0.7M over 5yr', 'Jumio, Smarty, ComplyAdvantage, Twilio, Docusign. Baked into both models.'],
        ]
        self._add_table(s, rows2, [2.2, 2.2, 2.8, 4.9],
                        self.ML, Inches(4.95), row_height=Inches(0.35))

    def _slide_07_scope(self):
        s = self._new_slide()
        self._section_label(s, 'Delivery Scope')
        self._heading(s, 'What You Get. ', 'Phase by Phase.', size=Pt(26))
        self._subtitle(s, 'Aligned with Jonny Mair & Drewe Noble. Scope refined through 31 journey-level workshops.')

        cw = Inches(2.85)
        gap = Inches(0.15)
        cx = self.ML
        cy = Inches(1.6)
        ch = Inches(4.8)

        phases = [
            ('MVP · INTERNAL RELEASE', 'Foundation', self.BLUE, [
                'Portfolio & Account Reporting',
                'Content Article Integration',
                'Customer Support Workspace',
                'Digital Assist (employee)',
                'CLO Essential',
            ], 'T24 · Fabasoft Folio · ContentStack\nSalesforce CRM · Twilio', 'Sentiment Analysis · Auto Reply'),
            ('ENHANCEMENT #1', 'Advisor Efficiency', self.GREEN, [
                'RM Dashboard (KPIs, AUM)',
                'Salesforce Client 360',
                'Notifications & NBAs',
                'Wealth Planning (Voyant)',
                'Live Chat + Secure Messaging',
                'CLO Premium',
                'Client Profile self-service',
            ], 'SF Marketforce · Snowflake\nOutlook', 'Smart Signals · Outreach\nMeeting Summary'),
            ('ENHANCEMENT #2', 'Digital Onboarding', self.PURPLE, [
                'Prospect portal & onboarding',
                'Private UK (SCIL, SWUSL)',
                'Entity clients + ScoBag, SCAL',
                'Swiss & SG data residency',
                'Multi-party agreement signing',
            ], 'ComplyAdvantage · Jumio\nSmarty · Docusign', 'Doc Processing'),
            ('FUTURE', 'Innovation & Beyond', self.AMBER, [
                'Credit & Lending',
                'Self-Service Investing',
                'Advisory Proposals',
                'Integrated payments',
                'New core + PMS integration',
            ], '', 'Review Prep · Custom Agents\nMulti-Agent Governance'),
        ]
        for title, subtitle, color, bullets, integrations, ai_cases in phases:
            self._card(s, cx, cy, cw, ch)
            self._colored_top_line(s, cx, cy, cw, color)
            self._txt(s, title, cx + Inches(0.1), cy + Inches(0.12), cw - Inches(0.2), Pt(10),
                      size=Pt(6.5), color=color, bold=True)
            self._txt(s, subtitle, cx + Inches(0.1), cy + Inches(0.28), cw - Inches(0.2), Pt(12),
                      size=Pt(10), color=self.DARK_TEXT, bold=True)
            self._bullet_list(s, bullets, cx + Inches(0.1), cy + Inches(0.55),
                            cw - Inches(0.2), Inches(2.0), size=Pt(7))
            if integrations:
                int_y = cy + Inches(2.8)
                self._txt(s, 'INTEGRATIONS', cx + Inches(0.1), int_y, cw - Inches(0.2), Pt(8),
                          size=Pt(6), color=self.AMBER, bold=True)
                self._txt(s, integrations, cx + Inches(0.1), int_y + Inches(0.14),
                          cw - Inches(0.2), Inches(0.4), size=Pt(6.5), color=self.SUB_TEXT)
            if ai_cases:
                ai_y = cy + Inches(3.6)
                self._txt(s, 'AI USE CASES', cx + Inches(0.1), ai_y, cw - Inches(0.2), Pt(8),
                          size=Pt(6), color=self.GREEN, bold=True)
                self._txt(s, ai_cases, cx + Inches(0.1), ai_y + Inches(0.14),
                          cw - Inches(0.2), Inches(0.4), size=Pt(6.5), color=self.SUB_TEXT)
            cx += cw + gap

        self._insight_card(s, 'Capacity Constraint',
            'Post-MVP, Schroders can execute Enhancement #1 OR #2 — not both simultaneously. Sequential model aligns with Enterprise licence ramp.',
            self.ML, Inches(6.6), self.CW, Inches(0.5), label_color=self.AMBER, bg=self.AMBER_LIGHT)

    def _slide_08_ai_advantage(self):
        s = self._new_slide()
        self._section_label(s, 'Intelligence Fabric')
        self._heading(s, 'Intelligent by ', 'Design.', size=Pt(26))
        self._subtitle(s, 'Two paths to AI capability. Build your own with Agent Factory. Activate ours from the productised library.')

        self._txt(s, 'IN SCOPE · 6 VALIDATED USE CASES', self.ML, Inches(1.55), self.CW, Pt(12),
                  size=Pt(7), color=self.DARK_TEXT, bold=True)

        # 6 use case cards (3×2)
        cw = Inches(3.8)
        gap = Inches(0.2)
        cy = Inches(1.8)
        ch = Inches(0.65)

        use_cases = [
            ('Sentiment Analysis', 'NLP dissatisfaction detection', 'MVP', self.BLUE),
            ('Auto Reply', 'Contextual routing & drafting', 'MVP', self.BLUE),
            ('Smart Signals', 'Proactive engagement triggers', 'ENH #1', self.GREEN),
            ('Personalised Outreach', 'Automated nurture sequences', 'ENH #1', self.GREEN),
            ('Meeting Summarisation', 'Automated capture & action items', 'ENH #1', self.GREEN),
            ('Doc Processing', 'Extraction, classification & routing', 'ENH #2', self.PURPLE),
        ]
        for i, (name, desc, phase, color) in enumerate(use_cases):
            col = i % 3
            row = i // 3
            cx = self.ML + col * (cw + gap)
            cy_i = Inches(1.8) + row * (ch + Inches(0.12))
            self._card(s, cx, cy_i, cw, ch, border=color)
            self._txt(s, name, cx + Inches(0.15), cy_i + Inches(0.08), cw - Inches(1.2), Pt(12),
                      size=Pt(10), color=self.DARK_TEXT, bold=True)
            self._txt(s, phase, cx + cw - Inches(0.9), cy_i + Inches(0.1), Inches(0.75), Pt(10),
                      size=Pt(6.5), color=color, bold=True, align=PP_ALIGN.RIGHT)
            self._txt(s, desc, cx + Inches(0.15), cy_i + Inches(0.35), cw - Inches(0.3), Pt(10),
                      size=Pt(8), color=self.SUB_TEXT)

        # Future capabilities
        fy = Inches(3.5)
        self._txt(s, 'SUBSEQUENT PHASES · ADDITIONAL AI CAPABILITIES', self.ML, fy, self.CW, Pt(12),
                  size=Pt(7), color=self.MUTED, bold=True)
        future = [
            ('Review Prep Copilot', '£2.5M/yr manual cost elimination'),
            ('Custom Agents', 'Agent Factory — your IP, your logic'),
            ('Multi-Agent Governance', 'Control layer for all AI agents'),
        ]
        fx = self.ML
        for name, desc in future:
            self._card(s, fx, fy + Inches(0.2), cw, Inches(0.55), border=self.MUTED)
            self._txt(s, name, fx + Inches(0.15), fy + Inches(0.25), cw - Inches(0.3), Pt(12),
                      size=Pt(9), color=self.DARK_TEXT, bold=True)
            self._txt(s, desc, fx + Inches(0.15), fy + Inches(0.45), cw - Inches(0.3), Pt(10),
                      size=Pt(7), color=self.SUB_TEXT)
            fx += cw + gap

        # Bottom cards
        by = Inches(4.55)
        bw = Inches(5.9)
        self._card(s, self.ML, by, bw, Inches(0.85), border=self.CYAN)
        self._txt(s, 'NO AI LOCK-IN', self.ML + Inches(0.15), by + Inches(0.05), bw - Inches(0.3), Pt(10),
                  size=Pt(8), color=self.CYAN, bold=True)
        self._txt(s, 'Multi-model architecture — Claude, GPT, Gemini, open-source. The platform selects the right model for each task. Switch providers without re-engineering agents.',
                  self.ML + Inches(0.15), by + Inches(0.25), bw - Inches(0.3), Inches(0.55),
                  size=Pt(8), color=self.SUB_TEXT)

        self._card(s, self.ML + bw + gap, by, bw, Inches(0.85), border=self.RED)
        self._txt(s, 'BIC CONSUMPTION MODEL', self.ML + bw + gap + Inches(0.15), by + Inches(0.05),
                  bw - Inches(0.3), Pt(10), size=Pt(8), color=self.RED, bold=True)
        self._txt(s, 'Predictable, capped AI spending. Productised agents are low-token and bundled. Complex use cases use BIC blocks — you control the spend.',
                  self.ML + bw + gap + Inches(0.15), by + Inches(0.25), bw - Inches(0.3), Inches(0.55),
                  size=Pt(8), color=self.SUB_TEXT)

    def _slide_09_implementation(self):
        s = self._new_slide()
        self._section_label(s, 'Implementation')
        self._heading(s, 'Beyond the ', 'Licence.')
        self._subtitle(s, 'Implementation, managed services, and marketplace — the full investment picture.')

        cw = Inches(3.8)
        gap = Inches(0.25)
        cx = self.ML
        cy = Inches(1.65)
        ch = Inches(2.8)

        cols = [
            ('MVP IMPLEMENTATION', self.BLUE, '[TBD]', 'Sam / Mikheil re-estimating based on refined scope', [
                'September range: £1.5M – £2.25M',
                'Same scope refined + marketplace vendors',
                'Must stay within boundaries',
            ]),
            ('MANAGED SERVICES', self.GREEN, '£675K/yr', 'Unchanged from September', [
                '4-person dedicated team',
                'Post-MVP ongoing operations',
                'Continuous innovation & optimisation',
                'Platform upgrades & maintenance',
            ]),
            ('MARKETPLACE VENDORS', self.PURPLE, '[TBD]', 'Third-party licensing for onboarding', [
                'Jumio (eID verification)',
                'ComplyAdvantage (AML screening)',
                'Smarty (address lookup)',
                'Docusign (e-signatures)',
                'Twilio (chat & messaging)',
            ]),
        ]
        for title, color, value, sub, bullets in cols:
            self._card(s, cx, cy, cw, ch)
            self._colored_top_line(s, cx, cy, cw, color)
            self._txt(s, title, cx + Inches(0.15), cy + Inches(0.15), cw - Inches(0.3), Pt(10),
                      size=Pt(7), color=color, bold=True)
            self._txt(s, value, cx + Inches(0.15), cy + Inches(0.4), cw - Inches(0.3), Inches(0.3),
                      size=Pt(18), color=self.DARK_TEXT, bold=True)
            self._txt(s, sub, cx + Inches(0.15), cy + Inches(0.72), cw - Inches(0.3), Pt(10),
                      size=Pt(7), color=self.SUB_TEXT)
            self._bullet_list(s, bullets, cx + Inches(0.15), cy + Inches(1.0),
                            cw - Inches(0.3), Inches(1.5), size=Pt(8))
            cx += cw + gap

        # Bottom cards
        by = Inches(4.65)
        bw = Inches(5.9)
        self._insight_card(s, 'Delivery Model',
            'Backbase-led implementation for MVP. Transition to UK delivery partner for ongoing phases. Managed services provides continuity. Grand Central reduces integration complexity.',
            self.ML, by, bw, Inches(0.65))
        self._insight_card(s, 'Key Dependencies',
            'Salesforce (CRM, Marketing, Client 360) · Snowflake (NBAs) · Fabasoft Folio (documents) · T24 Core (portfolio data) · ContentStack (content)',
            self.ML + bw + Inches(0.3), by, bw, Inches(0.65), label_color=self.GREEN, bg=self.GREEN_LIGHT)

    def _slide_10_not_included(self):
        s = self._new_slide()
        self._section_label(s, 'Scope Boundaries')
        self._heading(s, "What's ", 'Not Included.')
        self._subtitle(s, 'These capabilities are available when you\'re ready — they are not in the current commercial proposal.')

        hw = Inches(5.9)
        gap = Inches(0.3)
        cy = Inches(1.65)
        ch = Inches(3.0)

        # Future Capabilities
        self._card(s, self.ML, cy, hw, ch)
        self._txt(s, 'FUTURE CAPABILITIES', self.ML + Inches(0.2), cy + Inches(0.1), hw - Inches(0.4), Pt(10),
                  size=Pt(7), color=self.AMBER, bold=True)
        self._bullet_list(s, [
            'Credit & Lending',
            'Self-Service Investing, Trading & Digital Investing',
            'Elective Corporate Actions',
            'Fully integrated payments',
            'Digital Advice (goal-based planning, advisory proposals)',
            'Digital Sales — Lending (credit & lending origination)',
            'Digital Sales — Advisory (discretionary & advisory workflows)',
        ], self.ML + Inches(0.2), cy + Inches(0.35), hw - Inches(0.4), ch - Inches(0.4))

        # Future Integrations
        rx = self.ML + hw + gap
        self._card(s, rx, cy, hw, ch)
        self._txt(s, 'FUTURE INTEGRATIONS', rx + Inches(0.2), cy + Inches(0.1), hw - Inches(0.4), Pt(10),
                  size=Pt(7), color=self.PURPLE, bold=True)
        self._bullet_list(s, [
            'Core Connector (future core, post-migration)',
            'PMS Connector (future portfolio management system)',
            'Portfolio reporting rewired to new core',
            'Onboarding STP with new core',
            'Open additional deposit / XO accounts',
            'Feedzai (payment fraud detection)',
        ], rx + Inches(0.2), cy + Inches(0.35), hw - Inches(0.4), ch - Inches(0.4))

        # Assumptions row
        ay = Inches(4.85)
        aw = Inches(2.85)
        for i, (val, lbl) in enumerate([
            ('22K', 'Users (growing to 50K)'),
            ('225', 'Client Advisors'),
            ('£72B', 'AUM (current)'),
            ('5yr', 'Contract Term'),
        ]):
            self._stat_box(s, val, lbl, self.ML + i * (aw + Inches(0.2)), ay,
                          aw, Inches(0.55))

        # Note
        self._insight_card(s, 'Note',
            'Licence includes Backbase Managed Hosting. TVP: Discretionary, DFM & Execution only (Advisory is future). All prices based on signing by 31 May 2026.',
            self.ML, Inches(5.6), self.CW, Inches(0.55), label_color=self.AMBER, bg=self.AMBER_LIGHT)

    def _slide_11_next_steps(self):
        s = self._new_slide(dark=True)
        self._section_label(s, 'Next Steps', dark=True)
        self._heading(s, 'From Proposal to ', 'Partnership.', dark=True, size=Pt(32))
        self._subtitle(s, 'A clear path to decision, validation, and implementation.', dark=True)

        timeline = [
            ('March 31', 'Decision', self.BLUE,
             'Exclusivity decision. Proceed to validation phase with Backbase as preferred partner.'),
            ('Mid-April', 'Exclusivity Starts', self.CYAN,
             '8-week validation period. Model bank walkthrough, POC, legal & commercial in parallel.'),
            ('8 Weeks', 'Validation Phase', self.GREEN,
             'POC: Salesforce + Temenos R20/R21 connectivity proof. Clear exit criteria at each stage.'),
            ('End May', 'Target Signature', self.PURPLE,
             'Contract execution. All validation complete, commercials agreed, legal finalised.'),
            ('H2 2026', 'Implementation', self.AMBER,
             'Phase 1 MVP Foundation kicks off. Backbase-led delivery, transition to UK partner.'),
        ]
        ty = Inches(1.8)
        for date, title, color, desc in timeline:
            # Colored dot
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                     self.ML, ty + Inches(0.04), Inches(0.16), Inches(0.16))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            # Date
            self._txt(s, date, self.ML + Inches(0.25), ty, Inches(1.3), Inches(0.2),
                      size=Pt(9), color=color, bold=True)
            self._txt(s, f'— {title}', self.ML + Inches(1.55), ty, Inches(2.5), Inches(0.2),
                      size=Pt(9), color=self.WHITE, bold=True)
            self._txt(s, desc, self.ML + Inches(4.2), ty, Inches(8), Inches(0.25),
                      size=Pt(8), color=self.MUTED)
            ty += Inches(0.42)

        # Bottom 3-column
        by = Inches(4.3)
        bw = Inches(3.8)
        bgap = Inches(0.25)
        cols = [
            ('EXCLUSIVITY COVERS', ['Model bank walkthrough', 'POC (SF + Temenos)', 'Legal & commercial']),
            ('POC OBJECTIVES', ['Temenos R20/R21 connectivity', 'Salesforce integration', '2-week connector validation']),
            ('CONDITIONS', ['Based on signing by 31 May', 'Free to build any flow', 'No credits or restrictions']),
        ]
        for i, (title, items) in enumerate(cols):
            cx = self.ML + i * (bw + bgap)
            self._card(s, cx, by, bw, Inches(1.3), fill=self.CARD_BG, border=self.BORDER_DK, dark=True)
            self._txt(s, title, cx + Inches(0.15), by + Inches(0.08), bw - Inches(0.3), Pt(10),
                      size=Pt(7), color=self.BLUE, bold=True)
            self._bullet_list(s, items, cx + Inches(0.15), by + Inches(0.35),
                            bw - Inches(0.3), Inches(0.8), size=Pt(8), color=self.MUTED)

    def _slide_12_ai_roi(self):
        s = self._new_slide()
        self._section_label(s, 'Appendix A · AI Value Contribution')
        self._heading(s, 'AI Use Case ', 'ROI Mapping.', size=Pt(24))
        self._subtitle(s, 'Each in-scope AI use case maps to a specific value lever. Total AI contribution: £59M over 5 years on a ~£2.5M investment.')

        # 3 top stats
        sw = Inches(3.8)
        for i, (val, lbl, col) in enumerate([
            ('£59M', 'AI Value Over 5 Years', self.GREEN),
            ('23.6×', 'AI ROI Multiple', self.BLUE),
            ('~£2.5M', 'AI Investment (Incl. in Licence)', self.PURPLE),
        ]):
            self._stat_box(s, val, lbl, self.ML + i * (sw + Inches(0.25)), Inches(1.5),
                          sw, Inches(0.55), val_color=col)

        # 6 use case ROI cards (3×2)
        cw = Inches(3.8)
        gap = Inches(0.2)
        cases = [
            ('Sentiment Analysis', 'AUM Retention', '£12M', 'MVP', self.BLUE),
            ('Auto Reply', 'Ops Efficiency', '£8M', 'MVP', self.BLUE),
            ('Smart Signals', 'AUM Retention', '£12M', 'ENH #1', self.GREEN),
            ('Personalised Outreach', 'Revenue Uplift', '£15M', 'ENH #1', self.GREEN),
            ('Meeting Summarisation', 'CA Productivity', '£18M', 'ENH #1', self.GREEN),
            ('Doc Processing', 'Ops Efficiency', '£8M', 'ENH #2', self.PURPLE),
        ]
        for i, (name, lever, value, phase, color) in enumerate(cases):
            col = i % 3
            row = i // 3
            cx = self.ML + col * (cw + gap)
            cy = Inches(2.3) + row * (Inches(1.0) + Inches(0.15))
            self._card(s, cx, cy, cw, Inches(1.0), border=color)
            self._txt(s, name, cx + Inches(0.15), cy + Inches(0.08), cw - Inches(1.2), Pt(12),
                      size=Pt(10), color=self.DARK_TEXT, bold=True)
            self._txt(s, phase, cx + cw - Inches(0.85), cy + Inches(0.1), Inches(0.7), Pt(10),
                      size=Pt(6.5), color=color, bold=True, align=PP_ALIGN.RIGHT)
            self._txt(s, f'Value Lever: {lever}', cx + Inches(0.15), cy + Inches(0.38), cw - Inches(0.3), Pt(10),
                      size=Pt(8), color=self.SUB_TEXT)
            self._txt(s, value, cx + Inches(0.15), cy + Inches(0.6), cw - Inches(0.3), Inches(0.3),
                      size=Pt(18), color=color, bold=True)

        self._insight_card(s, 'Note',
            'AI contributions shown are cumulative 5-year values and represent the AI-amplified portion of each value lever. The £59M total is part of the broader £148M platform value at stake.',
            self.ML, Inches(5.7), self.CW, Inches(0.55))

    def _slide_13_product_licences(self):
        s = self._new_slide()
        self._section_label(s, 'Appendix B · Product Licence Map')
        self._heading(s, "What You're ", 'Paying For.', size=Pt(24))
        self._subtitle(s, 'Every licence cost maps to a specific product module and tier. Here is exactly what activates in each phase.')

        rows = [
            ['Module (SKU)', 'MVP (Year 1)', 'Enhancement (Year 2)', 'Scale (Year 3+)'],
            ['Digital Banking — Wealth', 'Premium', '→ Signature', 'Signature'],
            ['Digital Onboarding', '—', '★ Essential', 'Essential'],
            ['Digital Assist', 'Premium', 'Premium', 'Premium'],
            ['CLO', 'Essential', '→ Premium', 'Premium'],
            ['RM Workspace', '—', '★ Premium', '→ Signature'],
            ['Intelligence Fabric', '★ Agentic Runtime', '+ Foundation AI', 'Full AI + Custom'],
            ['Grand Central', 'SF Connector + T24', '+ Snowflake, Marketforce', '+ Core Connector'],
        ]
        self._add_table(s, rows, [3.0, 3.0, 3.0, 3.0],
                        self.ML, Inches(1.55), row_height=Inches(0.35))

        self._insight_card(s, 'Key Insight',
            'Every price increase maps to a specific module or connector activation. Nothing is arbitrary. The ramp from £2.5M → £4.6M reflects scope activation + marketplace, not inflation.',
            self.ML, Inches(4.6), self.CW, Inches(0.55))

        # Legend row
        ly = Inches(5.35)
        lw = Inches(3.8)
        lgap = Inches(0.25)

        # Legend
        self._card(s, self.ML, ly, lw, Inches(0.85))
        self._txt(s, 'LEGEND', self.ML + Inches(0.15), ly + Inches(0.05), lw - Inches(0.3), Pt(10),
                  size=Pt(7), color=self.BLUE, bold=True)
        self._txt(s, '★ New module activation\n→ Tier upgrade\n— Not yet active',
                  self.ML + Inches(0.15), ly + Inches(0.22), lw - Inches(0.3), Inches(0.6),
                  size=Pt(8), color=self.SUB_TEXT)

        self._card(s, self.ML + lw + lgap, ly, lw, Inches(0.85))
        self._txt(s, 'TIER DEFINITIONS', self.ML + lw + lgap + Inches(0.15), ly + Inches(0.05),
                  lw - Inches(0.3), Pt(10), size=Pt(7), color=self.GREEN, bold=True)
        self._txt(s, 'Essential — Core OOTB journeys\nPremium — Advanced features + config\nSignature — Full platform + bespoke',
                  self.ML + lw + lgap + Inches(0.15), ly + Inches(0.22), lw - Inches(0.3), Inches(0.6),
                  size=Pt(8), color=self.SUB_TEXT)

        self._card(s, self.ML + 2 * (lw + lgap), ly, lw, Inches(0.85))
        self._txt(s, 'MARKETPLACE (SEPARATE)', self.ML + 2 * (lw + lgap) + Inches(0.15), ly + Inches(0.05),
                  lw - Inches(0.3), Pt(10), size=Pt(7), color=self.PURPLE, bold=True)
        self._txt(s, 'Jumio (eID) · ComplyAdvantage (AML)\nSmarty (Address) · Docusign (e-sign)\nTwilio (Chat & Messaging)',
                  self.ML + 2 * (lw + lgap) + Inches(0.15), ly + Inches(0.22), lw - Inches(0.3), Inches(0.6),
                  size=Pt(8), color=self.SUB_TEXT)

    def _slide_14_delivery_phasing(self):
        s = self._new_slide()
        self._section_label(s, 'Appendix C · Delivery Roadmap')
        self._heading(s, 'Three-Year ', 'Delivery Phasing.', size=Pt(24))
        self._subtitle(s, 'Integrated delivery timeline across all horizons. Phased implementation with foundations running in parallel.')

        # Gantt chart as shapes
        gantt_left = self.ML + Inches(1.8)
        gantt_top = Inches(1.8)
        gantt_w = Inches(10.0)
        row_h = Inches(0.6)
        label_w = Inches(1.6)
        q_w = gantt_w / 12  # 12 quarters

        # Timeline header
        phases_header = [
            ('FIRST — MVP', 0, 4, self.BLUE),
            ('NEXT — FAST FOLLOW', 4, 4, self.GREEN),
            ('LATER — INTEGRATION', 8, 2, self.PURPLE),
            ('INNOVATION', 10, 2, self.AMBER),
        ]
        for label, start, span, color in phases_header:
            x = gantt_left + start * q_w
            w = span * q_w
            self._bar_rect(s, x, gantt_top - Inches(0.25), w, Inches(0.2), color)
            self._txt(s, label, x, gantt_top - Inches(0.25), w, Inches(0.2),
                      size=Pt(6), color=self.WHITE, bold=True, align=PP_ALIGN.CENTER)

        # Quarter labels
        for q in range(12):
            x = gantt_left + q * q_w
            self._txt(s, f'Q{q+1}', x, gantt_top, q_w, Inches(0.15),
                      size=Pt(6), color=self.MUTED, align=PP_ALIGN.CENTER)

        # Gantt rows
        rows = [
            ('Phase 1\nMVP Foundation', [
                (0, 1, self.CYAN, 'Setup'),
                (1, 3, self.BLUE, 'Build'),
                (3, 1, self.GREEN, 'Go-Live'),
            ]),
            ('Phase 2\nAdvisor Efficiency', [
                (4, 1, self.CYAN, 'Setup'),
                (5, 2.5, self.GREEN, 'Build & Deploy'),
                (7, 1, RGBColor(0x7B, 0x2F, 0xFF), 'BAU'),  # faded purple
            ]),
            ('Phase 3\nDigital Onboarding', [
                (6, 1, self.CYAN, 'Setup'),
                (7, 2, self.PURPLE, 'Build'),
                (9, 1, self.GREEN, 'Go-Live'),
            ]),
            ('Foundations\nPlatform & Core', [
                (0, 2.5, self.AMBER, 'Grand Central'),
                (3, 3, self.AMBER, 'Agentic AI'),
                (3, 3, self.AMBER, 'New Connectors'),
            ]),
        ]
        ry = gantt_top + Inches(0.2)
        for label, bars in rows:
            # Row label
            self._txt(s, label, self.ML, ry, label_w, row_h,
                      size=Pt(7), color=self.DARK_TEXT, bold=True)
            # Bars
            for start, span, color, text in bars:
                x = gantt_left + start * q_w + Pt(2)
                w = span * q_w - Pt(4)
                bh = Inches(0.22)
                self._bar_rect(s, x, ry + Inches(0.15), w, bh, color)
                if w > Inches(0.5):
                    self._txt(s, text, x + Pt(4), ry + Inches(0.16), w - Pt(8), bh,
                              size=Pt(5.5), color=self.WHITE, bold=True)
            ry += row_h

        # Legend
        leg_y = ry + Inches(0.3)
        leg_x = self.ML
        for label, color in [('Setup & CxO Workshops', self.CYAN), ('Adopt & Build', self.BLUE),
                             ('Testing & Go-Live', self.GREEN), ('Managed Services (BAU)', self.PURPLE),
                             ('Platform Component', self.AMBER)]:
            self._bar_rect(s, leg_x, leg_y, Pt(10), Pt(10), color)
            self._txt(s, label, leg_x + Pt(14), leg_y - Pt(1), Inches(1.6), Pt(12),
                      size=Pt(7), color=self.SUB_TEXT)
            leg_x += Inches(2.2)

    def _slide_15_contract_structure(self):
        s = self._new_slide()
        self._section_label(s, 'Appendix D · Contracting')
        self._heading(s, 'Contract Structure · ', 'Backbase Master Agreement', size=Pt(24))
        self._subtitle(s, 'Seven modular agreements under a Master Agreement + UK Critical Outsourcing Addendum. Modules shared for initial review marked with ✓.')

        hw = Inches(5.9)
        gap = Inches(0.3)
        cy = Inches(1.65)
        ch = Inches(3.5)

        # Mandatory
        self._card(s, self.ML, cy, hw, ch)
        self._txt(s, 'MANDATORY', self.ML + Inches(0.2), cy + Inches(0.1), hw - Inches(0.4), Pt(12),
                  size=Pt(8), color=self.BLUE, bold=True)

        mandatory = [
            ('A', 'Software Licence', '✓ Shared'),
            ('B', 'Third Party Add-ons', '✓ Shared'),
            ('C', 'Managed Hosting', '✓ Shared'),
            ('E', 'Professional Services', '✓ Shared'),
        ]
        my = cy + Inches(0.45)
        for letter, name, status in mandatory:
            self._card(s, self.ML + Inches(0.15), my, hw - Inches(0.3), Inches(0.4))
            self._txt(s, letter, self.ML + Inches(0.25), my + Inches(0.08), Inches(0.3), Inches(0.25),
                      size=Pt(10), color=self.BLUE, bold=True)
            self._txt(s, name, self.ML + Inches(0.6), my + Inches(0.08), Inches(3), Inches(0.25),
                      size=Pt(10), color=self.DARK_TEXT, bold=True)
            self._txt(s, status, self.ML + hw - Inches(1.2), my + Inches(0.1), Inches(1), Inches(0.25),
                      size=Pt(8), color=self.GREEN, align=PP_ALIGN.RIGHT)
            my += Inches(0.5)

        # Optional
        rx = self.ML + hw + gap
        self._card(s, rx, cy, hw, ch)
        self._txt(s, 'OPTIONAL', rx + Inches(0.2), cy + Inches(0.1), hw - Inches(0.4), Pt(12),
                  size=Pt(8), color=self.AMBER, bold=True)

        optional = [
            ('D', 'Managed Services', 'Phase 3+'),
            ('F', 'Training Services', 'As needed'),
            ('G', 'Grand Central', '✓ Shared'),
        ]
        oy = cy + Inches(0.45)
        for letter, name, status in optional:
            self._card(s, rx + Inches(0.15), oy, hw - Inches(0.3), Inches(0.4))
            self._txt(s, letter, rx + Inches(0.25), oy + Inches(0.08), Inches(0.3), Inches(0.25),
                      size=Pt(10), color=self.AMBER, bold=True)
            self._txt(s, name, rx + Inches(0.6), oy + Inches(0.08), Inches(3), Inches(0.25),
                      size=Pt(10), color=self.DARK_TEXT, bold=True)
            status_color = self.GREEN if '✓' in status else self.MUTED
            self._txt(s, status, rx + hw - Inches(1.2), oy + Inches(0.1), Inches(1), Inches(0.25),
                      size=Pt(8), color=status_color, align=PP_ALIGN.RIGHT)
            oy += Inches(0.5)

        # Bottom note
        self._insight_card(s, 'Note',
            'Backbase Master Agreement + UK Critical Outsourcing Addendum. Modules A, B, C, E, and G already shared with Schroders for initial review. Remaining modules shared as scope is confirmed.',
            self.ML, Inches(5.4), self.CW, Inches(0.6))

    # ══════════════════════════════════════���═══════════════
    #  GENERATE
    # ══════════════════════════════════════════════════════

    def generate(self):
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_W
        self.prs.slide_height = self.SLIDE_H

        self._slide_01_cover()
        self._slide_02_where_we_stand()
        self._slide_03_platform_vision()
        self._slide_04_commercial_models()
        self._slide_05_investment_value()
        self._slide_06_pricing_bridge()
        self._slide_07_scope()
        self._slide_08_ai_advantage()
        self._slide_09_implementation()
        self._slide_10_not_included()
        self._slide_11_next_steps()
        self._slide_12_ai_roi()
        self._slide_13_product_licences()
        self._slide_14_delivery_phasing()
        self._slide_15_contract_structure()

        return self.prs


if __name__ == '__main__':
    out_dir = Path(__file__).resolve().parent.parent / 'Engagement' / 'Schroders Group' / 'Output'
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / 'schroders_commercial_refresh_v3.pptx'

    gen = SchrodersCommercialV3Pptx()
    prs = gen.generate()
    prs.save(str(out_path))
    print(f'✓ Generated {out_path} ({out_path.stat().st_size / 1024:.0f} KB, {len(prs.slides)} slides)')
