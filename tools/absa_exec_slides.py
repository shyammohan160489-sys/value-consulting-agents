#!/usr/bin/env python3
"""
ABSA Executive Summary — New Slides Generator

Generates 7 supplementary slides for the ABSA RFP defense (7-8 April 2026).
Backbase branded. Google Slides compatible (13.333" x 7.5").

Usage:
    python3 tools/absa_exec_slides.py
"""

import sys, os
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from tools.pptx_presenter import PptxPresenter
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


class AbsaExecSlides(PptxPresenter):
    """7 supplementary slides for ABSA RFP defense."""

    def _notes(self, slide, text):
        """Add speaker notes to a slide."""
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = text

    def _icon_card(self, slide, icon, title, body, left, top, w, h,
                   icon_color=None, bg=None, dark=False):
        """Card with emoji icon, bold title, and body text."""
        self._card(slide, left, top, w, h, fill=bg, dark=dark)
        pad = Inches(0.2)
        # Icon
        self._txt(slide, icon, left + pad, top + Inches(0.12), Inches(0.4), Inches(0.35),
                  size=Pt(18), align=PP_ALIGN.LEFT)
        # Title
        self._txt(slide, title, left + pad, top + Inches(0.45), w - pad*2, Inches(0.25),
                  size=Pt(10), bold=True,
                  color=self.WHITE if dark else self.DARK_TEXT)
        # Body
        self._txt(slide, body, left + pad, top + Inches(0.72), w - pad*2, h - Inches(0.8),
                  size=Pt(8),
                  color=self.MUTED if dark else self.SUB_TEXT)

    # ══════════════════════════════════════════════════════
    #  SLIDE 1 — ABSA STRATEGIC CONTEXT
    # ══════════════════════════════════════════════════════

    def _slide_1_strategic_context(self):
        s = self._new_slide(dark=True)
        self._section_label(s, 'ABSA STRATEGIC CONTEXT', dark=True)
        self._heading(s, 'The Retail Imperative: ', accent='Why Now, Why Digital', dark=True)
        self._subtitle(s,
            'ABSA\'s CIB franchise delivers 21% ROE. Retail must close the gap to reach the Group 16% target.',
            dark=True)

        # Top stat boxes row
        stats = [
            ('R115.7bn', 'GROUP REVENUE', self.BLUE),
            ('15%', 'ROE (TARGET: 16%)', self.AMBER),
            ('53.8%', 'COST-TO-INCOME', self.RED),
            ('-300K', 'RETAIL CUSTOMERS LOST', self.RED),
            ('5.4M', 'DIGITALLY ACTIVE (OF 9.3M)', self.GREEN),
        ]
        box_w = Inches(2.2)
        gap = Inches(0.2)
        start_x = self.ML
        y = Inches(1.8)
        for i, (val, lbl, clr) in enumerate(stats):
            x = start_x + i * (box_w + gap)
            self._card(s, x, y, box_w, Inches(0.9), dark=True)
            self._txt(s, val, x, y + Inches(0.08), box_w, Inches(0.35),
                      size=Pt(22), color=clr, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, lbl, x, y + Inches(0.5), box_w, Inches(0.25),
                      size=Pt(7), color=self.MUTED, bold=True, align=PP_ALIGN.CENTER)

        # Three strategic pillars
        pillars = [
            ('Ringfence the Business',
             'Stop the bleed. 300K customers lost to Capitec and FNB. '
             'Digital self-service and superior CX are table stakes.',
             self.RED),
            ('Grow the Business',
             'CEO Kenny Fihla: "claw back lost retail market share." '
             'Segment-driven propositions for GenZ, HNWI, MSME, and Diaspora.',
             self.GREEN),
            ('Future-Proof the Business',
             'New Personal & Private Banking division (June 2025) needs a modern platform. '
             'R16.7bn IT investment must deliver proportional retail ROI.',
             self.BLUE),
        ]
        card_w = Inches(3.83)
        card_gap = Inches(0.26)
        cy = Inches(3.1)
        for i, (title, body, color) in enumerate(pillars):
            cx = self.ML + i * (card_w + card_gap)
            self._card(s, cx, cy, card_w, Inches(2.6), dark=True)
            self._colored_top_line(s, cx, cy, card_w, color)
            self._txt(s, title, cx + Inches(0.2), cy + Inches(0.2),
                      card_w - Inches(0.4), Inches(0.3),
                      size=Pt(13), bold=True, color=self.WHITE)
            self._txt(s, body, cx + Inches(0.2), cy + Inches(0.6),
                      card_w - Inches(0.4), Inches(1.8),
                      size=Pt(9), color=self.MUTED)

        # Source line
        self._txt(s, 'Sources: ABSA Group 2025 Annual Results, ABSA Media Statements March 2026, Moneyweb, PwC SA Major Banks Analysis',
                  self.ML, Inches(5.9), self.CW, Inches(0.3),
                  size=Pt(6), color=RGBColor(0x47, 0x55, 0x69))

        self._footer(s, 1, dark=True)

        self._notes(s,
            "TALKING POINTS:\n"
            "- Open with the numbers. ABSA's CIB is performing at 21% ROE. Retail is dragging the group average to 15%.\n"
            "- The 16% ROE target is a public commitment by CEO Kenny Fihla. It requires retail to step up.\n"
            "- Cost-to-income at 53.8% is 3pp above the SA industry average of 50.7%. That's roughly R3.5bn in excess operating cost.\n"
            "- The 300K customer loss is not hypothetical — it's reported in 2024 results. Capitec now has 21M customers.\n"
            "- The new Personal & Private Banking division (June 2025 reorg) signals ABSA is serious about retail transformation.\n"
            "- This is not a technology conversation. It's a business survival conversation."
        )
        return s

    # ══════════════════════════════════════════════════════
    #  SLIDE 2 — VALUE AT STAKE
    # ══════════════════════════════════════════════════════

    def _slide_2_value_at_stake(self):
        s = self._new_slide(dark=False)
        self._section_label(s, 'VALUE AT STAKE')
        self._heading(s, 'The Cost of Inaction: ', accent='What ABSA Leaves on the Table')
        self._subtitle(s,
            'Every quarter without transformation widens the gap to Capitec and FNB.')

        # Value leakage cards — 2x2 grid + 1
        items = [
            ('Cost-to-Income Gap', '~R3.5bn', 'excess annual operating cost',
             '53.8% vs industry average 50.7% on R115.7bn revenue base. '
             'Each 1pp improvement = R1.16bn in savings.',
             self.RED, self.RED_LIGHT),
            ('Customer Attrition', '~R1.5-2bn', 'estimated annual revenue leakage',
             '300K customers lost per year. Avg revenue per retail customer estimated at R5,000-6,500/yr. '
             'Capitec adds 300K+ customers per month.',
             self.AMBER, self.AMBER_LIGHT),
            ('Digital Adoption Gap', '42%', 'of customers still not digitally active',
             '5.4M digitally active out of 9.3M. Capitec: 12M active on 21M (57%). '
             'Each offline customer costs 5-8x more to serve.',
             self.BLUE, self.BLUE_LIGHT),
            ('ROE Bridge', '~R1.6bn', 'additional earnings needed for 16% target',
             'Current ROE 15%, target 16%. The gap must come from retail productivity — '
             'not CIB, which already delivers 21%.',
             self.PURPLE, self.PURPLE_LIGHT),
            ('NPS & Satisfaction', '15%', 'NPS vs Capitec ~45%',
             'ABSA improved from 8% to 15%, but remains 30pp behind Capitec. '
             'FNB named Best Mobile App 4 years running. Experience is the battleground.',
             self.GREEN, self.GREEN_LIGHT),
        ]

        # Row 1: 3 cards
        card_w = Inches(3.83)
        card_h = Inches(2.2)
        gap = Inches(0.26)
        y1 = self.CT
        for i in range(3):
            item = items[i]
            title, big_val, unit, body, clr, bg = item
            cx = self.ML + i * (card_w + gap)
            self._card(s, cx, y1, card_w, card_h, fill=bg)
            self._colored_top_line(s, cx, y1, card_w, clr)
            self._txt(s, title.upper(), cx + Inches(0.2), y1 + Inches(0.15),
                      card_w - Inches(0.4), Inches(0.2),
                      size=Pt(7), bold=True, color=clr)
            self._txt(s, big_val, cx + Inches(0.2), y1 + Inches(0.35),
                      card_w - Inches(0.4), Inches(0.35),
                      size=Pt(24), bold=True, color=clr)
            self._txt(s, unit, cx + Inches(0.2), y1 + Inches(0.7),
                      card_w - Inches(0.4), Inches(0.2),
                      size=Pt(8), bold=True, color=self.SUB_TEXT)
            self._txt(s, body, cx + Inches(0.2), y1 + Inches(1.0),
                      card_w - Inches(0.4), Inches(1.1),
                      size=Pt(8), color=self.SUB_TEXT)

        # Row 2: 2 cards centered
        y2 = y1 + card_h + Inches(0.25)
        offset_x = self.ML + (card_w + gap) * 0.5
        for i in range(2):
            item = items[3 + i]
            title, big_val, unit, body, clr, bg = item
            cx = offset_x + i * (card_w + gap)
            self._card(s, cx, y2, card_w, card_h, fill=bg)
            self._colored_top_line(s, cx, y2, card_w, clr)
            self._txt(s, title.upper(), cx + Inches(0.2), y2 + Inches(0.15),
                      card_w - Inches(0.4), Inches(0.2),
                      size=Pt(7), bold=True, color=clr)
            self._txt(s, big_val, cx + Inches(0.2), y2 + Inches(0.35),
                      card_w - Inches(0.4), Inches(0.35),
                      size=Pt(24), bold=True, color=clr)
            self._txt(s, unit, cx + Inches(0.2), y2 + Inches(0.7),
                      card_w - Inches(0.4), Inches(0.2),
                      size=Pt(8), bold=True, color=self.SUB_TEXT)
            self._txt(s, body, cx + Inches(0.2), y2 + Inches(1.0),
                      card_w - Inches(0.4), Inches(1.1),
                      size=Pt(8), color=self.SUB_TEXT)

        self._footer(s, 2)

        self._notes(s,
            "TALKING POINTS:\n"
            "- Walk through each card — these are ABSA's own published numbers, not Backbase projections.\n"
            "- The cost-to-income gap is the most tangible: 3pp above industry = ~R3.5bn. Digital self-service directly attacks this.\n"
            "- Customer attrition: 300K/year is conservative. Capitec is adding more than that per MONTH.\n"
            "- The digital adoption gap means 42% of ABSA's customers interact through expensive channels.\n"
            "- The ROE bridge shows this isn't optional — retail must contribute more for the 16% target.\n"
            "- Pause on NPS. ABSA improved. But the gap to Capitec is still 30pp. That's experience, not product.\n"
            "- Close with: 'Every quarter of delay compounds these gaps.'"
        )
        return s

    # ══════════════════════════════════════════════════════
    #  SLIDE 3 — WHY BACKBASE (VS ALTERNATIVES)
    # ══════════════════════════════════════════════════════

    def _slide_3_why_backbase(self):
        s = self._new_slide(dark=False)
        self._section_label(s, 'WHY BACKBASE')
        self._heading(s, 'What ABSA Needs ', accent='at This Scale')
        self._subtitle(s,
            'A transformation of 9.3M customers across multiple segments '
            'demands enterprise maturity, Africa experience, and a platform approach.')

        # Three columns comparing approaches
        col_w = Inches(3.83)
        col_gap = Inches(0.26)
        y = self.CT

        # Column headers
        headers = [
            ('Full-Stack Core Replacement', 'Higher risk, longer timelines',
             self.AMBER_LIGHT, self.AMBER),
            ('Front-End Builder', 'Fast start, limited scale',
             self.PURPLE_LIGHT, self.PURPLE),
            ('Engagement Banking Platform', 'Progressive modernization at scale',
             self.BLUE_LIGHT, self.BLUE),
        ]
        for i, (title, sub, bg, clr) in enumerate(headers):
            cx = self.ML + i * (col_w + col_gap)
            self._card(s, cx, y, col_w, Inches(0.85), fill=bg)
            self._colored_top_line(s, cx, y, col_w, clr)
            self._txt(s, title, cx + Inches(0.15), y + Inches(0.15),
                      col_w - Inches(0.3), Inches(0.25),
                      size=Pt(11), bold=True, color=self.DARK_TEXT)
            self._txt(s, sub, cx + Inches(0.15), y + Inches(0.45),
                      col_w - Inches(0.3), Inches(0.25),
                      size=Pt(8), color=self.SUB_TEXT)

        # Comparison rows
        criteria = [
            ('Integration Approach',
             'Requires core replacement or deep core commitment',
             'API layer on top of existing systems; no core dependency',
             'Engagement layer integrates via Grand Central iPaaS. Keep your core.'),
            ('Time to First Value',
             '18-24+ months to initial delivery',
             '3-6 months for front-end; back-office still disconnected',
             '6-9 months to live (MCB: 6 months). Full front-to-back journeys.'),
            ('Africa Track Record',
             'Greenfield/small bank references (Ethiopia, Zimbabwe)',
             'Early-stage; no Tier-1 bank at ABSA\'s scale',
             'MCB, I&M, 100+ banks globally. Multi-country proven.'),
            ('Platform Breadth',
             'Core + channels tightly coupled',
             'Front-end only. No process orchestration, AI, or employee tools',
             'Servicing + Sales + AI + Employee Assist + Process Orchestration'),
            ('Enterprise Readiness',
             'Enterprise delivery but complex migration path',
             '$3.3M funded startup. No managed services or 24/7 support',
             'Digital Factory model, 24/7 support, proven methodology'),
        ]

        row_h = Inches(0.78)
        ry = y + Inches(1.0)
        for ri, (criterion, col1, col2, col3) in enumerate(criteria):
            row_y = ry + ri * (row_h + Inches(0.05))
            row_bg = self.LIGHT_CARD if ri % 2 == 0 else self.ALT_ROW
            for ci, text in enumerate([col1, col2, col3]):
                cx = self.ML + ci * (col_w + col_gap)
                self._card(s, cx, row_y, col_w, row_h, fill=row_bg)
                # Criterion label
                if ci == 0:
                    self._txt(s, criterion.upper(), cx + Inches(0.15), row_y + Inches(0.05),
                              col_w - Inches(0.3), Inches(0.15),
                              size=Pt(6), bold=True, color=self.BLUE)
                self._txt(s, text, cx + Inches(0.15), row_y + Inches(0.2),
                          col_w - Inches(0.3), row_h - Inches(0.25),
                          size=Pt(7.5), color=self.SUB_TEXT)

        # Highlight the Backbase column with a subtle border
        highlight_y = y
        highlight_h = ry + len(criteria) * (row_h + Inches(0.05)) - y
        hx = self.ML + 2 * (col_w + col_gap)
        border_shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           hx - Pt(2), highlight_y - Pt(2),
                                           col_w + Pt(4), highlight_h + Pt(4))
        border_shape.fill.background()
        border_shape.line.color.rgb = self.BLUE
        border_shape.line.width = Pt(2)

        self._footer(s, 3)

        self._notes(s,
            "TALKING POINTS:\n"
            "- Frame this as 'what does a transformation at ABSA's scale require?' — not 'why are competitors bad.'\n"
            "- Column 1 maps to Intellect Design Arena. Don't name them — let ABSA connect the dots.\n"
            "- Column 2 maps to Plumery. Again, don't name — focus on the capability gap.\n"
            "- Column 3 is Backbase. Emphasize: engagement layer (no core risk), proven Africa references, full platform.\n"
            "- Key message: ABSA cannot afford a core replacement risk (Column 1) or bet on a startup (Column 2).\n"
            "- The Grand Central iPaaS point is critical — ABSA keeps their existing systems, we wrap around them.\n"
            "- MCB 6-month go-live is the most powerful proof point for time-to-value."
        )
        return s

    # ══════════════════════════════════════════════════════
    #  SLIDE 4 — BACKBASE ENTERPRISE MATURITY
    # ══════════════════════════════════════════════════════

    def _slide_4_enterprise_maturity(self):
        s = self._new_slide(dark=True)
        self._section_label(s, 'WHY BACKBASE', dark=True)
        self._heading(s, 'Built for Banks Like ABSA: ',
                      accent='Enterprise at Scale', dark=True)

        # Stat row
        stats_data = [
            ('100+', 'BANKS\nWORLDWIDE', self.BLUE),
            ('5', 'CONTINENTS\nPROVEN', self.CYAN),
            ('60+', 'OOTB\nCAPABILITIES', self.GREEN),
            ('6 MO', 'FIRST\nGO-LIVE', self.AMBER),
            ('24/7', 'MANAGED\nSUPPORT', self.PURPLE),
        ]
        box_w = Inches(2.2)
        gap = Inches(0.2)
        y = Inches(1.7)
        for i, (val, lbl, clr) in enumerate(stats_data):
            x = self.ML + i * (box_w + gap)
            self._card(s, x, y, box_w, Inches(0.95), dark=True)
            self._txt(s, val, x, y + Inches(0.08), box_w, Inches(0.35),
                      size=Pt(24), color=clr, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, lbl, x, y + Inches(0.48), box_w, Inches(0.4),
                      size=Pt(7), color=self.MUTED, bold=True, align=PP_ALIGN.CENTER)

        # Four capability cards
        capabilities = [
            ('Single Engagement Layer',
             'One platform for Retail, SME, Wealth, and Corporate. '
             'No point solutions. No integration spaghetti. '
             'Consistent CX across all segments and channels.',
             self.BLUE),
            ('Progressive Modernization',
             'Start with Phase 1 (Fix the Basics) and go live in 6-9 months. '
             'No rip-and-replace. Wrap around existing core systems. '
             'Each phase builds on the last.',
             self.GREEN),
            ('AI-Native Intelligence',
             'Intelligence Fabric and Customer Lifetime Orchestrator (CLO) are embedded, '
             'not bolted on. Hyper-personalization, next-best-action, churn prevention '
             'from Day 1 of Phase 3.',
             self.PURPLE),
            ('Digital Factory Model',
             'Joint operating model with dedicated Backbase experts. '
             'Governance across technology, security, UX, and data. '
             'Full ownership transfer — no vendor lock-in.',
             self.AMBER),
        ]
        card_w = Inches(2.85)
        card_h = Inches(2.8)
        card_gap = Inches(0.2)
        cy = Inches(3.0)
        for i, (title, body, clr) in enumerate(capabilities):
            cx = self.ML + i * (card_w + card_gap)
            self._card(s, cx, cy, card_w, card_h, dark=True)
            self._colored_top_line(s, cx, cy, card_w, clr)
            self._txt(s, title, cx + Inches(0.2), cy + Inches(0.2),
                      card_w - Inches(0.4), Inches(0.3),
                      size=Pt(11), bold=True, color=self.WHITE)
            self._txt(s, body, cx + Inches(0.2), cy + Inches(0.6),
                      card_w - Inches(0.4), card_h - Inches(0.8),
                      size=Pt(8.5), color=self.MUTED)

        self._footer(s, 4, dark=True)

        self._notes(s,
            "TALKING POINTS:\n"
            "- This slide answers 'can Backbase actually do this at our scale?'\n"
            "- 100+ banks worldwide — not a startup experiment. Name-drop MCB, I&M in Africa.\n"
            "- 60+ OOTB capabilities means ABSA doesn't build commoditized features from scratch.\n"
            "- Progressive modernization is the key differentiator: no big bang, no core replacement.\n"
            "- AI-Native: Intelligence Fabric is already in the platform. ABSA doesn't need a separate AI vendor.\n"
            "- Digital Factory: this is how we've delivered MCB. Joint team, shared governance, ownership transfer.\n"
            "- Close with: 'We don't sell software. We partner on transformation.'"
        )
        return s

    # ══════════════════════════════════════════════════════
    #  SLIDE 5 — BUSINESS CASE OPTION A (DIRECTIONAL)
    # ══════════════════════════════════════════════════════

    def _slide_5_business_case_directional(self):
        s = self._new_slide(dark=False)
        self._section_label(s, 'HIGH-LEVEL BUSINESS CASE')
        self._heading(s, 'Transformation Value: ', accent='Directional Ranges')
        self._subtitle(s,
            'Conservative benchmarks from comparable Backbase deployments. '
            'Subject to joint validation during discovery.')

        # Phase cards — horizontal
        phases = [
            ('Phase 1', 'Fix the Basics', 'Months 0-12',
             [('Cost-to-Income', '-1 to 2pp', '~R1.2-2.3bn savings'),
              ('Digital Adoption', '+15-20%', 'From 58% to 73-78%'),
              ('Onboarding Time', '-60 to 80%', 'STP account opening'),
              ('Contact Centre', '-25 to 35%', 'Call deflection via self-service')],
             self.GREEN, self.GREEN_LIGHT),
            ('Phase 2', 'Amplify Growth', 'Months 12-24',
             [('Retail Revenue', '+3 to 5%', 'On PPB book via segments'),
              ('Customer Attrition', '-20 to 30%', 'Reduce 300K/yr bleed'),
              ('Cross-sell Ratio', '+0.3-0.5x', 'Campaigns & bundles'),
              ('Digital Lending', '+15 to 25%', 'Conversion improvement')],
             self.BLUE, self.BLUE_LIGHT),
            ('Phase 3', 'Next Gen Innovation', 'Months 24-36',
             [('NPS', '+10 to 15pp', 'Close Capitec gap'),
              ('Cost per Interaction', '-40 to 50%', 'AI-driven servicing'),
              ('ROE Contribution', 'Bridge to 16%+', 'Retail pulls its weight'),
              ('Time-to-Market', '-50 to 60%', 'New product launches')],
             self.PURPLE, self.PURPLE_LIGHT),
        ]

        card_w = Inches(3.83)
        card_gap = Inches(0.26)
        cy = self.CT
        card_h = Inches(4.6)

        for pi, (phase, name, timeline, metrics, clr, bg) in enumerate(phases):
            cx = self.ML + pi * (card_w + card_gap)
            self._card(s, cx, cy, card_w, card_h, fill=bg)
            self._colored_top_line(s, cx, cy, card_w, clr)

            # Phase header
            self._txt(s, f'{phase}: {name}', cx + Inches(0.2), cy + Inches(0.15),
                      card_w - Inches(0.4), Inches(0.25),
                      size=Pt(12), bold=True, color=self.DARK_TEXT)
            self._txt(s, timeline, cx + Inches(0.2), cy + Inches(0.4),
                      card_w - Inches(0.4), Inches(0.2),
                      size=Pt(8), color=clr, bold=True)

            # Metric rows
            for mi, (metric, value, detail) in enumerate(metrics):
                my = cy + Inches(0.75) + mi * Inches(0.9)
                # Metric name
                self._txt(s, metric.upper(), cx + Inches(0.2), my,
                          card_w - Inches(0.4), Inches(0.15),
                          size=Pt(6.5), bold=True, color=clr)
                # Value
                self._txt(s, value, cx + Inches(0.2), my + Inches(0.15),
                          card_w - Inches(0.4), Inches(0.3),
                          size=Pt(16), bold=True, color=self.DARK_TEXT)
                # Detail
                self._txt(s, detail, cx + Inches(0.2), my + Inches(0.47),
                          card_w - Inches(0.4), Inches(0.2),
                          size=Pt(7.5), color=self.SUB_TEXT)

        # Caveat
        self._txt(s, 'Note: All ranges based on Backbase deployment benchmarks (MCB, NBB, Nordic Bank). '
                  'Specific targets to be validated during joint discovery phase.',
                  self.ML, Inches(6.5), self.CW, Inches(0.3),
                  size=Pt(6.5), color=self.MUTED)

        self._footer(s, 5)

        self._notes(s,
            "TALKING POINTS:\n"
            "- Position this as 'what we've seen at comparable banks' — not a guarantee.\n"
            "- Phase 1 is about cost reduction and operational efficiency. The C/I gap is the headline.\n"
            "- Phase 2 is about revenue growth. Customer attrition reduction is the most tangible.\n"
            "- Phase 3 is about competitive positioning. NPS and AI-driven cost reduction.\n"
            "- The caveat at the bottom is important — these are benchmarks, not commitments.\n"
            "- If asked for specifics, transition to the next slide (Option B).\n"
            "- Key message: 'We have evidence these ranges are achievable — we've done it before.'"
        )
        return s

    # ══════════════════════════════════════════════════════
    #  SLIDE 6 — BUSINESS CASE OPTION B (SPECIFIC)
    # ══════════════════════════════════════════════════════

    def _slide_6_business_case_specific(self):
        s = self._new_slide(dark=False)
        self._section_label(s, 'HIGH-LEVEL BUSINESS CASE')
        self._heading(s, 'ABSA-Specific Projections: ', accent='Using Your Numbers')
        self._subtitle(s,
            'Based on ABSA Group 2025 published results. '
            'For discussion — subject to joint validation.')

        # Table
        rows = [
            ['Metric', 'Current State', 'Year 1 Target', 'Year 3 Target', 'Basis'],
            ['Cost-to-Income Ratio', '53.8%', '52.5%', '50.5%',
             '1pp = ~R1.16bn savings'],
            ['Digitally Active Customers', '5.4M (58%)', '6.5M (70%)', '7.9M (85%)',
             'MCB: 77% post-Backbase'],
            ['Customer Base (SA)', '9.3M (declining)', '9.5M (stabilize)', '10.2M (grow)',
             'Stop 300K/yr attrition'],
            ['Digital Sales Conversion', '~5-8%', '12-15%', '20-25%',
             'NBB: 5min account opening'],
            ['Net Promoter Score', '15%', '22%', '30%+',
             'Close gap to Capitec ~45%'],
            ['ROE (Group)', '15.0%', '15.5%', '16%+',
             'Bridge the ~R1.6bn gap'],
        ]
        self._add_table(s, rows,
                        col_widths=[2.8, 1.8, 1.8, 1.8, 3.8],
                        left=self.ML, top=self.CT,
                        row_height=Inches(0.45),
                        body_size=Pt(9))

        # Value summary box below table
        summary_y = self.CT + Inches(0.45) * len(rows) + Inches(0.3)
        self._card(s, self.ML, summary_y, self.CW, Inches(1.6),
                   fill=self.BLUE_LIGHT)
        self._colored_top_line(s, self.ML, summary_y, self.CW, self.BLUE)
        self._txt(s, '3-YEAR CUMULATIVE VALUE POTENTIAL',
                  self.ML + Inches(0.3), summary_y + Inches(0.15),
                  self.CW - Inches(0.6), Inches(0.2),
                  size=Pt(8), bold=True, color=self.BLUE)

        # Three value columns
        values = [
            ('R3.8-5.4bn', 'Cost Reduction',
             'C/I improvement from 53.8% to 50.5% over 3 years',
             self.GREEN),
            ('R2.3-3.7bn', 'Revenue Protection',
             'Reducing customer attrition and growing digital sales',
             self.BLUE),
            ('R1.0-1.6bn', 'Revenue Growth',
             'Segment-driven propositions and cross-sell uplift',
             self.PURPLE),
        ]
        val_w = Inches(3.6)
        val_gap = Inches(0.3)
        val_start = self.ML + Inches(0.3)
        for i, (val, label, desc, clr) in enumerate(values):
            vx = val_start + i * (val_w + val_gap)
            vy = summary_y + Inches(0.45)
            self._txt(s, val, vx, vy, val_w, Inches(0.35),
                      size=Pt(20), bold=True, color=clr)
            self._txt(s, label.upper(), vx, vy + Inches(0.35), val_w, Inches(0.2),
                      size=Pt(7), bold=True, color=clr)
            self._txt(s, desc, vx, vy + Inches(0.55), val_w, Inches(0.3),
                      size=Pt(7.5), color=self.SUB_TEXT)

        # Caveat
        self._txt(s, 'Assumptions: Based on ABSA Group 2025 Annual Results (R115.7bn revenue, R62.2bn opex, 9.3M customers, 5.4M digitally active). '
                  'Projections are indicative and subject to joint discovery and validation. All figures in ZAR.',
                  self.ML, Inches(6.6), self.CW, Inches(0.3),
                  size=Pt(6), color=self.MUTED)

        self._footer(s, 6)

        self._notes(s,
            "TALKING POINTS:\n"
            "- This is the 'show me the money' slide. Use ONLY if ABSA asks for specifics.\n"
            "- Every number comes from ABSA's own published results — we're not making these up.\n"
            "- Cost-to-income is the most defensible: industry average is 50.7%, ABSA is at 53.8%.\n"
            "- Revenue protection (attrition) is the most emotionally resonant: 'you're losing R2bn/yr in customers.'\n"
            "- The 3-year cumulative value of R7-11bn is a strong number but needs the caveat.\n"
            "- If challenged: 'These are directional. We propose a 4-week joint discovery to validate.'\n"
            "- The 'Basis' column shows each projection is grounded in either ABSA data or Backbase reference."
        )
        return s

    # ══════════════════════════════════════════════════════
    #  SLIDE 7 — CASE STUDY RELEVANCE
    # ══════════════════════════════════════════════════════

    def _slide_7_case_relevance(self):
        s = self._new_slide(dark=True)
        self._section_label(s, 'PROVEN IN AFRICA & BEYOND', dark=True)
        self._heading(s, 'Why These References ', accent='Matter for ABSA', dark=True)
        self._subtitle(s,
            'Each case study maps directly to an ABSA challenge.',
            dark=True)

        # Three case study cards
        cases = [
            ('MCB (Mauritius)',
             'Phase 1: Fix the Basics',
             [('Challenge', 'Multi-market bank, legacy systems, manual processes'),
              ('ABSA Parallel', 'Pan-African footprint, operational complexity, cost pressure'),
              ('Approach', 'Progressive modernization, 6-month go-live, multi-country rollout'),
              ('Results', '+84% YoY app downloads, 77% digitally active, +266% digital transactions, 80% OOTB')],
             self.GREEN),
            ('NBB (Bahrain)',
             'Phase 2: Amplify Growth',
             [('Challenge', 'Paper-based processes, slow account opening, challenger bank pressure'),
              ('ABSA Parallel', '300K customer attrition, low digital conversion, FNB/Capitec competition'),
              ('Approach', 'End-to-end digital onboarding, segment-driven propositions'),
              ('Results', '$200M deposit increase, +45% mobile deposits, 5-min account creation')],
             self.BLUE),
            ('Tier-1 Nordic Bank',
             'Phase 3: Next Gen Innovation',
             [('Challenge', 'Wealth management limited to HNWI; retail customers underserved'),
              ('ABSA Parallel', 'HNWI, GenZ, Mass Market segments all need tailored digital propositions'),
              ('Approach', 'AI-powered CLO, democratized wealth management, conversational banking'),
              ('Results', 'AI agents deliver 90% of traditional advisor capability across all segments')],
             self.PURPLE),
        ]

        card_w = Inches(3.83)
        card_gap = Inches(0.26)
        card_h = Inches(4.2)
        cy = Inches(1.8)

        for i, (bank, phase, details, clr) in enumerate(cases):
            cx = self.ML + i * (card_w + card_gap)
            self._card(s, cx, cy, card_w, card_h, dark=True)
            self._colored_top_line(s, cx, cy, card_w, clr)

            # Bank name
            self._txt(s, bank, cx + Inches(0.2), cy + Inches(0.15),
                      card_w - Inches(0.4), Inches(0.25),
                      size=Pt(13), bold=True, color=self.WHITE)
            # Phase tag
            self._txt(s, phase.upper(), cx + Inches(0.2), cy + Inches(0.42),
                      card_w - Inches(0.4), Inches(0.2),
                      size=Pt(7), bold=True, color=clr)

            # Detail rows
            for di, (label, text) in enumerate(details):
                dy = cy + Inches(0.7) + di * Inches(0.85)
                self._txt(s, label.upper(), cx + Inches(0.2), dy,
                          card_w - Inches(0.4), Inches(0.15),
                          size=Pt(6.5), bold=True, color=clr)
                self._txt(s, text, cx + Inches(0.2), dy + Inches(0.15),
                          card_w - Inches(0.4), Inches(0.55),
                          size=Pt(8), color=self.MUTED)

        self._footer(s, 7, dark=True)

        self._notes(s,
            "TALKING POINTS:\n"
            "- Each column maps a case study to an ABSA phase. This isn't random — it's intentional.\n"
            "- MCB is the strongest proof point: same African context, similar bank size, multi-market.\n"
            "  Highlight: 6-month go-live, 80% OOTB, 77% digital adoption.\n"
            "- NBB answers the 'can you actually grow customers?' question. $200M deposit growth is tangible.\n"
            "- Nordic Bank shows the AI/innovation vision is real, not vaporware. 90% advisor parity via AI.\n"
            "- Key message: 'We've done each phase at a comparable institution. This is not theoretical.'\n"
            "- If asked for SA-specific references: mention I&M Bank and the broader Africa pipeline."
        )
        return s

    # ══════════════════════════════════════════════════════
    #  GENERATE
    # ══════════════════════════════════════════════════════

    def generate(self, output_path):
        self._init_presentation()
        self._slide_1_strategic_context()
        self._slide_2_value_at_stake()
        self._slide_3_why_backbase()
        self._slide_4_enterprise_maturity()
        self._slide_5_business_case_directional()
        self._slide_6_business_case_specific()
        self._slide_7_case_relevance()
        self.save(output_path)


if __name__ == '__main__':
    out = 'ABSA_Exec_Summary_New_Slides.pptx'
    AbsaExecSlides().generate(out)
