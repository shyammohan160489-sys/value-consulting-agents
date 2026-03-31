#!/usr/bin/env python3
"""
CIH Bank — Proposition Commerciale FR — PPTX Generator

French executive briefing deck matching CIH_Executive_Proposal_FR.html.
14 slides, Google Slides compatible (13.333" x 7.5").

Usage:
    python3 tools/cih_proposal_fr_pptx.py
"""

import sys
sys.path.insert(0, '/Users/shyam/cortex')

from pathlib import Path
from tools.pptx_presenter import PptxPresenter
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ASSETS = Path('/Users/shyam/cortex/Engagement/CIH Bank/assets')
CIH_LOGO = str(ASSETS / 'cih_bank_logo.png')
BLUE_GRADIENT = str(ASSETS / 'blue_gradient_bg.png')

# CIH template colors
HEADER_BOX_FILL = RGBColor(0xCF, 0xE2, 0xF3)
LIGHT_GRAY_BG   = RGBColor(0xF8, 0xFA, 0xFC)


class CIHProposalFR(PptxPresenter):
    """CIH Bank — French Proposition Commerciale executive deck."""

    ML = Inches(0.55)
    CW = Inches(12.2)

    # ── CIH helpers ──────────────────────────────────────

    def _title_bar(self, slide, title, subtitle=''):
        self._txt(slide, title, self.ML, Inches(0.45), self.CW, Inches(0.55),
                  size=Pt(28), bold=True, color=self.DARK_TEXT)
        if subtitle:
            self._txt(slide, subtitle, self.ML, Inches(1.05), self.CW, Inches(0.3),
                      size=Pt(12), color=self.SUB_TEXT)

    def _cih_table(self, slide, rows, col_widths, left, top, row_height=Inches(0.38)):
        return self._add_table(slide, rows, col_widths, left, top,
                               row_height=row_height,
                               header_bg=self.DARK_BG,
                               header_color=self.WHITE,
                               body_size=Pt(10))

    def _appendix_header(self, slide):
        self._txt(slide, 'ANNEXE', self.ML, Inches(0.3), self.CW, Pt(14),
                  size=Pt(8), color=self.MUTED, bold=True)

    # ── Generate all 14 slides ───────────────────────────

    def generate(self, output_path):
        self._init_presentation()
        n = 0

        # ════════════════════════════════════════════════════
        # SLIDE 1: COUVERTURE (dark hero)
        # ════════════════════════════════════════════════════
        n += 1
        s = self._new_slide(dark=True)

        # Pills
        pill_x = Inches(3.2)
        for label in ['PROPOSITION COMMERCIALE', 'MARS 2026', 'CONFIDENTIEL']:
            pw = Inches(2.0) if label == 'PROPOSITION COMMERCIALE' else Inches(1.3)
            self._bar_rect(s, pill_x, Inches(2.0), pw, Inches(0.28),
                           fill=RGBColor(0x1A, 0x3A, 0x5E))
            self._txt(s, label, pill_x, Inches(2.02), pw, Inches(0.26),
                      size=Pt(7), color=RGBColor(0xAA, 0xBB, 0xCC), bold=True,
                      align=PP_ALIGN.CENTER)
            pill_x += pw + Inches(0.15)

        # Title
        self._txt(s, 'Approfondir la Valeur Client',
                  Inches(1.5), Inches(2.65), Inches(10.3), Inches(0.7),
                  size=Pt(44), bold=True, color=self.WHITE, align=PP_ALIGN.CENTER)
        self._multi_text(s, [
            ('avec ', Pt(44), self.WHITE, True),
            ('CIH Bank.', Pt(44), self.GREEN, True),
        ], Inches(1.5), Inches(3.35), Inches(10.3), Inches(0.7),
           align=PP_ALIGN.CENTER)

        # Subtitle
        self._txt(s, 'Managed Hosting \u00b7 Digital Assist Premium \u00b7 Customer Lifecycle Orchestrator',
                  Inches(2.5), Inches(4.15), Inches(8.3), Inches(0.4),
                  size=Pt(13), color=RGBColor(0x64, 0x74, 0x8B), align=PP_ALIGN.CENTER)

        # 3 hero stats
        sx = Inches(3.0)
        for val, lbl, clr in [('500K', 'Utilisateurs wallet\nen attente de produits', self.RED),
                               ('3,5M', 'Utilisateurs actifs ciblés\nd\'ici l\'Année 5', self.BLUE),
                               ('3,6M\u20ac', 'Valeur annuelle\nen jeu', self.GREEN)]:
            self._txt(s, val, sx, Inches(4.85), Inches(2.0), Inches(0.6),
                      size=Pt(30), bold=True, color=clr, align=PP_ALIGN.CENTER)
            self._txt(s, lbl, sx, Inches(5.45), Inches(2.0), Inches(0.5),
                      size=Pt(8), color=RGBColor(0x64, 0x74, 0x8B), align=PP_ALIGN.CENTER)
            sx += Inches(2.6)

        self._footer(s, n, dark=True)

        # ════════════════════════════════════════════════════
        # SLIDE 2: NOTRE PARCOURS COMMUN
        # ════════════════════════════════════════════════════
        n += 1
        s = self._new_slide()
        self._section_label(s, 'NOTRE PARTENARIAT')
        self._title_bar(s, 'Deux ans de construction ensemble.')
        self._txt(s, 'CIH Bank et Backbase ont bâti une fondation bancaire digitale solide. Il est temps d\'approfondir la valeur client.',
                  self.ML, Inches(1.05), self.CW, Inches(0.3), size=Pt(12), color=self.SUB_TEXT)

        # Timeline
        tl_y = Inches(1.7)
        line_y = tl_y + Inches(0.09)
        self._bar_rect(s, Inches(1.5), line_y, Inches(10.3), Pt(2), fill=self.BORDER)
        tl_x = Inches(1.5)
        tl_gap = Inches(3.1)
        timeline_data = [
            ('JUILLET 2024', 'Signature du Contrat', 'Digital Banking Signature\nEngagement de 5 ans', self.GREEN),
            ('OCT-DÉC 2025', 'Ateliers Ignite', '4 ateliers stratégiques\nDéfinition des cas d\'usage', self.GREEN),
            ('JANV 2026', 'Restitution Ignite', '4 piliers définis\nRoadmap priorisée', self.GREEN),
            ('MARS 2026', 'Proposition Commerciale', 'MH + DA + CLO\nApprofondir la valeur client', self.BLUE),
        ]
        for i, (date, title, desc, dot_clr) in enumerate(timeline_data):
            cx = tl_x + tl_gap * i
            # dot
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, cx, tl_y, Inches(0.22), Inches(0.22))
            dot.fill.solid()
            dot.fill.fore_color.rgb = dot_clr
            dot.line.fill.background()
            # text
            self._txt(s, date, cx - Inches(0.5), tl_y + Inches(0.35), Inches(1.5), Pt(12),
                      size=Pt(7), color=self.MUTED, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, title, cx - Inches(0.5), tl_y + Inches(0.55), Inches(1.5), Inches(0.25),
                      size=Pt(10), color=self.DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, desc, cx - Inches(0.5), tl_y + Inches(0.8), Inches(1.5), Inches(0.5),
                      size=Pt(8), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)

        # Current investment card
        cy = Inches(3.5)
        self._card(s, Inches(1.5), cy, Inches(10.3), Inches(0.9),
                   fill=RGBColor(0xEF, 0xF6, 0xFF), border=RGBColor(0xCC, 0xDD, 0xFF))
        self._txt(s, 'Investissement Backbase Actuel', Inches(1.5), cy + Inches(0.08),
                  Inches(10.3), Inches(0.2), size=Pt(10), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
        inv_x = Inches(2.6)
        for val, lbl in [('6', 'modules licenciés'), ('2M+', 'utilisateurs actifs'), ('1,2M\u20ac', 'ARR actuel (A2)')]:
            self._txt(s, val, inv_x, cy + Inches(0.3), Inches(2.0), Inches(0.4),
                      size=Pt(24), bold=True, color=self.BLUE, align=PP_ALIGN.CENTER)
            self._txt(s, lbl, inv_x + Inches(0.0), cy + Inches(0.55), Inches(2.0), Inches(0.25),
                      size=Pt(10), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
            inv_x += Inches(2.7)

        self._footer(s, n)

        # ════════════════════════════════════════════════════
        # SLIDE 3: PRIORITÉS STRATÉGIQUES
        # ════════════════════════════════════════════════════
        n += 1
        s = self._new_slide()
        self._section_label(s, 'ATELIERS IGNITE')
        self._title_bar(s, 'Quatre piliers de la stratégie digitale de CIH.')
        self._txt(s, 'Identifiés au cours de 4 ateliers Ignite avec la direction de CIH (oct–déc 2025).',
                  self.ML, Inches(1.05), self.CW, Inches(0.3), size=Pt(12), color=self.SUB_TEXT)

        pillars = [
            ('Structurer l\'Offre', 'TVP, Family Banking, Pockets, packaging produit pour les jeunes clients et les familles.', self.GREEN, ['TVP', 'Family Banking', 'Pockets']),
            ('Orchestrer la Croissance Client', 'Convertir les utilisateurs wallet en clients multi-produits. Campagnes automatis\u00e9es.', self.BLUE, ['CLO', 'Engage', 'Segmentation']),
            ('\u00c9quiper Banque Directe', 'Outils de vente digitaux, live chat, rendez-vous vid\u00e9o, vente assist\u00e9e.', self.PURPLE, ['Digital Assist', 'Live Chat', 'Vid\u00e9o']),
            ('Digitaliser les Processus', 'Onboarding digital de bout en bout, Coach Financier, Marketplace, cr\u00e9dit immobilier.', self.AMBER, ['Onboarding', 'Marketplace', 'Cr\u00e9dit']),
        ]
        cx = self.ML
        cw = Inches(2.9)
        ct = Inches(1.55)
        ch = Inches(2.4)
        for title, desc, clr, pills in pillars:
            self._card(s, cx, ct, cw, ch)
            self._colored_top_line(s, cx, ct, cw, clr)
            self._txt(s, title, cx + Inches(0.12), ct + Inches(0.15), cw - Inches(0.24), Inches(0.3),
                      size=Pt(11), bold=True, color=self.DARK_TEXT)
            self._txt(s, desc, cx + Inches(0.12), ct + Inches(0.5), cw - Inches(0.24), Inches(0.8),
                      size=Pt(9), color=self.SUB_TEXT)
            # pills
            px = cx + Inches(0.12)
            # Compute light pill fill color
            hx = str(clr)
            r, g, b = int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16)
            pill_fill = RGBColor(min(r + 0xBB, 0xFF), min(g + 0xBB, 0xFF), min(b + 0xBB, 0xFF))
            for pill_text in pills:
                pw = Inches(0.1 + len(pill_text) * 0.065)
                self._bar_rect(s, px, ct + ch - Inches(0.4), pw, Inches(0.2), fill=pill_fill)
                self._txt(s, pill_text.upper(), px, ct + ch - Inches(0.4), pw, Inches(0.2),
                          size=Pt(6), color=clr, bold=True, align=PP_ALIGN.CENTER)
                px += pw + Inches(0.06)
            cx += cw + Inches(0.12)

        # KPI targets bar
        ky = Inches(4.2)
        self._card(s, self.ML, ky, self.CW, Inches(0.7))
        kx = Inches(1.5)
        for val, lbl, clr in [('500K', 'Utilisateurs wallet à convertir', self.BLUE),
                               ('25%', 'Origination produit digitale', self.PURPLE),
                               ('10K → 2K', 'Clients par agence', self.GREEN)]:
            self._txt(s, 'OBJECTIF KPI', kx, ky + Inches(0.06), Inches(3.2), Pt(10),
                      size=Pt(7), color=self.MUTED, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, val, kx, ky + Inches(0.2), Inches(3.2), Inches(0.25),
                      size=Pt(16), bold=True, color=clr, align=PP_ALIGN.CENTER)
            self._txt(s, lbl, kx, ky + Inches(0.46), Inches(3.2), Inches(0.2),
                      size=Pt(8), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
            kx += Inches(3.5)

        self._footer(s, n)

        # ════════════════════════════════════════════════════
        # SLIDE 4: L'OPPORTUNITÉ
        # ════════════════════════════════════════════════════
        n += 1
        s = self._new_slide()
        self._section_label(s, 'L\'OPPORTUNITÉ')
        self._title_bar(s, 'Trois forces convergentes créent une fenêtre unique.')
        self._txt(s, 'Le moment n\'a jamais été aussi propice pour approfondir le moteur de valeur client de CIH.',
                  self.ML, Inches(1.05), self.CW, Inches(0.3), size=Pt(12), color=self.SUB_TEXT)

        opps = [
            ('UTILISATEURS WALLET', '500K', 'Clients utilisant uniquement le wallet gratuit',
             'CIH intègre de nouveaux clients avec succès, mais la plupart n\'utilisent que le wallet car il est gratuit. Ils ne font pas de CIH leur banque principale — une opportunité massive d\'origination de produits.',
             self.RED),
            ('ORIGINATION PRODUIT', '25%', 'Objectif de vente digitale de produits',
             'CIH souhaite souscrire comptes d\'épargne, prêts personnels, cartes et assurances en digital — en convertissant les utilisateurs wallet en clients multi-produits.',
             self.BLUE),
            ('PLATEFORME PRÊTE', 'Juil. \'26', 'Managed Hosting Opérationnel',
             'La migration vers le Managed Hosting Backbase s\'achève en juillet, créant l\'infrastructure de niveau entreprise nécessaire pour exécuter CLO et Digital Assist à grande échelle.',
             self.GREEN),
        ]
        ox = self.ML
        ow = Inches(3.9)
        ot = Inches(1.55)
        oh = Inches(3.5)
        for label, val, subtitle, desc, clr in opps:
            self._card(s, ox, ot, ow, oh)
            self._colored_top_line(s, ox, ot, ow, clr)
            self._txt(s, label, ox + Inches(0.15), ot + Inches(0.15), ow - Inches(0.3), Pt(12),
                      size=Pt(8), color=clr, bold=True)
            self._txt(s, val, ox + Inches(0.15), ot + Inches(0.4), ow - Inches(0.3), Inches(0.5),
                      size=Pt(32), bold=True, color=clr)
            self._txt(s, subtitle, ox + Inches(0.15), ot + Inches(0.95), ow - Inches(0.3), Inches(0.25),
                      size=Pt(11), bold=True, color=self.DARK_TEXT)
            self._txt(s, desc, ox + Inches(0.15), ot + Inches(1.3), ow - Inches(0.3), Inches(1.8),
                      size=Pt(9), color=self.SUB_TEXT)
            ox += ow + Inches(0.15)

        self._footer(s, n)

        # ════════════════════════════════════════════════════
        # SLIDE 5: VALEUR EN JEU (dark hero)
        # ════════════════════════════════════════════════════
        n += 1
        s = self._new_slide(dark=True)
        self._txt(s, 'VALEUR EN JEU', Inches(1.5), Inches(1.5), Inches(10.3), Pt(16),
                  size=Pt(10), color=RGBColor(0x80, 0x90, 0xAA), bold=True, align=PP_ALIGN.CENTER)
        self._multi_text(s, [
            ('3,6M\u20ac', Pt(48), self.GREEN, True),
            (' par an', Pt(48), self.WHITE, True),
        ], Inches(1.5), Inches(2.0), Inches(10.3), Inches(0.8), align=PP_ALIGN.CENTER)
        self._txt(s, 'de valeur inexploitée.',
                  Inches(1.5), Inches(2.8), Inches(10.3), Inches(0.7),
                  size=Pt(48), bold=True, color=self.WHITE, align=PP_ALIGN.CENTER)

        self._txt(s, 'Origination de produits, revenus de vente croisée, génération de revenus Digital Assist,\nréduction du trafic en agence et efficacité opérationnelle — quantifiés de manière conservatrice.',
                  Inches(2.5), Inches(3.7), Inches(8.3), Inches(0.6),
                  size=Pt(13), color=RGBColor(0x64, 0x74, 0x8B), align=PP_ALIGN.CENTER)

        vx = Inches(2.5)
        for val, lbl, clr in [('1,1M\u20ac', 'Origination produit\n& vente croisée', self.GREEN),
                               ('300K\u20ac', 'Génération de revenus\nDA', self.PURPLE),
                               ('1,88M\u20ac', 'Réduction du trafic\nen agence', self.AMBER),
                               ('300K\u20ac', 'Économies ops\nMH', self.BLUE)]:
            self._txt(s, val, vx, Inches(4.6), Inches(2.0), Inches(0.4),
                      size=Pt(24), bold=True, color=clr, align=PP_ALIGN.CENTER)
            self._txt(s, lbl, vx, Inches(5.05), Inches(2.0), Inches(0.45),
                      size=Pt(8), color=RGBColor(0x64, 0x74, 0x8B), align=PP_ALIGN.CENTER)
            vx += Inches(2.2)

        self._txt(s, 'Méthodologie complète en Annexe →',
                  Inches(3.0), Inches(5.8), Inches(7.3), Inches(0.3),
                  size=Pt(9), color=RGBColor(0x44, 0x55, 0x66), align=PP_ALIGN.CENTER)
        self._footer(s, n, dark=True)

        # ════════════════════════════════════════════════════
        # SLIDE 6: CE QUE VOUS ACQUÉREZ
        # ════════════════════════════════════════════════════
        n += 1
        s = self._new_slide()
        self._section_label(s, 'PÉRIMÈTRE')
        self._title_bar(s, 'Ce que vous acquérez et pourquoi.')
        self._txt(s, 'Trois modules qui s\'ajoutent à votre plateforme Backbase existante. Pas de remplacement — ceci renforce ce qui est déjà en production.',
                  self.ML, Inches(1.05), self.CW, Inches(0.3), size=Pt(12), color=self.SUB_TEXT)

        modules = [
            ('MANAGED HOSTING ENTERPRISE 2', 'Infrastructure de niveau entreprise',
             'Migration de l\'hébergement interne vers le cloud géré Backbase. Capacité 180 RPS, SLA 99,9%, monitoring 24/7.',
             'POURQUOI : Fondation requise pour CLO et DA Premium à grande échelle. Élimine la charge d\'hébergement interne.',
             self.BLUE),
            ('DIGITAL ASSIST PREMIUM', 'Outils de vente pour 50 agents Banque Directe',
             'Passage d\'Essentials à Premium. Ajoute live chat, rendez-vous vidéo, vente assistée avec recommandations CLO.',
             'POURQUOI : Équiper les agents pour générer du revenu, pas seulement gérer le service. Chaque minute gagnée devient du temps de vente.',
             self.PURPLE),
            ('CUSTOMER LIFECYCLE ORCHESTRATOR', 'Moteur d\'origination produit piloté par l\'IA',
             'Identifie quels utilisateurs wallet sont prêts pour quels produits, puis exécute automatiquement des campagnes multi-canal.',
             'POURQUOI : Convertir 500K utilisateurs wallet en clients multi-produits. Le levier de croissance central.',
             self.GREEN),
        ]
        mx = self.ML
        mw = Inches(3.9)
        mt = Inches(1.55)
        mh = Inches(3.3)
        for label, subtitle, desc, why, clr in modules:
            self._card(s, mx, mt, mw, mh)
            self._colored_top_line(s, mx, mt, mw, clr)
            self._txt(s, label, mx + Inches(0.12), mt + Inches(0.12), mw - Inches(0.24), Inches(0.25),
                      size=Pt(8), color=clr, bold=True)
            self._txt(s, subtitle, mx + Inches(0.12), mt + Inches(0.4), mw - Inches(0.24), Inches(0.25),
                      size=Pt(10), color=self.DARK_TEXT, bold=True)
            self._txt(s, desc, mx + Inches(0.12), mt + Inches(0.7), mw - Inches(0.24), Inches(0.8),
                      size=Pt(9), color=self.SUB_TEXT)
            self._divider(s, mt + Inches(1.6))
            self._txt(s, why, mx + Inches(0.12), mt + Inches(1.7), mw - Inches(0.24), Inches(1.2),
                      size=Pt(8), color=clr, bold=True)
            mx += mw + Inches(0.15)

        # Platform note
        self._card(s, self.ML, Inches(5.1), self.CW, Inches(0.45),
                   fill=RGBColor(0xEF, 0xF6, 0xFF), border=RGBColor(0xCC, 0xDD, 0xFF))
        self._txt(s, 'Les trois modules fonctionnent sur la même plateforme Backbase que vous opérez déjà. Aucun nouveau fournisseur, aucune nouvelle intégration, aucune nouvelle formation.',
                  self.ML + Inches(0.15), Inches(5.15), self.CW - Inches(0.3), Inches(0.35),
                  size=Pt(9), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
        self._footer(s, n)

        # ════════════════════════════════════════════════════
        # SLIDE 7: POURQUOI LE CLO
        # ════════════════════════════════════════════════════
        n += 1
        s = self._new_slide()
        self._section_label(s, 'LE MOTEUR DE CROISSANCE')
        self._title_bar(s, 'Le CLO transforme les utilisateurs wallet en clients multi-produits.')
        self._txt(s, 'CIH intègre déjà des clients avec succès. Le défi n\'est pas l\'acquisition — c\'est la conversion vers des produits bancaires complets.',
                  self.ML, Inches(1.05), self.CW, Inches(0.3), size=Pt(12), color=self.SUB_TEXT)

        # 3 flow steps
        steps = [
            ('1. Ingestion des Données', 'Le CLO ingère les données transactionnelles, l\'utilisation de l\'app et les signaux démographiques depuis votre plateforme existante', self.MUTED),
            ('2. L\'IA Identifie l\'Opportunité', 'La segmentation comportementale identifie quels utilisateurs sont prêts pour l\'épargne, les prêts, les cartes ou l\'assurance', self.PURPLE),
            ('3. Exécution Multi-Canal', 'Offres personnalisées via in-app, push, email et canaux assistés par agents. Automatisé, mesuré, optimisé', self.GREEN),
        ]
        fx = self.ML
        fw = Inches(3.6)
        ft = Inches(1.55)
        fh = Inches(1.6)
        for i, (title, desc, clr) in enumerate(steps):
            self._card(s, fx, ft, fw, fh)
            self._colored_top_line(s, fx, ft, fw, clr)
            self._txt(s, title, fx + Inches(0.12), ft + Inches(0.15), fw - Inches(0.24), Inches(0.2),
                      size=Pt(10), bold=True, color=self.DARK_TEXT)
            self._txt(s, desc, fx + Inches(0.12), ft + Inches(0.45), fw - Inches(0.24), Inches(0.9),
                      size=Pt(8), color=self.SUB_TEXT)
            if i < 2:
                self._txt(s, '\u2192', fx + fw + Inches(0.05), ft + Inches(0.5), Inches(0.3), Inches(0.3),
                          size=Pt(18), bold=True, color=self.BLUE, align=PP_ALIGN.CENTER)
            fx += fw + Inches(0.4)

        # Impact cards
        iy = Inches(3.4)
        iw = Inches(5.9)
        ih = Inches(1.0)
        # Green card
        self._card(s, self.ML, iy, iw, ih, fill=self.GREEN_LIGHT, border=RGBColor(0x7D, 0xD3, 0xAB))
        self._txt(s, 'IMPACT PROJETÉ', self.ML + Inches(0.15), iy + Inches(0.08), iw - Inches(0.3), Pt(10),
                  size=Pt(7), color=self.GREEN, bold=True, align=PP_ALIGN.CENTER)
        self._txt(s, '25 000', self.ML + Inches(0.15), iy + Inches(0.25), iw - Inches(0.3), Inches(0.3),
                  size=Pt(22), bold=True, color=self.GREEN, align=PP_ALIGN.CENTER)
        self._txt(s, 'Utilisateurs wallet convertis en clients multi-produits par an (5% de 500K \u00d7 campagnes CLO)',
                  self.ML + Inches(0.15), iy + Inches(0.6), iw - Inches(0.3), Inches(0.3),
                  size=Pt(8), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
        # Blue card
        bx = self.ML + iw + Inches(0.2)
        self._card(s, bx, iy, iw, ih, fill=self.BLUE_LIGHT, border=RGBColor(0x7D, 0xAA, 0xFF))
        self._txt(s, 'IMPACT REVENUS', bx + Inches(0.15), iy + Inches(0.08), iw - Inches(0.3), Pt(10),
                  size=Pt(7), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)
        self._txt(s, '500K\u20ac+', bx + Inches(0.15), iy + Inches(0.25), iw - Inches(0.3), Inches(0.3),
                  size=Pt(22), bold=True, color=self.BLUE, align=PP_ALIGN.CENTER)
        self._txt(s, 'Revenu incrémental annuel de l\'origination produit seule (25K clients \u00d7 200\u20ac revenu moyen)',
                  bx + Inches(0.15), iy + Inches(0.6), iw - Inches(0.3), Inches(0.3),
                  size=Pt(8), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
        self._footer(s, n)

        # ════════════════════════════════════════════════════
        # SLIDE 8: CONDITIONS COMMERCIALES
        # ════════════════════════════════════════════════════
        n += 1
        s = self._new_slide()
        self._section_label(s, 'CONDITIONS COMMERCIALES')
        self._title_bar(s, 'Souscription Co-Termée Juillet 2029.')
        self._txt(s, 'Co-termé avec votre contrat Digital Banking existant (juillet 2029). Trois modules ajoutés à votre plateforme en place.',
                  self.ML, Inches(1.05), self.CW, Inches(0.3), size=Pt(12), color=self.SUB_TEXT)

        rows = [
            ['Module', 'Prix Catalogue', 'Remise', 'Souscription Annuelle'],
            ['Managed Hosting Enterprise 2 (180 RPS)', '\u20ac601 750', '10%', '\u20ac541 575'],
            ['Digital Assist Premium (50 utilisateurs)', '\u20ac350 140', '40%', '\u20ac210 084'],
            ['CLO — Plateforme de Base*', '\u20ac381 000', '30%', '\u20ac266 700'],
            ['CLO — Frais Utilisateur (500K \u00d7 0,70\u20ac)*', '\u20ac354 816', '30%', '\u20ac248 371'],
            ['TOTAL NOUVELLE SOUSCRIPTION ANNUELLE', '', '', '\u20ac1 266 730'],
        ]
        self._cih_table(s, rows, [5.5, 2.2, 1.3, 3.2], self.ML, Inches(1.55))

        # Stat boxes
        sy = Inches(4.0)
        sw = Inches(3.9)
        sh = Inches(0.65)
        self._stat_box(s, '3,8M\u20ac', 'TCV 3 ans', self.ML, sy, sw, sh,
                       val_color=self.GREEN, bg=self.GREEN_LIGHT)
        self._stat_box(s, 'Août \'26 — Juil. \'29', 'DURÉE', self.ML + sw + Inches(0.15), sy, sw, sh,
                       val_color=self.DARK_TEXT)
        self._stat_box(s, '2,8x', 'ROI (SCÉNARIO DE BASE)', self.ML + 2 * (sw + Inches(0.15)), sy, sw, sh,
                       val_color=self.GREEN, bg=self.GREEN_LIGHT)

        # Note
        self._insight_card(s, 'CO-TERMÉ AVEC DIGITAL BANKING',
                          'Les deux contrats se terminent en juillet 2029. Au renouvellement, tout se consolide en un seul accord.',
                          self.ML, Inches(4.9), self.CW, Inches(0.7),
                          label_color=self.BLUE, bg=self.BLUE_LIGHT)
        self._txt(s, '*La remise de 30% sur CLO s\'applique dans le cadre de l\'acquisition des trois modules.',
                  self.ML, Inches(5.65), self.CW, Inches(0.2),
                  size=Pt(7), color=self.MUTED)
        self._footer(s, n)

        # ════════════════════════════════════════════════════
        # SLIDE 9: FEUILLE DE ROUTE (was 11)
        # ════════════════════════════════════════════════════
        n += 1
        s = self._new_slide()
        self._section_label(s, 'FEUILLE DE ROUTE')
        self._title_bar(s, 'De la signature à la création de valeur.')
        self._txt(s, 'Déploiement phasé aligné avec la migration du Managed Hosting et la disponibilité opérationnelle de Banque Directe.',
                  self.ML, Inches(1.05), self.CW, Inches(0.3), size=Pt(12), color=self.SUB_TEXT)

        tl_y = Inches(1.8)
        self._bar_rect(s, Inches(1.0), tl_y + Inches(0.09), Inches(11.3), Pt(2), fill=self.BORDER)
        phases = [
            ('AVR 2026', 'Contrat & Lancement', 'Signature du contrat\nDébut migration MH', self.BLUE, True),
            ('JUIL 2026', 'MH + DA Opérationnels', 'Managed Hosting en production\nDA Premium actif\nBanque Directe équipée', self.MUTED, False),
            ('T4 2026', 'Lancement CLO', 'Plateforme CLO déployée\nPremières campagnes d\'origination\n500K utilisateurs wallet ciblés', self.MUTED, False),
            ('S1 2027', 'Montée en Charge', 'Campagnes de vente croisée\nSegmentation client complète\nTransformation modèle agence', self.MUTED, False),
            ('2028+', 'Pleine Valeur', 'Tous les leviers actifs\nPotentiel de 3,6M\u20ac/an\nOptimisation continue', self.MUTED, False),
        ]
        px = Inches(1.0)
        pgap = Inches(2.45)
        for date, title, desc, clr, is_current in phases:
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, px, tl_y, Inches(0.22), Inches(0.22))
            dot.fill.solid()
            dot.fill.fore_color.rgb = self.BLUE if is_current else self.BORDER
            dot.line.fill.background()
            self._txt(s, date, px - Inches(0.5), tl_y + Inches(0.35), Inches(1.5), Pt(12),
                      size=Pt(7), color=self.BLUE if is_current else self.MUTED, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, title, px - Inches(0.5), tl_y + Inches(0.55), Inches(1.5), Inches(0.25),
                      size=Pt(10), color=self.DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, desc, px - Inches(0.5), tl_y + Inches(0.8), Inches(1.5), Inches(0.8),
                      size=Pt(8), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
            px += pgap

        self._insight_card(s, 'NOTE CLÉ',
                          'Le déploiement CLO s\'appuie sur le contrat de services gérés existant de CIH pour la capacité d\'implémentation. Estimé à 6-8 semaines pour la mise en place de la plateforme CLO une fois le Managed Hosting stabilisé.',
                          self.ML, Inches(4.8), self.CW, Inches(0.7))
        self._footer(s, n)

        # ════════════════════════════════════════════════════
        # SLIDE 10: PROCHAINES ÉTAPES
        # ════════════════════════════════════════════════════
        n += 1
        s = self._new_slide()
        self._section_label(s, 'PLAN DE CLÔTURE MUTUEL')
        self._title_bar(s, 'Prochaines étapes pour avancer ensemble.')
        self._txt(s, 'Calendrier proposé pour s\'aligner avec le planning de migration du Managed Hosting.',
                  self.ML, Inches(1.05), self.CW, Inches(0.3), size=Pt(12), color=self.SUB_TEXT)

        steps = [
            ('SEMAINE 1 — MARS 2026', 'Revue Interne', 'CIH examine la proposition avec les parties prenantes internes. Backbase est disponible pour des appels de clarification.', self.BLUE),
            ('SEMAINE 2-3 — AVRIL 2026', 'Discussion Commerciale', 'Session conjointe pour finaliser les conditions commerciales et le périmètre.', self.PURPLE),
            ('SEMAINE 4-5 — AVRIL 2026', 'Signature du Contrat', 'Revue finale du contrat et signature. Aligné avec le lancement de la migration du Managed Hosting.', self.GREEN),
            ('EN CONTINU', 'Implémentation & Création de Valeur', 'Migration MH (juil.), activation DA Premium, déploiement CLO (T4). Revues trimestrielles de la valeur.', self.AMBER),
        ]
        positions = [(self.ML, Inches(1.55)), (self.ML + Inches(6.2), Inches(1.55)),
                     (self.ML, Inches(3.0)), (self.ML + Inches(6.2), Inches(3.0))]
        cw = Inches(5.9)
        ch = Inches(1.2)
        for i, ((date, title, desc, clr), (cx, cy)) in enumerate(zip(steps, positions)):
            self._card(s, cx, cy, cw, ch)
            # Left border
            self._bar_rect(s, cx, cy, Pt(4), ch, fill=clr)
            self._txt(s, date, cx + Inches(0.2), cy + Inches(0.08), cw - Inches(0.3), Pt(10),
                      size=Pt(7), color=clr, bold=True)
            self._txt(s, title, cx + Inches(0.2), cy + Inches(0.3), cw - Inches(0.3), Inches(0.2),
                      size=Pt(10), bold=True, color=self.DARK_TEXT)
            self._txt(s, desc, cx + Inches(0.2), cy + Inches(0.55), cw - Inches(0.3), Inches(0.5),
                      size=Pt(9), color=self.SUB_TEXT)

        # Validity notice
        self._card(s, self.ML, Inches(4.5), self.CW, Inches(0.4),
                   fill=self.AMBER_LIGHT, border=RGBColor(0xE8, 0xCC, 0x77))
        self._txt(s, 'Cette proposition est valide pour signature avant le 30/04/2026. Les tarifs sont susceptibles de changer après cette date.',
                  self.ML + Inches(0.15), Inches(4.55), self.CW - Inches(0.3), Inches(0.3),
                  size=Pt(9), color=self.AMBER, bold=True, align=PP_ALIGN.CENTER)

        # Closing statement
        self._txt(s, 'La plateforme est prête. Le business case est clair. Approfondissons la valeur client ensemble.',
                  self.ML, Inches(5.1), self.CW, Inches(0.4),
                  size=Pt(15), bold=True, color=self.DARK_TEXT, align=PP_ALIGN.CENTER)
        self._footer(s, n)

        # ════════════════════════════════════════════════════
        # SLIDE 11: ANNEXE — LEVIERS DE VALEUR
        # ════════════════════════════════════════════════════
        n += 1
        s = self._new_slide()
        self._appendix_header(s)
        self._section_label(s, 'MÉTHODOLOGIE DE VALEUR', top=Inches(0.5))
        self._txt(s, 'Cinq leviers de valeur, estimés de manière conservatrice.',
                  self.ML, Inches(0.75), self.CW, Inches(0.35),
                  size=Pt(20), bold=True, color=self.DARK_TEXT)
        self._txt(s, 'Chaque levier est calculé à partir de données spécifiques à CIH et de benchmarks de la banque de détail marocaine.',
                  self.ML, Inches(1.15), self.CW, Inches(0.25), size=Pt(10), color=self.SUB_TEXT)

        levers = [
            ('CLO — ORIGINATION PRODUIT', '500K\u20ac', '500K utilisateurs wallet \u00d7 5% conversion \u00d7 200\u20ac revenu moyen', 'Benchmark : 5% conservateur vs. 8-12% industrie', self.GREEN),
            ('CLO — VENTE CROISÉE', '600K\u20ac', '1,5% augmentation produits-par-client \u00d7 2M base active \u00d7 20\u20ac', 'Benchmark : 1,5% vs. 3-5% CLO mature', self.BLUE),
            ('DA — GÉNÉRATION DE REVENUS', '300K\u20ac', '50 agents BD \u00d7 20% temps vente en plus \u00d7 30K\u20ac/agent/an', 'DA Premium libère 20% du temps administratif', self.PURPLE),
            ('RÉDUCTION TRAFIC AGENCE', '1,88M\u20ac', '25% bascule digitale \u00d7 2,50\u20ac/transaction \u00d7 3M transactions', 'Benchmark : 2,50\u20ac coût transaction agence MENA', self.AMBER),
            ('MH — ÉCONOMIES OPS', '300K\u20ac', '3-4 ETP réaffectés à des missions stratégiques', 'Coût chargé 75-100K\u20ac par ETP infrastructure', self.MUTED),
            ('VALEUR ANNUELLE TOTALE', '3,58M\u20ac', 'Estimation conservatrice sur les 5 leviers', '', self.GREEN),
        ]
        lw = Inches(3.9)
        lh = Inches(1.7)
        lg = Inches(0.12)
        for i, (label, val, calc, bench, clr) in enumerate(levers):
            col = i % 3
            row = i // 3
            lx = self.ML + col * (lw + lg)
            ly = Inches(1.55) + row * (lh + lg)
            self._card(s, lx, ly, lw, lh)
            self._colored_top_line(s, lx, ly, lw, clr)
            self._txt(s, label, lx + Inches(0.12), ly + Inches(0.12), lw - Inches(0.24), Inches(0.2),
                      size=Pt(7), color=clr, bold=True)
            self._txt(s, val, lx + Inches(0.12), ly + Inches(0.35), lw - Inches(0.24), Inches(0.35),
                      size=Pt(22), bold=True, color=clr)
            self._txt(s, calc, lx + Inches(0.12), ly + Inches(0.8), lw - Inches(0.24), Inches(0.4),
                      size=Pt(8), color=self.SUB_TEXT)
            if bench:
                self._txt(s, bench, lx + Inches(0.12), ly + Inches(1.25), lw - Inches(0.24), Inches(0.35),
                          size=Pt(7), color=self.MUTED)
        self._footer(s, n)

        # ════════════════════════════════════════════════════
        # SLIDE 12: ANNEXE — ANALYSE ROI & SENSIBILITÉ
        # ════════════════════════════════════════════════════
        n += 1
        s = self._new_slide()
        self._appendix_header(s)
        self._section_label(s, 'ANALYSE ROI & SENSIBILITÉ', top=Inches(0.5))
        self._txt(s, 'À 50% de réalisation, l\'investissement génère des retours solides.',
                  self.ML, Inches(0.75), self.CW, Inches(0.35),
                  size=Pt(20), bold=True, color=self.DARK_TEXT)
        self._txt(s, 'Trois scénarios basés sur l\'investissement de \u20ac1 266 730/an co-termé à juillet 2029.',
                  self.ML, Inches(1.15), self.CW, Inches(0.25), size=Pt(10), color=self.SUB_TEXT)

        scenarios = [
            ('CONSERVATEUR (25%)', '895K\u20ac', '0,7x', 'Sous le seuil de rentabilité', self.AMBER, self.AMBER_LIGHT),
            ('SCÉNARIO DE BASE (50%)', '1,79M\u20ac', '1,4x', 'Retour solide', self.BLUE, self.BLUE_LIGHT),
            ('PLEIN POTENTIEL (100%)', '3,58M\u20ac', '2,8x', 'Pleine réalisation', self.GREEN, self.GREEN_LIGHT),
        ]
        scx = self.ML
        scw = Inches(3.9)
        sch = Inches(1.8)
        for label, val, roi, status, clr, bg in scenarios:
            self._card(s, scx, Inches(1.55), scw, sch, fill=bg)
            self._colored_top_line(s, scx, Inches(1.55), scw, clr)
            self._txt(s, label, scx, Inches(1.75), scw, Inches(0.2),
                      size=Pt(7), color=clr, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, val, scx, Inches(2.0), scw, Inches(0.4),
                      size=Pt(28), bold=True, color=clr, align=PP_ALIGN.CENTER)
            self._txt(s, 'Valeur annuelle réalisée', scx, Inches(2.45), scw, Inches(0.2),
                      size=Pt(9), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
            self._txt(s, roi, scx, Inches(2.7), scw, Inches(0.25),
                      size=Pt(18), bold=True, color=clr, align=PP_ALIGN.CENTER)
            self._txt(s, 'ROI — ' + status, scx, Inches(2.95), scw, Inches(0.2),
                      size=Pt(8), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
            scx += scw + Inches(0.15)

        # Assumptions card
        ay = Inches(3.6)
        self._card(s, self.ML, ay, self.CW, Inches(2.2))
        self._txt(s, 'HYPOTHÈSES CLÉS & MÉTHODOLOGIE', self.ML + Inches(0.15), ay + Inches(0.08),
                  self.CW - Inches(0.3), Inches(0.2), size=Pt(9), bold=True, color=self.DARK_TEXT)
        assumptions = [
            ('200\u20ac revenu moyen par produit :', 'Mixé entre épargne (120\u20ac), cartes (180\u20ac), prêts personnels (350\u20ac), assurance (150\u20ac).'),
            ('5% de conversion wallet :', 'Conservateur vs. 8-12% benchmarks industrie pour les campagnes ciblées par l\'IA.'),
            ('2,50\u20ac coût par transaction agence :', 'Moyenne banque de détail MENA (McKinsey). Inclut temps guichetier + frais généraux.'),
            ('1,5% augmentation vente croisée :', 'Moitié du benchmark marché mature (3-5%). Déploiement CLO initial.'),
            ('30K\u20ac revenu/agent :', 'Revenu annuel par agent Banque Directe grâce au temps de vente libéré.'),
            ('3-4 ETP économies infrastructure :', 'Équipe hébergement réaffectée. Coût chargé 75-100K\u20ac par ETP.'),
        ]
        for i, (title, desc) in enumerate(assumptions):
            col = i % 2
            row = i // 2
            ax = self.ML + Inches(0.15) + col * Inches(6.0)
            ayt = ay + Inches(0.35) + row * Inches(0.55)
            self._txt(s, title + ' ' + desc, ax, ayt, Inches(5.8), Inches(0.5),
                      size=Pt(8), color=self.SUB_TEXT)
        self._footer(s, n)

        # ════════════════════════════════════════════════════
        # SLIDE 13: ANNEXE — CALENDRIER DE VALEUR
        # ════════════════════════════════════════════════════
        n += 1
        s = self._new_slide()
        self._appendix_header(s)
        self._section_label(s, 'CALENDRIER DE RÉALISATION DE LA VALEUR', top=Inches(0.5))
        self._txt(s, 'Quand chaque levier entre en action.',
                  self.ML, Inches(0.75), self.CW, Inches(0.35),
                  size=Pt(20), bold=True, color=self.DARK_TEXT)
        self._txt(s, 'La valeur se construit progressivement à mesure que chaque module est mis en production.',
                  self.ML, Inches(1.15), self.CW, Inches(0.25), size=Pt(10), color=self.SUB_TEXT)

        rows = [
            ['Levier de Valeur', 'S2 2026', '2027', '2028', '2029-30', 'Plein Potentiel'],
            ['Origination Produit (CLO)', '—', '250K\u20ac', '500K\u20ac', '500K\u20ac', '500K\u20ac'],
            ['Vente Croisée (CLO)', '—', '200K\u20ac', '450K\u20ac', '600K\u20ac', '600K\u20ac'],
            ['Génération de Revenus DA', '75K\u20ac', '225K\u20ac', '300K\u20ac', '300K\u20ac', '300K\u20ac'],
            ['Réduction Trafic Agence', '188K\u20ac', '940K\u20ac', '1,50M\u20ac', '1,88M\u20ac', '1,88M\u20ac'],
            ['Économies Ops MH', '150K\u20ac', '300K\u20ac', '300K\u20ac', '300K\u20ac', '300K\u20ac'],
            ['VALEUR ANNUELLE TOTALE', '413K\u20ac', '1,92M\u20ac', '3,05M\u20ac', '3,58M\u20ac', '3,58M\u20ac'],
        ]
        self._cih_table(s, rows, [3.5, 1.6, 1.6, 1.6, 1.6, 1.7], self.ML, Inches(1.55))

        # Summary
        self._card(s, self.ML, Inches(4.6), self.CW, Inches(0.45),
                   fill=self.GREEN_LIGHT, border=RGBColor(0x7D, 0xD3, 0xAB))
        self._txt(s, 'Valeur cumulée sur 3 ans : 8,96M\u20ac   |   Investissement (co-termé juil. 2029) : 4,44M\u20ac   |   Valeur nette : 4,52M\u20ac',
                  self.ML + Inches(0.15), Inches(4.65), self.CW - Inches(0.3), Inches(0.3),
                  size=Pt(10), bold=True, color=self.GREEN, align=PP_ALIGN.CENTER)

        # Phasing note
        self._insight_card(s, 'LOGIQUE DE PHASAGE',
                          'Les bénéfices MH et DA commencent dès la mise en production (juil. 2026). Les campagnes CLO nécessitent 3-6 mois de configuration, l\'impact matériel débute en 2027. Plein potentiel atteint en 2028-29.',
                          self.ML, Inches(5.2), self.CW, Inches(0.7))
        self._footer(s, n)

        # ════════════════════════════════════════════════════
        # SLIDE 14: ANNEXE — VALEUR VS. INVESTISSEMENT
        # ════════════════════════════════════════════════════
        n += 1
        s = self._new_slide()
        self._appendix_header(s)
        self._section_label(s, 'VALEUR VS. INVESTISSEMENT', top=Inches(0.5))
        self._txt(s, 'Votre investissement est rentabilisé en 12 mois.',
                  self.ML, Inches(0.75), self.CW, Inches(0.35),
                  size=Pt(20), bold=True, color=self.DARK_TEXT)
        self._txt(s, 'Réalisation de la valeur année après année vs. investissement 1,27M\u20ac/an co-termé à juillet 2029.',
                  self.ML, Inches(1.15), self.CW, Inches(0.25), size=Pt(10), color=self.SUB_TEXT)

        rows = [
            ['', 'S2 2026', '2027', '2028', '2029', 'Total'],
            ['Valeur Réalisée', '413K\u20ac', '1,92M\u20ac', '3,05M\u20ac', '3,58M\u20ac', '8,96M\u20ac'],
            ['Investissement', '633K\u20ac', '1,27M\u20ac', '1,27M\u20ac', '1,27M\u20ac', '4,44M\u20ac'],
            ['Valeur Nette Annuelle', '-220K\u20ac', '+650K\u20ac', '+1,78M\u20ac', '+2,31M\u20ac', '+4,52M\u20ac'],
            ['Valeur Nette Cumulée', '-220K\u20ac', '+430K\u20ac', '+2,21M\u20ac', '+4,52M\u20ac', '4,52M\u20ac'],
        ]
        self._cih_table(s, rows, [2.8, 1.8, 1.8, 1.8, 1.8, 2.0], self.ML, Inches(1.55))

        # 3 bottom stat cards
        bsy = Inches(3.8)
        bsw = Inches(3.9)
        bsh = Inches(1.2)
        # Breakeven
        self._card(s, self.ML, bsy, bsw, bsh)
        self._colored_top_line(s, self.ML, bsy, bsw, self.AMBER)
        self._txt(s, 'SEUIL DE RENTABILITÉ', self.ML, bsy + Inches(0.15), bsw, Inches(0.15),
                  size=Pt(7), color=self.AMBER, bold=True, align=PP_ALIGN.CENTER)
        self._txt(s, 'Mois 12', self.ML, bsy + Inches(0.35), bsw, Inches(0.35),
                  size=Pt(26), bold=True, color=self.GREEN, align=PP_ALIGN.CENTER)
        self._txt(s, 'La valeur cumulée dépasse\nle coût cumulé mi-2027', self.ML, bsy + Inches(0.75), bsw, Inches(0.4),
                  size=Pt(8), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
        # Net value
        self._card(s, self.ML + bsw + Inches(0.15), bsy, bsw, bsh, fill=self.GREEN_LIGHT,
                   border=RGBColor(0x7D, 0xD3, 0xAB))
        self._colored_top_line(s, self.ML + bsw + Inches(0.15), bsy, bsw, self.GREEN)
        self._txt(s, 'VALEUR NETTE 3 ANS', self.ML + bsw + Inches(0.15), bsy + Inches(0.15), bsw, Inches(0.15),
                  size=Pt(7), color=self.GREEN, bold=True, align=PP_ALIGN.CENTER)
        self._txt(s, '4,52M\u20ac', self.ML + bsw + Inches(0.15), bsy + Inches(0.35), bsw, Inches(0.35),
                  size=Pt(26), bold=True, color=self.GREEN, align=PP_ALIGN.CENTER)
        self._txt(s, '8,96M\u20ac de valeur créée\ncontre 4,44M\u20ac investis', self.ML + bsw + Inches(0.15), bsy + Inches(0.75), bsw, Inches(0.4),
                  size=Pt(8), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
        # Multiple
        mx = self.ML + 2 * (bsw + Inches(0.15))
        self._card(s, mx, bsy, bsw, bsh, fill=self.BLUE_LIGHT,
                   border=RGBColor(0x7D, 0xAA, 0xFF))
        self._colored_top_line(s, mx, bsy, bsw, self.BLUE)
        self._txt(s, 'MULTIPLE DE VALEUR', mx, bsy + Inches(0.15), bsw, Inches(0.15),
                  size=Pt(7), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)
        self._txt(s, '2,0x', mx, bsy + Inches(0.35), bsw, Inches(0.35),
                  size=Pt(26), bold=True, color=self.BLUE, align=PP_ALIGN.CENTER)
        self._txt(s, 'Chaque 1\u20ac investi génère\n2,00\u20ac de valeur métier', mx, bsy + Inches(0.75), bsw, Inches(0.4),
                  size=Pt(8), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)

        # Closing note
        self._insight_card(s, 'SOUSCRIPTION CO-TERMÉE À JUILLET 2029',
                          '1,27M\u20ac/an co-termé avec votre contrat Digital Banking. La montée en valeur suit le déploiement phasé — MH et DA dès la mise en production, campagnes CLO à partir de 2027.',
                          self.ML, Inches(5.2), self.CW, Inches(0.7))
        self._footer(s, n)

        # ── Save ─────────────────────────────────────────
        self.save(output_path)


if __name__ == '__main__':
    deck = CIHProposalFR()
    deck.generate('/Users/shyam/cortex/Engagement/CIH Bank/Output/CIH_Proposition_Commerciale_FR.pptx')
