"""
Frontline 2026 Presenter — Google Slides-compatible PPTX builder.

Design system: Backbase Unified Frontline 2026
Canvas: 20" x 11.25" (Google Slides widescreen)
All text boxes include 15% width buffer for Google Slides compatibility.
Autofit is disabled on every text frame.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
import os
import json

# ──────────────────────────────────────────────
# Design Tokens
# ──────────────────────────────────────────────

CANVAS_W = Inches(20.0)
CANVAS_H = Inches(11.25)

# Colors
PRIMARY_NAVY = RGBColor(0x00, 0x1C, 0x3D)
ACTION_BLUE = RGBColor(0x1A, 0x5A, 0xFF)
SURFACE_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_GRAY = RGBColor(0xF5, 0xF7, 0xF9)
TEXT_MAIN = RGBColor(0x00, 0x1C, 0x3D)
TEXT_MUTED = RGBColor(0x5C, 0x6E, 0x84)
SUCCESS_GREEN = RGBColor(0x2E, 0xCC, 0x71)

# Fonts
FONT_PRIMARY = "Libre Franklin"
FONT_FALLBACK = "Helvetica"

# Typography sizes (from Master Template 2026 — 4-level hierarchy)
SIZE_H1 = Pt(45)       # Level 2: Slide headings
SIZE_H2 = Pt(24)       # Level 3: Supporting headings
SIZE_BODY = Pt(20)     # Level 4: General content
SIZE_LABEL = Pt(18)    # Level 1: Labels/captions (UPPERCASE)
SIZE_CAPTION = Pt(14)  # Smaller annotations
SIZE_LEGAL = Pt(10)    # Disclaimers
SIZE_AGENDA = Pt(24)   # Agenda items

# Semantic red (from Master Template)
SEMANTIC_RED = RGBColor(0xE0, 0x20, 0x20)

# Margins (safe zone — 0.5" per Master Template)
MARGIN_TOP = Inches(0.5)
MARGIN_BOTTOM = Inches(0.8)
MARGIN_LEFT = Inches(0.5)
MARGIN_RIGHT = Inches(0.5)
SAFE_ZONE = Inches(0.5)

# Google Slides text width buffer
TEXT_WIDTH_BUFFER = 1.15


class Frontline2026Presenter:
    """Generates PPTX files in the Backbase Unified Frontline 2026 style.

    All output is optimized for Google Slides import:
    - No autofit on text frames
    - 15% text width buffer
    - No gradients, shadows, or rotated text
    - Complex shapes grouped
    - 0.75" safe zone from edges
    """

    def __init__(self, output_path: str = None,
                 logo_path: str = None):
        self.prs = Presentation()
        self.prs.slide_width = CANVAS_W
        self.prs.slide_height = CANVAS_H
        self.output_path = output_path or "frontline_2026_output.pptx"
        self.slides = []

        # Resolve logo path — use correct wordmark with notched B
        base = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        if logo_path and os.path.exists(logo_path):
            self.logo_white = logo_path
            self.logo_dark = logo_path
        else:
            self.logo_white = os.path.join(base, "knowledge", "Ignite Inspire",
                                           "brand-assets", "backbase-wordmark-white.png")
            self.logo_dark = os.path.join(base, "knowledge", "Ignite Inspire",
                                          "brand-assets", "backbase-wordmark-dark.png")

    # ──────────────────────────────────────────
    # Core Helpers
    # ──────────────────────────────────────────

    def _add_blank_slide(self, dark=False):
        """Add a blank slide with chrome (top bar, blue accent, footer)."""
        layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(layout)
        self.slides.append(slide)
        self._add_slide_chrome(slide, dark=dark)
        return slide

    def _add_slide_chrome(self, slide, dark=False):
        """Add standard slide chrome: blue accent square + footer.

        Matches the actual Backbase 2026 Google Slides template.
        No top bar — removed per design feedback.
        """
        slide_num = len(self.slides)

        if not dark:
            # Inverted-L accent — matches HTML exactly
            # Shape:   ██        (top-right small block)
            #        ████        (bottom full-width block)
            # Outer bottom-right corner = axis intersection
            accent_left = Inches(0.35)
            accent_top = Inches(0.55)
            block_size = Inches(0.15)

            # Top-right small block
            self._add_rect(
                slide, accent_left + block_size, accent_top,
                block_size, block_size,
                fill_color=ACTION_BLUE, shape_name="BlueAccent_Top"
            )
            # Bottom full-width block
            self._add_rect(
                slide, accent_left, accent_top + block_size,
                block_size * 2, block_size,
                fill_color=ACTION_BLUE, shape_name="BlueAccent_Bottom"
            )

            # Axis lines — originate from bottom-RIGHT corner of the L
            line_color = RGBColor(0xE0, 0xE4, 0xE8)
            corner_x = accent_left + block_size * 2  # right edge of bottom block
            corner_y = accent_top + block_size * 2   # bottom edge of bottom block

            # Vertical axis — from corner downward
            self._add_rect(
                slide, corner_x, corner_y, Inches(0.01), Inches(9.5),
                fill_color=line_color, shape_name="AxisY"
            )
            # Horizontal axis — from corner rightward
            self._add_rect(
                slide, corner_x, corner_y, Inches(18.5), Inches(0.01),
                fill_color=line_color, shape_name="AxisX"
            )

        # Footer: logo on bottom-right, page number to its right
        footer_y = Inches(10.5)

        # Backbase logo — bottom-right (before page number)
        logo_path = self.logo_white if dark else self.logo_dark
        if os.path.exists(logo_path):
            slide.shapes.add_picture(
                logo_path,
                Inches(16.5), footer_y,
                height=Inches(0.35)
            )

        # Page number — far right, next to logo
        num_color = SURFACE_WHITE if dark else TEXT_MUTED
        self._add_textbox(
            slide, Inches(18.5), footer_y, Inches(0.75), Inches(0.4),
            text=f"{slide_num}",
            font_size=SIZE_CAPTION, font_color=num_color,
            alignment=PP_ALIGN.RIGHT, apply_buffer=False,
            shape_name=f"SlideNumber_{slide_num}"
        )

    def _set_slide_bg(self, slide, color: RGBColor):
        """Set solid background color on a slide."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_textbox(self, slide, left, top, width, height,
                     text="", font_name=FONT_PRIMARY, font_size=SIZE_BODY,
                     font_color=TEXT_MAIN, bold=False, italic=False,
                     alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                     apply_buffer=True, shape_name=None):
        """Add a text box with Google Slides compatibility rules applied.

        - Applies 15% width buffer by default
        - Disables autofit
        - Sets word wrap
        """
        if apply_buffer:
            width = int(width * TEXT_WIDTH_BUFFER)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        if shape_name:
            txBox.name = shape_name
        tf = txBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True

        try:
            tf.vertical_anchor = anchor
        except Exception:
            pass

        p = tf.paragraphs[0]
        p.alignment = alignment
        run = p.add_run()
        run.text = text

        font = run.font
        font.name = font_name
        font.size = font_size
        font.color.rgb = font_color
        font.bold = bold
        font.italic = italic

        return txBox, tf

    def _add_multiline_textbox(self, slide, left, top, width, height,
                               lines, font_name=FONT_PRIMARY, font_size=SIZE_BODY,
                               font_color=TEXT_MAIN, bold=False,
                               alignment=PP_ALIGN.LEFT, line_spacing=1.2,
                               apply_buffer=True):
        """Add a text box with multiple lines/paragraphs."""
        if apply_buffer:
            width = int(width * TEXT_WIDTH_BUFFER)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True

        for i, line_text in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.alignment = alignment
            p.space_after = Pt(line_spacing * font_size.pt / 72 * 20)

            run = p.add_run()
            run.text = line_text
            run.font.name = font_name
            run.font.size = font_size
            run.font.color.rgb = font_color
            run.font.bold = bold

        return txBox, tf

    def _add_rounded_rect(self, slide, left, top, width, height,
                          fill_color=ACTION_BLUE, text="",
                          font_name=FONT_PRIMARY, font_size=SIZE_BODY,
                          font_color=SURFACE_WHITE, bold=False,
                          alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                          shape_name=None):
        """Add a rounded rectangle shape with optional text."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
        )
        if shape_name:
            shape.name = shape_name
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()  # No border

        if text:
            tf = shape.text_frame
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.word_wrap = True
            try:
                tf.vertical_anchor = anchor
            except Exception:
                pass

            p = tf.paragraphs[0]
            p.alignment = alignment
            run = p.add_run()
            run.text = text
            run.font.name = font_name
            run.font.size = font_size
            run.font.color.rgb = font_color
            run.font.bold = bold

        return shape

    def _add_rect(self, slide, left, top, width, height,
                  fill_color=ACTION_BLUE, text="",
                  font_name=FONT_PRIMARY, font_size=SIZE_BODY,
                  font_color=SURFACE_WHITE, bold=False,
                  alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                  shape_name=None):
        """Add a rectangle shape with optional text."""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        if shape_name:
            shape.name = shape_name
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()

        if text:
            tf = shape.text_frame
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.word_wrap = True
            try:
                tf.vertical_anchor = anchor
            except Exception:
                pass

            p = tf.paragraphs[0]
            p.alignment = alignment
            run = p.add_run()
            run.text = text
            run.font.name = font_name
            run.font.size = font_size
            run.font.color.rgb = font_color
            run.font.bold = bold

        return shape

    # ──────────────────────────────────────────
    # Slide Layout Methods
    # ──────────────────────────────────────────

    def add_cover_slide(self, section_label: str, title: str, subtitle: str = ""):
        """Layout 1: Title/Cover — navy background, white text."""
        slide = self._add_blank_slide(dark=True)
        self._set_slide_bg(slide, PRIMARY_NAVY)

        # Section label (e.g., "INTRODUCTION")
        self._add_textbox(
            slide, Inches(1.5), Inches(3.67), Inches(6.3), Inches(0.4),
            text=section_label.upper(),
            font_size=SIZE_BODY, font_color=SURFACE_WHITE,
            bold=False
        )

        # Title — split by | for line breaks
        title_lines = [t.strip() for t in title.split("|")]
        self._add_multiline_textbox(
            slide, Inches(1.5), Inches(4.51), Inches(8.5), Inches(2.0),
            lines=title_lines,
            font_size=SIZE_H1, font_color=SURFACE_WHITE, bold=True
        )

        # Subtitle / date
        if subtitle:
            self._add_textbox(
                slide, Inches(1.5), Inches(7.31), Inches(8.5), Inches(0.4),
                text=subtitle,
                font_size=SIZE_BODY, font_color=SURFACE_WHITE
            )

        return slide

    def add_section_divider(self, section_label: str, title: str, tagline: str = ""):
        """Layout 2: Section Divider — navy background, introduces a new section."""
        slide = self._add_blank_slide(dark=True)
        self._set_slide_bg(slide, PRIMARY_NAVY)

        self._add_textbox(
            slide, Inches(1.67), Inches(3.32), Inches(6.3), Inches(0.4),
            text=section_label.upper(),
            font_size=SIZE_LABEL, font_color=SURFACE_WHITE
        )

        title_lines = [t.strip() for t in title.split("|")]
        self._add_multiline_textbox(
            slide, Inches(1.67), Inches(4.15), Inches(13.0), Inches(1.0),
            lines=title_lines,
            font_size=SIZE_H1, font_color=SURFACE_WHITE, bold=True
        )

        if tagline:
            self._add_textbox(
                slide, Inches(1.67), Inches(8.23), Inches(10.0), Inches(0.4),
                text=tagline,
                font_size=SIZE_BODY, font_color=SURFACE_WHITE
            )

        return slide

    def add_agenda_slide(self, section_label: str, title: str,
                         customer_name: str, agenda_items: list):
        """Layout 3: Agenda — left branding, right stacked items."""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)

        # Left side
        self._add_textbox(
            slide, Inches(1.48), Inches(1.55), Inches(5.13), Inches(0.3),
            text=section_label.upper(),
            font_size=SIZE_CAPTION, font_color=TEXT_MAIN
        )

        self._add_textbox(
            slide, Inches(1.41), Inches(2.22), Inches(5.62), Inches(0.76),
            text=title,
            font_size=SIZE_H1, font_color=TEXT_MAIN, bold=True
        )

        self._add_textbox(
            slide, Inches(1.47), Inches(5.25), Inches(5.62), Inches(0.4),
            text=customer_name,
            font_size=SIZE_BODY, font_color=TEXT_MAIN
        )

        # Right side — agenda items
        y_positions = [1.59, 3.46, 5.40, 7.30, 9.22]
        for i, item in enumerate(agenda_items[:5]):
            y = y_positions[i] if i < len(y_positions) else y_positions[-1] + (i - 4) * 1.9
            self._add_textbox(
                slide, Inches(9.69), Inches(y), Inches(8.9), Inches(0.4),
                text=item,
                font_size=SIZE_AGENDA, font_color=TEXT_MAIN, bold=True
            )

        return slide

    def add_content_slide(self, title: str, subtitle: str = "",
                          body_lines: list = None):
        """Layout 4: Challenge/Content — full-width with title bar."""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)

        self._add_textbox(
            slide, Inches(1.17), Inches(1.14), Inches(17.68), Inches(0.76),
            text=title,
            font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True
        )

        if subtitle:
            self._add_textbox(
                slide, Inches(1.17), Inches(2.0), Inches(17.58), Inches(0.4),
                text=subtitle,
                font_size=SIZE_BODY, font_color=TEXT_MUTED
            )

        if body_lines:
            self._add_multiline_textbox(
                slide, Inches(1.17), Inches(3.0), Inches(17.68), Inches(7.0),
                lines=body_lines,
                font_size=SIZE_BODY, font_color=TEXT_MAIN
            )

        return slide

    def add_split_comparison(self, title: str, section_label: str = "",
                             left_title: str = "", left_items: list = None,
                             right_title: str = "", right_items: list = None):
        """Layout 5: Split Comparison (From/To) — vertical divider, red 'from' heading."""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)

        # Vertical divider line between columns
        from pptx.enum.shapes import MSO_SHAPE
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(9.9), Inches(3.8), Inches(0.02), Inches(6.0)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0xE5, 0xE9, 0xF0)
        line.line.fill.background()
        line.name = "ColumnDivider"

        # Section label
        if section_label:
            self._add_textbox(
                slide, Inches(1.44), Inches(1.30), Inches(8.18), Inches(0.3),
                text=section_label.upper(),
                font_size=SIZE_CAPTION, font_color=TEXT_MAIN
            )

        # Title
        self._add_textbox(
            slide, Inches(1.42), Inches(1.93), Inches(17.2), Inches(0.76),
            text=title,
            font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True
        )

        # Left column — red heading for "from" state
        self._add_textbox(
            slide, Inches(1.36), Inches(4.27), Inches(8.07), Inches(0.4),
            text=left_title.upper(),
            font_size=SIZE_BODY, font_color=SEMANTIC_RED, bold=True
        )
        if left_items:
            self._add_multiline_textbox(
                slide, Inches(1.36), Inches(5.17), Inches(8.07), Inches(3.23),
                lines=left_items,
                font_size=SIZE_BODY, font_color=TEXT_MAIN
            )

        # Right column
        self._add_textbox(
            slide, Inches(10.55), Inches(4.27), Inches(8.07), Inches(0.4),
            text=right_title.upper(),
            font_size=SIZE_BODY, font_color=ACTION_BLUE, bold=True
        )
        if right_items:
            self._add_multiline_textbox(
                slide, Inches(10.55), Inches(5.17), Inches(8.07), Inches(3.23),
                lines=right_items,
                font_size=SIZE_BODY, font_color=TEXT_MAIN
            )

        return slide

    def add_showcase_slide(self, section_label: str, title: str,
                           description: str, image_path: str = None):
        """Layout 6: Product Showcase — left text, right image."""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)

        self._add_textbox(
            slide, Inches(1.48), Inches(1.82), Inches(5.13), Inches(0.3),
            text=section_label.upper(),
            font_size=SIZE_CAPTION, font_color=TEXT_MAIN
        )

        self._add_textbox(
            slide, Inches(1.41), Inches(2.49), Inches(6.23), Inches(0.76),
            text=title,
            font_size=SIZE_H1, font_color=TEXT_MAIN, bold=True
        )

        self._add_textbox(
            slide, Inches(1.41), Inches(6.31), Inches(6.23), Inches(1.5),
            text=description,
            font_size=SIZE_BODY, font_color=TEXT_MAIN
        )

        if image_path and os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path,
                Inches(9.16), Inches(1.40),
                Inches(9.21), Inches(8.45)
            )

        return slide

    def add_architecture_slide(self, title: str, subtitle: str = "",
                               customer_channels: list = None,
                               employee_workspaces: list = None,
                               platform_label: str = "AI-native Banking OS",
                               enablement_systems: list = None,
                               core_systems: list = None):
        """Layout 7: Architecture Diagram — layered platform view."""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)

        # Title
        self._add_textbox(
            slide, Inches(1.17), Inches(1.14), Inches(17.68), Inches(0.76),
            text=title,
            font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True
        )

        if subtitle:
            self._add_textbox(
                slide, Inches(1.17), Inches(2.0), Inches(17.58), Inches(0.4),
                text=subtitle,
                font_size=SIZE_BODY, font_color=TEXT_MUTED
            )

        # Layout constants
        content_left = Inches(1.5)
        content_width = Inches(17.0)
        row_height = Inches(0.93)
        box_gap = Inches(0.2)

        # Row 1: Customer Channels
        channels = customer_channels or ["Online", "Mobile", "Conversational"]
        y = Inches(3.5)
        self._add_textbox(
            slide, content_left, Inches(3.0), Inches(3.0), Inches(0.4),
            text="Customers", font_size=SIZE_BODY, font_color=TEXT_MAIN,
            bold=False, apply_buffer=False
        )
        box_w = int((content_width - box_gap * (len(channels) - 1)) / len(channels))
        for i, ch in enumerate(channels):
            x = int(content_left + i * (box_w + box_gap))
            self._add_rounded_rect(
                slide, x, y, box_w, row_height,
                fill_color=PRIMARY_NAVY, text=ch,
                font_size=SIZE_BODY, font_color=SURFACE_WHITE, bold=True
            )

        # Row 2: Employee Workspaces
        workspaces = employee_workspaces or ["Teller", "CSR", "RM / Advisor", "Operations"]
        y = Inches(4.7)
        self._add_textbox(
            slide, content_left, Inches(4.2), Inches(3.0), Inches(0.4),
            text="Employees", font_size=SIZE_BODY, font_color=TEXT_MAIN,
            bold=False, apply_buffer=False
        )
        box_w = int((content_width - box_gap * (len(workspaces) - 1)) / len(workspaces))
        for i, ws in enumerate(workspaces):
            x = int(content_left + i * (box_w + box_gap))
            self._add_rounded_rect(
                slide, x, y, box_w, row_height,
                fill_color=PRIMARY_NAVY, text=ws,
                font_size=SIZE_BODY, font_color=SURFACE_WHITE, bold=True
            )

        # Platform Bar
        y = Inches(6.0)
        self._add_rounded_rect(
            slide, content_left, y, content_width, Inches(0.7),
            fill_color=ACTION_BLUE, text=platform_label,
            font_size=SIZE_H2, font_color=SURFACE_WHITE, bold=True
        )

        # Row 3: Enablement Systems
        systems = enablement_systems or ["CRM", "CDP", "KYC", "AML", "Fraud", "Risk", "LMS"]
        y = Inches(7.0)
        self._add_textbox(
            slide, content_left, Inches(7.4), Inches(5.0), Inches(0.3),
            text="Enablement Systems", font_size=SIZE_CAPTION, font_color=TEXT_MUTED,
            apply_buffer=False
        )
        box_w = int((content_width - box_gap * (len(systems) - 1)) / len(systems))
        for i, sys_name in enumerate(systems):
            x = int(content_left + i * (box_w + box_gap))
            self._add_rounded_rect(
                slide, x, y, box_w, Inches(0.6),
                fill_color=BG_GRAY, text=sys_name,
                font_size=SIZE_CAPTION, font_color=TEXT_MAIN
            )

        # Row 4: Core Systems
        core = core_systems or ["Ledger", "Cards", "Payments", "Back Office"]
        y = Inches(7.9)
        self._add_textbox(
            slide, content_left, Inches(8.3), Inches(5.0), Inches(0.3),
            text="Core Systems", font_size=SIZE_CAPTION, font_color=TEXT_MUTED,
            apply_buffer=False
        )
        box_w = int((content_width - box_gap * (len(core) - 1)) / len(core))
        for i, c in enumerate(core):
            x = int(content_left + i * (box_w + box_gap))
            self._add_rounded_rect(
                slide, x, y, box_w, Inches(0.6),
                fill_color=BG_GRAY, text=c,
                font_size=SIZE_CAPTION, font_color=TEXT_MAIN
            )

        return slide

    def add_case_study_slide(self, title: str, body_lines: list = None,
                             legal_text: str = None):
        """Layout 8: Customer Case Study — with legal footer."""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)

        # Case label (top right)
        self._add_textbox(
            slide, Inches(16.0), Inches(1.0), Inches(3.0), Inches(0.3),
            text="CUSTOMER CASE",
            font_size=SIZE_CAPTION, font_color=RGBColor(0xCC, 0x00, 0x00),
            bold=True, alignment=PP_ALIGN.RIGHT
        )

        # Title
        self._add_textbox(
            slide, Inches(1.5), Inches(2.0), Inches(17.0), Inches(0.76),
            text=title,
            font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True
        )

        # Body
        if body_lines:
            self._add_multiline_textbox(
                slide, Inches(1.5), Inches(3.5), Inches(17.0), Inches(5.5),
                lines=body_lines,
                font_size=SIZE_BODY, font_color=TEXT_MAIN
            )

        # Legal footer
        footer = legal_text or (
            "Restricted use. This case study is intended solely for use "
            "in 1:1 discussions with prospective clients."
        )
        self._add_textbox(
            slide, Inches(1.5), Inches(10.0), Inches(17.0), Inches(0.5),
            text=footer,
            font_size=SIZE_LEGAL, font_color=TEXT_MUTED
        )

        return slide

    def add_stat_cards_slide(self, title: str, subtitle: str = "",
                             stats: list = None):
        """Stat cards layout — large numbers with labels.

        stats: list of dicts with keys: number, label, trend (optional)
        Example: [{"number": "3×", "label": "Faster velocity", "trend": "+12%"}]
        """
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)

        self._add_textbox(
            slide, Inches(1.17), Inches(1.14), Inches(17.68), Inches(0.76),
            text=title,
            font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True
        )

        if subtitle:
            self._add_textbox(
                slide, Inches(1.17), Inches(2.0), Inches(17.58), Inches(0.4),
                text=subtitle,
                font_size=SIZE_BODY, font_color=TEXT_MUTED
            )

        stats = stats or []
        n = len(stats)
        if n == 0:
            return slide

        # 2×3 grid for 6 cards, otherwise single row
        cols = 3 if n >= 5 else n
        rows_count = (n + cols - 1) // cols

        card_gap = Inches(0.3)
        total_w = Inches(16.5)
        card_w = int((total_w - card_gap * (cols - 1)) / cols)
        card_h = Inches(3.0)
        start_x = Inches(1.5)
        start_y = Inches(3.2)
        row_gap = Inches(0.3)

        for i, stat in enumerate(stats):
            col = i % cols
            row = i // cols
            x = int(start_x + col * (card_w + card_gap))
            y = int(start_y + row * (card_h + row_gap))

            # Card background
            self._add_rounded_rect(
                slide, x, y, card_w, card_h,
                fill_color=BG_GRAY
            )

            # Title label at top
            self._add_textbox(
                slide, x + Inches(0.2), y + Inches(0.2),
                card_w - Inches(0.4), Inches(0.3),
                text=stat.get("label", "").upper(),
                font_size=Pt(9), font_color=TEXT_MUTED, bold=True,
                alignment=PP_ALIGN.CENTER, apply_buffer=False
            )

            # Number
            self._add_textbox(
                slide, x + Inches(0.2), y + Inches(0.6),
                card_w - Inches(0.4), Inches(1.2),
                text=stat.get("number", ""),
                font_size=Pt(36), font_color=ACTION_BLUE, bold=True,
                alignment=PP_ALIGN.CENTER, apply_buffer=False
            )

            # Trend
            trend = stat.get("trend", "")
            if trend:
                is_positive = any(trend.startswith(c) for c in ["+", "\u2191", "2"])
                is_negative = any(trend.startswith(c) for c in ["-", "\u2212", "\u2193"])
                trend_color = SUCCESS_GREEN if is_positive or is_negative else TEXT_MUTED
                self._add_textbox(
                    slide, x + Inches(0.2), y + Inches(2.2),
                    card_w - Inches(0.4), Inches(0.4),
                    text=trend,
                    font_size=SIZE_CAPTION, font_color=trend_color,
                    bold=True, alignment=PP_ALIGN.CENTER, apply_buffer=False
                )

        return slide

    # ──────────────────────────────────────────
    # Rich Tile Components (matching HTML builder)
    # ──────────────────────────────────────────

    def add_tiles_slide(self, title: str, subtitle: str = "", section_label: str = "",
                        tiles: list = None, columns: int = 3):
        """Grid of colored tiles/cards. Each tile: {title, body, stat, pill, accent}
        accent: blue/red/green/amber/purple/cyan"""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)
        tiles = tiles or []

        # Section label
        if section_label:
            self._add_textbox(slide, Inches(1.17), Inches(1.14), Inches(5.0), Inches(0.3),
                text=section_label.upper(), font_size=SIZE_LABEL, font_color=TEXT_MUTED)

        # Title
        y_title = Inches(1.14) if not section_label else Inches(1.5)
        self._add_textbox(slide, Inches(1.17), y_title, Inches(17.68), Inches(0.76),
            text=title, font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True)

        if subtitle:
            self._add_textbox(slide, Inches(1.17), y_title + Inches(0.85), Inches(17.58), Inches(0.4),
                text=subtitle, font_size=SIZE_CAPTION, font_color=TEXT_MUTED)

        # Color map
        accent_colors = {
            'blue': ACTION_BLUE, 'red': SEMANTIC_RED, 'green': SUCCESS_GREEN,
            'amber': RGBColor(0xD9, 0x77, 0x06), 'purple': RGBColor(0x7B, 0x2F, 0xFF),
            'cyan': RGBColor(0x08, 0x91, 0xB2)
        }

        # Tile grid — compact height based on content
        n = len(tiles)
        cols = min(columns, n) if n > 0 else columns
        rows_count = (n + cols - 1) // cols if n > 0 else 1
        gap = Inches(0.2)
        grid_left = Inches(1.17)
        grid_width = Inches(17.68)
        tile_w = int((grid_width - gap * (cols - 1)) / cols)

        # Estimate tile height from content — compact, not fill-the-slide
        # Base: pill(0.28) + stat(0.55) + title(0.3) + body(0.6) + padding(0.4) ≈ 2.1"
        def _body_lines(b):
            if isinstance(b, list):
                return len(b)
            return len(b.split('\n'))
        has_multiline = any(_body_lines(t.get('body', '')) > 2 for t in tiles)
        tile_h = Inches(2.8) if has_multiline else Inches(2.1)
        # Cap so two rows fit
        if rows_count > 1:
            max_h = (Inches(7.2) - gap * (rows_count - 1)) / rows_count
            tile_h = min(tile_h, max_h)

        start_y = y_title + Inches(1.8)

        for i, t in enumerate(tiles):
            col = i % cols
            row = i // cols
            x = int(grid_left + col * (tile_w + gap))
            y = int(start_y + row * (tile_h + gap))
            accent = t.get('accent', 'blue')
            accent_color = accent_colors.get(accent, ACTION_BLUE)

            # Tile background — use regular rect, not rounded (avoids bubble look in Slides)
            self._add_rect(slide, x, y, tile_w, tile_h,
                fill_color=BG_GRAY, shape_name=f"Tile_{i}")
            # Top accent line — thin colored bar at top (like HTML)
            self._add_rect(slide, x, y, tile_w, Inches(0.04),
                fill_color=accent_color, shape_name=f"TileAccent_{i}")

            # Pill
            pill_text = t.get('pill', '')
            curr_y = y + Inches(0.18)
            if pill_text:
                self._add_textbox(slide, x + Inches(0.25), curr_y, tile_w - Inches(0.5), Inches(0.22),
                    text=pill_text.upper(), font_size=Pt(8), font_color=accent_color,
                    bold=True, apply_buffer=False)
                curr_y += Inches(0.26)

            # Stat number
            stat_text = t.get('stat', '')
            if stat_text:
                self._add_textbox(slide, x + Inches(0.25), curr_y, tile_w - Inches(0.5), Inches(0.5),
                    text=stat_text, font_size=Pt(24), font_color=accent_color,
                    bold=True, apply_buffer=False)
                curr_y += Inches(0.5)

            # Title
            tile_title = t.get('title', '')
            if tile_title:
                self._add_textbox(slide, x + Inches(0.25), curr_y, tile_w - Inches(0.5), Inches(0.28),
                    text=tile_title, font_size=Pt(11), font_color=TEXT_MAIN,
                    bold=True, apply_buffer=False)
                curr_y += Inches(0.3)

            # Body
            body = t.get('body', '')
            if isinstance(body, list):
                body = '\n'.join(body)
            if body:
                remaining = max(Inches(0.3), tile_h - (curr_y - y) - Inches(0.1))
                self._add_textbox(slide, x + Inches(0.25), curr_y, tile_w - Inches(0.5), remaining,
                    text=body, font_size=Pt(10), font_color=TEXT_MUTED,
                    apply_buffer=False)

        return slide

    def add_process_rows_slide(self, title: str, subtitle: str = "", section_label: str = "",
                                rows: list = None, footer_text: str = ""):
        """Before → After process comparison rows.
        Each row: {label, before, after, saving}"""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)
        rows = rows or []

        if section_label:
            self._add_textbox(slide, Inches(1.17), Inches(1.14), Inches(5.0), Inches(0.3),
                text=section_label.upper(), font_size=SIZE_LABEL, font_color=TEXT_MUTED)

        self._add_textbox(slide, Inches(1.17), Inches(1.5), Inches(17.68), Inches(0.76),
            text=title, font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True)

        if subtitle:
            self._add_textbox(slide, Inches(1.17), Inches(2.35), Inches(17.58), Inches(0.4),
                text=subtitle, font_size=SIZE_CAPTION, font_color=TEXT_MUTED)

        row_h = Inches(0.85)
        row_gap = Inches(0.12)
        start_y = Inches(3.2)

        # Column headers
        header_y = start_y - Inches(0.4)
        self._add_textbox(slide, Inches(1.5), header_y, Inches(6.0), Inches(0.3),
            text="PROCESS", font_size=Pt(9), font_color=TEXT_MUTED,
            bold=True, apply_buffer=False)
        self._add_textbox(slide, Inches(8.0), header_y, Inches(2.5), Inches(0.3),
            text="BEFORE", font_size=Pt(9), font_color=TEXT_MUTED,
            bold=True, apply_buffer=False, alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, Inches(11.3), header_y, Inches(2.5), Inches(0.3),
            text="AFTER", font_size=Pt(9), font_color=TEXT_MUTED,
            bold=True, apply_buffer=False, alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, Inches(14.5), header_y, Inches(3.5), Inches(0.3),
            text="IMPACT", font_size=Pt(9), font_color=TEXT_MUTED,
            bold=True, apply_buffer=False, alignment=PP_ALIGN.RIGHT)

        for i, r in enumerate(rows):
            y = int(start_y + i * (row_h + row_gap))

            # Row background
            self._add_rounded_rect(slide, Inches(1.17), y, Inches(17.68), row_h,
                fill_color=BG_GRAY, shape_name=f"ProcRow_{i}")

            # Label
            self._add_textbox(slide, Inches(1.5), y + Inches(0.2), Inches(6.0), Inches(0.45),
                text=r.get('label', ''), font_size=Pt(13), font_color=TEXT_MAIN,
                bold=True, apply_buffer=False)

            # Before (red)
            self._add_textbox(slide, Inches(8.0), y + Inches(0.2), Inches(2.5), Inches(0.45),
                text=r.get('before', ''), font_size=Pt(13), font_color=SEMANTIC_RED,
                bold=True, apply_buffer=False, alignment=PP_ALIGN.CENTER)

            # Arrow
            self._add_textbox(slide, Inches(10.5), y + Inches(0.2), Inches(0.8), Inches(0.45),
                text='\u2192', font_size=Pt(16), font_color=TEXT_MUTED,
                apply_buffer=False, alignment=PP_ALIGN.CENTER)

            # After (green)
            self._add_textbox(slide, Inches(11.3), y + Inches(0.2), Inches(2.5), Inches(0.45),
                text=r.get('after', ''), font_size=Pt(13), font_color=SUCCESS_GREEN,
                bold=True, apply_buffer=False, alignment=PP_ALIGN.CENTER)

            # Saving
            self._add_textbox(slide, Inches(14.5), y + Inches(0.2), Inches(3.5), Inches(0.45),
                text=r.get('saving', ''), font_size=Pt(16),
                font_color=ACTION_BLUE,
                bold=True, apply_buffer=False, alignment=PP_ALIGN.RIGHT)

        if footer_text:
            self._add_textbox(slide, Inches(1.17), Inches(9.5), Inches(17.68), Inches(0.4),
                text=footer_text, font_size=Pt(10), font_color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER, apply_buffer=False)

        return slide

    def add_bar_chart_slide(self, title: str, subtitle: str = "", section_label: str = "",
                             rows: list = None, summary_boxes: list = None,
                             footer_text: str = ""):
        """Horizontal bar chart rows with before/after bars + summary boxes.
        Each row: {label, before_mins, after_mins, saving}
        summary_boxes: [{label, stat, source}] or [{label, items: [{stat, label}], source}]
        """
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)
        rows = rows or []
        summary_boxes = summary_boxes or []

        if section_label:
            self._add_textbox(slide, Inches(1.17), Inches(1.14), Inches(5.0), Inches(0.3),
                text=section_label.upper(), font_size=SIZE_LABEL, font_color=TEXT_MUTED)
        self._add_textbox(slide, Inches(1.17), Inches(1.5), Inches(17.68), Inches(0.76),
            text=title, font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True)
        if subtitle:
            self._add_textbox(slide, Inches(1.17), Inches(2.35), Inches(17.58), Inches(0.4),
                text=subtitle, font_size=SIZE_CAPTION, font_color=TEXT_MUTED)

        max_mins = max((r.get("before_mins", 0) for r in rows), default=1) or 1
        bar_area_w = Inches(10.0)
        label_w = Inches(4.5)
        label_x = Inches(1.17)
        bar_x = Inches(5.8)
        start_y = Inches(3.2)
        row_h = Inches(0.9)

        # Soft colours
        before_color = RGBColor(0xF4, 0xA6, 0xA0)  # soft salmon
        after_color = RGBColor(0x93, 0xB5, 0xFF)    # soft blue
        before_text_c = RGBColor(0xC0, 0x56, 0x4F)
        after_text_c = ACTION_BLUE

        # Legend
        leg_y = start_y - Inches(0.35)
        self._add_rect(slide, Inches(6.0), leg_y, Inches(0.2), Inches(0.2),
            fill_color=before_color, shape_name="LegBefore")
        self._add_textbox(slide, Inches(6.3), leg_y, Inches(1.5), Inches(0.2),
            text="Current state", font_size=Pt(8), font_color=TEXT_MUTED,
            apply_buffer=False)
        self._add_rect(slide, Inches(8.0), leg_y, Inches(0.2), Inches(0.2),
            fill_color=after_color, shape_name="LegAfter")
        self._add_textbox(slide, Inches(8.3), leg_y, Inches(1.5), Inches(0.2),
            text="With Backbase", font_size=Pt(8), font_color=TEXT_MUTED,
            apply_buffer=False)

        for i, r in enumerate(rows):
            y = int(start_y + i * row_h)
            before = r.get("before_mins", 0)
            after = r.get("after_mins", 0)
            saving = r.get("saving", "")
            before_w = max(int(bar_area_w * before / max_mins), Inches(0.1))
            after_w = max(int(bar_area_w * after / max_mins), Inches(0.1))
            bar_h = Inches(0.25)

            # Label
            self._add_textbox(slide, label_x, y + Inches(0.15), label_w, Inches(0.4),
                text=r.get('label', ''), font_size=Pt(11), font_color=TEXT_MAIN,
                bold=True, apply_buffer=False, alignment=PP_ALIGN.RIGHT)

            # Before bar
            self._add_rounded_rect(slide, bar_x, y + Inches(0.05), before_w, bar_h,
                fill_color=before_color, shape_name=f"BarBefore_{i}")
            self._add_textbox(slide, int(bar_x + before_w + Inches(0.1)), y + Inches(0.05),
                Inches(1.5), bar_h,
                text=f"{before} mins", font_size=Pt(9), font_color=before_text_c,
                bold=True, apply_buffer=False)

            # After bar
            self._add_rounded_rect(slide, bar_x, y + Inches(0.35), after_w, bar_h,
                fill_color=after_color, shape_name=f"BarAfter_{i}")
            self._add_textbox(slide, int(bar_x + after_w + Inches(0.1)), y + Inches(0.35),
                Inches(2.5), bar_h,
                text=f"{after} mins ({saving})", font_size=Pt(9), font_color=after_text_c,
                bold=True, apply_buffer=False)

        # Summary boxes
        box_y = int(start_y + len(rows) * row_h + Inches(0.3))
        if summary_boxes:
            box_w = int(Inches(17.68) / len(summary_boxes) - Inches(0.15))
            for bi, box in enumerate(summary_boxes):
                bx = int(Inches(1.17) + bi * (box_w + Inches(0.15)))
                self._add_rounded_rect(slide, bx, box_y, box_w, Inches(1.3),
                    fill_color=RGBColor(0xF0, 0xF4, 0xFF), shape_name=f"SummBox_{bi}")
                shape = slide.shapes[-1]
                shape.line.color.rgb = RGBColor(0xCC, 0xD5, 0xFF)
                shape.line.width = Pt(1)

                # Label
                self._add_textbox(slide, bx + Inches(0.2), box_y + Inches(0.1),
                    box_w - Inches(0.4), Inches(0.25),
                    text=box.get("label", "").upper(), font_size=Pt(9),
                    font_color=ACTION_BLUE, bold=True,
                    apply_buffer=False, alignment=PP_ALIGN.CENTER)

                if "stat" in box:
                    # Single stat box
                    self._add_textbox(slide, bx + Inches(0.2), box_y + Inches(0.35),
                        box_w - Inches(0.4), Inches(0.5),
                        text=box["stat"], font_size=Pt(28), font_color=ACTION_BLUE,
                        bold=True, apply_buffer=False, alignment=PP_ALIGN.CENTER)
                    if box.get("source"):
                        self._add_textbox(slide, bx + Inches(0.2), box_y + Inches(0.9),
                            box_w - Inches(0.4), Inches(0.25),
                            text=box["source"], font_size=Pt(7), font_color=TEXT_MUTED,
                            apply_buffer=False, alignment=PP_ALIGN.CENTER)

                elif "items" in box:
                    # Multi-stat box (FTE freed)
                    items = box["items"]
                    item_w = int((box_w - Inches(0.4)) / len(items))
                    for ii, item in enumerate(items):
                        ix = int(bx + Inches(0.2) + ii * item_w)
                        self._add_textbox(slide, ix, box_y + Inches(0.35),
                            item_w, Inches(0.35),
                            text=item["stat"], font_size=Pt(22), font_color=ACTION_BLUE,
                            bold=True, apply_buffer=False, alignment=PP_ALIGN.CENTER)
                        self._add_textbox(slide, ix, box_y + Inches(0.7),
                            item_w, Inches(0.2),
                            text=item["label"], font_size=Pt(7), font_color=TEXT_MUTED,
                            bold=True, apply_buffer=False, alignment=PP_ALIGN.CENTER)
                    if box.get("source"):
                        self._add_textbox(slide, bx + Inches(0.2), box_y + Inches(0.95),
                            box_w - Inches(0.4), Inches(0.25),
                            text=box["source"], font_size=Pt(7), font_color=TEXT_MUTED,
                            apply_buffer=False, alignment=PP_ALIGN.CENTER)

        if footer_text:
            fy = box_y + Inches(1.5) if summary_boxes else int(start_y + len(rows) * row_h + Inches(0.3))
            self._add_textbox(slide, Inches(1.17), fy, Inches(17.68), Inches(0.4),
                text=footer_text, font_size=Pt(9), font_color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER, apply_buffer=False)

        return slide

    def add_pillar_rows_slide(self, title: str, subtitle: str = "", section_label: str = "",
                               columns: list = None, rows: list = None):
        """Three-column pillar flow (What → Why → What To Do).
        columns: [left_header, mid_header, right_header]
        rows: [{left, left_detail, mid, mid_detail, right, right_detail}]"""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)
        columns = columns or ["What's Happening", "Why", "What To Do"]
        rows = rows or []

        if section_label:
            self._add_textbox(slide, Inches(1.17), Inches(1.14), Inches(5.0), Inches(0.3),
                text=section_label.upper(), font_size=SIZE_LABEL, font_color=TEXT_MUTED)

        self._add_textbox(slide, Inches(1.17), Inches(1.5), Inches(17.68), Inches(0.76),
            text=title, font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True)

        if subtitle:
            self._add_textbox(slide, Inches(1.17), Inches(2.35), Inches(17.58), Inches(0.4),
                text=subtitle, font_size=SIZE_CAPTION, font_color=TEXT_MUTED)

        # Column headers — 3 columns fit in 17.68" with arrows
        col_w = Inches(5.4)
        arrow_w = Inches(0.74)
        start_x = Inches(1.17)
        header_y = Inches(3.2)

        self._add_textbox(slide, start_x, header_y, col_w, Inches(0.3),
            text=columns[0].upper(), font_size=Pt(10), font_color=SEMANTIC_RED,
            bold=True, alignment=PP_ALIGN.CENTER, apply_buffer=False)
        self._add_textbox(slide, start_x + col_w + arrow_w, header_y, col_w, Inches(0.3),
            text=columns[1].upper(), font_size=Pt(10), font_color=TEXT_MAIN,
            bold=True, alignment=PP_ALIGN.CENTER, apply_buffer=False)
        self._add_textbox(slide, start_x + (col_w + arrow_w) * 2, header_y, col_w, Inches(0.3),
            text=columns[2].upper(), font_size=Pt(10), font_color=ACTION_BLUE,
            bold=True, alignment=PP_ALIGN.CENTER, apply_buffer=False)

        # Rows
        row_h = Inches(1.2)
        row_gap = Inches(0.1)
        start_y = Inches(3.7)

        for i, r in enumerate(rows):
            y = int(start_y + i * (row_h + row_gap))

            # Left cell (red tint)
            self._add_rounded_rect(slide, start_x, y, col_w, row_h,
                fill_color=RGBColor(0xFD, 0xF0, 0xF0), shape_name=f"PillarL_{i}")
            self._add_textbox(slide, start_x + Inches(0.15), y + Inches(0.1), col_w - Inches(0.3), Inches(0.3),
                text=r.get('left', ''), font_size=Pt(12), font_color=TEXT_MAIN,
                bold=True, apply_buffer=False)
            self._add_textbox(slide, start_x + Inches(0.15), y + Inches(0.45), col_w - Inches(0.3), Inches(0.6),
                text=r.get('left_detail', ''), font_size=Pt(10), font_color=TEXT_MUTED,
                apply_buffer=False)

            # Arrow
            self._add_textbox(slide, start_x + col_w, y + Inches(0.3), arrow_w, Inches(0.3),
                text='\u2192', font_size=SIZE_CAPTION, font_color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER, apply_buffer=False)

            # Mid cell (neutral)
            mid_x = start_x + col_w + arrow_w
            self._add_rounded_rect(slide, mid_x, y, col_w, row_h,
                fill_color=BG_GRAY, shape_name=f"PillarM_{i}")
            self._add_textbox(slide, mid_x + Inches(0.15), y + Inches(0.1), col_w - Inches(0.3), Inches(0.3),
                text=r.get('mid', ''), font_size=Pt(12), font_color=TEXT_MAIN,
                bold=True, apply_buffer=False)
            self._add_textbox(slide, mid_x + Inches(0.15), y + Inches(0.45), col_w - Inches(0.3), Inches(0.6),
                text=r.get('mid_detail', ''), font_size=Pt(10), font_color=TEXT_MUTED,
                apply_buffer=False)

            # Arrow
            self._add_textbox(slide, mid_x + col_w, y + Inches(0.3), arrow_w, Inches(0.3),
                text='\u2192', font_size=SIZE_CAPTION, font_color=ACTION_BLUE,
                alignment=PP_ALIGN.CENTER, apply_buffer=False)

            # Right cell (blue tint)
            right_x = start_x + (col_w + arrow_w) * 2
            self._add_rounded_rect(slide, right_x, y, col_w, row_h,
                fill_color=RGBColor(0xEE, 0xF2, 0xFF), shape_name=f"PillarR_{i}")
            self._add_textbox(slide, right_x + Inches(0.15), y + Inches(0.1), col_w - Inches(0.3), Inches(0.3),
                text=r.get('right', ''), font_size=Pt(12), font_color=TEXT_MAIN,
                bold=True, apply_buffer=False)
            self._add_textbox(slide, right_x + Inches(0.15), y + Inches(0.45), col_w - Inches(0.3), Inches(0.6),
                text=r.get('right_detail', ''), font_size=Pt(10), font_color=TEXT_MUTED,
                apply_buffer=False)

        return slide

    def add_financial_table_slide(self, title: str, subtitle: str = "", section_label: str = "",
                                   headers: list = None, rows: list = None,
                                   total_row: list = None, footer_text: str = ""):
        """Financial table with headers, rows, total row."""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)
        headers = headers or []
        rows = rows or []

        if section_label:
            self._add_textbox(slide, Inches(1.17), Inches(1.14), Inches(5.0), Inches(0.3),
                text=section_label.upper(), font_size=SIZE_LABEL, font_color=TEXT_MUTED)

        self._add_textbox(slide, Inches(1.17), Inches(1.5), Inches(17.68), Inches(0.76),
            text=title, font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True)

        if subtitle:
            self._add_textbox(slide, Inches(1.17), Inches(2.35), Inches(17.58), Inches(0.4),
                text=subtitle, font_size=SIZE_CAPTION, font_color=TEXT_MUTED)

        # Table
        n_cols = len(headers)
        if n_cols == 0:
            return slide

        table_left = Inches(1.17)
        table_top = Inches(3.2)
        table_width = Inches(17.68)
        col_w = table_width / n_cols
        row_h = Inches(0.6)

        # Header row
        for j, h in enumerate(headers):
            x = int(table_left + j * col_w)
            self._add_textbox(slide, x, table_top, int(col_w), Inches(0.4),
                text=h.upper(), font_size=Pt(9), font_color=TEXT_MUTED,
                bold=True, apply_buffer=False)
        # Header underline
        self._add_rect(slide, table_left, table_top + Inches(0.4),
            table_width, Inches(0.02), fill_color=TEXT_MUTED, shape_name="TableHeaderLine")

        # Data rows
        data_start = table_top + Inches(0.55)
        all_rows = rows + ([{"cells": total_row, "total": True}] if total_row else [])

        for i, row in enumerate(all_rows):
            y = int(data_start + i * row_h)
            cells = row.get("cells", row) if isinstance(row, dict) else row
            is_total = row.get("total", False) if isinstance(row, dict) else False
            is_highlight = row.get("highlight", False) if isinstance(row, dict) else False

            # Highlight background
            if is_highlight:
                self._add_rect(slide, table_left, y, table_width, row_h,
                    fill_color=RGBColor(0xF0, 0xF4, 0xFF), shape_name=f"RowHL_{i}")

            # Total row top border
            if is_total:
                self._add_rect(slide, table_left, y, table_width, Inches(0.02),
                    fill_color=PRIMARY_NAVY, shape_name="TotalLine")

            for j, cell in enumerate(cells):
                x = int(table_left + j * col_w)
                font_color = PRIMARY_NAVY if is_total else TEXT_MAIN
                self._add_textbox(slide, x, y + Inches(0.1), int(col_w), Inches(0.4),
                    text=str(cell), font_size=Pt(13),
                    font_color=font_color, bold=is_total,
                    apply_buffer=False)

            # Row bottom line
            if not is_total:
                self._add_rect(slide, table_left, y + row_h, table_width, Inches(0.005),
                    fill_color=RGBColor(0xE5, 0xE9, 0xF0), shape_name=f"RowLine_{i}")

        if footer_text:
            footer_y = int(data_start + len(all_rows) * row_h + Inches(0.3))
            self._add_textbox(slide, Inches(1.17), footer_y, Inches(17.68), Inches(0.4),
                text=footer_text, font_size=Pt(10), font_color=TEXT_MUTED,
                alignment=PP_ALIGN.CENTER, apply_buffer=False, italic=True)

        return slide

    def add_statement_slide(self, text: str, highlight_words: list = None):
        """Large centered statement text. highlight_words shown in action blue."""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)

        # For PPTX we can't do inline color highlighting easily,
        # so render the full text in navy, bold, centered
        self._add_textbox(slide, Inches(2.0), Inches(3.0), Inches(16.0), Inches(5.0),
            text=text, font_size=Pt(36), font_color=TEXT_MAIN,
            bold=True, alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE, apply_buffer=False)

        return slide

    # ──────────────────────────────────────────
    # Alert / Flag Cards (severity-coded)
    # ──────────────────────────────────────────

    def add_alert_cards_slide(self, title: str, subtitle: str = "", section_label: str = "",
                               alerts: list = None):
        """Color-coded alert cards with left border accent.
        Each alert: {severity: red|amber|blue|green, title, body}"""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)
        alerts = alerts or []

        severity_colors = {
            'red': (SEMANTIC_RED, RGBColor(0xFD, 0xF0, 0xF0)),
            'amber': (RGBColor(0xD9, 0x77, 0x06), RGBColor(0xFE, 0xF9, 0xE7)),
            'blue': (ACTION_BLUE, RGBColor(0xEE, 0xF2, 0xFF)),
            'green': (SUCCESS_GREEN, RGBColor(0xEA, 0xFA, 0xF1)),
        }

        if section_label:
            self._add_textbox(slide, Inches(1.17), Inches(1.14), Inches(5.0), Inches(0.3),
                text=section_label.upper(), font_size=SIZE_LABEL, font_color=TEXT_MUTED)

        self._add_textbox(slide, Inches(1.17), Inches(1.5), Inches(17.68), Inches(0.76),
            text=title, font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True)

        if subtitle:
            self._add_textbox(slide, Inches(1.17), Inches(2.35), Inches(17.58), Inches(0.4),
                text=subtitle, font_size=SIZE_CAPTION, font_color=TEXT_MUTED)

        card_h = Inches(1.2)
        card_gap = Inches(0.15)
        start_y = Inches(3.2)

        for i, a in enumerate(alerts):
            sev = a.get('severity', 'blue')
            accent_color, bg_color = severity_colors.get(sev, severity_colors['blue'])
            y = int(start_y + i * (card_h + card_gap))

            # Card background
            self._add_rounded_rect(slide, Inches(1.17), y, Inches(17.68), card_h,
                fill_color=bg_color, shape_name=f"Alert_{i}")
            # Left accent bar
            self._add_rect(slide, Inches(1.17), y, Inches(0.06), card_h,
                fill_color=accent_color, shape_name=f"AlertBar_{i}")

            # Severity pill
            label = sev.upper()
            self._add_textbox(slide, Inches(1.5), y + Inches(0.15), Inches(1.0), Inches(0.3),
                text=label, font_size=Pt(9), font_color=accent_color,
                bold=True, apply_buffer=False)

            # Title
            self._add_textbox(slide, Inches(2.6), y + Inches(0.15), Inches(15.0), Inches(0.3),
                text=a.get('title', ''), font_size=Pt(14), font_color=TEXT_MAIN,
                bold=True, apply_buffer=False)

            # Body
            self._add_textbox(slide, Inches(1.5), y + Inches(0.5), Inches(16.0), Inches(0.6),
                text=a.get('body', ''), font_size=Pt(11), font_color=TEXT_MAIN,
                apply_buffer=False)

        return slide

    # ──────────────────────────────────────────
    # Architecture Stack (layered rows with pills)
    # ──────────────────────────────────────────

    def add_architecture_stack_slide(self, title: str, subtitle: str = "", section_label: str = "",
                                      layers: list = None):
        """Layered architecture diagram with pills.
        Each layer: {label, items: [{name}], bg_hex, dark, accent_hex}"""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)
        layers = layers or []

        if section_label:
            self._add_textbox(slide, Inches(1.17), Inches(1.14), Inches(5.0), Inches(0.3),
                text=section_label.upper(), font_size=SIZE_LABEL, font_color=TEXT_MUTED)

        self._add_textbox(slide, Inches(1.17), Inches(1.5), Inches(17.68), Inches(0.76),
            text=title, font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True)

        if subtitle:
            self._add_textbox(slide, Inches(1.17), Inches(2.35), Inches(17.58), Inches(0.4),
                text=subtitle, font_size=SIZE_CAPTION, font_color=TEXT_MUTED)

        layer_h = Inches(1.2)
        layer_gap = Inches(0.12)
        start_y = Inches(3.2)

        for li, layer in enumerate(layers):
            y = int(start_y + li * (layer_h + layer_gap))
            is_dark = layer.get('dark', False)
            bg_hex = layer.get('bg_hex', '#F0F4FF' if not is_dark else '#001C3D')
            bg_r, bg_g, bg_b = int(bg_hex[1:3], 16), int(bg_hex[3:5], 16), int(bg_hex[5:7], 16)
            bg_color = RGBColor(bg_r, bg_g, bg_b)
            text_c = SURFACE_WHITE if is_dark else TEXT_MAIN
            accent_hex = layer.get('accent_hex', '#1A5AFF')
            acc_r, acc_g, acc_b = int(accent_hex[1:3], 16), int(accent_hex[3:5], 16), int(accent_hex[5:7], 16)
            accent_color = RGBColor(acc_r, acc_g, acc_b)

            # Layer background
            self._add_rounded_rect(slide, Inches(1.17), y, Inches(17.68), layer_h,
                fill_color=bg_color, shape_name=f"Layer_{li}")
            # Left accent bar
            self._add_rect(slide, Inches(1.17), y, Inches(0.05), layer_h,
                fill_color=accent_color, shape_name=f"LayerBar_{li}")

            # Label
            self._add_textbox(slide, Inches(1.5), y + Inches(0.1), Inches(6.0), Inches(0.25),
                text=layer.get('label', '').upper(), font_size=Pt(10), font_color=text_c,
                bold=True, apply_buffer=False)

            # Pills
            items = layer.get('items', [])
            pill_w = Inches(2.8)
            pill_h = Inches(0.4)
            pill_gap = Inches(0.15)
            pills_per_row = 5
            pill_start_x = Inches(1.5)
            pill_start_y = y + Inches(0.45)

            for pi, item in enumerate(items):
                col = pi % pills_per_row
                row = pi // pills_per_row
                px = int(pill_start_x + col * (pill_w + pill_gap))
                py = int(pill_start_y + row * (pill_h + pill_gap))
                name = item if isinstance(item, str) else item.get('name', '')

                pill_fill = SURFACE_WHITE if not is_dark else RGBColor(0x00, 0x2E, 0x5A)
                self._add_rounded_rect(slide, px, py, pill_w, pill_h,
                    fill_color=pill_fill, shape_name=f"Pill_{li}_{pi}")
                self._add_textbox(slide, px + Inches(0.1), py, pill_w - Inches(0.2), pill_h,
                    text=name, font_size=Pt(11), font_color=text_c,
                    anchor=MSO_ANCHOR.MIDDLE, apply_buffer=False)

        return slide

    # ──────────────────────────────────────────
    # Pillar → Platform Map (pillars left, boxes right)
    # ──────────────────────────────────────────

    def add_pillar_platform_map_slide(self, title: str, subtitle: str = "", section_label: str = "",
                                       pillars: list = None, footer_left: str = "", footer_right: str = ""):
        """Pillar-to-platform mapping: pillars on left with arrows → platform boxes on right.
        Each pillar: {name, sub?, accent (hex), section_above?, items: [{name, sub}]}
        """
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)
        pillars = pillars or []

        if section_label:
            self._add_textbox(slide, Inches(1.17), Inches(1.14), Inches(5.0), Inches(0.3),
                text=section_label.upper(), font_size=SIZE_LABEL, font_color=TEXT_MUTED)
        self._add_textbox(slide, Inches(1.17), Inches(1.5), Inches(17.68), Inches(0.76),
            text=title, font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True)
        if subtitle:
            self._add_textbox(slide, Inches(1.17), Inches(2.35), Inches(17.58), Inches(0.4),
                text=subtitle, font_size=SIZE_CAPTION, font_color=TEXT_MUTED)

        pillar_x = Inches(1.17)
        pillar_w = Inches(3.0)
        arrow_x = Inches(4.4)
        box_start_x = Inches(5.2)
        box_w = Inches(2.8)
        box_gap = Inches(0.15)
        row_h = Inches(1.3)
        row_gap = Inches(0.15)
        start_y = Inches(3.2)

        for ri, pillar in enumerate(pillars):
            y = int(start_y + ri * (row_h + row_gap))
            accent_hex = pillar.get('accent', '#1A5AFF')
            ar, ag, ab = int(accent_hex[1:3], 16), int(accent_hex[3:5], 16), int(accent_hex[5:7], 16)
            accent_c = RGBColor(ar, ag, ab)

            # Row background
            self._add_rounded_rect(slide, pillar_x, y, Inches(17.68), row_h,
                fill_color=BG_GRAY, shape_name=f"PillarRow_{ri}")
            # Left accent bar
            self._add_rect(slide, pillar_x, y, Inches(0.06), row_h,
                fill_color=accent_c, shape_name=f"PillarBar_{ri}")

            # Section label above items (e.g. "Single Pane of Glass")
            sec_above = pillar.get('section_above', '')
            if sec_above:
                self._add_textbox(slide, box_start_x, y + Inches(0.05), Inches(12.0), Inches(0.2),
                    text=sec_above.upper(), font_size=Pt(8), font_color=TEXT_MUTED,
                    bold=True, apply_buffer=False)

            # Pillar label
            plabel = pillar.get('name', '')
            psub = pillar.get('sub', '')
            label_text = f"{plabel}\n{psub}" if psub else plabel
            self._add_textbox(slide, pillar_x + Inches(0.2), y + Inches(0.15), pillar_w - Inches(0.3), Inches(0.5),
                text=plabel, font_size=Pt(12), font_color=accent_c, bold=True, apply_buffer=False)
            if psub:
                self._add_textbox(slide, pillar_x + Inches(0.2), y + Inches(0.45), pillar_w - Inches(0.3), Inches(0.3),
                    text=psub, font_size=Pt(9), font_color=TEXT_MUTED, apply_buffer=False)

            # Arrow
            self._add_textbox(slide, arrow_x, y + Inches(0.35), Inches(0.5), Inches(0.3),
                text='\u2192', font_size=Pt(16), font_color=accent_c,
                apply_buffer=False, alignment=PP_ALIGN.CENTER)

            # Platform boxes
            items = pillar.get('items', [])
            items_y = y + Inches(0.3) if not sec_above else y + Inches(0.4)
            box_h = Inches(0.65)

            if len(items) == 1:
                # Full-width single box (orchestration layer, Grand Central)
                self._add_rounded_rect(slide, box_start_x, items_y, Inches(13.5), box_h,
                    fill_color=SURFACE_WHITE, shape_name=f"Box_{ri}_0")
                # Border effect via outline
                shape = slide.shapes[-1]
                shape.line.color.rgb = accent_c
                shape.line.width = Pt(1.5)
                self._add_textbox(slide, box_start_x + Inches(0.15), items_y + Inches(0.05),
                    Inches(13.0), Inches(0.25),
                    text=items[0].get('name', ''), font_size=Pt(11), font_color=TEXT_MAIN,
                    bold=True, apply_buffer=False)
                self._add_textbox(slide, box_start_x + Inches(0.15), items_y + Inches(0.3),
                    Inches(13.0), Inches(0.25),
                    text=items[0].get('sub', ''), font_size=Pt(8), font_color=TEXT_MUTED,
                    apply_buffer=False)
            else:
                for bi, item in enumerate(items):
                    bx = int(box_start_x + bi * (box_w + box_gap))
                    self._add_rounded_rect(slide, bx, items_y, box_w, box_h,
                        fill_color=SURFACE_WHITE, shape_name=f"Box_{ri}_{bi}")
                    shape = slide.shapes[-1]
                    shape.line.color.rgb = accent_c
                    shape.line.width = Pt(1)
                    self._add_textbox(slide, bx + Inches(0.1), items_y + Inches(0.08),
                        box_w - Inches(0.2), Inches(0.25),
                        text=item.get('name', ''), font_size=Pt(11), font_color=TEXT_MAIN,
                        bold=True, apply_buffer=False, alignment=PP_ALIGN.CENTER)
                    self._add_textbox(slide, bx + Inches(0.1), items_y + Inches(0.35),
                        box_w - Inches(0.2), Inches(0.2),
                        text=item.get('sub', ''), font_size=Pt(8), font_color=TEXT_MUTED,
                        apply_buffer=False, alignment=PP_ALIGN.CENTER)

        # Footer boxes
        footer_y = int(start_y + len(pillars) * (row_h + row_gap) + Inches(0.15))
        if footer_left:
            self._add_rounded_rect(slide, Inches(1.17), footer_y, Inches(8.5), Inches(0.6),
                fill_color=RGBColor(0xF0, 0xF4, 0xFF), shape_name="FooterLeft")
            self._add_textbox(slide, Inches(1.4), footer_y + Inches(0.05), Inches(8.0), Inches(0.5),
                text=f"Schroders retains: {footer_left}", font_size=Pt(9), font_color=TEXT_MAIN,
                apply_buffer=False)
        if footer_right:
            self._add_rounded_rect(slide, Inches(10.0), footer_y, Inches(8.85), Inches(0.6),
                fill_color=RGBColor(0xFF, 0xF5, 0xF5), shape_name="FooterRight")
            self._add_textbox(slide, Inches(10.2), footer_y + Inches(0.05), Inches(8.4), Inches(0.5),
                text=f"What gets retired: {footer_right}", font_size=Pt(9),
                font_color=RGBColor(0xDC, 0x26, 0x26), apply_buffer=False)

        return slide

    # ──────────────────────────────────────────
    # Architecture Boxes (Current → Target side-by-side)
    # ──────────────────────────────────────────

    def add_architecture_boxes_slide(self, title: str, subtitle: str = "", section_label: str = "",
                                      left_title: str = "", right_title: str = "",
                                      left_layers: list = None, right_layers: list = None,
                                      legend: list = None, footer_text: str = ""):
        """Side-by-side architecture with box diagrams (current vs target).
        Each layer: {label, items: [{name, sub, bg?, border?, color?, dark?}], banner?, banner_bg?, ...}
        """
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)
        left_layers = left_layers or []
        right_layers = right_layers or []

        if section_label:
            self._add_textbox(slide, Inches(1.17), Inches(1.14), Inches(5.0), Inches(0.3),
                text=section_label.upper(), font_size=SIZE_LABEL, font_color=TEXT_MUTED)
        self._add_textbox(slide, Inches(1.17), Inches(1.5), Inches(17.68), Inches(0.76),
            text=title, font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True)
        if subtitle:
            self._add_textbox(slide, Inches(1.17), Inches(2.35), Inches(17.58), Inches(0.4),
                text=subtitle, font_size=SIZE_CAPTION, font_color=TEXT_MUTED)

        half_w = Inches(8.5)
        left_x = Inches(1.17)
        right_x = Inches(10.3)
        start_y = Inches(3.2)

        # Column headers
        self._add_textbox(slide, left_x, start_y - Inches(0.4), half_w, Inches(0.3),
            text=left_title.upper(), font_size=Pt(11), font_color=SEMANTIC_RED,
            bold=True, apply_buffer=False, alignment=PP_ALIGN.CENTER)
        self._add_textbox(slide, right_x, start_y - Inches(0.4), half_w, Inches(0.3),
            text=right_title.upper(), font_size=Pt(11), font_color=ACTION_BLUE,
            bold=True, apply_buffer=False, alignment=PP_ALIGN.CENTER)

        # Divider line
        self._add_rect(slide, Inches(9.85), start_y - Inches(0.3), Inches(0.02), Inches(7.0),
            fill_color=RGBColor(0xE0, 0xE4, 0xE8), shape_name="Divider")

        def _render_layers(layers, base_x, base_y, side):
            cy = base_y
            for li, layer in enumerate(layers):
                label = layer.get('label', '')
                items = layer.get('items', [])
                is_banner = layer.get('banner', False)

                if label:
                    label_color = SEMANTIC_RED if side == 'left' else ACTION_BLUE
                    self._add_textbox(slide, base_x, cy, half_w, Inches(0.22),
                        text=label.upper(), font_size=Pt(8), font_color=label_color,
                        bold=True, apply_buffer=False)
                    cy += Inches(0.25)

                if is_banner:
                    # Full-width banner box
                    b_bg = layer.get('banner_bg', '#FFF5F5')
                    b_border = layer.get('banner_border', '#FCA5A5')
                    b_text_c = layer.get('banner_text_color', '#DC2626')
                    br, bg_val, bb = int(b_bg[1:3], 16), int(b_bg[3:5], 16), int(b_bg[5:7], 16)
                    bdr, bdg, bdb = int(b_border[1:3], 16), int(b_border[3:5], 16), int(b_border[5:7], 16)
                    btr, btg, btb = int(b_text_c[1:3], 16), int(b_text_c[3:5], 16), int(b_text_c[5:7], 16)

                    shape = self._add_rounded_rect(slide, base_x, cy, half_w, Inches(0.55),
                        fill_color=RGBColor(br, bg_val, bb), shape_name=f"Banner_{side}_{li}")
                    shape.line.color.rgb = RGBColor(bdr, bdg, bdb)
                    shape.line.width = Pt(1.5)
                    shape.line.dash_style = 2  # dash

                    item = items[0] if items else {}
                    self._add_textbox(slide, base_x + Inches(0.2), cy + Inches(0.05),
                        half_w - Inches(0.4), Inches(0.2),
                        text=item.get('name', ''), font_size=Pt(10), font_color=RGBColor(btr, btg, btb),
                        bold=True, apply_buffer=False, alignment=PP_ALIGN.CENTER)
                    self._add_textbox(slide, base_x + Inches(0.2), cy + Inches(0.28),
                        half_w - Inches(0.4), Inches(0.2),
                        text=item.get('sub', ''), font_size=Pt(7), font_color=TEXT_MUTED,
                        apply_buffer=False, alignment=PP_ALIGN.CENTER)
                    cy += Inches(0.65)
                else:
                    # Row of item boxes
                    n = len(items)
                    if n == 0:
                        continue
                    item_gap = Inches(0.12)
                    item_w = int((half_w - (n - 1) * item_gap) / n) if n > 0 else half_w
                    item_h = Inches(0.55)

                    for ii, item in enumerate(items):
                        ix = int(base_x + ii * (item_w + item_gap))
                        i_bg = item.get('bg', '#EBF0FF' if side == 'right' else '#F5F7F9')
                        i_border = item.get('border', '#1A5AFF' if side == 'right' else '#D1D5DB')
                        i_color = item.get('color', '#001C3D')
                        is_dark = item.get('dark', False)

                        ibr, ibg, ibb = int(i_bg[1:3], 16), int(i_bg[3:5], 16), int(i_bg[5:7], 16)
                        idr, idg, idb = int(i_border[1:3], 16), int(i_border[3:5], 16), int(i_border[5:7], 16)
                        icr, icg, icb = int(i_color[1:3], 16), int(i_color[3:5], 16), int(i_color[5:7], 16)

                        fill_c = RGBColor(0x00, 0x1C, 0x3D) if is_dark else RGBColor(ibr, ibg, ibb)
                        text_c = SURFACE_WHITE if is_dark else RGBColor(icr, icg, icb)

                        shape = self._add_rounded_rect(slide, ix, cy, item_w, item_h,
                            fill_color=fill_c, shape_name=f"Item_{side}_{li}_{ii}")
                        shape.line.color.rgb = RGBColor(idr, idg, idb)
                        shape.line.width = Pt(1)

                        self._add_textbox(slide, ix + Inches(0.05), cy + Inches(0.05),
                            item_w - Inches(0.1), Inches(0.2),
                            text=item.get('name', ''), font_size=Pt(9), font_color=text_c,
                            bold=True, apply_buffer=False, alignment=PP_ALIGN.CENTER)
                        self._add_textbox(slide, ix + Inches(0.05), cy + Inches(0.28),
                            item_w - Inches(0.1), Inches(0.2),
                            text=item.get('sub', ''), font_size=Pt(7),
                            font_color=SURFACE_WHITE if is_dark else TEXT_MUTED,
                            apply_buffer=False, alignment=PP_ALIGN.CENTER)

                    cy += Inches(0.7)

            return cy

        _render_layers(left_layers, left_x, start_y, 'left')
        _render_layers(right_layers, right_x, start_y, 'right')

        # Arrow between columns
        self._add_textbox(slide, Inches(9.5), Inches(5.5), Inches(0.7), Inches(0.4),
            text='\u2192', font_size=Pt(24), font_color=TEXT_MUTED,
            apply_buffer=False, alignment=PP_ALIGN.CENTER)

        # Legend
        if legend:
            leg_y = Inches(9.2)
            leg_x = Inches(1.5)
            for i, item in enumerate(legend):
                lx = int(leg_x + i * Inches(3.2))
                lbl = item.get('label', '')
                color_hex = item.get('color', '#EBF0FF')
                border_hex = item.get('border', '')
                cr, cg, cb = int(color_hex[1:3], 16), int(color_hex[3:5], 16), int(color_hex[5:7], 16)
                swatch = self._add_rect(slide, lx, leg_y, Inches(0.2), Inches(0.2),
                    fill_color=RGBColor(cr, cg, cb), shape_name=f"Legend_{i}")
                if border_hex:
                    dr, dg, db = int(border_hex[1:3], 16), int(border_hex[3:5], 16), int(border_hex[5:7], 16)
                    swatch.line.color.rgb = RGBColor(dr, dg, db)
                    swatch.line.width = Pt(1)
                self._add_textbox(slide, lx + Inches(0.3), leg_y, Inches(2.5), Inches(0.2),
                    text=lbl, font_size=Pt(8), font_color=TEXT_MUTED, apply_buffer=False)

        # Footer
        if footer_text:
            self._add_rounded_rect(slide, Inches(1.17), Inches(9.6), Inches(17.68), Inches(0.6),
                fill_color=RGBColor(0xF0, 0xF4, 0xFF), shape_name="ArchFooter")
            self._add_textbox(slide, Inches(1.4), Inches(9.65), Inches(17.2), Inches(0.5),
                text=footer_text, font_size=Pt(9), font_color=TEXT_MAIN,
                apply_buffer=False, alignment=PP_ALIGN.CENTER)

        return slide

    # ──────────────────────────────────────────
    # Comparison Matrix (criteria left, options across)
    # ──────────────────────────────────────────

    def add_comparison_matrix_slide(self, title: str, subtitle: str = "", section_label: str = "",
                                     options: list = None, criteria: list = None,
                                     summaries: list = None, footer_text: str = ""):
        """Comparison matrix: criteria rows × option columns."""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)
        options = options or []
        criteria = criteria or []
        summaries = summaries or []

        accent_map = {"red": RGBColor(0xDC, 0x26, 0x26), "blue": ACTION_BLUE,
                       "green": RGBColor(0x16, 0xA3, 0x4A), "amber": RGBColor(0xD9, 0x77, 0x06),
                       "cyan": RGBColor(0x08, 0x91, 0xB2), "purple": RGBColor(0x7C, 0x3A, 0xED)}

        if section_label:
            self._add_textbox(slide, Inches(1.17), Inches(1.14), Inches(5.0), Inches(0.3),
                text=section_label.upper(), font_size=SIZE_LABEL, font_color=TEXT_MUTED)
        self._add_textbox(slide, Inches(1.17), Inches(1.5), Inches(17.68), Inches(0.76),
            text=title, font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True)
        if subtitle:
            self._add_textbox(slide, Inches(1.17), Inches(2.35), Inches(17.58), Inches(0.4),
                text=subtitle, font_size=SIZE_CAPTION, font_color=TEXT_MUTED)

        n = len(options)
        label_w = Inches(2.8)
        start_x = Inches(1.17)
        col_w = Inches(4.8) if n == 3 else Inches(14.0 / n)
        cols_start = start_x + label_w + Inches(0.2)
        start_y = Inches(3.2)

        # Option headers
        for i, opt in enumerate(options):
            x = int(cols_start + i * col_w)
            ac = accent_map.get(opt.get('accent', 'blue'), ACTION_BLUE)
            self._add_textbox(slide, x, start_y, col_w, Inches(0.35),
                text=opt['name'].upper(), font_size=Pt(10), font_color=ac,
                bold=True, alignment=PP_ALIGN.CENTER, apply_buffer=False)
            # Underline
            line_w = col_w - Inches(0.3)
            self._add_rect(slide, x + Inches(0.15), start_y + Inches(0.35),
                line_w, Inches(0.025), fill_color=ac)

        # Criteria rows
        row_y = start_y + Inches(0.6)
        row_h = Inches(0.85)

        for ri, cr in enumerate(criteria):
            y = int(row_y + ri * row_h)
            # Alternating bg
            if ri % 2 == 0:
                self._add_rect(slide, start_x, y, Inches(17.68), row_h,
                    fill_color=BG_GRAY, shape_name=f"RowBg_{ri}")

            # Criteria label
            self._add_textbox(slide, start_x + Inches(0.15), y + Inches(0.1),
                label_w - Inches(0.3), row_h - Inches(0.2),
                text=cr['label'], font_size=Pt(10), font_color=TEXT_MAIN,
                bold=True, apply_buffer=False)

            # Cell values
            for ci, cell_text in enumerate(cr.get('cells', [])):
                x = int(cols_start + ci * col_w)
                # Strip ** markers for PPTX (plain text)
                clean = cell_text.replace('**', '')
                self._add_textbox(slide, x + Inches(0.1), y + Inches(0.05),
                    col_w - Inches(0.2), row_h - Inches(0.1),
                    text=clean, font_size=Pt(9), font_color=TEXT_MAIN,
                    apply_buffer=False)

        # Summary boxes
        if summaries:
            sum_y = int(row_y + len(criteria) * row_h + Inches(0.2))
            for si, sm in enumerate(summaries):
                x = int(cols_start + si * col_w)
                ac = accent_map.get(sm.get('accent', 'blue'), ACTION_BLUE)
                # Box with light accent bg
                r_val = ac.red if hasattr(ac, 'red') else 0x1A
                g_val = ac.green if hasattr(ac, 'green') else 0x5A
                b_val = ac.blue if hasattr(ac, 'blue') else 0xFF
                light_bg = RGBColor(min(255, r_val + (255 - r_val) * 9 // 10),
                                     min(255, g_val + (255 - g_val) * 9 // 10),
                                     min(255, b_val + (255 - b_val) * 9 // 10))
                shape = self._add_rounded_rect(slide, x + Inches(0.1), sum_y,
                    col_w - Inches(0.2), Inches(0.7), fill_color=light_bg)
                shape.line.color.rgb = ac
                shape.line.width = Pt(1.5)
                self._add_textbox(slide, x + Inches(0.2), sum_y + Inches(0.05),
                    col_w - Inches(0.4), Inches(0.25),
                    text=sm.get('label', '').upper(), font_size=Pt(8), font_color=ac,
                    bold=True, alignment=PP_ALIGN.CENTER, apply_buffer=False)
                self._add_textbox(slide, x + Inches(0.2), sum_y + Inches(0.3),
                    col_w - Inches(0.4), Inches(0.3),
                    text=sm.get('sub', ''), font_size=Pt(11), font_color=TEXT_MAIN,
                    bold=True, alignment=PP_ALIGN.CENTER, apply_buffer=False)

        # Footer
        if footer_text:
            foot_y = int(row_y + len(criteria) * row_h + Inches(1.2))
            self._add_rounded_rect(slide, Inches(1.17), foot_y, Inches(17.68), Inches(0.6),
                fill_color=RGBColor(0xF0, 0xF4, 0xFF))
            clean_footer = footer_text.replace('**', '')
            self._add_textbox(slide, Inches(1.4), foot_y + Inches(0.05), Inches(17.2), Inches(0.5),
                text=clean_footer, font_size=Pt(9), font_color=TEXT_MAIN,
                apply_buffer=False, alignment=PP_ALIGN.CENTER)

        return slide

    # ──────────────────────────────────────────
    # Clean Stat Cards (lightweight)
    # ──────────────────────────────────────────

    def add_context_stats_slide(self, title: str, subtitle: str = "", section_label: str = "",
                                 body: str = "", stats: list = None):
        """Context section with body text and clean stat cards.
        Each stat: {number, label}"""
        slide = self._add_blank_slide()
        self._set_slide_bg(slide, SURFACE_WHITE)
        stats = stats or []

        if section_label:
            self._add_textbox(slide, Inches(1.17), Inches(1.14), Inches(5.0), Inches(0.3),
                text=section_label.upper(), font_size=SIZE_LABEL, font_color=TEXT_MUTED)

        self._add_textbox(slide, Inches(1.17), Inches(1.5), Inches(17.68), Inches(0.76),
            text=title, font_size=SIZE_H2, font_color=TEXT_MAIN, bold=True)

        if subtitle:
            self._add_textbox(slide, Inches(1.17), Inches(2.35), Inches(17.58), Inches(0.4),
                text=subtitle, font_size=SIZE_CAPTION, font_color=TEXT_MUTED)

        body_y = Inches(3.2)
        if body:
            self._add_textbox(slide, Inches(1.17), body_y, Inches(16.0), Inches(2.0),
                text=body, font_size=SIZE_BODY, font_color=TEXT_MAIN, apply_buffer=False)
            body_y = Inches(5.5)

        # Stat cards
        n = len(stats)
        card_gap = Inches(0.2)
        total_w = Inches(17.68)
        card_w = int((total_w - card_gap * (n - 1)) / n) if n > 0 else total_w
        card_h = Inches(2.5)

        for i, s in enumerate(stats):
            x = int(Inches(1.17) + i * (card_w + card_gap))
            self._add_rounded_rect(slide, x, body_y, card_w, card_h,
                fill_color=BG_GRAY, shape_name=f"CtxStat_{i}")
            self._add_textbox(slide, x + Inches(0.3), body_y + Inches(0.4), card_w - Inches(0.6), Inches(1.0),
                text=str(s.get('number', '')), font_size=Pt(36), font_color=ACTION_BLUE,
                bold=True, apply_buffer=False)
            self._add_textbox(slide, x + Inches(0.3), body_y + Inches(1.5), card_w - Inches(0.6), Inches(0.5),
                text=s.get('label', ''), font_size=SIZE_CAPTION, font_color=TEXT_MUTED,
                apply_buffer=False)

        return slide

    # ──────────────────────────────────────────
    # Save
    # ──────────────────────────────────────────

    def save(self, path: str = None):
        """Save the presentation to disk. Updates slide numbers to include total."""
        total = len(self.slides)
        # Update all slide number text boxes with "N / Total" format
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.name and shape.name.startswith("SlideNumber_"):
                    num = shape.name.split("_")[1]
                    tf = shape.text_frame
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.text = f"{num} / {total}"

        out = path or self.output_path
        self.prs.save(out)
        return out


# ──────────────────────────────────────────────
# Quick test / CLI usage
# ──────────────────────────────────────────────

if __name__ == "__main__":
    p = Frontline2026Presenter("frontline_2026_sample.pptx")

    p.add_cover_slide("Introduction", "AI-Native | Banking OS", "March 2026")

    p.add_agenda_slide(
        "AI-NATIVE BANKING OS", "Agenda", "Backbase x <Customer>",
        ["Unified Frontline", "Modernizing Segments + Channels",
         "Banking OS", "Customer Cases", "Next Steps"]
    )

    p.add_section_divider(
        "INTRODUCTION TO BACKBASE", "Unified Frontline",
        "What problem are we trying to solve?"
    )

    p.add_content_slide(
        "Challenge \u2022 50% of Banking work lives between systems",
        "The expensive whitespace between systems: handoffs, exceptions, manual coordination."
    )

    p.add_split_comparison(
        "Trusted partner for 100+ Banks",
        section_label="UNIFIED FRONTLINE",
        left_title="Fragmented Frontline",
        left_items=[
            "20-40 disconnected apps, workflows + tools",
            "Manual steps everywhere \u2192 high cost-to-serve",
            "Banker spends 60% of time on non-value activities"
        ],
        right_title="Unified Frontline",
        right_items=[
            "Consistent, modern experiences across all channels",
            "3\u00d7 faster change velocity with reusable components",
            "AI-augmented workflows reduce manual effort by 40%"
        ]
    )

    p.add_architecture_slide(
        "Banking OS for the AI era",
        "Where we play; one platform to unify + orchestrate the banks frontline business."
    )

    p.add_stat_cards_slide(
        "Impact \u2022 Measurable business outcomes",
        "Results from leading banking transformations.",
        stats=[
            {"number": "3\u00d7", "label": "Faster time-to-market", "trend": "+200%"},
            {"number": "40%", "label": "Reduction in cost-to-serve", "trend": "\u2193 OpEx"},
            {"number": "60%", "label": "Less manual processing", "trend": "+Automation"},
        ]
    )

    p.add_case_study_slide(
        "Global Bank \u2014 Unified Frontline Transformation",
        body_lines=[
            "Replaced 35 legacy applications with a single unified platform.",
            "Reduced customer onboarding time from 14 days to 2 days.",
            "Achieved 98% straight-through processing for standard requests.",
            "Enabled AI-assisted advisory for relationship managers."
        ]
    )

    out = p.save()
    print(f"Saved: {out}")
