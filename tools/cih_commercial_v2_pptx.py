#!/usr/bin/env python3
"""
CIH Bank — Commercial Proposal PPTX Generator (v2)

Redesigned to match CIH's own internal template style:
  - 13.333" x 7.5" widescreen (Google Slides compatible)
  - Large, readable text (H1 20pt, H2 14pt, body 12pt)
  - Colored header boxes (#CFE2F3 light blue)
  - Clean white content slides with professional tables
  - CIH Bank logo on cover
  - Backbase logo + page number footer
  - Blue gradient section dividers

Usage:
    python3 tools/cih_commercial_v2_pptx.py
"""

import sys
sys.path.insert(0, '/Users/shyam/cortex')

from pathlib import Path
from tools.pptx_presenter import PptxPresenter
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ASSETS = Path('/Users/shyam/cortex/Engagement/CIH Bank/assets')
CIH_LOGO = str(ASSETS / 'cih_bank_logo.png')
BLUE_GRADIENT = str(ASSETS / 'blue_gradient_bg.png')

# Colors matching CIH reference template
HEADER_BOX_FILL = RGBColor(0xCF, 0xE2, 0xF3)   # light blue header boxes
GREEN_BADGE     = RGBColor(0x6A, 0xA8, 0x4F)    # "recommended" badge
LIGHT_GRAY_BG   = RGBColor(0xF8, 0xFA, 0xFC)    # alternating row


class CIHCommercialV2(PptxPresenter):
    """CIH Bank commercial proposal — styled to match CIH internal template."""

    # Override margins — slightly tighter to match reference
    ML = Inches(0.55)
    CW = Inches(12.2)

    # ── Helpers specific to CIH template style ──────────

    def _title_bar(self, slide, title, subtitle=''):
        """H1 title at top of content slide — 28pt bold, left-aligned."""
        self._txt(slide, title, self.ML, Inches(0.45), self.CW, Inches(0.55),
                  size=Pt(28), bold=True, color=self.DARK_TEXT)
        if subtitle:
            self._txt(slide, subtitle, self.ML, Inches(1.05), self.CW, Inches(0.3),
                      size=Pt(12), color=self.SUB_TEXT)

    def _section_header_box(self, slide, text, top=Inches(1.45), width=None):
        """Colored rectangle header box (light blue fill, bold text) — CIH style."""
        w = width or Inches(6.0)
        h = Inches(0.4)
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, self.ML, top, w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = HEADER_BOX_FILL
        box.line.fill.background()
        self._txt(slide, text, self.ML + Inches(0.15), top + Inches(0.05),
                  w - Inches(0.3), Inches(0.3),
                  size=Pt(13), bold=True, color=self.DARK_TEXT)

    def _body_text(self, slide, text, left, top, width, height):
        """Body text — 16pt, readable, CIH style."""
        self._txt(slide, text, left, top, width, height,
                  size=Pt(16), color=self.DARK_TEXT)

    def _cih_table(self, slide, rows, col_widths, left, top, row_height=Inches(0.38)):
        """CIH-style table — dark header, alternating rows, 10pt text."""
        return self._add_table(slide, rows, col_widths, left, top,
                               row_height=row_height,
                               header_bg=self.DARK_BG,
                               header_color=self.WHITE,
                               body_size=Pt(10))

    def _divider_slide(self, title, subtitle='', slide_number=1):
        """Blue gradient section divider — white text centered."""
        s = self._new_slide(dark=True)
        # Blue gradient background image
        if Path(BLUE_GRADIENT).exists():
            s.shapes.add_picture(BLUE_GRADIENT, 0, 0, self.SLIDE_W, self.SLIDE_H)
        self._txt(s, title,
                  Inches(1.5), Inches(2.5), Inches(10), Inches(1.0),
                  size=Pt(32), bold=True, color=self.WHITE,
                  align=PP_ALIGN.LEFT)
        if subtitle:
            self._txt(s, subtitle,
                      Inches(1.5), Inches(3.6), Inches(8), Inches(0.8),
                      size=Pt(14), color=RGBColor(0xCC, 0xDD, 0xFF),
                      align=PP_ALIGN.LEFT)
        self._footer(s, slide_number, dark=True)
        return s

    # ── Slide builders ──────────────────────────────────

    def generate(self, output_path):
        self._init_presentation()
        n = 0

        # ── 1. COVER ────────────────────────────────────
        n += 1
        s = self._new_slide(dark=False)
        # White background with subtle left-aligned content
        bg = s.background.fill
        bg.solid()
        bg.fore_color.rgb = self.WHITE

        # CIH Bank logo — top right
        if Path(CIH_LOGO).exists():
            s.shapes.add_picture(CIH_LOGO, Inches(9.5), Inches(0.8),
                                 Inches(2.1), Inches(0.72))

        # Title — left aligned, large
        self._txt(s, 'Commercial Proposal',
                  Inches(0.8), Inches(3.0), Inches(7), Inches(0.6),
                  size=Pt(32), bold=True, color=self.DARK_TEXT)
        self._txt(s, 'for CIH Bank',
                  Inches(0.8), Inches(3.6), Inches(7), Inches(0.5),
                  size=Pt(28), bold=True, color=self.BLUE)

        # Date
        self._txt(s, 'March 2026',
                  Inches(0.8), Inches(4.4), Inches(4), Inches(0.3),
                  size=Pt(16), color=self.SUB_TEXT)

        # Scope line
        self._txt(s, 'Managed Hosting  ·  Digital Assist Premium  ·  Customer Lifecycle Orchestrator',
                  Inches(0.8), Inches(5.0), Inches(9), Inches(0.3),
                  size=Pt(11), color=self.MUTED)

        # Confidential badge
        self._txt(s, 'CONFIDENTIAL',
                  Inches(0.8), Inches(6.6), Inches(2), Inches(0.2),
                  size=Pt(8), bold=True, color=self.MUTED)

        self._footer(s, n)

        # ── 2. SECTION: Our Partnership ─────────────────
        n += 1
        self._divider_slide('Our Partnership Journey', 'Five years of digital growth with CIH Bank', n)

        # ── 3. PARTNERSHIP TIMELINE ─────────────────────
        n += 1
        s = self._new_slide()
        self._title_bar(s, 'Partnership Timeline')
        self._section_header_box(s, 'Key Milestones', width=Inches(3.5))

        phases = [
            ('Jul 2019', 'Backbase Selected', 'Digital banking platform chosen'),
            ('2020–21', 'Core Launch', 'Retail & SME banking live'),
            ('2022–23', 'Scale', '2M+ active digital users'),
            ('Jul 2024', 'Renewal', '5-year renewal through Jul 2029'),
            ('2026', 'Next Phase', 'MH + DA + CLO upsell'),
        ]
        line_y = Inches(2.5)
        self._divider(s, line_y, color=self.BLUE)
        step_w = self.CW / len(phases)
        for i, (date, title, desc) in enumerate(phases):
            cx = self.ML + i * step_w + step_w / 2
            dot = Inches(0.12)
            color = self.GREEN if i < 4 else self.BLUE
            self._bar_rect(s, cx - dot / 2, line_y - dot / 2, dot, dot, fill=color)
            self._txt(s, date, cx - Inches(0.8), line_y + Inches(0.25),
                      Inches(1.6), Inches(0.2),
                      size=Pt(9), bold=True, color=self.DARK_TEXT, align=PP_ALIGN.CENTER)
            self._txt(s, title, cx - Inches(0.8), line_y + Inches(0.45),
                      Inches(1.6), Inches(0.2),
                      size=Pt(11), bold=True, color=color, align=PP_ALIGN.CENTER)
            self._txt(s, desc, cx - Inches(0.8), line_y + Inches(0.7),
                      Inches(1.6), Inches(0.3),
                      size=Pt(9), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)

        # Summary stats
        gap = Inches(0.25)
        sw = (self.CW - 3 * gap) / 4
        stats = [
            ('2M+', 'Active Users', self.GREEN, self.GREEN_LIGHT),
            ('7 yrs', 'Partnership', self.BLUE, self.BLUE_LIGHT),
            ('€5.85M', 'Current TCV', self.PURPLE, self.PURPLE_LIGHT),
            ('Jul 2029', 'Contract End', self.AMBER, self.AMBER_LIGHT),
        ]
        for i, (val, lbl, vc, bg) in enumerate(stats):
            self._stat_box(s, val, lbl,
                           self.ML + i * (sw + gap), Inches(4.5), sw, Inches(0.6),
                           val_color=vc, bg=bg)
        self._footer(s, n)

        # ── 4. STRATEGIC PRIORITIES ─────────────────────
        n += 1
        s = self._new_slide()
        self._title_bar(s, "CIH Bank's Strategic Priorities")
        self._section_header_box(s, 'Context', width=Inches(2.5))

        # Body paragraph — CIH context
        self._body_text(s,
            'CIH Bank aims to strengthen its operating model by moving to managed hosting, '
            'reducing infrastructure burden on internal teams and guaranteeing SLA uptime.\n\n'
            'With 2M+ active users today and a target of 3.5M by 2029, the bank needs to '
            'convert 500K wallet-only users into product holders through digital origination '
            'and cross-sell campaigns.\n\n'
            'Additionally, the Banque Directe remote advisory team of 50 agents needs '
            'AI-assisted tools to increase selling time and revenue per agent.',
            self.ML, Inches(2.0), Inches(7.5), Inches(3.5))

        # Priority cards on the right
        priorities = [
            ('Grow Digital Revenue', 'Convert wallet-only users into product holders', self.GREEN),
            ('Modernize Infrastructure', 'Managed hosting for uptime & efficiency', self.BLUE),
            ('Empower Banque Directe', 'AI tools for 50 remote advisors', self.PURPLE),
        ]
        card_x = Inches(8.5)
        card_w = Inches(4.3)
        for i, (title, desc, color) in enumerate(priorities):
            cy = Inches(2.0) + i * Inches(1.1)
            self._card(s, card_x, cy, card_w, Inches(0.85))
            self._colored_top_line(s, card_x, cy, card_w, color)
            self._txt(s, title, card_x + Inches(0.15), cy + Inches(0.1),
                      card_w - Inches(0.3), Inches(0.2),
                      size=Pt(11), bold=True, color=color)
            self._txt(s, desc, card_x + Inches(0.15), cy + Inches(0.4),
                      card_w - Inches(0.3), Inches(0.35),
                      size=Pt(9), color=self.SUB_TEXT)
        self._footer(s, n)

        # ── 5. SECTION: Scope & Pricing ─────────────────
        n += 1
        self._divider_slide('Scope & Pricing', 'Three modules, one growth platform', n)

        # ── 6. MANAGED HOSTING ──────────────────────────
        n += 1
        s = self._new_slide()
        self._title_bar(s, 'Managed Hosting Cost',
                        'Digital Banking & up to 3 Digital Onboarding flows')

        rows = [
            ['Packages', 'MH Essentials', 'MH Enterprise 1', 'MH Enterprise 2\n(Recommended)'],
            ['Description', '30 RPS min\n99.5% SLA', '90 RPS min\n99.9% SLA', '180 RPS min\n99.9% SLA'],
            ['Annual Cost (€)', '€343,200', '€553,750', '€601,750'],
            ['CIH Discount', '—', '—', '10%'],
            ['Annual Cost after Discount', '—', '—', '€541,575'],
        ]
        self._cih_table(s, rows, [3.5, 2.8, 2.8, 3.1], self.ML, Inches(1.6))

        # Recommended badge
        badge_x = self.ML + Inches(3.5 + 2.8 + 2.8 + 0.3)
        self._bar_rect(s, badge_x, Inches(1.35), Inches(2.2), Inches(0.22), fill=GREEN_BADGE)
        self._txt(s, 'RECOMMENDED', badge_x, Inches(1.36),
                  Inches(2.2), Inches(0.2),
                  size=Pt(8), bold=True, color=self.WHITE, align=PP_ALIGN.CENTER)

        # Add-ons section
        self._section_header_box(s, 'Available Add-Ons', top=Inches(4.4), width=Inches(4))
        addon_rows = [
            ['Add-On', 'Monthly Cost (€)'],
            ['Additional 90 RPS capacity', '€4,000'],
            ['10 DevOps Seats', '€2,400'],
            ['Additional Runtime Environment', '€2,500'],
        ]
        self._cih_table(s, addon_rows, [7.0, 3.0], self.ML, Inches(4.9))
        self._footer(s, n)

        # ── 7. DIGITAL ASSIST PREMIUM ───────────────────
        n += 1
        s = self._new_slide()
        self._title_bar(s, 'Digital Assist Premium — License Cost',
                        '50-seat AI-assisted advisor platform for Banque Directe')

        da_rows = [
            ['Module', 'Year 1', 'Year 2', 'Year 3', 'Year 4', 'Year 5'],
            ['Digital Assist — Premium\n(50 users)', '€350,140', '€350,140', '€350,140', '€350,140', '€350,140'],
            ['CIH Discount', '40%', '40%', '40%', '40%', '40%'],
            ['Annual Cost after Discount', '€210,084', '€210,084', '€210,084', '€210,084', '€210,084'],
        ]
        self._cih_table(s, da_rows, [3.8, 1.5, 1.5, 1.5, 1.5, 1.5], self.ML, Inches(1.6))

        # MH subtotal
        self._section_header_box(s, 'Combined Annual: MH + DA', top=Inches(3.8), width=Inches(5.5))
        combo_rows = [
            ['Component', 'Annual Cost (€)'],
            ['Managed Hosting Enterprise 2', '€541,575'],
            ['Digital Assist Premium (50 seats)', '€210,084'],
            ['Subtotal (MH + DA)', '€751,659'],
        ]
        self._cih_table(s, combo_rows, [7.0, 3.0], self.ML, Inches(4.3))
        self._footer(s, n)

        # ── 8. CLO PRICING ──────────────────────────────
        n += 1
        s = self._new_slide()
        self._title_bar(s, 'Customer Lifecycle Orchestrator — License Cost',
                        'Propensity-driven campaigns targeting 500K wallet-only users')

        clo_rows = [
            ['Module', 'Year 1', 'Year 2', 'Year 3', 'Year 4', 'Year 5'],
            ['CLO (Base Price)', '€381,000', '€381,000', '€381,000', '€381,000', '€381,000'],
            ['Number of Users', '500,000', '500,000', '500,000', '500,000', '500,000'],
            ['Per-User Fee', '€0.70', '€0.70', '€0.70', '€0.70', '€0.70'],
            ['User Fee Total', '€354,816', '€354,816', '€354,816', '€354,816', '€354,816'],
            ['CLO Total (before discount)', '€735,816', '€735,816', '€735,816', '€735,816', '€735,816'],
            ['CIH Discount (30%)', '€515,071', '€515,071', '€515,071', '€515,071', '€515,071'],
        ]
        self._cih_table(s, clo_rows, [3.8, 1.5, 1.5, 1.5, 1.5, 1.5], self.ML, Inches(1.6))
        self._footer(s, n)

        # ── 9. TOTAL LICENSE COST ───────────────────────
        n += 1
        s = self._new_slide()
        self._title_bar(s, 'Total Annual License Cost — All Modules')

        total_rows = [
            ['Module', 'List Price (€)', 'Discount', 'Annual Fee (€)'],
            ['Managed Hosting Enterprise 2\n(180 RPS)', '€601,750', '10%', '€541,575'],
            ['Digital Assist Premium\n(50 users)', '€350,140', '40%', '€210,084'],
            ['CLO — Base Platform', '€381,000', '30%', '€266,700'],
            ['CLO — User Fee\n(500K × €0.70)', '€354,816', '30%', '€248,371'],
            ['TOTAL NEW ANNUAL COMMITMENT', '', '', '€1,266,730'],
        ]
        self._cih_table(s, total_rows, [4.0, 2.5, 1.5, 2.8], self.ML, Inches(1.5))

        # Key stats below
        gap = Inches(0.25)
        sw = (self.CW - 2 * gap) / 3
        for i, (val, lbl, vc, bg) in enumerate([
            ('€1.27M', 'Annual Commitment', self.BLUE, self.BLUE_LIGHT),
            ('€3.8M', '3-Year TCV', self.GREEN, self.GREEN_LIGHT),
            ('Co-terms Jul 2029', 'Contract Alignment', self.DARK_TEXT, self.AMBER_LIGHT),
        ]):
            self._stat_box(s, val, lbl,
                           self.ML + i * (sw + gap), Inches(4.5), sw, Inches(0.6),
                           val_color=vc, bg=bg)

        self._insight_card(s, 'NOTE',
                           'Pricing valid subject to contract signature by Q2 2026. '
                           'All prices in EUR, exclusive of applicable taxes.',
                           self.ML, Inches(5.4), self.CW, Inches(0.45),
                           label_color=self.BLUE, bg=self.BLUE_LIGHT)
        self._footer(s, n)

        # ── 10. SECTION: Value & ROI ────────────────────
        n += 1
        self._divider_slide('Value & ROI', 'Quantified business impact across five value levers', n)

        # ── 11. VALUE LEVERS ────────────────────────────
        n += 1
        s = self._new_slide()
        self._title_bar(s, 'Value Methodology — Five Value Levers',
                        'Each lever uses CIH-specific data and Moroccan retail banking benchmarks')

        levers = [
            ['Value Lever', 'Annual Value', 'Calculation', 'Benchmark Source'],
            ['CLO — Product Origination', '€500K', '500K × 5% conversion × €200 avg revenue', '5% conservative (vs 8–12% industry)'],
            ['CLO — Cross-Sell Uplift', '€600K', '1.5% products/customer × 2M × €20/product', '1.5% vs 3–5% mature CLO'],
            ['DA — Revenue Generation', '€300K', '50 agents × 20% more selling time × €30K/yr', 'DA frees 20% admin → revenue'],
            ['Branch Traffic Reduction', '€1.88M', '25% digital shift × €2.50/txn × 3M txns/yr', 'MENA branch transaction cost'],
            ['MH — Ops Savings', '€300K', '3–4 FTEs redeployed × €75–100K/FTE', 'Morocco infra FTE benchmarks'],
            ['TOTAL ANNUAL VALUE', '€3.58M', '', 'Conservative estimate'],
        ]
        self._cih_table(s, levers, [3.0, 1.5, 4.0, 3.5], self.ML, Inches(1.6))
        self._footer(s, n)

        # ── 12. ROI SENSITIVITY ─────────────────────────
        n += 1
        s = self._new_slide()
        self._title_bar(s, 'ROI & Sensitivity Analysis',
                        'Three scenarios based on annual investment of €1.27M (Option 1)')

        gap = Inches(0.2)
        card_w = (self.CW - 2 * gap) / 3
        scenarios = [
            ('CONSERVATIVE\n(25% Realization)', '€895K', '0.7x', self.AMBER, self.AMBER_LIGHT),
            ('BASE CASE\n(50% Realization)', '€1.79M', '1.4x', self.BLUE, self.BLUE_LIGHT),
            ('FULL POTENTIAL\n(100% Realization)', '€3.58M', '2.8x', self.GREEN, self.GREEN_LIGHT),
        ]
        for i, (label, value, roi, color, bg) in enumerate(scenarios):
            x = self.ML + i * (card_w + gap)
            y = Inches(1.6)
            self._card(s, x, y, card_w, Inches(1.8), fill=bg)
            self._colored_top_line(s, x, y, card_w, color)
            self._txt(s, label, x + Inches(0.2), y + Inches(0.15),
                      card_w - Inches(0.4), Inches(0.5),
                      size=Pt(10), bold=True, color=color, align=PP_ALIGN.CENTER)
            self._txt(s, value, x + Inches(0.2), y + Inches(0.65),
                      card_w - Inches(0.4), Inches(0.35),
                      size=Pt(24), bold=True, color=color, align=PP_ALIGN.CENTER)
            self._txt(s, 'Annual value realized', x + Inches(0.2), y + Inches(1.0),
                      card_w - Inches(0.4), Inches(0.2),
                      size=Pt(8), color=self.SUB_TEXT, align=PP_ALIGN.CENTER)
            self._txt(s, f'ROI: {roi}', x + Inches(0.2), y + Inches(1.3),
                      card_w - Inches(0.4), Inches(0.3),
                      size=Pt(16), bold=True, color=color, align=PP_ALIGN.CENTER)

        # Assumptions table
        self._section_header_box(s, 'Key Assumptions', top=Inches(3.7), width=Inches(3.5))
        assumptions = [
            ['Assumption', 'Value', 'Source'],
            ['Avg product revenue', '€200', 'Moroccan retail banking (blended)'],
            ['Wallet conversion rate', '5%', 'Conservative vs 8–12% (Forrester)'],
            ['Branch transaction cost', '€2.50', 'MENA average (McKinsey)'],
            ['Cross-sell uplift', '1.5%', 'Half of 3–5% mature benchmark'],
            ['Revenue per agent', '€30K', 'Moroccan market rates'],
            ['Infra FTE cost', '€75–100K', 'Fully loaded, Morocco'],
        ]
        self._cih_table(s, assumptions, [3.0, 1.5, 7.5], self.ML, Inches(4.2),
                        row_height=Inches(0.3))
        self._footer(s, n)

        # ── 13. VALUE TIMELINE ──────────────────────────
        n += 1
        s = self._new_slide()
        self._title_bar(s, 'Value Realization Timeline',
                        'Value builds progressively as each module goes live')

        timeline_rows = [
            ['Value Lever', 'H2 2026', '2027', '2028', '2029–30', 'Full Potential'],
            ['Product Origination (CLO)', '—', '€250K', '€500K', '€500K', '€500K'],
            ['Cross-Sell Uplift (CLO)', '—', '€200K', '€450K', '€600K', '€600K'],
            ['DA Revenue Generation', '€75K', '€225K', '€300K', '€300K', '€300K'],
            ['Branch Traffic Reduction', '€188K', '€940K', '€1.50M', '€1.88M', '€1.88M'],
            ['MH Ops Savings', '€150K', '€300K', '€300K', '€300K', '€300K'],
            ['TOTAL ANNUAL VALUE', '€413K', '€1.92M', '€3.05M', '€3.58M', '€3.58M'],
        ]
        self._cih_table(s, timeline_rows, [3.5, 1.7, 1.7, 1.7, 1.7, 1.7],
                        self.ML, Inches(1.6))

        # Value vs Investment
        self._section_header_box(s, 'Value vs. Investment', top=Inches(4.5), width=Inches(4))
        vi_rows = [
            ['', 'H2 2026', '2027', '2028', '2029', '2030', 'Total'],
            ['Value Realized', '€413K', '€1.92M', '€3.05M', '€3.58M', '€3.58M', '€12.5M'],
            ['Investment', '€575K', '€1.15M', '€1.15M', '€1.15M', '€1.15M', '€5.2M'],
            ['Net Value', '−€162K', '+€770K', '+€1.90M', '+€2.43M', '+€2.43M', '+€7.4M'],
        ]
        self._cih_table(s, vi_rows, [2.0, 1.5, 1.5, 1.5, 1.5, 1.5, 1.5],
                        self.ML, Inches(5.0), row_height=Inches(0.32))
        self._footer(s, n)

        # ── 14. SECTION: Implementation ─────────────────
        n += 1
        self._divider_slide('Implementation & Next Steps', 'From signing to full value', n)

        # ── 15. ROADMAP ─────────────────────────────────
        n += 1
        s = self._new_slide()
        self._title_bar(s, 'Implementation Roadmap',
                        'Phased deployment with value milestones at each stage')

        phases = [
            ('Q3 2026', 'FOUNDATION', ['MH migration kickoff', 'DA Premium deployment', 'CLO integration begins'], self.BLUE),
            ('Q4 2026', 'ACTIVATE', ['MH fully live', 'DA agents onboarded', 'First CLO campaigns'], self.GREEN),
            ('H1 2027', 'SCALE', ['Product origination campaigns', 'Cross-sell optimization', 'Branch traffic reduction'], self.PURPLE),
            ('H2 2027+', 'COMPOUND', ['Full CLO maturity', 'AI-driven personalization', 'Value target: €3.6M/yr'], self.GREEN),
        ]
        card_w4 = (self.CW - 3 * Inches(0.15)) / 4
        for i, (date, title, items, color) in enumerate(phases):
            x = self.ML + i * (card_w4 + Inches(0.15))
            y = Inches(1.6)
            self._card(s, x, y, card_w4, Inches(2.5))
            self._colored_top_line(s, x, y, card_w4, color)
            self._txt(s, date, x + Inches(0.15), y + Inches(0.12),
                      card_w4 - Inches(0.3), Inches(0.2),
                      size=Pt(9), bold=True, color=self.MUTED)
            self._txt(s, title, x + Inches(0.15), y + Inches(0.35),
                      card_w4 - Inches(0.3), Inches(0.25),
                      size=Pt(12), bold=True, color=color)
            self._bullet_list(s, items,
                              x + Inches(0.15), y + Inches(0.7),
                              card_w4 - Inches(0.3), Inches(1.5),
                              size=Pt(10))
        self._footer(s, n)

        # ── 16. COMMERCIAL TERMS ────────────────────────
        n += 1
        s = self._new_slide()
        self._title_bar(s, 'Commercial Terms')
        self._section_header_box(s, 'License Conditions', width=Inches(3.5))

        terms = (
            '•  5-year license on an annual subscription model, co-terming with the '
            'existing Backbase contract through July 2029.\n\n'
            '•  The annual subscription includes maintenance and support fees.\n\n'
            '•  All prices are quoted in Euro and are exclusive of applicable taxes.\n\n'
            '•  Commercial and general terms of the existing Master Agreement apply.\n\n'
            '•  License payments are due via Euro wire transfer within 30 days of invoice.\n\n'
            '•  This revised offer is conditional upon contract signature by Q2 2026.'
        )
        self._body_text(s, terms, self.ML, Inches(2.0), Inches(10), Inches(4.5))
        self._footer(s, n)

        # ── 17. NEXT STEPS ──────────────────────────────
        n += 1
        s = self._new_slide()
        self._title_bar(s, 'Next Steps — Mutual Close Plan')

        steps_data = [
            ('1', 'Commercial Alignment', 'Agree on option, confirm terms', 'March 2026', self.BLUE),
            ('2', 'Technical Scoping', 'MH migration plan, DA integration, CLO config', 'April 2026', self.PURPLE),
            ('3', 'Contract Execution', 'Sign and mobilize', 'May 2026', self.GREEN),
            ('4', 'Kickoff & Deployment', 'Phased go-live starting Q3 2026', 'June 2026', self.GREEN),
        ]
        for i, (num, title, desc, date, color) in enumerate(steps_data):
            y = Inches(1.5) + i * Inches(0.95)
            self._bar_rect(s, self.ML, y, Inches(0.4), Inches(0.4), fill=color)
            self._txt(s, num, self.ML, y + Inches(0.03), Inches(0.4), Inches(0.4),
                      size=Pt(16), bold=True, color=self.WHITE, align=PP_ALIGN.CENTER)
            self._txt(s, title, self.ML + Inches(0.6), y + Inches(0.02),
                      Inches(4), Inches(0.25),
                      size=Pt(14), bold=True, color=self.DARK_TEXT)
            self._txt(s, desc, self.ML + Inches(0.6), y + Inches(0.3),
                      Inches(6), Inches(0.25),
                      size=Pt(11), color=self.SUB_TEXT)
            self._txt(s, date, Inches(10), y + Inches(0.1),
                      Inches(3), Inches(0.25),
                      size=Pt(11), bold=True, color=color, align=PP_ALIGN.RIGHT)
        self._footer(s, n)

        # ── 18. CLOSING ─────────────────────────────────
        n += 1
        s = self._new_slide(dark=True)
        if Path(BLUE_GRADIENT).exists():
            s.shapes.add_picture(BLUE_GRADIENT, 0, 0, self.SLIDE_W, self.SLIDE_H)
        self._txt(s, 'Thank you.',
                  Inches(1.5), Inches(2.2), Inches(8), Inches(0.8),
                  size=Pt(36), bold=True, color=self.WHITE)
        self._txt(s, 'We look forward to deepening our partnership\nwith CIH Bank.',
                  Inches(1.5), Inches(3.2), Inches(8), Inches(0.6),
                  size=Pt(16), color=RGBColor(0xCC, 0xDD, 0xFF))

        # CIH logo on closing
        if Path(CIH_LOGO).exists():
            s.shapes.add_picture(CIH_LOGO, Inches(1.5), Inches(5.0),
                                 Inches(1.8), Inches(0.62))
        self._footer(s, n, dark=True)

        # ── SAVE ────────────────────────────────────────
        self.save(output_path)


if __name__ == '__main__':
    deck = CIHCommercialV2()
    deck.generate('/Users/shyam/cortex/Engagement/CIH Bank/Output/CIH_Commercial_Proposal_v2.pptx')
