#!/usr/bin/env python3
"""
Schroders Playback PPTX Generator

Converts the 18-scene interactive HTML presentation into a
Google Slides-compatible .pptx file with dark theme styling.

Usage:
    python3 tools/schroders_pptx_generator.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


class SchrodersPptxGenerator:
    """Generates a dark-themed PPTX from Schroders playback data."""

    # ── Dimensions ─────────────────────────────────────────────
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    # ── Colors ─────────────────────────────────────────────────
    DARK_BG   = RGBColor(0x09, 0x1C, 0x35)
    DARKER_BG = RGBColor(0x06, 0x0E, 0x1A)
    CARD_BG   = RGBColor(0x0E, 0x22, 0x40)
    BLUE      = RGBColor(0x33, 0x66, 0xFF)
    PURPLE    = RGBColor(0x7B, 0x2F, 0xFF)
    RED       = RGBColor(0xFF, 0x72, 0x62)
    GREEN     = RGBColor(0x26, 0xBC, 0x71)
    AMBER     = RGBColor(0xFF, 0xAC, 0x09)
    CYAN      = RGBColor(0x69, 0xFE, 0xFF)
    GRAY      = RGBColor(0x8A, 0x9B, 0xB5)
    MID       = RGBColor(0x3A, 0x49, 0x5D)
    WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BG  = RGBColor(0xF5, 0xF6, 0xFA)
    BORDER    = RGBColor(0x1A, 0x3A, 0x5E)
    DARK_TEXT  = RGBColor(0x1A, 0x23, 0x40)

    # RAG
    RAG_ABSENT     = RGBColor(0xDC, 0x26, 0x26)
    RAG_FRAGMENTED = RGBColor(0xEA, 0x58, 0x0C)
    RAG_DEFINED    = RGBColor(0x08, 0x91, 0xB2)
    RAG_ORCHESTRATED = RGBColor(0x05, 0x96, 0x69)

    FONT = 'Libre Franklin'
    ML = Inches(0.6)     # margin left
    CW = Inches(12.1)    # content width
    CT = Inches(1.6)     # content top (below heading)

    # ── DATA ───────────────────────────────────────────────────

    COVER_STATS = [
        {'\u00a372B': 'Assets Under Management', 'd': 'Cazenove Capital'},
        {'30K': 'Clients', 'd': 'Discretionary + Advisory'},
        {'225': 'Client Advisors', 'd': 'Spending 65% on admin'},
    ]

    HEATMAP_ROWS = [
        ['Digital Prospecting',       'Y','','', 1, 3, 'Critical'],
        ['Client Onboarding',         'Y','Y','Y', 0, 3, 'Critical'],
        ['KYC/AML Automation',        '','Y','Y', 1, 3, 'Critical'],
        ['Unified Advisor Cockpit',   'Y','Y','', 0, 3, 'Critical'],
        ['Client Self-Service',       'Y','','', 0, 3, 'High'],
        ['Client Review & Reporting', 'Y','Y','', 1, 3, 'Critical'],
        ['Wealth Planning Tools',     'Y','','', 1, 3, 'High'],
        ['Client Engagement',         'Y','','', 0, 3, 'High'],
        ['Document Management',       '','Y','Y', 1, 2, 'Medium'],
        ['Analytics & Intelligence',  '','Y','', 0, 3, 'High'],
        ['API & Integration Layer',   '','','Y', 1, 3, 'Critical'],
        ['Master Client Record',      '','Y','Y', 1, 3, 'Critical'],
    ]

    PEER_METRICS = [
        ('Digital Onboarding Time', '9.6 days (Schroders)', '1 day (Best)', 0.85),
        ('Advisor Review Preparation', '4\u20136 hrs (Schroders)', '1.5 hrs (Best)', 0.75),
        ('Client Self-Service Adoption', '<30% (Schroders)', '85%+ (Best)', 0.25),
        ('Speed to Market \u2014 Change Delivery', '12\u201316 wks (Schroders)', '<6 wks (Best)', 0.80),
        ('Technology & Innovation Agility', 'Regressive (Schroders)', 'AI-native (Best)', 0.90),
        ('Digital Investment as % of Revenue', '<2% (Schroders est.)', '5\u20138% (Leaders)', 0.20),
    ]

    ROADMAP_PHASES = [
        ('Phase 1 \u2014 0\u201312 Months', 'Foundation & Core Journeys', '\u00a311M Year 1',
         ['API & Integration Layer (Grand Central)', 'Master Client Record \u2014 golden source',
          'Digital Onboarding with KYC/AML', 'Unified Advisor Cockpit (Client 360\u00b0)',
          'Salesforce connector & data sync']),
        ('Phase 2 \u2014 12\u201324 Months', 'Growth & Engagement', '\u00a356.4M cumulative by Year 3',
         ['Client Self-Service Portal & mobile app', 'Digital Review & reporting automation',
          'Engagement Suite (NBA, Smart Signals)', 'Secure messaging & document exchange']),
        ('Phase 3 \u2014 24\u201336 Months', 'Intelligence & Scale', '\u00a3133.5M full run-rate by Year 5',
         ['AI agents (review, health monitor, planning)', 'Digital Wealth Planning & Voyant integration',
          'Family / intergenerational planning', 'International advisory expansion support']),
    ]

    BENEFIT_BARS = [('\u00a311M', 11), ('\u00a317.6M', 17.6), ('\u00a327.8M', 27.8),
                    ('\u00a336.8M', 36.8), ('\u00a336.1M', 36.1)]

    STRATEGIC_FIT = [
        ('\U0001f3d7\ufe0f', 'Platform, Not Point Solutions',
         ['Single unified platform replacing 6\u20137 siloed systems',
          'One codebase, one data layer, one client record',
          'Build once, deploy across all channels',
          'Eliminates integration tax']),
        ('\U0001f4b0', 'Wealth Management DNA',
         ['Purpose-built for wealth \u2014 not retail banking retrofit',
          'Pre-built journeys: onboarding, reviews, planning',
          'Proven with leading private banks globally',
          'FCA-ready compliance and audit trail built in']),
        ('\U0001f680', 'Innovation Velocity',
         ['<40 days for changes vs 3\u20136 months (competitors)',
          '12 months to production vs 18\u201324 months',
          'Full SDK access \u2014 not locked into vendor roadmap',
          'AI-native architecture with Agentic Studio']),
        ('\u26a1', 'Predictable Costs',
         ['45% lower total cost of ownership vs alternatives',
          'No hidden costs: AI-ready data, integration, omni-channel included',
          'Grand Central eliminates bespoke integration effort',
          'Phased approach: start small, scale fast']),
    ]

    PROOF_POINTS = [
        ('45%', 'Lower Cost', 'vs. competitive alternatives'),
        ('+13%', 'AUM Fee Growth', 'Proven client results'),
        ('+50', 'NPS Improvement', 'Client satisfaction uplift'),
        ('83%', 'Cost-to-Serve Reduction', 'Operational efficiency'),
        ('70%', 'Increase in Referrals', 'Organic growth driver'),
        ('2yr', 'Avg. Payback Period', 'Across wealth deployments'),
    ]

    RECAP_CRITICAL = [
        ('AI Pricing Model', 'Michael', 'in-progress',
         'How is AI priced \u2014 per-prompt, module licence, or MVP bundled?',
         '\u201CWe will share the details.\u201D Session with Michael & Peter next week.',
         'Addressed in scope toggle \u2014 Premium vs Signature packaging.'),
        ('Data Validation & Escalation', 'Jonny', 'to-scope',
         'Address updates in high-risk jurisdictions \u2014 what controls trigger escalation?',
         'Validation in journeys; escalation via Case Manager. Needs joint scoping.',
         'Architecture walkthrough today covers orchestration and escalation rules.'),
        ('Payment Method Verification', 'Tim', 'to-scope',
         'SortCode/CHAPS compliance \u2014 can system verify receiving bank accepts payment type?',
         '\u201CWe\u2019ll have to validate.\u201D COP checks confirmed; CHAPS TBC.',
         'Technical deep-dive with payments team to confirm.'),
        ('APP Fraud Protection', 'Liz', 'to-scope',
         'Restrict payments to own accounts or pre-validated counterparties only.',
         'Configurable restrictions \u2014 dial up/down by account type. Needs scoping.',
         'Payment controls configurable per client segment.'),
        ('Opportunity Status in Dashboard', 'Jonny', 'to-scope',
         'Can advisor dashboard show opportunity pipeline from Salesforce?',
         '\u201CNot today.\u201D Flagged for employee experience session scoping.',
         'SF connector extension \u2014 Phase 1 scope, custom widget.'),
        ('SIP Setup End-to-End Visibility', 'Liz', 'to-scope',
         'Intelliflow \u2192 SIP \u2192 T24 \u2014 no end-to-end visibility today.',
         'Account aggregation capabilities exist; needs scoping.',
         'Orchestration layer covers this in Phase 2.'),
    ]

    RECAP_CONFIDENCE = [
        ('SF Connector Data Model', 'Jonny', 'to-scope'),
        ('Custom SF Data Objects', 'Jonny', 'to-scope'),
        ('Correspondent Banking', 'Tim', 'resolved'),
        ('Configure vs Customise Costs', 'Paul', 'in-progress'),
        ('Chat Multi-user & Routing', 'Briony', 'resolved'),
        ('Chat Archiving & History', 'Paul', 'resolved'),
        ('Client Segmentation Rules', 'Liz', 'resolved'),
        ('SF Marketing Cloud Integration', 'Andreas', 'to-scope'),
    ]

    RECAP_NICE = [
        ('Live Chat Sentiment Analysis', 'Briony', 'to-scope'),
        ('Chat Close/Reopen Workflow', 'Paul', 'resolved'),
        ('Financial Planning Integration', 'Liz', 'to-scope'),
        ('Churn/Risk Scoring Models', 'Group', 'to-scope'),
    ]

    ROI_LEVERS = [
        ('Growth Engine', '46%', '\u00a361M', BLUE,
         [('Onboarding \u2192 NNB Fees', '\u00a338M'),
          ('Web & Mobile NNB Fees', '\u00a315M'),
          ('Fee Recognition Acceleration', '\u00a38.5M')]),
        ('Operations Factory', '16%', '\u00a322M \u2014 99 FTE', GREEN,
         [('Onboarding Ops Efficiency', '\u00a34.6M'),
          ('Client Review Efficiency', '\u00a310M'),
          ('Web & Mobile Efficiency', '\u00a30.7M'),
          ('Admin & Planning Automation', '\u00a36.7M')]),
        ('Tailored Propositions', '33%', '\u00a344M', PURPLE,
         [('Client Experience & Engagement Uplift', '\u00a322M'),
          ('Share of Wallet & Cross-sell', '\u00a314M'),
          ('Reduced Attrition', '\u00a38M')]),
        ('Architecture as Strategy', '5%', '\u00a36.5M', CYAN,
         [('Integration Layer Consolidation', '\u00a33M'),
          ('Data Sync & Re-keying Elimination', '\u00a32M'),
          ('Platform Consolidation Savings', '\u00a31.5M')]),
    ]

    AI_SPOTLIGHT = [
        ('L1', 'Automation', 'Digital form & doc automation',
         'Replace 90-min PDF forms with digital flows',
         'ONBOARDING \u00b7 CLM', '\u00a311M fee recognition + \u00a36M ops savings', 'Phase 1'),
        ('L1', 'Automation', 'Auto-aggregated review decks',
         'Replace 2h manual data pull from 7 systems',
         'ADVISOR COCKPIT', '\u00a316M review efficiency', 'Phase 2'),
        ('L2', 'Predictive', 'Smart lead scoring & targeting',
         'AI ranks prospects by conversion probability',
         'SMART SIGNALS', 'Part of \u00a330M lead conversion', 'Phase 2'),
        ('L2', 'Predictive', 'Churn prediction & alerts',
         'Identify at-risk clients before outflow',
         'INTELLIGENCE FABRIC', '\u00a38M attrition reduction', 'Phase 2\u20133'),
        ('L3', 'Embedded', 'AI-generated proposals',
         'Contextual proposals from client profile + market data',
         'DIGITAL PROPOSALS', 'Part of \u00a320M review\u2192NNB', 'Phase 3'),
        ('L3', 'Embedded', 'Conversational AI advisor',
         'Natural language interface for queries & instructions',
         'EMPLOYEE WORKSPACE', 'Part of \u00a336M CX uplift', 'Phase 3'),
    ]

    # ── HELPERS ─────────────────────────────────────────────────

    def _new_slide(self, dark=True):
        layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.DARK_BG if dark else self.LIGHT_BG
        return slide

    def _txt(self, slide, text, left, top, width, height, size=Pt(12),
             color=None, bold=False, align=PP_ALIGN.LEFT, font=None):
        """Add a simple text box. Returns the text frame."""
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color or self.WHITE
        p.font.name = font or self.FONT
        p.alignment = align
        return tf

    def _section_label(self, slide, text, top=Inches(0.45)):
        self._txt(slide, text.upper(), self.ML, top, self.CW, Pt(18),
                  size=Pt(10), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)

    def _heading(self, slide, main, accent='', top=Inches(0.85), size=Pt(36)):
        tb = slide.shapes.add_textbox(self.ML, top, self.CW, Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r1 = p.add_run()
        r1.text = main
        r1.font.size = size
        r1.font.bold = True
        r1.font.color.rgb = self.WHITE
        r1.font.name = self.FONT
        if accent:
            r2 = p.add_run()
            r2.text = accent
            r2.font.size = size
            r2.font.bold = True
            r2.font.color.rgb = self.BLUE
            r2.font.name = self.FONT

    def _subtitle(self, slide, text, top=Inches(1.45)):
        self._txt(slide, text, self.ML, top, self.CW, Inches(0.4),
                  size=Pt(14), color=self.GRAY, align=PP_ALIGN.CENTER)

    def _card(self, slide, left, top, w, h, fill=None, border=None):
        """Rounded rectangle card. Returns shape."""
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        if fill:
            s.fill.solid()
            s.fill.fore_color.rgb = fill
        else:
            s.fill.solid()
            s.fill.fore_color.rgb = self.CARD_BG
        s.line.color.rgb = border or self.BORDER
        s.line.width = Pt(1)
        return s

    def _stat_card(self, slide, value, label, detail, left, top, w, h,
                   val_color=None, highlight=False):
        border = self.BLUE if highlight else self.BORDER
        self._card(slide, left, top, w, h, border=border)
        # accent line
        a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left + Inches(0.15), top + Pt(4),
                                   w - Inches(0.3), Pt(2))
        a.fill.solid()
        a.fill.fore_color.rgb = self.BLUE if highlight else self.BORDER
        a.line.fill.background()
        # value
        self._txt(slide, value, left, top + Inches(0.15), w, Inches(0.5),
                  size=Pt(30), color=val_color or self.BLUE, bold=True,
                  align=PP_ALIGN.CENTER)
        # label
        self._txt(slide, label.upper(), left, top + Inches(0.65), w, Pt(14),
                  size=Pt(8), color=self.GRAY, bold=True, align=PP_ALIGN.CENTER)
        # detail
        if detail:
            self._txt(slide, detail, left, top + Inches(0.85), w, Pt(12),
                      size=Pt(9), color=self.GRAY, align=PP_ALIGN.CENTER)

    def _info_card(self, slide, title, bullets, left, top, w, h,
                   border_color=None, icon=''):
        self._card(slide, left, top, w, h, border=border_color or self.BORDER)
        # left accent bar
        if border_color:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         left, top + Inches(0.1),
                                         Pt(4), h - Inches(0.2))
            bar.fill.solid()
            bar.fill.fore_color.rgb = border_color
            bar.line.fill.background()
        y = top + Inches(0.12)
        if icon:
            self._txt(slide, icon, left + Inches(0.15), y, Inches(0.4), Pt(20),
                      size=Pt(16), align=PP_ALIGN.LEFT)
            y += Pt(22)
        self._txt(slide, title, left + Inches(0.15), y, w - Inches(0.3), Pt(18),
                  size=Pt(12), bold=True, align=PP_ALIGN.LEFT)
        y += Pt(22)
        for b in bullets:
            self._txt(slide, '\u2022 ' + b, left + Inches(0.15), y,
                      w - Inches(0.3), Pt(14),
                      size=Pt(9), color=self.GRAY, align=PP_ALIGN.LEFT)
            y += Pt(15)

    def _insight_box(self, slide, text, top=Inches(6.6)):
        self._card(slide, self.ML, top, self.CW, Inches(0.55))
        self._txt(slide, text, self.ML + Inches(0.2), top + Inches(0.08),
                  self.CW - Inches(0.4), Inches(0.4),
                  size=Pt(10), color=self.GRAY, align=PP_ALIGN.LEFT)

    def _status_pill(self, status):
        if status == 'resolved':
            return '\u2705 Resolved'
        elif status == 'in-progress':
            return '\U0001f504 In Progress'
        return '\U0001f4cb To Scope'

    def _status_color(self, status):
        if status == 'resolved':
            return self.GREEN
        elif status == 'in-progress':
            return self.AMBER
        return self.GRAY

    # ── SLIDE BUILDERS ──────────────────────────────────────────

    def _slide_cover(self):
        s = self._new_slide()
        self._section_label(s, 'Playback \u2014 26 & 27 February 2026')
        # Mega title
        tb = s.shapes.add_textbox(self.ML, Inches(1.2), self.CW, Inches(2.5))
        tf = tb.text_frame
        tf.word_wrap = True
        for line, color in [('ACCELERATING', self.WHITE),
                            ('SCHRODERS', self.BLUE),
                            ('ADVANTAGE', self.BLUE)]:
            p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = line
            r.font.size = Pt(56)
            r.font.bold = True
            r.font.color.rgb = color
            r.font.name = self.FONT
        # Subtitle
        self._txt(s, 'Schroders Wealth Management \u2014 Cazenove Capital',
                  self.ML, Inches(3.8), self.CW, Inches(0.35),
                  size=Pt(14), color=self.GRAY, align=PP_ALIGN.CENTER)
        self._txt(s, '10 workshops  \u00b7  July\u2013December 2025  \u00b7  One digital wealth platform',
                  self.ML, Inches(4.15), self.CW, Inches(0.3),
                  size=Pt(11), color=self.MID, align=PP_ALIGN.CENTER)
        # Stat cards
        cw = Inches(3.6)
        gap = Inches(0.5)
        sx = Inches(1.5)
        for i, (val_label) in enumerate([
            ('\u00a372B', 'Assets Under Management', 'Cazenove Capital'),
            ('30K', 'Clients', 'Discretionary + Advisory'),
            ('225', 'Client Advisors', 'Spending 65% on admin'),
        ]):
            self._stat_card(s, val_label[0], val_label[1], val_label[2],
                           sx + i * (cw + gap), Inches(5.0), cw, Inches(1.15))

    def _slide_agenda(self):
        s = self._new_slide()
        self._section_label(s, 'Your Playback Journey')
        self._heading(s, 'Day 1 Recap & ', 'Day 2 Preview')
        col_w = Inches(5.5)
        lx = Inches(0.8)
        rx = Inches(7.0)
        ty = Inches(1.9)
        # Column headers
        self._card(s, lx, ty, col_w, Pt(28), fill=RGBColor(0x3A, 0x20, 0x20),
                   border=self.RED)
        self._txt(s, 'Yesterday \u2014 Functional & Technical',
                  lx + Inches(0.1), ty + Pt(3), col_w - Inches(0.2), Pt(20),
                  size=Pt(11), bold=True, color=self.RED, align=PP_ALIGN.CENTER)
        self._card(s, rx, ty, col_w, Pt(28), fill=RGBColor(0x12, 0x30, 0x20),
                   border=self.GREEN)
        self._txt(s, 'Today \u2014 Strategic Value Story',
                  rx + Inches(0.1), ty + Pt(3), col_w - Inches(0.2), Pt(20),
                  size=Pt(11), bold=True, color=self.GREEN, align=PP_ALIGN.CENTER)
        # Items
        day1 = [('\U0001f4f1', 'Client Experience Deep-Dive', 'Web & App demos, self-service journeys'),
                ('\u2699\ufe0f', 'Broader Capability Demo', 'CLM, onboarding flows, collaboration'),
                ('\U0001f50c', 'Integrations & Security', 'Salesforce, document vault, identity'),
                ('\U0001f4ca', 'Analytics Capabilities', 'Smart signals, next best actions, segmentation'),
                ('\U0001f6e0\ufe0f', 'Implementation Discussion', 'Grand Central, dual-core coexistence')]
        day2 = [('\U0001f3af', 'Value Story & Business Case', 'This presentation \u2014 your strategic context'),
                ('\U0001f5fa\ufe0f', 'Journey Transformation Deep-Dive', 'Onboarding, Reviews, Wealth Planning'),
                ('\U0001f4b0', 'ROI & Peer Validation', '\u00a3133.5M opportunity, peer benchmarks'),
                ('\U0001f680', 'Roadmap & What\u2019s Next', 'Phased approach, decision framework')]
        for col_items, cx in [(day1, lx), (day2, rx)]:
            for j, (icon, title, sub) in enumerate(col_items):
                iy = ty + Pt(38) + j * Inches(0.65)
                self._card(s, cx, iy, col_w, Inches(0.55))
                self._txt(s, icon, cx + Inches(0.12), iy + Inches(0.08),
                          Inches(0.35), Inches(0.3), size=Pt(16))
                self._txt(s, title, cx + Inches(0.5), iy + Inches(0.05),
                          col_w - Inches(0.65), Pt(16),
                          size=Pt(11), bold=True)
                self._txt(s, sub, cx + Inches(0.5), iy + Inches(0.30),
                          col_w - Inches(0.65), Pt(14),
                          size=Pt(9), color=self.GRAY)

    def _slide_open_items_overview(self):
        s = self._new_slide()
        self._section_label(s, 'Day 1 Recap')
        self._heading(s, 'We Heard You \u2014 ', 'Open Items Tracker')
        self._subtitle(s, '18 items from yesterday\u2019s session. Here\u2019s where we stand.')
        # Summary pills
        pw = Inches(3.2)
        gap = Inches(0.3)
        sx = Inches(2.0)
        for i, (count, label, col) in enumerate([
            ('6', 'Deal-Critical', self.RED),
            ('8', 'Confidence-Building', self.AMBER),
            ('4', 'Nice-to-Close', self.CYAN),
        ]):
            x = sx + i * (pw + gap)
            self._card(s, x, Inches(1.9), pw, Inches(0.6))
            self._txt(s, count, x + Inches(0.15), Inches(1.95), Inches(0.5), Inches(0.45),
                      size=Pt(24), bold=True, color=col, align=PP_ALIGN.CENTER)
            self._txt(s, label, x + Inches(0.7), Inches(2.0), pw - Inches(0.85), Inches(0.4),
                      size=Pt(12), color=self.GRAY)
        # Status bar
        self._txt(s, '\u25cf 5 Resolved    \u25cf 2 In Progress    \u25cf 11 To Scope',
                  self.ML, Inches(2.7), self.CW, Pt(16),
                  size=Pt(11), color=self.GRAY, align=PP_ALIGN.CENTER)
        # Tier summary cards
        tiers = [
            ('Tier 1 \u2014 Deal-Critical', '6 items \u00b7 1 in progress, 5 to scope', self.RED),
            ('Tier 2 \u2014 Confidence-Building', '8 items \u00b7 4 resolved, 1 in progress, 3 to scope', self.AMBER),
            ('Tier 3 \u2014 Nice-to-Close', '4 items \u00b7 1 resolved, 3 to scope', self.CYAN),
        ]
        for i, (name, count, col) in enumerate(tiers):
            ty = Inches(3.3) + i * Inches(1.2)
            self._card(s, Inches(1.5), ty, Inches(10.3), Inches(1.0))
            # colored left bar
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     Inches(1.5), ty + Inches(0.05),
                                     Pt(5), Inches(0.9))
            bar.fill.solid()
            bar.fill.fore_color.rgb = col
            bar.line.fill.background()
            self._txt(s, name, Inches(1.8), ty + Inches(0.15),
                      Inches(5), Pt(18), size=Pt(14), bold=True)
            self._txt(s, count, Inches(1.8), ty + Inches(0.50),
                      Inches(8), Pt(14), size=Pt(10), color=self.GRAY)

    def _slide_deal_critical(self):
        s = self._new_slide()
        self._section_label(s, 'Day 1 Recap \u2014 Tier 1')
        self._heading(s, 'Deal-Critical Items ', '(6)', size=Pt(30))
        cw = Inches(5.8)
        gap = Inches(0.3)
        ch = Inches(2.3)
        for i, (title, raised, status, q, resp, act) in enumerate(self.RECAP_CRITICAL):
            col = i % 2
            row = i // 2
            x = self.ML + col * (cw + gap)
            y = Inches(1.8) + row * (ch + Inches(0.15))
            sc = self._status_color(status)
            self._card(s, x, y, cw, ch, border=sc)
            # title + status
            self._txt(s, title, x + Inches(0.15), y + Inches(0.08),
                      cw - Inches(1.5), Pt(16), size=Pt(12), bold=True)
            self._txt(s, self._status_pill(status),
                      x + cw - Inches(1.3), y + Inches(0.08),
                      Inches(1.2), Pt(14), size=Pt(8), color=sc, bold=True,
                      align=PP_ALIGN.RIGHT)
            self._txt(s, 'Raised by ' + raised,
                      x + Inches(0.15), y + Inches(0.35),
                      Inches(3), Pt(12), size=Pt(8), color=self.GRAY)
            # Q / R / A
            labels = [('QUESTION', self.AMBER, q),
                      ('DAY 1 RESPONSE', self.GRAY, resp),
                      ('DAY 2 ACTION', self.GREEN, act)]
            ly = y + Inches(0.6)
            for lbl, lc, txt in labels:
                self._txt(s, lbl, x + Inches(0.15), ly,
                          Inches(1.5), Pt(10), size=Pt(7), color=lc, bold=True)
                ly += Pt(11)
                self._txt(s, txt, x + Inches(0.15), ly,
                          cw - Inches(0.3), Pt(28), size=Pt(8), color=self.GRAY)
                ly += Pt(30)

    def _slide_confidence_nice(self):
        s = self._new_slide()
        self._section_label(s, 'Day 1 Recap \u2014 Tiers 2 & 3')
        self._heading(s, 'Confidence-Building & ', 'Nice-to-Close', size=Pt(28))
        # Tier 2 header
        self._txt(s, 'TIER 2 \u2014 CONFIDENCE-BUILDING (8 items)',
                  self.ML, Inches(1.65), self.CW, Pt(16),
                  size=Pt(10), color=self.AMBER, bold=True)
        cw = Inches(2.85)
        gap = Inches(0.15)
        for i, (title, raised, status) in enumerate(self.RECAP_CONFIDENCE):
            col = i % 4
            row = i // 4
            x = self.ML + col * (cw + gap)
            y = Inches(1.95) + row * Inches(0.72)
            sc = self._status_color(status)
            self._card(s, x, y, cw, Inches(0.62))
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x, y + Inches(0.05), Pt(3), Inches(0.52))
            bar.fill.solid()
            bar.fill.fore_color.rgb = self.AMBER
            bar.line.fill.background()
            self._txt(s, title, x + Inches(0.12), y + Inches(0.05),
                      cw - Inches(0.25), Pt(14), size=Pt(10), bold=True)
            self._txt(s, raised, x + Inches(0.12), y + Inches(0.28),
                      Inches(1.5), Pt(12), size=Pt(8), color=self.GRAY)
            self._txt(s, self._status_pill(status),
                      x + cw - Inches(1.2), y + Inches(0.38),
                      Inches(1.1), Pt(12), size=Pt(7), color=sc, bold=True,
                      align=PP_ALIGN.RIGHT)
        # Tier 3 header
        self._txt(s, 'TIER 3 \u2014 NICE-TO-CLOSE (4 items)',
                  self.ML, Inches(3.65), self.CW, Pt(16),
                  size=Pt(10), color=self.CYAN, bold=True)
        for i, (title, raised, status) in enumerate(self.RECAP_NICE):
            x = self.ML + i * (Inches(2.85) + Inches(0.15))
            y = Inches(3.95)
            sc = self._status_color(status)
            self._card(s, x, y, Inches(2.85), Inches(0.62))
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     x, y + Inches(0.05), Pt(3), Inches(0.52))
            bar.fill.solid()
            bar.fill.fore_color.rgb = self.CYAN
            bar.line.fill.background()
            self._txt(s, title, x + Inches(0.12), y + Inches(0.05),
                      Inches(2.6), Pt(14), size=Pt(10), bold=True)
            self._txt(s, raised, x + Inches(0.12), y + Inches(0.28),
                      Inches(1.5), Pt(12), size=Pt(8), color=self.GRAY)
            self._txt(s, self._status_pill(status),
                      x + Inches(1.6), y + Inches(0.38),
                      Inches(1.1), Pt(12), size=Pt(7), color=sc, bold=True,
                      align=PP_ALIGN.RIGHT)
        # Summary insight
        self._insight_box(s,
            'Summary: 5 items fully resolved \u2022 2 addressed in today\u2019s sessions \u2022 '
            '11 to be scoped in architecture & delivery discussions.',
            top=Inches(4.85))

    def _slide_vision(self):
        s = self._new_slide()
        self._section_label(s, 'Where You Are & Where You\u2019re Going')
        self._heading(s, 'The Aspiration vs. ', 'The Reality')
        cw = Inches(5.7)
        gap = Inches(0.5)
        lx = Inches(0.8)
        rx = lx + cw + gap
        ty = self.CT
        # Aspiration
        self._card(s, lx, ty, cw, Inches(4.5), border=self.BLUE)
        self._txt(s, 'THE ASPIRATION', lx + Inches(0.2), ty + Inches(0.12),
                  cw, Pt(14), size=Pt(10), color=self.BLUE, bold=True)
        asp = [('\U0001f30d', 'International advisory expansion \u2014 Middle East & Asia markets'),
               ('\U0001f4c8', '7%+ NNA growth year-on-year \u2014 \u00a35B net new assets in 2024'),
               ('\U0001f4cb', 'Deepen planning relationships \u2014 20% \u2192 50%+ penetration'),
               ('\u2b50', 'Differentiate on experience \u2014 digitally amplified, relationship-led')]
        ay = ty + Inches(0.5)
        for icon, txt in asp:
            self._txt(s, icon + '  ' + txt, lx + Inches(0.2), ay,
                      cw - Inches(0.4), Inches(0.6),
                      size=Pt(11), color=self.WHITE)
            ay += Inches(0.65)
        self._txt(s, 'UK wealth AUM projected to reach $12.7T by 2029',
                  lx + Inches(0.2), ty + Inches(3.8), cw - Inches(0.4), Pt(12),
                  size=Pt(9), color=self.CYAN, bold=True)
        # Reality
        self._card(s, rx, ty, cw, Inches(4.5), border=self.RED)
        self._txt(s, 'THE REALITY', rx + Inches(0.2), ty + Inches(0.12),
                  cw, Pt(14), size=Pt(10), color=self.RED, bold=True)
        real = [('\u23f3', '4.5 hours to onboard a single client', 'Benchmark: 1 hour \u2014 4.5\u00d7 gap'),
                ('\U0001f6ab', '35% of onboarding applications rejected', 'Benchmark: 5% \u2014 7\u00d7 gap'),
                ('\U0001f5c2\ufe0f', '6\u20137 systems to prepare one client review', 'Benchmark: 1 unified cockpit'),
                ('\U0001f4c4', 'Only 20% of clients have a financial plan', 'Industry target: 50%+')]
        ry = ty + Inches(0.5)
        for icon, txt, bench in real:
            self._txt(s, icon + '  ' + txt, rx + Inches(0.2), ry,
                      cw - Inches(0.4), Pt(16), size=Pt(11), bold=True)
            self._txt(s, bench, rx + Inches(0.55), ry + Pt(18),
                      cw - Inches(0.7), Pt(12), size=Pt(9), color=self.RED)
            ry += Inches(0.85)

    def _slide_journey_overview(self):
        s = self._new_slide()
        self._section_label(s, 'Three Journeys That Matter Most')
        self._heading(s, 'The Client Lifecycle \u2014 ', 'Before & After')
        cw = Inches(3.7)
        gap = Inches(0.35)
        sx = Inches(1.0)
        journeys = [
            ('Client Onboarding', '9.6 days \u2192 3 days', '69%', 'faster', '\u00a366M 5-year value'),
            ('Client Review', '4\u20136 hrs \u2192 45 min', '85%', 'faster', '\u00a336M 5-year value'),
            ('Wealth Planning', '~6.5 hrs \u2192 ~1.5 hrs', '77%', 'faster', '\u00a36.4M 5-year value'),
        ]
        for i, (name, before_after, pct, desc, value) in enumerate(journeys):
            x = sx + i * (cw + gap)
            y = Inches(1.9)
            self._card(s, x, y, cw, Inches(3.8))
            self._txt(s, name, x + Inches(0.2), y + Inches(0.15),
                      cw - Inches(0.4), Pt(18), size=Pt(14), bold=True,
                      align=PP_ALIGN.CENTER)
            self._txt(s, before_after, x + Inches(0.2), y + Inches(0.55),
                      cw - Inches(0.4), Pt(14), size=Pt(11), color=self.GRAY,
                      align=PP_ALIGN.CENTER)
            # Big percentage
            self._txt(s, pct, x + Inches(0.2), y + Inches(1.0),
                      cw - Inches(0.4), Inches(0.8),
                      size=Pt(56), color=self.GREEN, bold=True,
                      align=PP_ALIGN.CENTER)
            self._txt(s, desc, x + Inches(0.2), y + Inches(2.0),
                      cw - Inches(0.4), Pt(16), size=Pt(14), color=self.GRAY,
                      align=PP_ALIGN.CENTER)
            self._txt(s, value, x + Inches(0.2), y + Inches(2.8),
                      cw - Inches(0.4), Pt(16), size=Pt(12), color=self.BLUE,
                      bold=True, align=PP_ALIGN.CENTER)

    def _journey_slide(self, label_num, journey_name, headline,
                       current_time, future_time,
                       current_steps, future_steps,
                       current_calls, future_calls, value_text):
        s = self._new_slide()
        self._section_label(s, f'Journey {label_num} \u2014 {journey_name}')
        self._heading(s, headline)
        cw = Inches(5.7)
        gap = Inches(0.5)
        lx = Inches(0.8)
        rx = lx + cw + gap
        # Current column
        self._card(s, lx, Inches(1.7), cw, Inches(4.8), border=self.RED)
        self._txt(s, f'Current State ({current_time})',
                  lx + Inches(0.15), Inches(1.78), cw, Pt(16),
                  size=Pt(11), color=self.RED, bold=True)
        cy = Inches(2.1)
        for step, actor, time in current_steps:
            self._txt(s, f'\u25b8 {step}', lx + Inches(0.15), cy,
                      cw - Inches(0.3), Pt(14), size=Pt(10), bold=True)
            self._txt(s, f'   {actor} \u00b7 {time}', lx + Inches(0.15), cy + Pt(14),
                      cw - Inches(0.3), Pt(12), size=Pt(8), color=self.GRAY)
            cy += Inches(0.38)
        for call in current_calls:
            self._txt(s, '\u26a0 ' + call, lx + Inches(0.15), cy,
                      cw - Inches(0.3), Pt(12), size=Pt(9), color=self.RED)
            cy += Pt(14)
        # Future column
        self._card(s, rx, Inches(1.7), cw, Inches(4.8), border=self.GREEN)
        self._txt(s, f'Future State ({future_time})',
                  rx + Inches(0.15), Inches(1.78), cw, Pt(16),
                  size=Pt(11), color=self.GREEN, bold=True)
        fy = Inches(2.1)
        for step, actor, time in future_steps:
            self._txt(s, f'\u25b8 {step}', rx + Inches(0.15), fy,
                      cw - Inches(0.3), Pt(14), size=Pt(10), bold=True)
            self._txt(s, f'   {actor} \u00b7 {time}', rx + Inches(0.15), fy + Pt(14),
                      cw - Inches(0.3), Pt(12), size=Pt(8), color=self.GRAY)
            fy += Inches(0.38)
        for call in future_calls:
            self._txt(s, '\u2705 ' + call, rx + Inches(0.15), fy,
                      cw - Inches(0.3), Pt(12), size=Pt(9), color=self.GREEN)
            fy += Pt(14)
        # Value badge
        self._txt(s, value_text, self.ML, Inches(6.8), self.CW, Pt(16),
                  size=Pt(12), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)

    def _slide_onboarding(self):
        self._journey_slide(
            1, 'Client Onboarding', 'From 9.6 Days to 3 Days',
            '9.6 days', '3 days',
            [('Initial Meeting', 'Advisor', '90 min'),
             ('Document Collection', 'Client + Advisor', '3\u20135 days'),
             ('Manual Data Entry', 'Operations', '2 hrs'),
             ('Compliance Check', 'Compliance', '2\u20133 days'),
             ('Mandate Setup', 'Operations', '1 day'),
             ('Welcome Pack', 'Admin', '1 day')],
            [('Digital Pre-boarding', 'Client', '15 min'),
             ('Advisor Review & Customise', 'Advisor', '30 min'),
             ('Automated Compliance', 'System', 'Real-time'),
             ('Digital Mandate & Signature', 'Client', '15 min'),
             ('Account Live', 'System', 'Same day')],
            ['4.5 hrs advisor effort', '35% rejection rate', '5 hrs back-office'],
            ['78% less advisor effort', '69% faster SLA', '<5% rejection'],
            '\u00a366M 5-year value at stake')

    def _slide_client_review(self):
        self._journey_slide(
            2, 'Client Review Preparation', 'From 4\u20136 Hours to 45 Minutes',
            '4\u20136 hrs', '45 min',
            [('Pull Data \u2014 Custody', 'Advisor', '30 min'),
             ('Pull Data \u2014 CRM', 'Advisor', '20 min'),
             ('Pull Data \u2014 4 More Systems', 'Advisor', '65 min'),
             ('Manual Consolidation', 'Advisor', '60\u201390 min'),
             ('Cross-check & Format', 'Advisor', '30\u201345 min')],
            [('AI Agent Auto-prepares Pack', 'System', 'Overnight'),
             ('Advisor Reviews Insights', 'Advisor', '30 min'),
             ('Customise Talking Points', 'Advisor', '15 min'),
             ('One-click PDF & Brief', 'System', 'Instant')],
            ['6\u20137 systems accessed', 'High error risk'],
            ['1 unified cockpit', '85% time reduction'],
            '\u00a336M 5-year value at stake')

    def _slide_wealth_planning(self):
        self._journey_slide(
            3, 'Wealth Planning', 'From ~6.5 Hours to ~1.5 Hours',
            '~6.5 hrs', '~1.5 hrs',
            [('Paper Fact-find Form', 'Advisor + Client', '90 min'),
             ('Manual Voyant Data Entry', 'Admin', '60 min'),
             ('Scenario Modelling', 'Advisor', '30 min'),
             ('Reasons-Why Letter', 'Advisor', '3 hrs'),
             ('Print & Post', 'Admin', '2\u20133 days')],
            [('Digital Fact-find', 'Client pre-fills', '20 min'),
             ('Auto-populate Voyant', 'System', 'Instant'),
             ('Live Scenario Planning', 'Advisor + Client', '30 min'),
             ('AI-generated Letter', 'System + Advisor', '30 min'),
             ('Digital Delivery', 'System', 'Instant')],
            ['90-min fact-finding', '3-hr letters'],
            ['78% faster fact-find', '83% faster letters'],
            '\u00a36.4M 5-year value at stake')

    def _slide_heatmap(self):
        s = self._new_slide(dark=False)
        self._section_label(s, 'Capability Assessment')
        # Override heading for dark text
        tb = s.shapes.add_textbox(self.ML, Inches(0.85), self.CW, Inches(0.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for txt, col in [('12 Core Capabilities \u2014 ', self.DARK_TEXT),
                         ('Current vs. Target', self.BLUE)]:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(30)
            r.font.bold = True
            r.font.color.rgb = col
            r.font.name = self.FONT
        # Table
        headers = ['Capability', 'Front', 'Mid', 'Back', 'Current', 'Target', 'Priority']
        col_ws = [Inches(2.8), Inches(0.7), Inches(0.7), Inches(0.7),
                  Inches(1.0), Inches(1.0), Inches(1.0)]
        tbl_shape = s.shapes.add_table(
            len(self.HEATMAP_ROWS) + 1, len(headers),
            Inches(1.8), Inches(1.7),
            sum(col_ws), Inches(5.0)
        )
        tbl = tbl_shape.table
        for i, w in enumerate(col_ws):
            tbl.columns[i].width = w
        # Header
        for i, h in enumerate(headers):
            cell = tbl.cell(0, i)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xE8, 0xEE)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.bold = True
                p.font.color.rgb = self.MID
                p.font.name = self.FONT
                p.alignment = PP_ALIGN.CENTER
        # Rows
        rag_colors = {0: self.RAG_ABSENT, 1: self.RAG_FRAGMENTED,
                      2: self.RAG_DEFINED, 3: self.RAG_ORCHESTRATED}
        pri_colors = {'Critical': self.RAG_ABSENT, 'High': self.RAG_FRAGMENTED,
                      'Medium': self.RAG_DEFINED}
        for r, row in enumerate(self.HEATMAP_ROWS):
            for c in range(7):
                cell = tbl.cell(r + 1, c)
                val = row[c]
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(8)
                    p.font.name = self.FONT
                    p.font.color.rgb = self.DARK_TEXT
                    p.alignment = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
                # Color coding for current/target
                if c == 4 and isinstance(val, int):
                    cell.fill.fore_color.rgb = rag_colors.get(val, self.WHITE)
                    for p in cell.text_frame.paragraphs:
                        p.font.color.rgb = self.WHITE
                        p.font.bold = True
                if c == 5 and isinstance(val, int):
                    cell.fill.fore_color.rgb = rag_colors.get(val, self.WHITE)
                    for p in cell.text_frame.paragraphs:
                        p.font.color.rgb = self.WHITE
                        p.font.bold = True
                if c == 6:
                    for p in cell.text_frame.paragraphs:
                        p.font.color.rgb = pri_colors.get(str(val), self.MID)
                        p.font.bold = True

    def _slide_business_case(self):
        s = self._new_slide()
        self._section_label(s, 'Strategic Fit \u2014 Project Nova')
        self._heading(s, 'The Business Case \u2014 ', '\u00a3133.5M Over Five Years')
        # 3 ROI cards
        cw = Inches(3.7)
        gap = Inches(0.35)
        sx = Inches(1.0)
        cards = [
            ('15.3\u00d7', 'Return on Investment', '\u00a3133.5M benefits on \u00a38.7M investment', False),
            ('<12', 'Months to Payback', 'Cumulative positive within Year 1', True),
            ('\u00a395M', 'Net Present Value', 'NPV at 8% discount rate', False),
        ]
        for i, (v, l, d, hl) in enumerate(cards):
            self._stat_card(s, v, l, d, sx + i * (cw + gap), Inches(1.8),
                           cw, Inches(1.1), highlight=hl)
        # 4 lever cards
        lw = Inches(2.85)
        lg = Inches(0.15)
        for i, (name, pct, val, col, _) in enumerate(self.ROI_LEVERS):
            x = self.ML + i * (lw + lg)
            y = Inches(3.2)
            self._card(s, x, y, lw, Inches(1.2), border=col)
            self._txt(s, pct, x + Inches(0.12), y + Inches(0.08),
                      Inches(0.8), Pt(24), size=Pt(20), bold=True, color=col)
            self._txt(s, name, x + Inches(0.12), y + Inches(0.45),
                      lw - Inches(0.25), Pt(14), size=Pt(10), bold=True)
            self._txt(s, val, x + Inches(0.12), y + Inches(0.7),
                      lw - Inches(0.25), Pt(12), size=Pt(9), color=self.GRAY)
            self._txt(s, '\u25bc Detail on next slide',
                      x + Inches(0.12), y + Inches(0.92),
                      lw - Inches(0.25), Pt(12), size=Pt(7), color=self.BLUE)

    def _slide_lever_detail(self):
        s = self._new_slide()
        self._section_label(s, 'Value Lever Breakdown')
        self._heading(s, 'Four Strategic Fit Pillars \u2014 ', 'Sub-Lever Detail', size=Pt(28))
        lw = Inches(2.85)
        lg = Inches(0.15)
        for i, (name, pct, val, col, subs) in enumerate(self.ROI_LEVERS):
            x = self.ML + i * (lw + lg)
            h = Inches(0.55) + len(subs) * Inches(0.42)
            y = Inches(1.7)
            self._card(s, x, y, lw, h, border=col)
            self._txt(s, f'{name} ({val})', x + Inches(0.12), y + Inches(0.08),
                      lw - Inches(0.25), Pt(16), size=Pt(11), bold=True, color=col)
            sy = y + Inches(0.45)
            for sub_name, sub_val in subs:
                self._txt(s, sub_name, x + Inches(0.12), sy,
                          lw - Inches(0.25), Pt(12), size=Pt(9))
                self._txt(s, sub_val, x + Inches(0.12), sy + Pt(13),
                          lw - Inches(0.25), Pt(12), size=Pt(9),
                          color=self.BLUE, bold=True)
                sy += Inches(0.42)

    def _slide_peer_validation(self):
        s = self._new_slide()
        self._section_label(s, 'Peer Validation')
        self._heading(s, 'Where You Stand vs. ', 'Best-in-Class')
        bar_w = Inches(10.5)
        bx = Inches(1.4)
        for i, (name, s_label, b_label, pos) in enumerate(self.PEER_METRICS):
            by = Inches(1.85) + i * Inches(0.82)
            # Metric name
            self._txt(s, name, bx, by, bar_w, Pt(14),
                      size=Pt(10), bold=True)
            # Bar bg
            bar_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        bx, by + Pt(18), bar_w, Pt(18))
            bar_bg.fill.solid()
            bar_bg.fill.fore_color.rgb = RGBColor(0x0E, 0x22, 0x40)
            bar_bg.line.fill.background()
            # Range fill
            rf = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    bx + Emu(int(bar_w * 0.1)),
                                    by + Pt(18),
                                    Emu(int(bar_w * 0.8)), Pt(18))
            rf.fill.solid()
            rf.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x6E)
            rf.line.fill.background()
            # Marker
            mk = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    bx + Emu(int(bar_w * pos)),
                                    by + Pt(14), Pt(3), Pt(26))
            mk.fill.solid()
            mk.fill.fore_color.rgb = self.BLUE
            mk.line.fill.background()
            # Labels
            self._txt(s, s_label, bx, by + Pt(38), Inches(4), Pt(12),
                      size=Pt(8), color=self.RED, bold=True)
            self._txt(s, b_label, bx + bar_w - Inches(3), by + Pt(38),
                      Inches(3), Pt(12), size=Pt(8), color=self.GREEN,
                      bold=True, align=PP_ALIGN.RIGHT)
        self._insight_box(s,
            'Heavy under-investment in digital has led to regressive innovation \u2014 '
            'widening gap vs. peers deploying AI-native platforms.')

    def _slide_ai_maturity(self):
        s = self._new_slide()
        self._section_label(s, 'AI for Wealth')
        self._heading(s, 'AI Maturity Model \u2014 ', 'Wealth Management', size=Pt(28))
        self._subtitle(s, 'Columns = wealth lifecycle stages  \u00b7  Rows = AI maturity levels')
        # Column headers
        cols = ['Prospecting', 'Onboarding', 'Client Review', 'Client Servicing', 'Growth & Retention']
        cw = Inches(2.2)
        gap = Inches(0.1)
        sx = Inches(1.2)
        hy = Inches(1.85)
        for i, c in enumerate(cols):
            x = sx + i * (cw + gap)
            self._card(s, x, hy, cw, Pt(24), fill=RGBColor(0x14, 0x2A, 0x50))
            self._txt(s, c, x, hy + Pt(3), cw, Pt(18),
                      size=Pt(9), bold=True, color=self.BLUE, align=PP_ALIGN.CENTER)
        # Level rows
        levels = [
            ('L1 \u2014 Automation', [
                'Campaign gen', 'Form & doc automation', 'Review decks', 'Self-service', 'Data sync']),
            ('L2 \u2014 Predictive', [
                'Lead scoring', 'KYC triage', 'Review triggers', 'Next Best Action', 'Churn prediction']),
            ('L3 \u2014 Embedded', [
                'Prospect lounge', 'Orchestration', 'AI proposals', 'Conversational AI', 'Dynamic pricing']),
            ('L4 \u2014 Autonomous', [
                'Self-optimising', 'Straight-through', 'Auto-rebalance', 'Agentic ops', 'Self-healing']),
        ]
        rh = Inches(0.85)
        rg = Inches(0.08)
        for ri, (level_name, tiles) in enumerate(levels):
            ry = Inches(2.25) + ri * (rh + rg)
            # Row label
            self._txt(s, level_name, Inches(0.1), ry + Inches(0.15),
                      Inches(1.05), Inches(0.5),
                      size=Pt(8), bold=True, color=self.GRAY)
            is_future = ri == 3
            for ci, tile in enumerate(tiles):
                tx = sx + ci * (cw + gap)
                fill = RGBColor(0x0A, 0x18, 0x2D) if is_future else self.CARD_BG
                bdr = self.MID if is_future else self.BORDER
                self._card(s, tx, ry, cw, rh, fill=fill, border=bdr)
                tc = self.MID if is_future else self.WHITE
                self._txt(s, tile, tx + Inches(0.08), ry + Inches(0.12),
                          cw - Inches(0.16), Pt(14),
                          size=Pt(9), bold=True, color=tc)
                if is_future:
                    self._txt(s, 'FUTURE', tx + Inches(0.08), ry + Inches(0.48),
                              Inches(0.8), Pt(10),
                              size=Pt(6), color=self.MID, bold=True)
        # Legend
        self._txt(s, '\U0001f534 Not started    \U0001f7e0 Basic    \U0001f535 Developing    \U0001f7e2 Integrated',
                  self.ML, Inches(6.2), self.CW, Pt(14),
                  size=Pt(9), color=self.GRAY, align=PP_ALIGN.CENTER)
        self._txt(s, 'Click-through detail on next slide \u2192',
                  self.ML, Inches(6.5), self.CW, Pt(14),
                  size=Pt(9), color=self.BLUE, align=PP_ALIGN.CENTER)

    def _slide_ai_use_cases(self):
        s = self._new_slide()
        self._section_label(s, 'AI for Wealth \u2014 Selected Use Cases')
        self._heading(s, 'Spotlight: Six ', 'High-Impact AI Capabilities', size=Pt(28))
        cw = Inches(3.7)
        gap = Inches(0.3)
        ch = Inches(2.2)
        for i, (level, level_name, title, desc, product, value, phase) in enumerate(self.AI_SPOTLIGHT):
            col = i % 3
            row = i // 3
            x = Inches(0.65) + col * (cw + gap)
            y = Inches(1.75) + row * (ch + Inches(0.15))
            lvl_col = self.AMBER if level == 'L1' else self.CYAN if level == 'L2' else self.PURPLE
            self._card(s, x, y, cw, ch, border=lvl_col)
            # Level badge
            self._txt(s, f'{level} \u00b7 {level_name}', x + Inches(0.12), y + Inches(0.08),
                      cw - Inches(0.25), Pt(12), size=Pt(7), color=lvl_col, bold=True)
            self._txt(s, title, x + Inches(0.12), y + Inches(0.30),
                      cw - Inches(0.25), Pt(16), size=Pt(12), bold=True)
            self._txt(s, desc, x + Inches(0.12), y + Inches(0.58),
                      cw - Inches(0.25), Inches(0.4), size=Pt(9), color=self.GRAY)
            self._txt(s, product, x + Inches(0.12), y + Inches(1.0),
                      cw - Inches(0.25), Pt(12), size=Pt(7), color=self.BLUE, bold=True)
            # Value + phase
            self._txt(s, value, x + Inches(0.12), y + Inches(1.3),
                      cw - Inches(0.25), Pt(14), size=Pt(10), color=self.GREEN, bold=True)
            self._txt(s, phase, x + Inches(0.12), y + Inches(1.6),
                      cw - Inches(0.25), Pt(12), size=Pt(8), color=self.GRAY)

    def _slide_roadmap(self):
        s = self._new_slide()
        self._section_label(s, 'Implementation Roadmap')
        self._heading(s, 'Three Phases to ', 'Full Value')
        cw = Inches(3.7)
        gap = Inches(0.3)
        sx = Inches(1.0)
        phase_colors = [self.BLUE, self.PURPLE, self.GREEN]
        for i, (title, sub, value, items) in enumerate(self.ROADMAP_PHASES):
            x = sx + i * (cw + gap)
            y = Inches(1.8)
            col = phase_colors[i]
            h = Inches(3.4)
            self._card(s, x, y, cw, h, border=col)
            self._txt(s, title, x + Inches(0.12), y + Inches(0.08),
                      cw - Inches(0.25), Pt(16), size=Pt(11), bold=True, color=col)
            self._txt(s, sub, x + Inches(0.12), y + Inches(0.35),
                      cw - Inches(0.25), Pt(14), size=Pt(9), color=self.GRAY)
            iy = y + Inches(0.65)
            for item in items:
                self._txt(s, '\u2022 ' + item, x + Inches(0.12), iy,
                          cw - Inches(0.25), Pt(14), size=Pt(9))
                iy += Pt(16)
            self._txt(s, value, x + Inches(0.12), y + h - Inches(0.35),
                      cw - Inches(0.25), Pt(16), size=Pt(11),
                      color=col, bold=True)
        # Benefit bars
        by = Inches(5.5)
        self._txt(s, 'Annual Benefit Ramp:', self.ML, by, Inches(2), Pt(14),
                  size=Pt(10), bold=True, color=self.GRAY)
        bar_max = 37
        bar_w_max = Inches(8.5)
        bx = Inches(3.5)
        for i, (label, val) in enumerate(self.BENEFIT_BARS):
            y = by + Inches(0.05) + i * Inches(0.35)
            self._txt(s, f'Y{i+1}', bx - Inches(0.4), y, Inches(0.35), Pt(14),
                      size=Pt(9), color=self.GRAY, align=PP_ALIGN.RIGHT)
            bw = Emu(int(bar_w_max * val / bar_max))
            bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     bx, y + Pt(2), bw, Pt(20))
            bar.fill.solid()
            bar.fill.fore_color.rgb = self.BLUE
            bar.line.fill.background()
            self._txt(s, label, bx + bw + Inches(0.1), y, Inches(1), Pt(14),
                      size=Pt(9), color=self.BLUE, bold=True)

    def _slide_what_to_expect(self):
        s = self._new_slide()
        self._section_label(s, 'Looking Ahead')
        self._heading(s, 'What to Expect ', 'Today')
        self._subtitle(s, 'The sessions ahead will bring the numbers to life with live demonstrations.')
        sessions = [
            ('\U0001f50d', 'Morning Session', 'Product & Integration Deep-Dive',
             ['Live platform demonstrations', 'Architecture & integration approach',
              'Grand Central in action']),
            ('\U0001f91d', 'Afternoon Session', 'Delivery & Collaboration',
             ['Implementation approach & timeline', 'Roles, responsibilities & governance',
              'Change management considerations']),
            ('\U0001f3af', 'Closing Session', 'Alignment & Discussion',
             ['Open Q&A with full team', 'Key decisions and alignment',
              'Agreed next steps']),
        ]
        cw = Inches(3.7)
        gap = Inches(0.35)
        sx = Inches(1.0)
        for i, (icon, title, sub, items) in enumerate(sessions):
            x = sx + i * (cw + gap)
            y = Inches(2.0)
            self._card(s, x, y, cw, Inches(3.2))
            self._txt(s, icon, x + Inches(0.15), y + Inches(0.12),
                      Inches(0.4), Pt(28), size=Pt(22))
            self._txt(s, title, x + Inches(0.15), y + Inches(0.55),
                      cw - Inches(0.3), Pt(18), size=Pt(14), bold=True)
            self._txt(s, sub, x + Inches(0.15), y + Inches(0.85),
                      cw - Inches(0.3), Pt(14), size=Pt(10), color=self.GRAY)
            iy = y + Inches(1.2)
            for item in items:
                self._txt(s, '\u2022 ' + item, x + Inches(0.15), iy,
                          cw - Inches(0.3), Pt(14), size=Pt(10))
                iy += Pt(16)

    def _slide_close(self):
        s = self._new_slide()
        self._section_label(s, 'Schroders Wealth Management \u00d7 Backbase')
        tb = s.shapes.add_textbox(self.ML, Inches(1.5), self.CW, Inches(4))
        tf = tb.text_frame
        tf.word_wrap = True
        lines = [('THE RELATIONSHIP', self.WHITE),
                 ('STAYS PERSONAL.', self.WHITE),
                 ('EVERYTHING AROUND IT', self.BLUE),
                 ('BECOMES DIGITAL.', self.BLUE)]
        for txt, color in lines:
            p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(46)
            r.font.bold = True
            r.font.color.rgb = color
            r.font.name = self.FONT
        self._txt(s, 'February 2026', self.ML, Inches(6.2), self.CW, Pt(16),
                  size=Pt(12), color=self.GRAY, align=PP_ALIGN.CENTER)

    def _slide_strategic_fit(self):
        s = self._new_slide()
        self._txt(s, 'APPENDIX', self.ML, Inches(0.2), Inches(1), Pt(12),
                  size=Pt(8), color=self.MID, bold=True)
        self._section_label(s, 'Why Backbase for Schroders')
        self._heading(s, 'Strategic Fit \u2014 ', 'Four Arguments')
        cw = Inches(2.85)
        gap = Inches(0.15)
        for i, (icon, title, bullets) in enumerate(self.STRATEGIC_FIT):
            x = self.ML + i * (cw + gap)
            h = Inches(3.6)
            self._info_card(s, title, bullets, x, Inches(1.7), cw, h,
                           border_color=self.BLUE, icon=icon)

    def _slide_platform_proof(self):
        s = self._new_slide()
        self._txt(s, 'APPENDIX', self.ML, Inches(0.2), Inches(1), Pt(12),
                  size=Pt(8), color=self.MID, bold=True)
        self._section_label(s, 'Proven Platform Advantage')
        self._heading(s, 'Results from ', 'Backbase Clients')
        cw = Inches(3.7)
        gap = Inches(0.3)
        sx = Inches(1.0)
        for i, (val, label, detail) in enumerate(self.PROOF_POINTS):
            col = i % 3
            row = i // 3
            x = sx + col * (cw + gap)
            y = Inches(1.8) + row * Inches(1.45)
            self._stat_card(s, val, label, detail, x, y, cw, Inches(1.2))
        self._insight_box(s,
            'Grand Central architecture enables dual-core coexistence \u2014 '
            'connect to existing custody, PMS, and CRM without rip-and-replace. '
            'Unlike horizontal CRM platforms, Backbase is purpose-built for financial services.',
            top=Inches(5.0))

    def _scope_slide(self, scope_name, sub, investment, ret, npv,
                     invest_detail, roi_detail, npv_detail,
                     fte, downside, breakeven, pillars,
                     fte_d, down_d, break_d, pill_d,
                     caps, recommendation, icon, desc):
        s = self._new_slide()
        self._txt(s, 'APPENDIX \u2014 EXECUTIVE DECISION PAPER',
                  self.ML, Inches(0.2), Inches(4), Pt(12),
                  size=Pt(8), color=self.MID, bold=True)
        self._section_label(s, 'Investment Scenarios')
        self._heading(s, f'Choose Your Ambition \u2014 ', scope_name, size=Pt(30))
        self._txt(s, sub, self.ML, Inches(1.4), self.CW, Pt(14),
                  size=Pt(11), color=self.GRAY, align=PP_ALIGN.CENTER)
        # 3 KPI cards
        cw = Inches(3.7)
        gap = Inches(0.35)
        sx = Inches(1.0)
        for i, (v, l, d, hl) in enumerate([
            (investment, 'Total Investment (5yr)', invest_detail, False),
            (ret, 'Total Return (5yr)', roi_detail, True),
            (npv, 'Net Present Value', npv_detail, False),
        ]):
            self._stat_card(s, v, l, d, sx + i * (cw + gap), Inches(1.75),
                           cw, Inches(1.0), highlight=hl)
        # 4 detail cards
        dw = Inches(2.85)
        dg = Inches(0.15)
        for i, (v, l, d) in enumerate([
            (fte, 'FTE Capacity Freed', fte_d),
            (downside, 'Downside (60%)', down_d),
            (breakeven, 'Breakeven', break_d),
            (pillars, 'Pillars Included', pill_d),
        ]):
            x = self.ML + i * (dw + dg)
            y = Inches(3.0)
            self._card(s, x, y, dw, Inches(0.8))
            self._txt(s, l.upper(), x + Inches(0.1), y + Inches(0.05),
                      dw - Inches(0.2), Pt(10), size=Pt(7),
                      color=self.GRAY, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, v, x + Inches(0.1), y + Inches(0.25),
                      dw - Inches(0.2), Pt(24), size=Pt(18),
                      bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, d, x + Inches(0.1), y + Inches(0.55),
                      dw - Inches(0.2), Pt(12), size=Pt(7),
                      color=self.GRAY, align=PP_ALIGN.CENTER)
        # Capabilities
        self._txt(s, 'CAPABILITIES INCLUDED:', self.ML, Inches(4.0),
                  Inches(3), Pt(14), size=Pt(8), color=self.GRAY, bold=True)
        ccw = Inches(2.85)
        for i, (cap_name, cap_status) in enumerate(caps):
            col = i % 4
            row = i // 4
            x = self.ML + col * (ccw + Inches(0.15))
            y = Inches(4.25) + row * Inches(0.42)
            icon_txt = '\u2705' if cap_status == 'included' else '\U0001f513' if cap_status == 'unlocked' else '\U0001f512'
            self._txt(s, f'{icon_txt} {cap_name}', x, y,
                      ccw, Pt(14), size=Pt(9),
                      color=self.WHITE if cap_status != 'locked' else self.GRAY)
        # Recommendation
        self._insight_box(s, 'Recommendation: ' + recommendation, top=Inches(5.3))

    def _slide_scope_premium(self):
        self._scope_slide(
            'Digital Wealth Premium', 'Channel Replacement',
            '\u00a35.8M', '\u00a383M', '\u00a358M',
            'Licence + implementation + change', '14.3\u00d7 ROI \u2022 14-month payback',
            'NPV at 8% discount rate',
            '~65', '8.6\u00d7', '7.0%', '2 of 4',
            'Across onboarding, review & admin',
            'Even at 60% of projected benefits',
            'Of projected benefits needed to break even',
            'Growth Engine + Operations Factory',
            [('Prospect Portal & Digital Onboarding', 'included'),
             ('Unified Advisor Workspace (Client 360\u00b0)', 'included'),
             ('Web & Mobile Client App', 'included'),
             ('Omni-channel Orchestration', 'included'),
             ('Next Best Action & Smart Signals', 'locked'),
             ('Grand Central (iPaaS Integration Layer)', 'locked'),
             ('Agentic Studio & AI Workflows', 'locked'),
             ('Personalised Experiences & Entitlements', 'locked')],
            'Start with Premium to modernise channels and build the digital foundation. '
            'Proven 14-month payback. Phase 1 delivers immediate capacity relief while '
            'creating the architecture for an AI-first future.',
            '\U0001f4e6', 'Replace existing channels with a unified digital wealth platform.')

    def _slide_scope_signature(self):
        self._scope_slide(
            'Signature + Grand Central', 'AI-First Architecture',
            '\u00a312.5M', '\u00a3133.5M', '\u00a395M',
            'Full suite + iPaaS + AI platform', '10.7\u00d7 ROI \u2022 <12-month payback',
            'NPV at 8% discount rate',
            '99', '6.4\u00d7', '9.4%', '4 of 4',
            'Full FTE impact across all journeys',
            'Even at 60% of projected benefits',
            'Of projected benefits needed to break even',
            'All four Strategic Fit pillars',
            [('Prospect Portal & Digital Onboarding', 'included'),
             ('Unified Advisor Workspace (Client 360\u00b0)', 'included'),
             ('Web & Mobile Client App', 'included'),
             ('Omni-channel Orchestration', 'included'),
             ('Next Best Action & Smart Signals', 'unlocked'),
             ('Grand Central (iPaaS Integration Layer)', 'unlocked'),
             ('Agentic Studio & AI Workflows', 'unlocked'),
             ('Personalised Experiences & Entitlements', 'unlocked')],
            'Go bold with Signature to capture the full \u00a3133.5M opportunity. '
            'Grand Central eliminates integration debt. AI agents drive exponential returns. '
            'This is the transformational play.',
            '\U0001f680', 'Full AI-first architecture with iPaaS integration layer.')

    # ── ORCHESTRATOR ────────────────────────────────────────────

    def generate(self, output_path: str):
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_W
        self.prs.slide_height = self.SLIDE_H

        self._slide_cover()               # 1
        self._slide_agenda()               # 2
        self._slide_open_items_overview()  # 3
        self._slide_deal_critical()        # 4
        self._slide_confidence_nice()      # 5
        self._slide_vision()               # 6
        self._slide_journey_overview()     # 7
        self._slide_onboarding()           # 8
        self._slide_client_review()        # 9
        self._slide_wealth_planning()      # 10
        self._slide_heatmap()              # 11
        self._slide_business_case()        # 12
        self._slide_lever_detail()         # 13
        self._slide_peer_validation()      # 14
        self._slide_ai_maturity()          # 15
        self._slide_ai_use_cases()         # 16
        self._slide_roadmap()              # 17
        self._slide_what_to_expect()       # 18
        self._slide_close()                # 19
        self._slide_strategic_fit()        # 20
        self._slide_platform_proof()       # 21
        self._slide_scope_premium()        # 22
        self._slide_scope_signature()      # 23

        self.prs.save(output_path)
        print(f'Saved: {output_path} ({Path(output_path).stat().st_size // 1024} KB)')


if __name__ == '__main__':
    OUTPUT_DIR = Path(__file__).parent.parent / (
        'Engagement/Schroders Group/Past working folder/'
        'schroders/2026-02_wealth_playback/outputs'
    )
    output_file = OUTPUT_DIR / 'Schroders_Playback_Day2.pptx'
    SchrodersPptxGenerator().generate(str(output_file))
