#!/usr/bin/env python3
"""
Schroders Executive Read-Out v3 — PPTX Generator

10-slide simplified deck for COO meeting (March 9, 2026).
Google Slides compatible. Backbase brand (Libre Franklin, #3366FF).

Usage:
    python3 tools/schroders_executive_v3_pptx.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


class SchrordersV3PptxGenerator:
    """Generates a 10-slide executive read-out PPTX."""

    # ── Dimensions ─────────────────────────────────────────
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    # ── Colors (Backbase Brand) ────────────────────────────
    DARK_BG    = RGBColor(0x09, 0x1C, 0x35)
    CARD_BG    = RGBColor(0x0E, 0x22, 0x40)
    BLUE       = RGBColor(0x33, 0x66, 0xFF)
    PURPLE     = RGBColor(0x7B, 0x2F, 0xFF)
    RED        = RGBColor(0xDC, 0x26, 0x26)
    GREEN      = RGBColor(0x05, 0x96, 0x69)
    AMBER      = RGBColor(0xD9, 0x77, 0x06)
    CYAN       = RGBColor(0x08, 0x91, 0xB2)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT_BG   = RGBColor(0xF5, 0xFA, 0xFF)
    LIGHT_CARD = RGBColor(0xFF, 0xFF, 0xFF)
    DARK_TEXT   = RGBColor(0x09, 0x1C, 0x35)
    SUB_TEXT    = RGBColor(0x64, 0x74, 0x8B)
    MUTED      = RGBColor(0x94, 0xA3, 0xB8)
    BORDER     = RGBColor(0xE2, 0xE8, 0xF0)
    BORDER_DK  = RGBColor(0x1A, 0x3A, 0x5E)
    BLUE_LIGHT = RGBColor(0xEF, 0xF6, 0xFF)
    GREEN_LIGHT = RGBColor(0xEC, 0xFD, 0xF5)
    AMBER_LIGHT = RGBColor(0xFF, 0xFB, 0xEB)
    PURPLE_LIGHT = RGBColor(0xF5, 0xF3, 0xFF)

    FONT = 'Libre Franklin'
    ML = Inches(0.6)     # margin left
    MR = Inches(0.6)     # margin right
    CW = Inches(12.1)    # content width
    CT = Inches(1.6)     # content top

    # ── HELPERS ─────────────────────────────────────────────

    def _new_slide(self, dark=False):
        layout = self.prs.slide_layouts[6]  # blank
        slide = self.prs.slides.add_slide(layout)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = self.DARK_BG if dark else self.LIGHT_BG
        return slide

    def _txt(self, slide, text, left, top, width, height, size=Pt(12),
             color=None, bold=False, align=PP_ALIGN.LEFT, font=None):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = size
        p.font.bold = bold
        p.font.color.rgb = color or self.DARK_TEXT
        p.font.name = font or self.FONT
        p.alignment = align
        return tf

    def _section_label(self, slide, text, top=Inches(0.4), color=None, dark=False):
        self._txt(slide, text.upper(), self.ML, top, self.CW, Pt(16),
                  size=Pt(9), color=color or (self.BLUE if not dark else RGBColor(0x69, 0xFE, 0xFF)),
                  bold=True, align=PP_ALIGN.LEFT)

    def _heading(self, slide, main, accent='', top=Inches(0.75), size=Pt(32),
                 dark=False):
        tb = slide.shapes.add_textbox(self.ML, top, self.CW, Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = main
        r1.font.size = size
        r1.font.bold = True
        r1.font.color.rgb = self.WHITE if dark else self.DARK_TEXT
        r1.font.name = self.FONT
        if accent:
            r2 = p.add_run()
            r2.text = accent
            r2.font.size = size
            r2.font.bold = True
            r2.font.color.rgb = self.BLUE if not dark else RGBColor(0x69, 0xFE, 0xFF)
            r2.font.name = self.FONT

    def _subtitle(self, slide, text, top=Inches(1.25), dark=False):
        self._txt(slide, text, self.ML, top, self.CW, Inches(0.35),
                  size=Pt(12), color=self.MUTED if dark else self.SUB_TEXT)

    def _card(self, slide, left, top, w, h, fill=None, border=None, dark=False):
        s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = fill or (self.CARD_BG if dark else self.LIGHT_CARD)
        s.line.color.rgb = border or (self.BORDER_DK if dark else self.BORDER)
        s.line.width = Pt(1)
        return s

    def _stat_box(self, slide, value, label, left, top, w, h,
                  val_color=None, bg=None):
        """Compact stat box — value on top, label below."""
        self._card(slide, left, top, w, h, fill=bg or self.BLUE_LIGHT)
        self._txt(slide, str(value), left, top + Inches(0.08), w, Inches(0.35),
                  size=Pt(20), color=val_color or self.BLUE, bold=True,
                  align=PP_ALIGN.CENTER)
        self._txt(slide, label.upper(), left, top + Inches(0.42), w, Pt(12),
                  size=Pt(7), color=self.SUB_TEXT, bold=True,
                  align=PP_ALIGN.CENTER)

    def _multi_text(self, slide, runs, left, top, width, height, align=PP_ALIGN.LEFT):
        """Add textbox with multiple styled runs."""
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        for text, size, color, bold in runs:
            r = p.add_run()
            r.text = str(text)
            r.font.size = size
            r.font.bold = bold
            r.font.color.rgb = color
            r.font.name = self.FONT
        return tf

    def _divider(self, slide, top):
        """Thin horizontal line."""
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   self.ML, top, self.CW, Pt(1))
        s.fill.solid()
        s.fill.fore_color.rgb = self.BORDER
        s.line.fill.background()

    # ── SLIDE BUILDERS ──────────────────────────────────────

    def _slide_01_cover(self):
        """Dark hero cover slide."""
        s = self._new_slide(dark=True)
        # Section label
        self._section_label(s, 'Executive Read-Out · March 9th, 2026',
                           top=Inches(0.5), dark=True)
        # Mega title
        tb = s.shapes.add_textbox(self.ML, Inches(1.3), self.CW, Inches(2.0))
        tf = tb.text_frame
        tf.word_wrap = True
        for line, color in [('Accelerating the', self.WHITE),
                            ('Schroders Advantage.', self.BLUE)]:
            p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = line
            r.font.size = Pt(48)
            r.font.bold = True
            r.font.color.rgb = color
            r.font.name = self.FONT
        # Subtitle
        self._txt(s, 'A digitally enabled, insight-led client and advisor platform.\nFrom relationship depth to digital scale.',
                  self.ML, Inches(3.2), Inches(8), Inches(0.6),
                  size=Pt(14), color=self.MUTED)

        # 4 pill tags
        tags = ['Advisor Productivity', 'Client Retention & AUM Growth',
                'Innovation Agility', 'No AI Lock-In']
        tx = self.ML
        for tag in tags:
            pw = Inches(len(tag) * 0.1 + 0.3)
            self._card(s, tx, Inches(4.1), pw, Inches(0.3),
                      fill=RGBColor(0x0E, 0x22, 0x40), border=self.BORDER_DK,
                      dark=True)
            self._txt(s, tag, tx + Inches(0.1), Inches(4.12), pw - Inches(0.2), Inches(0.25),
                     size=Pt(8), color=self.MUTED, bold=True)
            tx += pw + Inches(0.15)

        # 3 hero stats
        stat_w = Inches(3.4)
        gap = Inches(0.4)
        sx = self.ML
        for val, label, vcol in [
            ('£148M', 'Value at Stake', self.GREEN),
            ('~7.0×', 'Return on Investment', self.BLUE),
            ('Month 8', 'Breakeven', self.WHITE),
        ]:
            self._card(s, sx, Inches(5.0), stat_w, Inches(1.2),
                      fill=self.CARD_BG, dark=True)
            self._txt(s, val, sx, Inches(5.1), stat_w, Inches(0.5),
                     size=Pt(28), color=vcol, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, label.upper(), sx, Inches(5.65), stat_w, Pt(14),
                     size=Pt(8), color=self.MUTED, bold=True, align=PP_ALIGN.CENTER)
            sx += stat_w + gap

        # Footer
        self._txt(s, 'v3.0 · Schroders × Backbase',
                  self.ML, Inches(6.9), self.CW, Pt(12),
                  size=Pt(8), color=RGBColor(0x3A, 0x49, 0x5D))

    def _slide_02_context(self):
        """The Case Is Clear — context and maturity."""
        s = self._new_slide()
        self._section_label(s, 'Context · Where We Stand')
        self._heading(s, 'The Case Is Clear. ', 'The Moment Is Now.')
        self._subtitle(s, '10 workshops · 12 capabilities assessed · Master analysis validated')

        # Today stats row
        self._txt(s, 'TODAY · SCHRODERS WEALTH', self.ML, Inches(1.7), self.CW, Pt(12),
                 size=Pt(8), color=self.BLUE, bold=True)

        stats = [('225', 'Client Advisors'), ('35K+', 'Clients'), ('£72B', 'AUM')]
        sx = self.ML
        for val, label in stats:
            self._stat_box(s, val, label, sx, Inches(2.0), Inches(3.4), Inches(0.7))
            sx += Inches(3.8)

        # Maturity assessment
        self._divider(s, Inches(3.0))
        self._txt(s, 'CAPABILITY MATURITY ASSESSMENT', self.ML, Inches(3.2), self.CW, Pt(12),
                 size=Pt(8), color=self.BLUE, bold=True)

        # Current vs Target
        self._card(s, self.ML, Inches(3.5), Inches(5.5), Inches(1.6))
        self._txt(s, '12', self.ML + Inches(0.3), Inches(3.55), Inches(1), Inches(0.5),
                 size=Pt(32), color=self.BLUE, bold=True)
        self._txt(s, 'Core Capabilities', self.ML + Inches(1.3), Inches(3.65), Inches(3), Pt(14),
                 size=Pt(12), color=self.DARK_TEXT, bold=True)
        # Current score
        self._txt(s, 'Current: 0.6 / 4.0', self.ML + Inches(0.3), Inches(4.2), Inches(4), Pt(12),
                 size=Pt(11), color=self.RED, bold=True)
        self._txt(s, 'Target: 3.5 / 4.0', self.ML + Inches(0.3), Inches(4.55), Inches(4), Pt(12),
                 size=Pt(11), color=self.GREEN, bold=True)

        # Value at stake card
        self._card(s, Inches(6.8), Inches(3.5), Inches(5.5), Inches(1.6),
                  fill=self.BLUE_LIGHT, border=self.BLUE)
        self._txt(s, '£148M', Inches(7.0), Inches(3.6), Inches(4), Inches(0.5),
                 size=Pt(36), color=self.GREEN, bold=True)
        self._txt(s, 'Five-Year Value (AI from Year 2)', Inches(7.0), Inches(4.15), Inches(5), Pt(14),
                 size=Pt(11), color=self.DARK_TEXT)
        self._txt(s, '~7.0×', Inches(7.0), Inches(4.55), Inches(2), Pt(14),
                 size=Pt(14), color=self.BLUE, bold=True)
        self._txt(s, 'Return on Investment', Inches(8.2), Inches(4.58), Inches(3), Pt(14),
                 size=Pt(10), color=self.SUB_TEXT)

        # Bottom insight
        self._card(s, self.ML, Inches(5.6), self.CW, Inches(1.0),
                  fill=self.LIGHT_CARD, border=self.BORDER)
        self._txt(s, 'Every peer is investing. SJP: £150M+ tech transformation. Quilter: £40M+ platform modernisation. JP Morgan: $2B+ AI investment. The question is not whether, but how quickly.',
                  self.ML + Inches(0.3), Inches(5.7), self.CW - Inches(0.6), Inches(0.8),
                  size=Pt(10), color=self.SUB_TEXT)

    def _slide_03_platform(self):
        """What You Get — 3-card platform overview."""
        s = self._new_slide()
        self._section_label(s, 'The Platform')
        self._heading(s, 'One Platform. ', 'Every Channel. Every Data Source.')
        self._subtitle(s, 'Not a channel replacement. A platform transformation — connecting advisors, clients, systems and intelligence.')

        # 3 cards
        cw = Inches(3.7)
        gap = Inches(0.3)
        ch = Inches(4.4)
        cx = self.ML

        cards = [
            ('Digital Channels', self.BLUE, 'Client & Employee Experiences', [
                'Client portal (web & mobile)',
                'Prospect portal & onboarding',
                'RM Workspace & Digital Assist',
                'Unified employee dashboard',
                'Replaces Temenos e-Services',
            ]),
            ('Grand Central', self.GREEN, 'Data Connectivity Backbone', [
                'T24, Salesforce, Folio connectivity',
                'Configuration-based, not code',
                'API-first, event-driven',
                'Pre-built wealth connectors',
                'Backbase Marketplace',
                'Pre-existing connectors with partners',
                'Group-wide potential',
            ]),
            ('Intelligence Fabric', self.PURPLE, 'AI Orchestration Layer', [
                'Multi-model AI (Claude, GPT, Gemini)',
                'Custom Agent Factory',
                'Productised agent library',
                'RAG, observability, auditability',
                'Multi-agent governance',
            ]),
        ]

        for title, color, sub, bullets in cards:
            self._card(s, cx, Inches(1.8), cw, ch, border=color)
            # Colored top bar
            bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     cx, Inches(1.8), cw, Pt(4))
            bar.fill.solid()
            bar.fill.fore_color.rgb = color
            bar.line.fill.background()
            # Title
            self._txt(s, title.upper(), cx + Inches(0.2), Inches(2.0), cw - Inches(0.4), Pt(14),
                     size=Pt(10), color=color, bold=True)
            self._txt(s, sub, cx + Inches(0.2), Inches(2.3), cw - Inches(0.4), Pt(12),
                     size=Pt(11), color=self.DARK_TEXT, bold=True)
            # Bullets
            by = Inches(2.7)
            for b in bullets:
                self._txt(s, '→ ' + b, cx + Inches(0.2), by, cw - Inches(0.4), Pt(14),
                         size=Pt(9), color=self.SUB_TEXT)
                by += Inches(0.3)
            cx += cw + gap

        # Bottom insight
        self._card(s, self.ML, Inches(6.4), self.CW, Inches(0.7),
                  fill=self.BLUE_LIGHT, border=self.BLUE)
        self._txt(s, 'Salesforce sees CRM data. Temenos sees core data. Backbase sees everything — and makes it available to every channel, every agent, every decision.',
                  self.ML + Inches(0.3), Inches(6.5), self.CW - Inches(0.6), Inches(0.5),
                  size=Pt(10), color=self.DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)

    def _slide_04_ai(self):
        """What You Get: AI — Agent Factory + Productised Library."""
        s = self._new_slide()
        self._section_label(s, 'Intelligence Fabric')
        self._heading(s, 'Intelligent by Design.')
        self._subtitle(s, 'Two paths to AI. Build your own with Agent Factory. Activate ours from the productised library. Same infrastructure, same governance.')

        # 2-column
        hw = Inches(5.8)
        gap = Inches(0.3)

        # Left — Agent Factory
        self._card(s, self.ML, Inches(1.8), hw, Inches(2.3), border=self.PURPLE)
        self._txt(s, 'CUSTOM AGENTS · AGENT FACTORY', self.ML + Inches(0.2), Inches(1.9), hw, Pt(12),
                 size=Pt(8), color=self.PURPLE, bold=True)
        self._txt(s, 'Design. Test. Deploy. Iterate.', self.ML + Inches(0.2), Inches(2.15), hw, Pt(12),
                 size=Pt(11), color=self.DARK_TEXT, bold=True)
        bullets_l = ['Build agents from Schroders\' own policies and data',
                     'Full development lifecycle within the platform',
                     'Your IP, your logic, your competitive advantage',
                     'Available from Day 1 on Enterprise']
        by = Inches(2.55)
        for b in bullets_l:
            self._txt(s, '→ ' + b, self.ML + Inches(0.2), by, hw - Inches(0.4), Pt(12),
                     size=Pt(9), color=self.SUB_TEXT)
            by += Inches(0.25)

        # Right — Productised
        rx = self.ML + hw + gap
        self._card(s, rx, Inches(1.8), hw, Inches(2.3), border=self.BLUE)
        self._txt(s, 'PRODUCTISED AGENTS · BACKBASE LIBRARY', rx + Inches(0.2), Inches(1.9), hw, Pt(12),
                 size=Pt(8), color=self.BLUE, bold=True)
        self._txt(s, 'Ready-made. Continuously Improved.', rx + Inches(0.2), Inches(2.15), hw, Pt(12),
                 size=Pt(11), color=self.DARK_TEXT, bold=True)
        bullets_r = ['Pre-built agents for wealth management tasks',
                     'Maintained and improved by Backbase R&D',
                     'Low-effort agents bundled with the platform',
                     'Complex use cases via BIC blocks as library expands']
        by = Inches(2.55)
        for b in bullets_r:
            self._txt(s, '→ ' + b, rx + Inches(0.2), by, hw - Inches(0.4), Pt(12),
                     size=Pt(9), color=self.SUB_TEXT)
            by += Inches(0.25)

        # 7 validated use cases
        self._txt(s, 'VALIDATED USE CASES · SCOPED WITH YOUR TEAM', self.ML, Inches(4.3), self.CW, Pt(12),
                 size=Pt(8), color=self.BLUE, bold=True)

        use_cases = [
            ('Review Prep', '£2.5M/yr manual cost', self.PURPLE),
            ('Meeting Summary', 'Automated capture', self.PURPLE),
            ('Sentiment Analysis', 'Message intelligence', self.PURPLE),
            ('Auto Reply', 'Contextual drafting', self.BLUE),
            ('Doc Processing', 'Extraction & routing', self.BLUE),
            ('Smart Signals', 'Proactive alerts', self.BLUE),
            ('Personalised Outreach', 'Tailored comms', self.BLUE),
        ]
        ucw = Inches(1.6)
        ucx = self.ML
        for name, desc, color in use_cases:
            self._card(s, ucx, Inches(4.55), ucw, Inches(0.8), border=color)
            self._txt(s, name.upper(), ucx + Inches(0.08), Inches(4.6), ucw - Inches(0.16), Pt(10),
                     size=Pt(7), color=color, bold=True)
            self._txt(s, desc, ucx + Inches(0.08), Inches(4.85), ucw - Inches(0.16), Pt(10),
                     size=Pt(7), color=self.SUB_TEXT)
            ucx += ucw + Inches(0.1)

        # Bottom cards: No AI Lock-In + Multi-Agent Governance
        self._card(s, self.ML, Inches(5.6), Inches(5.8), Inches(1.1), border=self.CYAN)
        self._txt(s, 'NO AI LOCK-IN', self.ML + Inches(0.2), Inches(5.65), Inches(5), Pt(10),
                 size=Pt(8), color=self.CYAN, bold=True)
        self._txt(s, 'Multi-model architecture — Claude, GPT, Gemini, open-source. Switch providers without re-engineering agents.',
                  self.ML + Inches(0.2), Inches(5.9), Inches(5.4), Inches(0.6),
                  size=Pt(9), color=self.SUB_TEXT)

        self._card(s, Inches(6.7), Inches(5.6), Inches(5.8), Inches(1.1), border=self.PURPLE)
        self._txt(s, 'MULTI-AGENT GOVERNANCE', Inches(6.9), Inches(5.65), Inches(5), Pt(10),
                 size=Pt(8), color=self.PURPLE, bold=True)
        self._txt(s, 'Backbase as the control layer for all AI agents — including Salesforce and other systems. One governance framework. Full auditability.',
                  Inches(6.9), Inches(5.9), Inches(5.4), Inches(0.6),
                  size=Pt(9), color=self.SUB_TEXT)

    def _slide_05_commercial(self):
        """Enterprise or Pay As You Grow — side by side."""
        s = self._new_slide()
        self._section_label(s, 'Commercial Model')
        self._heading(s, 'Enterprise or ', 'Pay As You Grow.')
        self._subtitle(s, 'Same platform. Same delivery. Same value. Two commercial structures.')

        hw = Inches(5.8)
        gap = Inches(0.3)

        # ── Enterprise (left) ──
        self._card(s, self.ML, Inches(1.7), hw, Inches(4.0),
                  fill=self.BLUE_LIGHT, border=self.BLUE)
        # Recommended badge
        badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   self.ML + hw - Inches(1.5), Inches(1.75),
                                   Inches(1.4), Inches(0.25))
        badge.fill.solid()
        badge.fill.fore_color.rgb = self.BLUE
        badge.line.fill.background()
        self._txt(s, 'RECOMMENDED', self.ML + hw - Inches(1.5), Inches(1.76),
                 Inches(1.4), Inches(0.22),
                 size=Pt(7), color=self.WHITE, bold=True, align=PP_ALIGN.CENTER)

        self._txt(s, 'OPTION A · ENTERPRISE', self.ML + Inches(0.2), Inches(1.85),
                 hw, Pt(12), size=Pt(8), color=self.BLUE, bold=True)

        # Year boxes
        ent_years = [('Y1', '£2.3M', '60% discount'), ('Y2', '£2.6M', '55% discount'),
                     ('Y3', '£3.6M', '55% discount'), ('Y4', '£4.0M', '50% discount'),
                     ('Y5', '£4.4M', '45% discount')]
        bw = Inches(1.0)
        bx = self.ML + Inches(0.2)
        for yr, price, disc in ent_years:
            self._card(s, bx, Inches(2.2), bw, Inches(0.8), fill=self.LIGHT_CARD)
            self._txt(s, yr, bx, Inches(2.22), bw, Pt(8),
                     size=Pt(7), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, price, bx, Inches(2.4), bw, Pt(20),
                     size=Pt(16), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, disc, bx, Inches(2.72), bw, Pt(8),
                     size=Pt(6), color=self.GREEN, bold=True, align=PP_ALIGN.CENTER)
            bx += bw + Inches(0.08)

        self._txt(s, 'Full platform from Day 1. Scaling entry — pay less upfront, more as value compounds. All AI capabilities immediately.',
                  self.ML + Inches(0.2), Inches(3.15), hw - Inches(0.4), Inches(0.5),
                  size=Pt(9), color=self.SUB_TEXT)

        # Enterprise summary stats
        for i, (val, label) in enumerate([('£16.9M', '5-YR LICENCE'), ('~7.0×', 'ROI'), ('Month 8', 'BREAKEVEN')]):
            ex = self.ML + Inches(0.2) + i * Inches(1.8)
            self._txt(s, val, ex, Inches(3.75), Inches(1.6), Pt(18),
                     size=Pt(14), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, label, ex, Inches(4.05), Inches(1.6), Pt(10),
                     size=Pt(7), color=self.SUB_TEXT, bold=True, align=PP_ALIGN.CENTER)

        # ── PAG (right) ──
        rx = self.ML + hw + gap
        self._card(s, rx, Inches(1.7), hw, Inches(4.0),
                  fill=self.AMBER_LIGHT, border=self.AMBER)
        self._txt(s, 'OPTION B · PAY AS YOU GROW', rx + Inches(0.2), Inches(1.85),
                 hw, Pt(12), size=Pt(8), color=self.AMBER, bold=True)

        pag_years = [('Y1', '£2.7M', '30% discount'), ('Y2', '£3.5M', '30% discount'),
                     ('Y3+', '£4.8M', '30% discount')]
        bx = rx + Inches(0.2)
        bw_pag = Inches(1.6)
        for yr, price, disc in pag_years:
            self._card(s, bx, Inches(2.2), bw_pag, Inches(0.8), fill=self.LIGHT_CARD)
            self._txt(s, yr, bx, Inches(2.22), bw_pag, Pt(8),
                     size=Pt(7), color=self.AMBER, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, price, bx, Inches(2.4), bw_pag, Pt(20),
                     size=Pt(16), color=self.AMBER, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, disc, bx, Inches(2.72), bw_pag, Pt(8),
                     size=Pt(6), color=self.GREEN, bold=True, align=PP_ALIGN.CENTER)
            bx += bw_pag + Inches(0.15)

        self._txt(s, 'Modular activation. License grows as modules go live. AI unlocks progressively. Agentic Runtime from Year 1.',
                  rx + Inches(0.2), Inches(3.15), hw - Inches(0.4), Inches(0.5),
                  size=Pt(9), color=self.SUB_TEXT)

        for i, (val, label) in enumerate([('£20.6M', '5-YR LICENCE'), ('~5.9×', 'ROI'), ('Month 9', 'BREAKEVEN')]):
            ex = rx + Inches(0.2) + i * Inches(1.8)
            self._txt(s, val, ex, Inches(3.75), Inches(1.6), Pt(18),
                     size=Pt(14), color=self.AMBER, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, label, ex, Inches(4.05), Inches(1.6), Pt(10),
                     size=Pt(7), color=self.SUB_TEXT, bold=True, align=PP_ALIGN.CENTER)

        # Key insight
        self._card(s, self.ML, Inches(5.85), self.CW, Inches(0.5),
                  fill=self.BLUE_LIGHT, border=self.BLUE)
        self._txt(s, 'Enterprise is £0.4M less in Year 1 and £3.7M less over five years. Every price increase maps to a specific module activation.',
                  self.ML + Inches(0.3), Inches(5.92), self.CW - Inches(0.6), Inches(0.35),
                  size=Pt(10), color=self.DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)

        # Tier swimlane
        self._txt(s, 'MODULE ACTIVATION BY YEAR', self.ML, Inches(6.45), self.CW, Pt(10),
                 size=Pt(7), color=self.BLUE, bold=True)

        tier_data = [
            ('Digital Wealth', 'Premium', '→ Signature', 'Signature'),
            ('Digital Onboarding', '—', '★ Essential', 'Essential'),
            ('Digital Assist', 'Premium', 'Premium', 'Premium'),
            ('CLO', 'Essential', '→ Premium', 'Premium'),
            ('RM Workspace', '—', '★ Premium', '→ Signature'),
            ('Intelligence Fabric', '★ Agentic', '+ Foundation AI', 'Full AI + Custom'),
        ]
        ty = Inches(6.6)
        tw = Inches(2.0)
        cx_cols = [self.ML, self.ML + tw + Inches(0.1),
                   self.ML + 2*(tw + Inches(0.1)),
                   self.ML + 3*(tw + Inches(0.1))]
        # Headers
        for ci, hdr in enumerate(['Module', 'Year 1', 'Year 2', 'Year 3+']):
            col = [self.DARK_TEXT, self.BLUE, self.GREEN, self.PURPLE][ci]
            self._txt(s, hdr, cx_cols[ci], ty, tw, Pt(8),
                     size=Pt(6), color=col, bold=True)
        ty += Pt(10)
        for mod, y1, y2, y3 in tier_data:
            for ci, val in enumerate([mod, y1, y2, y3]):
                col = self.DARK_TEXT if ci == 0 else self.SUB_TEXT
                self._txt(s, val, cx_cols[ci], ty, tw, Pt(8),
                         size=Pt(5), color=col)
            ty += Pt(8)

    def _slide_06_phases(self):
        """Phased Delivery — 3-column merged phases."""
        s = self._new_slide()
        self._section_label(s, 'Delivery Roadmap')
        self._heading(s, 'Three Phases. ', 'One Destination.')
        self._subtitle(s, 'Phase 1 builds the foundation. Phase 2 activates AI and external clients. Phase 3 scales across jurisdictions.')

        cw = Inches(3.7)
        gap = Inches(0.3)
        ch = Inches(5.0)

        phases = [
            ('Phase 1', 'Year 1 — Foundation', self.BLUE,
             '~£2.3–3.9M', '£1.5M change', '~£11M value', '7 AI use cases',
             ['Internal release first',
              'Prospect + Client Portal (Web + Mobile)',
              'Employee Portal + basic advisor tools',
              'Digital Wealth, Onboarding, Assist, CLO',
              'Agentic Runtime — data foundation',
              'Grand Central: T24, Salesforce, Folio',
              'Temenos e-Services continues (external)']),
            ('Phase 2', 'Year 2 — Expand', self.GREEN,
             '~£2.6–5.1M', '£1.0M change', '~£31M cumul.', '+3 AI Agents',
             ['External client migration (go-live)',
              'Temenos e-Services deprecated',
              'RM Workspace Premium for advisors',
              'Full Intelligence Fabric activated',
              'Sentiment Analysis, Auto Reply, Doc Processing',
              'ComplyAdvantage + Onfido replace legacy',
              'Snowflake + ContentStack integrated']),
            ('Phase 3', 'Year 3+ — Scale', self.PURPLE,
             '~£3.6–6.9M', '~£0.6M change', '£148M cumul.', '7+ AI use cases',
             ['Full platform — all modules active',
              'Digital Advice, Sales (Lending, Advisory, Discretionary)',
              'Custom Agents, Smart Signals, Outreach',
              'Review Prep Copilot, Meeting Summarisation',
              'Multi-Agent Governance',
              'Aladdin Wealth Connector',
              'Geographic: UK → Guernsey → Switzerland → Singapore → US']),
        ]

        cx = self.ML
        for phase, subtitle, color, cost, change, value, ai, bullets in phases:
            self._card(s, cx, Inches(1.75), cw, ch, border=color)
            # Phase badge
            badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       cx + Inches(0.15), Inches(1.85),
                                       Inches(0.8), Inches(0.22))
            badge.fill.solid()
            badge.fill.fore_color.rgb = color
            badge.line.fill.background()
            self._txt(s, phase.upper(), cx + Inches(0.15), Inches(1.86),
                     Inches(0.8), Inches(0.2),
                     size=Pt(7), color=self.WHITE, bold=True, align=PP_ALIGN.CENTER)
            # Subtitle
            self._txt(s, subtitle, cx + Inches(1.05), Inches(1.87), cw - Inches(1.2), Pt(12),
                     size=Pt(9), color=color, bold=True)
            # Summary stats (2×2)
            sy = Inches(2.2)
            mini_w = Inches(1.6)
            for row in [(cost, change), (value, ai)]:
                mx = cx + Inches(0.15)
                for val in row:
                    self._txt(s, val, mx, sy, mini_w, Pt(12),
                             size=Pt(8), color=self.DARK_TEXT, bold=True)
                    mx += mini_w
                sy += Pt(14)
            # Divider
            div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     cx + Inches(0.15), Inches(2.7),
                                     cw - Inches(0.3), Pt(1))
            div.fill.solid()
            div.fill.fore_color.rgb = self.BORDER
            div.line.fill.background()
            # Bullets
            by = Inches(2.85)
            for b in bullets:
                self._txt(s, '→ ' + b, cx + Inches(0.15), by, cw - Inches(0.3), Pt(12),
                         size=Pt(8), color=self.SUB_TEXT)
                by += Inches(0.28)
            cx += cw + gap

    def _slide_07_investment(self):
        """Cost vs Value — static bar chart as a table representation."""
        s = self._new_slide()
        self._section_label(s, 'Where Value Meets Investment')
        self._heading(s, 'Year-on-Year: Cost, Value & ', 'AI Amplification.')
        self._subtitle(s, 'Enterprise model shown. Licence from sales pricing. Value from master analysis (validated).')

        # YoY data (Enterprise)
        yoy = [
            {'year': 'Year 1', 'license': 2.3, 'change': 1.5, 'base': 11.0, 'ai': 0},
            {'year': 'Year 2', 'license': 2.6, 'change': 1.0, 'base': 17.6, 'ai': 2.0},
            {'year': 'Year 3', 'license': 3.6, 'change': 0.6, 'base': 27.8, 'ai': 4.0},
            {'year': 'Year 4', 'license': 4.0, 'change': 0.6, 'base': 36.8, 'ai': 5.7},
            {'year': 'Year 5', 'license': 4.4, 'change': 0.6, 'base': 36.1, 'ai': 6.7},
        ]

        # Stacked bar chart using shapes
        chart_left = Inches(1.2)
        chart_top = Inches(2.0)
        chart_w = Inches(11.0)
        chart_h = Inches(3.8)
        max_val = 45
        bar_w = Inches(0.45)
        group_w = chart_w / 5
        colors = [self.RED, self.AMBER, self.GREEN, self.PURPLE]

        # Y-axis labels
        for val in [0, 10, 20, 30, 40]:
            y_pos = chart_top + chart_h - (val / max_val) * chart_h
            self._txt(s, f'£{val}M', Inches(0.3), y_pos - Pt(5), Inches(0.8), Pt(12),
                     size=Pt(7), color=self.MUTED, align=PP_ALIGN.RIGHT)
            # Grid line
            line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      chart_left, y_pos, chart_w, Pt(0.5))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
            line.line.fill.background()

        for i, d in enumerate(yoy):
            gx = chart_left + i * group_w + (group_w - 4 * bar_w - 3 * Inches(0.05)) / 2
            values = [d['license'], d['change'], d['base'], d['ai']]

            for j, (val, color) in enumerate(zip(values, colors)):
                if val == 0:
                    continue
                bar_x = gx + j * (bar_w + Inches(0.05))
                bar_h = (val / max_val) * chart_h
                bar_y = chart_top + chart_h - bar_h
                bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         bar_x, bar_y, bar_w, bar_h)
                bar.fill.solid()
                bar.fill.fore_color.rgb = color
                bar.line.fill.background()
                # Value label above bar
                self._txt(s, f'£{val:.1f}M', bar_x - Inches(0.05), bar_y - Pt(10),
                         bar_w + Inches(0.1), Pt(10),
                         size=Pt(6), color=color, bold=True, align=PP_ALIGN.CENTER)

            # X-axis label
            self._txt(s, d['year'], gx, chart_top + chart_h + Pt(4),
                     Inches(2), Pt(12),
                     size=Pt(9), color=self.DARK_TEXT, bold=True)

        # Legend
        lx = Inches(1.2)
        ly = Inches(6.2)
        for label, color in [('Licence', self.RED), ('Change', self.AMBER),
                              ('Base Value', self.GREEN), ('AI Value', self.PURPLE)]:
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, lx, ly + Pt(2), Pt(8), Pt(8))
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()
            self._txt(s, label, lx + Pt(12), ly, Inches(1), Pt(12),
                     size=Pt(8), color=self.SUB_TEXT)
            lx += Inches(1.5)

        # Summary stats
        self._card(s, Inches(7.5), Inches(6.1), Inches(5.2), Inches(0.5),
                  fill=self.BLUE_LIGHT, border=self.BLUE)
        for i, (val, label) in enumerate([('£148M', 'CUMULATIVE VALUE'),
                                           ('£21.2M', 'TOTAL INVESTMENT'),
                                           ('~7.0×', 'ROI')]):
            self._txt(s, val, Inches(7.7) + i * Inches(1.7), Inches(6.12), Inches(1.5), Pt(14),
                     size=Pt(12), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, label, Inches(7.7) + i * Inches(1.7), Inches(6.35), Inches(1.5), Pt(8),
                     size=Pt(6), color=self.SUB_TEXT, bold=True, align=PP_ALIGN.CENTER)

    def _slide_08_financials(self):
        """Five-Year Economics — summary table."""
        s = self._new_slide()
        self._section_label(s, 'Five-Year Economics')
        self._heading(s, 'The Numbers. ', 'Side by Side.')
        self._subtitle(s, 'Enterprise (recommended) vs Pay As You Grow. All figures in £M.')

        # Table
        from pptx.util import Inches, Pt
        rows = 7  # header + 5 years + total
        cols = 8
        tbl_shape = s.shapes.add_table(rows, cols,
                                        self.ML, Inches(1.9),
                                        self.CW, Inches(3.0))
        tbl = tbl_shape.table

        # Column widths
        col_widths = [Inches(0.8), Inches(1.2), Inches(1.7), Inches(1.7),
                      Inches(1.2), Inches(1.5), Inches(2.0), Inches(2.0)]
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w

        # Headers
        headers = ['Year', 'Phase', 'Ent. Licence', 'PAG Licence', 'Change',
                   'Ent. Total', 'Value (Base+AI)', 'Cumulative']
        for i, h in enumerate(headers):
            cell = tbl.cell(0, i)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.bold = True
                p.font.color.rgb = self.WHITE
                p.font.name = self.FONT
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.DARK_BG

        # Data
        data = [
            ['Y1', 'Foundation', '£2.3M', '£2.7M', '£1.5M', '£3.8M', '£11.0M', '£11M'],
            ['Y2', 'Expand', '£2.6M', '£3.5M', '£1.0M', '£3.6M', '£19.6M', '£31M'],
            ['Y3', 'Scale', '£3.6M', '£4.8M', '£0.6M', '£4.2M', '£31.8M', '£62M'],
            ['Y4', 'Optimise', '£4.0M', '£4.8M', '£0.6M', '£4.6M', '£42.5M', '£105M'],
            ['Y5', 'Compound', '£4.4M', '£4.8M', '£0.6M', '£5.0M', '£42.8M', '£148M'],
            ['Total', '', '£16.9M', '£20.6M', '£4.3M', '£21.2M', '£147.7M', '£148M'],
        ]
        for r, row_data in enumerate(data):
            for c, val in enumerate(row_data):
                cell = tbl.cell(r + 1, c)
                cell.text = val
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(8)
                    p.font.name = self.FONT
                    p.font.color.rgb = self.DARK_TEXT
                    if r == 5:  # Total row
                        p.font.bold = True
                    if c == 2:  # Ent licence
                        p.font.color.rgb = self.BLUE
                    elif c == 3:  # PAG licence
                        p.font.color.rgb = self.AMBER
                    elif c == 6 or c == 7:  # Value
                        p.font.color.rgb = self.GREEN
                if r == 5:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.BLUE_LIGHT

        # ROI comparison cards
        cards_data = [
            ('Enterprise', self.BLUE, self.BLUE_LIGHT,
             [('~7.0×', 'ROI'), ('Month 8', 'Breakeven'), ('£16.9M', '5yr Licence')]),
            ('Pay As You Grow', self.AMBER, self.AMBER_LIGHT,
             [('~5.9×', 'ROI'), ('Month 9', 'Breakeven'), ('£20.6M', '5yr Licence')]),
        ]
        cx = self.ML
        for name, color, bg, stats in cards_data:
            self._card(s, cx, Inches(5.2), Inches(5.8), Inches(1.6), fill=bg, border=color)
            self._txt(s, name.upper(), cx + Inches(0.2), Inches(5.25), Inches(3), Pt(10),
                     size=Pt(8), color=color, bold=True)
            sx = cx + Inches(0.2)
            for val, label in stats:
                self._txt(s, val, sx, Inches(5.55), Inches(1.6), Pt(18),
                         size=Pt(16), color=color, bold=True, align=PP_ALIGN.CENTER)
                self._txt(s, label, sx, Inches(5.9), Inches(1.6), Pt(10),
                         size=Pt(7), color=self.SUB_TEXT, bold=True, align=PP_ALIGN.CENTER)
                sx += Inches(1.8)
            cx += Inches(6.1)

        # Insight
        self._card(s, self.ML, Inches(6.95), self.CW, Inches(0.35),
                  fill=self.BLUE_LIGHT, border=self.BLUE)
        self._txt(s, 'Enterprise saves £3.7M over 5 years vs PAG. Scaling discount (60%→45%). Every price increase maps to a specific module activation.',
                  self.ML + Inches(0.3), Inches(6.97), self.CW - Inches(0.6), Inches(0.3),
                  size=Pt(8), color=self.DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)

    def _slide_09_recommendation(self):
        """Recommendation — 4 expectations mapped."""
        s = self._new_slide()
        self._section_label(s, 'Recommendation')
        self._heading(s, 'Your Expectations. ', 'Our Answers.')
        self._subtitle(s, 'Every recommendation maps to a specific platform capability and proven reference.')

        recs = [
            ('Exit analysis, mobilise immediately.',
             'Enterprise model. Phase 1 Month 1. POC validates T24/Salesforce connectivity in 1 week.',
             'Ila Bank — 6-month mobilisation'),
            ('No AI vendor lock-in.',
             'Multi-model architecture (Claude, GPT, Gemini). Agent Factory for custom agents. Switch LLMs without re-engineering.',
             'EQ Bank — multi-model AI'),
            ('Full integration, not point solutions.',
             'Grand Central: T24, Salesforce, Folio. 15+ connectors. Config not code. £2M integration cost avoided.',
             'NatWest — T24 Grand Central'),
            ('Single pane of glass for advisors.',
             'RM Workspace + Intelligence Fabric. Review Prep Copilot saves 225 advisors 2.5hrs/review — £2.5M/yr.',
             'RM Workspace — unified advisor experience'),
        ]

        ty = Inches(1.8)
        for i, (expectation, recommendation, reference) in enumerate(recs):
            row_h = Inches(1.15)
            # Expectation card (left)
            self._card(s, self.ML, ty, Inches(4.0), row_h, border=self.BORDER)
            self._txt(s, 'EXPECTATION', self.ML + Inches(0.15), ty + Inches(0.08),
                     Inches(3.5), Pt(8),
                     size=Pt(6), color=self.BLUE, bold=True)
            self._txt(s, f'"{expectation}"', self.ML + Inches(0.15), ty + Inches(0.28),
                     Inches(3.7), Inches(0.7),
                     size=Pt(9), color=self.DARK_TEXT, bold=True)

            # Arrow
            self._txt(s, '→', Inches(4.8), ty + Inches(0.35), Inches(0.3), Inches(0.3),
                     size=Pt(18), color=self.BLUE, bold=True, align=PP_ALIGN.CENTER)

            # Recommendation card (right)
            self._card(s, Inches(5.2), ty, Inches(7.2), row_h,
                      fill=self.BLUE_LIGHT, border=self.BLUE)
            self._txt(s, 'RECOMMENDATION', Inches(5.35), ty + Inches(0.08),
                     Inches(6.5), Pt(8),
                     size=Pt(6), color=self.BLUE, bold=True)
            self._txt(s, recommendation, Inches(5.35), ty + Inches(0.28),
                     Inches(6.8), Inches(0.45),
                     size=Pt(8), color=self.DARK_TEXT)
            self._txt(s, f'Ref: {reference}', Inches(5.35), ty + Inches(0.78),
                     Inches(6.5), Pt(10),
                     size=Pt(7), color=self.MUTED)

            ty += row_h + Inches(0.15)

        # Bottom card
        self._card(s, self.ML, Inches(6.2), self.CW, Inches(0.9),
                  fill=self.LIGHT_CARD, border=self.BLUE)
        self._txt(s, 'COST PREDICTABILITY', self.ML + Inches(0.15), Inches(6.28),
                 Inches(3), Pt(8), size=Pt(6), color=self.BLUE, bold=True)
        self._txt(s, 'Scaling discount (60%→45% of list). Capped BIC for AI. Every price increase maps to a specific module. £16.9M over 5 years — predictable, auditable.',
                  self.ML + Inches(0.15), Inches(6.48), self.CW - Inches(0.3), Inches(0.5),
                  size=Pt(9), color=self.SUB_TEXT)

    def _slide_10_close(self):
        """Close Plan + CTA — dark hero."""
        s = self._new_slide(dark=True)
        self._section_label(s, 'Mutual Close Plan', top=Inches(0.4), dark=True)
        self._heading(s, 'Target: Signed by ', 'End of April 2026.',
                     top=Inches(0.75), dark=True)

        # Timeline
        milestones = [
            ('March 9', 'COO Meeting', 'Platform Value & AI Use Cases', self.BLUE, True),
            ('Mid-March', 'Scope & Confidence', 'Product Catalogue refinement', self.WHITE, False),
            ('Late March', 'Exclusivity', 'Formal procurement starts', self.PURPLE, False),
            ('Early April', '1-Week POC', 'T24/Salesforce connectivity proof', self.GREEN, False),
            ('Mid-April', 'Procurement', 'Due Diligence · Infosec · MSA Review', self.WHITE, False),
            ('End of April', 'Signature', 'Contract signed · Implementation begins', self.GREEN, False),
        ]

        mx = self.ML
        mw = Inches(1.8)
        my = Inches(1.5)

        for date, title, sub, color, is_current in milestones:
            self._card(s, mx, my, mw, Inches(1.3),
                      fill=self.CARD_BG, border=color, dark=True)
            # Dot indicator
            if is_current:
                dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                          mx + mw/2 - Pt(4), my - Pt(8),
                                          Pt(8), Pt(8))
                dot.fill.solid()
                dot.fill.fore_color.rgb = self.BLUE
                dot.line.fill.background()

            self._txt(s, date, mx + Inches(0.1), my + Inches(0.08), mw - Inches(0.2), Pt(10),
                     size=Pt(8), color=color, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, title, mx + Inches(0.1), my + Inches(0.35), mw - Inches(0.2), Pt(14),
                     size=Pt(10), color=self.WHITE, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, sub, mx + Inches(0.1), my + Inches(0.65), mw - Inches(0.2), Inches(0.5),
                     size=Pt(7), color=self.MUTED, align=PP_ALIGN.CENTER)
            mx += mw + Inches(0.15)

        # Connecting line
        line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   self.ML + mw/2, my - Pt(5),
                                   Inches(11.0), Pt(2))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x1A, 0x3A, 0x5E)
        line.line.fill.background()

        # Hero CTA
        tb = s.shapes.add_textbox(self.ML, Inches(3.3), self.CW, Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        for text, color in [('Build the Foundation.\n', self.WHITE),
                            ('Accelerate the Advantage.', self.BLUE)]:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(36)
            r.font.bold = True
            r.font.color.rgb = color
            r.font.name = self.FONT

        # 4 hero stats
        stat_w = Inches(2.6)
        gap = Inches(0.3)
        sx = self.ML + Inches(0.3)
        for val, label, vcol in [
            ('~£2.3M', 'Y1 Entry (Enterprise)', self.BLUE),
            ('£16.9M', '5-Year Licence', self.WHITE),
            ('£148M', 'Value at Stake', self.GREEN),
            ('~7.0×', 'ROI', self.BLUE),
        ]:
            self._card(s, sx, Inches(5.0), stat_w, Inches(1.0),
                      fill=self.CARD_BG, dark=True)
            self._txt(s, val, sx, Inches(5.05), stat_w, Inches(0.45),
                     size=Pt(22), color=vcol, bold=True, align=PP_ALIGN.CENTER)
            self._txt(s, label.upper(), sx, Inches(5.5), stat_w, Pt(12),
                     size=Pt(7), color=self.MUTED, bold=True, align=PP_ALIGN.CENTER)
            sx += stat_w + gap

        # Footer
        self._txt(s, 'One platform. Every channel. Every data source. From relationship depth to digital scale.',
                  self.ML, Inches(6.3), self.CW, Inches(0.3),
                  size=Pt(11), color=self.MUTED, align=PP_ALIGN.CENTER)
        self._txt(s, 'Schroders × Backbase · March 2026',
                  self.ML, Inches(6.7), self.CW, Pt(12),
                  size=Pt(8), color=RGBColor(0x3A, 0x49, 0x5D), align=PP_ALIGN.CENTER)

    # ── GENERATE ──────────────────────────────────────────

    def generate(self, output_path):
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_W
        self.prs.slide_height = self.SLIDE_H

        self._slide_01_cover()          # 1  — Dark hero
        self._slide_02_context()        # 2  — Context
        self._slide_03_platform()       # 3  — What You Get: Platform
        self._slide_04_ai()             # 4  — What You Get: AI
        self._slide_05_commercial()     # 5  — Enterprise vs PAG
        self._slide_06_phases()         # 6  — Phased Delivery
        self._slide_07_investment()     # 7  — Cost vs Value (chart)
        self._slide_08_financials()     # 8  — Five-Year Economics
        self._slide_09_recommendation() # 9  — Recommendation
        self._slide_10_close()          # 10 — Close Plan + CTA

        self.prs.save(output_path)
        size_kb = Path(output_path).stat().st_size // 1024
        print(f'✓ Saved: {output_path} ({size_kb} KB, 10 slides)')


if __name__ == '__main__':
    OUTPUT_DIR = Path(__file__).parent.parent / 'Engagement/Schroders Group/Output'
    output_file = OUTPUT_DIR / 'Schroders_Executive_Readout_v3.pptx'
    SchrordersV3PptxGenerator().generate(str(output_file))
